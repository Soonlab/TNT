"""
Fig 7C — Convergence test heatmap (9 baseline × 4 cascade Δ)

Native PowerPoint object via python-pptx. Each heatmap cell is an individual
rectangle shape; cell colors come from a diverging RdBu_r colormap clipped at
[-0.5, +0.5] so that the null pattern (0/36 P<0.05) is visually obvious as a
near-white grid. All textboxes are TEXT_BOX (double-click editable in PPT).
Arial throughout; shadow effects stripped from every shape.

Literature motif references:
  - Mariathasan et al. Nature 2018 (PMID 29443960) — diverging blue-white-red
    correlation grid in extended-data figures of the TGFβ atlas.
  - Thorsson et al. Immunity 2018 (PMID 29628290) — pan-cancer signature
    correlation panels with light-grey cell borders + RdBu diverging palette.
  - Rooney, Shukla, Wu et al. Cell 2015 (PMID 25594174) — immune correlation
    matrix with cell-by-cell r annotations.

Project figure rules applied:
  Rule 0 — literature search performed (see README).
  Rule 1 — one panel per slide (single slide deck).
  Rule 2 — no figure title text.
  Rule 3 — every visual element is a native python-pptx shape/textbox.
  Rule 4 — applies to main + supplementary figures.
  Rule 5 — Arial font everywhere.
  Rule 6 — shadow effects removed via XML.
"""

from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# --------------------------------------------------------------------------
# Paths
# --------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parent.parent
TSV_IN = ROOT / "tables" / "convergence_36pair_used.tsv"
PPTX_OUT = ROOT / "figures" / "Fig7C_convergence_heatmap.pptx"

# --------------------------------------------------------------------------
# Feature lists (column names → display labels)
# Order taken from the existing 36-pair convergence test
# (DSB→DNA Repair→HRR→E2F→G2-M→MycV2→MHC-II→MSI→amp; cascades CD8→IGH→MHC-II→Treg)
# --------------------------------------------------------------------------
BASELINES = [
    ("DNA Double-Strand Break Repair R-HSA-5693532", "DSB repair (Reactome)"),
    ("DNA Repair R-HSA-73894",                       "DNA repair (Reactome)"),
    ("HDR Thru Homologous Recombination (HRR) R-HSA-5685942",
                                                     "HRR (Reactome)"),
    ("E2F Targets",                                  "E2F targets (Hallmark)"),
    ("G2-M Checkpoint",                              "G2-M checkpoint (Hallmark)"),
    ("Myc Targets V2",                               "Myc targets V2 (Hallmark)"),
    ("MHC_II",                                       "MHC-II (baseline)"),
    ("MSI_pct",                                      "MSI (%)"),
    ("frac_amp",                                     "Genomic amp fraction"),
]
CASCADES = [
    ("CD8_cytotoxic_delta", "Δ CD8-cytotoxic"),
    ("IGH_n_delta",         "Δ IGH clonotypes"),
    ("MHC_II_delta",        "Δ MHC-II"),
    ("Treg_delta",          "Δ Treg"),
]
HEADLINE_PAIR = ("DNA Double-Strand Break Repair R-HSA-5693532", "CD8_cytotoxic_delta")

# Palette anchors (RdBu_r endpoints — Brewer-style)
BLUE_DEEP = (33, 102, 172)
RED_DEEP  = (178,  24,  43)
WHITE     = (255, 255, 255)
GOLD      = RGBColor(218, 165, 32)
GREY_BORD = RGBColor(200, 200, 200)
BLACK     = RGBColor(0, 0, 0)

# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------
def diverging(r: float, vmin: float = -0.5, vmax: float = 0.5) -> RGBColor:
    """Two-arm diverging ramp (deep-blue ↔ white ↔ deep-red)."""
    r = max(min(r, vmax), vmin)
    if r >= 0:
        t = r / vmax
        R = int(WHITE[0] - t * (WHITE[0] - RED_DEEP[0]))
        G = int(WHITE[1] - t * (WHITE[1] - RED_DEEP[1]))
        B = int(WHITE[2] - t * (WHITE[2] - RED_DEEP[2]))
    else:
        t = -r / (-vmin)
        R = int(WHITE[0] - t * (WHITE[0] - BLUE_DEEP[0]))
        G = int(WHITE[1] - t * (WHITE[1] - BLUE_DEEP[1]))
        B = int(WHITE[2] - t * (WHITE[2] - BLUE_DEEP[2]))
    return RGBColor(R, G, B)

def kill_shadow(shape) -> None:
    """Strip every effect (shadow/glow) from the shape via empty effectLst."""
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    etree.SubElement(spPr, qn("a:effectLst"))

def set_text(
    tb,
    text: str,
    size: float = 8,
    bold: bool = False,
    color: RGBColor = BLACK,
    align=PP_ALIGN.CENTER,
    anchor=MSO_ANCHOR.MIDDLE,
) -> None:
    """Set textbox text with Arial styling. Idempotent across runs."""
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    kill_shadow(tb)

def fmt_r(r: float) -> str:
    """Format r as ±0.XX with Unicode minus."""
    s = f"{r:+.2f}"
    return s.replace("-", "−")

# --------------------------------------------------------------------------
# Load data
# --------------------------------------------------------------------------
df = pd.read_csv(TSV_IN, sep="\t")
rmap = {(row.baseline_feature, row.cascade_feature): row for row in df.itertuples()}

# Sanity: report headline pair values
hp = rmap[HEADLINE_PAIR]
print(
    f"[sanity] headline DSB×ΔCD8cyt: "
    f"plain r={hp.spearman_r:+.3f} P={hp.spearman_p:.3f}  |  "
    f"partial r={hp.partial_r:+.3f} P={hp.partial_P:.3f}"
)

# --------------------------------------------------------------------------
# Build slide
# --------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Grid geometry
GRID_LEFT = Inches(3.30)
GRID_TOP  = Inches(1.70)
CELL_W    = Inches(0.92)
CELL_H    = Inches(0.50)
N_ROWS    = len(BASELINES)
N_COLS    = len(CASCADES)
GRID_W    = CELL_W * N_COLS
GRID_H    = CELL_H * N_ROWS

# --- Column axis title ---
col_title = slide.shapes.add_textbox(
    GRID_LEFT, GRID_TOP - Inches(1.05),
    GRID_W, Inches(0.28),
)
set_text(
    col_title,
    "Cascade Δ features (radiation-phase dynamics)",
    size=10, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
)

# --- Column labels (top, one per cascade) ---
for j, (_, lbl) in enumerate(CASCADES):
    tb = slide.shapes.add_textbox(
        GRID_LEFT + CELL_W * j, GRID_TOP - Inches(0.72),
        CELL_W, Inches(0.62),
    )
    set_text(tb, lbl, size=9, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

# --- Row axis title (above row labels, right-aligned to align with labels) ---
row_title = slide.shapes.add_textbox(
    GRID_LEFT - Inches(2.20), GRID_TOP - Inches(0.34),
    Inches(2.10), Inches(0.26),
)
set_text(
    row_title,
    "Baseline tumor-intrinsic features",
    size=10, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM,
)

# --- Row labels (left of grid, one per baseline) ---
for i, (_, lbl) in enumerate(BASELINES):
    tb = slide.shapes.add_textbox(
        GRID_LEFT - Inches(2.20), GRID_TOP + CELL_H * i,
        Inches(2.10), CELL_H,
    )
    set_text(tb, lbl, size=9, bold=False, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# --- Heatmap cells (rectangle + value textbox) ---
for i, (b_col, _) in enumerate(BASELINES):
    for j, (c_col, _) in enumerate(CASCADES):
        row = rmap.get((b_col, c_col))
        if row is None:
            continue
        r_val = row.partial_r
        P_val = row.partial_P
        q_val = row.BH_q

        x = GRID_LEFT + CELL_W * j
        y = GRID_TOP  + CELL_H * i

        # cell rectangle
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, CELL_W, CELL_H)
        rect.fill.solid()
        rect.fill.fore_color.rgb = diverging(r_val, -0.5, 0.5)

        if (b_col, c_col) == HEADLINE_PAIR:
            rect.line.color.rgb = GOLD
            rect.line.width = Pt(2.5)
        else:
            rect.line.color.rgb = GREY_BORD
            rect.line.width = Pt(0.5)
        kill_shadow(rect)

        # significance suffix
        suffix = ""
        if q_val < 0.05:
            suffix = "**"
        elif P_val < 0.05:
            suffix = "*"

        text_color = RGBColor(255, 255, 255) if abs(r_val) >= 0.40 else BLACK
        ann = slide.shapes.add_textbox(x, y, CELL_W, CELL_H)
        set_text(
            ann, fmt_r(r_val) + suffix,
            size=8, bold=False, color=text_color,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

# --------------------------------------------------------------------------
# Colorbar (right of grid)
# --------------------------------------------------------------------------
CB_LEFT = GRID_LEFT + GRID_W + Inches(0.55)
CB_TOP  = GRID_TOP
CB_W    = Inches(0.26)
CB_H    = GRID_H

# colorbar header
cblab = slide.shapes.add_textbox(
    CB_LEFT - Inches(0.25), CB_TOP - Inches(0.34),
    Inches(1.6), Inches(0.26),
)
set_text(
    cblab, "Partial Spearman r",
    size=9, bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
)

# 100 stacked rectangles for smooth ramp
N_STOPS = 100
stop_h = CB_H / N_STOPS
for k in range(N_STOPS):
    # k=0 top -> r = +0.5, k=N_STOPS-1 bottom -> r = -0.5
    frac = k / (N_STOPS - 1)
    r_at = 0.5 - frac * 1.0
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        CB_LEFT, CB_TOP + stop_h * k,
        CB_W,
        stop_h + Emu(900),  # 1-EMU overlap to avoid hairline gaps in PPT renderers
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = diverging(r_at, -0.5, 0.5)
    rect.line.fill.background()  # no outline on stops
    kill_shadow(rect)

# colorbar outer border
border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, CB_LEFT, CB_TOP, CB_W, CB_H)
border.fill.background()
border.line.color.rgb = RGBColor(64, 64, 64)
border.line.width = Pt(0.75)
kill_shadow(border)

# colorbar tick labels (and short ticks)
for r_val in [0.50, 0.25, 0.0, -0.25, -0.50]:
    frac = (0.5 - r_val) / 1.0
    ty = CB_TOP + CB_H * frac
    # tick mark
    tick = slide.shapes.add_connector(
        1,  # straight connector
        CB_LEFT + CB_W, ty,
        CB_LEFT + CB_W + Inches(0.08), ty,
    )
    tick.line.color.rgb = RGBColor(64, 64, 64)
    tick.line.width = Pt(0.5)
    kill_shadow(tick)
    # numeric label
    txt = "0" if r_val == 0 else fmt_r(r_val)
    tb = slide.shapes.add_textbox(
        CB_LEFT + CB_W + Inches(0.12), ty - Inches(0.10),
        Inches(0.55), Inches(0.20),
    )
    set_text(tb, txt, size=8, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# --------------------------------------------------------------------------
# Bottom legend / annotations
# --------------------------------------------------------------------------
LEG_LEFT = GRID_LEFT - Inches(2.20)
LEG_TOP  = GRID_TOP + GRID_H + Inches(0.30)
LEG_W    = Inches(8.30)

leg1 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP, LEG_W, Inches(0.22))
set_text(
    leg1,
    "n = 11–12 paired subjects · partial Spearman, response-group adjusted · "
    "BH-corrected across 36 pairs",
    size=8, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
)

leg2 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP + Inches(0.24), LEG_W, Inches(0.22))
set_text(
    leg2,
    "Cell value = partial Spearman r.   "
    "* nominal P < 0.05    ** BH q < 0.05    "
    "Gold border: headline pair (DSB repair × Δ CD8-cytotoxic).",
    size=8, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
)

leg3 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP + Inches(0.48), LEG_W, Inches(0.22))
set_text(
    leg3,
    "Convergence result: 0/36 partial P < 0.05    ·    0/36 BH q < 0.05    "
    "(static and dynamic layers observationally independent).",
    size=8, bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
)

# --------------------------------------------------------------------------
# Save
# --------------------------------------------------------------------------
PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(PPTX_OUT)
print(f"[ok] wrote {PPTX_OUT}")
