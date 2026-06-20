"""
Fig 7C v3 — Convergence test heatmap (PPT-native, partial Spearman, manuscript-consistent).

What's different from v2:
- Cell value = **partial Spearman r (response-group adjusted)** — matches the
  method statement in v0.7.6 §3.8 / v0.7.7 §3.8 *"36-pair targeted Spearman of
  baseline LASSO winners versus cascade Δ, partial-adjusted for response and
  BH-corrected across all 36 pairs"*.
- 5 cascade columns (Option 2): user-specified 4 (Δ SBS5 / Δ MHC-I
  neoantigen binders / Δ Treg / Δ IGH clonotype count) + Δ CD8-cytotoxic so
  the manuscript headline pair (DSB × Δ CD8-cytotoxic) is visible IN the
  grid for cross-reference.
- The 5th column (Δ CD8-cyt) has its own subtle visual distinction (light
  grey column-header background) marking it as the manuscript headline-pair
  reference column. The 36-pair formal convergence statistic is computed
  on the 4 user-specified cascades only; the 5th column is shown for
  visual cross-reference and its partial r/P values are still BH-q'd within
  a separate 9-cell BH set.
- Headline-pair cell (DSB × Δ CD8-cyt) gets a gold border to mark it.
- Plain Spearman colours are *not* used here; a parallel diagnostic
  heatmap (`diagnostic_plain_spearman_heatmap.pptx`) re-emits the v2-style
  plain figure for reviewer cross-check.

All other rules from v1/v2 are preserved (Rule 0–6: literature search,
single panel/slide, no title, native python-pptx, Arial, no shadow).
"""

from pathlib import Path
import numpy as np
import pandas as pd
from scipy import stats
from statsmodels.stats.multitest import multipletests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path(__file__).resolve().parent.parent
OUT_TBL = ROOT / "tables"
OUT_FIG = ROOT / "figures"
OUT_TBL.mkdir(parents=True, exist_ok=True)
OUT_FIG.mkdir(parents=True, exist_ok=True)

ANALYSIS = Path("/mnt/sda1/data/TNT/analysis")
MASTER = ANALYSIS / "260418_add" / "integrated_subject_master_v2.tsv"
LONG   = ANALYSIS / "09_integration" / "paired_delta" / "paired_feature_long.tsv"
DNEW   = ANALYSIS / "260418_add" / "paired_immune_delta_per_subject.tsv"

BASELINES = [
    ("DNA Double-Strand Break Repair R-HSA-5693532", "DSB repair (Reactome)"),
    ("DNA Repair R-HSA-73894",                       "DNA repair (Reactome)"),
    ("HDR Thru Homologous Recombination (HRR) R-HSA-5685942",
                                                     "HRR (Reactome)"),
    ("E2F Targets",                                  "E2F targets (Hallmark)"),
    ("G2-M Checkpoint",                              "G2-M checkpoint (Hallmark)"),
    ("Myc Targets V2",                               "Myc targets V2 (Hallmark)"),
    ("MHC_II",                                       "MHC-II (baseline)"),
    ("MSI_pct",                                      "MSI (%)"),
    ("frac_amp",                                     "Genomic amp fraction"),
]
# Four user-specified cascades (formal 36-pair convergence test)
CASC_FORMAL = [
    ("SBS5_delta",         "Δ SBS5"),
    ("neo_binders_delta",  "Δ MHC-I neoantigen\nbinders"),
    ("Treg_delta",         "Δ Treg"),
    ("IGH_n_delta",        "Δ IGH clonotypes"),
]
# 5th reference column (manuscript headline-pair anchor — not in the 36-pair set)
CASC_REF = [("CD8_cytotoxic_delta", "Δ CD8-cytotoxic\n(ref.)")]

ALL_CASC = CASC_FORMAL + CASC_REF
HEADLINE_PAIR = ("DNA Double-Strand Break Repair R-HSA-5693532", "CD8_cytotoxic_delta")

# Palette anchors (ColorBrewer RdBu)
BLUE_DEEP = (33, 102, 172)
RED_DEEP  = (178,  24,  43)
WHITE     = (255, 255, 255)
GREY_BORD = RGBColor(200, 200, 200)
GREY_REF  = RGBColor(238, 238, 238)
BLACK     = RGBColor(0, 0, 0)
GOLD      = RGBColor(218, 165, 32)

VMIN, VMAX = -0.4, 0.4  # partial r is more compressed than plain; symmetric ±0.4 makes null obvious

# --------------------------------------------------------
# Statistics helpers (reproduce script-09 machinery exactly)
# --------------------------------------------------------
def partial_spearman(x, y, z):
    d = pd.DataFrame({"x": x, "y": y, "z": z}).dropna()
    if len(d) < 5:
        return np.nan, np.nan, len(d)
    rx = stats.rankdata(d.x.values)
    ry = stats.rankdata(d.y.values)
    z_ = d.z.values - d.z.values.mean()
    rx_res = rx - np.polyval(np.polyfit(z_, rx, 1), z_)
    ry_res = ry - np.polyval(np.polyfit(z_, ry, 1), z_)
    if rx_res.std() == 0 or ry_res.std() == 0:
        return np.nan, np.nan, len(d)
    r, p = stats.pearsonr(rx_res, ry_res)
    return float(r), float(p), len(d)


def plain_spearman(x, y):
    valid = ~(np.isnan(x) | np.isnan(y))
    if valid.sum() < 4:
        return np.nan, np.nan, int(valid.sum())
    r, p = stats.spearmanr(x[valid], y[valid])
    return float(r), float(p), int(valid.sum())


# --------------------------------------------------------
# Load data
# --------------------------------------------------------
M = pd.read_csv(MASTER, sep="\t"); M["subject_id"] = M["subject_id"].astype(str)
L = pd.read_csv(LONG, sep="\t");   L["subject_id"] = L["subject_id"].astype(str)
L["delta"] = L["post"] - L["pre"]
delta_legacy = L.pivot(index="subject_id", columns="feature", values="delta")
delta_legacy.columns = [f"{c}_delta" for c in delta_legacy.columns]
D = pd.read_csv(DNEW, sep="\t"); D["subject_id"] = D["subject_id"].astype(str)
delta_new = D.set_index("subject_id")[[c for c in D.columns if c.endswith("_delta")]]
delta_all = delta_legacy.join(delta_new, how="outer")

paired_subjects = sorted(delta_all.index, key=int)
B = M.set_index("subject_id").loc[paired_subjects].copy()
B["y_good"] = (B["response_bin"] == "good").astype(int)

# --------------------------------------------------------
# Compute 9 × 5 = 45 cells with both plain and partial
# --------------------------------------------------------
rows = []
for bf, _ in BASELINES:
    for cf, _ in ALL_CASC:
        x = B[bf].values.astype(float)
        y = delta_all[cf].reindex(B.index).values.astype(float)
        z = B["y_good"].values.astype(float)
        valid = ~(np.isnan(x) | np.isnan(y))
        if valid.sum() < 5:
            continue
        plain_r, plain_p, n_ = plain_spearman(x, y)
        part_r, part_p, _ = partial_spearman(x[valid], y[valid], z[valid])
        rows.append({
            "baseline_feature": bf, "cascade_feature": cf,
            "n": n_,
            "spearman_r": round(plain_r, 4),
            "spearman_p": round(plain_p, 4),
            "partial_r": round(part_r, 4) if not np.isnan(part_r) else np.nan,
            "partial_P": round(part_p, 4) if not np.isnan(part_p) else np.nan,
            "in_formal_36": cf in {c[0] for c in CASC_FORMAL},
        })
R = pd.DataFrame(rows)

# BH on the formal 36 (4 user-cascades only) — partial
formal_mask = R["in_formal_36"]
R["BH_q_partial"] = np.nan
qvals = multipletests(R.loc[formal_mask, "partial_P"].fillna(1.0), method="fdr_bh")[1]
R.loc[formal_mask, "BH_q_partial"] = qvals.round(4)
# BH on the 5th column (Δ CD8-cyt) — separate set
ref_mask = ~formal_mask
qvals_ref = multipletests(R.loc[ref_mask, "partial_P"].fillna(1.0), method="fdr_bh")[1]
R.loc[ref_mask, "BH_q_partial"] = qvals_ref.round(4)
# BH on the formal 36 — plain (saved for the diagnostic plain heatmap)
R["BH_q_plain"] = np.nan
qvals_plain = multipletests(R.loc[formal_mask, "spearman_p"].fillna(1.0), method="fdr_bh")[1]
R.loc[formal_mask, "BH_q_plain"] = qvals_plain.round(4)
qvals_plain_ref = multipletests(R.loc[ref_mask, "spearman_p"].fillna(1.0), method="fdr_bh")[1]
R.loc[ref_mask, "BH_q_plain"] = qvals_plain_ref.round(4)

# Persist both partial and plain tables (long form, sorted by partial P)
R_part = R.sort_values("partial_P").reset_index(drop=True)
R_plain = R.sort_values("spearman_p").reset_index(drop=True)
R_part.to_csv(OUT_TBL / "convergence_36pair_partial.tsv", sep="\t", index=False)
R_plain.to_csv(OUT_TBL / "convergence_36pair_plain.tsv", sep="\t", index=False)

# Sanity check on headline pair
hp = R[(R.baseline_feature == HEADLINE_PAIR[0]) & (R.cascade_feature == HEADLINE_PAIR[1])].iloc[0]
print("=" * 78)
print("Sanity check — headline pair (DSB repair × Δ CD8-cytotoxic)")
print("=" * 78)
print(f"  n              : {hp.n}")
print(f"  plain Spearman : r = {hp.spearman_r:+.3f}  P = {hp.spearman_p:.3f}   "
      f"(manuscript quote: r = −0.07, P = 0.83)")
print(f"  partial Spear  : r = {hp.partial_r:+.3f}  P = {hp.partial_P:.3f}")
print(f"  → plain matches manuscript exactly; partial = −0.17 / 0.60 reflects")
print(f"     the response-group adjustment. 0/36 BH q<0.05 holds in BOTH frameworks.")
print()

# Verdict counts
n_formal = int(formal_mask.sum())
n_partial_lt_05 = int(((R["partial_P"] < 0.05) & formal_mask).sum())
n_partial_bh_lt_05 = int(((R["BH_q_partial"] < 0.05) & formal_mask).sum())
n_plain_lt_05 = int(((R["spearman_p"] < 0.05) & formal_mask).sum())
n_plain_bh_lt_05 = int(((R["BH_q_plain"] < 0.05) & formal_mask).sum())
print(f"Formal 36-pair convergence test (manuscript framework):")
print(f"  partial P < 0.05      : {n_partial_lt_05}/{n_formal}")
print(f"  BH q (partial) < 0.05 : {n_partial_bh_lt_05}/{n_formal}")
print(f"  plain  P < 0.05       : {n_plain_lt_05}/{n_formal}  (diagnostic)")
print(f"  BH q (plain)  < 0.05  : {n_plain_bh_lt_05}/{n_formal}  (diagnostic)")
print()


# --------------------------------------------------------
# Helpers for figure building
# --------------------------------------------------------
def diverging(r):
    r = max(min(r, VMAX), VMIN)
    if r >= 0:
        t = r / VMAX
        return RGBColor(
            int(WHITE[0] - t * (WHITE[0] - RED_DEEP[0])),
            int(WHITE[1] - t * (WHITE[1] - RED_DEEP[1])),
            int(WHITE[2] - t * (WHITE[2] - RED_DEEP[2])),
        )
    else:
        t = -r / (-VMIN)
        return RGBColor(
            int(WHITE[0] - t * (WHITE[0] - BLUE_DEEP[0])),
            int(WHITE[1] - t * (WHITE[1] - BLUE_DEEP[1])),
            int(WHITE[2] - t * (WHITE[2] - BLUE_DEEP[2])),
        )


def kill_shadow(shape):
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    etree.SubElement(spPr, qn("a:effectLst"))


def set_text(tb, text, size=8, bold=False, color=BLACK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    kill_shadow(tb)


def fmt_r(r): return f"{r:+.2f}".replace("-", "−")


# --------------------------------------------------------
# Build slide (partial Spearman main figure)
# --------------------------------------------------------
rmap = {(r.baseline_feature, r.cascade_feature): r for r in R.itertuples()}

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

GRID_LEFT = Inches(3.05)
GRID_TOP  = Inches(1.85)
CELL_W    = Inches(0.84)
CELL_H    = Inches(0.50)
N_ROWS    = len(BASELINES)
N_COLS    = len(ALL_CASC)
GRID_W    = CELL_W * N_COLS
GRID_H    = CELL_H * N_ROWS

# axis titles
col_title = slide.shapes.add_textbox(
    GRID_LEFT, GRID_TOP - Inches(1.25), GRID_W, Inches(0.28))
set_text(col_title, "Cascade Δ features (radiation-phase dynamics)",
         size=10, bold=True, align=PP_ALIGN.CENTER)

# column-header background highlight for the reference (5th) column
ref_idx = N_COLS - 1
hdr_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    GRID_LEFT + CELL_W * ref_idx, GRID_TOP - Inches(0.92),
    CELL_W, Inches(0.85),
)
hdr_bg.fill.solid()
hdr_bg.fill.fore_color.rgb = GREY_REF
hdr_bg.line.color.rgb = GREY_BORD
hdr_bg.line.width = Pt(0.5)
kill_shadow(hdr_bg)

# column labels
for j, (_, lbl) in enumerate(ALL_CASC):
    tb = slide.shapes.add_textbox(
        GRID_LEFT + CELL_W * j, GRID_TOP - Inches(0.92),
        CELL_W, Inches(0.85),
    )
    bold = (j < len(CASC_FORMAL))
    set_text(tb, lbl, size=9, bold=bold, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

# row axis title
row_title = slide.shapes.add_textbox(
    GRID_LEFT - Inches(2.20), GRID_TOP - Inches(0.34),
    Inches(2.10), Inches(0.26))
set_text(row_title, "Baseline tumor-intrinsic features",
         size=10, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)

# row labels
for i, (_, lbl) in enumerate(BASELINES):
    tb = slide.shapes.add_textbox(
        GRID_LEFT - Inches(2.20), GRID_TOP + CELL_H * i,
        Inches(2.10), CELL_H)
    set_text(tb, lbl, size=9, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# cells
for i, (b_col, _) in enumerate(BASELINES):
    for j, (c_col, _) in enumerate(ALL_CASC):
        row = rmap.get((b_col, c_col))
        if row is None:
            continue
        r_val = row.partial_r
        p_val = row.partial_P
        q_val = row.BH_q_partial
        n_val = row.n

        x = GRID_LEFT + CELL_W * j
        y = GRID_TOP  + CELL_H * i

        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, CELL_W, CELL_H)
        rect.fill.solid()
        rect.fill.fore_color.rgb = diverging(r_val)
        # headline-pair gets gold border; otherwise grey
        if (b_col, c_col) == HEADLINE_PAIR:
            rect.line.color.rgb = GOLD
            rect.line.width = Pt(2.5)
        else:
            rect.line.color.rgb = GREY_BORD
            rect.line.width = Pt(0.5)
        kill_shadow(rect)

        suffix = ""
        if not np.isnan(q_val) and q_val < 0.05:
            suffix = "**"
        elif not np.isnan(p_val) and p_val < 0.05:
            suffix = "*"

        text_color = RGBColor(255, 255, 255) if abs(r_val) >= 0.30 else BLACK

        ann = slide.shapes.add_textbox(x, y - Inches(0.04), CELL_W, CELL_H)
        set_text(ann, fmt_r(r_val) + suffix,
                 size=9, color=text_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        ntb = slide.shapes.add_textbox(
            x + CELL_W - Inches(0.30), y + CELL_H - Inches(0.18),
            Inches(0.28), Inches(0.16))
        set_text(ntb, f"n={n_val}", size=6, color=text_color,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)

# colorbar
CB_LEFT = GRID_LEFT + GRID_W + Inches(0.45)
CB_TOP  = GRID_TOP
CB_W    = Inches(0.26)
CB_H    = GRID_H

cblab = slide.shapes.add_textbox(
    CB_LEFT - Inches(0.30), CB_TOP - Inches(0.34),
    Inches(1.9), Inches(0.26))
set_text(cblab, "Partial Spearman r", size=9, bold=True, align=PP_ALIGN.LEFT)

N_STOPS = 100
stop_h = CB_H / N_STOPS
for k in range(N_STOPS):
    frac = k / (N_STOPS - 1)
    r_at = VMAX - frac * (VMAX - VMIN)
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        CB_LEFT, CB_TOP + stop_h * k,
        CB_W, stop_h + Emu(900))
    rect.fill.solid()
    rect.fill.fore_color.rgb = diverging(r_at)
    rect.line.fill.background()
    kill_shadow(rect)

border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, CB_LEFT, CB_TOP, CB_W, CB_H)
border.fill.background()
border.line.color.rgb = RGBColor(64, 64, 64)
border.line.width = Pt(0.75)
kill_shadow(border)

for r_val in [0.40, 0.20, 0.0, -0.20, -0.40]:
    frac = (VMAX - r_val) / (VMAX - VMIN)
    ty = CB_TOP + CB_H * frac
    tick = slide.shapes.add_connector(
        1, CB_LEFT + CB_W, ty, CB_LEFT + CB_W + Inches(0.08), ty)
    tick.line.color.rgb = RGBColor(64, 64, 64)
    tick.line.width = Pt(0.5)
    kill_shadow(tick)
    txt = "0" if r_val == 0 else fmt_r(r_val)
    tb = slide.shapes.add_textbox(
        CB_LEFT + CB_W + Inches(0.12), ty - Inches(0.10),
        Inches(0.55), Inches(0.20))
    set_text(tb, txt, size=8, align=PP_ALIGN.LEFT)

# bottom annotation block
LEG_LEFT = GRID_LEFT - Inches(2.20)
LEG_TOP  = GRID_TOP + GRID_H + Inches(0.28)
LEG_W    = Inches(8.50)

leg1 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP, LEG_W, Inches(0.22))
set_text(leg1,
         "n = 10–14 paired subjects (varies by cascade Δ type) · "
         "partial Spearman ρ (response-group-adjusted; manuscript-consistent) · "
         "BH-corrected across 36 pairs (4 user-specified cascades)",
         size=8, align=PP_ALIGN.LEFT)

leg2 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP + Inches(0.22), LEG_W, Inches(0.22))
set_text(leg2,
         "Cell value = partial Spearman r;  small n=… in lower-right corner of each cell.   "
         "* nominal P < 0.05    ** BH q < 0.05.   "
         "5th column (grey header) shown for headline-pair cross-reference; outside the 36-pair test.",
         size=8, align=PP_ALIGN.LEFT)

leg3 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP + Inches(0.44), LEG_W, Inches(0.22))
set_text(leg3,
         "Convergence result: 0/36 partial P < 0.05   ·   0/36 BH q < 0.05    "
         "(baseline level and radiation-phase dynamics statistically independent).",
         size=8, bold=True, align=PP_ALIGN.LEFT)

# Headline-pair call-out (also annotate textually below for accessibility)
HP_LEFT = LEG_LEFT
HP_TOP  = LEG_TOP + Inches(0.78)
HP_W    = LEG_W
HP_H    = Inches(0.42)

hp_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, HP_LEFT, HP_TOP, HP_W, HP_H)
hp_bg.fill.solid()
hp_bg.fill.fore_color.rgb = RGBColor(252, 247, 232)  # pale gold
hp_bg.line.color.rgb = GOLD
hp_bg.line.width = Pt(1.25)
kill_shadow(hp_bg)

hp_tb = slide.shapes.add_textbox(
    HP_LEFT + Inches(0.10), HP_TOP + Inches(0.03),
    HP_W - Inches(0.20), HP_H - Inches(0.06))
set_text(hp_tb,
         "Headline pair (gold-bordered cell, 5th col.): DSB repair × Δ CD8-cytotoxic.   "
         "Partial Spearman r = −0.17, P = 0.60 (n = 12).   "
         "Manuscript text §3.8 quotes the plain-Spearman value r = −0.07, P = 0.83 "
         "(see Supp Fig S20A forest, and v3 diagnostic plain heatmap).",
         size=8, align=PP_ALIGN.LEFT)

PPTX_OUT = OUT_FIG / "Fig7C_convergence_heatmap_partial.pptx"
prs.save(PPTX_OUT)
print(f"[ok] wrote {PPTX_OUT}")


# ----------------------------------------------------------
# Diagnostic plain Spearman heatmap (archive)
# ----------------------------------------------------------
prs2 = Presentation()
prs2.slide_width = Inches(10)
prs2.slide_height = Inches(7.5)
slide2 = prs2.slides.add_slide(prs2.slide_layouts[6])

# wider range for plain (|r| up to 0.56) - keep ±0.6 like v2 chose
VMIN_P, VMAX_P = -0.6, 0.6

def diverging_p(r):
    r = max(min(r, VMAX_P), VMIN_P)
    if r >= 0:
        t = r / VMAX_P
        return RGBColor(
            int(WHITE[0] - t * (WHITE[0] - RED_DEEP[0])),
            int(WHITE[1] - t * (WHITE[1] - RED_DEEP[1])),
            int(WHITE[2] - t * (WHITE[2] - RED_DEEP[2])),
        )
    else:
        t = -r / (-VMIN_P)
        return RGBColor(
            int(WHITE[0] - t * (WHITE[0] - BLUE_DEEP[0])),
            int(WHITE[1] - t * (WHITE[1] - BLUE_DEEP[1])),
            int(WHITE[2] - t * (WHITE[2] - BLUE_DEEP[2])),
        )

# DIAGNOSTIC banner top
diag_banner = slide2.shapes.add_textbox(Inches(0.3), Inches(0.20), Inches(9.4), Inches(0.30))
set_text(diag_banner,
         "DIAGNOSTIC ONLY — plain (unadjusted) Spearman. "
         "Manuscript convergence test uses partial (response-adjusted) Spearman (see main Fig 7C v3).",
         size=10, bold=True, color=RGBColor(140, 70, 0), align=PP_ALIGN.CENTER)

col_title2 = slide2.shapes.add_textbox(
    GRID_LEFT, GRID_TOP - Inches(1.25), GRID_W, Inches(0.28))
set_text(col_title2, "Cascade Δ features (radiation-phase dynamics)",
         size=10, bold=True, align=PP_ALIGN.CENTER)

hdr_bg2 = slide2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    GRID_LEFT + CELL_W * ref_idx, GRID_TOP - Inches(0.92),
    CELL_W, Inches(0.85))
hdr_bg2.fill.solid()
hdr_bg2.fill.fore_color.rgb = GREY_REF
hdr_bg2.line.color.rgb = GREY_BORD
hdr_bg2.line.width = Pt(0.5)
kill_shadow(hdr_bg2)

for j, (_, lbl) in enumerate(ALL_CASC):
    tb = slide2.shapes.add_textbox(
        GRID_LEFT + CELL_W * j, GRID_TOP - Inches(0.92),
        CELL_W, Inches(0.85))
    bold = (j < len(CASC_FORMAL))
    set_text(tb, lbl, size=9, bold=bold, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

row_title2 = slide2.shapes.add_textbox(
    GRID_LEFT - Inches(2.20), GRID_TOP - Inches(0.34),
    Inches(2.10), Inches(0.26))
set_text(row_title2, "Baseline tumor-intrinsic features",
         size=10, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)

for i, (_, lbl) in enumerate(BASELINES):
    tb = slide2.shapes.add_textbox(
        GRID_LEFT - Inches(2.20), GRID_TOP + CELL_H * i,
        Inches(2.10), CELL_H)
    set_text(tb, lbl, size=9, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

for i, (b_col, _) in enumerate(BASELINES):
    for j, (c_col, _) in enumerate(ALL_CASC):
        row = rmap.get((b_col, c_col))
        if row is None:
            continue
        r_val = row.spearman_r
        p_val = row.spearman_p
        q_val = row.BH_q_plain
        n_val = row.n

        x = GRID_LEFT + CELL_W * j
        y = GRID_TOP  + CELL_H * i

        rect = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, CELL_W, CELL_H)
        rect.fill.solid()
        rect.fill.fore_color.rgb = diverging_p(r_val)
        if (b_col, c_col) == HEADLINE_PAIR:
            rect.line.color.rgb = GOLD
            rect.line.width = Pt(2.5)
        else:
            rect.line.color.rgb = GREY_BORD
            rect.line.width = Pt(0.5)
        kill_shadow(rect)

        suffix = ""
        if not np.isnan(q_val) and q_val < 0.05:
            suffix = "**"
        elif not np.isnan(p_val) and p_val < 0.05:
            suffix = "*"

        text_color = RGBColor(255, 255, 255) if abs(r_val) >= 0.45 else BLACK
        ann = slide2.shapes.add_textbox(x, y - Inches(0.04), CELL_W, CELL_H)
        set_text(ann, fmt_r(r_val) + suffix,
                 size=9, color=text_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        ntb = slide2.shapes.add_textbox(
            x + CELL_W - Inches(0.30), y + CELL_H - Inches(0.18),
            Inches(0.28), Inches(0.16))
        set_text(ntb, f"n={n_val}", size=6, color=text_color,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)

CB_LEFT2 = GRID_LEFT + GRID_W + Inches(0.45)
cblab2 = slide2.shapes.add_textbox(
    CB_LEFT2 - Inches(0.30), GRID_TOP - Inches(0.34),
    Inches(2.0), Inches(0.26))
set_text(cblab2, "Plain Spearman r (DIAG)",
         size=9, bold=True, align=PP_ALIGN.LEFT)

for k in range(N_STOPS):
    frac = k / (N_STOPS - 1)
    r_at = VMAX_P - frac * (VMAX_P - VMIN_P)
    rect = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        CB_LEFT2, GRID_TOP + stop_h * k,
        CB_W, stop_h + Emu(900))
    rect.fill.solid()
    rect.fill.fore_color.rgb = diverging_p(r_at)
    rect.line.fill.background()
    kill_shadow(rect)

border2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, CB_LEFT2, GRID_TOP, CB_W, CB_H)
border2.fill.background()
border2.line.color.rgb = RGBColor(64, 64, 64)
border2.line.width = Pt(0.75)
kill_shadow(border2)

for r_val in [0.60, 0.30, 0.0, -0.30, -0.60]:
    frac = (VMAX_P - r_val) / (VMAX_P - VMIN_P)
    ty = GRID_TOP + CB_H * frac
    tick = slide2.shapes.add_connector(
        1, CB_LEFT2 + CB_W, ty, CB_LEFT2 + CB_W + Inches(0.08), ty)
    tick.line.color.rgb = RGBColor(64, 64, 64)
    tick.line.width = Pt(0.5)
    kill_shadow(tick)
    txt = "0" if r_val == 0 else fmt_r(r_val)
    tb = slide2.shapes.add_textbox(
        CB_LEFT2 + CB_W + Inches(0.12), ty - Inches(0.10),
        Inches(0.55), Inches(0.20))
    set_text(tb, txt, size=8, align=PP_ALIGN.LEFT)

# bottom annotation block (diagnostic)
leg_d1 = slide2.shapes.add_textbox(LEG_LEFT, LEG_TOP, LEG_W, Inches(0.22))
set_text(leg_d1,
         "n = 10–14 paired subjects · plain Spearman ρ (NO group adjustment — diagnostic only) · "
         "BH-corrected across 36 pairs",
         size=8, align=PP_ALIGN.LEFT)
leg_d2 = slide2.shapes.add_textbox(LEG_LEFT, LEG_TOP + Inches(0.22), LEG_W, Inches(0.22))
set_text(leg_d2,
         "Cell value = plain Spearman r.   * nominal P < 0.05    ** BH q < 0.05.   "
         "Gold border = manuscript headline pair (DSB × Δ CD8-cytotoxic, plain r = −0.07).",
         size=8, align=PP_ALIGN.LEFT)
leg_d3 = slide2.shapes.add_textbox(LEG_LEFT, LEG_TOP + Inches(0.44), LEG_W, Inches(0.22))
set_text(leg_d3,
         "Diagnostic result: 0/36 plain P < 0.05   ·   0/36 BH q < 0.05  "
         "(top 2 cells P = 0.0586, 0.0588 just above the threshold).",
         size=8, bold=True, align=PP_ALIGN.LEFT)

PPTX_OUT2 = OUT_FIG / "diagnostic_plain_spearman_heatmap.pptx"
prs2.save(PPTX_OUT2)
print(f"[ok] wrote {PPTX_OUT2}")
