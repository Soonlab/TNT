"""
Supp Fig S11 — partial Spearman sensitivity heatmap (3-layer L3).

9 baseline × 4 cascade Δ (same 36-pair grid as Fig 7C). Demoted from main
figure; explicitly labelled as sensitivity analysis with collider-bias caveat.
"""
import sys
from pathlib import Path
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "scripts"))
from _shared_pptx_helpers import (
    diverging, kill_shadow, set_text, fmt_r,
    GREY_BORD, BLACK,
)

TSV_IN = ROOT / "tables" / "partial_spearman_36pair.tsv"
PPTX_OUT = ROOT / "figures" / "SuppFig_S11_partial_spearman_sensitivity.pptx"

BASELINES = [
    ("DNA Double-Strand Break Repair R-HSA-5693532", "DSB repair (Reactome)"),
    ("DNA Repair R-HSA-73894",                       "DNA repair (Reactome)"),
    ("HDR Thru Homologous Recombination (HRR) R-HSA-5685942",
                                                     "HRR (Reactome)"),
    ("E2F Targets",                                  "E2F targets (Hallmark)"),
    ("G2-M Checkpoint",                              "G2-M checkpoint (Hallmark)"),
    ("Myc Targets V2",                               "Myc targets V2 (Hallmark)"),
    ("MHC_II",                                       "MHC-II (baseline)"),
    ("MSI_pct",                                      "MSI (%)"),
    ("frac_amp",                                     "Genomic amp fraction"),
]
CASCADES = [
    ("SBS5_delta",         "Δ SBS5"),
    ("neo_binders_delta",  "Δ MHC-I neoantigen\nbinders"),
    ("Treg_delta",         "Δ Treg"),
    ("IGH_n_delta",        "Δ IGH clonotypes"),
]
VMIN, VMAX = -0.4, 0.4

df = pd.read_csv(TSV_IN, sep="\t")
rmap = {(r.baseline_feature, r.cascade_feature): r for r in df.itertuples()}

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

GRID_LEFT = Inches(3.05)
GRID_TOP  = Inches(1.85)
CELL_W    = Inches(0.84)
CELL_H    = Inches(0.48)
N_ROWS    = len(BASELINES)
N_COLS    = len(CASCADES)
GRID_W    = CELL_W * N_COLS
GRID_H    = CELL_H * N_ROWS

# Axis titles
col_title = slide.shapes.add_textbox(
    GRID_LEFT, GRID_TOP - Inches(1.20), GRID_W, Inches(0.28))
set_text(col_title, "Cascade Δ features (radiation-phase dynamics)",
         size=10, bold=True, align=PP_ALIGN.CENTER)

for j, (_, lbl) in enumerate(CASCADES):
    tb = slide.shapes.add_textbox(
        GRID_LEFT + CELL_W * j, GRID_TOP - Inches(0.92),
        CELL_W, Inches(0.85))
    set_text(tb, lbl, size=9, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

row_title = slide.shapes.add_textbox(
    GRID_LEFT - Inches(2.20), GRID_TOP - Inches(0.34),
    Inches(2.10), Inches(0.26))
set_text(row_title, "Baseline tumor-intrinsic features",
         size=10, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)

for i, (_, lbl) in enumerate(BASELINES):
    tb = slide.shapes.add_textbox(
        GRID_LEFT - Inches(2.20), GRID_TOP + CELL_H * i,
        Inches(2.10), CELL_H)
    set_text(tb, lbl, size=9, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# Cells
for i, (b_col, _) in enumerate(BASELINES):
    for j, (c_col, _) in enumerate(CASCADES):
        row = rmap.get((b_col, c_col))
        if row is None:
            continue
        r_val = row.partial_r
        p_val = row.partial_P
        q_val = row.BH_q_partial

        x = GRID_LEFT + CELL_W * j
        y = GRID_TOP  + CELL_H * i

        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, CELL_W, CELL_H)
        rect.fill.solid()
        rect.fill.fore_color.rgb = diverging(r_val, VMIN, VMAX)
        rect.line.color.rgb = GREY_BORD
        rect.line.width = Pt(0.5)
        kill_shadow(rect)

        suffix = ""
        if not pd.isna(q_val) and q_val < 0.05:
            suffix = "**"
        elif p_val < 0.05:
            suffix = "*"

        text_color = RGBColor(255, 255, 255) if abs(r_val) >= 0.30 else BLACK
        ann = slide.shapes.add_textbox(x, y - Inches(0.04), CELL_W, CELL_H)
        set_text(ann, fmt_r(r_val) + suffix,
                 size=9, color=text_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        if p_val < 0.10:
            ptb = slide.shapes.add_textbox(
                x, y + CELL_H - Inches(0.20), CELL_W, Inches(0.18))
            set_text(ptb, f"P={p_val:.3f}", size=7, color=text_color,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

# Colorbar
CB_LEFT = GRID_LEFT + GRID_W + Inches(0.40)
CB_TOP  = GRID_TOP
CB_W    = Inches(0.26)
CB_H    = GRID_H

cblab = slide.shapes.add_textbox(
    CB_LEFT - Inches(0.30), CB_TOP - Inches(0.34),
    Inches(1.9), Inches(0.26))
set_text(cblab, "Partial Spearman r", size=9, bold=True, align=PP_ALIGN.LEFT)

N_STOPS = 100
stop_h = CB_H / N_STOPS
for k in range(N_STOPS):
    frac = k / (N_STOPS - 1)
    r_at = VMAX - frac * (VMAX - VMIN)
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        CB_LEFT, CB_TOP + stop_h * k,
        CB_W, stop_h + Emu(900))
    rect.fill.solid()
    rect.fill.fore_color.rgb = diverging(r_at, VMIN, VMAX)
    rect.line.fill.background()
    kill_shadow(rect)

border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, CB_LEFT, CB_TOP, CB_W, CB_H)
border.fill.background()
border.line.color.rgb = RGBColor(64, 64, 64)
border.line.width = Pt(0.75)
kill_shadow(border)

for r_val in [0.40, 0.20, 0.0, -0.20, -0.40]:
    frac = (VMAX - r_val) / (VMAX - VMIN)
    ty = CB_TOP + CB_H * frac
    tick = slide.shapes.add_connector(
        1, CB_LEFT + CB_W, ty, CB_LEFT + CB_W + Inches(0.08), ty)
    tick.line.color.rgb = RGBColor(64, 64, 64)
    tick.line.width = Pt(0.5)
    kill_shadow(tick)
    txt = "0" if r_val == 0 else fmt_r(r_val)
    tb = slide.shapes.add_textbox(
        CB_LEFT + CB_W + Inches(0.12), ty - Inches(0.10),
        Inches(0.55), Inches(0.20))
    set_text(tb, txt, size=8, align=PP_ALIGN.LEFT)

# Bottom annotation
LEG_LEFT = GRID_LEFT - Inches(2.20)
LEG_TOP  = GRID_TOP + GRID_H + Inches(0.22)
LEG_W    = Inches(8.50)

leg_head = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP, LEG_W, Inches(0.22))
set_text(leg_head,
         "Sensitivity analysis — partial Spearman with response-group adjustment.",
         size=9, bold=True, align=PP_ALIGN.LEFT)

leg1 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP + Inches(0.22), LEG_W, Inches(0.20))
set_text(leg1,
         "Partial Spearman ρ · BH-corrected across 36 pairs.   "
         "P value shown for cells with P < 0.10.",
         size=8, align=PP_ALIGN.LEFT)

leg2 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP + Inches(0.42), LEG_W, Inches(0.20))
set_text(leg2,
         "Partial r range [−0.48, +0.46]  "
         "(plain Spearman range [−0.54, +0.56] compressed by group adjustment).",
         size=8, align=PP_ALIGN.LEFT)

leg3 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP + Inches(0.62), LEG_W, Inches(0.20))
set_text(leg3,
         "0/36 partial P < 0.05   ·   0/36 BH q < 0.05    (robust to group adjustment).",
         size=8, bold=True, align=PP_ALIGN.LEFT)

# Caveat box
CAV_TOP = LEG_TOP + Inches(0.88)
CAV_H   = Inches(0.50)
cav_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                LEG_LEFT, CAV_TOP, LEG_W, CAV_H)
cav_bg.fill.solid()
cav_bg.fill.fore_color.rgb = RGBColor(252, 240, 230)
cav_bg.line.color.rgb = RGBColor(197, 62, 31)
cav_bg.line.width = Pt(1.0)
kill_shadow(cav_bg)

cav_tb = slide.shapes.add_textbox(
    LEG_LEFT + Inches(0.08), CAV_TOP + Inches(0.04),
    LEG_W - Inches(0.16), CAV_H - Inches(0.08))
set_text(cav_tb,
         "Caveat: Conditioning on response group (an outcome variable in our analytical design) "
         "may introduce collider bias under some causal structures. This sensitivity analysis is "
         "presented as supportive pair-level null evidence, not the primary test.  "
         "Primary analysis: Fig 7C (plain Spearman).  Omnibus pattern test: Supp Fig S12.",
         size=8, align=PP_ALIGN.LEFT)

PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(PPTX_OUT)
print(f"[ok] wrote {PPTX_OUT}")
