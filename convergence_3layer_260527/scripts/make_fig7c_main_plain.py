"""
Fig 7C (Main) — promoted plain Spearman heatmap (3-layer L1).

9 baseline × 4 cascade Δ (no 5th Δ CD8-cyt column — that pair is now handled
by a separate headline-pair textbox).
"""
import sys
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "scripts"))
from _shared_pptx_helpers import (
    diverging, kill_shadow, set_text, fmt_r,
    GREY_BORD, BLACK, GOLD,
)

TSV_IN = ROOT / "tables" / "plain_spearman_36pair.tsv"
PPTX_OUT = ROOT / "figures" / "Fig7C_main_plain_spearman.pptx"

BASELINES = [
    ("DNA Double-Strand Break Repair R-HSA-5693532", "DSB repair (Reactome)"),
    ("DNA Repair R-HSA-73894",                       "DNA repair (Reactome)"),
    ("HDR Thru Homologous Recombination (HRR) R-HSA-5685942",
                                                     "HRR (Reactome)"),
    ("E2F Targets",                                  "E2F targets (Hallmark)"),
    ("G2-M Checkpoint",                              "G2-M checkpoint (Hallmark)"),
    ("Myc Targets V2",                               "Myc targets V2 (Hallmark)"),
    ("MHC_II",                                       "MHC-II (baseline)"),
    ("MSI_pct",                                      "MSI (%)"),
    ("frac_amp",                                     "Genomic amp fraction"),
]
CASCADES = [
    ("SBS5_delta",         "Δ SBS5"),
    ("neo_binders_delta",  "Δ MHC-I neoantigen\nbinders"),
    ("Treg_delta",         "Δ Treg"),
    ("IGH_n_delta",        "Δ IGH clonotypes"),
]
VMIN, VMAX = -0.6, 0.6  # plain |r| up to 0.56

df = pd.read_csv(TSV_IN, sep="\t")
rmap = {(r.baseline_feature, r.cascade_feature): r for r in df.itertuples()}

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

GRID_LEFT = Inches(3.30)
GRID_TOP  = Inches(1.85)
CELL_W    = Inches(0.92)
CELL_H    = Inches(0.50)
N_ROWS    = len(BASELINES)
N_COLS    = len(CASCADES)
GRID_W    = CELL_W * N_COLS
GRID_H    = CELL_H * N_ROWS

# Column axis title
col_title = slide.shapes.add_textbox(
    GRID_LEFT, GRID_TOP - Inches(1.20), GRID_W, Inches(0.28))
set_text(col_title, "Cascade Δ features (radiation-phase dynamics)",
         size=10, bold=True, align=PP_ALIGN.CENTER)

# Column labels
for j, (_, lbl) in enumerate(CASCADES):
    tb = slide.shapes.add_textbox(
        GRID_LEFT + CELL_W * j, GRID_TOP - Inches(0.92),
        CELL_W, Inches(0.85))
    set_text(tb, lbl, size=9, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

# Row axis title
row_title = slide.shapes.add_textbox(
    GRID_LEFT - Inches(2.20), GRID_TOP - Inches(0.34),
    Inches(2.10), Inches(0.26))
set_text(row_title, "Baseline tumor-intrinsic features",
         size=10, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)

# Row labels
for i, (_, lbl) in enumerate(BASELINES):
    tb = slide.shapes.add_textbox(
        GRID_LEFT - Inches(2.20), GRID_TOP + CELL_H * i,
        Inches(2.10), CELL_H)
    set_text(tb, lbl, size=9, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# Cells
for i, (b_col, _) in enumerate(BASELINES):
    for j, (c_col, _) in enumerate(CASCADES):
        row = rmap.get((b_col, c_col))
        if row is None:
            continue
        r_val = row.spearman_r
        p_val = row.spearman_p
        q_val = row.BH_q_plain
        n_val = row.n

        x = GRID_LEFT + CELL_W * j
        y = GRID_TOP  + CELL_H * i

        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, CELL_W, CELL_H)
        rect.fill.solid()
        rect.fill.fore_color.rgb = diverging(r_val, VMIN, VMAX)
        rect.line.color.rgb = GREY_BORD
        rect.line.width = Pt(0.5)
        kill_shadow(rect)

        suffix = ""
        if not pd.isna(q_val) and q_val < 0.05:
            suffix = "**"
        elif p_val < 0.05:
            suffix = "*"

        text_color = RGBColor(255, 255, 255) if abs(r_val) >= 0.45 else BLACK
        ann = slide.shapes.add_textbox(x, y - Inches(0.04), CELL_W, CELL_H)
        set_text(ann, fmt_r(r_val) + suffix,
                 size=9, color=text_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        if p_val < 0.10:
            ptb = slide.shapes.add_textbox(
                x, y + CELL_H - Inches(0.20), CELL_W, Inches(0.18))
            set_text(ptb, f"P={p_val:.3f}", size=7, color=text_color,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

# Colorbar
CB_LEFT = GRID_LEFT + GRID_W + Inches(0.45)
CB_TOP  = GRID_TOP
CB_W    = Inches(0.26)
CB_H    = GRID_H

cblab = slide.shapes.add_textbox(
    CB_LEFT - Inches(0.20), CB_TOP - Inches(0.34),
    Inches(1.7), Inches(0.26))
set_text(cblab, "Plain Spearman r", size=9, bold=True, align=PP_ALIGN.LEFT)

N_STOPS = 100
stop_h = CB_H / N_STOPS
for k in range(N_STOPS):
    frac = k / (N_STOPS - 1)
    r_at = VMAX - frac * (VMAX - VMIN)
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        CB_LEFT, CB_TOP + stop_h * k,
        CB_W, stop_h + Emu(900))
    rect.fill.solid()
    rect.fill.fore_color.rgb = diverging(r_at, VMIN, VMAX)
    rect.line.fill.background()
    kill_shadow(rect)

border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, CB_LEFT, CB_TOP, CB_W, CB_H)
border.fill.background()
border.line.color.rgb = RGBColor(64, 64, 64)
border.line.width = Pt(0.75)
kill_shadow(border)

for r_val in [0.60, 0.30, 0.0, -0.30, -0.60]:
    frac = (VMAX - r_val) / (VMAX - VMIN)
    ty = CB_TOP + CB_H * frac
    tick = slide.shapes.add_connector(
        1, CB_LEFT + CB_W, ty, CB_LEFT + CB_W + Inches(0.08), ty)
    tick.line.color.rgb = RGBColor(64, 64, 64)
    tick.line.width = Pt(0.5)
    kill_shadow(tick)
    txt = "0" if r_val == 0 else fmt_r(r_val)
    tb = slide.shapes.add_textbox(
        CB_LEFT + CB_W + Inches(0.12), ty - Inches(0.10),
        Inches(0.55), Inches(0.20))
    set_text(tb, txt, size=8, align=PP_ALIGN.LEFT)

# Bottom annotation block
LEG_LEFT = GRID_LEFT - Inches(2.20)
LEG_TOP  = GRID_TOP + GRID_H + Inches(0.28)
LEG_W    = Inches(8.50)

leg1 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP, LEG_W, Inches(0.22))
set_text(leg1,
         "Cell value = plain Spearman r;  BH-corrected across 36 pairs.   "
         "P value shown for cells with P < 0.10.   "
         "* nominal P < 0.05    ** BH q < 0.05.",
         size=8, align=PP_ALIGN.LEFT)

leg2 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP + Inches(0.24), LEG_W, Inches(0.22))
set_text(leg2,
         "0/36 nominal P < 0.05   ·   0/36 BH q < 0.05    "
         "(no individual baseline-cascade pair detected).",
         size=8, bold=True, align=PP_ALIGN.LEFT)

leg3 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP + Inches(0.48), LEG_W, Inches(0.22))
set_text(leg3,
         "See Supp Fig S12 for omnibus block-wise sign-coherence test "
         "(12/12 cells in predicted direction, P = 2.4 × 10⁻⁴).   "
         "See Supp Fig S11 for partial-Spearman sensitivity (response-group adjusted).",
         size=8, align=PP_ALIGN.LEFT)

# Headline-pair call-out (top-right of slide)
HP_LEFT = Inches(7.65)
HP_TOP  = Inches(0.30)
HP_W    = Inches(2.20)
HP_H    = Inches(1.15)

hp_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, HP_LEFT, HP_TOP, HP_W, HP_H)
hp_bg.fill.solid()
hp_bg.fill.fore_color.rgb = RGBColor(252, 247, 232)
hp_bg.line.color.rgb = GOLD
hp_bg.line.width = Pt(1.25)
kill_shadow(hp_bg)

hp_tb = slide.shapes.add_textbox(
    HP_LEFT + Inches(0.08), HP_TOP + Inches(0.06),
    HP_W - Inches(0.16), HP_H - Inches(0.12))
set_text(hp_tb,
         "Headline pair (manuscript-quoted)\n"
         "DSB repair × Δ CD8-cytotoxic\n"
         "r = −0.07,  P = 0.83 (n = 12)\n"
         "[not in this 4-cascade panel;\nsee Supp Fig S20A forest]",
         size=8, align=PP_ALIGN.LEFT)

PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(PPTX_OUT)
print(f"[ok] wrote {PPTX_OUT}")
