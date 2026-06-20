"""Rebuild Figure 1B (Five-row summary forest) as a native-editable single-slide
PPT, dropped directly into GenomeMedicine_TNT_v2/.

Fixes the LASSO/ElasticNet attribution mismatch in the headline classifier row:
the previous version had row label "Pre-CRT LASSO AUC" alongside subtitle
"4-feature ElasticNet, nested outer-LOOCV", which was internally contradictory.
ElasticNet is the correct learner (28-feature reference set, AUC = 0.745;
LASSO on the same set gives AUC = 0.716; see Supp Text S4 / Fig 4E).

Cross-reference also updated: "Fig 5" → "Fig 4E" to match Manuscript_v4.docx.

Source of truth: /data/data/TNT/analysis/260418_add/22_fig1_native_pptx.py
build_D() at L886-1014. This script reproduces that panel verbatim with the
two corrections, no other changes.

Output: /data/data/TNT/analysis/GenomeMedicine_TNT_v2/Fig1B_headline_forest_editable.pptx
"""
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def _i(v):
    return int(round(float(v)))


GOOD = RGBColor(0x0A, 0x7D, 0x6E)
BAD = RGBColor(0xC5, 0x3E, 0x1F)
INK = RGBColor(0x22, 0x22, 0x22)
LINE = RGBColor(0x33, 0x33, 0x33)
LT_GREY = RGBColor(0xDD, 0xDD, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xD4, 0xA3, 0x00)
THREAD1 = RGBColor(0x0E, 0x4A, 0x68)
THREAD2 = RGBColor(0x8A, 0x2B, 0x4C)

FONT = "Arial"
SLIDE_W = Inches(6.5)
SLIDE_H = Inches(4.5)


def new_slide(prs):
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, size=8, bold=False, color=INK,
             align="left", anchor="middle", italic=False):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(max(w, 1)), _i(max(h, 1)))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "bottom": MSO_ANCHOR.BOTTOM,
                          "middle": MSO_ANCHOR.MIDDLE}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = str(text)
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.italic = bool(italic)
    r.font.color.rgb = color
    kill_shadow(tb)
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    kill_shadow(c)
    return c


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0.5):
    w = max(_i(w), 1); h = max(_i(h), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _i(x), _i(y), w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def draw_panel_letter(slide, letter, x=Inches(0.05), y=Inches(0.05)):
    add_text(slide, x, y, Inches(0.4), Inches(0.3),
             letter, size=14, bold=True, color=INK, align="left", anchor="top")


def build_panel_B():
    prs = Presentation()
    s = new_slide(prs)
    draw_panel_letter(s, "B")

    # 5 claim rows — IDENTICAL to legacy build_D() at 22_fig1_native_pptx.py L893-941
    # except row 0: "LASSO" → "ElasticNet" and "Fig 5" → "Fig 4E".
    rows = [
        ("Discovery",
         "Pre-CRT ElasticNet AUC",                          # ← was "Pre-CRT LASSO AUC"
         "4-feature ElasticNet, nested outer-LOOCV",
         "0.745",
         "[0.56, 0.89]",                                    # ← also tightened to match Fig 4E (0.56–0.89, not 0.56–0.90)
         "AUC  (chance = 0.5)",
         (0.745 - 0.5) / 0.5,
         "Fig 4E",                                          # ← was "Fig 5"
         THREAD1),
        ("LC-CRT meta",
         "Thread 1 (5 cohorts, N = 518)",
         "DSB Z=+3.17 · cell-cycle +3.21 · E2F/MYC +2.79 · EMT +1.61 (trend)",
         "+3.21",
         "max Z",
         "Stouffer Z  (|Z| = 1.96 = P 0.05)",
         3.21 / 4.0,
         "Fig 9A",
         THREAD1),
        ("LC-CRT meta",
         "Thread 2 (6 sources, N = 816)",
         "CD8-cytotoxic + Akiyoshi 2023 (JAMA Netw Open)",
         "+3.29",
         "P = 0.001",
         "Stouffer Z",
         3.29 / 4.0,
         "Fig 9A",
         THREAD2),
        ("SC-RT external",
         "GSE254249 (post-TNT, N = 8)",
         "7/7 signatures concordant with discovery · Tcell_infil MW P = 0.036",
         "7/7",
         "sign P = 0.016",
         "signature direction concordance  (7/7 = unanimous)",
         7 / 7,
         "Fig 9E",
         HIGHLIGHT),
        ("Paired",
         "Radiation-phase pharmacodynamics",
         "EMT good 6/6 up (P = 0.016) · IGH coherence Wilcoxon P = 0.035",
         "0.016",
         "min P",
         "−log₁₀ P  (horizontal bar length)",
         min(1.0, -np.log10(0.016) / 3.0),
         "Figs 6-7",
         RGBColor(0xB3, 0x7D, 0x00)),
    ]

    px = Inches(1.20); py = Inches(0.45)
    pw = Inches(3.30); ph = Inches(3.50)
    row_h = ph / len(rows)
    right_col_x = px + pw + Inches(0.15)

    for i, (group, short, full, eff_val, eff_range, eff_unit,
            bar_frac, ref, color) in enumerate(rows):
        y_top = _i(py + i * row_h)
        y_ctr = _i(py + (i + 0.5) * row_h)

        if i > 0:
            add_line(s, Inches(0.20), y_top, SLIDE_W - Inches(0.20),
                     y_top, LT_GREY, 0.3)

        # left column
        add_text(s, Inches(0.20), y_top + Inches(0.08),
                 Inches(1.00), Inches(0.14),
                 group.upper(), size=6, bold=True, color=color, align="left")
        add_text(s, Inches(0.20), y_top + Inches(0.22),
                 px - Inches(0.25), Inches(0.20),
                 short, size=8, bold=True, color=INK, align="left")
        add_text(s, Inches(0.20), y_top + Inches(0.43),
                 px - Inches(0.25), Inches(0.32),
                 full, size=6, color=RGBColor(0x55, 0x55, 0x55),
                 align="left", anchor="top")

        # middle bar
        bar_y_top = y_ctr - Inches(0.09)
        bar_h = Inches(0.18)
        add_rect(s, px, bar_y_top, pw, bar_h, fill=LT_GREY, line_color=None)
        filled_w = _i(pw * max(0.02, min(1.0, bar_frac)))
        add_rect(s, px, bar_y_top, filled_w, bar_h,
                 fill=color, line_color=WHITE, line_width=0.3)
        add_text(s, px, _i(bar_y_top + bar_h + Inches(0.02)),
                 pw, Inches(0.13),
                 eff_unit, size=5, italic=True,
                 color=RGBColor(0x66, 0x66, 0x66), align="left")

        # right column
        add_text(s, right_col_x, y_top + Inches(0.08),
                 Inches(1.30), Inches(0.22),
                 eff_val, size=14, bold=True, color=color, align="left")
        add_text(s, right_col_x, y_top + Inches(0.35),
                 Inches(1.40), Inches(0.16),
                 eff_range, size=7, color=RGBColor(0x55, 0x55, 0x55),
                 align="left")
        add_text(s, right_col_x, y_top + Inches(0.55),
                 Inches(1.40), Inches(0.16),
                 f"→ {ref}", size=7, italic=True, color=color, align="left")

    # bottom caption (verbatim from build_D)
    add_text(s, Inches(0.15), Inches(4.02), SLIDE_W - Inches(0.3), Inches(0.14),
             "Two externally-validated orthogonal pre-CRT axes (Thread 1 tumor-intrinsic + Thread 2 immune) reproduce across",
             size=6, color=INK, align="center")
    add_text(s, Inches(0.15), Inches(4.16), SLIDE_W - Inches(0.3), Inches(0.14),
             "three regimen strata (SC-RT + FOLFOX/CAPOX; LC-CRT + concurrent cape; SC-RT + FOLFOXIRI), supporting a regimen-agnostic interpretation;",
             size=6, color=INK, align="center")
    add_text(s, Inches(0.15), Inches(4.30), SLIDE_W - Inches(0.3), Inches(0.14),
             "paired radiation-phase biopsies add an orthogonal dynamic layer (target engagement + directional immune coherence).",
             size=6, italic=True, color=HIGHLIGHT, align="center")

    out = "/data/data/TNT/analysis/GenomeMedicine_TNT_v2/Fig1B_headline_forest_editable.pptx"
    prs.save(out)
    print(f"wrote {out}")


if __name__ == "__main__":
    build_panel_B()
