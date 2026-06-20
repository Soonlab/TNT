"""Rebuild Supplementary Figure S5D as a native-editable single-slide PPT
showing the 3-model nested-LOOCV comparison on the 28-feature reference set.

Replaces the legacy panel sourced from 10_ml_predictor/ml_loocv_results.tsv,
which was non-nested (leakage-prone) LOOCV mislabeled as "nested LOOCV" in
the manuscript caption — see audit at:
  - scripts/supp_figures_v1.py L470 (reads ml_loocv_results.tsv = non-nested)
  - 10_ml_predictor/ml_loocv_results.tsv (LASSO 0.7549 / ElasticNet 0.6699 / RF 0.5768)
  - Manuscript_v4.docx caption para 114 (mislabels these as "nested LOOCV")

The new panel uses the honest nested-LOOCV values from
260418_add/nested_cv_28feat_three_model.tsv:
  LASSO        AUC = 0.7157 [0.5238, 0.8824]
  ElasticNet   AUC = 0.7451 [0.5633, 0.8947]   ← headline (Fig 4E reference)
  RandomForest AUC = 0.5752 [0.3758, 0.7664]

Computed by scripts/57_compute (inline above) under outer LOOCV + inner 5-fold
SelectKBest tuning, with random_state=0 set on the saga solver to ensure
reproducibility (saga is non-deterministic without random_state).

Output: GenomeMedicine_TNT_v2/SuppFig_S5D_three_model_nested_28feat_editable.pptx
"""
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def _i(v):
    return int(round(float(v)))


GOOD_DEEP = RGBColor(0x0A, 0x7D, 0x6E)
BAD_DEEP  = RGBColor(0xC5, 0x3E, 0x1F)
GOLD_DEEP = RGBColor(0xB3, 0x7D, 0x00)
GREY_DEEP = RGBColor(0x55, 0x55, 0x55)
BLACK     = RGBColor(0x22, 0x22, 0x22)
LT_GREY   = RGBColor(0xDD, 0xDD, 0xDD)
LINE_GREY = RGBColor(0xAA, 0xAA, 0xAA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xD4, 0xA3, 0x00)

FONT = "Arial"
SLIDE_W = Inches(6.5)
SLIDE_H = Inches(4.5)


def new_slide(prs):
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, size=8, bold=False, color=BLACK,
             align="left", anchor="middle", italic=False):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(max(w,1)), _i(max(h,1)))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "bottom": MSO_ANCHOR.BOTTOM,
                          "middle": MSO_ANCHOR.MIDDLE}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = str(text)
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.italic = bool(italic)
    r.font.color.rgb = color
    kill_shadow(tb)
    return tb


def add_line(slide, x1, y1, x2, y2, color=BLACK, width=0.6, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        try:
            c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass
    kill_shadow(c)
    return c


def add_circle(slide, cx, cy, r, fill=None, line_color=None, line_width=0.6):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx) - r, _i(cy) - r, 2*r, 2*r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def draw_panel_letter(slide, letter):
    add_text(slide, Inches(0.05), Inches(0.05), Inches(0.4), Inches(0.3),
             letter, size=14, bold=True, color=BLACK, align="left", anchor="top")


# ----- data -----
DATA = pd.read_csv(
    "/data/data/TNT/analysis/260418_add/nested_cv_28feat_three_model.tsv",
    sep="\t")

# Plot order: high → low AUC for visual ranking
DATA = DATA.sort_values("AUC", ascending=False).reset_index(drop=True)

# ----- build slide -----
prs = Presentation()
s = new_slide(prs)
draw_panel_letter(s, "D")

# Plot area
plot_x0 = Inches(1.55); plot_x1 = Inches(5.95)
plot_y0 = Inches(0.85); plot_y1 = Inches(2.85)
plot_w = plot_x1 - plot_x0           # EMU int
plot_h = plot_y1 - plot_y0           # EMU int
plot_h_in = (plot_y1 - plot_y0) / Inches(1)   # height in inches as float

# X-axis: AUC, range 0.30 to 1.00
AUC_MIN, AUC_MAX = 0.30, 1.00
def x_for(auc):
    return plot_x0 + plot_w * (auc - AUC_MIN) / (AUC_MAX - AUC_MIN)

# Background grid + axis ticks
for tick in [0.3, 0.4, 0.5, 0.6, 0.7, 0.8, 0.9, 1.0]:
    add_line(s, x_for(tick), plot_y0, x_for(tick), plot_y1,
             color=LT_GREY, width=0.4, dashed=False)
    add_text(s, x_for(tick) - Inches(0.18), plot_y1 + Inches(0.04),
             Inches(0.36), Inches(0.16),
             f"{tick:.1f}", size=7, color=GREY_DEEP, align="center")

# Chance line (AUC = 0.5)
add_line(s, x_for(0.5), plot_y0 - Inches(0.05),
         x_for(0.5), plot_y1 + Inches(0.02),
         color=GREY_DEEP, width=1.0, dashed=True)
add_text(s, x_for(0.5) - Inches(0.45), plot_y0 - Inches(0.22),
         Inches(0.9), Inches(0.16),
         "chance", size=6, italic=True, color=GREY_DEEP, align="center")

# Top + bottom plot frame
add_line(s, plot_x0, plot_y0, plot_x1, plot_y0, color=BLACK, width=0.7)
add_line(s, plot_x0, plot_y1, plot_x1, plot_y1, color=BLACK, width=0.7)
add_line(s, plot_x0, plot_y0, plot_x0, plot_y1, color=BLACK, width=0.7)

# X-axis label
add_text(s, plot_x0, plot_y1 + Inches(0.24),
         plot_w, Inches(0.18),
         "Outer-LOOCV AUC  (95 % bootstrap CI, 2000 resamples)",
         size=9, color=BLACK, align="center", anchor="top")

# Rows: one per model
n_rows = len(DATA)
row_h_in = plot_h_in / (n_rows + 1)   # leave one slot of breathing room

for i, r in DATA.iterrows():
    yc = plot_y0 + Inches(row_h_in * (i + 0.7))
    auc, lo, hi, mdl = r.AUC, r.CI_low, r.CI_high, r.model

    # Color: ElasticNet (headline) = GOOD_DEEP, LASSO = GOLD_DEEP, RF = GREY_DEEP
    if mdl == "ElasticNet":
        c = GOOD_DEEP
    elif mdl == "LASSO":
        c = GOLD_DEEP
    else:
        c = GREY_DEEP

    # Row label (model name)
    add_text(s, Inches(0.40), yc - Inches(0.10),
             Inches(1.10), Inches(0.20),
             mdl, size=10, bold=True, color=c, align="right", anchor="middle")

    # CI bar
    add_line(s, x_for(lo), yc, x_for(hi), yc, color=c, width=2.5)
    # Cap ticks
    add_line(s, x_for(lo), yc - Inches(0.06),
             x_for(lo), yc + Inches(0.06), color=c, width=1.5)
    add_line(s, x_for(hi), yc - Inches(0.06),
             x_for(hi), yc + Inches(0.06), color=c, width=1.5)
    # Point estimate
    add_circle(s, x_for(auc), yc, Inches(0.10),
               fill=c, line_color=BLACK, line_width=0.7)

    # AUC value to the right of the CI
    add_text(s, x_for(hi) + Inches(0.10), yc - Inches(0.10),
             Inches(0.95), Inches(0.20),
             f"AUC = {auc:.3f}", size=9, bold=(mdl == "ElasticNet"),
             color=c, align="left", anchor="middle")
    # CI text below
    add_text(s, x_for(hi) + Inches(0.10), yc + Inches(0.06),
             Inches(0.95), Inches(0.16),
             f"[{lo:.2f}, {hi:.2f}]", size=7, color=GREY_DEEP,
             align="left", anchor="middle")

    # Headline marker for ElasticNet
    if mdl == "ElasticNet":
        add_text(s, Inches(0.20), yc + Inches(0.18),
                 Inches(1.20), Inches(0.16),
                 "headline (Fig 4E)", size=6, italic=True,
                 color=HIGHLIGHT, align="right", anchor="top")

# ----- Annotations / legend block -----
# Title
add_text(s, Inches(0.20), Inches(0.30),
         SLIDE_W - Inches(0.4), Inches(0.20),
         "Three-model nested-LOOCV comparison on the 28-feature reference set",
         size=10, bold=True, color=BLACK, align="center", anchor="top")
add_text(s, Inches(0.20), Inches(0.50),
         SLIDE_W - Inches(0.4), Inches(0.18),
         "(outer 35-fold leave-one-out + inner 5-fold SelectKBest tuning; "
         "bootstrap 95 % CI from 2000 resamples)",
         size=7, italic=True, color=GREY_DEEP, align="center", anchor="top")

# Footnote: contrast vs the older non-nested LOOCV reference values
add_text(s, Inches(0.20), Inches(3.30),
         SLIDE_W - Inches(0.4), Inches(0.18),
         "For comparison, an earlier non-nested LOOCV pass (feature selection on full training set; leakage-prone) gave",
         size=7, color=GREY_DEEP, align="center", anchor="top")
add_text(s, Inches(0.20), Inches(3.46),
         SLIDE_W - Inches(0.4), Inches(0.18),
         "LASSO 0.755 / ElasticNet 0.670 / RandomForest 0.577 (Supp Text S4 §S4.1, reported for transparency only).",
         size=7, color=GREY_DEEP, italic=True, align="center", anchor="top")

# Bottom interpretive block
add_text(s, Inches(0.20), Inches(3.78),
         SLIDE_W - Inches(0.4), Inches(0.18),
         "ElasticNet on the contamination-cleaned 28-feature set is the headline classifier (AUC 0.745, Fig 4E ROC).",
         size=8, color=BLACK, align="center", anchor="top")
add_text(s, Inches(0.20), Inches(3.96),
         SLIDE_W - Inches(0.4), Inches(0.18),
         "LASSO and ElasticNet 95 % CIs overlap; RandomForest's lower performance is consistent with an approximately linear signal in the selected feature space.",
         size=7, italic=True, color=GREY_DEEP, align="center", anchor="top")
add_text(s, Inches(0.20), Inches(4.14),
         SLIDE_W - Inches(0.4), Inches(0.18),
         "All three lower bounds approach 0.5 — the predictor is reported as a discovery-stage signal pending external TNT-matched validation (§3.10–§3.11).",
         size=7, italic=True, color=GREY_DEEP, align="center", anchor="top")

OUT = "/data/data/TNT/analysis/GenomeMedicine_TNT_v2/SuppFig_S5D_three_model_nested_28feat_editable.pptx"
prs.save(OUT)
print(f"wrote {OUT}")
