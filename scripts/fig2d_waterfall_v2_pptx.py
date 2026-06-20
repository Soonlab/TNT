"""Fig 2D v2 — native PowerPoint reconstruction of the TMB + MSI waterfall.

Every visible element is a native PPT object — users can double-click and edit
each without ungroup:
  - Text (axis labels, tick labels, sample IDs, group labels, MSI value
    annotations, legend, MSI-H cutoff callout, title) = add_textbox (TEXT_BOX)
  - MSI bars, response/timepoint strips, legend swatches = MSO_SHAPE.RECTANGLE
  - TMB markers = MSO_SHAPE.OVAL
  - Axis spines, tick marks, TMB cutoff line, group/sub-group separators
    = MSO_CONNECTOR.STRAIGHT (dashed where applicable)

Output: manuscript/ppt/Fig2D_TMB_MSI_waterfall_native_editable.pptx
"""
from pathlib import Path
import importlib.util
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

spec = importlib.util.spec_from_file_location(
    'nph', '/mnt/sda1/data/TNT/analysis/scripts/44_native_ppt_helpers.py')
nph = importlib.util.module_from_spec(spec); spec.loader.exec_module(nph)

BASE = Path('/mnt/sda1/data/TNT/analysis')
OUT_DIR = BASE/'manuscript/ppt'; OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_FILE = OUT_DIR/'Fig2D_TMB_MSI_waterfall_native_editable.pptx'

def hex2rgb(h): return (int(h[1:3],16), int(h[3:5],16), int(h[5:7],16))
GOOD = hex2rgb('#2E86AB'); BAD = hex2rgb('#E63946')
TMB_COLOR = hex2rgb('#6a4c93')
PRE = hex2rgb('#457b9d'); POST = hex2rgb('#f4a261')
SEP = hex2rgb('#1d3557')
CUTOFF_RED = hex2rgb('#d62828')

UNMATCHED = {13,15,16,17,18,19,33}

# ---------- data ----------
msi_files = list((BASE/'02_wes_tmb_msi/msi').glob('*.tsv'))
msi = pd.concat([pd.read_csv(f, sep='\t') for f in msi_files])
tmb = pd.read_csv(BASE/'02_wes_tmb_msi/tmb_per_sample.tsv', sep='\t')
mer = msi.merge(tmb[['sample_id','TMB_nonsyn_per_Mb']], on='sample_id')
mer = mer[(mer.timepoint!='normal') & (~mer.subject_id.isin(UNMATCHED))].copy()

def build_order(sub):
    paired = sub.groupby('subject_id').filter(lambda x: len(x)==2)
    paired_subj = sorted(paired.subject_id.unique())
    single_subj = sorted([s for s in sub.subject_id.unique() if s not in paired_subj])
    order = []
    for s in paired_subj:
        rows = sub[sub.subject_id==s].sort_values('timepoint',
                    key=lambda x: x.map({'pre':0,'post':1}))
        order.extend(rows.sample_id.tolist())
    for s in single_subj:
        order.extend(sub[sub.subject_id==s].sample_id.tolist())
    return order, len(paired_subj)*2

bad_order, bad_paired_n = build_order(mer[mer.response_bin=='bad'])
good_order, good_paired_n = build_order(mer[mer.response_bin=='good'])
order = bad_order + good_order
df = mer.set_index('sample_id').loc[order].reset_index()
n = len(df)
sep_idx     = len(bad_order)                # boundary between poor and good
bad_sub_idx = bad_paired_n                  # within-poor paired→single boundary
good_sub_idx = len(bad_order) + good_paired_n

# ---------- presentation ----------
prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(8.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
slide = prs.slides.add_slide(BLANK)

# Title
nph.add_textbox(slide, 'Sample-level MSI and TMB waterfall — Poor (left) vs Good (right); paired PR/PO adjacent',
                Inches(0.5), Inches(0.12), Inches(15.0), Inches(0.4),
                size=12, bold=True, color=SEP,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Panel letter
nph.add_textbox(slide, 'D', Inches(0.10), Inches(0.10), Inches(0.4), Inches(0.4),
                size=22, bold=True, color=SEP)

# ---------- layout frame ----------
LEFT_MARGIN  = Inches(1.30)
RIGHT_MARGIN = Inches(1.25)
PLOT_X0 = LEFT_MARGIN
PLOT_W  = SW - LEFT_MARGIN - RIGHT_MARGIN
COL_W   = PLOT_W // n
STRIP_H = Inches(0.18)
RESP_Y  = Inches(0.65)
TIME_Y  = RESP_Y + STRIP_H + Inches(0.02)
PLOT_Y0 = TIME_Y + STRIP_H + Inches(0.08)
PLOT_H  = Inches(5.0)
LABEL_Y0 = PLOT_Y0 + PLOT_H + Inches(0.05)

# ---------- strip headers ----------
nph.add_textbox(slide, 'Response',
                Inches(0.3), RESP_Y, LEFT_MARGIN - Inches(0.35), STRIP_H,
                size=9, bold=True, color=SEP,
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
nph.add_textbox(slide, 'Timepoint',
                Inches(0.3), TIME_Y, LEFT_MARGIN - Inches(0.35), STRIP_H,
                size=9, bold=True, color=SEP,
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# ---------- strips (rectangles) ----------
for i, row in df.iterrows():
    x0 = PLOT_X0 + i*COL_W
    resp_col = BAD if row.response_bin == 'bad' else GOOD
    tp_col   = PRE if row.timepoint  == 'pre'  else POST
    for (y, col) in [(RESP_Y, resp_col), (TIME_Y, tp_col)]:
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y, COL_W, STRIP_H)
        r.fill.solid(); r.fill.fore_color.rgb = RGBColor(*col)
        r.line.color.rgb = RGBColor(255,255,255); r.line.width = Pt(0.3)
        r.shadow.inherit = False

# ---------- main axes ----------
msi_max = 0.25  # y-max
tmb_max = 11    # right-axis max
# Left spine + bottom spine
nph.add_line(slide, PLOT_X0, PLOT_Y0, PLOT_X0, PLOT_Y0+PLOT_H,
             color=SEP, width_pt=0.9)  # left MSI spine
nph.add_line(slide, PLOT_X0, PLOT_Y0+PLOT_H, PLOT_X0+PLOT_W, PLOT_Y0+PLOT_H,
             color=SEP, width_pt=0.9)  # bottom
# Right TMB spine
nph.add_line(slide, PLOT_X0+PLOT_W, PLOT_Y0, PLOT_X0+PLOT_W, PLOT_Y0+PLOT_H,
             color=TMB_COLOR, width_pt=0.9)

# Left Y ticks (MSI %)
def msi_y(v): return PLOT_Y0 + PLOT_H - int(PLOT_H * v / msi_max)
for tv in [0, 0.05, 0.10, 0.15, 0.20, 0.25]:
    ty = msi_y(tv)
    nph.add_line(slide, PLOT_X0 - Emu(45000), ty, PLOT_X0, ty,
                 color=SEP, width_pt=0.9)
    nph.add_textbox(slide, f'{tv:.2f}',
                    PLOT_X0 - Inches(0.55), ty - Inches(0.12),
                    Inches(0.50), Inches(0.24),
                    size=9, align=PP_ALIGN.RIGHT, color=SEP)
# Left Y label (rotated) — center x at 0.35", unrotated box 1.5×0.3
nph.add_textbox(slide, 'MSI percentage (%)',
                Inches(0.35) - Inches(0.75), PLOT_Y0 + PLOT_H//2 - Inches(0.15),
                Inches(1.5), Inches(0.3),
                size=10.5, bold=True, color=SEP,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, rotation=-90)

# Right Y ticks (TMB)
def tmb_y(v): return PLOT_Y0 + PLOT_H - int(PLOT_H * v / tmb_max)
for tv in [0, 2, 4, 6, 8, 10]:
    ty = tmb_y(tv)
    nph.add_line(slide, PLOT_X0+PLOT_W, ty, PLOT_X0+PLOT_W + Emu(45000), ty,
                 color=TMB_COLOR, width_pt=0.9)
    nph.add_textbox(slide, str(tv),
                    PLOT_X0+PLOT_W + Inches(0.06), ty - Inches(0.12),
                    Inches(0.40), Inches(0.24),
                    size=9, align=PP_ALIGN.LEFT, color=TMB_COLOR)
# Right Y label (rotated) — center x at SW-0.35, unrotated box 2.0×0.3
nph.add_textbox(slide, 'TMB (nonsyn / Mb)',
                SW - Inches(0.35) - Inches(1.0), PLOT_Y0 + PLOT_H//2 - Inches(0.15),
                Inches(2.0), Inches(0.3),
                size=10.5, bold=True, color=TMB_COLOR,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, rotation=90)

# TMB 10/Mb dashed cutoff
cy = tmb_y(10)
nph.add_line(slide, PLOT_X0, cy, PLOT_X0+PLOT_W, cy,
             color=TMB_COLOR, width_pt=0.8, dash=True)
nph.add_textbox(slide, 'TMB 10/Mb cutoff',
                PLOT_X0 + Inches(0.05), cy - Inches(0.25),
                Inches(1.6), Inches(0.22),
                size=8.5, color=TMB_COLOR,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# ---------- MSI bars + annotations ----------
MIN_BAR_FRAC = 0.012  # visual floor
for i, row in df.iterrows():
    x0 = PLOT_X0 + i*COL_W + Emu(10000)
    bar_w = COL_W - Emu(20000)
    real_v = float(row.MSI_pct) if pd.notna(row.MSI_pct) else 0.0
    v_plot = max(real_v / msi_max, MIN_BAR_FRAC)
    bar_h = int(PLOT_H * v_plot)
    y0 = PLOT_Y0 + PLOT_H - bar_h
    col = BAD if row.response_bin=='bad' else GOOD
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y0, bar_w, bar_h)
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor(*col)
    r.line.color.rgb = RGBColor(255,255,255); r.line.width = Pt(0.25)
    r.shadow.inherit = False
    if real_v > 0.005:
        nph.add_textbox(slide, f'{real_v:.2f}',
                        x0 - Inches(0.15), y0 - Inches(0.25),
                        Inches(0.45), Inches(0.22),
                        size=7.5, color=SEP,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---------- TMB markers ----------
MARKER_D = Emu(110000)
for i, row in df.iterrows():
    cx = PLOT_X0 + i*COL_W + COL_W//2
    cy_m = tmb_y(float(row.TMB_nonsyn_per_Mb))
    m = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               cx - MARKER_D//2, cy_m - MARKER_D//2,
                               MARKER_D, MARKER_D)
    m.fill.solid(); m.fill.fore_color.rgb = RGBColor(*TMB_COLOR)
    m.line.color.rgb = RGBColor(255,255,255); m.line.width = Pt(0.5)
    m.shadow.inherit = False

# ---------- separators ----------
# Main group boundary (solid)
sx = PLOT_X0 + sep_idx*COL_W
nph.add_line(slide, sx, PLOT_Y0, sx, PLOT_Y0+PLOT_H,
             color=SEP, width_pt=1.8)
# Sub-boundaries (dotted)
for sub_i in [bad_sub_idx, good_sub_idx]:
    sx = PLOT_X0 + sub_i*COL_W
    nph.add_line(slide, sx, PLOT_Y0, sx, PLOT_Y0+PLOT_H,
                 color=SEP, width_pt=0.8, dash=True)

# ---------- group labels (Poor / Good) ----------
def group_label(cx_idx, label, color):
    cx = PLOT_X0 + cx_idx*COL_W
    tb_w = Inches(1.75)
    tb = nph.add_textbox(slide, label,
                         cx - tb_w//2, PLOT_Y0 + Inches(0.12),
                         tb_w, Inches(0.34),
                         size=11, bold=True, color=color,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                         fill=(255,255,255), border=color)
    return tb

group_label(sep_idx/2, 'Poor (TRG 2–3)', BAD)
group_label((sep_idx + n)/2, 'Good (TRG 0–1)', GOOD)

# ---------- MSI-H callout ----------
callout_y = PLOT_Y0 + Inches(1.55)
nph.add_textbox(slide,
    'MSI-H clinical cutoff = 20 %   (all samples shown ≤ 0.19 %; axis rescaled)',
    PLOT_X0 + PLOT_W//2 - Inches(3.5), callout_y,
    Inches(7.0), Inches(0.35),
    size=10, bold=True, color=CUTOFF_RED,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    fill=(255,255,255), border=CUTOFF_RED)

# ---------- sample labels (rotated 90°) ----------
for i, row in df.iterrows():
    cx = PLOT_X0 + i*COL_W + COL_W//2
    nph.add_textbox(slide, row.sample_id,
                    cx - Inches(0.28), LABEL_Y0,
                    Inches(0.56), Inches(0.6),
                    size=8.5, color=SEP,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    rotation=-90)

# ---------- legend ----------
LEG_Y = LABEL_Y0 + Inches(0.7)
LEG_X = Inches(0.7)
SW_W  = Inches(0.20); SW_H = Inches(0.20)
def legend_item(x, y, fill, label, shape='rect'):
    if shape == 'rect':
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Emu(40000), SW_W, SW_H)
    elif shape == 'oval':
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Emu(40000), SW_W, SW_H)
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(*fill)
    s.line.color.rgb = RGBColor(255,255,255); s.line.width = Pt(0.4)
    s.shadow.inherit = False
    nph.add_textbox(slide, label,
                    x + SW_W + Inches(0.05), y,
                    Inches(1.9), Inches(0.32),
                    size=9.5, color=SEP,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

lx = LEG_X
for fill, lbl in [(GOOD,'Good (TRG 0–1)'), (BAD,'Poor (TRG 2–3)'),
                  (PRE,'Pre-CRT'), (POST,'Post-CRT')]:
    legend_item(lx, LEG_Y, fill, lbl, 'rect')
    lx += Inches(1.80)
# TMB marker
legend_item(lx, LEG_Y, TMB_COLOR, 'TMB (right axis)', 'oval')
lx += Inches(1.90)
# Group boundary line + label
line_y = LEG_Y + Inches(0.18)
nph.add_line(slide, lx, line_y, lx + Inches(0.25), line_y,
             color=SEP, width_pt=1.8)
nph.add_textbox(slide, 'Poor | Good boundary',
                lx + Inches(0.28), LEG_Y, Inches(1.9), Inches(0.32),
                size=9.5, color=SEP,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
lx += Inches(2.20)
nph.add_line(slide, lx, line_y, lx + Inches(0.25), line_y,
             color=SEP, width_pt=0.9, dash=True)
nph.add_textbox(slide, 'Paired | single boundary',
                lx + Inches(0.28), LEG_Y, Inches(1.9), Inches(0.32),
                size=9.5, color=SEP,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT_FILE)
print(f'Saved native editable PPT: {OUT_FILE}')
print(f'  41 matched tumors: bad={len(bad_order)} (paired {bad_paired_n}, single {len(bad_order)-bad_paired_n}),'
      f' good={len(good_order)} (paired {good_paired_n}, single {len(good_order)-good_paired_n})')
