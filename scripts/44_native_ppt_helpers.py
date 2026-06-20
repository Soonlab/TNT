"""Helper library for building native PowerPoint figures with python-pptx.

Provides reusable primitives (axes, ticks, lines, markers, CI bars, legend boxes)
so we can reconstruct forest plots, boxplots, and similar structured figures
with every element as a native PPT object — fully editable.

All drawing functions take a slide and a 'canvas' dict describing the plot area
in EMU coordinates. Data coordinates are mapped to EMU via the canvas mapping.
"""
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Coordinate helpers ----------
def make_canvas(x_emu, y_emu, w_emu, h_emu, xlim, ylim, invert_y=False):
    """Create a canvas mapping with optional inverted Y (for forest plots with rows)."""
    return dict(x0=x_emu, y0=y_emu, w=w_emu, h=h_emu,
                xmin=xlim[0], xmax=xlim[1],
                ymin=ylim[0], ymax=ylim[1], invert_y=invert_y)

def x_to_emu(c, x):
    t = (x - c['xmin']) / (c['xmax'] - c['xmin'])
    return c['x0'] + int(t * c['w'])

def y_to_emu(c, y):
    t = (y - c['ymin']) / (c['ymax'] - c['ymin'])
    if c.get('invert_y', False):
        t = 1.0 - t
    return c['y0'] + int(t * c['h'])

# ---------- Drawing primitives ----------
def add_line(slide, x1, y1, x2, y2, color=(0,0,0), width_pt=0.8, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = RGBColor(*color)
    conn.line.width = Pt(width_pt)
    if dash:
        try:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass
    return conn

def add_marker(slide, cx, cy, size_emu, color=(0,0,0), shape=MSO_SHAPE.OVAL, edge=(0,0,0), edge_w=0.7):
    half = size_emu // 2
    s = slide.shapes.add_shape(shape, cx-half, cy-half, size_emu, size_emu)
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(*color)
    s.line.color.rgb = RGBColor(*edge); s.line.width = Pt(edge_w)
    s.shadow.inherit = False
    return s

def add_rect(slide, x, y, w, h, fill=None, edge=(0,0,0), edge_w=0.7):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None:
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor(*fill)
    else:
        s.fill.background()
    s.line.color.rgb = RGBColor(*edge); s.line.width = Pt(edge_w)
    s.shadow.inherit = False
    return s

def add_diamond(slide, cx, cy, w_emu, h_emu, fill=(0xE6,0x39,0x46), edge=(0,0,0), edge_w=0.7):
    s = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, cx - w_emu//2, cy - h_emu//2, w_emu, h_emu)
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(*fill)
    s.line.color.rgb = RGBColor(*edge); s.line.width = Pt(edge_w)
    s.shadow.inherit = False
    return s

def add_textbox(slide, text, x, y, w, h, size=10, bold=False, italic=False,
                color=(0x33,0x33,0x33), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                font='Arial', fill=None, border=None, rotation=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = RGBColor(*color)
        r.font.name = font
    if fill is not None:
        tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor(*fill)
    if border is not None:
        tb.line.color.rgb = RGBColor(*border); tb.line.width = Pt(0.75)
    if rotation:
        tb.rotation = rotation
    return tb

# ---------- High-level components ----------
def draw_axis(slide, canvas, x_ticks=None, y_ticks=None,
              x_label='', y_label='',
              show_top=False, show_right=False,
              tick_size_pt=8, label_size_pt=10, spine_w=0.8,
              show_x_tick_labels=True, show_y_tick_labels=True,
              x_tick_formatter=None, y_tick_formatter=None,
              x_tick_rotation=0):
    """Draw axes with editable tick labels and axis labels."""
    c = canvas
    # Spines
    add_line(slide, c['x0'], c['y0']+c['h'], c['x0']+c['w'], c['y0']+c['h'], width_pt=spine_w)  # bottom
    add_line(slide, c['x0'], c['y0'],        c['x0'],         c['y0']+c['h'], width_pt=spine_w)  # left
    if show_top:
        add_line(slide, c['x0'], c['y0'], c['x0']+c['w'], c['y0'], width_pt=spine_w)
    if show_right:
        add_line(slide, c['x0']+c['w'], c['y0'], c['x0']+c['w'], c['y0']+c['h'], width_pt=spine_w)
    # x ticks
    tick_len = Emu(30000)
    if x_ticks:
        for t in x_ticks:
            xe = x_to_emu(c, t)
            add_line(slide, xe, c['y0']+c['h'], xe, c['y0']+c['h']+tick_len, width_pt=spine_w)
            if show_x_tick_labels:
                label = (x_tick_formatter(t) if x_tick_formatter else
                         (f'{t:g}' if abs(t) < 1e6 else f'{t:.1e}'))
                add_textbox(slide, label,
                            xe - Emu(240000), c['y0']+c['h']+Emu(40000),
                            Emu(480000), Emu(180000),
                            size=tick_size_pt, align=PP_ALIGN.CENTER, rotation=x_tick_rotation)
    if y_ticks:
        for t in y_ticks:
            ye = y_to_emu(c, t)
            add_line(slide, c['x0'], ye, c['x0']-tick_len, ye, width_pt=spine_w)
            if show_y_tick_labels:
                label = (y_tick_formatter(t) if y_tick_formatter else
                         (f'{t:g}' if abs(t) < 1e6 else f'{t:.1e}'))
                add_textbox(slide, label,
                            c['x0']-Emu(650000), ye-Emu(100000),
                            Emu(560000), Emu(200000),
                            size=tick_size_pt, align=PP_ALIGN.RIGHT)
    # axis labels
    if x_label:
        add_textbox(slide, x_label,
                    c['x0'], c['y0']+c['h']+Emu(260000),
                    c['w'], Emu(260000),
                    size=label_size_pt, bold=True, align=PP_ALIGN.CENTER)
    if y_label:
        add_textbox(slide, y_label,
                    c['x0']-Emu(1050000), c['y0'] + c['h']//2 - Emu(500000),
                    Emu(1000000), Emu(1000000),
                    size=label_size_pt, bold=True, align=PP_ALIGN.CENTER, rotation=-90)

def draw_vline(slide, canvas, x, color=(0,0,0), width_pt=0.6, dash=False):
    xe = x_to_emu(canvas, x)
    add_line(slide, xe, canvas['y0'], xe, canvas['y0']+canvas['h'],
             color=color, width_pt=width_pt, dash=dash)

def draw_hline(slide, canvas, y, color=(0,0,0), width_pt=0.6, dash=False):
    ye = y_to_emu(canvas, y)
    add_line(slide, canvas['x0'], ye, canvas['x0']+canvas['w'], ye,
             color=color, width_pt=width_pt, dash=dash)

def draw_forest_row(slide, canvas, y_data, med, lo, hi,
                    label='', right_label='',
                    color=(0x2E,0x86,0xAB), marker_size=Emu(180000),
                    ci_width_pt=1.4, label_size=10, color_right=None):
    """Draw one forest-plot row: horizontal CI bar + central marker + row labels."""
    c = canvas
    ye = y_to_emu(c, y_data)
    x_lo = x_to_emu(c, lo); x_hi = x_to_emu(c, hi); x_med = x_to_emu(c, med)
    # CI line
    add_line(slide, x_lo, ye, x_hi, ye, color=color, width_pt=ci_width_pt)
    # small caps at ends
    cap = Emu(50000)
    add_line(slide, x_lo, ye-cap, x_lo, ye+cap, color=color, width_pt=ci_width_pt)
    add_line(slide, x_hi, ye-cap, x_hi, ye+cap, color=color, width_pt=ci_width_pt)
    # marker
    add_marker(slide, x_med, ye, marker_size, color=color, edge=(0,0,0), edge_w=0.6)
    # left row label
    if label:
        add_textbox(slide, label,
                    c['x0']-Emu(1400000), ye-Emu(140000),
                    Emu(1300000), Emu(280000),
                    size=label_size, align=PP_ALIGN.RIGHT)
    # right row label (annotation)
    if right_label:
        add_textbox(slide, right_label,
                    c['x0']+c['w']+Emu(120000), ye-Emu(140000),
                    Emu(3000000), Emu(280000),
                    size=label_size, color=color_right or (0x55,0x55,0x55))
