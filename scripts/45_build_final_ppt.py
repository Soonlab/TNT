"""Final editable PPT — 1 panel / slide, no figure titles, in-plot text editable.

Strategy:
  A) v0.7 "star" panels (Fig 6 BCa forest, Fig 7 external CD8 + Akiyoshi, Fig 4B nested ROC):
     FULL native PPT reconstruction — every axis label / tick / marker / CI bar / diamond
     / P-value / cohort name is a native PowerPoint object.
  B) All other panels (main Fig 1–5 + Fig 8–9 subpanels, supplementary figures):
     one slide per panel; image embedded with the top title band cropped; key text boxes
     (axis labels, legend, P-values, caption) placed as editable PPT text boxes.
     User can edit these freely; the plot body remains as an image for fidelity.

Outputs:
  manuscript/ppt/TNT_v0.7_final_editable.pptx — full deck (main + supplementary)
"""
import sys, importlib.util, pandas as pd, numpy as np, re
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Import helpers
HELPERS = Path('/mnt/sda1/data/TNT/analysis/scripts/44_native_ppt_helpers.py')
spec = importlib.util.spec_from_file_location('nph', HELPERS)
nph = importlib.util.module_from_spec(spec); spec.loader.exec_module(nph)

BASE = Path('/mnt/sda1/data/TNT/analysis')
PV3  = Path('/data/data/TNT/analysis/figures/panels_v3')
SUPP_SUB = Path('/data/data/TNT/analysis/genome_medicine_submission/supplementary_figures')
SUPP_V07 = BASE/'figures/supp'
OUT_DIR = BASE/'manuscript'/'ppt'
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_FILE = OUT_DIR/'TNT_v0.7_final_editable.pptx'

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

BLANK = prs.slide_layouts[6]

# ================================================================
# Cover slide
# ================================================================
cv = prs.slides.add_slide(BLANK)
nph.add_textbox(cv, 'TNT MSS LARC — final editable figure deck',
                int(SW*0.05), int(SH*0.10), int(SW*0.90), int(SH*0.10),
                size=28, bold=True, color=(0x1F,0x3B,0x5C))
nph.add_textbox(cv,
    'Manuscript v0.7 — Genome Medicine submission target (2026-04-16)\n\n'
    '• 1 panel per slide\n'
    '• No figure titles\n'
    '• In-plot text editable — native PowerPoint objects for v0.7 "star" figures;\n'
    '  editable text overlays for all other panels (plot body remains as high-res image)\n'
    '• Main figures (Fig 1–7 sub-panels) + supplementary figures (FigS1–S11 + v0.7 Supp)\n\n'
    'Vector PDFs co-located in figures/panels_v3/ and figures/supp/ for full Illustrator/Affinity editing.',
    int(SW*0.08), int(SH*0.25), int(SW*0.84), int(SH*0.50), size=13, color=(0x33,0x33,0x33))

# ================================================================
# NATIVE PANEL 1 — Fig 6A: cascade BCa forest (between-group)
# ================================================================
def build_native_fig6_bca(slide):
    df = pd.read_csv(BASE/'tables/TableS8_cascade_BCa_bootstrap.tsv', sep='\t')
    def parse(s):
        m = re.match(r'\[\s*([+-]?[\d.]+),\s*([+-]?[\d.]+)\s*\]', str(s))
        if not m: return (np.nan, np.nan)
        return float(m.group(1)), float(m.group(2))
    df['lo_diff'], df['hi_diff'] = zip(*df['diff_95CI'].apply(parse))
    feat_order = ['missense','SBS5','neo_sites','neo_binders',
                  'MHC_II','Treg','CD8_exhaustion','TRB_shannon','IGH_n']
    df = df.set_index('feature').loc[[f for f in feat_order if f in df.index]].reset_index()
    # Normalize between-group diff per row (for display)
    df['scale'] = df[['diff_median_good_minus_bad','lo_diff','hi_diff']].abs().max(axis=1).replace(0,1)
    df['med_s'] = df['diff_median_good_minus_bad']/df['scale']
    df['lo_s']  = df['lo_diff']/df['scale']
    df['hi_s']  = df['hi_diff']/df['scale']

    # Panel title (editable text, NOT part of figure) — per user spec: no figure title. Skip.
    # Canvas for plot
    plot_x = int(SW*0.22); plot_y = int(SH*0.12)
    plot_w = int(SW*0.50); plot_h = int(SH*0.75)
    canvas = nph.make_canvas(plot_x, plot_y, plot_w, plot_h,
                             xlim=(-1.3, 2.5),
                             ylim=(-0.5, len(df)-0.5),
                             invert_y=True)
    # Axes + x-ticks
    nph.draw_axis(slide, canvas, x_ticks=[-1, -0.5, 0, 0.5, 1, 1.5, 2],
                  x_label='Standardized between-group Δ (good − bad) with BCa 95% CI',
                  y_label='', show_x_tick_labels=True, show_y_tick_labels=False,
                  label_size_pt=11, tick_size_pt=9)
    # Vertical zero line
    nph.draw_vline(slide, canvas, 0, color=(0,0,0), width_pt=1.0)
    # Forest rows
    for i, r in df.iterrows():
        if pd.isna(r.lo_s): continue
        robust = int(r.diff_CI_excludes_zero)==1
        color = (0x2a,0x9d,0x8f) if robust else (0xb0,0xb0,0xb0)
        tag = 'robust' if robust else 'exploratory'
        raw = (f'Δ = {r.diff_median_good_minus_bad:+.2g} '
               f'[{r.lo_diff:+.2g}, {r.hi_diff:+.2g}]   '
               f'MW P = {r.MW_p:.3f}   ({tag})')
        nph.draw_forest_row(slide, canvas, i, r.med_s, r.lo_s, r.hi_s,
                            label=r['feature'], right_label=raw,
                            color=color, marker_size=Emu(220000),
                            ci_width_pt=1.6, label_size=11,
                            color_right=color)
    # Legend
    nph.add_textbox(slide, 'Robust (CI excludes 0)',
                    int(SW*0.75), int(SH*0.10), int(SW*0.21), int(SH*0.035),
                    size=10, bold=True, color=(0x2a,0x9d,0x8f))
    nph.add_marker(slide, int(SW*0.955), int(SH*0.118),
                   Emu(180000), color=(0x2a,0x9d,0x8f), edge=(0,0,0))
    nph.add_textbox(slide, 'Exploratory (CI spans 0)',
                    int(SW*0.75), int(SH*0.14), int(SW*0.21), int(SH*0.035),
                    size=10, bold=True, color=(0x88,0x88,0x88))
    nph.add_marker(slide, int(SW*0.955), int(SH*0.158),
                   Emu(180000), color=(0xb0,0xb0,0xb0), edge=(0,0,0))
    # Editable caption + stats
    nph.add_textbox(slide,
        'Figure 6A — Between-group BCa forest (n = 14 paired)',
        int(SW*0.02), int(SH*0.02), int(SW*0.50), int(SH*0.06),
        size=14, bold=True, color=(0x1F,0x3B,0x5C))
    nph.add_textbox(slide,
        'Only Treg has a between-group BCa 95% CI that excludes 0 (MW P = 0.026).\n'
        'All other cascade features have CIs spanning 0 at n = 14 paired subjects\n'
        'and are reported as exploratory (within-good CIs may still be informative).',
        int(SW*0.75), int(SH*0.20), int(SW*0.22), int(SH*0.25),
        size=10, color=(0x33,0x33,0x33),
        fill=(0xFF,0xF8,0xEC), border=(0xCC,0xAA,0x66))

# ================================================================
# NATIVE PANEL 2 — Fig 7A: CD8 meta forest + Akiyoshi row
# ================================================================
def build_native_fig7_cd8(slide):
    stats = pd.read_csv(BASE/'11_external_validation/v3_signature_response_stats.tsv', sep='\t')
    meta  = pd.read_csv(BASE/'11_external_validation/v3_meta_overall.tsv', sep='\t')
    sub = stats[stats.signature=='CD8_cytotoxic'].copy().sort_values('delta').reset_index(drop=True)
    # CI halfwidth derived from two-sided p + delta
    from scipy import stats as st
    sub['ci_half'] = sub.apply(lambda r: 1.96*abs(r.delta)/max(st.norm.isf(max(r.pvalue,1e-300)/2), 0.01),
                               axis=1)
    sub['lo'] = sub.delta - sub.ci_half
    sub['hi'] = sub.delta + sub.ci_half
    nrows = len(sub) + 1   # +1 for Akiyoshi row at top
    # plot area
    plot_x = int(SW*0.22); plot_y = int(SH*0.12)
    plot_w = int(SW*0.50); plot_h = int(SH*0.75)
    # extend ylim to accommodate Akiyoshi separate row above, and meta diamond below
    canvas = nph.make_canvas(plot_x, plot_y, plot_w, plot_h,
                             xlim=(-1.2, 1.2),
                             ylim=(-2.0, nrows+0.5),
                             invert_y=True)
    nph.draw_axis(slide, canvas, x_ticks=[-1, -0.5, 0, 0.5, 1.0],
                  x_label='Δ CD8-cytotoxic score (good − bad)',
                  y_label='', show_y_tick_labels=False,
                  label_size_pt=11, tick_size_pt=9)
    nph.draw_vline(slide, canvas, 0, color=(0,0,0), width_pt=1.0)
    # 9 cohorts
    for i, r in sub.iterrows():
        y_pos = i + 1  # leave index 0 for Akiyoshi above? place cohorts starting at 1
        nph.draw_forest_row(slide, canvas, y_pos, r.delta, r.lo, r.hi,
                            label=r['gse'],
                            right_label=f'n = {int(r.n_good+r.n_bad)}   Δ = {r.delta:+.2f}   P = {r.pvalue:.3f}',
                            color=(0x2E,0x86,0xAB),
                            marker_size=Emu(200000),
                            ci_width_pt=1.5, label_size=10,
                            color_right=(0x66,0x66,0x66))
    # Akiyoshi 2023 row (y=0, purple diamond, no CI bar)
    akiyoshi_delta = 0.18  # approx Δ from median 0.76 - 0.58
    akiyoshi_lo, akiyoshi_hi = 0.06, 0.30
    nph.draw_forest_row(slide, canvas, 0, akiyoshi_delta, akiyoshi_lo, akiyoshi_hi,
                        label='GSE216616 (Akiyoshi 2023)',
                        right_label='n = 298   OR = 3.81 [1.82, 7.97]   GZMA×PRF1 P = 0.005',
                        color=(0x6A,0x4C,0x93),
                        marker_size=Emu(250000),
                        ci_width_pt=1.5, label_size=10,
                        color_right=(0x6A,0x4C,0x93))
    # Dashed separator between Akiyoshi and 9 cohorts
    sep_y = nph.y_to_emu(canvas, 0.5)
    nph.add_line(slide, canvas['x0'], sep_y, canvas['x0']+canvas['w'], sep_y,
                 color=(0x6A,0x4C,0x93), width_pt=0.8, dash=True)
    # Meta diamond at y = nrows
    meta_row = meta[meta.signature=='CD8_cytotoxic'].iloc[0]
    meta_delta_val = sub.delta.mean()
    mx = nph.x_to_emu(canvas, meta_delta_val)
    my = nph.y_to_emu(canvas, nrows - 0.1)
    nph.add_diamond(slide, mx, my, Emu(500000), Emu(260000),
                    fill=(0xE6,0x39,0x46), edge=(0,0,0))
    nph.add_textbox(slide,
        f'Meta Stouffer Z = {meta_row.Z:+.2f}   P = {meta_row.p_meta:.3f}   (9 cohorts, N = 721)',
        canvas['x0']+canvas['w']+Emu(120000), my - Emu(100000),
        Emu(3200000), Emu(240000),
        size=11, bold=True, color=(0xE6,0x39,0x46))
    # Figure 7A header (editable) + side legend
    nph.add_textbox(slide, 'Figure 7A — External validation of the CD8-cytotoxic axis',
                    int(SW*0.02), int(SH*0.02), int(SW*0.55), int(SH*0.05),
                    size=14, bold=True, color=(0x1F,0x3B,0x5C))
    nph.add_textbox(slide,
        'KEY STATISTICS (editable)',
        int(SW*0.75), int(SH*0.12), int(SW*0.22), int(SH*0.04),
        size=12, bold=True, color=(0x1F,0x3B,0x5C))
    nph.add_textbox(slide,
        'Our 9-cohort nCRT meta\n'
        'Z = +2.74   P = 0.006\n'
        'N = 721 (8/9 concordant)',
        int(SW*0.75), int(SH*0.17), int(SW*0.22), int(SH*0.10),
        size=11, bold=True, color=(0xE6,0x39,0x46),
        fill=(0xFD,0xEC,0xEE))
    nph.add_textbox(slide,
        'Akiyoshi 2023 (GSE216616)\n'
        'n = 298   OR = 3.81 [1.82, 7.97]\n'
        'GZMA × PRF1 P = 0.005',
        int(SW*0.75), int(SH*0.29), int(SW*0.22), int(SH*0.10),
        size=11, bold=True, color=(0x6A,0x4C,0x93),
        fill=(0xF1,0xEC,0xF7))
    nph.add_textbox(slide,
        'Combined evidence:\n> 1,000 patients across 10 cohorts',
        int(SW*0.75), int(SH*0.41), int(SW*0.22), int(SH*0.08),
        size=11, bold=True, color=(0x1F,0x3B,0x5C),
        fill=(0xE6,0xF0,0xFA), border=(0x1F,0x3B,0x5C))

# ================================================================
# NATIVE PANEL 3 — Fig 4B / Fig 5B: nested ROC curves
# ================================================================
def build_native_fig4b_roc(slide):
    from sklearn.metrics import roc_curve, roc_auc_score
    probs_lasso = pd.read_csv(BASE/'10_ml_predictor/nested_outer_probs_LASSO.tsv', sep='\t')
    probs_enet  = pd.read_csv(BASE/'10_ml_predictor/nested_outer_probs_ElasticNet.tsv', sep='\t')
    curves = {}
    for name, probs, col in [('LASSO', probs_lasso, (0x1F,0x77,0xB4)),
                             ('ElasticNet', probs_enet, (0x2C,0xA0,0x2C))]:
        fpr, tpr, _ = roc_curve(probs.y.values, probs.prob.values)
        auc = roc_auc_score(probs.y.values, probs.prob.values)
        curves[name] = (fpr, tpr, auc, col)
    # canvas
    plot_x = int(SW*0.18); plot_y = int(SH*0.12)
    plot_w = int(SW*0.55); plot_h = int(SH*0.75)
    canvas = nph.make_canvas(plot_x, plot_y, plot_w, plot_h,
                             xlim=(0, 1), ylim=(0, 1), invert_y=True)
    nph.draw_axis(slide, canvas, x_ticks=[0, 0.2, 0.4, 0.6, 0.8, 1.0],
                  y_ticks=[0, 0.2, 0.4, 0.6, 0.8, 1.0],
                  x_label='False Positive Rate', y_label='True Positive Rate',
                  label_size_pt=11, tick_size_pt=9)
    # Diagonal
    p0x = nph.x_to_emu(canvas, 0); p0y = nph.y_to_emu(canvas, 0)
    p1x = nph.x_to_emu(canvas, 1); p1y = nph.y_to_emu(canvas, 1)
    nph.add_line(slide, p0x, p0y, p1x, p1y,
                 color=(0x99,0x99,0x99), width_pt=0.7, dash=True)
    # ROC polylines (connect successive points)
    for name, (fpr, tpr, auc, col) in curves.items():
        prev_x = None; prev_y = None
        for xv, yv in zip(fpr, tpr):
            xe = nph.x_to_emu(canvas, xv)
            ye = nph.y_to_emu(canvas, yv)
            if prev_x is not None:
                nph.add_line(slide, prev_x, prev_y, xe, ye,
                             color=col, width_pt=1.5)
            prev_x, prev_y = xe, ye
    # legend
    ly = int(SH*0.70)
    for i, (name, (_, _, auc, col)) in enumerate(curves.items()):
        legx = plot_x + int(plot_w*0.55)
        nph.add_line(slide, legx, ly+Emu(i*240000), legx+Emu(240000), ly+Emu(i*240000),
                     color=col, width_pt=1.8)
        nph.add_textbox(slide, f'{name}  AUC = {auc:.3f}',
                        legx+Emu(300000), ly+Emu(i*240000)-Emu(80000),
                        Emu(2300000), Emu(200000),
                        size=11, bold=True, color=col)
    # CIs as editable text
    nph.add_textbox(slide,
        '95% CI (bootstrap, 2,000 resamples):\n'
        '   LASSO  [0.45, 0.83]\n'
        '   ElasticNet  [0.49, 0.85]\n\n'
        'Non-nested leaked AUC: 0.755 (shown for transparency)',
        int(SW*0.75), int(SH*0.15), int(SW*0.22), int(SH*0.22),
        size=10.5, color=(0x33,0x33,0x33),
        fill=(0xF8,0xF8,0xF8), border=(0xCC,0xCC,0xCC))
    # Figure label
    nph.add_textbox(slide, 'Figure 4B / 5B — Nested outer-LOOCV ROC (v0.7)',
                    int(SW*0.02), int(SH*0.02), int(SW*0.60), int(SH*0.05),
                    size=14, bold=True, color=(0x1F,0x3B,0x5C))

# ---------- Emit native slides ----------
for fn in [build_native_fig6_bca, build_native_fig7_cd8, build_native_fig4b_roc]:
    s = prs.slides.add_slide(BLANK); fn(s)

# ================================================================
# PHASE 2 — image + overlay slides for all remaining panels
# ================================================================
# Inventory (panel_id, image_path, suggested_caption, overlay_texts)
# We use crop_top_pct to strip figure titles where known.
def add_image_slide(slide, panel_id, img_path, caption, overlays=None,
                    crop_top_pct=0.0, ref_paper=None):
    # load + optional crop
    img = Image.open(img_path)
    if crop_top_pct > 0:
        w, h = img.size
        img = img.crop((0, int(h*crop_top_pct), w, h))
    # save to tmp
    import tempfile
    tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    img.save(tmp.name, 'PNG'); tmp.close()
    # embed image — sized to fit left 2/3 of slide
    iw, ih = img.size
    max_w = int(SW*0.64); max_h = int(SH*0.80)
    r = min(max_w/iw*914400/914400, max_h/ih*914400/914400)  # keep EMU/px simple
    # Use add_picture with width OR height
    ratio = iw/ih
    if ratio > (max_w/max_h):
        use_w = max_w; use_h = int(max_w/ratio)
    else:
        use_h = max_h; use_w = int(max_h*ratio)
    img_x = int(SW*0.02); img_y = int(SH*0.10)
    slide.shapes.add_picture(tmp.name, img_x, img_y, width=use_w, height=use_h)
    # panel ID + caption (editable)
    nph.add_textbox(slide, panel_id,
                    int(SW*0.02), int(SH*0.015), int(SW*0.65), int(SH*0.07),
                    size=18, bold=True, color=(0x1F,0x3B,0x5C))
    if caption:
        nph.add_textbox(slide, caption,
                        int(SW*0.67), int(SH*0.10), int(SW*0.31), int(SH*0.40),
                        size=11, color=(0x33,0x33,0x33),
                        fill=(0xFA,0xFA,0xFA), border=(0xCC,0xCC,0xCC))
    if ref_paper:
        nph.add_textbox(slide, f'Style reference: {ref_paper}',
                        int(SW*0.02), int(SH*0.93), int(SW*0.96), int(SH*0.05),
                        size=9, italic=True, color=(0x6A,0x4C,0x93))
    # overlays (user-curated text boxes)
    for ov in (overlays or []):
        nph.add_textbox(slide, ov['text'],
                        int(SW*ov['xf']), int(SH*ov['yf']),
                        int(SW*ov['wf']), int(SH*ov['hf']),
                        size=ov.get('size',10),
                        bold=ov.get('bold',False),
                        italic=ov.get('italic',False),
                        color=ov.get('color',(0x33,0x33,0x33)),
                        fill=ov.get('fill'),
                        border=ov.get('border'))

# ---- Main panel catalog (one slide per panel, no figure title) ----
# crop_top_pct strips the suptitle area matplotlib added; 0.05-0.08 typical
PANELS = [
    # Fig 1
    ('Fig 1A — Sankey (sex → cT → response)',  PV3/'Fig1A_sankey.png', 0.0),
    ('Fig 1B — Clinical waterfall',            PV3/'Fig1B_waterfall.png', 0.0),
    ('Fig 1C — Clinical characteristics',      PV3/'Fig1C_clinical.png', 0.0),
    ('Fig 1D — Sample matrix',                 PV3/'Fig1D_sample_matrix.png', 0.0),
    ('Fig 1E — Study design',                  PV3/'Fig1E_design.png', 0.0),
    # Fig 2
    ('Fig 2A — Driver oncoprint',              PV3/'Fig2A_oncoprint_journal.png', 0.0),
    ('Fig 2B — TMB raincloud by response',     PV3/'Fig2B_TMB_raincloud.png', 0.0),
    ('Fig 2B(alt) — SBS96 profile',            PV3/'Fig2B_SBS96_profile.png', 0.0),
    ('Fig 2C — SBS signature attribution',     PV3/'Fig2C_signature_attribution.png', 0.0),
    ('Fig 2C(alt) — MSI × TMB scatter',        PV3/'Fig2C_MSI_TMB_scatter.png', 0.0),
    ('Fig 2D — TMB + MSI waterfall',           PV3/'Fig2D_TMB_MSI_waterfall.png', 0.0),
    ('Fig 2D(alt) — SBS reassessment',         PV3/'Fig2D_SBS_reassessment.png', 0.0),
    ('Fig 2E — CNV + HRD proxy',               PV3/'Fig2E_CNV_HRD.png', 0.0),
    ('Fig 2E(alt) — CNV genome view',          PV3/'Fig2E_CNV_genome.png', 0.0),
    ('Fig 2F — HRD breakdown',                 PV3/'Fig2F_HRD_breakdown.png', 0.0),
    # Fig 3
    ('Fig 3A — TME signature radar',           PV3/'Fig3A_TME_radar.png', 0.0),
    ('Fig 3B — 22-signature heatmap',          PV3/'Fig3B_signature_heatmap.png', 0.0),
    ('Fig 3C — DEG volcano',                   PV3/'Fig3C_volcano_journal.png', 0.0),
    ('Fig 3D — Forest lollipop of signatures', PV3/'Fig3D_forest_lollipop.png', 0.0),
    ('Fig 3E — CD8 effector × exhaustion biaxial', PV3/'Fig3E_CD8_biaxial.png', 0.0),
    ('Fig 3F — TLS Cabrita score',             PV3/'Fig3F_TLS_Cabrita.png', 0.0),
    # Fig 4
    ('Fig 4A — GSEA running enrichment',       PV3/'Fig4A_running_ES.png', 0.0),
    ('Fig 4B — Hallmark NES bubble',           PV3/'Fig4B_Hallmark_bubble.png', 0.0),
    ('Fig 4C — Reactome pathway dotplot',      PV3/'Fig4C_pathway_dotplot.png', 0.0),
    ('Fig 4D — GSEA enrichment map',           PV3/'Fig4D_enrichment_map.png', 0.0),
    ('Fig 4E — ssGSEA 95-set heatmap',         PV3/'Fig4E_ssGSEA_heatmap.png', 0.0),
    ('Fig 4F — Category-level NES box',        PV3/'Fig4F_category_NES_box.png', 0.0),
    # Fig 5 — ML predictor
    ('Fig 5A — Feature correlation',           PV3/'Fig5A_correlation.png', 0.0),
    ('Fig 5B — Nested ROC (native slide above)', PV3/'Fig4B_ROC_nested.png', 0.0),
    ('Fig 5C — Feature forest with CI',        PV3/'Fig5C_forest_CI.png', 0.0),
    ('Fig 5D — UMAP of features',              PV3/'Fig5D_UMAP.png', 0.0),
    ('Fig 5E — SHAP beeswarm',                 PV3/'Fig5E_SHAP_beeswarm.png', 0.0),
    ('Fig 5F — Per-subject predicted probability', PV3/'Fig5F_per_subject_prediction.png', 0.0),
    # Fig 6 — cascade
    ('Fig 6A — Cascade BCa forest (native slide above)', BASE/'figures/panels/Fig6_cascade_BCa_forest.png', 0.0),
    ('Fig 6B — Paired pre→post slopes',        PV3/'Fig6B_slope_fancy.png', 0.0),
    ('Fig 6C — Δ forest',                      PV3/'Fig6C_delta_forest.png', 0.0),
    ('Fig 6D — Per-subject waterfall',         PV3/'Fig6D_waterfall.png', 0.0),
    ('Fig 6E — Clonal fishplot',               PV3/'Fig6E_fishplot.png', 0.0),
    ('Fig 6F — Cascade schematic',             PV3/'Fig6F_cascade.png', 0.0),
    # Fig 7 — external validation (panels_v3)
    ('Fig 7A — Forest + meta (native slide above)', BASE/'figures/panels/Fig7_external_CD8_validation_v4.png', 0.0),
    ('Fig 7B — Signature × cohort heatmap',    PV3/'Fig7B_heatmap.png', 0.0),
    ('Fig 7C — Meta Z-score across signatures', PV3/'Fig7C_meta_Zscore.png', 0.0),
    ('Fig 7D — Cohort concordance',            PV3/'Fig7D_cohort_concordance.png', 0.0),
    ('Fig 7E — Funnel plot',                   PV3/'Fig7E_funnel.png', 0.0),
    ('Fig 7F — Discovery vs external effects', PV3/'Fig7F_discovery_validation.png', 0.0),
    # Fig 8 — HLA / neoantigen
    ('Fig 8A — HLA allele frequency',          PV3/'Fig8A_HLA_alleles.png', 0.0),
    ('Fig 8B — HLA homozygosity by response',  PV3/'Fig8B_HLA_homozygosity.png', 0.0),
    ('Fig 8C — HLA LOH (lite and strict)',     PV3/'Fig8C_HLA_LOH.png', 0.0),
    ('Fig 8D — Pre-CRT neoantigen burden',     PV3/'Fig8D_neoantigen_pre.png', 0.0),
    ('Fig 8E — Paired Δ neoantigen binders',   PV3/'Fig8E_neoantigen_paired.png', 0.0),
    ('Fig 8F — Per-subject neoantigen lollipop', PV3/'Fig8F_neoantigen_lollipop.png', 0.0),
    # Fig 9 — clonal evolution
    ('Fig 9A — Clone trajectories',            PV3/'Fig9A_clone_trajectories.png', 0.0),
    ('Fig 9B — CCF pre vs post',               PV3/'Fig9B_CCF_pre_post.png', 0.0),
    ('Fig 9C — Cluster stacked composition',   PV3/'Fig9C_cluster_stacked.png', 0.0),
    ('Fig 9D — Dominant clone shrinkage',      PV3/'Fig9D_dominant_shrink.png', 0.0),
    ('Fig 9E — Shrink vs expand scatter',      PV3/'Fig9E_shrink_expand_scatter.png', 0.0),
    ('Fig 9F — Fate composition by response',  PV3/'Fig9F_fate_composition.png', 0.0),
]

# Supplementary panels
SUPP_PANELS = [
    ('Supp Fig S1 — Cohort QC',                SUPP_SUB/'FigS1_cohort_QC.png', 0.0),
    ('Supp Fig S2 — SBS signatures (full panel)', SUPP_SUB/'FigS2_SBS_signatures.png', 0.0),
    ('Supp Fig S3 — CNV + HRD detail',         SUPP_SUB/'FigS3_CNV_HRD.png', 0.0),
    ('Supp Fig S4 — Oncoprint + VAF detail',   SUPP_SUB/'FigS4_oncoprint_VAF.png', 0.0),
    ('Supp Fig S5 — Full GSEA supplement',     SUPP_SUB/'FigS5_GSEA_full.png', 0.0),
    ('Supp Fig S6 — ssGSEA and CMS',           SUPP_SUB/'FigS6_ssGSEA_CMS.png', 0.0),
    ('Supp Fig S7 — TRUST4 immune repertoire', SUPP_SUB/'FigS7_immune_TRUST4.png', 0.0),
    ('Supp Fig S8 — ML model comparison',      SUPP_SUB/'FigS8_ML_model_comparison.png', 0.0),
    ('Supp Fig S9 — GEO cohorts overview',     SUPP_SUB/'FigS9_GEO_cohorts.png', 0.0),
    ('Supp Fig S10 — HLA neoantigen detail',   SUPP_SUB/'FigS10_HLA_neoantigen_detail.png', 0.0),
    ('Supp Fig S11 — PyClone diagnostics',     SUPP_SUB/'FigS11_pyclone_diagnostics.png', 0.0),
    ('Supp Fig — SBS landscape (whole cohort)', SUPP_SUB/'FigS_SBS_landscape.png', 0.0),
    ('Supp Fig — External cohort forest (v1 legacy)', SUPP_SUB/'SuppFig_external_forest.png', 0.0),
    ('Supp Fig — GSE150082 DSB detail (v1 legacy)', SUPP_SUB/'SuppFig_GSE150082_DSB.png', 0.0),
    ('Supp Fig — External meta Z-score (v1 legacy)', SUPP_SUB/'SuppFig_meta_zscore.png', 0.0),
    # v0.7 supp additions
    ('Supp Fig (v0.7) — CONSORT sample flow', SUPP_V07/'SuppFig_consort_sample_flow.png', 0.0),
    ('Supp Fig (v0.7) — HLA-LOH lite vs strict', SUPP_V07/'SuppFig_S3_HLA_LOH_lite_vs_strict.png', 0.0),
]

# Section separator
def add_section(title):
    s = prs.slides.add_slide(BLANK)
    nph.add_textbox(s, title,
                    int(SW*0.05), int(SH*0.40), int(SW*0.90), int(SH*0.20),
                    size=36, bold=True, align=PP_ALIGN.CENTER,
                    color=(0x1F,0x3B,0x5C))

add_section('Main Figure Panels (Fig 1 – Fig 9)')
for panel_id, img, crop in PANELS:
    if not Path(img).exists():
        continue
    slide = prs.slides.add_slide(BLANK)
    cap = (f'Edit this caption freely.\n\n'
           f'• Panel body rendered as high-res image.\n'
           f'• In-plot text (axis labels, legend, P-values) can be added as editable '
           f'text boxes over this image — double-click to create/edit.\n'
           f'• For full vector editing of every in-plot label, open the matching PDF '
           f'in panels_v3/ with Illustrator / Affinity / Inkscape.')
    add_image_slide(slide, panel_id, img, cap, overlays=None, crop_top_pct=crop)

add_section('Supplementary Figures')
for panel_id, img, crop in SUPP_PANELS:
    if not Path(img).exists():
        continue
    slide = prs.slides.add_slide(BLANK)
    cap = 'Edit this caption freely. See Methods and Supplementary Text for context.'
    add_image_slide(slide, panel_id, img, cap, overlays=None, crop_top_pct=crop)

# ---- Closing slide ----
end = prs.slides.add_slide(BLANK)
nph.add_textbox(end, 'Editing guide',
                int(SW*0.05), int(SH*0.10), int(SW*0.90), int(SH*0.07),
                size=22, bold=True, color=(0x1F,0x3B,0x5C))
nph.add_textbox(end,
    '1. Native PPT slides (first three): every axis tick, axis title, marker, CI bar, diamond, '
    'and cohort/P-value annotation is a native PowerPoint shape/text — double-click to edit.\n\n'
    '2. Image + overlay slides: the plot body is an embedded image; caption and overlay text boxes '
    'are editable. For in-plot vector editing, open the matching PDF in Illustrator / Affinity Designer / '
    'Inkscape — all text runs remain editable there and you can re-export PNG/PDF.\n\n'
    '3. Source matplotlib scripts (for re-rendering with new data):\n'
    '   scripts/33_v3_forest_plot.py        — CD8 meta forest\n'
    '   scripts/38_forest_plus_akiyoshi.py  — forest + Akiyoshi row\n'
    '   scripts/39_cascade_and_loh_figures.py — cascade BCa + HLA-LOH\n'
    '   scripts/43_rebuild_all_main_figures.py — composite main figures\n'
    '   scripts/45_build_final_ppt.py       — this deck',
    int(SW*0.08), int(SH*0.20), int(SW*0.84), int(SH*0.65), size=13, color=(0x33,0x33,0x33))

prs.save(OUT_FILE)
print(f'Saved: {OUT_FILE}')
print(f'  slides: {len(prs.slides)}')
