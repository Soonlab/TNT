"""Fig 2E (alt) v2 — native PowerPoint reconstruction of genome-wide CNV.

All visible elements are native PPT objects — double-click to edit each:
  - Text (title, chromosome labels, sample IDs, Poor/Good group labels,
    colorbar labels, axis/legend text) = add_textbox (TEXT_BOX)
  - Chromosome ideogram bars, group sidebars, CNV segment cells,
    colorbar gradient blocks, group-sep wipe rectangles = MSO_SHAPE.RECTANGLE
  - Chromosome dividers, group separator (thick black), colorbar border
    = MSO_CONNECTOR.STRAIGHT

CNV cells use segment-level rectangles (one rect per contiguous same-state
segment per sample) rather than 5-Mb bins, so the file stays compact
(~1.5k cell rectangles instead of ~21k bins).

Output: manuscript/ppt/Fig2E_CNV_genome_native_editable.pptx
"""
from pathlib import Path
import importlib.util
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

spec = importlib.util.spec_from_file_location(
    'nph', '/mnt/sda1/data/TNT/analysis/scripts/44_native_ppt_helpers.py')
nph = importlib.util.module_from_spec(spec); spec.loader.exec_module(nph)

ROOT = Path('/mnt/sda1/data/TNT/analysis')
OUT_DIR = ROOT/'manuscript/ppt'; OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_FILE = OUT_DIR/'Fig2E_CNV_genome_native_editable.pptx'

def hex2rgb(h): return (int(h[1:3],16), int(h[3:5],16), int(h[5:7],16))
GOOD = hex2rgb('#2E86AB'); BAD = hex2rgb('#E63946'); SEP = hex2rgb('#1d3557')
CHR_ALT = [(0xcc,0xcc,0xcc), (0x88,0x88,0x88)]

# CN → RGB (discrete 5-level palette, matches matplotlib cmap poles)
CN_COLORS = {
    'homodel':  hex2rgb('#08306b'),  # 0
    'loss':     hex2rgb('#6baed6'),  # 1
    'neutral':  (255, 255, 255),     # 2 — don't draw
    'gain':     hex2rgb('#fc9272'),  # 3
    'amp':      hex2rgb('#67000d'),  # 4+
}
def cn_class(cn):
    if cn < 0.5: return 'homodel'
    if cn < 1.5: return 'loss'
    if cn < 2.5: return 'neutral'
    if cn < 3.5: return 'gain'
    return 'amp'

# ---------- reference layout ----------
CHR_SIZE = {'1':248956422,'2':242193529,'3':198295559,'4':190214555,'5':181538259,
            '6':170805979,'7':159345973,'8':145138636,'9':138394717,'10':133797422,
            '11':135086622,'12':133275309,'13':114364328,'14':107043718,'15':101991189,
            '16':90338345,'17':83257441,'18':80373285,'19':58617616,'20':64444167,
            '21':46709983,'22':50818468,'X':156040895}
chr_list = list(CHR_SIZE.keys())
chr_starts = {}; cum = 0
for c in chr_list:
    chr_starts[c] = cum; cum += CHR_SIZE[c]
total_bp = cum

# ---------- data: sample order and segments ----------
wes_inv = pd.read_csv(ROOT/'00_cohort/wes_inventory.tsv', sep='\t')
pre = wes_inv[wes_inv.timepoint=='pre']
meta = pre[['sample_id','subject_id','response_bin']].rename(columns={'sample_id':'sample'}).copy()
meta['resp_rank'] = meta.response_bin.map({'bad':0,'good':1})
meta = meta.sort_values(['resp_rank','subject_id']).reset_index(drop=True)
sample_order = meta['sample'].tolist()
n_bad = int((meta.response_bin=='bad').sum())
n_good = int((meta.response_bin=='good').sum())
N = len(sample_order)

CNV_DIR = ROOT/'04_wes_cnv_clonal/cnvkit'
segs = {}
for s in sample_order:
    f = CNV_DIR/f'{s}_DNA.call.cns'
    df = pd.read_csv(f, sep='\t')
    df = df[df.chromosome.astype(str).isin(chr_list)].copy()
    rows = []
    for _, r in df.iterrows():
        chrom = str(r.chromosome)
        gs = chr_starts[chrom] + int(r.start)
        ge = chr_starts[chrom] + int(r.end)
        cls = cn_class(float(r.cn))
        if cls == 'neutral': continue
        rows.append((gs, ge, cls, float(r.cn)))
    segs[s] = rows

seg_count = sum(len(v) for v in segs.values())
print(f'Total non-neutral segments across {N} samples: {seg_count}')

# ---------- presentation ----------
prs = Presentation()
prs.slide_width  = Inches(17)
prs.slide_height = Inches(9.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
slide = prs.slides.add_slide(BLANK)

# Panel letter + title
nph.add_textbox(slide, 'E', Inches(0.10), Inches(0.10), Inches(0.40), Inches(0.40),
                size=22, bold=True, color=SEP)
nph.add_textbox(slide,
    f'Genome-wide copy number landscape — Poor (top, n={n_bad}) / Good (bottom, n={n_good}); '
    f'chr 1–22 + X',
    Inches(0.6), Inches(0.12), Inches(16.0), Inches(0.38),
    size=13, bold=True, color=SEP,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---------- layout frame ----------
SIDEBAR_X = Inches(0.55)
SIDEBAR_W = Inches(0.22)
GROUP_LABEL_X = Inches(0.05)
GROUP_LABEL_W = Inches(0.50)
PLOT_X0 = Inches(1.15)
PLOT_W  = Inches(13.2)
IDEO_Y  = Inches(0.70)
IDEO_H  = Inches(0.35)
PLOT_Y0 = IDEO_Y + IDEO_H + Inches(0.05)
PLOT_H  = Inches(7.3)
ROW_H   = PLOT_H / N

# ---------- chromosome ideogram ----------
for i, c in enumerate(chr_list):
    x0 = PLOT_X0 + int(PLOT_W * chr_starts[c] / total_bp)
    w  = int(PLOT_W * CHR_SIZE[c] / total_bp)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, IDEO_Y, w, IDEO_H)
    rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(*CHR_ALT[i%2])
    rect.line.color.rgb = RGBColor(255,255,255); rect.line.width = Pt(0.4)
    rect.shadow.inherit = False
    nph.add_textbox(slide, c,
                    x0, IDEO_Y, w, IDEO_H,
                    size=10, bold=True, color=(0,0,0),
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---------- group sidebars + labels ----------
# Poor on top
poor_y0 = PLOT_Y0
poor_h  = int(ROW_H * n_bad)
side_poor = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   SIDEBAR_X, poor_y0, SIDEBAR_W, poor_h)
side_poor.fill.solid(); side_poor.fill.fore_color.rgb = RGBColor(*BAD)
side_poor.line.color.rgb = RGBColor(255,255,255); side_poor.line.width = Pt(0.4)
side_poor.shadow.inherit = False
# Good below
good_y0 = poor_y0 + poor_h
good_h  = int(ROW_H * n_good)
side_good = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   SIDEBAR_X, good_y0, SIDEBAR_W, good_h)
side_good.fill.solid(); side_good.fill.fore_color.rgb = RGBColor(*GOOD)
side_good.line.color.rgb = RGBColor(255,255,255); side_good.line.width = Pt(0.4)
side_good.shadow.inherit = False

# Group text labels (left of sidebar)
nph.add_textbox(slide, f'Poor\n(TRG 2–3)\nn = {n_bad}',
                GROUP_LABEL_X, poor_y0 + poor_h//2 - Inches(0.45),
                GROUP_LABEL_W, Inches(0.9),
                size=11, bold=True, color=BAD,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
nph.add_textbox(slide, f'Good\n(TRG 0–1)\nn = {n_good}',
                GROUP_LABEL_X, good_y0 + good_h//2 - Inches(0.45),
                GROUP_LABEL_W, Inches(0.9),
                size=11, bold=True, color=GOOD,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---------- CNV segment cells ----------
def bp_to_x(bp_pos):
    return PLOT_X0 + int(PLOT_W * bp_pos / total_bp)

for i, s in enumerate(sample_order):
    y_row = PLOT_Y0 + int(ROW_H * i)
    # draw each non-neutral segment as a rectangle
    for (gs, ge, cls, cn) in segs[s]:
        x0 = bp_to_x(gs)
        x1 = bp_to_x(ge)
        w  = max(x1 - x0, Emu(15000))
        col = CN_COLORS[cls]
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x0, y_row, w, int(ROW_H))
        r.fill.solid(); r.fill.fore_color.rgb = RGBColor(*col)
        r.line.fill.background()  # no border (would obscure thin cells)
        r.shadow.inherit = False

# ---------- chromosome dividers (white vertical lines over heatmap) ----------
for c in chr_list[1:]:
    xd = bp_to_x(chr_starts[c])
    nph.add_line(slide, xd, PLOT_Y0, xd, PLOT_Y0 + PLOT_H,
                 color=(255,255,255), width_pt=0.7)

# ---------- group horizontal separator (bold black) ----------
sep_y = PLOT_Y0 + poor_h
nph.add_line(slide, PLOT_X0, sep_y, PLOT_X0 + PLOT_W, sep_y,
             color=(0,0,0), width_pt=2.2)

# ---------- plot border (thin) ----------
for (x1,y1,x2,y2) in [
    (PLOT_X0, PLOT_Y0, PLOT_X0+PLOT_W, PLOT_Y0),                 # top
    (PLOT_X0, PLOT_Y0+PLOT_H, PLOT_X0+PLOT_W, PLOT_Y0+PLOT_H),   # bottom
    (PLOT_X0, PLOT_Y0, PLOT_X0, PLOT_Y0+PLOT_H),                 # left
    (PLOT_X0+PLOT_W, PLOT_Y0, PLOT_X0+PLOT_W, PLOT_Y0+PLOT_H),   # right
]:
    nph.add_line(slide, x1, y1, x2, y2, color=SEP, width_pt=0.8)

# ---------- sample labels (right side) ----------
for i, s in enumerate(sample_order):
    y_row = PLOT_Y0 + int(ROW_H * i)
    nph.add_textbox(slide, s,
                    PLOT_X0 + PLOT_W + Inches(0.05), y_row - Inches(0.04),
                    Inches(0.70), ROW_H + Inches(0.08),
                    size=8, color=SEP,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# ---------- colorbar (discrete 5 blocks) ----------
CBAR_X = PLOT_X0 + PLOT_W + Inches(0.95)
CBAR_Y = PLOT_Y0 + Inches(2.0)
CBAR_W = Inches(0.30)
CBAR_BLOCK_H = Inches(0.45)
cn_order = [('amp','4+\n(amp)'), ('gain','3\n(gain)'), ('neutral','2\n(neutral)'),
            ('loss','1\n(loss)'), ('homodel','0\n(homo del)')]
nph.add_textbox(slide, 'Copy number',
                CBAR_X - Inches(0.4), CBAR_Y - Inches(0.45),
                Inches(1.4), Inches(0.3),
                size=10, bold=True, color=SEP,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for k, (cls, lbl) in enumerate(cn_order):
    y0 = CBAR_Y + k*CBAR_BLOCK_H
    col = CN_COLORS[cls]
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  CBAR_X, y0, CBAR_W, CBAR_BLOCK_H)
    rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(*col)
    rect.line.color.rgb = RGBColor(*SEP); rect.line.width = Pt(0.5)
    rect.shadow.inherit = False
    nph.add_textbox(slide, lbl,
                    CBAR_X + CBAR_W + Inches(0.05), y0,
                    Inches(1.1), CBAR_BLOCK_H,
                    size=9, color=SEP,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# ---------- save ----------
prs.save(OUT_FILE)
# Shape count
from pptx import Presentation as _P
_p = _P(OUT_FILE)
total = sum(1 for s in _p.slides[0].shapes)
print(f'Saved: {OUT_FILE}')
print(f'  Total native shapes: {total}')
