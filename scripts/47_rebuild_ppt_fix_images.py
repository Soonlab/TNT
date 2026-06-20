"""Fix image-display bug in the merged deck.

The previous script cloned slide shape XML without copying the backing image
parts or remapping rId values — result: "this picture can't be displayed".

Fix: for every slide being cloned,
  1. Walk the source slide's relationships
  2. For each image/chart rel, relate the target part into the destination slide's
     part (python-pptx registers the Part in the destination package)
  3. Build {old_rId: new_rId} map
  4. Clone shape XML and rewrite every r:embed and r:link attribute using the map
"""
import subprocess, shutil
from pathlib import Path
from copy import deepcopy
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

BASE = Path('/mnt/sda1/data/TNT/analysis')
PV3  = Path('/data/data/TNT/analysis/figures/panels_v3')
SUPP_SUB = Path('/data/data/TNT/analysis/genome_medicine_submission/supplementary_figures')
SUPP_V07 = BASE/'figures/supp'
OUT_DIR = BASE/'manuscript'/'ppt'
TMP = Path('/tmp/tnt_ppt_build')
TMP.mkdir(parents=True, exist_ok=True)
BATCH_DIR = TMP/'pptx_batch'
OUT_FILE = OUT_DIR/'TNT_v0.7_final_fully_editable.pptx'

MAIN_PANELS = [
    ('Fig 1A — Study design (v0.7.1)',         PV3/'Fig1A_study_design_v2.pdf'),
    ('Fig 1B — Cohort Sankey',                 PV3/'Fig1A_sankey.pdf'),
    ('Fig 1C — Sample × assay matrix',         PV3/'Fig1D_sample_matrix.pdf'),
    ('Fig 1D — Headline 3-narrative preview',  PV3/'Fig1D_preview_forest.pdf'),
    ('Fig 2A — Driver oncoprint',              PV3/'Fig2A_oncoprint_journal.pdf'),
    ('Fig 2B — TMB raincloud by response',     PV3/'Fig2B_TMB_raincloud.pdf'),
    ('Fig 2B(alt) — SBS96 profile',            PV3/'Fig2B_SBS96_profile.pdf'),
    ('Fig 2C — SBS signature attribution',     PV3/'Fig2C_signature_attribution.pdf'),
    ('Fig 2C(alt) — MSI × TMB scatter',        PV3/'Fig2C_MSI_TMB_scatter.pdf'),
    ('Fig 2D — TMB + MSI waterfall',           PV3/'Fig2D_TMB_MSI_waterfall.pdf'),
    ('Fig 2D(alt) — SBS reassessment',         PV3/'Fig2D_SBS_reassessment.pdf'),
    ('Fig 2E — CNV + HRD proxy',               PV3/'Fig2E_CNV_HRD.pdf'),
    ('Fig 2E(alt) — CNV genome view',          PV3/'Fig2E_CNV_genome.pdf'),
    ('Fig 2F — HRD breakdown',                 PV3/'Fig2F_HRD_breakdown.pdf'),
    ('Fig 3A — TME signature radar',           PV3/'Fig3A_TME_radar.pdf'),
    ('Fig 3B — 22-signature heatmap',          PV3/'Fig3B_signature_heatmap.pdf'),
    ('Fig 3C — DEG volcano',                   PV3/'Fig3C_volcano_journal.pdf'),
    ('Fig 3D — Forest lollipop of signatures', PV3/'Fig3D_forest_lollipop.pdf'),
    ('Fig 3E — CD8 effector × exhaustion biaxial', PV3/'Fig3E_CD8_biaxial.pdf'),
    ('Fig 3F — TLS Cabrita score',             PV3/'Fig3F_TLS_Cabrita.pdf'),
    ('Fig 4A — GSEA running enrichment',       PV3/'Fig4A_running_ES.pdf'),
    ('Fig 4B — Hallmark NES bubble',           PV3/'Fig4B_Hallmark_bubble.pdf'),
    ('Fig 4C — Reactome pathway dotplot',      PV3/'Fig4C_pathway_dotplot.pdf'),
    ('Fig 4D — GSEA enrichment map',           PV3/'Fig4D_enrichment_map.pdf'),
    ('Fig 4E — ssGSEA 95-set heatmap',         PV3/'Fig4E_ssGSEA_heatmap.pdf'),
    ('Fig 4F — Category-level NES box',        PV3/'Fig4F_category_NES_box.pdf'),
    ('Fig 5A — Feature correlation',           PV3/'Fig5A_correlation.pdf'),
    ('Fig 5B — Nested LOOCV ROC (v0.7, leakage-free)', PV3/'Fig4B_ROC_nested.pdf'),
    ('Fig 5C — Feature forest with CI',        PV3/'Fig5C_forest_CI.pdf'),
    ('Fig 5D — UMAP of features',              PV3/'Fig5D_UMAP.pdf'),
    ('Fig 5E — SHAP beeswarm',                 PV3/'Fig5E_SHAP_beeswarm.pdf'),
    ('Fig 5F — Per-subject predicted probability', PV3/'Fig5F_per_subject_prediction.pdf'),
    ('Fig 6A — Cascade BCa forest (v0.7)',     BASE/'figures/panels/Fig6_cascade_BCa_forest.pdf'),
    ('Fig 6B — Paired pre→post slopes',        PV3/'Fig6B_slope_fancy.pdf'),
    ('Fig 6C — Δ forest',                      PV3/'Fig6C_delta_forest.pdf'),
    ('Fig 6D — Per-subject waterfall',         PV3/'Fig6D_waterfall.pdf'),
    ('Fig 6E — Clonal fishplot',               PV3/'Fig6E_fishplot.pdf'),
    ('Fig 6F — Cascade schematic',             PV3/'Fig6F_cascade.pdf'),
    ('Fig 7A — CD8 meta + Akiyoshi (v0.7)',    BASE/'figures/panels/Fig7_external_CD8_validation_v4.pdf'),
    ('Fig 7B — Signature × cohort heatmap',    PV3/'Fig7B_heatmap.pdf'),
    ('Fig 7C — Meta Z-score across signatures', PV3/'Fig7C_meta_Zscore.pdf'),
    ('Fig 7D — Cohort concordance',            PV3/'Fig7D_cohort_concordance.pdf'),
    ('Fig 7E — Funnel plot',                   PV3/'Fig7E_funnel.pdf'),
    ('Fig 7F — Discovery vs external effects', PV3/'Fig7F_discovery_validation.pdf'),
    ('Fig 8A — HLA allele frequency',          PV3/'Fig8A_HLA_alleles.pdf'),
    ('Fig 8B — HLA homozygosity by response',  PV3/'Fig8B_HLA_homozygosity.pdf'),
    ('Fig 8C — HLA LOH',                       PV3/'Fig8C_HLA_LOH.pdf'),
    ('Fig 8D — Pre-CRT neoantigen burden',     PV3/'Fig8D_neoantigen_pre.pdf'),
    ('Fig 8E — Paired Δ neoantigen binders',   PV3/'Fig8E_neoantigen_paired.pdf'),
    ('Fig 8F — Per-subject neoantigen lollipop', PV3/'Fig8F_neoantigen_lollipop.pdf'),
    ('Fig 9A — Clone trajectories',            PV3/'Fig9A_clone_trajectories.pdf'),
    ('Fig 9B — CCF pre vs post',               PV3/'Fig9B_CCF_pre_post.pdf'),
    ('Fig 9C — Cluster stacked composition',   PV3/'Fig9C_cluster_stacked.pdf'),
    ('Fig 9D — Dominant clone shrinkage',      PV3/'Fig9D_dominant_shrink.pdf'),
    ('Fig 9E — Shrink vs expand scatter',      PV3/'Fig9E_shrink_expand_scatter.pdf'),
    ('Fig 9F — Fate composition by response',  PV3/'Fig9F_fate_composition.pdf'),
]
SUPP_PANELS = [
    ('Supp Fig S1 — Cohort QC',                SUPP_SUB/'FigS1_cohort_QC.pdf'),
    ('Supp Fig S2 — SBS signatures full',      SUPP_SUB/'FigS2_SBS_signatures.pdf'),
    ('Supp Fig S3 — CNV + HRD detail',         SUPP_SUB/'FigS3_CNV_HRD.pdf'),
    ('Supp Fig S4 — Oncoprint + VAF detail',   SUPP_SUB/'FigS4_oncoprint_VAF.pdf'),
    ('Supp Fig S5 — Full GSEA supplement',     SUPP_SUB/'FigS5_GSEA_full.pdf'),
    ('Supp Fig S6 — ssGSEA and CMS',           SUPP_SUB/'FigS6_ssGSEA_CMS.pdf'),
    ('Supp Fig S7 — TRUST4 immune repertoire', SUPP_SUB/'FigS7_immune_TRUST4.pdf'),
    ('Supp Fig S8 — ML model comparison',      SUPP_SUB/'FigS8_ML_model_comparison.pdf'),
    ('Supp Fig S9 — GEO cohorts overview',     SUPP_SUB/'FigS9_GEO_cohorts.pdf'),
    ('Supp Fig S10 — HLA neoantigen detail',   SUPP_SUB/'FigS10_HLA_neoantigen_detail.pdf'),
    ('Supp Fig S11 — PyClone diagnostics',     SUPP_SUB/'FigS11_pyclone_diagnostics.pdf'),
    ('Supp Fig — SBS landscape cohort',        SUPP_SUB/'FigS_SBS_landscape.pdf'),
    ('Supp Fig — External cohort forest (legacy)', SUPP_SUB/'SuppFig_external_forest.pdf'),
    ('Supp Fig — GSE150082 DSB detail',        SUPP_SUB/'SuppFig_GSE150082_DSB.pdf'),
    ('Supp Fig — External meta Z-score',       SUPP_SUB/'SuppFig_meta_zscore.pdf'),
    ('Supp Fig (v0.7) — CONSORT sample flow',  SUPP_V07/'SuppFig_consort_sample_flow.pdf'),
    ('Supp Fig (v0.7) — HLA-LOH lite vs strict', SUPP_V07/'SuppFig_S3_HLA_LOH_lite_vs_strict.pdf'),
]

ALL = MAIN_PANELS + SUPP_PANELS

R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

# Namespaces of attributes whose values are rIds that need remapping when shapes
# cross packages. Cover every standard case matplotlib PDFs produce.
RELATIONAL_ATTRS = [
    f'{{{R_NS}}}embed',
    f'{{{R_NS}}}link',
    f'{{{R_NS}}}id',
]

IMAGE_RELTYPES = (
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image',
)

def clone_slide_with_media(src_prs, dst_prs):
    """Clone src_prs.slides[0] into dst_prs with all embedded images preserved."""
    src_slide = src_prs.slides[0]
    layout = dst_prs.slide_layouts[6]
    new_slide = dst_prs.slides.add_slide(layout)

    # strip the blank layout's placeholders to avoid conflicts
    for shp in list(new_slide.shapes):
        shp.element.getparent().remove(shp.element)

    # ---- Step 1: transfer image parts (copy blobs to dst package) and build rId map ----
    # Must copy blobs into dst package (SHA-deduplicated) rather than reusing src Parts,
    # because src ImageParts carry src-package partnames that collide across sources.
    dst_pkg = dst_prs.part.package
    id_map = {}
    for rel in src_slide.part.rels.values():
        if rel.reltype in IMAGE_RELTYPES:
            blob = rel.target_part.blob
            new_img_part = dst_pkg.get_or_add_image_part(BytesIO(blob))
            new_rId = new_slide.part.relate_to(new_img_part, rel.reltype)
            id_map[rel.rId] = new_rId

    # ---- Step 2: clone each shape and rewrite rIds ----
    spTree = new_slide.shapes._spTree
    for shp in src_slide.shapes:
        new_el = deepcopy(shp.element)
        # rewrite every rId attribute we care about
        for elem in new_el.iter():
            for attr in RELATIONAL_ATTRS:
                if attr in elem.attrib:
                    old = elem.attrib[attr]
                    if old in id_map:
                        elem.attrib[attr] = id_map[old]
        spTree.insert_element_before(new_el, 'p:extLst')

    return new_slide

def add_title_bar(slide, text, prs):
    tb = slide.shapes.add_textbox(Emu(200000), Emu(100000),
                                  prs.slide_width - Emu(400000), Emu(400000))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = RGBColor(0x1F, 0x3B, 0x5C); r.font.name = 'Arial'

def main():
    # Reuse previously converted per-panel PPTXs in BATCH_DIR
    # (all 74 were successfully converted last run).
    if not BATCH_DIR.exists() or not list(BATCH_DIR.glob('*.pptx')):
        print('Batch-converted PPTXs not found; run script 46 first.')
        return

    master = Presentation()
    master.slide_width  = Inches(13.33)
    master.slide_height = Inches(7.5)
    BLANK = master.slide_layouts[6]

    cover = master.slides.add_slide(BLANK)
    tb = cover.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.2))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = 'TNT MSS LARC — v0.7 fully-editable figure deck'
    r.font.size = Pt(30); r.font.bold = True
    r.font.color.rgb = RGBColor(0x1F, 0x3B, 0x5C); r.font.name = 'Arial'
    tb2 = cover.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.5))
    tf = tb2.text_frame; tf.word_wrap = True
    for i, line in enumerate([
        '1 panel per slide. No figure titles.',
        'All in-plot text (axis labels, ticks, legend, P-values, cohort names) is a NATIVE PowerPoint text run.',
        'All embedded images are properly relocated to the master deck.',
        '',
        'Double-click anywhere to edit text. Images are editable via crop/resize/replace.',
    ]):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(14); r.font.name = 'Arial'
        r.font.color.rgb = RGBColor(0x33,0x33,0x33)

    ok = 0; skipped = []
    for idx, (title, pdf) in enumerate(ALL, 1):
        pptx_path = BATCH_DIR / (Path(pdf).stem + '.pptx')
        if not pptx_path.exists():
            skipped.append((title, str(pdf)))
            continue
        try:
            src = Presentation(str(pptx_path))
            new_slide = clone_slide_with_media(src, master)
            add_title_bar(new_slide, title, master)
            ok += 1
        except Exception as e:
            import traceback
            skipped.append((title, f'{Path(pdf).name} — {type(e).__name__}: {e}'))
            if idx <= 3:
                traceback.print_exc()
        if idx % 10 == 0:
            print(f'  {idx}/{len(ALL)} (ok={ok}, skipped={len(skipped)})')

    if skipped:
        s = master.slides.add_slide(BLANK)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5),
                                  Inches(12.3), Inches(1.0))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = f'Skipped panels ({len(skipped)})'
        r.font.size = Pt(22); r.font.bold = True
        r.font.color.rgb = RGBColor(0xB2,0x1F,0x1F)
        tb2 = s.shapes.add_textbox(Inches(0.5), Inches(1.6),
                                   Inches(12.3), Inches(5.5))
        tf = tb2.text_frame; tf.word_wrap = True
        for i, (t, detail) in enumerate(skipped):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            r = p.add_run(); r.text = f'  • {t}   ({detail})'
            r.font.size = Pt(11); r.font.name = 'Arial'
            r.font.color.rgb = RGBColor(0x66,0x66,0x66)

    master.save(OUT_FILE)
    print(f'\nSaved: {OUT_FILE}')
    print(f'  slides: {len(master.slides)}   ok: {ok}   skipped: {len(skipped)}')

if __name__ == '__main__':
    main()
