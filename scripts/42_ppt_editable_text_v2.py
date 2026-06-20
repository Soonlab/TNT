"""Rebuild editable PPT (v2) — text overlays for key strings.

Strategy: embed the full PNG as background, then place editable PPT text boxes
at the positions of key strings (title, AUC, meta Z, Akiyoshi OR, P-values).
Users can edit these text boxes directly in PowerPoint without having to touch
the image pixels.

For each figure, a small YAML-like Python list encodes:
    (text, x_rel, y_rel, w_rel, h_rel, font_size, bold, color)
with x/y as fractions of slide width/height, so positions adapt to the slide.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

BASE = Path('/mnt/sda1/data/TNT/analysis')
OUT  = BASE/'manuscript'/'ppt'/'TNT_v0.7_figures_editable_v2.pptx'

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

SW = prs.slide_width   # EMU
SH = prs.slide_height

def add_textbox(slide, text, xf, yf, wf, hf, size=12, bold=False, italic=False,
                color=(0x33,0x33,0x33), align=PP_ALIGN.LEFT, fill=None, border=None,
                font='Arial'):
    """Add a text box using fractional coordinates of slide."""
    tb = slide.shapes.add_textbox(int(SW*xf), int(SH*yf), int(SW*wf), int(SH*hf))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = RGBColor(*color)
        r.font.name = font
    if fill is not None:
        tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor(*fill)
    if border is not None:
        tb.line.color.rgb = RGBColor(*border); tb.line.width = Pt(0.75)
    return tb

# ---------- Figures and their editable overlays ----------
# Each entry: (slide_title, image_path, x_img, y_img, w_img, h_img, ref_paper, overlays[])
# overlays = list of (text, x, y, w, h, font_size, bold, color)

FIG_CATALOG = [
    # -------- Fig 4B (new nested ROC) --------
    dict(
        title='Figure 4B — Nested outer-LOOCV ROC (leakage-free, v0.7)',
        image=str(BASE/'figures/panels_v3/Fig4B_ROC_nested.png'),
        image_box=(0.025, 0.12, 0.55, 0.78),
        reference='McGranahan N et al. Cell 2017 Fig 5 — ROC with 95 % bootstrap CI band.',
        overlays=[
            dict(text='Figure 4B', xf=0.03, yf=0.04, wf=0.30, hf=0.06,
                 size=24, bold=True, color=(0x1F,0x3B,0x5C)),
            dict(text='Nested outer-LOOCV (leakage-free)', xf=0.03, yf=0.09, wf=0.55, hf=0.04,
                 size=12, italic=True, color=(0x55,0x55,0x55)),
            # Editable AUC callouts (right-side panel)
            dict(text='AUC summary (editable)', xf=0.60, yf=0.12, wf=0.36, hf=0.05,
                 size=14, bold=True, color=(0x1F,0x3B,0x5C)),
            dict(text='LASSO   AUC = 0.650   [0.45, 0.83]', xf=0.60, yf=0.19, wf=0.36, hf=0.05,
                 size=13, bold=True, color=(0x1F,0x77,0xB4), fill=(0xE6,0xF0,0xFA)),
            dict(text='ElasticNet   AUC = 0.686   [0.49, 0.85]', xf=0.60, yf=0.25, wf=0.36, hf=0.05,
                 size=13, bold=True, color=(0x2C,0xA0,0x2C), fill=(0xE6,0xF5,0xE6)),
            dict(text='Non-nested (leaked) 0.755 — shown for transparency.',
                 xf=0.60, yf=0.33, wf=0.36, hf=0.05,
                 size=10, italic=True, color=(0x88,0x55,0x44)),
            dict(text='Permutation P — pending (time-prohibitive; 1000 × nested LOOCV).',
                 xf=0.60, yf=0.39, wf=0.36, hf=0.05,
                 size=9, italic=True, color=(0x99,0x99,0x99)),
            dict(text='Key message (editable):\n'
                      'Leakage-free AUC is modest with 95 % CI touching 0.5.\n'
                      'Pre-CRT tumor-intrinsic classifier = discovery-stage, awaits TNT-matched validation.\n'
                      'Immune (CD8-cytotoxic) axis is the externally-reproducible arm (Fig 7).',
                 xf=0.60, yf=0.48, wf=0.36, hf=0.28,
                 size=10, color=(0x33,0x33,0x33), fill=(0xFF,0xF8,0xEC), border=(0xCC,0xAA,0x66)),
        ],
    ),
    # -------- Fig 6 (new cascade BCa forest) --------
    dict(
        title='Figure 6 — Cascade BCa bootstrap forest (v0.7)',
        image=str(BASE/'figures/panels/Fig6_cascade_BCa_forest.png'),
        image_box=(0.025, 0.12, 0.70, 0.78),
        reference='Mariathasan S et al. Nature 2018 Fig 4b; Tumeh PC Nature 2014; Chen DS & Mellman I Nature 2017.',
        overlays=[
            dict(text='Figure 6', xf=0.03, yf=0.04, wf=0.22, hf=0.06,
                 size=24, bold=True, color=(0x1F,0x3B,0x5C)),
            dict(text='Within-group + between-group paired Δ with BCa 95 % CI (n = 14)',
                 xf=0.03, yf=0.09, wf=0.70, hf=0.04,
                 size=12, italic=True, color=(0x55,0x55,0x55)),
            # Right-side editable annotations
            dict(text='Robust (CI excludes 0)', xf=0.75, yf=0.14, wf=0.22, hf=0.04,
                 size=12, bold=True, color=(0x2a,0x9d,0x8f)),
            dict(text='• Treg   diff = +1.21 [+0.06, +1.97]\n   MW P = 0.026',
                 xf=0.75, yf=0.19, wf=0.22, hf=0.08,
                 size=11, color=(0x2a,0x9d,0x8f), fill=(0xE8,0xF5,0xF2)),
            dict(text='Exploratory (CI spans 0)', xf=0.75, yf=0.30, wf=0.22, hf=0.04,
                 size=12, bold=True, color=(0x88,0x88,0x88)),
            dict(text='• MHC-II, CD8 exhaustion\n• IGH clonotype count\n'
                      '• SBS5 mutation clearance\n• Neoantigen binders / sites\n• TRB Shannon diversity',
                 xf=0.75, yf=0.35, wf=0.22, hf=0.20,
                 size=10, color=(0x66,0x66,0x66), fill=(0xF5,0xF5,0xF5)),
            dict(text='Caveat (editable):\n'
                      'Between-group BCa CIs reflect n = 6–7 per group.\n'
                      'Within-good CI excludes 0 for SBS5, neo-binders, '
                      'Treg — within-group robustness is consistent even '
                      'where between-group inference is under-powered.',
                 xf=0.75, yf=0.58, wf=0.22, hf=0.28,
                 size=9.5, color=(0x33,0x33,0x33), fill=(0xFF,0xF8,0xEC),
                 border=(0xCC,0xAA,0x66)),
        ],
    ),
    # -------- Fig 7 (new external + Akiyoshi) --------
    dict(
        title='Figure 7 — External CD8-cytotoxic validation (v0.7)',
        image=str(BASE/'figures/panels/Fig7_external_CD8_validation_v4.png'),
        image_box=(0.025, 0.12, 0.63, 0.82),
        reference='Ayers M et al. JCI 2017; Rooney MS et al. Cell 2015; Akiyoshi T et al. JAMA Netw Open 2023.',
        overlays=[
            dict(text='Figure 7', xf=0.03, yf=0.04, wf=0.22, hf=0.06,
                 size=24, bold=True, color=(0x1F,0x3B,0x5C)),
            dict(text='External validation — CD8-cytotoxic axis across 10 cohorts, N > 1,000',
                 xf=0.03, yf=0.09, wf=0.65, hf=0.04,
                 size=12, italic=True, color=(0x55,0x55,0x55)),
            # Right-side editable stats
            dict(text='KEY STATISTICS', xf=0.69, yf=0.13, wf=0.28, hf=0.04,
                 size=13, bold=True, color=(0x1F,0x3B,0x5C)),
            dict(text='Our 9-cohort meta-analysis', xf=0.69, yf=0.18, wf=0.28, hf=0.04,
                 size=11, bold=True, color=(0xE6,0x39,0x46)),
            dict(text='   Stouffer Z = +2.74,  P = 0.006\n   N = 721 (8/9 cohorts concordant)',
                 xf=0.69, yf=0.22, wf=0.28, hf=0.07,
                 size=11, color=(0xE6,0x39,0x46), fill=(0xFD,0xEC,0xEE)),
            dict(text='Akiyoshi 2023 (GSE216616)', xf=0.69, yf=0.31, wf=0.28, hf=0.04,
                 size=11, bold=True, color=(0x6A,0x4C,0x93)),
            dict(text='   n = 298,  OR = 3.81 [1.82, 7.97]\n   GZMA × PRF1 P = 0.005\n   Hallmark IFN-γ enriched in good',
                 xf=0.69, yf=0.35, wf=0.28, hf=0.10,
                 size=11, color=(0x6A,0x4C,0x93), fill=(0xF1,0xEC,0xF7)),
            dict(text='Combined evidence', xf=0.69, yf=0.47, wf=0.28, hf=0.04,
                 size=11, bold=True, color=(0x1F,0x3B,0x5C)),
            dict(text='   > 1,000 independent patients\n   across 10 cohorts',
                 xf=0.69, yf=0.51, wf=0.28, hf=0.08,
                 size=11, color=(0x1F,0x3B,0x5C), fill=(0xE6,0xF0,0xFA),
                 border=(0x1F,0x3B,0x5C)),
            dict(text='Key message (editable):\n'
                      'The CD8-cytotoxic effector axis (GZMA/B, PRF1, '
                      'IFNG, CD8A/B, NKG7, CXCL9/10) is the externally '
                      '-reproducible arm of our discovery. Tumor-intrinsic '
                      'DSB/HDR/E2F axes are cohort-heterogeneous and remain '
                      'discovery-stage pending TNT-matched validation.',
                 xf=0.69, yf=0.62, wf=0.28, hf=0.30,
                 size=9.5, color=(0x33,0x33,0x33), fill=(0xFF,0xF8,0xEC),
                 border=(0xCC,0xAA,0x66)),
        ],
    ),
]

# ---------- Build slides ----------
# Cover
cover = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(cover, 'TNT MSS LARC — v0.7 editable figure deck (v2, overlayed text)',
            0.04, 0.08, 0.92, 0.08, size=26, bold=True, color=(0x1F,0x3B,0x5C))
add_textbox(cover,
            'v2 differences from v1:\n'
            '  • Key numeric strings (AUC, meta Z, OR, P values) are rendered as '
            'native PowerPoint text boxes OVERLAID on the figure image.\n'
            '  • The underlying figure PNG still carries the plot elements, but the '
            '"KEY STATISTICS" side panels and the big figure title are editable text.\n'
            '  • Edit by double-clicking any text box in PowerPoint — the image is untouched.\n'
            '  • For full in-plot text editability (axis labels, legend, data labels), '
            'the source PDFs are vector and can be opened in Illustrator / Affinity Designer.',
            0.05, 0.22, 0.90, 0.30, size=13)
add_textbox(cover,
            'Scope of v2 overlay treatment: Fig 4B nested ROC, Fig 6 cascade BCa forest, '
            'Fig 7 external CD8 validation (all three v0.7 additions). Fig 1–3, Fig 4A–F, Fig 5, '
            'and supplementary figures remain as in v1; their vector PDFs in '
            '/figures/panels_v3/ retain native text for editing.',
            0.05, 0.55, 0.90, 0.20, size=11, italic=True, color=(0x55,0x55,0x55))

# Figure slides
for fc in FIG_CATALOG:
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    # image
    x, y, w, h = fc['image_box']
    sl.shapes.add_picture(fc['image'],
                          int(SW*x), int(SH*y),
                          width=int(SW*w), height=int(SH*h))
    # reference citation (top)
    add_textbox(sl,
                f"Reference (style source): {fc['reference']}",
                0.03, 0.93, 0.94, 0.05, size=9, italic=True,
                color=(0x6A,0x4C,0x93))
    # overlays
    for ov in fc['overlays']:
        add_textbox(sl, ov['text'], ov['xf'], ov['yf'], ov['wf'], ov['hf'],
                    size=ov.get('size', 11),
                    bold=ov.get('bold', False),
                    italic=ov.get('italic', False),
                    color=ov.get('color', (0x33,0x33,0x33)),
                    fill=ov.get('fill'),
                    border=ov.get('border'))

# Appendix
app = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(app, 'How to edit figure text in this deck',
            0.04, 0.05, 0.92, 0.07, size=22, bold=True, color=(0x1F,0x3B,0x5C))
add_textbox(app,
            '1. TITLE / KEY STATISTICS / CAPTION boxes = native PowerPoint text boxes.\n'
            '   → Double-click to edit; all text, fonts, colors, and fills are editable.\n\n'
            '2. IN-PLOT TEXT (axis labels, cohort names, legend) lives inside the image.\n'
            '   → For in-plot editing, open the matching PDF in:\n'
            '       Adobe Illustrator   (Open → File → commit changes → Save As PDF)\n'
            '       Affinity Designer   (Open → edit text → Export PDF/PNG)\n'
            '       Inkscape             (Open → Ctrl+T → edit → Save As PDF)\n'
            '   PDFs are vector and keep every text run as an editable object.\n\n'
            '3. Source PDFs for v0.7 additions:\n'
            '     figures/panels_v3/Fig4B_ROC_nested.pdf\n'
            '     figures/panels/Fig6_cascade_BCa_forest.pdf\n'
            '     figures/panels/Fig7_external_CD8_validation_v4.pdf\n'
            '     figures/supp/SuppFig_consort_sample_flow.pdf\n'
            '     figures/supp/SuppFig_S3_HLA_LOH_lite_vs_strict.pdf\n\n'
            '4. To re-render from scratch with new numbers, source matplotlib scripts are at:\n'
            '     scripts/33_v3_forest_plot.py       (CD8 forest)\n'
            '     scripts/38_forest_plus_akiyoshi.py  (forest + Akiyoshi row)\n'
            '     scripts/39_cascade_and_loh_figures.py (cascade BCa + HLA-LOH)\n'
            '     scripts/40_build_editable_ppt.py    (this deck)',
            0.05, 0.15, 0.90, 0.75, size=12)

prs.save(OUT)
print(f'Saved v2 editable PPT: {OUT}')
print(f'  slides: {len(prs.slides)}')
