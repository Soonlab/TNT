"""Build an editable PowerPoint with all manuscript v0.7 figures.

Each slide contains:
  - Figure PNG embedded (high-res)
  - Editable title text box
  - Reference paper citation (top-tier Nature/Cell journal) as editable text
  - Editable caption placeholder
  - Style notes as editable sidebar

Figures reflect the style conventions of the cited reference papers in each domain.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

BASE = Path('/mnt/sda1/data/TNT/analysis')
FIG  = BASE/'figures'
OUT  = BASE/'manuscript'/'ppt'
OUT.mkdir(parents=True, exist_ok=True)

# ---------- Figure catalog ----------
# (slide_title, figure_path, reference_paper, style_notes, caption_seed)
SLIDES = [
    # --------- Fig 1 ---------
    ('Figure 1 — Cohort overview',
     FIG/'panels/Fig1_cohort.png',
     'Chen P-L et al. Cancer Discov 2016, "Analysis of Immune Signatures in Longitudinal Tumor Samples" — used cohort pie + stacked-bar layout.',
     'Arial-like sans, 7.2×5.5 inches, 600 dpi, 0.6-pt spines, good=#2E86AB / bad=#E63946.',
     'Thirty-five MSS LARC patients (good n=18, bad n=17) analysed by matched WES (77 samples) and RNA-seq (56 samples). cT stage, age/sex, sample matrix, and study design panels.'),
    # --------- Fig 2 ---------
    ('Figure 2 — WES landscape',
     FIG/'panels/Fig2_WES_landscape.png',
     'Alexandrov LB et al. Nature 2020, "The repertoire of mutational signatures in human cancer" — SBS stacked-bar + driver oncoprint convention.',
     'Driver oncoprint rows = genes (APC, TP53, KRAS, FBXW7, KMT2D), cols = tumors; TMB/MSI strip above.',
     'TMB (good 1.85 vs bad 1.40 /Mb, P=0.186), MSI all <0.19 %, driver oncoprint dominated by APC-TP53-KRAS, SBS5/SBS1 dominant, SBS3 absent, LST HRD-proxy trend higher in bad.'),
    # --------- Fig 3 ---------
    ('Figure 3 — RNA immune signatures',
     FIG/'panels/Fig3_RNA_signatures.png',
     'Mariathasan S et al. Nature 2018, "TGF-β attenuates tumour response to PD-L1 blockade" (IMvigor210) — immune signature heatmap and stratified boxplot convention.',
     'Z-score heatmap rows = 22 immune signatures, cols = samples, col annotation bar = response; boxplots with jittered points.',
     '22 immune signatures and pathway scores; pre-CRT CD8-proliferation higher in good (P=0.035); MHC II modestly lower in good (P=0.074); post-CRT immune activation in good responders.'),
    # --------- Fig 4 ---------
    ('Figure 4 — GSEA integration',
     FIG/'panels/Fig4_GSEA_integration.png',
     'Cristescu R et al. Science 2018, "Pan-tumor genomic biomarkers for PD-1 checkpoint blockade" — ssGSEA-style integration and Hallmark GSEA bar plots.',
     'Hallmark GSEA bars ordered by NES, cap colour by P<10⁻¹⁰; Reactome heatmap; ssGSEA top-box jitter panels.',
     'E2F targets (NES 2.78, P<10⁻¹⁰), G2M, MYC, Reactome DSB/HDR up in good; EMT down. ssGSEA DSB P=0.007, HDR P=0.020.'),
    # --------- Fig 4B (new nested ROC) ---------
    ('Figure 4B — Nested outer-LOOCV ROC (NEW, v0.7)',
     FIG/'panels_v3/Fig4B_ROC_nested.png',
     'McGranahan N et al. Cell 2017, "Allele-Specific HLA Loss and Immune Escape" Fig 5 — ROC with 95 % bootstrap CI band + permutation inset convention.',
     'Nested outer-LOOCV + inner 5-fold hyperparam tuning (no leakage); bootstrap 95 % CI band; AUC reported with CI; contrast with leaked non-nested 0.755 for transparency.',
     'LASSO nested AUC 0.650 [0.45, 0.83]; ElasticNet nested AUC 0.686 [0.49, 0.85]. Non-nested (leaked) 0.755 shown as comparator. 95 % CI touching 0.5 — pre-CRT classifier is a modest discovery-stage predictor pending TNT-matched validation.'),
    # --------- Fig 5 ---------
    ('Figure 5 — Neoantigen clearance across radiation phase',
     FIG/'panels/Fig5_neoantigen_cascade.png',
     'McGranahan N et al. Science 2016, "Clonal neoantigens elicit T cell immunoreactivity and sensitivity to immune checkpoint blockade" — paired pre/post slopegraph + neoantigen burden boxplot convention; Riaz N et al. Cell 2017 (Nivolumab) for paired dynamics.',
     'Panels A–B boxplot with jitter, individual patient points; Panel C slopegraph (pre → post per subject), lines colored by response; Panel D optional HLA-LOH detail moved to Supp S3.',
     'Pre-CRT missense P=0.082, PCN P=0.15, strong binders P=0.55 (trend for good > bad). Paired Δ binders good median −312 [−626, −123] vs bad −100; within-good CI excludes zero, between-group exploratory.'),
    # --------- Fig 6 (new cascade BCa forest) ---------
    ('Figure 6 — Cascade BCa bootstrap forest (NEW, v0.7)',
     FIG/'panels/Fig6_cascade_BCa_forest.png',
     'Mariathasan S et al. Nature 2018 Fig 4b; Tumeh PC et al. Nature 2014 — paired delta with 95 % CI forest; Chen DS, Mellman I Nature 2017 cancer-immunity-cycle figure as mental model.',
     'Two-panel forest: A = within-good (blue circle) and within-bad (red square) medians with BCa 95 % CIs, standardized per feature; B = between-group diff BCa 95 % CI with teal diamond for CIs excluding 0 (robust), grey diamond for CIs spanning 0 (exploratory).',
     'Treg is the only feature with between-group BCa CI excluding zero (MW P=0.026). All other cascade features (SBS5, neo-binder, MHC II, CD8 exhaustion, IGH, TRB Shannon) are reported as exploratory with CIs spanning zero at n = 14 paired subjects.'),
    # --------- Fig 7 (new Akiyoshi row) ---------
    ('Figure 7 — External CD8-cytotoxic validation (10 cohorts, N > 1,000) (NEW, v0.7)',
     FIG/'panels/Fig7_external_CD8_validation_v4.png',
     'Ayers M et al. J Clin Invest 2017, "IFN-γ–related mRNA profile predicts clinical response to PD-1 blockade" Fig 2; Rooney MS et al. Cell 2015 "Molecular and Genetic Properties of Tumors Associated with Local Immune Cytolytic Activity" — per-cohort forest + meta diamond style. Akiyoshi T et al. JAMA Netw Open 2023 independent 298-patient finding convergence.',
     'Panel A: 9 cohort forest + red meta diamond + separate purple Akiyoshi row (n=298, OR 3.81 [1.82, 7.97]) with dashed separator, labelled "> 1,000 patients across 10 cohorts"; Panel B: Stouffer Z bar across signatures; Panel C: CD8 vs Tumor_cellcycle decoupling scatter.',
     'Stouffer Z = +2.74, P = 0.006 across 9 nCRT cohorts (N = 721), 8/9 concordant. Akiyoshi 2023 convergence brings total > 1,000 patients. Tumor-intrinsic axes cohort-heterogeneous (P > 0.19). EMT correct direction.'),
    # --------- Supp S1 CONSORT ---------
    ('Supp Fig S1 — CONSORT sample flow (NEW, v0.7)',
     FIG/'supp/SuppFig_consort_sample_flow.png',
     'CONSORT 2010 Statement (Schulz KF et al. BMJ 2010) + STROBE Statement (von Elm E et al. Ann Intern Med 2007) — sample-flow diagram style. For molecular cohort adaptation see Rosenthal R et al. Nature 2019 Fig 1 CONSORT panel.',
     'Boxes colored by branch (WES blue, RNA-seq pink, paired green). Each box = analysis step; arrows = sample transitions; number of patients/tumors/paired sets annotated.',
     '35 enrolled → 18 good + 17 bad; 77 WES samples → 49 PASS Mutect2 VCFs → 41 matched T-N pairs → 14 paired pre-/post-CRT (WES) → 12 PyClone convergent / 11 neoantigen Δ; 56 RNA-seq → 33 pre-CRT for DEG/GSEA → 12 paired Δ (6 good + 6 bad for between-group).'),
    # --------- Supp S3 HLA-LOH ---------
    ('Supp Fig S3 — HLA-LOH lite vs strict (NEW, v0.7)',
     FIG/'supp/SuppFig_S3_HLA_LOH_lite_vs_strict.png',
     'McGranahan N et al. Cell 2017 "Allele-Specific HLA Loss and Immune Escape in Lung Cancer Evolution" Fig 2 and Fig 4 — LOHHLA subject panels + allele-imbalance convention. Rosenthal R et al. Nature 2019 "Neoantigen-directed immune escape" Fig 5 — LOH prevalence stratified barplot.',
     'Three panels: A grouped barplot lite vs strict LOH subject counts stratified by response; B per-subject strict-LOH pre→post resolution lines (subj 3: 2→0 loci, subj 4: 1→0); C text panel documenting thresholds (|Δratio|≥0.20, Bonferroni P<0.01, depth ≥30).',
     'Under stricter Bonferroni-corrected IMGT-allele criteria, pre-CRT LOH events drop from 10 lite → 2 strict; both strict-LOH subjects are eventual good responders and both show complete pre→post resolution. Reported as anecdotal corroboration, not between-group statistical claim.'),
    # --------- Supp Text S3 diagnostic ---------
    ('Supp Text S3 — External meta v3 diagnostic (signature and label correction)',
     None,
     'Reporting-transparency precedent: Mariathasan S et al. Nature 2018 supplementary methods + Rosenthal R et al. Nature 2019 methods transparency on response classifier.',
     'Text-only slide. Describes the v1 classifier bug ("Non-responder" substring mismatch) and signature confounding (cell-cycle genes labelled CD8_proliferation), and the v3 correction (pure CD8-cytotoxic signature + manual per-cohort TRG scale + N rose from 179 to 721).',
     'Before (v1): CD8_proliferation Z=+0.06, P=0.48. After (v3): CD8_cytotoxic Z=+2.74, P=0.006 (8/9 concordant). EMT direction also recovered. The v1 artifact is retained in the main Methods and this Supp Text for reviewer-accessible transparency.'),
]

# ---------- Build PPT ----------
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_title(slide, text, y=0.15, size=22, bold=True, color=(0x1F,0x3B,0x5C)):
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(y), Inches(12.6), Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = RGBColor(*color)
    r.font.name = 'Arial'
    return tb

def add_text(slide, text, x, y, w, h, size=10, bold=False, italic=False, color=(0x33,0x33,0x33), align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = RGBColor(*color)
        r.font.name = 'Arial'
    return tb

# ---------- Cover slide ----------
cover = prs.slides.add_slide(prs.slide_layouts[6])
add_title(cover, 'TNT MSS LARC — Manuscript v0.7 Figure Set (Genome Medicine)',
          y=0.8, size=26)
add_text(cover,
         'Editable figure deck\n'
         'Version 0.7 — 2026-04-15\n\n'
         'Each slide embeds the final PNG (600 dpi) and carries:\n'
         '  • Editable title and caption\n'
         '  • Reference paper citation (top-tier Nature/Cell/Science family) describing the style source\n'
         '  • Style notes reproducing the reference convention\n\n'
         'Paired PDF vector versions are co-located in /mnt/sda1/data/TNT/analysis/figures/.\n'
         'Main figures in figures/panels/; supplementary figures in figures/supp/.',
         0.8, 1.8, 11.7, 4.5, size=14)
add_text(cover,
         'Figures in this deck (in manuscript order):\n'
         '  Fig 1  Cohort overview\n'
         '  Fig 2  WES landscape\n'
         '  Fig 3  RNA immune signatures\n'
         '  Fig 4  GSEA integration\n'
         '  Fig 4B Nested outer-LOOCV ROC (NEW, v0.7)\n'
         '  Fig 5  Neoantigen clearance\n'
         '  Fig 6  Cascade BCa bootstrap forest (NEW, v0.7)\n'
         '  Fig 7  External CD8-cytotoxic validation + Akiyoshi 2023 convergence (NEW, v0.7)\n'
         '  Supp Fig S1  CONSORT sample flow (NEW, v0.7)\n'
         '  Supp Fig S3  HLA-LOH lite vs strict (NEW, v0.7)\n'
         '  Supp Text S3  External meta v3 diagnostic',
         0.8, 5.2, 11.7, 2.2, size=12, italic=True, color=(0x55,0x55,0x55))

# ---------- Content slides ----------
for slide_title, img_path, reference, style, caption in SLIDES:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, slide_title)
    # reference citation (top band)
    add_text(slide,
             f'Reference (style source): {reference}',
             0.35, 0.72, 12.6, 0.5,
             size=10, italic=True, color=(0x6A,0x4C,0x93))
    # embed figure (if exists)
    if img_path and Path(img_path).exists():
        # fit inside 7.5 x 4.8 inches, left aligned
        pic = slide.shapes.add_picture(str(img_path),
                                       Inches(0.35), Inches(1.25),
                                       height=Inches(4.8))
        # if too wide, shrink by width instead
        if pic.width > Inches(8.0):
            slide.shapes._spTree.remove(pic._element)
            pic = slide.shapes.add_picture(str(img_path),
                                           Inches(0.35), Inches(1.25),
                                           width=Inches(8.0))
    else:
        # text-only slide (Supp Text S3)
        add_text(slide,
                 '[Text-only supplementary. See Supp_Text_S3_external_meta_diagnostic.md]',
                 0.35, 3.0, 8.0, 1.0, size=14, italic=True, color=(0x88,0x88,0x88),
                 align=PP_ALIGN.CENTER)
    # side panel: style notes + editable caption
    panel_x = 8.6
    add_text(slide, 'STYLE NOTES', panel_x, 1.25, 4.3, 0.35, size=11, bold=True, color=(0x1F,0x3B,0x5C))
    add_text(slide, style, panel_x, 1.6, 4.3, 2.2, size=9.5)
    add_text(slide, 'EDITABLE CAPTION', panel_x, 3.95, 4.3, 0.35, size=11, bold=True, color=(0x1F,0x3B,0x5C))
    add_text(slide, caption, panel_x, 4.3, 4.3, 2.6, size=9.5)
    # footer
    add_text(slide,
             'TNT MSS LARC v0.7 — Genome Medicine target — 2026-04-15',
             0.35, 7.05, 12.6, 0.3,
             size=8, italic=True, color=(0x99,0x99,0x99))

# ---------- Appendix: reference summary ----------
app = prs.slides.add_slide(prs.slide_layouts[6])
add_title(app, 'Appendix — Top-tier style references per figure')
ref_lines = (
    'Fig 1  Chen P-L et al. Cancer Discov 2016 — cohort/schematic\n'
    'Fig 2  Alexandrov LB et al. Nature 2020 — SBS signatures, driver oncoprint\n'
    'Fig 3  Mariathasan S et al. Nature 2018 — immune signature heatmap (IMvigor210)\n'
    'Fig 4  Cristescu R et al. Science 2018 — GSEA integration; Subramanian A 2005 — GSEA itself\n'
    'Fig 4B McGranahan N et al. Cell 2017 Fig 5 — ROC + bootstrap CI band + permutation inset\n'
    'Fig 5  McGranahan N et al. Science 2016 — paired pre/post neoantigen slopegraph; Riaz N Cell 2017 — nivolumab paired dynamics\n'
    'Fig 6  Mariathasan S Nature 2018 Fig 4b — paired delta forest; Tumeh PC Nature 2014 — pre/post immune dynamics; Chen DS & Mellman I Nature 2017 — cancer-immunity cycle\n'
    'Fig 7  Ayers M et al. J Clin Invest 2017 — per-cohort forest; Rooney MS et al. Cell 2015 — cytolytic activity (GZMA × PRF1); Akiyoshi T JAMA Netw Open 2023 — convergent 298-patient validation\n'
    'Supp S1  CONSORT 2010 (Schulz KF et al. BMJ 2010); Rosenthal R et al. Nature 2019 — adapted CONSORT for molecular cohorts\n'
    'Supp S3  McGranahan N Cell 2017 Fig 2–4 — LOHHLA allele imbalance; Rosenthal R Nature 2019 — neoantigen-directed immune escape LOH stratification\n'
    'Supp Text S3  Reporting-transparency precedent: Mariathasan S Nature 2018 supplementary; Rosenthal R Nature 2019 methods transparency'
)
add_text(app, ref_lines, 0.4, 0.9, 12.5, 6.4, size=12, color=(0x33,0x33,0x33))

OUT_FILE = OUT/'TNT_v0.7_figures_editable.pptx'
prs.save(OUT_FILE)
print(f'Saved PPT: {OUT_FILE}')
print(f'  slides: {len(prs.slides)}')
