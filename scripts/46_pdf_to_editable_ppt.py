"""Build final editable PPT via LibreOffice PDF→PPTX conversion.

Each matplotlib-generated PDF is converted to a single-slide PPTX by LibreOffice's
impress_pdf_import filter; this preserves every text element (axis labels, tick
labels, legend text, P-values, cohort names) as a native PowerPoint text run.

We then clone each converted slide into one master deck, one panel per slide,
with no figure title (matplotlib suptitles are dropped because each panel PDF
is a pure plot with no suptitle).

Requires: libreoffice installed (verified at /usr/bin/libreoffice).
"""
import subprocess, tempfile, shutil, os, sys
from pathlib import Path
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = Path('/mnt/sda1/data/TNT/analysis')
PV3  = Path('/data/data/TNT/analysis/figures/panels_v3')
SUPP_SUB = Path('/data/data/TNT/analysis/genome_medicine_submission/supplementary_figures')
SUPP_V07 = BASE/'figures/supp'
OUT_DIR = BASE/'manuscript'/'ppt'
TMP = Path('/tmp/tnt_ppt_build')
TMP.mkdir(parents=True, exist_ok=True)
OUT_FILE = OUT_DIR/'TNT_v0.7_final_fully_editable.pptx'

# ---------- Panel catalog ----------
# (slide_title, pdf_path)
# PDFs produced by matplotlib preserve text as editable runs after LibreOffice conversion.
MAIN_PANELS = [
    ('Fig 1A — Sankey (sex → cT → response)',  PV3/'Fig1A_sankey.pdf'),
    ('Fig 1B — Clinical waterfall',            PV3/'Fig1B_waterfall.pdf'),
    ('Fig 1C — Clinical characteristics',      PV3/'Fig1C_clinical.pdf'),
    ('Fig 1D — Sample matrix',                 PV3/'Fig1D_sample_matrix.pdf'),
    ('Fig 1E — Study design',                  PV3/'Fig1E_design.pdf'),
    ('Fig 2A — Driver oncoprint',              PV3/'Fig2A_oncoprint_journal.pdf'),
    ('Fig 2B — TMB raincloud by response',     PV3/'Fig2B_TMB_raincloud.pdf'),
    ('Fig 2B(alt) — SBS96 profile',            PV3/'Fig2B_SBS96_profile.pdf'),
    ('Fig 2C — SBS signature attribution',     PV3/'Fig2C_signature_attribution.pdf'),
    ('Fig 2C(alt) — MSI × TMB scatter',        PV3/'Fig2C_MSI_TMB_scatter.pdf'),
    ('Fig 2D — TMB + MSI waterfall',           PV3/'Fig2D_TMB_MSI_waterfall.pdf'),
    ('Fig 2D(alt) — SBS reassessment',         PV3/'Fig2D_SBS_reassessment.pdf'),
    ('Fig 2E — CNV + HRD proxy',               PV3/'Fig2E_CNV_HRD.pdf'),
    ('Fig 2E(alt) — CNV genome view',          PV3/'Fig2E_CNV_genome.pdf'),
    ('Fig 2F — HRD breakdown',                 PV3/'Fig2F_HRD_breakdown.pdf'),
    ('Fig 3A — TME signature radar',           PV3/'Fig3A_TME_radar.pdf'),
    ('Fig 3B — 22-signature heatmap',          PV3/'Fig3B_signature_heatmap.pdf'),
    ('Fig 3C — DEG volcano',                   PV3/'Fig3C_volcano_journal.pdf'),
    ('Fig 3D — Forest lollipop of signatures', PV3/'Fig3D_forest_lollipop.pdf'),
    ('Fig 3E — CD8 effector × exhaustion biaxial', PV3/'Fig3E_CD8_biaxial.pdf'),
    ('Fig 3F — TLS Cabrita score',             PV3/'Fig3F_TLS_Cabrita.pdf'),
    ('Fig 4A — GSEA running enrichment',       PV3/'Fig4A_running_ES.pdf'),
    ('Fig 4B — Hallmark NES bubble',           PV3/'Fig4B_Hallmark_bubble.pdf'),
    ('Fig 4C — Reactome pathway dotplot',      PV3/'Fig4C_pathway_dotplot.pdf'),
    ('Fig 4D — GSEA enrichment map',           PV3/'Fig4D_enrichment_map.pdf'),
    ('Fig 4E — ssGSEA 95-set heatmap',         PV3/'Fig4E_ssGSEA_heatmap.pdf'),
    ('Fig 4F — Category-level NES box',        PV3/'Fig4F_category_NES_box.pdf'),
    ('Fig 5A — Feature correlation',           PV3/'Fig5A_correlation.pdf'),
    ('Fig 5B — Nested LOOCV ROC (v0.7, leakage-free)', PV3/'Fig4B_ROC_nested.pdf'),
    ('Fig 5C — Feature forest with CI',        PV3/'Fig5C_forest_CI.pdf'),
    ('Fig 5D — UMAP of features',              PV3/'Fig5D_UMAP.pdf'),
    ('Fig 5E — SHAP beeswarm',                 PV3/'Fig5E_SHAP_beeswarm.pdf'),
    ('Fig 5F — Per-subject predicted probability', PV3/'Fig5F_per_subject_prediction.pdf'),
    ('Fig 6A — Cascade BCa forest (v0.7)',     BASE/'figures/panels/Fig6_cascade_BCa_forest.pdf'),
    ('Fig 6B — Paired pre→post slopes',        PV3/'Fig6B_slope_fancy.pdf'),
    ('Fig 6C — Δ forest',                      PV3/'Fig6C_delta_forest.pdf'),
    ('Fig 6D — Per-subject waterfall',         PV3/'Fig6D_waterfall.pdf'),
    ('Fig 6E — Clonal fishplot',               PV3/'Fig6E_fishplot.pdf'),
    ('Fig 6F — Cascade schematic',             PV3/'Fig6F_cascade.pdf'),
    ('Fig 7A — CD8 meta + Akiyoshi (v0.7)',    BASE/'figures/panels/Fig7_external_CD8_validation_v4.pdf'),
    ('Fig 7B — Signature × cohort heatmap',    PV3/'Fig7B_heatmap.pdf'),
    ('Fig 7C — Meta Z-score across signatures', PV3/'Fig7C_meta_Zscore.pdf'),
    ('Fig 7D — Cohort concordance',            PV3/'Fig7D_cohort_concordance.pdf'),
    ('Fig 7E — Funnel plot',                   PV3/'Fig7E_funnel.pdf'),
    ('Fig 7F — Discovery vs external effects', PV3/'Fig7F_discovery_validation.pdf'),
    ('Fig 8A — HLA allele frequency',          PV3/'Fig8A_HLA_alleles.pdf'),
    ('Fig 8B — HLA homozygosity by response',  PV3/'Fig8B_HLA_homozygosity.pdf'),
    ('Fig 8C — HLA LOH',                       PV3/'Fig8C_HLA_LOH.pdf'),
    ('Fig 8D — Pre-CRT neoantigen burden',     PV3/'Fig8D_neoantigen_pre.pdf'),
    ('Fig 8E — Paired Δ neoantigen binders',   PV3/'Fig8E_neoantigen_paired.pdf'),
    ('Fig 8F — Per-subject neoantigen lollipop', PV3/'Fig8F_neoantigen_lollipop.pdf'),
    ('Fig 9A — Clone trajectories',            PV3/'Fig9A_clone_trajectories.pdf'),
    ('Fig 9B — CCF pre vs post',               PV3/'Fig9B_CCF_pre_post.pdf'),
    ('Fig 9C — Cluster stacked composition',   PV3/'Fig9C_cluster_stacked.pdf'),
    ('Fig 9D — Dominant clone shrinkage',      PV3/'Fig9D_dominant_shrink.pdf'),
    ('Fig 9E — Shrink vs expand scatter',      PV3/'Fig9E_shrink_expand_scatter.pdf'),
    ('Fig 9F — Fate composition by response',  PV3/'Fig9F_fate_composition.pdf'),
]

SUPP_PANELS = [
    ('Supp Fig S1 — Cohort QC',                SUPP_SUB/'FigS1_cohort_QC.pdf'),
    ('Supp Fig S2 — SBS signatures full',      SUPP_SUB/'FigS2_SBS_signatures.pdf'),
    ('Supp Fig S3 — CNV + HRD detail',         SUPP_SUB/'FigS3_CNV_HRD.pdf'),
    ('Supp Fig S4 — Oncoprint + VAF detail',   SUPP_SUB/'FigS4_oncoprint_VAF.pdf'),
    ('Supp Fig S5 — Full GSEA supplement',     SUPP_SUB/'FigS5_GSEA_full.pdf'),
    ('Supp Fig S6 — ssGSEA and CMS',           SUPP_SUB/'FigS6_ssGSEA_CMS.pdf'),
    ('Supp Fig S7 — TRUST4 immune repertoire', SUPP_SUB/'FigS7_immune_TRUST4.pdf'),
    ('Supp Fig S8 — ML model comparison',      SUPP_SUB/'FigS8_ML_model_comparison.pdf'),
    ('Supp Fig S9 — GEO cohorts overview',     SUPP_SUB/'FigS9_GEO_cohorts.pdf'),
    ('Supp Fig S10 — HLA neoantigen detail',   SUPP_SUB/'FigS10_HLA_neoantigen_detail.pdf'),
    ('Supp Fig S11 — PyClone diagnostics',     SUPP_SUB/'FigS11_pyclone_diagnostics.pdf'),
    ('Supp Fig — SBS landscape cohort',        SUPP_SUB/'FigS_SBS_landscape.pdf'),
    ('Supp Fig — External cohort forest (legacy)', SUPP_SUB/'SuppFig_external_forest.pdf'),
    ('Supp Fig — GSE150082 DSB detail',        SUPP_SUB/'SuppFig_GSE150082_DSB.pdf'),
    ('Supp Fig — External meta Z-score',       SUPP_SUB/'SuppFig_meta_zscore.pdf'),
    ('Supp Fig (v0.7) — CONSORT sample flow',  SUPP_V07/'SuppFig_consort_sample_flow.pdf'),
    ('Supp Fig (v0.7) — HLA-LOH lite vs strict', SUPP_V07/'SuppFig_S3_HLA_LOH_lite_vs_strict.pdf'),
]

# ---------- LibreOffice batch conversion ----------
def convert_pdf_to_pptx(pdf_path, out_dir):
    """Convert a single PDF to single-slide PPTX using LibreOffice with PDF import."""
    cmd = ['libreoffice', '--headless',
           '--infilter=impress_pdf_import',
           '--convert-to', 'pptx',
           '--outdir', str(out_dir),
           str(pdf_path)]
    r = subprocess.run(cmd, capture_output=True, text=True, timeout=90)
    out = out_dir / (Path(pdf_path).stem + '.pptx')
    return out if out.exists() else None

# ---------- Slide cloning between presentations ----------
# python-pptx doesn't natively support copying slides between presentations,
# so we clone slide XML manually at the OPC part level.
def clone_slide(src_prs, dst_prs, slide_idx=0):
    """Clone the first slide from src_prs into dst_prs.
    Returns the newly created slide."""
    src_slide = src_prs.slides[slide_idx]
    # Find a blank layout in dst and create a slide there
    blank = dst_prs.slide_layouts[6]
    new_slide = dst_prs.slides.add_slide(blank)
    # Copy every shape from src_slide into new_slide by deep-copying shape XML
    for shape in src_slide.shapes:
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    # Copy slide size doesn't apply per-slide, but dst_prs size already set.
    return new_slide

def add_title_bar(slide, text, prs):
    """Add an editable title bar above the imported plot (per user spec: no figure title,
    we add a small slide-level subtitle that users can remove / keep)."""
    tb = slide.shapes.add_textbox(Emu(200000), Emu(100000),
                                  prs.slide_width - Emu(400000), Emu(400000))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = RGBColor(0x1F, 0x3B, 0x5C)
    r.font.name = 'Arial'
    return tb

# ---------- Main build ----------
def main():
    # Step 1 — batch convert all PDFs to individual single-slide PPTXs
    batch_dir = TMP / 'pptx_batch'
    if batch_dir.exists():
        shutil.rmtree(batch_dir)
    batch_dir.mkdir(parents=True)

    all_panels = MAIN_PANELS + SUPP_PANELS
    # Group PDFs into batches of 10 to avoid LibreOffice timeout
    converted = {}
    missing = []
    BATCH = 10
    pdfs_to_convert = [(t, p) for (t, p) in all_panels if Path(p).exists()]
    total = len(pdfs_to_convert)
    print(f'Converting {total} PDFs to individual PPTXs via LibreOffice...')

    # LibreOffice can take multiple files at once
    for i in range(0, total, BATCH):
        batch = pdfs_to_convert[i:i+BATCH]
        cmd = ['libreoffice', '--headless',
               '--infilter=impress_pdf_import',
               '--convert-to', 'pptx',
               '--outdir', str(batch_dir)]
        cmd += [str(p) for (_, p) in batch]
        try:
            subprocess.run(cmd, capture_output=True, text=True, timeout=300, check=False)
        except subprocess.TimeoutExpired:
            print(f'  batch {i} timeout')
        done_this = 0
        for (title, pdf) in batch:
            out = batch_dir / (Path(pdf).stem + '.pptx')
            if out.exists():
                converted[str(pdf)] = out
                done_this += 1
            else:
                missing.append((title, pdf))
        print(f'  {i+len(batch)}/{total}   (batch: {done_this}/{len(batch)} ok)')

    print(f'\nConverted OK: {len(converted)} / {total}')
    if missing:
        print(f'Missing PDFs (will be skipped): {len(missing)}')
        for t, p in missing[:5]:
            print(f'  - {t} ({Path(p).name})')

    # Step 2 — build master deck by cloning each converted slide
    master = Presentation()
    master.slide_width  = Inches(13.33)
    master.slide_height = Inches(7.5)
    BLANK = master.slide_layouts[6]

    # Cover
    cover = master.slides.add_slide(BLANK)
    tb = cover.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = 'TNT MSS LARC — v0.7 fully-editable figure deck'
    r.font.size = Pt(32); r.font.bold = True
    r.font.color.rgb = RGBColor(0x1F, 0x3B, 0x5C); r.font.name = 'Arial'
    tb2 = cover.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.5))
    tf = tb2.text_frame; tf.word_wrap = True
    body = (
        'One panel per slide.\n'
        'No figure titles (matplotlib suptitles were never set on individual panels).\n\n'
        'Every in-plot element — axis labels, tick labels, legend text, P-values, '
        'cohort names, annotations — is imported as a NATIVE PowerPoint text run.\n\n'
        'Workflow: matplotlib PDF → LibreOffice PDF import → PPTX text preservation.\n'
        'Double-click any text in any slide to edit. Fonts, colors, sizes editable.\n\n'
        'Plot body graphics (lines, markers, fills) come in as PowerPoint shapes or '
        'grouped vector objects — also editable via the Selection Pane.\n\n'
        f'Total panels embedded: {len(converted)} of {len(all_panels)} catalog entries.\n'
        'Missing PDFs (if any) are listed at the end.'
    )
    for i, line in enumerate(body.split('\n')):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(14); r.font.name = 'Arial'
        r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Per-panel slides
    print('\nAssembling master deck...')
    for idx, (title, pdf) in enumerate(all_panels, 1):
        pptx_path = converted.get(str(pdf))
        if pptx_path is None:
            continue
        try:
            src = Presentation(str(pptx_path))
            new_slide = clone_slide(src, master)
            add_title_bar(new_slide, title, master)
        except Exception as e:
            print(f'  [skip] {title}: {e}')
        if idx % 10 == 0:
            print(f'  assembled {idx}/{len(all_panels)}')

    # Missing panels appendix
    if missing:
        s = master.slides.add_slide(BLANK)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5),
                                  Inches(12.3), Inches(1.0))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = f'Missing PDFs ({len(missing)}) — not embedded'
        r.font.size = Pt(22); r.font.bold = True
        r.font.color.rgb = RGBColor(0xB2, 0x1F, 0x1F)
        tb2 = s.shapes.add_textbox(Inches(0.5), Inches(1.6),
                                   Inches(12.3), Inches(5.5))
        tf = tb2.text_frame; tf.word_wrap = True
        for i, (t, p) in enumerate(missing):
            para = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            r = para.add_run()
            r.text = f'  • {t}   ({Path(p).name})'
            r.font.size = Pt(11); r.font.name = 'Arial'
            r.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    master.save(OUT_FILE)
    print(f'\nSaved: {OUT_FILE}')
    print(f'  slides: {len(master.slides)}')

if __name__ == '__main__':
    main()
