#!/usr/bin/env python3
"""Generate editable PPTX of the graphical abstract by embedding SVG-rendered PNG + overlay text boxes."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

OUT = Path("/data/data/TNT/analysis/figures/graphical_abstract")
PNG = OUT / "GA_v2.png"

prs = Presentation()
prs.slide_width = Inches(18)
prs.slide_height = Inches(8.1)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# full-page background image
slide.shapes.add_picture(str(PNG), Emu(0), Emu(0), prs.slide_width, prs.slide_height)

# add a few editable overlay text boxes for key numbers users might want to tweak
def add_label(left, top, width, height, text, size=11, bold=False, color=RGBColor(0x1F,0x24,0x30), align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

# editable key stats
add_label(0.4, 0.25, 14, 0.5,
    "Radiation-phase molecular response predicts final TNT outcome in MSS LARC",
    size=28, bold=True)
add_label(0.4, 0.72, 12, 0.3,
    "Integrated WES + RNA-seq of 35 locally advanced rectal cancers (pre / post CRT biopsies) — Kim et al., 2026",
    size=14, color=RGBColor(0x6B,0x72,0x80))

# take-home bar
add_label(0.5, 8.45, 14, 0.6,
    "Take-home:  Radiation-phase molecular response — mutation / neoantigen clearance and a reproducible CD8-cytotoxic axis — defines a mid-treatment biomarker window to personalise consolidation chemotherapy in MSS locally advanced rectal cancer.",
    size=14, bold=True, color=RGBColor(0xFF,0xFF,0xFF))

# meta p-value
add_label(12.3, 3.85, 3, 0.3,
    "Z = +2.74   P = 0.006",
    size=11, bold=True, color=RGBColor(0x2A,0x2E,0x3A))

pptx_path = OUT / "GA_v2_editable.pptx"
prs.save(str(pptx_path))
print("wrote:", pptx_path)
