#!/usr/bin/env python3
"""
27_fig9_merge_abcde.py (v0.7.5)

Merge the standalone Fig 9 Panel E deck (1 slide) into the existing
Fig 9 A/B/C/D deck (4 slides) to produce a single 5-panel deck
suitable for submission composite.

Also render each panel to PDF, then stitch into a single composite
Fig9_external_validation.pdf placed into
genome_medicine_submission/main_figures/.
"""
import os, shutil, subprocess
from pptx import Presentation
from pptx.util import Inches, Emu
from copy import deepcopy
from lxml import etree

PPT_DIR = "/data/data/TNT/analysis/260418_add/ppt"
SRC_ABCD = f"{PPT_DIR}/Fig9_external_validation_native_editable.pptx"
SRC_E    = f"{PPT_DIR}/Fig9E_scrt_validation_native_editable.pptx"
OUT_MERGED = f"{PPT_DIR}/Fig9_external_validation_v075_native_editable.pptx"
OUT_SUBMIT_DIR = "/data/data/TNT/analysis/genome_medicine_submission/main_figures"

# ---- 1. Append Panel E slide to the A-D deck ----
prs_abcd = Presentation(SRC_ABCD)
prs_e    = Presentation(SRC_E)

# Ensure the A-D deck has matching slide dimensions as E (both 6.5 x 4.5)
prs_abcd.slide_width  = Emu(int(6.5 * 914400))
prs_abcd.slide_height = Emu(int(4.5 * 914400))

src_slide = prs_e.slides[0]
blank_layout = prs_abcd.slide_layouts[6]
new_slide = prs_abcd.slides.add_slide(blank_layout)

# Copy all shapes from panel E into the new slide
for shp in src_slide.shapes:
    el = deepcopy(shp.element)
    new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')

prs_abcd.save(OUT_MERGED)
print(f"wrote merged deck: {OUT_MERGED}  ({len(prs_abcd.slides)} slides)")

# ---- 2. Render merged deck to PDF ----
subprocess.run([
    "libreoffice", "--headless", "--convert-to", "pdf",
    "--outdir", PPT_DIR, OUT_MERGED
], check=True, capture_output=True, timeout=120)
MERGED_PDF = f"{PPT_DIR}/Fig9_external_validation_v075_native_editable.pdf"
print(f"wrote merged PDF: {MERGED_PDF}")

# ---- 3. Copy the 5-slide PDF as the Fig 9 composite ----
os.makedirs(OUT_SUBMIT_DIR, exist_ok=True)
out_pdf = f"{OUT_SUBMIT_DIR}/Fig9_external_validation.pdf"
shutil.copy(MERGED_PDF, out_pdf)
print(f"copied to {out_pdf}")

# Also keep a png thumbnail of the last panel (E) for quick preview
subprocess.run([
    "pdftoppm", "-png", "-r", "200", "-f", "5", "-l", "5",
    out_pdf, f"{OUT_SUBMIT_DIR}/Fig9_external_validation_panelE"
], check=True, capture_output=True, timeout=60)

# Render all pages as PNGs into a gallery location
subprocess.run([
    "pdftoppm", "-png", "-r", "200",
    out_pdf, f"{OUT_SUBMIT_DIR}/Fig9_external_validation"
], check=True, capture_output=True, timeout=120)
print(f"also wrote PNGs for Fig9 panels to {OUT_SUBMIT_DIR}/")
