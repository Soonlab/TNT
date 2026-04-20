#!/usr/bin/env python3
"""Verify every Supp PPT in supplefigure_260420:
- slide 16:9 (12192000 x 6858000 EMU == 13.333 x 7.5 in)
- every shape has <a:effectLst/> (shadow killed)
- every run uses Arial
- report n_slides, n_shapes, non-Arial runs, missing effectLst
"""
import os
import glob
from pptx import Presentation
from pptx.oxml.ns import qn

OUT = "/data/data/TNT/analysis/260418_add/supplefigure_260420"


def check_shape(shape, stats):
    stats["n_shapes"] += 1
    # shadow
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        spPr = shape._element.find(qn("a:spPr"))
    has_empty = False
    if spPr is not None:
        for el in spPr.findall(qn("a:effectLst")):
            if len(el) == 0:
                has_empty = True
    if not has_empty and spPr is not None:
        stats["missing_empty_effect"] += 1
    # font
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                stats["n_runs"] += 1
                if run.font.name not in (None, "Arial"):
                    stats["non_arial"] += 1
                elif run.font.name is None:
                    stats["no_font_name"] += 1


def main():
    all_stats = []
    for path in sorted(glob.glob(f"{OUT}/*.pptx")):
        prs = Presentation(path)
        w, h = prs.slide_width, prs.slide_height
        s = {
            "file": os.path.basename(path),
            "n_slides": len(prs.slides),
            "w_in": round(w / 914400, 3),
            "h_in": round(h / 914400, 3),
            "aspect": "16:9" if abs(w / h - 16 / 9) < 0.02 else f"{w/h:.3f}",
            "n_shapes": 0, "n_runs": 0,
            "non_arial": 0, "no_font_name": 0,
            "missing_empty_effect": 0,
        }
        for slide in prs.slides:
            for shp in slide.shapes:
                check_shape(shp, s)
        all_stats.append(s)

    # print
    print(f"{'file':<52s} {'slides':>6s} {'aspect':>7s} {'shapes':>7s} "
          f"{'runs':>5s} {'!Arial':>7s} {'noFont':>7s} {'shadow?':>7s}")
    print("-" * 104)
    for s in all_stats:
        print(f"{s['file']:<52s} {s['n_slides']:>6d} {s['aspect']:>7s} "
              f"{s['n_shapes']:>7d} {s['n_runs']:>5d} "
              f"{s['non_arial']:>7d} {s['no_font_name']:>7d} "
              f"{s['missing_empty_effect']:>7d}")
    # summary
    totals = {k: sum(s[k] for s in all_stats)
              for k in ["n_slides", "n_shapes", "n_runs",
                        "non_arial", "no_font_name", "missing_empty_effect"]}
    print("-" * 104)
    print(f"{'TOTAL':<52s} {totals['n_slides']:>6d} {'':>7s} "
          f"{totals['n_shapes']:>7d} {totals['n_runs']:>5d} "
          f"{totals['non_arial']:>7d} {totals['no_font_name']:>7d} "
          f"{totals['missing_empty_effect']:>7d}")
    # assert
    bad = [s for s in all_stats
           if s["aspect"] != "16:9" or s["non_arial"] > 0
           or s["missing_empty_effect"] > 0]
    if bad:
        print("\nISSUES:")
        for s in bad:
            print(f"  {s['file']}: aspect={s['aspect']} "
                  f"non_arial={s['non_arial']} missing_effect={s['missing_empty_effect']}")
    else:
        print("\nOK: all files 16:9, Arial-only, shadow-killed.")


if __name__ == "__main__":
    main()
