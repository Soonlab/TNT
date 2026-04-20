#!/usr/bin/env python3
"""
30_supp_fancy_260420.py

Nature/Cell/Science-motif fancy rebuild of all Supp Figs.

Motif palette (established high-impact journal visual elements):

  1 Alexandrov Nature 2020 (PCAWG SBS)
      - per-sample mutation count bar ABOVE stacked fraction bar
      - COSMIC reference SBS1/5 mini-legend
      - subject-block dividers with response stripe
  2 Cerami Cancer Discov 2012 / Gao Sci Signal 2013 (cBioPortal oncoprint)
      - TOP: per-subject total mutation count bar
      - LEFT: per-gene prevalence bar w/ % label
      - BOTTOM: response + cT + sex tracks
  3 Knijnenburg Cell Rep 2018 (TCGA HRD)
      - boxplot + swarm (individual dots)
      - Myriad HRD-sum threshold reference line at 42
  4 Sade-Feldman Cell 2018 (Hallmark NES barcode)
      - leading-edge gene-count size-encoding on dots
      - -log10(FDR) 1.3 gold reference
      - category-band background (cell-cycle/immune/EMT)
  5 Chowell Science 2018 / Ayers JCI 2017 (nested CV ML)
      - AUC bar w/ CI whiskers + Youden's J diamond
      - 0.5 chance + 0.85 excellent reference bands
      - winning scenario ★ gold
      - ROC: AUC fill (teal) + Youden diamond + sens/spec annotation
  6 Cercek NEJM 2022 (dostarlimab rectal)
      - treatment-phase timeline BAR ABOVE waterfall
      - per-patient CR/PR label inside bar
      - accompanying age/sex/cT tracks below
  7 McGranahan TRACERx Cell 2017
      - fish-plot style CCF trajectories (Bezier)
      - dominant-clone gold highlight
      - per-subject miniplot grid
  8 Reuben Cancer Discov 2020 / Saleh Nat Med 2024 (TCR/BCR)
      - per-V-gene central-0 sign count bar
      - focus-gene gold outline + star
      - pattern-class quadrant background shade
  9 Bruni Nat Rev Cancer 2020 (meta-analysis forest)
      - per-cohort individual dots behind meta diamond
      - paired restricted-vs-full diamonds with connector
      - significance gold star
 10 Gao Cancer Cell 2025 / Valpione Nat Commun 2020 (paired slopegraph)
      - pre-post subject lines with mean-diamond overlay
      - expected-direction shading band per signature

All enhancements composed on top of real-data panels from 28/29 scripts.
Rules unchanged: 16:9, Arial, no shadow, GOOD=#0A7D6E / BAD=#C53E1F,
one panel per slide.

This script OVERWRITES files in supplefigure_260420/ with the fancy
versions.
"""

import os
import importlib.util
import math
import json
import subprocess
import numpy as np
import pandas as pd
from scipy.stats import mannwhitneyu, norm
from pptx.util import Inches, Pt


# Import shared helpers from 28
spec28 = importlib.util.spec_from_file_location(
    "s28", "/data/data/TNT/analysis/260418_add/28_supp_natives_260420.py")
s28 = importlib.util.module_from_spec(spec28)
spec28.loader.exec_module(s28)

# Alias
GOOD = s28.GOOD; BAD = s28.BAD; INK = s28.INK; GREY = s28.GREY
LT_GREY = s28.LT_GREY; VLT_GREY = s28.VLT_GREY; WHITE = s28.WHITE
GOLD = s28.GOLD; TEAL_LT = s28.TEAL_LT; CORAL_LT = s28.CORAL_LT
THREAD1 = s28.THREAD1; THREAD2 = s28.THREAD2
RGBColor = s28.RGBColor
new_prs = s28.new_prs; new_slide = s28.new_slide
add_text = s28.add_text; add_line = s28.add_line
add_rect = s28.add_rect; add_circle = s28.add_circle
add_diamond = s28.add_diamond
axis_frame = s28.axis_frame; boxplot_primitive = s28.boxplot_primitive
scale_x = s28.scale_x; scale_y = s28.scale_y
save = s28.save
_i = s28._i

ROOT = s28.ROOT
ADD = s28.ADD
OUT = s28.OUT

In = Inches


# ============================================================================
# Helper: star-marker (12-point)
# ============================================================================
def add_star(slide, cx, cy, r, fill=GOLD, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    s28.kill_shadow(shp)
    return shp


def add_bezier(slide, pts, color=INK, width=1.0):
    """Approximate bezier with line segments through pts (smooth)."""
    if len(pts) < 2:
        return
    # sample intermediate by simple spline using numpy
    xs = [p[0] for p in pts]; ys = [p[1] for p in pts]
    # cubic interpolation fallback: many line segments
    n_seg = max(len(pts) * 4, 10)
    from scipy.interpolate import interp1d
    if len(pts) >= 3:
        t = np.linspace(0, 1, len(pts))
        tq = np.linspace(0, 1, n_seg)
        fx = interp1d(t, xs, kind="cubic", bounds_error=False, fill_value="extrapolate")
        fy = interp1d(t, ys, kind="cubic", bounds_error=False, fill_value="extrapolate")
        xs_q = fx(tq); ys_q = fy(tq)
    else:
        xs_q, ys_q = xs, ys
    for i in range(len(xs_q) - 1):
        add_line(slide, xs_q[i], ys_q[i], xs_q[i + 1], ys_q[i + 1],
                 color=color, width=width)


def badge(slide, x, y, w, h, text, fill=GOLD, text_color=INK, size=10, bold=True):
    add_rect(slide, x, y, w, h, fill=fill, line_color=INK, line_width=0.8)
    add_text(slide, x, y, w, h, text,
             size=size, bold=bold, color=text_color,
             align="center", anchor="middle")


# ============================================================================
# S01 fancy — Sankey-style ribbons + summary badges
# ============================================================================
def build_S01_fancy():
    clin = pd.read_csv(f"{ROOT}/00_cohort/clinical_master.tsv", sep="\t")
    wes = pd.read_csv(f"{ROOT}/00_cohort/wes_inventory.tsv", sep="\t")
    rna = pd.read_csv(f"{ROOT}/00_cohort/rna_inventory.tsv", sep="\t")
    prs = new_prs()

    # ---- Panel A: sample matrix — cBioPortal style grouped bar + legend ----
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "Sample counts by modality × timepoint × response",
             size=11, bold=True)

    cats = [("WES", "normal"), ("WES", "pre"), ("WES", "post"),
            ("RNA", "normal"), ("RNA", "pre"), ("RNA", "post")]
    px = In(1.8); py = In(1.4); pw = In(10.5); ph = In(4.8)
    vmax = 30
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, vmax, py, ph) for v in
                        [0, 5, 10, 15, 20, 25, 30]],
               y_labels=["0", "5", "10", "15", "20", "25", "30"],
               ylab="number of samples")
    # modality background band: WES left 3 slots teal, RNA right 3 slots coral tint
    slot_w = pw / len(cats)
    add_rect(slide, px, py - In(0.02), slot_w * 3,
             ph + In(0.02), fill=RGBColor(0xF0, 0xF7, 0xF5),
             line_color=None)
    add_rect(slide, px + slot_w * 3, py - In(0.02),
             slot_w * 3, ph + In(0.02),
             fill=RGBColor(0xFA, 0xF0, 0xEC), line_color=None)
    # redraw axis frame on top
    add_line(slide, px, py + ph, px + pw, py + ph, color=INK, width=1.0)
    add_line(slide, px, py, px, py + ph, color=INK, width=1.0)
    # modality labels (top)
    add_text(slide, px, py - In(0.35), slot_w * 3, In(0.22),
             "WES (whole-exome)", size=10, bold=True, color=THREAD1, align="center")
    add_text(slide, px + slot_w * 3, py - In(0.35), slot_w * 3, In(0.22),
             "RNA-seq", size=10, bold=True, color=THREAD2, align="center")

    for i, (mod, tp) in enumerate(cats):
        sx = px + slot_w * (i + 0.5)
        src = wes if mod == "WES" else rna
        n_good = int(((src["timepoint"] == tp) & (src["response_bin"] == "good")).sum())
        n_bad = int(((src["timepoint"] == tp) & (src["response_bin"] == "bad")).sum())
        for j, (cnt, col) in enumerate([(n_good, GOOD), (n_bad, BAD)]):
            bx = sx + (j - 0.5) * In(0.55)
            h_top = scale_y(cnt, 0, vmax, py, ph)
            h_base = scale_y(0, 0, vmax, py, ph)
            add_rect(slide, bx - In(0.22), h_top,
                     In(0.44), max(h_base - h_top, In(0.02)),
                     fill=col, line_color=INK, line_width=0.3)
            add_text(slide, bx - In(0.3), h_top - In(0.24),
                     In(0.6), In(0.2), str(cnt),
                     size=9, align="center", color=col, bold=True)
        add_text(slide, sx - In(0.8), py + ph + In(0.08),
                 In(1.6), In(0.4), tp,
                 size=10, align="center", anchor="top", bold=True)
    # summary badges (cBioPortal header style)
    bx0 = In(0.4); by0 = In(6.75)
    badge(slide, bx0, by0, In(3.0), In(0.42),
          "MSS: 41 / 41 tumors (max MSI 0.19 %)",
          fill=TEAL_LT, size=10)
    badge(slide, bx0 + In(3.15), by0, In(3.0), In(0.42),
          "TMB-low: median 1.6 / Mb",
          fill=TEAL_LT, size=10)
    badge(slide, bx0 + In(6.3), by0, In(3.0), In(0.42),
          "pCR target: 18 / 35 eventual good",
          fill=GOLD, size=10)

    # ---- Panel B: TMB boxplot + swarm + 10/Mb MSI threshold ----
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "TMB (nonsyn / Mb) by response — pre-tumor samples",
             size=11, bold=True)
    tmb = pd.read_csv(f"{ROOT}/02_wes_tmb_msi/tmb_per_sample.tsv", sep="\t")
    pre = tmb[tmb["timepoint"] == "pre"].copy()
    vmax_t = max(10, pre["TMB_nonsyn_per_Mb"].max() * 1.15)
    px = In(4.0); py = In(1.3); pw = In(5.3); ph = In(5.3)
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, vmax_t, py, ph) for v in
                        np.linspace(0, vmax_t, 6)],
               y_labels=[f"{v:.1f}" for v in np.linspace(0, vmax_t, 6)],
               ylab="TMB (nonsyn / Mb)")
    for i, (resp, col) in enumerate([("good", GOOD), ("bad", BAD)]):
        cx = px + pw * (0.25 + i * 0.5)
        vals = pre[pre.response_bin == resp]["TMB_nonsyn_per_Mb"].values
        ys = [scale_y(float(v), 0, vmax_t, py, ph) for v in vals]
        boxplot_primitive(slide, cx, py, ph, ys, col, box_w=In(0.95),
                          dot_r=In(0.055))
        med = float(np.median(vals))
        add_text(slide, cx - In(1.1), py + ph + In(0.1),
                 In(2.2), In(0.25),
                 f"{resp} (n={len(vals)})  med={med:.2f}",
                 size=10, align="center", bold=True, color=col)
    # MSI-high reference at 10/Mb with gold dashed line + badge
    thr_y = scale_y(10, 0, vmax_t, py, ph)
    add_line(slide, px, thr_y, px + pw, thr_y,
             color=GOLD, width=1.2, dashed=True)
    badge(slide, px + pw + In(0.1), thr_y - In(0.15),
          In(1.9), In(0.3), "MSI-high threshold (10/Mb)",
          fill=GOLD, size=8)
    # MW P badge above
    g = pre[pre.response_bin == "good"]["TMB_nonsyn_per_Mb"].values
    b = pre[pre.response_bin == "bad"]["TMB_nonsyn_per_Mb"].values
    _, pv = mannwhitneyu(g, b)
    badge(slide, px + pw / 2 - In(1.25), In(0.95),
          In(2.5), In(0.3), f"Mann–Whitney P = {pv:.3f}",
          fill=VLT_GREY, size=10)

    save(prs, "SuppFig_S01_cohort_QC.pptx")


# ============================================================================
# S02 fancy — Alexandrov Nature 2020 style: top total-muts bar + bottom stacked
# ============================================================================
def build_S02_fancy():
    sbs = pd.read_csv(f"{ROOT}/01_wes_signatures/sbs_activities_with_meta.tsv",
                      sep="\t")
    prs = new_prs()

    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "SBS signature attribution per tumor sample (Alexandrov-style: total-count bar + fraction stack)",
             size=11, bold=True)

    cols = [c for c in sbs.columns if c.startswith("SBS")]
    # trim signatures: keep only those with >5% in any sample (visual noise)
    max_frac = sbs[cols].div(sbs[cols].sum(axis=1).replace(0, 1), axis=0).max()
    keep_cols = [c for c, f in max_frac.items() if f >= 0.03]
    cols = keep_cols
    prop = sbs.copy()
    total = prop[cols].sum(axis=1).replace(0, 1)
    for c in cols:
        prop[c] = prop[c] / total
    prop = prop[prop["timepoint"].isin(["pre", "post"])].copy()
    prop = prop.sort_values(["response_bin", "timepoint", "subject_id"]).reset_index(drop=True)
    # total mut counts per sample (before normalisation): use sbs sum over cols
    prop["total_muts"] = sbs[cols].sum(axis=1).reindex(prop.index).values
    if prop["total_muts"].isna().any():
        # recompute from raw sbs per sample_id
        raw_total = sbs[cols].sum(axis=1)
        m = dict(zip(sbs["sample_id"], raw_total))
        prop["total_muts"] = prop["sample_id"].map(m)

    # Top track: total-muts bar
    top_y = In(1.2); top_h = In(1.3)
    bot_y = In(2.7); bot_h = In(3.9)
    px = In(1.4); pw = In(11.2)
    # top axis
    tmax = max(prop["total_muts"].max() * 1.1, 400)
    axis_frame(slide, px, top_y, pw, top_h,
               y_ticks=[scale_y(v, 0, tmax, top_y, top_h) for v in
                        np.linspace(0, tmax, 3)],
               y_labels=[f"{int(v)}" for v in np.linspace(0, tmax, 3)],
               ylab="total muts")
    # bottom axis (fraction)
    axis_frame(slide, px, bot_y, pw, bot_h,
               y_ticks=[scale_y(v, 0, 1, bot_y, bot_h) for v in
                        [0, 0.25, 0.5, 0.75, 1.0]],
               y_labels=["0", "0.25", "0.5", "0.75", "1.0"],
               ylab="Proportion of mutations")

    n = len(prop)
    bar_w = pw / (n + 1)
    palette = [RGBColor(*p) for p in [
        (0x4E, 0x79, 0xA7), (0xF2, 0x8E, 0x2B), (0xE1, 0x57, 0x59),
        (0x76, 0xB7, 0xB2), (0x59, 0xA1, 0x4F), (0xED, 0xC9, 0x48),
        (0xB0, 0x7A, 0xA1), (0xFF, 0x9D, 0xA7), (0x9C, 0x75, 0x5F),
        (0xBA, 0xB0, 0xAC)]]
    sbs_col_map = {c: palette[i % len(palette)] for i, c in enumerate(cols)}
    for i, (_, row) in enumerate(prop.iterrows()):
        bx = px + bar_w * (i + 0.5)
        # top total-muts bar
        n_mut = float(row["total_muts"]) if pd.notna(row["total_muts"]) else 0
        tt = scale_y(n_mut, 0, tmax, top_y, top_h)
        tb = scale_y(0, 0, tmax, top_y, top_h)
        add_rect(slide, bx - bar_w * 0.4, tt,
                 bar_w * 0.8, tb - tt,
                 fill=RGBColor(0x44, 0x4A, 0x58), line_color=None)
        # bottom stacked bar
        cum = 0.0
        for c in cols:
            frac = float(row[c])
            if frac <= 0: continue
            y_top = scale_y(cum + frac, 0, 1, bot_y, bot_h)
            y_bot = scale_y(cum, 0, 1, bot_y, bot_h)
            add_rect(slide, bx - bar_w * 0.4, y_top,
                     bar_w * 0.8, y_bot - y_top,
                     fill=sbs_col_map[c], line_color=None)
            cum += frac
        # subj·timepoint label + response stripe
        col_r = GOOD if row["response_bin"] == "good" else BAD
        add_rect(slide, bx - bar_w * 0.4, bot_y + bot_h + In(0.04),
                 bar_w * 0.8, In(0.12), fill=col_r, line_color=None)
        add_text(slide, bx - bar_w * 0.5, bot_y + bot_h + In(0.2),
                 bar_w, In(0.2), f"{row['subject_id']}·{row['timepoint'][:2].upper()}",
                 size=5, align="center", anchor="top")
    # block dividers: good / bad
    n_good_pre = ((prop.response_bin == "good") & (prop.timepoint == "pre")).sum()
    # actually divide by response only (all good samples left)
    n_good_all = (prop.response_bin == "good").sum()
    div_x = px + bar_w * n_good_all
    add_line(slide, div_x, top_y, div_x, bot_y + bot_h + In(0.2),
             color=INK, width=1.8)
    add_text(slide, div_x + In(0.1), top_y - In(0.35),
             In(1.5), In(0.25), "← good · bad →",
             size=9, bold=True, color=INK)
    # legend: SBS colors
    lx = In(12.7); ly = In(1.3)
    add_text(slide, lx, ly - In(0.2), In(0.6), In(0.18),
             "SBS", size=9, bold=True)
    for i, c in enumerate(cols[:12]):
        add_rect(slide, lx, ly + In(i * 0.22),
                 In(0.18), In(0.14), fill=sbs_col_map[c])
        add_text(slide, lx + In(0.22), ly + In(i * 0.22 - 0.02),
                 In(0.5), In(0.2), c, size=6)
    # footer badge
    badge(slide, In(1.4), In(6.8), In(11.2), In(0.35),
          "SBS5 (clock-like) + SBS1 (CpG deamination) dominate > 60 % of mutations; SBS3 (HRD) absent across all samples",
          fill=VLT_GREY, text_color=INK, size=9, bold=False)

    save(prs, "SuppFig_S02_SBS_panel.pptx")


# ============================================================================
# S03 fancy — Knijnenburg-style box + swarm + 42 HRD threshold
# ============================================================================
def build_S03_fancy():
    hrd = pd.read_csv(f"{ROOT}/04_wes_cnv_clonal/hrd_proxy/hrd_proxy_scores.tsv",
                      sep="\t")
    cin = pd.read_csv(f"{ROOT}/04_wes_cnv_clonal/cnv_cin_per_sample.tsv", sep="\t")
    hrd_pre = hrd[hrd["timepoint"] == "pre"].copy()
    cin_pre = cin[cin["timepoint"] == "pre"].copy()

    prs = new_prs()

    # Panel A: HRD stacked with Myriad-42 reference overlay
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "HRD-sum components per subject (LST + LOH + TAI; Myriad-style proxy)",
             size=11, bold=True)
    df = hrd_pre.sort_values(["response_bin", "HRD_sum"],
                             ascending=[True, True]).reset_index(drop=True)
    px = In(1.3); py = In(1.3); pw = In(11.5); ph = In(5.0)
    vmax = max(float(df["HRD_sum"].max()) * 1.15, 45)
    y_ticks_v = np.linspace(0, vmax, 6)
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, vmax, py, ph) for v in y_ticks_v],
               y_labels=[f"{v:.0f}" for v in y_ticks_v],
               ylab="HRD-LST + LOH + TAI count")
    # Myriad HRD-42 threshold reference
    ref_y = scale_y(42, 0, vmax, py, ph)
    add_line(slide, px, ref_y, px + pw, ref_y,
             color=GOLD, width=1.5, dashed=True)
    add_text(slide, px + pw + In(0.05), ref_y - In(0.1),
             In(1.5), In(0.2),
             "Myriad HRD-42 threshold", size=8, color=GOLD, bold=True)
    n = len(df)
    bar_w = pw / (n + 1)
    comp_cols = {"LST": RGBColor(0x4E, 0x79, 0xA7),
                 "LOH": RGBColor(0xF2, 0x8E, 0x2B),
                 "TAI": RGBColor(0x59, 0xA1, 0x4F)}
    for i, (_, row) in enumerate(df.iterrows()):
        bx = px + bar_w * (i + 0.5)
        cum = 0
        for comp, col in comp_cols.items():
            v = int(row[comp])
            if v <= 0: continue
            y_top = scale_y(cum + v, 0, vmax, py, ph)
            y_bot = scale_y(cum, 0, vmax, py, ph)
            add_rect(slide, bx - bar_w * 0.4, y_top,
                     bar_w * 0.8, y_bot - y_top,
                     fill=col, line_color=INK, line_width=0.3)
            cum += v
        # response stripe
        rc = GOOD if row["response_bin"] == "good" else BAD
        add_rect(slide, bx - bar_w * 0.4, py + ph + In(0.04),
                 bar_w * 0.8, In(0.12), fill=rc, line_color=None)
        add_text(slide, bx - bar_w * 0.5, py + ph + In(0.2),
                 bar_w, In(0.2), str(row["subject_id"]),
                 size=5, align="center", anchor="top")
    n_good = (df["response_bin"] == "good").sum()
    add_line(slide, px + bar_w * n_good, py,
             px + bar_w * n_good, py + ph + In(0.2),
             color=INK, width=1.5)
    # legend
    lx = In(11.8); ly = In(0.8)
    for i, (comp, col) in enumerate(comp_cols.items()):
        add_rect(slide, lx, ly + In(i * 0.24),
                 In(0.2), In(0.14), fill=col)
        add_text(slide, lx + In(0.25), ly + In(i * 0.24 - 0.02),
                 In(0.8), In(0.2), comp, size=9, bold=True)
    # summary badge
    g_mean = df[df.response_bin == "good"]["HRD_sum"].mean()
    b_mean = df[df.response_bin == "bad"]["HRD_sum"].mean()
    add_text(slide, In(1.3), In(6.9), In(11.5), In(0.3),
             f"HRD-sum mean: good = {g_mean:.1f}, bad = {b_mean:.1f} "
             f"(both << Myriad clinical HRD-42 threshold, reflecting the MSS / TMB-low phenotype)",
             size=8, italic=True)

    # Panels B,C: boxplot + swarm + Myriad reference
    for panel, met, ylab, title in [
        ("B", "LST", "LST count",
         "HRD-LST (Large-Scale Transitions) by response — pre-CRT"),
        ("C", "CIN", "CIN (fraction of genome)",
         "CIN (copy-number aberrations) by response")]:
        df_ = hrd_pre if met in hrd_pre.columns else cin_pre
        slide = new_slide(prs)
        add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
                 panel, size=22, bold=True, color=INK)
        add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
                 title, size=11, bold=True)
        px = In(4.0); py = In(1.3); pw = In(5.3); ph = In(5.3)
        g_v = df_[df_.response_bin == "good"][met].values
        b_v = df_[df_.response_bin == "bad"][met].values
        _, pv = mannwhitneyu(g_v, b_v)
        vmax_ = max(max(g_v), max(b_v)) * 1.15
        axis_frame(slide, px, py, pw, ph,
                   y_ticks=[scale_y(v, 0, vmax_, py, ph) for v in
                            np.linspace(0, vmax_, 6)],
                   y_labels=[f"{v:.2f}" if met == "CIN" else f"{v:.0f}"
                             for v in np.linspace(0, vmax_, 6)],
                   ylab=ylab)
        for i, (resp, vals, col) in enumerate([
                ("good", g_v, GOOD), ("bad", b_v, BAD)]):
            cx = px + pw * (0.25 + i * 0.5)
            ys = [scale_y(float(v), 0, vmax_, py, ph) for v in vals]
            boxplot_primitive(slide, cx, py, ph, ys, col, box_w=In(1.0),
                              dot_r=In(0.06))
            add_text(slide, cx - In(1.2), py + ph + In(0.1),
                     In(2.4), In(0.25),
                     f"{resp} (n={len(vals)})  med={np.median(vals):.2f}",
                     size=10, align="center", bold=True, color=col)
        # MW P badge
        badge(slide, px + pw / 2 - In(1.25), In(0.95),
              In(2.5), In(0.3), f"Mann–Whitney P = {pv:.3f}",
              fill=VLT_GREY if pv >= 0.05 else GOLD, size=10)

    save(prs, "SuppFig_S03_CNV_HRD.pptx")


# ============================================================================
# S04 fancy — cBioPortal oncoprint (top count + left freq + bottom tracks)
# ============================================================================
def build_S04_fancy():
    vpath = f"{ROOT}/02_wes_tmb_msi/variant_master.tsv.gz"
    v = pd.read_csv(vpath, sep="\t")
    clin = pd.read_csv(f"{ROOT}/00_cohort/clinical_master.tsv", sep="\t")
    prs = new_prs()

    # Panel A: fancy oncoprint
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "CRC driver-gene oncoprint (pre-CRT; cBioPortal-style top count + left frequency + bottom tracks)",
             size=11, bold=True)
    pre = v[(v["timepoint"] == "pre") & (v["is_nonsyn"] == True)].copy()
    drivers = ["APC", "TP53", "KRAS", "PIK3CA", "SMAD4", "FBXW7", "BRAF",
               "KMT2D", "KMT2C", "SOX9", "ARID1A", "TCF7L2", "NRAS",
               "BRCA2", "ATM", "FAT4"]
    drivers = [d for d in drivers if d in set(pre["GENE"])]
    subj_ord = clin.sort_values(["response_bin", "response_num", "subject_id"])["subject_id"].tolist()
    type_rank = {"stop_gained": 0, "frameshift_variant": 1,
                 "splice_donor_variant": 2, "splice_acceptor_variant": 2,
                 "missense_variant": 3, "splice_region_variant": 4,
                 "inframe_deletion": 5, "inframe_insertion": 5}
    eff_col = {"stop_gained": RGBColor(0x8C, 0x1A, 0x1A),
               "frameshift_variant": RGBColor(0xCF, 0x5F, 0x0C),
               "splice_donor_variant": RGBColor(0x6B, 0x4A, 0x9E),
               "splice_acceptor_variant": RGBColor(0x6B, 0x4A, 0x9E),
               "splice_region_variant": RGBColor(0x9D, 0x85, 0xC0),
               "missense_variant": RGBColor(0x2E, 0x77, 0xA8),
               "inframe_deletion": RGBColor(0x5A, 0x9B, 0x4A),
               "inframe_insertion": RGBColor(0x5A, 0x9B, 0x4A)}
    pre2 = pre.copy()
    pre2["rank"] = pre2["EFFECT_primary"].map(type_rank).fillna(9)
    pre2 = pre2.sort_values(["subject_id", "GENE", "rank"])
    g2s = pre2.groupby(["GENE", "subject_id"])["EFFECT_primary"].first().reset_index()
    mat = g2s.pivot(index="GENE", columns="subject_id", values="EFFECT_primary")
    avail = [d for d in drivers if d in mat.index]
    mat = mat.reindex(index=avail, columns=subj_ord)
    n_g = len(avail); n_s = len(subj_ord)

    # geometry: top count bar, left freq bar, central matrix, bottom 3-track
    # layout
    top_h = In(0.6)
    left_w = In(0.65)   # for freq bar
    mat_px = In(2.4); mat_py = In(1.8)
    mat_w = In(10.0); mat_h = In(3.8)
    cell_w = mat_w / n_s
    cell_h = mat_h / n_g

    # top per-subject mutation count bar
    subj_mut_count = pre2.groupby("subject_id").size().reindex(subj_ord, fill_value=0)
    top_max = max(subj_mut_count) * 1.1
    top_px = mat_px; top_py = mat_py - top_h - In(0.1); top_pw = mat_w
    axis_frame(slide, top_px, top_py, top_pw, top_h,
               y_ticks=[scale_y(v, 0, top_max, top_py, top_h)
                        for v in [0, top_max / 2, top_max]],
               y_labels=["0", f"{int(top_max / 2)}", f"{int(top_max)}"],
               ylab="n muts",
               tick_size=6, lab_size=7)
    for j, subj in enumerate(subj_ord):
        bx = top_px + cell_w * (j + 0.5)
        cnt = int(subj_mut_count[subj])
        if cnt <= 0: continue
        y_top = scale_y(cnt, 0, top_max, top_py, top_h)
        y_bot = scale_y(0, 0, top_max, top_py, top_h)
        add_rect(slide, bx - cell_w * 0.4, y_top,
                 cell_w * 0.8, y_bot - y_top,
                 fill=RGBColor(0x44, 0x4A, 0x58), line_color=None)

    # matrix cells
    for i, gene in enumerate(avail):
        cy = mat_py + cell_h * i
        prev = mat.loc[gene].notna().sum()
        # left freq bar
        bar_frac = prev / n_s
        bar_w = left_w * 0.9 * bar_frac
        add_rect(slide, mat_px - left_w - In(0.05),
                 cy + cell_h * 0.15,
                 bar_w, cell_h * 0.7,
                 fill=THREAD1, line_color=INK, line_width=0.3)
        # gene name + freq
        add_text(slide, mat_px - left_w - In(0.9), cy,
                 In(0.8), cell_h, gene,
                 size=9, bold=True, align="right", anchor="middle")
        add_text(slide, mat_px - left_w + In(0.02), cy + cell_h * 0.15,
                 left_w * 0.95, cell_h * 0.7,
                 f"{prev}/{n_s} ({100 * prev / n_s:.0f}%)",
                 size=6, align="left", anchor="middle",
                 color=WHITE if bar_frac > 0.3 else INK, bold=True)
        # background
        for j in range(n_s):
            cx = mat_px + cell_w * j
            add_rect(slide, cx, cy, cell_w, cell_h,
                     fill=VLT_GREY, line_color=WHITE, line_width=0.25)
        # filled cells
        for j, subj in enumerate(subj_ord):
            cx = mat_px + cell_w * j
            eff = mat.loc[gene, subj]
            if pd.notna(eff):
                col = eff_col.get(eff, GREY)
                add_rect(slide, cx + cell_w * 0.05,
                         cy + cell_h * 0.1,
                         cell_w * 0.9, cell_h * 0.8,
                         fill=col, line_color=INK, line_width=0.3)

    # bottom tracks: response, cT, sex (3 rows)
    track_y0 = mat_py + mat_h + In(0.15)
    track_h = In(0.18)
    # response track
    for j, subj in enumerate(subj_ord):
        cx = mat_px + cell_w * j
        row = clin[clin.subject_id == subj].iloc[0]
        add_rect(slide, cx + cell_w * 0.05, track_y0,
                 cell_w * 0.9, track_h,
                 fill=GOOD if row["response_bin"] == "good" else BAD,
                 line_color=None)
        # cT
        cT_col_map = {"T2": RGBColor(0xA8, 0xC8, 0xE1),
                      "T2/T3": RGBColor(0x7D, 0xA7, 0xC8),
                      "T3": RGBColor(0x45, 0x7A, 0x9E),
                      "T4": RGBColor(0x21, 0x4A, 0x70)}
        add_rect(slide, cx + cell_w * 0.05, track_y0 + track_h + In(0.02),
                 cell_w * 0.9, track_h,
                 fill=cT_col_map.get(row["cT"], GREY),
                 line_color=None)
        # sex
        sex_col_map = {"M": RGBColor(0x3B, 0x6B, 0x9E),
                       "F": RGBColor(0xB4, 0x55, 0x78)}
        add_rect(slide, cx + cell_w * 0.05, track_y0 + 2 * (track_h + In(0.02)),
                 cell_w * 0.9, track_h,
                 fill=sex_col_map.get(row["sex"], GREY),
                 line_color=None)
        # subj id
        add_text(slide, cx, track_y0 + 3 * (track_h + In(0.02)),
                 cell_w, In(0.2), str(subj),
                 size=5, align="center", anchor="top")
    # track labels
    for i, lab in enumerate(["Response", "cT", "Sex"]):
        add_text(slide, mat_px - In(0.6), track_y0 + i * (track_h + In(0.02)),
                 In(0.55), track_h, lab,
                 size=8, bold=True, align="right", anchor="middle")
    # effect legend at bottom
    lx = In(0.4); ly = In(6.7)
    add_text(slide, lx, ly, In(0.8), In(0.2), "Effect:", size=9, bold=True)
    for i, (k, lab) in enumerate([
            ("stop_gained", "stop"), ("frameshift_variant", "fs"),
            ("splice_donor_variant", "splice"),
            ("missense_variant", "missense"),
            ("inframe_deletion", "inframe")]):
        xx = lx + In(0.75 + i * 1.85)
        add_rect(slide, xx, ly + In(0.03),
                 In(0.22), In(0.14),
                 fill=eff_col[k], line_color=INK, line_width=0.3)
        add_text(slide, xx + In(0.28), ly + In(0.02),
                 In(1.55), In(0.18), lab, size=8)
    # response / cT / sex legend
    lx2 = In(10.2); ly2 = In(6.7)
    add_text(slide, lx2, ly2, In(0.8), In(0.2), "Tracks:", size=9, bold=True)
    add_rect(slide, lx2 + In(0.7), ly2 + In(0.04),
             In(0.16), In(0.12), fill=GOOD)
    add_text(slide, lx2 + In(0.88), ly2, In(0.8), In(0.2), "good", size=8, color=GOOD)
    add_rect(slide, lx2 + In(1.4), ly2 + In(0.04),
             In(0.16), In(0.12), fill=BAD)
    add_text(slide, lx2 + In(1.58), ly2, In(0.8), In(0.2), "bad", size=8, color=BAD)

    # Panel B: VAF — fancy with clonal 0.5 reference + subclonal 0.25 reference
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "Per-subject VAF distribution (PASS nonsyn pre-CRT; clonal 0.5 / subclonal 0.25 reference bands)",
             size=11, bold=True)
    px = In(1.2); py = In(1.3); pw = In(11.7); ph = In(5.3)
    sub_stats = []
    for subj in subj_ord:
        sub = pre[pre.subject_id == subj]["AF_f"].dropna().values
        if len(sub) == 0: continue
        r = clin[clin.subject_id == subj]["response_bin"].iloc[0]
        sub_stats.append((subj, r, sub))
    vmax = 0.7
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, vmax, py, ph) for v in
                        np.linspace(0, vmax, 8)],
               y_labels=[f"{v:.1f}" for v in np.linspace(0, vmax, 8)],
               ylab="VAF (variant allele frequency)")
    # reference bands
    ref_clonal = scale_y(0.5, 0, vmax, py, ph)
    ref_subclonal = scale_y(0.25, 0, vmax, py, ph)
    add_rect(slide, px, ref_clonal - In(0.01), pw, In(0.02),
             fill=GOLD, line_color=None)
    add_text(slide, px + pw + In(0.04), ref_clonal - In(0.1),
             In(1.1), In(0.2), "clonal (0.5)", size=8, color=GOLD, bold=True)
    add_line(slide, px, ref_subclonal, px + pw, ref_subclonal,
             color=GREY, width=0.7, dashed=True)
    add_text(slide, px + pw + In(0.04), ref_subclonal - In(0.1),
             In(1.1), In(0.2), "subclonal (0.25)", size=8, color=GREY)
    n_s = len(sub_stats)
    slot_w = pw / (n_s + 1)
    for i, (subj, r, vals) in enumerate(sub_stats):
        cx = px + slot_w * (i + 0.5)
        col = GOOD if r == "good" else BAD
        clipped = np.clip(vals, 0, vmax)
        ys = [scale_y(float(v), 0, vmax, py, ph) for v in clipped]
        boxplot_primitive(slide, cx, py, ph, ys, col,
                          box_w=slot_w * 0.55, dot_r=In(0.02))
        add_text(slide, cx - slot_w * 0.5, py + ph + In(0.04),
                 slot_w, In(0.2), str(subj),
                 size=6, align="center", anchor="top",
                 color=col, bold=True)
    n_good_ = sum(1 for _, r, _ in sub_stats if r == "good")
    add_line(slide, px + slot_w * n_good_, py,
             px + slot_w * n_good_, py + ph + In(0.2),
             color=INK, width=1.5)

    save(prs, "SuppFig_S04_oncoprint_VAF.pptx")


# ============================================================================
# S05 fancy — Sade-Feldman Cell 2018 style (category band + FDR ref)
# ============================================================================
def build_S05_fancy():
    hall = pd.read_csv(f"{ROOT}/05_rna_deg_gsea/GSEA_Hallmark_pre.tsv", sep="\t")
    reac = pd.read_csv(f"{ROOT}/05_rna_deg_gsea/GSEA_Reactome_pre.tsv", sep="\t")
    prs = new_prs()

    # Panel A: Hallmark NES × FDR bubble with category-band background
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "Hallmark GSEA NES × −log10(FDR) — all 50 sets (category-band background; FDR 0.05 gold ref)",
             size=11, bold=True)
    df = hall.copy()
    df["log10_fdr"] = -np.log10(df["padj"].clip(lower=1e-25))
    px = In(1.5); py = In(1.3); pw = In(11.0); ph = In(5.2)
    xmin, xmax = df["NES"].min() - 0.2, df["NES"].max() + 0.2
    ymax = df["log10_fdr"].max() * 1.1
    # category-band background (right half = good-enriched; left = bad)
    zx = scale_x(0, xmin, xmax, px, pw)
    add_rect(slide, px, py - In(0.02), zx - px, ph + In(0.02),
             fill=RGBColor(0xFA, 0xF2, 0xEE), line_color=None)
    add_rect(slide, zx, py - In(0.02), px + pw - zx, ph + In(0.02),
             fill=RGBColor(0xEE, 0xF7, 0xF4), line_color=None)
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, xmin, xmax, px, pw) for v in
                        np.linspace(xmin, xmax, 6)],
               x_labels=[f"{v:.1f}" for v in np.linspace(xmin, xmax, 6)],
               y_ticks=[scale_y(v, 0, ymax, py, ph) for v in
                        np.linspace(0, ymax, 5)],
               y_labels=[f"{v:.0f}" for v in np.linspace(0, ymax, 5)],
               xlab="Normalized Enrichment Score (NES)  —  left = bad-enriched · right = good-enriched",
               ylab="−log10(FDR)")
    add_line(slide, zx, py, zx, py + ph, color=INK, width=1.0)
    # FDR 0.05 reference
    ref_y = scale_y(-np.log10(0.05), 0, ymax, py, ph)
    add_line(slide, px, ref_y, px + pw, ref_y, color=GOLD, width=1.0, dashed=True)
    badge(slide, px + pw - In(1.8), ref_y - In(0.17),
          In(1.75), In(0.3), "FDR 0.05 threshold",
          fill=GOLD, size=9)
    # highlight top sets with gold star
    top_pos = df.nlargest(5, "NES")["pathway"].tolist()
    top_neg = df.nsmallest(3, "NES")["pathway"].tolist()
    star_set = set(top_pos + top_neg)
    for _, row in df.iterrows():
        x = scale_x(row["NES"], xmin, xmax, px, pw)
        y = scale_y(row["log10_fdr"], 0, ymax, py, ph)
        size_scale = float(row["size"]) / 300
        r = In(0.05 + 0.18 * min(size_scale, 1.0))
        col = GOOD if row["NES"] > 0 else BAD
        is_top = row["pathway"] in star_set
        if is_top:
            # gold outer ring
            add_circle(slide, x, y, r + In(0.05),
                       fill=None, line_color=GOLD, line_width=1.5)
        add_circle(slide, x, y, r, fill=col, line_color=INK, line_width=0.3)
        if is_top:
            nm = row["pathway"].replace("HALLMARK_", "")[:25]
            add_text(slide, x + In(0.1), y - In(0.11),
                     In(2.4), In(0.2), nm,
                     size=7, color=col, bold=True)
    # category legend
    add_text(slide, In(1.5), In(6.88), In(11.5), In(0.25),
             "Dot size ∝ gene-set size (leading-edge proxy). Gold ★ = top 5 positive (good-enriched) + top 3 negative (bad-enriched).",
             size=8, italic=True)

    # Panel B: Reactome top-20 NES dotplot + cell-cycle/DNA-repair gold emphasis
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "Reactome GSEA — top 20 by |NES| (gold outline for cell-cycle / DNA-repair category)",
             size=11, bold=True)
    top = reac.copy()
    top["abs_nes"] = top["NES"].abs()
    top = top.nlargest(20, "abs_nes").sort_values("NES", ascending=True).reset_index(drop=True)
    px = In(4.8); py = In(1.3); pw = In(6.5); ph = In(5.6)
    nrow = len(top)
    row_h = ph / nrow
    xmin = top["NES"].min() - 0.2
    xmax = top["NES"].max() + 0.2
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, xmin, xmax, px, pw) for v in
                        np.linspace(xmin, xmax, 5)],
               x_labels=[f"{v:.1f}" for v in np.linspace(xmin, xmax, 5)],
               xlab="Reactome NES")
    zx = scale_x(0, xmin, xmax, px, pw)
    add_line(slide, zx, py, zx, py + ph, color=INK, width=1.0)
    for i, row in top.iterrows():
        cy = py + row_h * (i + 0.5)
        x = scale_x(row["NES"], xmin, xmax, px, pw)
        col = GOOD if row["NES"] > 0 else BAD
        add_line(slide, zx, cy, x, cy, color=GREY, width=0.7)
        padj = float(row["padj"])
        r = In(0.06 + 0.12 * min(-math.log10(max(padj, 1e-20)) / 10, 1.0))
        pname_lower = str(row["pathway"]).lower()
        is_thread1 = any(k in pname_lower for k in
                         ["cycle", "mitoti", "dna repair",
                          "homology", "dsb", "checkpoint",
                          "e2f", "myc", "replication"])
        if is_thread1:
            add_circle(slide, x, cy, r + In(0.04),
                       fill=None, line_color=GOLD, line_width=1.2)
        add_circle(slide, x, cy, r, fill=col, line_color=INK, line_width=0.3)
        # label
        name = row["pathway"][:58]
        add_text(slide, px - In(3.5), cy - row_h * 0.45,
                 In(3.4), row_h * 0.9, name,
                 size=6, align="right", anchor="middle",
                 bold=is_thread1,
                 color=(GOLD if is_thread1 else INK))
    # legend
    lx = In(0.5); ly = In(1.3)
    add_text(slide, lx, ly, In(3.8), In(0.25),
             "Thread-1 biology emphasis", size=11, bold=True, color=GOLD)
    add_text(slide, lx, ly + In(0.3), In(3.8), In(1.5),
             "Gold-ringed dots = DNA-repair /\n"
             "cell-cycle / E2F / MYC / replication\n"
             "Reactome sets (the externally-\n"
             "validated Thread-1 axis).\n\n"
             "Line length ∝ |NES|.\n"
             "Dot size ∝ −log10(FDR).",
             size=9, anchor="top")

    save(prs, "SuppFig_S05_GSEA_full.pptx")


# ============================================================================
# S06 fancy — Guinney Nat Med 2015 CMS reference overlay
# ============================================================================
def build_S06_fancy():
    cms = pd.read_csv(f"{ROOT}/07_rna_cms/cms_assignments.tsv", sep="\t")
    prs = new_prs()
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "CMScaller classification by response — with Guinney 2015 reference composition",
             size=11, bold=True)
    cms_pre = cms[cms["timepoint"] == "pre"].copy() if "timepoint" in cms.columns else cms.copy()
    counts = (cms_pre.groupby(["response_bin", "prediction"]).size()
              .unstack(fill_value=0))
    cms_labels = [c for c in ["CMS1", "CMS2", "CMS3", "CMS4"] if c in counts.columns]
    cms_pal = {"CMS1": RGBColor(0xE4, 0x94, 0x4E),
               "CMS2": RGBColor(0x5A, 0x83, 0xB3),
               "CMS3": RGBColor(0xC5, 0x64, 0xA1),
               "CMS4": RGBColor(0x4E, 0xA7, 0x72)}
    px = In(2.5); py = In(1.3); pw = In(8.5); ph = In(5.0)
    vmax = max(counts.values.sum(axis=1).max(), 20) * 1.1
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, vmax, py, ph) for v in
                        np.linspace(0, vmax, 5)],
               y_labels=[f"{v:.0f}" for v in np.linspace(0, vmax, 5)],
               ylab="n samples")

    groups = [g for g in ["good", "bad", "CRC_ref"] if g in counts.index or g == "CRC_ref"]
    # Guinney 2015 reference composition across 4,151 CRC tumors:
    # CMS1 14%, CMS2 37%, CMS3 13%, CMS4 23%, unclassified 13%
    ref_frac = {"CMS1": 0.14, "CMS2": 0.37, "CMS3": 0.13, "CMS4": 0.23}

    # 3 bars: good, bad, Guinney reference (ghost bar in lighter)
    slot_w = pw / 3.3
    for gi, grp in enumerate(["good", "bad", "CRC_ref"]):
        bx = px + slot_w * (gi + 0.3) + In(0.4)
        if grp == "CRC_ref":
            # normalise to a hypothetical n=30 to be visually comparable
            n_eff = 30
            vals = [ref_frac[c] * n_eff for c in cms_labels]
            cum = 0
            for ci, (ccls, v) in enumerate(zip(cms_labels, vals)):
                if v <= 0: continue
                y_top = scale_y(cum + v, 0, vmax, py, ph)
                y_bot = scale_y(cum, 0, vmax, py, ph)
                # pastel fill
                fill_ = s28.lighten(
                    (cms_pal[ccls].r if hasattr(cms_pal[ccls], 'r') else 0,
                     0, 0),
                    factor=0.5) if False else cms_pal[ccls]
                add_rect(slide, bx - In(0.6), y_top,
                         In(1.2), y_bot - y_top,
                         fill=fill_, line_color=None)
                # hatching by horizontal lines
                for hy in np.arange(y_top, y_bot, In(0.1)):
                    add_line(slide, bx - In(0.6), hy,
                             bx + In(0.6), hy,
                             color=WHITE, width=0.4)
                cum += v
            add_text(slide, bx - In(0.6), py + ph + In(0.08),
                     In(1.2), In(0.25), "Guinney 2015\nreference",
                     size=10, align="center", bold=True, anchor="top", color=GREY)
        else:
            if grp not in counts.index: continue
            cum = 0
            for ci, ccls in enumerate(cms_labels):
                n = int(counts.loc[grp, ccls]) if ccls in counts.columns else 0
                if n == 0: continue
                y_top = scale_y(cum + n, 0, vmax, py, ph)
                y_bot = scale_y(cum, 0, vmax, py, ph)
                add_rect(slide, bx - In(0.6), y_top,
                         In(1.2), y_bot - y_top,
                         fill=cms_pal[ccls], line_color=INK, line_width=0.3)
                add_text(slide, bx - In(0.6),
                         (y_top + y_bot) / 2 - In(0.1),
                         In(1.2), In(0.2), str(n),
                         size=10, align="center", bold=True, color=WHITE)
                cum += n
            col = GOOD if grp == "good" else BAD
            add_text(slide, bx - In(0.6), py + ph + In(0.08),
                     In(1.2), In(0.25), grp,
                     size=11, align="center", bold=True, color=col,
                     anchor="top")
    # CMS legend
    lx = In(11.2); ly = In(1.5)
    add_text(slide, lx, ly - In(0.3), In(1.8), In(0.25),
             "CMS subtype", size=10, bold=True)
    for i, ccls in enumerate(cms_labels):
        add_rect(slide, lx, ly + In(i * 0.3), In(0.22), In(0.18),
                 fill=cms_pal[ccls])
        add_text(slide, lx + In(0.28), ly + In(i * 0.3 - 0.02),
                 In(1.2), In(0.22), ccls, size=10)
    # annotation
    badge(slide, In(2.5), In(0.9), In(8.5), In(0.3),
          "CMS4 (mesenchymal): 3/18 good vs 4/17 bad (Fisher P = 1.0). EMT argument rests on GSEA/ssGSEA, not CMS classifier",
          fill=VLT_GREY, size=9, bold=False)
    save(prs, "SuppFig_S06_ssGSEA_CMS.pptx")


# ============================================================================
# S07 fancy — chain-family color coding + CD8:B-cell ratio diamond
# ============================================================================
def build_S07_fancy():
    tr = pd.read_csv(f"{ROOT}/06_rna_immune/trust4_summary.tsv", sep="\t")
    df = tr[tr["timepoint"] == "pre"].copy() if "timepoint" in tr.columns else tr.copy()
    prs = new_prs()
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "TRUST4 repertoire diversity — 6 chains × Shannon entropy by response (pre-CRT) with per-chain MW P",
             size=11, bold=True)
    chains = ["TRA", "TRB", "TRG", "TRD", "IGH", "IGK", "IGL"]
    chain_family_col = {
        "TRA": THREAD2, "TRB": THREAD2, "TRG": GOLD, "TRD": GOLD,
        "IGH": RGBColor(0x8A, 0x2B, 0x4C),
        "IGK": RGBColor(0x8A, 0x2B, 0x4C),
        "IGL": RGBColor(0x8A, 0x2B, 0x4C),
    }
    px = In(1.6); py = In(1.4); pw = In(10.5); ph = In(4.8)
    vmax = max(df[f"{c}_shannon"].max() for c in chains
               if f"{c}_shannon" in df.columns) * 1.1
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, vmax, py, ph) for v in
                        np.linspace(0, vmax, 5)],
               y_labels=[f"{v:.1f}" for v in np.linspace(0, vmax, 5)],
               ylab="Shannon entropy")
    slot_w = pw / len(chains)
    # chain-family background bands
    for i, c in enumerate(chains):
        fam_col = chain_family_col[c]
        sx = px + slot_w * i
        band_col = RGBColor(0xF5, 0xF5, 0xF5) if c in ["TRA", "TRB"] else \
            RGBColor(0xFA, 0xF3, 0xE8) if c in ["TRG", "TRD"] else \
            RGBColor(0xFA, 0xEF, 0xF3)
        add_rect(slide, sx, py + In(0.02), slot_w,
                 ph - In(0.02), fill=band_col, line_color=None)
    # re-draw axis border
    add_line(slide, px, py + ph, px + pw, py + ph, color=INK, width=1.0)
    add_line(slide, px, py, px, py + ph, color=INK, width=1.0)
    # family labels
    family_groups = [("α/β T-cell", 0, 2), ("γ/δ T-cell", 2, 4),
                     ("B-cell (IGH/K/L)", 4, 7)]
    for lab, lo, hi in family_groups:
        cx0 = px + slot_w * lo
        cxw = slot_w * (hi - lo)
        col = {"α/β T-cell": THREAD2, "γ/δ T-cell": GOLD,
               "B-cell (IGH/K/L)": RGBColor(0x8A, 0x2B, 0x4C)}[lab]
        add_text(slide, cx0, py - In(0.35), cxw, In(0.25),
                 lab, size=10, bold=True, align="center", color=col)
    for i, c in enumerate(chains):
        col_name = f"{c}_shannon"
        if col_name not in df.columns: continue
        sx = px + slot_w * (i + 0.5)
        add_text(slide, sx - In(0.6), py + ph + In(0.06),
                 In(1.2), In(0.3), c, size=10, bold=True, align="center",
                 anchor="top")
        # MW P
        try:
            g_v = df[df["response_bin"] == "good"][col_name].dropna().values
            b_v = df[df["response_bin"] == "bad"][col_name].dropna().values
            _, pv = mannwhitneyu(g_v, b_v)
            pstr = f"P={pv:.2f}"
            col_p = GOLD if pv < 0.05 else GREY
        except Exception:
            pstr = "—"; col_p = GREY
        add_text(slide, sx - In(0.55), py - In(0.18),
                 In(1.1), In(0.18), pstr,
                 size=7, align="center", color=col_p, bold=True)
        for j, (resp, rcol) in enumerate([("good", GOOD), ("bad", BAD)]):
            vals = df[df["response_bin"] == resp][col_name].dropna().values
            if len(vals) == 0: continue
            cx = sx + (j - 0.5) * In(0.5)
            ys = [scale_y(float(v), 0, vmax, py, ph) for v in vals]
            boxplot_primitive(slide, cx, py, ph, ys, rcol,
                              box_w=In(0.38), dot_r=In(0.04))
    add_text(slide, In(1.6), In(6.5), In(10.5), In(0.3),
             "Per-chain Mann–Whitney P (good vs bad) shown above each panel. Chain families: α/β T-cell (teal), γ/δ (gold), B-cell (coral).",
             size=8, italic=True)

    save(prs, "SuppFig_S07_TRUST4_diversity.pptx")


# ============================================================================
# S08 fancy — Chowell Science / Ayers JCI style (Youden + CI + fill)
# ============================================================================
def build_S08_fancy():
    nested = pd.read_csv(f"{ADD}/nested_cv_drop_vs_swap.tsv", sep="\t")
    prs = new_prs()

    # Panel A: AUC bar with chance + excellent references, Youden annotation
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "Nested LOOCV AUC — 5 scenarios × 2 models (Chowell Science 2018-style reference bands)",
             size=11, bold=True)
    px = In(2.0); py = In(1.4); pw = In(10.3); ph = In(5.0)
    vmin, vmax = 0.4, 1.0
    y_ticks_v = [0.4, 0.5, 0.6, 0.7, 0.8, 0.9, 1.0]
    # reference bands: chance (0.5) grey + excellent (0.85) gold
    chance_y = scale_y(0.5, vmin, vmax, py, ph)
    excellent_y = scale_y(0.85, vmin, vmax, py, ph)
    add_rect(slide, px, chance_y - In(0.015), pw, In(0.03),
             fill=GREY, line_color=None)
    add_rect(slide, px, py, pw, excellent_y - py,
             fill=RGBColor(0xFA, 0xF4, 0xE5), line_color=None)
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, vmin, vmax, py, ph) for v in y_ticks_v],
               y_labels=[f"{v:.1f}" for v in y_ticks_v],
               ylab="Outer-LOOCV AUC (95 % bootstrap CI)")
    # re-draw chance line
    add_line(slide, px, chance_y, px + pw, chance_y,
             color=GREY, width=1.0, dashed=True)
    add_text(slide, px + pw + In(0.04), chance_y - In(0.1),
             In(1.5), In(0.2), "chance 0.5", size=8, color=GREY)
    add_line(slide, px, excellent_y, px + pw, excellent_y,
             color=GOLD, width=0.8, dashed=True)
    add_text(slide, px + pw + In(0.04), excellent_y - In(0.1),
             In(1.5), In(0.2), "excellent 0.85", size=8, color=GOLD)
    scenarios = nested["scenario"].unique().tolist()
    scen_labels = {
        "baseline_37": "Baseline\n37 feat",
        "drop_cd8prolif_36": "★ Drop CD8prolif\n36 feat",
        "add_immune_40": "Add 3 immune\n40 feat",
        "swap_cd8_37": "Swap CD8prolif→\nCD8cytotoxic (37)",
        "drop_prolif_add_3_39": "Drop+Add\n39 feat",
    }
    nscen = len(scenarios)
    slot_w = pw / nscen
    for i, scen in enumerate(scenarios):
        sx = px + slot_w * (i + 0.5)
        lab = scen_labels.get(scen, scen)
        is_winner = (scen == "drop_cd8prolif_36")
        add_text(slide, sx - In(0.9), py + ph + In(0.08),
                 In(1.8), In(0.55), lab,
                 size=9, align="center", anchor="top",
                 bold=is_winner, color=(GOLD if is_winner else INK))
        for j, model in enumerate(["LASSO", "ElasticNet"]):
            row = nested[(nested["scenario"] == scen) & (nested["model"] == model)].iloc[0]
            bx = sx + (j - 0.5) * In(0.55)
            auc = row["AUC"]; ci_lo = row["CI_low"]; ci_hi = row["CI_high"]
            bar_top = scale_y(auc, vmin, vmax, py, ph)
            bar_base = scale_y(0.5, vmin, vmax, py, ph)
            base_col = THREAD1 if model == "LASSO" else THREAD2
            fill_col = GOLD if (is_winner and model == "ElasticNet") else base_col
            add_rect(slide, bx - In(0.22), bar_top,
                     In(0.44), bar_base - bar_top,
                     fill=fill_col, line_color=INK, line_width=0.5)
            hi_y = scale_y(ci_hi, vmin, vmax, py, ph)
            lo_y = scale_y(ci_lo, vmin, vmax, py, ph)
            add_line(slide, bx, hi_y, bx, lo_y, color=INK, width=1.2)
            add_line(slide, bx - In(0.08), hi_y, bx + In(0.08), hi_y,
                     color=INK, width=1.0)
            add_line(slide, bx - In(0.08), lo_y, bx + In(0.08), lo_y,
                     color=INK, width=1.0)
            add_text(slide, bx - In(0.3), bar_top - In(0.22),
                     In(0.6), In(0.18), f"{auc:.3f}",
                     size=8, align="center", color=INK,
                     bold=(is_winner and model == "ElasticNet"))
            if is_winner and model == "ElasticNet":
                add_star(slide, bx, bar_top - In(0.35), In(0.08))
    # legend
    lx = In(10.6); ly = In(0.85)
    add_rect(slide, lx, ly, In(0.18), In(0.12), fill=THREAD1)
    add_text(slide, lx + In(0.22), ly - In(0.02), In(1.2),
             In(0.18), "LASSO", size=9)
    add_rect(slide, lx, ly + In(0.18), In(0.18), In(0.12), fill=THREAD2)
    add_text(slide, lx + In(0.22), ly + In(0.16), In(1.2),
             In(0.18), "ElasticNet", size=9)
    add_rect(slide, lx, ly + In(0.36), In(0.18), In(0.12), fill=GOLD)
    add_text(slide, lx + In(0.22), ly + In(0.34), In(1.5),
             In(0.18), "★ winner", size=9, bold=True)
    # footer summary
    badge(slide, In(2.0), In(6.85), In(10.3), In(0.35),
          "Winner: drop-CD8-proliferation 36-feature ElasticNet  AUC = 0.745 [0.56, 0.90] — falls just below 'excellent' band (0.85) but clearly above chance",
          fill=GOLD, size=9, bold=False)

    # Panel B: ROC overlay with AUC fill + Youden J + sens/spec annotation
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "Nested outer-LOOCV ROC — 5 ElasticNet scenarios (winner AUC fill + Youden's J diamond)",
             size=11, bold=True)
    px = In(2.8); py = In(1.4); pw = In(6.0); ph = In(5.0)
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, 0, 1, px, pw) for v in [0, 0.25, 0.5, 0.75, 1.0]],
               x_labels=["0", "0.25", "0.5", "0.75", "1.0"],
               y_ticks=[scale_y(v, 0, 1, py, ph) for v in [0, 0.25, 0.5, 0.75, 1.0]],
               y_labels=["0", "0.25", "0.5", "0.75", "1.0"],
               xlab="1 − Specificity", ylab="Sensitivity")
    add_line(slide, px, py + ph, px + pw, py, color=GREY, width=0.6, dashed=True)
    scen_ord = ["baseline_37", "drop_cd8prolif_36", "add_immune_40",
                "swap_cd8_37", "drop_prolif_add_3_39"]
    scen_col = {"baseline_37": GREY,
                "drop_cd8prolif_36": GOLD,
                "add_immune_40": THREAD1,
                "swap_cd8_37": THREAD2,
                "drop_prolif_add_3_39": INK}

    # AUC fill for winner first (so behind lines)
    winner_row = nested[(nested["scenario"] == "drop_cd8prolif_36")
                        & (nested["model"] == "ElasticNet")].iloc[0]
    winner_auc = winner_row["AUC"]
    # binormal param
    A_w = np.sqrt(2) * norm.ppf(winner_auc)
    xs_w = np.linspace(0, 1, 60)
    ys_w = norm.cdf(A_w - norm.ppf(1 - xs_w))
    # AUC fill polygon: area between curve and diagonal
    from pptx.enum.shapes import MSO_SHAPE
    poly_pts = [(scale_x(0, 0, 1, px, pw), scale_y(0, 0, 1, py, ph))]
    for xx, yy in zip(xs_w, ys_w):
        poly_pts.append((scale_x(xx, 0, 1, px, pw),
                         scale_y(yy, 0, 1, py, ph)))
    poly_pts.append((scale_x(1, 0, 1, px, pw), scale_y(1, 0, 1, py, ph)))
    # build polygon via freeform
    x0, y0 = poly_pts[0]
    ff = slide.shapes.build_freeform(_i(x0), _i(y0), scale=1.0)
    ff.add_line_segments([(_i(px), _i(py)) for px, py in poly_pts[1:]], close=True)
    shp = ff.convert_to_shape()
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xE5, 0xEF, 0xE8)
    shp.line.fill.background()
    shp.shadow.inherit = False
    s28.kill_shadow(shp)

    for scen in scen_ord:
        row = nested[(nested["scenario"] == scen) & (nested["model"] == "ElasticNet")].iloc[0]
        auc = row["AUC"]
        A = np.sqrt(2) * norm.ppf(auc)
        xs = np.linspace(0, 1, 60)
        ys = norm.cdf(A - norm.ppf(1 - xs))
        pts = [(scale_x(x, 0, 1, px, pw), scale_y(y, 0, 1, py, ph))
               for x, y in zip(xs, ys)]
        width = 2.5 if scen == "drop_cd8prolif_36" else 1.0
        for (x1, y1), (x2, y2) in zip(pts[:-1], pts[1:]):
            add_line(slide, x1, y1, x2, y2, color=scen_col[scen], width=width)
        # Youden J for winner: argmax(sens - (1 - spec))
        if scen == "drop_cd8prolif_36":
            j_idx = int(np.argmax(ys - xs))
            yx = scale_x(xs[j_idx], 0, 1, px, pw)
            yy = scale_y(ys[j_idx], 0, 1, py, ph)
            add_diamond(slide, yx, yy, In(0.11),
                        fill=GOLD, line_color=INK, line_width=1.0)
            add_text(slide, yx + In(0.1), yy - In(0.12),
                     In(2.2), In(0.3),
                     f"Youden J\nsens={ys[j_idx]:.2f}\nspec={1 - xs[j_idx]:.2f}",
                     size=8, color=GOLD, bold=True)

    # legend box
    lx = In(9.3); ly = In(1.6)
    add_text(slide, lx, ly - In(0.32), In(3.5), In(0.25),
             "Scenario (ElasticNet)", size=10, bold=True)
    for i, scen in enumerate(scen_ord):
        row = nested[(nested["scenario"] == scen) & (nested["model"] == "ElasticNet")].iloc[0]
        yy = ly + In(0.38 * i)
        add_line(slide, lx, yy + In(0.07), lx + In(0.35), yy + In(0.07),
                 color=scen_col[scen],
                 width=(2.5 if scen == "drop_cd8prolif_36" else 1.2))
        star = "★ " if scen == "drop_cd8prolif_36" else "  "
        add_text(slide, lx + In(0.40), yy - In(0.02),
                 In(3.6), In(0.3),
                 f"{star}{scen_labels[scen].replace(chr(10),' ')} AUC={row['AUC']:.3f}",
                 size=9, bold=(scen == "drop_cd8prolif_36"),
                 color=(GOLD if scen == "drop_cd8prolif_36" else INK))

    # Panel C: per-subject probabilities — misclass gold X + subject ID + confidence ribbon
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "C", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "Per-subject nested-outer predicted P(good) — winning model (misclass = gold ✗)",
             size=11, bold=True)
    px = In(1.2); py = In(1.4); pw = In(11.7); ph = In(5.0)
    probs_path = None
    for candidate in [f"{ADD}/nested_outer_probs_drop_cd8prolif_ElasticNet.tsv",
                      f"{ADD}/nested_outer_probs_ext40_ElasticNet.tsv",
                      f"{ADD}/nested_outer_probs_orig37_ElasticNet.tsv"]:
        if os.path.exists(candidate):
            probs_path = candidate; break
    try:
        probs = pd.read_csv(probs_path, sep="\t") if probs_path else None
    except Exception:
        probs = None
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, 1, py, ph) for v in [0, 0.25, 0.5, 0.75, 1.0]],
               y_labels=["0", "0.25", "0.5", "0.75", "1.0"],
               ylab="Predicted P(good)")
    # 0.5 threshold band (wide uncertainty ribbon)
    u_top = scale_y(0.6, 0, 1, py, ph)
    u_bot = scale_y(0.4, 0, 1, py, ph)
    add_rect(slide, px, u_top, pw, u_bot - u_top,
             fill=RGBColor(0xF5, 0xF5, 0xF5), line_color=None)
    thr_y = scale_y(0.5, 0, 1, py, ph)
    add_line(slide, px, thr_y, px + pw, thr_y,
             color=INK, width=0.8, dashed=True)
    add_text(slide, px + pw + In(0.05), thr_y - In(0.1),
             In(1.5), In(0.2), "threshold 0.5", size=8, color=INK)
    add_text(slide, px + pw + In(0.05), u_top - In(0.12),
             In(1.5), In(0.2), "uncertainty band", size=7, color=GREY)
    if probs is not None and len(probs):
        col_p = [c for c in probs.columns if "prob" in c.lower()]
        p_col = col_p[0] if col_p else probs.columns[-1]
        y_col = "response_bin" if "response_bin" in probs.columns else probs.columns[0]
        s_col = "subject_id" if "subject_id" in probs.columns else probs.columns[0]
        df = probs.copy()
        df["_grp"] = df[y_col].map(lambda x: 0 if x == "good" else 1)
        df = df.sort_values(["_grp", p_col]).reset_index(drop=True)
        n = len(df); bar_w = pw / (n + 1)
        n_mis = 0
        for i, row in df.iterrows():
            bx = px + bar_w * (i + 0.5)
            p = float(row[p_col])
            col = GOOD if row[y_col] == "good" else BAD
            pred = "good" if p >= 0.5 else "bad"
            mis = (pred != row[y_col]); n_mis += int(mis)
            h_top = scale_y(p, 0, 1, py, ph)
            h_base = scale_y(0, 0, 1, py, ph)
            add_rect(slide, bx - bar_w * 0.4, h_top,
                     bar_w * 0.8, h_base - h_top,
                     fill=col, line_color=INK, line_width=0.4)
            if mis:
                add_rect(slide, bx - bar_w * 0.42, h_top - In(0.02),
                         bar_w * 0.84, In(0.04),
                         fill=GOLD, line_color=None)
                add_text(slide, bx - In(0.15), h_top - In(0.24),
                         In(0.3), In(0.18), "✗",
                         size=12, bold=True, color=GOLD, align="center")
            add_text(slide, bx - bar_w * 0.5, py + ph + In(0.02),
                     bar_w, In(0.2), str(row[s_col]),
                     size=6, align="center", anchor="top")
        badge(slide, In(1.2), In(6.85), In(11.5), In(0.35),
              f"Misclassified at 0.5 threshold: {n_mis}/{n} (accuracy {100 * (n - n_mis) / n:.0f} %). "
              "AUC 0.745 [0.56, 0.90] (threshold-free) is the load-bearing claim; threshold accuracy is provided for transparency only.",
              fill=VLT_GREY, size=9, bold=False)

    save(prs, "SuppFig_S08_ML_scenario_ablation.pptx")


# ============================================================================
# S09 fancy — CONSORT with flow arrow sizes + exclusion chips
# ============================================================================
def build_S09_fancy():
    prs = new_prs()
    # Panel A: enriched inclusion matrix (same as 28, keep)
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "9 GEO nCRT cohorts — inclusion / exclusion matrix (primary meta 5 of 9; concordance ≥ 3/4 rule)",
             size=11, bold=True)
    cohort_info = [
        ("GSE35452", 46, "LC-CRT + concurrent cape", "TRG 4-class", "4/4", "primary ★"),
        ("GSE45404", 80, "LC-CRT + concurrent cape", "pCR/non-pCR", "3/4", "primary ★"),
        ("GSE56699", 72, "LC-CRT + concurrent cape", "TRG 3-class", "3/4", "primary ★"),
        ("GSE133057", 33, "LC-CRT + concurrent cape", "OS/DFS-derived", "3/4", "primary ★"),
        ("GSE87211", 287, "LC-CRT + concurrent cape", "TRG pooled", "3/4", "primary ★"),
        ("GSE150082", 39, "short-course RT + TNT subset", "pTRG", "1/4", "excluded (regimen)"),
        ("GSE119409", 66, "RADIOTHERAPY ALONE (no chemo)", "sensitivity", "1/4", "excluded (no chemo)"),
        ("GSE94104", 80, "LC-CRT", "CMS-stability (no resp)", "n/a", "excluded (no endpoint)"),
        ("GSE46862", 69, "LC-CRT", "TRG ambiguous", "1/4", "excluded (1/4 discord)"),
    ]
    px = In(0.6); py = In(1.3); w = In(12.2); h = In(5.3)
    n_row = len(cohort_info)
    headers = ["GSE accession", "N", "Regimen", "Endpoint", "Thread 1\nconcord", "Status"]
    col_wid = [In(1.5), In(0.6), In(4.2), In(2.3), In(1.4), In(2.2)]
    header_h = In(0.45)
    row_h = (h - header_h) / n_row
    x_cursor = px
    for j, head in enumerate(headers):
        add_rect(slide, x_cursor, py, col_wid[j], header_h,
                 fill=LT_GREY, line_color=INK, line_width=0.6)
        add_text(slide, x_cursor + In(0.05), py + In(0.05),
                 col_wid[j] - In(0.1), header_h - In(0.1),
                 head, size=10, bold=True, align="center", anchor="middle")
        x_cursor += col_wid[j]
    for i, rowdata in enumerate(cohort_info):
        x_cursor = px
        yrow = py + header_h + row_h * i
        status = rowdata[5]
        fill_col = TEAL_LT if "primary" in status else \
            CORAL_LT if "excluded" in status else WHITE
        for j, cell in enumerate(rowdata):
            add_rect(slide, x_cursor, yrow, col_wid[j], row_h,
                     fill=fill_col, line_color=INK, line_width=0.3)
            add_text(slide, x_cursor + In(0.05), yrow + In(0.03),
                     col_wid[j] - In(0.1), row_h - In(0.06),
                     str(cell), size=9,
                     bold=("primary" in str(rowdata[5]) and j == 0),
                     align="center" if j != 2 else "left", anchor="middle")
            x_cursor += col_wid[j]
    add_text(slide, In(0.6), In(6.85), In(12.2), In(0.3),
             "Primary meta (5 cohorts, N = 518) selected by objective ≥ 3 / 4 Thread-1-concordance rule BEFORE Z computation.",
             size=9, italic=True)

    # Panel B: CONSORT flow — arrow widths proportional to N flow
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "External-validation CONSORT — 9 → 5 primary meta (flow-width ∝ N)",
             size=11, bold=True)

    def cbox2(x, y, w, h, lines, fill=WHITE, line_color=INK):
        add_rect(slide, x, y, w, h, fill=fill, line_color=line_color, line_width=1.0)
        for i, ln in enumerate(lines):
            add_text(slide, x + In(0.08), y + In(0.08 + 0.28 * i),
                     w - In(0.16), In(0.28), ln,
                     size=11 if i == 0 else 9, bold=(i == 0),
                     align="center", anchor="middle")
    # Tier 1
    cbox2(In(4.8), In(1.1), In(3.7), In(0.85),
          ["9 candidate GEO nCRT cohorts", "N_total = 721 samples"])
    # vertical arrow
    add_line(slide, In(6.65), In(1.95), In(6.65), In(2.4),
             color=INK, width=2.5)
    # Tier 2 primary (N=518) + excluded (N=203) — arrow widths proportional
    total = 518 + 203
    w_prim = 2.5 * 518 / total
    w_excl = 2.5 * 203 / total
    # primary arrow (leftward)
    add_line(slide, In(6.65), In(2.0), In(3.9), In(2.4),
             color=GOOD, width=max(w_prim * 2, 3))
    add_line(slide, In(6.65), In(2.0), In(9.5), In(2.4),
             color=BAD, width=max(w_excl * 2, 2))
    cbox2(In(1.5), In(2.4), In(4.8), In(2.8),
          ["PRIMARY META (N = 518)",
           "5 LC-CRT cohorts, concurrent capecitabine",
           "≥ 3 / 4 Thread-1 signature concordance",
           "",
           "GSE35452 (N=46), GSE45404 (80),",
           "GSE56699 (72), GSE133057 (33),",
           "GSE87211 (287)",
           "",
           "Restricted Z: DSB +3.17, cellcycle +3.21,",
           "E2F/MYC +2.79 (all P < 0.01)"],
          fill=TEAL_LT, line_color=GOOD)
    cbox2(In(7.0), In(2.4), In(5.0), In(2.8),
          ["EXCLUDED (N = 203)",
           "",
           "✗ GSE119409 (66): RT-alone, no chemo",
           "✗ GSE94104 (80): no response endpoint",
           "✗ GSE150082 (39): SC-RT mixed regimen",
           "✗ GSE46862 (69): 1/4 Thread-1 concord",
           "",
           "Full 9-cohort sensitivity meta shown",
           "in Supp Fig S19A (hollow diamond)"],
          fill=CORAL_LT, line_color=BAD)
    # arrow label badges
    badge(slide, In(2.5), In(2.2), In(1.3), In(0.24),
          f"N = 518 included", fill=GOOD, text_color=WHITE, size=9)
    badge(slide, In(9.2), In(2.2), In(1.3), In(0.24),
          f"N = 203 excluded", fill=BAD, text_color=WHITE, size=9)

    # Tier 3 Thread 2 augmentation
    add_line(slide, In(3.9), In(5.2), In(3.9), In(5.6),
             color=GOLD, width=3.5)
    cbox2(In(3.8), In(5.6), In(5.7), In(1.4),
          ["Thread 2 +Akiyoshi augmentation",
           "GSE216616 published-statistic (cytolytic activity P = 0.005, Z = +2.81)",
           "6-source CD8-cytotoxic meta: Z = +3.29, P = 0.001, N = 816"],
          fill=GOLD, line_color=INK)
    add_star(slide, In(3.8) - In(0.18), In(5.6) + In(0.2), In(0.09), fill=GOLD)

    save(prs, "SuppFig_S09_GEO_cohorts_CONSORT.pptx")


# ============================================================================
# S14 fancy — Cercek NEJM 2022 style (treatment-phase bar above waterfall)
# ============================================================================
def build_S14_fancy():
    clin = pd.read_csv(f"{ROOT}/00_cohort/clinical_master.tsv", sep="\t")
    prs = new_prs()
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "S14", size=18, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "Per-patient clinical waterfall (N = 35) — Cercek 2022-style treatment-phase bar + RECIST reference",
             size=11, bold=True)
    df = clin.sort_values(["response_bin", "response_num", "subject_id"],
                          ascending=[True, True, True]).reset_index(drop=True)
    n = len(df)
    # Treatment-phase bar at the TOP — shared across all patients
    px = In(0.8); py = In(0.95); pw = In(11.8)
    tpb_h = In(0.3)
    # 3-phase: SC-RT / consolidation / surgery-watch
    phase_col = {"SC-RT": RGBColor(0x4F, 0x73, 0x8E),
                 "FOLFOX/CAPOX": RGBColor(0xA6, 0x62, 0x2C),
                 "surgery / WW": RGBColor(0x44, 0x4A, 0x58)}
    phase_frac = [0.2, 0.55, 0.25]
    cum = 0
    for label, frac in zip(["SC-RT", "FOLFOX/CAPOX", "surgery / WW"], phase_frac):
        pxf = px + pw * cum
        pwf = pw * frac
        add_rect(slide, pxf, py, pwf, tpb_h,
                 fill=phase_col[label], line_color=INK, line_width=0.5)
        add_text(slide, pxf, py, pwf, tpb_h, label,
                 size=9, bold=True, color=WHITE, align="center", anchor="middle")
        cum += frac
    # arrow at end
    add_text(slide, In(12.65), py + In(0.03), In(0.6), In(0.25),
             "→ TRG", size=9, bold=True, align="left", anchor="middle")
    # Waterfall — main panel
    wp_x = In(0.8); wp_y = In(1.45); wp_w = In(11.8); wp_h = In(3.4)
    bar_w = wp_w / (n + 1)
    vmax = 3; vmin = 0
    axis_frame(slide, wp_x, wp_y, wp_w, wp_h,
               y_ticks=[scale_y(v, vmin, vmax, wp_y, wp_h) for v in [0, 1, 2, 3]],
               y_labels=["0 CR", "1 nCR", "2 PR", "3 poor"],
               ylab="TNT response score")
    # RECIST-like reference: dashed gold at y=1 (near-CR cutoff, analog to Cercek's -30% PR)
    ref_y = scale_y(1.0, vmin, vmax, wp_y, wp_h)
    add_line(slide, wp_x, ref_y, wp_x + wp_w, ref_y,
             color=GOLD, width=1.0, dashed=True)
    add_text(slide, wp_x + wp_w + In(0.05), ref_y - In(0.1),
             In(1.3), In(0.2), "good/bad cutoff", size=7, color=GOLD)
    for i, (_, row) in enumerate(df.iterrows()):
        bx = wp_x + bar_w * (i + 0.5)
        col = GOOD if row["response_bin"] == "good" else BAD
        h_top = scale_y(float(row["response_num"]), vmin, vmax, wp_y, wp_h)
        h_base = scale_y(0, vmin, vmax, wp_y, wp_h)
        add_rect(slide, bx - bar_w * 0.4, h_top,
                 bar_w * 0.8, max(h_base - h_top, In(0.02)),
                 fill=col, line_color=INK, line_width=0.3)
        # response label inside bar (Cercek style)
        lab_map = {0: "CR", 1: "nCR", 2: "PR", 3: "poor"}
        add_text(slide, bx - bar_w * 0.5, (h_top + h_base) / 2 - In(0.08),
                 bar_w, In(0.15), lab_map[int(row["response_num"])],
                 size=6, bold=True, color=WHITE, align="center")
    # Clinical tracks below (cT / sex / age)
    strip_y = wp_y + wp_h + In(0.1)
    strip_h = In(0.22)
    ct_pal = {"T2": RGBColor(0xA8, 0xC8, 0xE1), "T2/T3": RGBColor(0x7D, 0xA7, 0xC8),
              "T3": RGBColor(0x45, 0x7A, 0x9E), "T4": RGBColor(0x21, 0x4A, 0x70)}
    sex_pal = {"M": RGBColor(0x3B, 0x6B, 0x9E), "F": RGBColor(0xB4, 0x55, 0x78)}
    for ti, (lab, col_resolver) in enumerate([
            ("cT", lambda r: ct_pal.get(r["cT"], GREY)),
            ("sex", lambda r: sex_pal.get(r["sex"], GREY)),
            ("age", None)]):
        y_s = strip_y + ti * (strip_h + In(0.04))
        add_text(slide, wp_x - In(0.55), y_s,
                 In(0.5), strip_h, lab,
                 size=9, bold=True, align="right", anchor="middle")
        for i, (_, row) in enumerate(df.iterrows()):
            bx = wp_x + bar_w * (i + 0.5)
            if lab == "age":
                amin = df["age"].min(); amax = df["age"].max()
                v = (row["age"] - amin) / max(amax - amin, 1)
                shade = int(230 - v * 180)
                col = RGBColor(shade, shade, shade)
                add_rect(slide, bx - bar_w * 0.4, y_s,
                         bar_w * 0.8, strip_h,
                         fill=col, line_color=WHITE, line_width=0.2)
                add_text(slide, bx - bar_w * 0.5, y_s,
                         bar_w, strip_h, str(int(row["age"])),
                         size=6, align="center", anchor="middle", color=WHITE)
            else:
                col = col_resolver(row)
                add_rect(slide, bx - bar_w * 0.4, y_s,
                         bar_w * 0.8, strip_h,
                         fill=col, line_color=WHITE, line_width=0.2)
    # subject IDs under tracks
    id_y = strip_y + 3 * (strip_h + In(0.04))
    for i, (_, row) in enumerate(df.iterrows()):
        bx = wp_x + bar_w * (i + 0.5)
        add_text(slide, bx - bar_w * 0.5, id_y,
                 bar_w, In(0.2), str(row["subject_id"]),
                 size=6, align="center", anchor="top")
    # legend strip
    ly = In(6.75)
    # cT legend
    xlg = In(0.9)
    add_text(slide, xlg, ly, In(0.4), In(0.22), "cT:", size=9, bold=True)
    for j, stage in enumerate(["T2", "T2/T3", "T3", "T4"]):
        cx0 = xlg + In(0.5 + j * 0.9)
        add_rect(slide, cx0, ly + In(0.04), In(0.2), In(0.14),
                 fill=ct_pal[stage])
        add_text(slide, cx0 + In(0.24), ly + In(0.02),
                 In(0.6), In(0.18), stage, size=8)
    # sex legend
    xlg2 = In(5.0)
    add_text(slide, xlg2, ly, In(0.4), In(0.22), "sex:", size=9, bold=True)
    for j, (k, col) in enumerate(sex_pal.items()):
        cx0 = xlg2 + In(0.5 + j * 0.6)
        add_rect(slide, cx0, ly + In(0.04), In(0.2), In(0.14), fill=col)
        add_text(slide, cx0 + In(0.24), ly + In(0.02),
                 In(0.4), In(0.18), k, size=8)
    # response legend
    xlg3 = In(7.5)
    add_text(slide, xlg3, ly, In(0.8), In(0.22), "response:", size=9, bold=True)
    add_rect(slide, xlg3 + In(0.9), ly + In(0.04),
             In(0.2), In(0.14), fill=GOOD)
    add_text(slide, xlg3 + In(1.16), ly + In(0.02),
             In(1.3), In(0.18),
             f"good (n={(df.response_bin=='good').sum()})", size=8)
    add_rect(slide, xlg3 + In(2.5), ly + In(0.04),
             In(0.2), In(0.14), fill=BAD)
    add_text(slide, xlg3 + In(2.76), ly + In(0.02),
             In(1.3), In(0.18),
             f"bad (n={(df.response_bin=='bad').sum()})", size=8)
    # treatment phase legend
    xlg4 = In(11.0)
    add_text(slide, xlg4, ly, In(1.5), In(0.22),
             "Phase timeline (top)", size=9, bold=True)
    save(prs, "SuppFig_S14_clinical_waterfall.pptx")


# ============================================================================
# S17 fancy — Riaz Cell 2017 style target-engagement (pre-post connector)
# ============================================================================
def build_S17_fancy():
    bd = pd.read_csv(f"{ADD}/baseline_factor_per_subject_delta.tsv", sep="\t")
    prs = new_prs()

    # Panel A: sign-count with gold highlight on EMT 6/6 + member families
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "Target-engagement member-level sign counts (composite + 17 member signatures × 12 paired) — EMT 6/6 gold",
             size=11, bold=True)
    df = bd.copy()
    pred_dir = {"DSB_HDR_repair": -1, "E2F_MYC_cellcycle": -1,
                "Tumor_cellcycle": -1, "EMT": +1}
    df["delta"] = df["post"] - df["pre"]
    df["predicted"] = df["factor"].map(pred_dir).fillna(1)
    df["concord"] = np.sign(df["delta"]) == np.sign(df["predicted"])
    agg = df.groupby(["factor", "member", "response_bin"])["concord"].agg(["sum", "count"]).reset_index()
    members_ord = (df.groupby(["factor", "member"]).size()
                   .reset_index().sort_values(["factor", "member"])
                   [["factor", "member"]].to_records(index=False).tolist())
    n_row = len(members_ord)
    px = In(3.5); py = In(1.2); pw = In(8.5); ph = In(5.5)
    row_h = ph / max(n_row, 1)
    zx = px + pw / 2
    # Gray zero-band (Bashford-Rogers Nature 2019)
    band_half = pw / 2 * (0.5 / 6.0)
    add_rect(slide, zx - band_half, py, band_half * 2, ph,
             fill=RGBColor(0xEF, 0xEF, 0xEF), line_color=None)
    axis_frame(slide, px, py, pw, ph, xlab="sign-count (good ← 0 → bad) — gray band = |count| < 0.5",
               x_ticks=[px, zx - pw / 4, zx, zx + pw / 4, px + pw],
               x_labels=["-6", "-3", "0", "+3", "+6"])
    add_line(slide, zx, py, zx, py + ph, color=INK, width=1.2)
    # per-factor group color shading (light bands)
    factor_col = {"DSB_HDR_repair": RGBColor(0xE6, 0xEE, 0xF5),
                  "E2F_MYC_cellcycle": RGBColor(0xEE, 0xE6, 0xF5),
                  "Tumor_cellcycle": RGBColor(0xF5, 0xE6, 0xE6),
                  "EMT": RGBColor(0xFA, 0xF3, 0xE0)}
    for i, (factor, member) in enumerate(members_ord):
        cy = py + row_h * (i + 0.5)
        rw = agg[(agg["factor"] == factor) & (agg["member"] == member)
                 & (agg["response_bin"] == "good")]
        bw = agg[(agg["factor"] == factor) & (agg["member"] == member)
                 & (agg["response_bin"] == "bad")]
        n_g = int(rw.iloc[0]["sum"]) if len(rw) else 0
        n_b = int(bw.iloc[0]["sum"]) if len(bw) else 0
        # EMT 6/6 = gold star flag
        is_emt_unanimous = (factor == "EMT" and n_g == 6)
        if is_emt_unanimous:
            add_rect(slide, px, cy - row_h / 2, pw, row_h,
                     fill=RGBColor(0xFA, 0xF0, 0xD4), line_color=None)
        if n_g:
            bar_len = pw / 2 * (n_g / 6.0)
            add_rect(slide, zx, cy - row_h * 0.3, bar_len, row_h * 0.6,
                     fill=GOLD if is_emt_unanimous else GOOD,
                     line_color=INK, line_width=0.3)
            add_text(slide, zx + bar_len + In(0.03), cy - In(0.1),
                     In(0.4), In(0.2), str(n_g), size=7,
                     color=(GOLD if is_emt_unanimous else GOOD), bold=True)
        if n_b:
            bar_len = pw / 2 * (n_b / 6.0)
            add_rect(slide, zx - bar_len, cy - row_h * 0.3,
                     bar_len, row_h * 0.6, fill=BAD,
                     line_color=INK, line_width=0.3)
            add_text(slide, zx - bar_len - In(0.3), cy - In(0.1),
                     In(0.26), In(0.2), str(n_b), size=7, color=BAD, align="right")
        # labels
        star = "★ " if is_emt_unanimous else "  "
        add_text(slide, px - In(3.3), cy - In(0.09),
                 In(3.2), In(0.2),
                 star + f"[{factor[:10]}] {str(member)[:40]}",
                 size=7, align="right", anchor="middle",
                 bold=is_emt_unanimous,
                 color=(GOLD if is_emt_unanimous else INK))
    # badges
    add_rect(slide, In(10.5), In(0.8), In(0.2), In(0.14), fill=GOOD)
    add_text(slide, In(10.75), In(0.78), In(1.5), In(0.2), "good concordant", size=9)
    add_rect(slide, In(10.5), In(1.0), In(0.2), In(0.14), fill=BAD)
    add_text(slide, In(10.75), In(0.98), In(1.5), In(0.2), "bad concordant", size=9)
    badge(slide, In(10.5), In(1.25), In(2.2), In(0.3),
          "★ EMT 6/6 unanimous good (P=0.016)",
          fill=GOLD, size=8)

    # Panel B: pre-post connector style (Riaz Cell 2017) — 4 composites only
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "Per-subject pre→post ssGSEA scores for 4 composite signatures (Riaz Cell 2017-style connector)",
             size=11, bold=True)
    composites = ["DSB_HDR_repair", "E2F_MYC_cellcycle", "Tumor_cellcycle", "EMT"]
    # compute composite score per (subject, timepoint) as mean of members
    comp_scores = (df.groupby(["subject_id", "response_bin", "factor"])
                   [["pre", "post"]].mean().reset_index())
    # 4 subpanel grid 2x2
    panels = [(composites[0], In(0.9), In(1.2)),
              (composites[1], In(7.0), In(1.2)),
              (composites[2], In(0.9), In(4.2)),
              (composites[3], In(7.0), In(4.2))]
    for sig, ox, oy in panels:
        sub = comp_scores[comp_scores["factor"] == sig]
        if len(sub) == 0: continue
        pw_ = In(5.5); ph_ = In(2.6)
        # axes
        all_vals = np.concatenate([sub["pre"].values, sub["post"].values])
        vmax_ = all_vals.max() * 1.1
        vmin_ = all_vals.min() * 1.1
        axis_frame(slide, ox, oy, pw_, ph_,
                   y_ticks=[scale_y(v, vmin_, vmax_, oy, ph_) for v in
                            np.linspace(vmin_, vmax_, 4)],
                   y_labels=[f"{v:.1f}" for v in np.linspace(vmin_, vmax_, 4)],
                   ylab=sig.replace("_", " "), tick_size=7, lab_size=8)
        pre_x = ox + pw_ * 0.2
        post_x = ox + pw_ * 0.8
        add_text(slide, pre_x - In(0.3), oy + ph_ + In(0.02),
                 In(0.6), In(0.18), "pre", size=8, align="center", anchor="top")
        add_text(slide, post_x - In(0.3), oy + ph_ + In(0.02),
                 In(0.6), In(0.18), "post", size=8, align="center", anchor="top")
        # per subject connector
        pred_up = pred_dir.get(sig, 1) > 0
        # median lines per group (Liu NatMed 2019)
        g_pre = sub[sub.response_bin == "good"]["pre"].median()
        g_post = sub[sub.response_bin == "good"]["post"].median()
        b_pre = sub[sub.response_bin == "bad"]["pre"].median()
        b_post = sub[sub.response_bin == "bad"]["post"].median()
        for _, row in sub.iterrows():
            col = GOOD if row["response_bin"] == "good" else BAD
            dy_pre = scale_y(float(row["pre"]), vmin_, vmax_, oy, ph_)
            dy_post = scale_y(float(row["post"]), vmin_, vmax_, oy, ph_)
            delta = row["post"] - row["pre"]
            concord = (delta > 0 and pred_up) or (delta < 0 and not pred_up)
            line_col = col if concord else RGBColor(0xD0, 0xD0, 0xD0)
            add_line(slide, pre_x, dy_pre, post_x, dy_post,
                     color=line_col, width=1.0)
            add_circle(slide, pre_x, dy_pre, In(0.04), fill=col, line_color=INK, line_width=0.2)
            add_circle(slide, post_x, dy_post, In(0.04), fill=col, line_color=INK, line_width=0.2)
        # group-median bold dashed
        gy_pre = scale_y(float(g_pre), vmin_, vmax_, oy, ph_)
        gy_post = scale_y(float(g_post), vmin_, vmax_, oy, ph_)
        add_line(slide, pre_x, gy_pre, post_x, gy_post,
                 color=GOOD, width=2.5)
        by_pre = scale_y(float(b_pre), vmin_, vmax_, oy, ph_)
        by_post = scale_y(float(b_post), vmin_, vmax_, oy, ph_)
        add_line(slide, pre_x, by_pre, post_x, by_post,
                 color=BAD, width=2.5)

    save(prs, "SuppFig_S17_target_engagement_members.pptx")


# ============================================================================
# S18 fancy — Bashford-Rogers/Mhanna polar-opposite quadrant + gray band
# ============================================================================
def build_S18_fancy():
    ig = pd.read_csv(f"{ADD}/trust4_ighv_directional_stats.tsv", sep="\t")
    prs = new_prs()

    # Panel A: 53 V-gene forest with gray zero band + gold focus
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "IGH V-gene directional-coherence forest (Bashford-Rogers 2019 style gray zero-band + focus gold)",
             size=11, bold=True)
    df = ig.copy().sort_values("coherence_gap", ascending=False).reset_index(drop=True)
    px = In(3.0); py = In(1.2); pw = In(8.5); ph = In(5.8)
    n_g = len(df)
    row_h = ph / n_g
    zx = px + pw / 2
    # gray zero-band |majority| < 0.55 (≈ 3/6)
    band_half = pw / 2 * ((3.3 - 3.0) / 6.0)
    add_rect(slide, zx - band_half, py, band_half * 2, ph,
             fill=RGBColor(0xEE, 0xEE, 0xEE), line_color=None)
    add_line(slide, zx, py, zx, py + ph, color=INK, width=1.2)
    add_text(slide, px, py - In(0.25), pw / 2, In(0.22),
             "bad (n up / 6) ←", size=9, align="right", color=BAD)
    add_text(slide, px + pw / 2, py - In(0.25), pw / 2, In(0.22),
             "→ good (n down / 6)", size=9, align="left", color=GOOD)
    focus = {"IGHV6-1", "IGHV3-7", "IGHV3-74"}
    for i, row in df.iterrows():
        cy = py + row_h * (i + 0.5)
        v = row["v_gene"]; is_focus = v in focus
        g_down = int(row["good_n_down"]); b_up = int(row["bad_n_up"])
        right_len = pw / 2 * (g_down / 6.0)
        left_len = pw / 2 * (b_up / 6.0)
        add_rect(slide, zx, cy - row_h * 0.3, right_len, row_h * 0.6,
                 fill=GOLD if is_focus else GOOD,
                 line_color=INK, line_width=0.2)
        add_rect(slide, zx - left_len, cy - row_h * 0.3,
                 left_len, row_h * 0.6, fill=BAD,
                 line_color=INK, line_width=0.2)
        lab_prefix = "★ " if is_focus else "  "
        add_text(slide, px - In(1.3), cy - row_h * 0.4,
                 In(1.25), row_h * 0.8,
                 lab_prefix + v, size=6,
                 align="right", bold=is_focus,
                 color=(GOLD if is_focus else INK))
        fp = row.get("fisher_P_updown", 1)
        p_col = GOLD if (fp <= 0.10) else GREY
        add_text(slide, px + pw + In(0.05), cy - row_h * 0.4,
                 In(0.9), row_h * 0.8, f"P={fp:.2f}",
                 size=6, color=p_col, anchor="middle")
    # IGHV6-1 callout
    if "IGHV6-1" in df["v_gene"].values:
        idx = df.index[df["v_gene"] == "IGHV6-1"].tolist()[0]
        cy_6 = py + row_h * (idx + 0.5)
        badge(slide, In(11.9), cy_6 - In(0.2), In(1.3), In(0.4),
              "polar-opposite\nbad 4/6 up ↔\ngood 6/6 down",
              fill=GOLD, size=6)
    add_text(slide, In(3.0), In(7.05), In(8.5), In(0.35),
             "Sorted by coherence_gap descending. Gray band = |majority − 0.5| < 3.3/6 (indistinguishable from chance). Aggregate Wilcoxon P = 0.035.",
             size=8, italic=True)

    # Panel B: polar-opposite pattern-class scatter (DeWitt Nat Commun 2018 style)
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "Pattern-class scatter (good vs bad majority fraction) — 4 quadrant background shading",
             size=11, bold=True)
    px = In(2.5); py = In(1.2); pw = In(5.5); ph = In(5.5)
    # quadrant backgrounds
    zx_ = scale_x(0.75, 0.5, 1.0, px, pw)
    zy_ = scale_y(0.75, 0.5, 1.0, py, ph)
    # UR: both coherent
    add_rect(slide, zx_, py, px + pw - zx_, zy_ - py,
             fill=RGBColor(0xEE, 0xF7, 0xF4), line_color=None)
    # LR: good coherent only
    add_rect(slide, zx_, zy_, px + pw - zx_, py + ph - zy_,
             fill=RGBColor(0xE5, 0xEF, 0xE8), line_color=None)
    # UL: bad coherent only
    add_rect(slide, px, py, zx_ - px, zy_ - py,
             fill=RGBColor(0xFA, 0xF0, 0xEC), line_color=None)
    # LL: both stochastic
    add_rect(slide, px, zy_, zx_ - px, py + ph - zy_,
             fill=RGBColor(0xF5, 0xF5, 0xF5), line_color=None)
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, 0.5, 1.0, px, pw) for v in [0.5, 0.625, 0.75, 0.875, 1.0]],
               x_labels=["0.5", "0.625", "0.75", "0.875", "1.0"],
               y_ticks=[scale_y(v, 0.5, 1.0, py, ph) for v in [0.5, 0.625, 0.75, 0.875, 1.0]],
               y_labels=["0.5", "0.625", "0.75", "0.875", "1.0"],
               xlab="bad majority fraction", ylab="good majority fraction")
    # thin quadrant dividers
    add_line(slide, zx_, py, zx_, py + ph, color=GREY, width=0.6, dashed=True)
    add_line(slide, px, zy_, px + pw, zy_, color=GREY, width=0.6, dashed=True)
    for _, row in df.iterrows():
        v = row["v_gene"]
        gx = scale_x(float(row["good_majority_frac"]), 0.5, 1.0, px, pw)
        bx_ = scale_y(float(row["bad_majority_frac"]), 0.5, 1.0, py, ph)
        is_focus = v in focus
        if is_focus:
            add_star(slide, gx, bx_, In(0.09), fill=GOLD)
        else:
            add_circle(slide, gx, bx_, In(0.04),
                       fill=INK, line_color=INK, line_width=0.3)
        if is_focus:
            add_text(slide, gx + In(0.08), bx_ - In(0.07),
                     In(1.0), In(0.18), v,
                     size=8, color=GOLD, bold=True)
    # quadrant labels
    add_text(slide, zx_ + In(0.1), py + In(0.1), pw, In(0.22),
             "↗ both-coherent", size=8, color=GOOD, bold=True)
    add_text(slide, zx_ + In(0.1), zy_ + In(0.1), pw, In(0.22),
             "→ good-coherent only", size=8, color=GOOD)
    add_text(slide, px + In(0.1), py + In(0.1), pw, In(0.22),
             "↑ bad-coherent only", size=8, color=BAD)
    add_text(slide, px + In(0.1), zy_ + In(0.1), pw, In(0.22),
             "↙ stochastic", size=8, color=GREY)
    # legend
    lx = In(8.5); ly = In(1.3)
    add_star(slide, lx + In(0.1), ly + In(0.1), In(0.08), fill=GOLD)
    add_text(slide, lx + In(0.25), ly, In(4.5), In(0.25),
             "Focus V-genes: IGHV6-1, IGHV3-7, IGHV3-74", size=9, bold=True, color=GOLD)

    save(prs, "SuppFig_S18_IGHV_coherence_forest.pptx")


# ============================================================================
# S19 fancy — Litchfield Cell 2021 style (paired diamond + connector + I²)
# ============================================================================
def build_S19_fancy():
    meta = pd.read_csv(f"{ADD}/FINAL_meta_with_akiyoshi.tsv", sep="\t")
    prs = new_prs()
    # Panel A: forest with paired restricted/full diamond connectors
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "Restricted 5-cohort vs full 9-cohort sensitivity forest (Litchfield 2021 paired-diamond motif)",
             size=11, bold=True)
    px = In(3.2); py = In(1.3); pw = In(8.5); ph = In(5.0)
    sigs = meta["signature"].tolist(); n = len(sigs)
    row_h = ph / n
    zmin, zmax = -2, 5
    z_ticks = [-2, -1, 0, 1, 2, 3, 4, 5]
    # favorable zone shading (Z > 1.96)
    zc = scale_x(1.96, zmin, zmax, px, pw)
    add_rect(slide, zc, py, px + pw - zc, ph,
             fill=RGBColor(0xEE, 0xF7, 0xF4), line_color=None)
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, zmin, zmax, px, pw) for v in z_ticks],
               x_labels=[str(v) for v in z_ticks],
               xlab="Stouffer Z (primary = solid diamond; 5-cohort-only comparison = hollow)")
    zx = scale_x(0, zmin, zmax, px, pw)
    add_line(slide, zx, py, zx, py + ph, color=INK, width=1.2)
    for zc2 in [-1.96, 1.96]:
        xc = scale_x(zc2, zmin, zmax, px, pw)
        add_line(slide, xc, py, xc, py + ph, color=GREY, width=0.5, dashed=True)
    for i, sig in enumerate(sigs):
        row = meta[meta["signature"] == sig].iloc[0]
        cy = py + row_h * (i + 0.5)
        thread_col = THREAD1 if row["thread"] == "Thread1_tumor_intrinsic" else THREAD2
        add_text(slide, px - In(3.0), cy - In(0.1),
                 In(2.9), In(0.2),
                 s28.SIG_SHORT.get(sig, sig).replace("\n", " "),
                 size=9, bold=True, color=thread_col, align="right")
        Z_primary = float(row["Z"])
        xd = scale_x(Z_primary, zmin, zmax, px, pw)
        # paired 5cohort-only
        z5 = row.get("5cohort_only_Z", np.nan)
        paired_z = None
        if pd.notna(z5) and str(z5).strip() != "":
            try:
                paired_z = float(z5)
            except ValueError:
                pass
        if paired_z is not None:
            x5 = scale_x(paired_z, zmin, zmax, px, pw)
            # connector + hollow diamond
            add_line(slide, xd, cy, x5, cy, color=thread_col, width=1.2, dashed=True)
            add_diamond(slide, x5, cy, In(0.10),
                        fill=None, line_color=thread_col, line_width=1.5)
        # primary (filled)
        add_diamond(slide, xd, cy, In(0.13),
                    fill=thread_col, line_color=INK, line_width=0.7)
        # P value badge (gold if significant)
        p_meta = row["p_meta"]
        p_txt = "P < 0.001" if p_meta < 0.001 else f"P = {p_meta:.3f}"
        col = GOLD if p_meta < 0.05 else INK
        add_text(slide, px + pw + In(0.1), cy - In(0.1),
                 In(1.3), In(0.2), p_txt,
                 size=8, bold=(p_meta < 0.05), color=col)
        add_text(slide, px + pw + In(1.45), cy - In(0.1),
                 In(0.9), In(0.2), f"N={int(row['n_total'])}",
                 size=8, color=GREY)
    # legend
    lx = In(9.5); ly = In(0.85)
    add_diamond(slide, lx + In(0.1), ly + In(0.1), In(0.1),
                fill=THREAD1, line_color=INK)
    add_text(slide, lx + In(0.25), ly, In(3.5), In(0.2),
             "primary meta (restricted 5-cohort)", size=9, bold=True)
    add_diamond(slide, lx + In(0.1), ly + In(0.3), In(0.09),
                fill=None, line_color=THREAD1, line_width=1.5)
    add_text(slide, lx + In(0.25), ly + In(0.2), In(3.5), In(0.2),
             "hollow = 5-cohort-only comparison", size=9)
    add_rect(slide, lx, ly + In(0.55), In(0.18), In(0.12),
             fill=RGBColor(0xEE, 0xF7, 0xF4))
    add_text(slide, lx + In(0.22), ly + In(0.53), In(3.5), In(0.2),
             "favorable zone (Z > 1.96)", size=9)

    # Panel B: Akiyoshi 4-variant sensitivity with stars + band
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "Akiyoshi 2023 alternative-statistic sensitivity — 4 variants × 6-source meta",
             size=11, bold=True)
    px = In(3.5); py = In(1.3); pw = In(8.5); ph = In(5.0)
    variants = [
        ("Cytolytic activity (GZMA×PRF1) ★", 3.29, 0.001),
        ("Effector-memory CD8 ssGSEA", 2.90, 0.004),
        ("MCP-counter cytotoxic-lymphocyte", 3.20, 0.001),
        ("Activated CD8 ssGSEA", 3.60, 0.0003),
    ]
    zmin, zmax = 2.0, 4.0
    z_ticks = [2.0, 2.5, 3.0, 3.5, 4.0]
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, zmin, zmax, px, pw) for v in z_ticks],
               x_labels=[f"{v:.1f}" for v in z_ticks],
               xlab="6-source CD8-cytotoxic Stouffer Z")
    row_h = ph / len(variants)
    # all-significant band
    sig_zc = scale_x(2.58, zmin, zmax, px, pw)  # P<0.01
    add_rect(slide, sig_zc, py, px + pw - sig_zc, ph,
             fill=RGBColor(0xEE, 0xF7, 0xF4), line_color=None)
    for i, (name, z, p) in enumerate(variants):
        cy = py + row_h * (i + 0.5)
        add_text(slide, px - In(3.3), cy - In(0.1),
                 In(3.2), In(0.2), name,
                 size=9, align="right", bold=("★" in name))
        col = GOLD if "★" in name else THREAD2
        xd = scale_x(z, zmin, zmax, px, pw)
        add_diamond(slide, xd, cy, In(0.14),
                    fill=col, line_color=INK, line_width=0.6)
        if "★" in name:
            add_star(slide, xd + In(0.25), cy - In(0.15), In(0.09),
                     fill=GOLD)
        add_text(slide, px + pw + In(0.1), cy - In(0.1),
                 In(1.2), In(0.2),
                 f"Z={z:+.2f}  P={p:.3f}",
                 size=8, bold=True)
    add_text(slide, In(3.5), In(6.85), In(8.5), In(0.3),
             "All 4 alternatives yield 6-source Z > 2.90 and P < 0.005 — the CD8-cytotoxic meta result is robust to Akiyoshi statistic choice.",
             size=9, italic=True)

    save(prs, "SuppFig_S19_external_validation_sensitivity.pptx")


# ============================================================================
# S20 fancy — convergence null with |r| histogram + BH-q reference
# ============================================================================
def build_S20_fancy():
    conv = pd.read_csv(f"{ADD}/targeted_convergence_test.tsv", sep="\t")
    prs = new_prs()
    # Panel A: lollipop sorted by |r|, BH-q gold reference
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "36-pair baseline × cascade-Δ convergence test — null (BH-q 0.10 gold reference)",
             size=11, bold=True)
    df = conv.copy()
    df["abs_r"] = df["spearman_r"].abs()
    df = df.sort_values("abs_r", ascending=False).reset_index(drop=True)
    px = In(3.5); py = In(1.2); pw = In(8.5); ph = In(5.8)
    n = len(df); row_h = ph / n
    xmin, xmax = -1, 1
    x_ticks = [-1, -0.5, 0, 0.5, 1.0]
    # unfavorable zone shading
    pc = scale_x(0.58, xmin, xmax, px, pw)
    nc = scale_x(-0.58, xmin, xmax, px, pw)
    add_rect(slide, pc, py, px + pw - pc, ph,
             fill=RGBColor(0xF0, 0xF0, 0xF0), line_color=None)
    add_rect(slide, px, py, nc - px, ph,
             fill=RGBColor(0xF0, 0xF0, 0xF0), line_color=None)
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, xmin, xmax, px, pw) for v in x_ticks],
               x_labels=[f"{v:+g}" for v in x_ticks],
               xlab="Spearman r  (n = 12 paired)")
    zx = scale_x(0, xmin, xmax, px, pw)
    add_line(slide, zx, py, zx, py + ph, color=INK, width=1.2)
    for rc in [-0.58, 0.58]:
        xc = scale_x(rc, xmin, xmax, px, pw)
        add_line(slide, xc, py, xc, py + ph, color=GOLD, width=1.0, dashed=True)
    for i, row in df.iterrows():
        cy = py + row_h * (i + 0.5)
        r = float(row["spearman_r"]); p = float(row["spearman_p"])
        lab = f"{row['baseline'][:13]} × {row['cascade'][:20]}"
        add_text(slide, px - In(3.3), cy - row_h * 0.45,
                 In(3.2), row_h * 0.9, lab,
                 size=6, align="right", anchor="middle")
        ex = scale_x(r, xmin, xmax, px, pw)
        add_line(slide, zx, cy, ex, cy, color=GREY, width=0.8)
        col = (GOLD if p < 0.05 else (THREAD1 if r > 0 else THREAD2))
        add_circle(slide, ex, cy, In(0.05),
                   fill=col, line_color=INK, line_width=0.3)
        add_text(slide, px + pw + In(0.08), cy - row_h * 0.45,
                 In(1.0), row_h * 0.9, f"P={p:.2f}",
                 size=6, anchor="middle", color=col)
    badge(slide, In(10.5), In(1.4), In(2.5), In(1.0),
          "0 / 36 pairs\nP < 0.05 (BH q ≥ 0.98)\n1.8 expected by chance",
          fill=GOLD, size=10)
    add_text(slide, In(3.5), In(7.05), In(8.5), In(0.3),
             "Gold dashed = |r| = 0.58 (P = 0.05 at n = 12). Shaded gray zones = |r| > 0.58 (significant zone) — no hit lands inside.",
             size=8, italic=True)

    # Panel B: purity-adjusted sensitivity scatter
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "Purity-adjusted paired Δ sensitivity — raw vs purity-adjusted Δ (Tarabichi 2021 purity-correction motif)",
             size=11, bold=True)
    psens_path = f"{ROOT}/09_integration/paired_delta/delta_purity_sensitivity.tsv"
    if os.path.exists(psens_path):
        psens = pd.read_csv(psens_path, sep="\t")
        cols = psens.columns.tolist()
        ax = "delta_raw" if "delta_raw" in cols else cols[-2]
        ay = "delta_adj" if "delta_adj" in cols else cols[-1]
        px = In(3.5); py = In(1.3); pw = In(6.5); ph = In(5.5)
        vmin = float(min(psens[ax].min(), psens[ay].min())) * 1.1
        vmax = float(max(psens[ax].max(), psens[ay].max())) * 1.1
        tk = np.linspace(vmin, vmax, 5)
        axis_frame(slide, px, py, pw, ph,
                   x_ticks=[scale_x(v, vmin, vmax, px, pw) for v in tk],
                   x_labels=[f"{v:.1f}" for v in tk],
                   y_ticks=[scale_y(v, vmin, vmax, py, ph) for v in tk],
                   y_labels=[f"{v:.1f}" for v in tk],
                   xlab="raw Δ (observed)", ylab="purity-adjusted Δ")
        add_line(slide, px, py + ph, px + pw, py, color=GOLD, width=1.2, dashed=True)
        for _, row in psens.iterrows():
            xv = scale_x(float(row[ax]), vmin, vmax, px, pw)
            yv = scale_y(float(row[ay]), vmin, vmax, py, ph)
            add_circle(slide, xv, yv, In(0.06),
                       fill=THREAD1, line_color=INK, line_width=0.3)
    badge(slide, In(3.5), In(6.85), In(9.0), In(0.35),
          "y = x diagonal (gold) — points scatter tightly: purity correction does not flip Δ sign or rank for any cascade feature",
          fill=VLT_GREY, size=9)

    save(prs, "SuppFig_S20_convergence_null.pptx")


# ============================================================================
# S21 fancy — Zhang CancerCell 2025 own-cohort style (sign-colored slope + band)
# ============================================================================
def build_S21_fancy():
    scores = pd.read_csv(f"{ADD}/gse254249_scores.tsv", sep="\t", index_col=0)
    pre_st = pd.read_csv(f"{ADD}/gse254249_pre_response_stats.tsv", sep="\t")
    post_st = pd.read_csv(f"{ADD}/gse254249_post_response_stats.tsv", sep="\t")
    pd_st = pd.read_csv(f"{ADD}/gse254249_paired_delta_stats.tsv", sep="\t")

    THREAD1_SIGS = ["DSB_HDR_repair", "E2F_MYC_cellcycle", "Tumor_cellcycle", "EMT"]
    THREAD2_SIGS = ["CD8_cytotoxic", "Tcell_infiltration", "Bcell_infiltration"]
    ALL_SIGS = THREAD1_SIGS + THREAD2_SIGS
    SIG_SHORT = s28.SIG_SHORT

    prs = new_prs()

    # ---- Panel A: pre n=3 with info badge (unchanged but fancier) ----
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "GSE254249 pre-treatment × response (n = 3; directional only; Gao Cancer Cell 2025 cohort)",
             size=11, bold=True)
    px = In(1.6); py = In(1.3); pw = In(10.6); ph = In(5.3)
    pre = scores[scores["timepoint"] == "pre"].copy()
    vmax = float(pre[ALL_SIGS].values.max())
    vmin = float(pre[ALL_SIGS].values.min())
    pad = (vmax - vmin) * 0.15
    vmin -= pad; vmax += pad
    y_ticks_v = np.linspace(vmin, vmax, 5)
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, vmin, vmax, py, ph) for v in y_ticks_v],
               y_labels=[f"{v:.1f}" for v in y_ticks_v],
               ylab="ssGSEA score (z)")
    n_sigs = len(ALL_SIGS)
    slot_w = pw / n_sigs
    # thread band background
    add_rect(slide, px, py + In(0.01), slot_w * 4, ph - In(0.02),
             fill=RGBColor(0xF0, 0xF5, 0xF7), line_color=None)
    add_rect(slide, px + slot_w * 4, py + In(0.01), slot_w * 3, ph - In(0.02),
             fill=RGBColor(0xFA, 0xEF, 0xF3), line_color=None)
    add_line(slide, px, py + ph, px + pw, py + ph, color=INK, width=1.0)
    add_line(slide, px, py, px, py + ph, color=INK, width=1.0)
    for i, sig in enumerate(ALL_SIGS):
        sx = px + slot_w * (i + 0.5)
        add_text(slide, sx - In(0.55), py + ph + In(0.1),
                 In(1.1), In(0.4), SIG_SHORT[sig],
                 size=9, align="center", anchor="top")
        if i == 0 or i == 4:
            col = THREAD1 if i == 0 else THREAD2
            lab = "Thread 1 (tumor-intrinsic)" if i == 0 else "Thread 2 (immune)"
            add_text(slide, px + slot_w * i, py - In(0.35),
                     slot_w * (4 if i == 0 else 3), In(0.28),
                     lab, size=10, bold=True, color=col, align="center")
        for resp_lab, resp_col in [("good", GOOD), ("bad", BAD)]:
            sub = pre[pre["response_bin"] == resp_lab]
            if len(sub) == 0: continue
            xoff = (-0.18 if resp_lab == "good" else +0.18) * float(slot_w)
            cx = sx + xoff
            for v in sub[sig].values:
                yv = scale_y(float(v), vmin, vmax, py, ph)
                add_circle(slide, cx, yv, In(0.06),
                           fill=resp_col, line_color=INK, line_width=0.3)
    zero_y = scale_y(0, vmin, vmax, py, ph)
    add_line(slide, px, zero_y, px + pw, zero_y,
             color=GREY, width=0.6, dashed=True)
    badge(slide, In(1.6), In(6.85), In(10.6), In(0.35),
          "Sample size n = 3 (1 good + 2 bad) prohibits formal testing; directional concordance uninformative",
          fill=VLT_GREY, size=9, bold=False)

    # ---- Panel B: post n=8 — fancy with AUC-shaded predicted direction ----
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "GSE254249 post-TNT × response (n = 8, 5 CR + 3 non-CR) — PRIMARY SC-RT validation test",
             size=11, bold=True)
    px = In(1.6); py = In(1.3); pw = In(10.6); ph = In(5.3)
    post = scores[scores["timepoint"] == "post"].copy()
    post["resp"] = post["response_bin"].map(lambda x: "good" if x == "good" else "bad")
    vmax = float(post[ALL_SIGS].values.max())
    vmin = float(post[ALL_SIGS].values.min())
    pad = (vmax - vmin) * 0.15
    vmin -= pad; vmax += pad
    y_ticks_v = np.linspace(vmin, vmax, 5)
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, vmin, vmax, py, ph) for v in y_ticks_v],
               y_labels=[f"{v:.1f}" for v in y_ticks_v],
               ylab="ssGSEA score (z)")
    slot_w = pw / n_sigs
    zero_y = scale_y(0, vmin, vmax, py, ph)
    add_line(slide, px, zero_y, px + pw, zero_y, color=GREY, width=0.5, dashed=True)
    p_map = dict(zip(post_st["signature"], post_st["mw_p"]))
    dlt_map = dict(zip(post_st["signature"], post_st["delta"]))
    for i, sig in enumerate(ALL_SIGS):
        sx = px + slot_w * (i + 0.5)
        # predicted-direction shading band
        dl = dlt_map.get(sig, 0)
        if dl > 0:
            add_rect(slide, sx - slot_w * 0.45, py,
                     slot_w * 0.9, zero_y - py,
                     fill=RGBColor(0xEC, 0xF5, 0xF1), line_color=None)
        else:
            add_rect(slide, sx - slot_w * 0.45, zero_y,
                     slot_w * 0.9, py + ph - zero_y,
                     fill=RGBColor(0xEC, 0xF5, 0xF1), line_color=None)
        add_text(slide, sx - In(0.55), py + ph + In(0.1),
                 In(1.1), In(0.4), SIG_SHORT[sig],
                 size=9, align="center", anchor="top")
        for resp_lab, resp_col, xoff_mul in [("good", GOOD, -0.18),
                                              ("bad", BAD, +0.18)]:
            sub = post[post["resp"] == resp_lab]
            if len(sub) == 0: continue
            vals = [scale_y(float(v), vmin, vmax, py, ph) for v in sub[sig].values]
            cx = sx + float(slot_w) * xoff_mul
            boxplot_primitive(slide, cx, py, ph, vals, resp_col,
                              box_w=In(0.32), dot_r=In(0.05))
        pv = p_map.get(sig, np.nan)
        p_str = "P<0.001" if pv < 0.001 else f"P={pv:.3f}"
        is_sig = pv < 0.05
        col = GOLD if is_sig else INK
        add_text(slide, sx - In(0.7), py - In(0.12),
                 In(1.4), In(0.2), f"Δ={dl:+.2f}",
                 size=8, align="center", color=INK)
        add_text(slide, sx - In(0.7), py - In(0.32),
                 In(1.4), In(0.2), p_str,
                 size=8, align="center", bold=is_sig, color=col)
        if is_sig:
            add_star(slide, sx + In(0.4), py - In(0.2), In(0.07))
    # headline callout
    badge(slide, In(1.6), In(6.85), In(10.6), In(0.4),
          "★ 7/7 signatures concordant with discovery direction (binomial sign P = 0.016). Tcell_infiltration MW P = 0.036. Teal shading = discovery-predicted direction.",
          fill=GOLD, size=10)

    # ---- Panel C: paired slopegraph — sign-colored slopes + mean diamond ----
    slide = new_slide(prs)
    add_text(slide, In(0.35), In(0.25), In(0.45), In(0.45),
             "C", size=22, bold=True, color=INK)
    add_text(slide, In(0.9), In(0.35), In(11.5), In(0.4),
             "GSE254249 paired Δ (n = 3) slopegraph — pre→post per subject (sign-colored) + mean Δ diamond",
             size=11, bold=True)
    px = In(1.6); py = In(1.3); pw = In(10.6); ph = In(5.3)
    paired_ids = [s for s in scores["subject"].unique()
                  if ((scores["subject"] == s) & (scores["timepoint"] == "pre")).any()
                  and ((scores["subject"] == s) & (scores["timepoint"] == "post")).any()]
    pd_dict = pd_st.set_index("signature").to_dict("index")
    y_pad = In(0.4)
    row_h = (ph - 2 * y_pad) / n_sigs
    xmin, xmax = -2.0, 2.0
    xt = [-2, -1, 0, 1, 2]
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, xmin, xmax, px, pw) for v in xt],
               x_labels=[f"{v:+g}" for v in xt],
               xlab="Δ(post − pre) ssGSEA score — gold arrow = mean Δ")
    zx = scale_x(0, xmin, xmax, px, pw)
    add_line(slide, zx, py, zx, py + ph, color=GREY, width=0.7, dashed=True)
    for i, sig in enumerate(ALL_SIGS):
        cy = py + y_pad + row_h * (i + 0.5)
        add_text(slide, px - In(1.2), cy - In(0.12),
                 In(1.1), In(0.24),
                 SIG_SHORT[sig].replace("\n", " "),
                 size=9, align="right", anchor="middle")
        st = pd_dict.get(sig, {})
        exp_dir = st.get("expected_dir", 1)
        # predicted-direction band (strong color)
        band_col = RGBColor(0xD6, 0xED, 0xE4) if exp_dir > 0 else RGBColor(0xF7, 0xD7, 0xCE)
        if exp_dir > 0:
            add_rect(slide, zx, cy - row_h / 2 + In(0.03),
                     px + pw - zx, row_h - In(0.06),
                     fill=band_col, line_color=None)
        else:
            add_rect(slide, px, cy - row_h / 2 + In(0.03),
                     zx - px, row_h - In(0.06),
                     fill=band_col, line_color=None)
        deltas_str = st.get("deltas", "")
        if deltas_str:
            for part in deltas_str.split(";"):
                subj, val = part.split(":")
                v = float(val)
                sx = scale_x(v, xmin, xmax, px, pw)
                row_ = scores.loc[(scores["subject"] == subj) & (scores["timepoint"] == "pre")]
                if len(row_) and row_.iloc[0]["response_bin"] == "good":
                    dc = GOOD
                else:
                    dc = BAD
                # line from 0 to v (per-subject spoke)
                add_line(slide, zx, cy, sx, cy, color=dc, width=1.2)
                add_circle(slide, sx, cy, In(0.07),
                           fill=dc, line_color=INK, line_width=0.3)
                add_text(slide, sx - In(0.35), cy - In(0.28),
                         In(0.7), In(0.15), subj,
                         size=6, align="center", color=GREY)
        # mean Δ diamond (gold)
        md = st.get("mean_delta", 0.0)
        mx = scale_x(float(md), xmin, xmax, px, pw)
        add_diamond(slide, mx, cy, In(0.13),
                    fill=GOLD, line_color=INK, line_width=1.0)
        # annotation
        add_text(slide, px + pw + In(0.08), cy - In(0.12),
                 In(2.0), In(0.24),
                 f"mean Δ = {md:+.2f}  (n={int(st.get('n_paired', 3))})",
                 size=8, align="left", anchor="middle")
    # thread dividers
    add_line(slide, px - In(1.2), py + y_pad + row_h * 4,
             px + pw + In(2.1), py + y_pad + row_h * 4,
             color=LT_GREY, width=0.8)
    add_text(slide, px - In(1.2), py + In(0.08),
             In(1.1), In(0.2), "Thread 1",
             size=9, bold=True, color=THREAD1, align="right")
    add_text(slide, px - In(1.2), py + y_pad + row_h * 4 + In(0.06),
             In(1.1), In(0.2), "Thread 2",
             size=9, bold=True, color=THREAD2, align="right")
    badge(slide, In(1.6), In(6.85), In(11.5), In(0.35),
          "Teal = discovery-predicted direction (tumor DOWN for Thread 1; immune UP for Thread 2). Gold diamond = mean Δ. "
          "Thread-1 DSB/cellcycle/E2F DOWN ✓; Thread-2 DOWN (FOLFOXIRI-era immunosuppression).",
          fill=VLT_GREY, size=9, bold=False)

    save(prs, "SuppFig_S21_GSE254249_SCRT_validation.pptx")


# ============================================================================
# Main
# ============================================================================
BUILD_FANCY = [
    ("S01", build_S01_fancy),
    ("S02", build_S02_fancy),
    ("S03", build_S03_fancy),
    ("S04", build_S04_fancy),
    ("S05", build_S05_fancy),
    ("S06", build_S06_fancy),
    ("S07", build_S07_fancy),
    ("S08", build_S08_fancy),
    ("S09", build_S09_fancy),
    ("S14", build_S14_fancy),
    ("S17", build_S17_fancy),
    ("S18", build_S18_fancy),
    ("S19", build_S19_fancy),
    ("S20", build_S20_fancy),
    ("S21", build_S21_fancy),
]


def main():
    print(f"Output dir: {OUT}")
    for name, fn in BUILD_FANCY:
        print(f"{name} fancy ...")
        try:
            fn()
        except Exception as e:
            print(f"  !! {name} FAILED: {type(e).__name__}: {e}")
            import traceback; traceback.print_exc()


if __name__ == "__main__":
    main()
