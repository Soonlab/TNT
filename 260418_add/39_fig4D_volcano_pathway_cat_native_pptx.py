#!/usr/bin/env python3
"""
39_fig4D_volcano_pathway_cat_native_pptx.py

Build Fig 4D (DEG GENE-level volcano, coloured by Hallmark pathway
category) as native-editable PowerPoint — main + supp.

Matches matplotlib panel produced by
38_fig4D_volcano_gene_level_pathway_coloured.py. Distinct from Fig 3B
(which is pathway-level Hallmark NES × FDR bubble).

Design summary:
  - Every dot = one gene from the DEG (22,838 total, 228 at P < 0.01)
  - Significant genes coloured by Hallmark gene-set membership
        Cell cycle = E2F/G2M/MYC/Mitotic Spindle (n=484 genes)
        DNA repair = DNA_REPAIR + UV response (n=449)
        Immune     = IFN α/γ + Inflammatory + Complement + IL2/6 +
                     Allograft + TNFα + Coagulation (n=1041)
        EMT/Stromal= EMT + Myogenesis + Apical Junction + Angiogenesis +
                     Hedgehog (n=601)
        Other      = no pathway membership — coloured by direction
  - Gold ring emphasis for |log2FC| >= 1

Rules:
  1. One panel per slide; 2. no title; 3. native shapes + textboxes;
  4. main + supp; 5. Arial; 6. shadows killed.

Motif-informed from:
  - Mariathasan 2018 Nature (IMvigor210 DEG volcano)
  - Sade-Feldman 2018 Cell (pathway-category volcano variant)
  - Love 2014 DESeq2 / Blighe EnhancedVolcano (quadruple classifier)
  - Litchfield 2021 Cell (gold emphasis ring for large effects)
"""
from pathlib import Path
import numpy as np
import pandas as pd
import gseapy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn
from lxml import etree


def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def _i(v): return int(v)


GOOD      = RGBColor(0x0A, 0x7D, 0x6E)
BAD       = RGBColor(0xC5, 0x3E, 0x1F)
INK       = RGBColor(0x0E, 0x2A, 0x47)
LINE      = RGBColor(0x33, 0x33, 0x33)
GREY      = RGBColor(0x77, 0x77, 0x77)
LT_GREY   = RGBColor(0xC0, 0xC6, 0xCF)
NS_CLOUD  = RGBColor(0xDD, 0xE2, 0xE8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xD4, 0xA3, 0x00)

CAT_COLOR = {
    'Cell cycle':  RGBColor(0x05, 0x7A, 0x64),
    'DNA repair':  RGBColor(0x00, 0x56, 0x7D),
    'Immune':      RGBColor(0xC1, 0x14, 0x56),
    'EMT/Stromal': RGBColor(0xB0, 0x32, 0x19),
    'Other-good':  GOOD,
    'Other-poor':  BAD,
}

FONT = "Arial"
OUT = Path("/data/data/TNT/analysis/260418_add/ppt")
DATA_ROOT = Path("/data/data/TNT/analysis")
OUT.mkdir(parents=True, exist_ok=True)

# ---- load DEG + Hallmark gene sets ----
deg = pd.read_csv(DATA_ROOT/"05_rna_deg_gsea/DEG_good_vs_bad_pre.tsv", sep="\t")
deg['-log10p'] = -np.log10(deg.pvalue.replace(0, 1e-300))
deg = deg.dropna(subset=['log2FoldChange', '-log10p']).copy()
deg['log2FC_clip'] = deg.log2FoldChange.clip(-5, 5)
deg['nlp_clip']    = deg['-log10p'].clip(0, 8)

hm_sets = gseapy.get_library(name='MSigDB_Hallmark_2020', organism='Human')
def union(*keys): return set().union(*[set(hm_sets.get(k, [])) for k in keys])

CELL_CYCLE = union('E2F Targets','G2-M Checkpoint','MYC Targets V1',
                   'MYC Targets V2','Mitotic Spindle')
DNA_REPAIR = union('DNA Repair','UV Response Up','UV Response Dn')
IMMUNE     = union('Interferon Alpha Response','Interferon Gamma Response',
                   'Inflammatory Response','Complement',
                   'IL-2/STAT5 Signaling','IL-6/JAK/STAT3 Signaling',
                   'Allograft Rejection','TNF-alpha Signaling via NF-kB',
                   'Coagulation')
EMT_STROMAL = union('Epithelial Mesenchymal Transition','Myogenesis',
                    'Apical Junction','Angiogenesis','Hedgehog Signaling')

def gcat(g):
    if g in CELL_CYCLE:  return 'Cell cycle'
    if g in DNA_REPAIR:  return 'DNA repair'
    if g in IMMUNE:      return 'Immune'
    if g in EMT_STROMAL: return 'EMT/Stromal'
    return 'Other'

deg['cat'] = deg.gene.apply(gcat)

P_THRESH = 0.01
FC_THRESH = 1.0
NLP_THRESH = -np.log10(P_THRESH)

sig_mask = deg.pvalue < P_THRESH
sig = deg[sig_mask].copy()
ns  = deg[~sig_mask]

def marker_col(r):
    if r['cat'] != 'Other':
        return CAT_COLOR[r['cat']]
    return CAT_COLOR['Other-good'] if r.log2FoldChange > 0 else CAT_COLOR['Other-poor']

sig['m_color'] = sig.apply(marker_col, axis=1)
fc_big = sig[sig.log2FoldChange.abs() >= FC_THRESH]

cat_counts = {c: int((sig['cat'] == c).sum()) for c in
              ['Cell cycle','DNA repair','Immune','EMT/Stromal','Other']}
n_other_up = int(((sig['cat']=='Other') & (sig.log2FoldChange>0)).sum())
n_other_dn = int(((sig['cat']=='Other') & (sig.log2FoldChange<0)).sum())
print(f"sig n={len(sig)}  cats={cat_counts}  other up-good={n_other_up} up-poor={n_other_dn}")


# ---- PPT primitives ----
def new_slide(prs, w, h):
    prs.slide_width = w; prs.slide_height = h
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, size=8, bold=False,
             color=INK, align="left", anchor="middle", italic=False,
             rotation=0):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))
    tf = tb.text_frame
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    tf.word_wrap = True
    tf.vertical_anchor = {"top":MSO_ANCHOR.TOP,"middle":MSO_ANCHOR.MIDDLE,
                          "bottom":MSO_ANCHOR.BOTTOM}.get(anchor, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = {"left":PP_ALIGN.LEFT,"right":PP_ALIGN.RIGHT,
                   "center":PP_ALIGN.CENTER}.get(align, PP_ALIGN.LEFT)
    r = p.add_run(); r.text = str(text)
    r.font.name = FONT; r.font.size = Pt(size)
    r.font.bold = bool(bold); r.font.italic = bool(italic)
    r.font.color.rgb = color
    if rotation: tb.rotation = rotation
    kill_shadow(tb)
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.5, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   _i(x1),_i(y1),_i(x2),_i(y2))
    c.line.color.rgb = color; c.line.width = Pt(width)
    if dashed: c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    kill_shadow(c)
    return c


def add_circle(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx-r),_i(cy-r),_i(2*r),_i(2*r))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None: shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color; shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_hollow(slide, cx, cy, r, edge, lw=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx-r),_i(cy-r),_i(2*r),_i(2*r))
    shp.fill.background()
    shp.line.color.rgb = edge; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_triangle_arrow(slide, x1, y1, x2, y2, color, lw=1.8):
    add_line(slide, x1, y1, x2, y2, color=color, width=lw)
    dx, dy = x2 - x1, y2 - y1
    ang = np.degrees(np.arctan2(dy, dx))
    head = Emu(60000)
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                 _i(x2-head/2),_i(y2-head/2),_i(head),_i(head))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.rotation = ang
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def build_slide(prs, slide_w, slide_h, *, ax_x, ax_y, ax_w, ax_h,
                ns_sample_n, label_top_n, variant="main", legend_cols=2):
    s = new_slide(prs, slide_w, slide_h)
    x_lo, x_hi = -5.2, 5.2
    y_lo, y_hi = 0.0, 8.5
    def tx(v): return _i(ax_x + (v-x_lo)/(x_hi-x_lo) * ax_w)
    def ty(v): return _i(ax_y + ax_h - (v-y_lo)/(y_hi-y_lo) * ax_h)

    # spines
    add_line(s, ax_x, ax_y+ax_h, ax_x+ax_w, ax_y+ax_h, LINE, 0.6)
    add_line(s, ax_x, ax_y,      ax_x,      ax_y+ax_h, LINE, 0.6)

    for v in range(-5, 6):
        xx = tx(v)
        add_line(s, xx, ax_y+ax_h, xx, ax_y+ax_h+Inches(0.05), LINE, 0.4)
        add_text(s, xx-Inches(0.20), ax_y+ax_h+Inches(0.06),
                 Inches(0.40), Inches(0.18),
                 f"{v:+d}" if v != 0 else "0",
                 size=7, color=INK, align="center")
    for v in range(0, 9):
        yy = ty(v)
        add_line(s, ax_x-Inches(0.05), yy, ax_x, yy, LINE, 0.4)
        add_text(s, ax_x-Inches(0.45), yy-Inches(0.09),
                 Inches(0.38), Inches(0.18),
                 f"{v}", size=7, color=INK, align="right")

    add_text(s, ax_x, ax_y+ax_h+Inches(0.28),
             ax_w, Inches(0.22),
             "log₂ fold change (good vs poor)",
             size=9, bold=True, color=INK, align="center")
    yl_W = Inches(1.40); yl_H = Inches(0.22)
    yl_cx = ax_x - Inches(0.55); yl_cy = ax_y + ax_h/2
    lb = add_text(s, yl_cx-yl_W/2, yl_cy-yl_H/2, yl_W, yl_H,
                  "−log₁₀(p-value)",
                  size=9, bold=True, color=INK, align="center")
    lb.rotation = -90

    # reference lines
    add_line(s, tx(0), ax_y, tx(0), ax_y+ax_h, INK, 0.6)
    add_line(s, ax_x, ty(NLP_THRESH), ax_x+ax_w, ty(NLP_THRESH),
             GREY, 0.6, dashed=True)
    add_text(s, ax_x+Inches(0.05), ty(NLP_THRESH)-Inches(0.20),
             Inches(0.9), Inches(0.18),
             f"P = {P_THRESH}", size=7, color=GREY, italic=True, align="left")

    # non-sig cloud
    rng = np.random.default_rng(42)
    ns_sample = ns.sample(min(ns_sample_n, len(ns)), random_state=42)
    r_ns = Emu(16000) if variant == "main" else Emu(14000)
    for _, g in ns_sample.iterrows():
        add_circle(s, tx(g.log2FC_clip), ty(g.nlp_clip), r_ns,
                   fill=NS_CLOUD, line_color=None)

    # sig dots — category-coloured (pathway first, "Other" direction last on TOP)
    # draw "Other" (206 dots) first so pathway dots land on top
    r_sig_other = Emu(22000) if variant == "main" else Emu(18000)
    r_sig_cat   = Emu(32000) if variant == "main" else Emu(26000)
    for _, g in sig[sig['cat'] == 'Other'].iterrows():
        add_circle(s, tx(g.log2FC_clip), ty(g.nlp_clip), r_sig_other,
                   fill=g['m_color'], line_color=INK, line_width=0.35)
    for _, g in sig[sig['cat'] != 'Other'].iterrows():
        add_circle(s, tx(g.log2FC_clip), ty(g.nlp_clip), r_sig_cat,
                   fill=g['m_color'], line_color=INK, line_width=0.5)

    # |log2FC| >= 1 gold emphasis
    r_ring = Emu(46000) if variant == "main" else Emu(38000)
    for _, g in fc_big.iterrows():
        add_hollow(s, tx(g.log2FC_clip), ty(g.nlp_clip), r_ring,
                   HIGHLIGHT, lw=1.1)

    # labels: top pathway sig + top "Other"
    pa = sig[sig['cat'] != 'Other'].copy()
    pa['score'] = pa['-log10p'] * pa.log2FoldChange.abs()
    pa = pa.sort_values('score', ascending=False).head(15 if variant=="main" else 25)
    ot = sig[sig['cat'] == 'Other'].copy()
    ot['score'] = ot['-log10p'] * ot.log2FoldChange.abs()
    ot = ot.sort_values('score', ascending=False).head(15 if variant=="main" else 35)
    label_df = pd.concat([pa, ot], ignore_index=True).head(label_top_n)
    rng2 = np.random.default_rng(5)
    for _, g in label_df.iterrows():
        is_right = g.log2FoldChange > 0
        bw = Inches(0.85)
        tx_ = (tx(g.log2FC_clip)+Inches(0.10)) if is_right else (tx(g.log2FC_clip)-bw-Inches(0.10))
        y_off = Inches(-0.08) + Emu(int(rng2.uniform(-50000,50000)))
        add_text(s, tx_, ty(g.nlp_clip) + y_off,
                 bw, Inches(0.16),
                 g.gene,
                 size=7 if variant=="main" else 7.5,
                 bold=True, color=g['m_color'],
                 align="left" if is_right else "right")

    # direction arrows
    y_arrow = ty(7.85)
    add_triangle_arrow(s, tx(0.7), y_arrow, tx(4.6), y_arrow, GOOD, lw=2.2)
    add_text(s, tx(0.7), y_arrow-Inches(0.30),
             tx(4.6)-tx(0.7), Inches(0.22),
             "Up in Good responders",
             size=9, bold=True, color=GOOD, align="center")
    add_triangle_arrow(s, tx(-0.7), y_arrow, tx(-4.6), y_arrow, BAD, lw=2.2)
    add_text(s, tx(-4.6), y_arrow-Inches(0.30),
             tx(-0.7)-tx(-4.6), Inches(0.22),
             "Up in Poor responders",
             size=9, bold=True, color=BAD, align="center")

    # legend
    leg_y = ax_y + ax_h + Inches(0.60)
    entries = [
        (CAT_COLOR['Cell cycle'],  None, 0.26, f"Cell cycle  (n = {cat_counts['Cell cycle']})"),
        (CAT_COLOR['DNA repair'],  None, 0.26, f"DNA repair  (n = {cat_counts['DNA repair']})"),
        (CAT_COLOR['Immune'],      None, 0.26, f"Immune  (n = {cat_counts['Immune']})"),
        (CAT_COLOR['EMT/Stromal'], None, 0.26, f"EMT / Stromal  (n = {cat_counts['EMT/Stromal']})"),
        (CAT_COLOR['Other-good'],  None, 0.22, f"Other, up-good  (n = {n_other_up})"),
        (CAT_COLOR['Other-poor'],  None, 0.22, f"Other, up-poor  (n = {n_other_dn})"),
        (None, HIGHLIGHT,          0.30, f"|log₂ FC| ≥ {FC_THRESH:.0f}  (n = {len(fc_big)})"),
        (NS_CLOUD, None,           0.18, f"Non-significant (P ≥ {P_THRESH})"),
    ]
    entry_w = ax_w / legend_cols
    row_gap = Inches(0.27)
    for idx, (fill_, edge_, r_in, label) in enumerate(entries):
        r_i = idx // legend_cols
        c_i = idx % legend_cols
        cx = ax_x + entry_w * c_i + Inches(0.15)
        cy = leg_y + r_i * row_gap + Inches(0.09)
        if fill_ is not None:
            add_circle(s, cx, cy, Emu(int(r_in*100000)),
                       fill=fill_, line_color=INK, line_width=0.3)
        else:
            add_hollow(s, cx, cy, Emu(int(r_in*100000)), edge_, lw=1.0)
        add_text(s, cx+Inches(0.15), leg_y+r_i*row_gap,
                 entry_w-Inches(0.30), Inches(0.18),
                 label, size=7.5, color=INK, align="left")


prs_main = Presentation()
build_slide(prs_main,
            slide_w=Inches(9.0), slide_h=Inches(7.2),
            ax_x=Inches(1.30), ax_y=Inches(0.45),
            ax_w=Inches(7.30), ax_h=Inches(4.70),
            ns_sample_n=500,
            label_top_n=25,
            variant="main",
            legend_cols=3)
out_main = OUT / "Fig4D_volcano_gene_pathwaycat_native_editable.pptx"
prs_main.save(out_main)
print(f"wrote MAIN → {out_main}")

prs_supp = Presentation()
build_slide(prs_supp,
            slide_w=Inches(11.5), slide_h=Inches(8.0),
            ax_x=Inches(1.45), ax_y=Inches(0.55),
            ax_w=Inches(9.40), ax_h=Inches(5.90),
            ns_sample_n=1200,
            label_top_n=55,
            variant="supp",
            legend_cols=4)
out_supp = OUT / "SuppFig_volcano_gene_pathwaycat_detailed_native_editable.pptx"
prs_supp.save(out_supp)
print(f"wrote SUPP → {out_supp}")

from pptx import Presentation as _Pr
for p in (out_main, out_supp):
    pr = _Pr(str(p))
    print(f"{p.name}: {sum(len(list(s.shapes)) for s in pr.slides)} shapes")
