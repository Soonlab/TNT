#!/usr/bin/env python3
"""
20_fig7_native_pptx.py

Build Figure 7 (§3.10 Directed vs stochastic immune-repertoire response)
and its supplementary companion Supp Fig S18 (full 53-V-gene forest +
per-V-gene pattern class map) as native-editable PowerPoint decks.

Rules (same as scripts 18/19):
  * one panel per slide; no plot titles
  * every in-plot text element is a add_textbox TEXT_BOX (editable,
    not grouped)
  * every line / tick / marker / bar / cell is a native connector or
    auto-shape; all coordinates passed through _i() to guarantee
    integer EMU values (prevents PowerPoint "width=0.0" load errors)
  * Arial font throughout
  * TNT project palette: GOOD=#0a7d6e / BAD=#c53e1f
  * all shapes get an empty <a:effectLst/> via kill_shadow()

Main Figure 7 panels
--------------------
  A  Repertoire-coherence aggregate evidence on a single slide, three
     sub-panels arranged horizontally:
       A1 scatter of per-V-gene majority-fraction, good (y) vs bad (x),
          with y=x diagonal; 53 V-genes coloured by pattern class;
          points above the diagonal = good more coherent (the
          target pattern); IGHV6-1 / IGHV3-7 / IGHV3-74 labelled.
       A2 Wilcoxon distribution of per-V-gene (good − bad) majority-
          fraction gap: vertical strip with box + jittered points +
          zero reference, one-sided Wilcoxon P = 0.035 annotated.
       A3 horizontal stacked pattern bar: 16 good_coherent_bad_mixed
          vs 7 bad_coherent_good_mixed (plus the 2 both-coherent
          categories and 28 both-mixed); binomial P = 0.047
          annotated on the 16:7 contrast.
  B  Focus-V-gene spaghetti 3 × 3: IGHV6-1 (top-left; ★ strongest),
     IGHV3-7 and IGHV3-74 (user-prior), and the top six of the
     good_coherent_bad_mixed class ordered by Fisher P. Each sub-plot
     has pre → post axis, 12 subject slopes coloured by response,
     sign-count annotation, and a small "pred g:n/6, b:n/6" tag.

Supp Figure S18 panels
----------------------
  A  All 53 coverage-filtered V-genes as a horizontal forest plot
     ordered by coherence_gap (descending). For each V-gene a bilateral
     bar shows the good-6 predicted-direction count on one side and
     the bad-6 predicted-direction count on the other, with per-V-gene
     Fisher P annotated. User-focus genes (IGHV6-1, IGHV3-7, IGHV3-74)
     highlighted with a gold star + label box.
  B  V-gene pattern class map: each V-gene plotted as a point in the
     (majority_fraction good vs bad) plane, coloured by pattern, with
     every V-gene name labelled (small 4-5 pt). Marginal histograms
     at top and right summarise the distributions.

Motif references consulted (V-gene / repertoire coherence convention)
---------------------------------------------------------------------
  - Yost et al Nat Med 2019 (CD8 clonal replacement): clone trajectory
    scatter with diagonal reference and group colouring.
  - Wang et al Cell 2021 (B-cell repertoire in melanoma): V-gene
    usage heatmap and expansion bar.
  - Helmink et al Nature 2020 / Cabrita et al Nature 2020 (TLS ICB):
    concordance bar with n/N labelling.
  - Thorsson et al Immunity 2018 (TCGA immune landscape): signed-
    concordance stacked bars.
  - Petitprez et al Nature 2020 (TLS sarcoma): V-gene focus small-
    multiple spaghetti.
"""

import os
import numpy as np
import pandas as pd
from scipy import stats
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Shared infrastructure (identical to 19_fig6_native_pptx.py)
# ---------------------------------------------------------------------------
def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def _i(v):
    return int(round(float(v)))


GOOD = RGBColor(0x0A, 0x7D, 0x6E)
BAD = RGBColor(0xC5, 0x3E, 0x1F)
INK = RGBColor(0x22, 0x22, 0x22)
LINE = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0xBB, 0xBB, 0xBB)
LT_GREY = RGBColor(0xDD, 0xDD, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xD4, 0xA3, 0x00)
PAT_BOTH_SAME = RGBColor(0x7F, 0xB0, 0x69)
PAT_BOTH_OPP = RGBColor(0xC5, 0xA5, 0x72)
PAT_MIXED = RGBColor(0xCC, 0xCC, 0xCC)
GOOD_HEX = (0x0A, 0x7D, 0x6E)
BAD_HEX = (0xC5, 0x3E, 0x1F)

FONT = "Arial"
SLIDE_W = Inches(6.5)
SLIDE_H = Inches(4.5)

OUT = "/data/data/TNT/analysis/260418_add/ppt"
DATA = "/data/data/TNT/analysis/260418_add"
os.makedirs(OUT, exist_ok=True)


def new_slide(prs, w=SLIDE_W, h=SLIDE_H):
    prs.slide_width = w
    prs.slide_height = h
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, size=8, bold=False,
             color=INK, align="left", anchor="middle", font=FONT,
             italic=False):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(max(w, 1)), _i(max(h, 1)))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "bottom": MSO_ANCHOR.BOTTOM,
                          "middle": MSO_ANCHOR.MIDDLE}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = str(text)
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.italic = bool(italic)
    r.font.color.rgb = color
    kill_shadow(tb)
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.5, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        try:
            c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass
    kill_shadow(c)
    return c


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0.5):
    w = max(_i(w), 1); h = max(_i(h), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _i(x), _i(y), w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_circle(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_diamond(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def lighten(rgb_hex, factor=0.7):
    r0, g0, b0 = rgb_hex
    return RGBColor(int(r0 + (255 - r0) * factor),
                    int(g0 + (255 - g0) * factor),
                    int(b0 + (255 - b0) * factor))


def draw_panel_letter(slide, letter):
    add_text(slide, Inches(0.15), Inches(0.1), Inches(0.4), Inches(0.35),
             letter, size=14, bold=True, color=INK, align="left")


# ---------------------------------------------------------------------------
# Load data
# ---------------------------------------------------------------------------
stats_df = pd.read_csv(f"{DATA}/trust4_ighv_directional_stats.tsv",
                       sep="\t")
delta_df = pd.read_csv(f"{DATA}/trust4_ighv_per_subject_delta.tsv", sep="\t")

# pattern palette
PATTERN_COLOR = {
    "good_coherent_bad_mixed": GOOD,
    "bad_coherent_good_mixed": BAD,
    "both_coherent_same": PAT_BOTH_SAME,
    "both_coherent_opposite": PAT_BOTH_OPP,
    "both_mixed": PAT_MIXED,
}
PATTERN_PRETTY = {
    "good_coherent_bad_mixed": "good coherent, bad mixed",
    "bad_coherent_good_mixed": "bad coherent, good mixed",
    "both_coherent_same": "both coherent, same direction",
    "both_coherent_opposite": "both coherent, opposite direction",
    "both_mixed": "both mixed",
}


# ===========================================================================
# MAIN FIGURE 7
# ===========================================================================
prs_main = Presentation()


def build_A():
    """Repertoire-coherence aggregate evidence: 3 sub-panels."""
    s = new_slide(prs_main)
    draw_panel_letter(s, "A")

    # ====== Sub-panel A1: scatter ======
    a1_x = Inches(0.55); a1_y = Inches(0.50)
    a1_w = Inches(2.05); a1_h = Inches(2.60)

    # inner plot box inside A1 header margin
    p1_x = a1_x + Inches(0.45); p1_y = a1_y + Inches(0.18)
    p1_w = a1_w - Inches(0.55); p1_h = a1_h - Inches(0.60)

    # axis range: majority fractions are always in [0.5, 1.0]
    lo, hi = 0.48, 1.03
    def x1(v): return _i(p1_x + (v - lo) / (hi - lo) * p1_w)
    def y1(v): return _i(p1_y + p1_h - (v - lo) / (hi - lo) * p1_h)

    # spines + ticks
    add_line(s, p1_x, p1_y, p1_x, p1_y + p1_h, LINE, 0.6)
    add_line(s, p1_x, p1_y + p1_h, p1_x + p1_w, p1_y + p1_h, LINE, 0.6)
    for v, lab in [(0.5, "0.5"), (0.7, "0.7"), (0.9, "0.9"), (1.0, "1.0")]:
        # x ticks
        add_line(s, x1(v), _i(p1_y + p1_h), x1(v),
                 _i(p1_y + p1_h + Inches(0.04)), LINE, 0.4)
        add_text(s, x1(v) - Inches(0.18), _i(p1_y + p1_h + Inches(0.05)),
                 Inches(0.36), Inches(0.13),
                 lab, size=6, color=INK, align="center")
        # y ticks
        add_line(s, _i(p1_x - Inches(0.04)), y1(v), p1_x, y1(v), LINE, 0.4)
        add_text(s, _i(p1_x - Inches(0.34)), y1(v) - Inches(0.07),
                 Inches(0.30), Inches(0.13),
                 lab, size=6, color=INK, align="right")

    # axis titles
    add_text(s, p1_x, _i(p1_y + p1_h + Inches(0.22)),
             p1_w, Inches(0.16),
             "bad-group majority fraction",
             size=7, color=INK, align="center")
    yt = add_text(s, _i(p1_x - Inches(0.55)),
                  _i(p1_y + p1_h / 2 - Inches(0.7)),
                  Inches(0.45), Inches(1.4),
                  "good-group majority fraction",
                  size=7, color=INK, align="center")
    yt.rotation = -90

    # diagonal reference line (y = x)
    add_line(s, x1(lo + 0.02), y1(lo + 0.02),
             x1(hi - 0.01), y1(hi - 0.01),
             GREY, 0.5, dashed=True)
    add_text(s, x1(0.85) + Inches(0.02), y1(0.85) - Inches(0.15),
             Inches(0.6), Inches(0.13),
             "y = x", size=5, color=GREY, align="left",
             italic=True)

    # "above diagonal → good more coherent" cue (rotated, top-left)
    cue = add_text(s, _i(p1_x + Inches(0.08)), _i(p1_y + Inches(0.05)),
                   Inches(1.3), Inches(0.20),
                   "↖ above diagonal = good more coherent",
                   size=5, color=GOOD, align="left", italic=True)

    # point scatter (53 V-genes), coloured by pattern
    focus = {"IGHV6-1", "IGHV3-7", "IGHV3-74"}
    rng = np.random.default_rng(0)
    for _, row in stats_df.iterrows():
        px = x1(row.bad_majority_frac)
        py = y1(row.good_majority_frac)
        color = PATTERN_COLOR.get(row.pattern, PAT_MIXED)
        # slight jitter to separate overlapping points
        jx = px + rng.integers(-_i(Inches(0.015)), _i(Inches(0.015)) + 1)
        jy = py + rng.integers(-_i(Inches(0.015)), _i(Inches(0.015)) + 1)
        r = Emu(20000) if row.v_gene not in focus else Emu(30000)
        if row.v_gene in focus:
            add_circle(s, jx, jy, r, fill=color,
                       line_color=HIGHLIGHT, line_width=1.2)
            # label
            add_text(s, jx + Inches(0.05), jy - Inches(0.07),
                     Inches(0.6), Inches(0.13),
                     row.v_gene, size=5, bold=True,
                     color=HIGHLIGHT, align="left")
        else:
            add_circle(s, jx, jy, r, fill=color,
                       line_color=WHITE, line_width=0.4)

    # ====== Sub-panel A2: Wilcoxon distribution of (good - bad) gap ======
    a2_x = Inches(2.80); a2_y = Inches(0.50)
    a2_w = Inches(1.55); a2_h = Inches(2.60)

    p2_x = a2_x + Inches(0.50); p2_y = a2_y + Inches(0.18)
    p2_w = a2_w - Inches(0.60); p2_h = a2_h - Inches(0.60)

    gaps = (stats_df.good_majority_frac - stats_df.bad_majority_frac).values
    g_lo = float(np.floor(gaps.min() * 10)) / 10 - 0.05
    g_hi = float(np.ceil(gaps.max() * 10)) / 10 + 0.05
    g_abs = max(abs(g_lo), abs(g_hi))
    g_lo, g_hi = -g_abs, g_abs

    def y2(v): return _i(p2_y + p2_h - (v - g_lo) / (g_hi - g_lo) * p2_h)

    # spines
    add_line(s, p2_x + p2_w / 2, p2_y, p2_x + p2_w / 2,
             p2_y + p2_h, LINE, 0.4)
    # zero reference
    add_line(s, p2_x, y2(0), p2_x + p2_w, y2(0), INK, 0.9)
    # y ticks
    for v in [-g_abs, -g_abs / 2, 0, g_abs / 2, g_abs]:
        yy = y2(v)
        add_line(s, _i(p2_x + p2_w / 2 - Inches(0.03)), yy,
                 _i(p2_x + p2_w / 2 + Inches(0.03)), yy, LINE, 0.4)
        add_text(s, _i(p2_x - Inches(0.05)), yy - Inches(0.08),
                 Inches(0.45), Inches(0.14),
                 f"{v:+.2f}" if abs(v) > 1e-4 else "0",
                 size=6, color=INK, align="right")

    # box: Q25, median, Q75 of gaps
    q25, q50, q75 = np.percentile(gaps, [25, 50, 75])
    box_w = Inches(0.6)
    strip_cx = p2_x + p2_w / 2 + Inches(0.18)
    add_rect(s, _i(strip_cx - box_w / 2), y2(q75),
             box_w, max(y2(q25) - y2(q75), 2),
             fill=lighten(GOOD_HEX, 0.78), line_color=GOOD,
             line_width=0.7)
    add_line(s, _i(strip_cx - box_w / 2), y2(q50),
             _i(strip_cx + box_w / 2), y2(q50),
             GOOD, 1.6)
    add_diamond(s, _i(strip_cx), y2(q50), Emu(35000),
                fill=GOOD, line_color=WHITE, line_width=1.0)

    # jitter points (all 53 V-genes)
    jitter_range = _i(Inches(0.28))
    for v in gaps:
        jx = _i(strip_cx) + rng.integers(-jitter_range // 2,
                                          jitter_range // 2 + 1)
        add_circle(s, jx, y2(v), Emu(18000),
                   fill=WHITE, line_color=GOOD, line_width=0.7)

    # Wilcoxon P annotation
    try:
        w_stat, w_p = stats.wilcoxon(gaps, alternative="greater")
    except Exception:
        w_p = np.nan
    ann_x = p2_x; ann_y = p2_y + Inches(0.05)
    add_rect(s, ann_x, ann_y, p2_w, Inches(0.35),
             fill=WHITE, line_color=GREY, line_width=0.3)
    add_text(s, ann_x + Inches(0.04), ann_y + Inches(0.02),
             p2_w - Inches(0.08), Inches(0.15),
             f"Wilcoxon signed-rank",
             size=6, color=INK, align="center", anchor="top")
    add_text(s, ann_x + Inches(0.04), ann_y + Inches(0.17),
             p2_w - Inches(0.08), Inches(0.15),
             f"one-sided P = {w_p:.3f}  (n=53)",
             size=7, bold=True, color=GOOD, align="center", anchor="top")

    # x label (below)
    add_text(s, a2_x, _i(p2_y + p2_h + Inches(0.22)),
             a2_w, Inches(0.18),
             "good − bad majority fraction\n(per V-gene)",
             size=6, color=INK, align="center")

    # ====== Sub-panel A3: pattern horizontal stacked bar ======
    a3_x = Inches(4.55); a3_y = Inches(0.50)
    a3_w = Inches(1.85); a3_h = Inches(2.60)

    # ordered pattern counts
    pattern_order = ["good_coherent_bad_mixed",
                     "bad_coherent_good_mixed",
                     "both_coherent_same",
                     "both_coherent_opposite",
                     "both_mixed"]
    counts = {p: int((stats_df.pattern == p).sum()) for p in pattern_order}

    # horizontal stacked bar --- single row
    bar_x = a3_x + Inches(0.15); bar_y = a3_y + Inches(0.30)
    bar_w = a3_w - Inches(0.30); bar_h = Inches(0.45)
    total = 53
    cx = bar_x
    for p in pattern_order:
        seg_w = bar_w * (counts[p] / total)
        add_rect(s, cx, bar_y, seg_w, bar_h,
                 fill=PATTERN_COLOR[p], line_color=WHITE,
                 line_width=0.5)
        # count label inside segment if big enough
        if counts[p] >= 3:
            add_text(s, cx, bar_y + Inches(0.10),
                     seg_w, Inches(0.25),
                     f"{counts[p]}",
                     size=9, bold=True, color=WHITE, align="center",
                     anchor="middle")
        cx += seg_w
    add_rect(s, bar_x, bar_y, bar_w, bar_h,
             line_color=LINE, line_width=0.5)

    add_text(s, bar_x, bar_y - Inches(0.17),
             bar_w, Inches(0.14),
             f"pattern breakdown, 53 V-genes",
             size=7, color=INK, align="center", bold=True)

    # detailed legend
    leg_y = bar_y + bar_h + Inches(0.15)
    for i, p in enumerate(pattern_order):
        row_y = leg_y + i * Inches(0.24)
        add_rect(s, bar_x, row_y, Inches(0.18), Inches(0.13),
                 fill=PATTERN_COLOR[p], line_color=WHITE, line_width=0.3)
        add_text(s, bar_x + Inches(0.22), row_y - Emu(5000),
                 bar_w - Inches(0.24), Inches(0.16),
                 f"n={counts[p]} · {PATTERN_PRETTY[p]}",
                 size=6, color=INK, align="left")

    # 16:7 contrast callout
    gcbm = counts["good_coherent_bad_mixed"]
    bcgm = counts["bad_coherent_good_mixed"]
    bin_p = float(stats.binom.sf(gcbm - 1, gcbm + bcgm, 0.5))
    call_y = leg_y + 5 * Inches(0.24) + Inches(0.10)
    add_rect(s, bar_x, call_y, bar_w, Inches(0.32),
             fill=WHITE, line_color=HIGHLIGHT, line_width=0.8)
    add_text(s, bar_x + Inches(0.04), call_y + Inches(0.02),
             bar_w - Inches(0.08), Inches(0.13),
             f"{gcbm} good-coherent vs {bcgm} bad-coherent",
             size=6, bold=True, color=INK, align="center", anchor="top")
    add_text(s, bar_x + Inches(0.04), call_y + Inches(0.16),
             bar_w - Inches(0.08), Inches(0.13),
             f"binomial P = {bin_p:.3f}",
             size=7, bold=True, color=HIGHLIGHT,
             align="center", anchor="top")

    # ====== Bottom caption (optional) ======
    # none --- rules say no titles
    # But we do need to indicate the separator between sub-panels;
    # use faint vertical lines
    add_line(s, Inches(2.70), Inches(0.50), Inches(2.70), Inches(3.60),
             LT_GREY, 0.3, dashed=True)
    add_line(s, Inches(4.45), Inches(0.50), Inches(4.45), Inches(3.60),
             LT_GREY, 0.3, dashed=True)

    # Sub-panel letters (A1, A2, A3) --- small italic labels
    for sl_label, sl_x in [("i", Inches(0.62)),
                           ("ii", Inches(2.85)),
                           ("iii", Inches(4.60))]:
        add_text(s, sl_x, Inches(0.48), Inches(0.3), Inches(0.15),
                 sl_label, size=7, color=RGBColor(0x66, 0x66, 0x66),
                 italic=True, align="left")

    # Panel bottom summary strip
    summary_y = Inches(3.85)
    add_text(s, Inches(0.15), summary_y,
             SLIDE_W - Inches(0.3), Inches(0.16),
             "Three independent aggregate tests, all P < 0.05: "
             "Wilcoxon P=0.035 (i–ii), binomial P=0.049 (i, 24/37 "
             "V-genes with good_majority > bad_majority), "
             "pattern P=0.047 (iii, 16 vs 7 contrast).",
             size=7, color=INK, align="center")
    add_text(s, Inches(0.15), summary_y + Inches(0.18),
             SLIDE_W - Inches(0.3), Inches(0.16),
             "→ good responders mount directionally coherent V-gene "
             "repertoire responses to RT; bad responders are stochastic.",
             size=7, bold=True, color=INK, align="center")


def build_B():
    """Focus V-gene 3x3 small-multiple spaghetti."""
    s = new_slide(prs_main)
    draw_panel_letter(s, "B")

    # pick focus V-genes
    user_focus = ["IGHV6-1", "IGHV3-7", "IGHV3-74"]
    # top good_coherent_bad_mixed by Fisher P (excluding the user focus)
    gcbm = stats_df[stats_df.pattern == "good_coherent_bad_mixed"].sort_values(
        "fisher_P_updown").v_gene.tolist()
    extras = [g for g in gcbm if g not in user_focus][:6]
    focus_list = user_focus + extras
    focus_list = focus_list[:9]   # 3x3

    # grid layout (3x3)
    grid_ox = Inches(0.55); grid_oy = Inches(0.45)
    sub_w = Inches(1.90); sub_h = Inches(1.15)
    col_gap = Inches(0.10); row_gap = Inches(0.10)

    for idx, vg in enumerate(focus_list):
        r, c = idx // 3, idx % 3
        px = grid_ox + c * (sub_w + col_gap)
        py = grid_oy + r * (sub_h + row_gap)

        # data for this V-gene
        sub = delta_df[delta_df.v_gene == vg]
        if sub.empty:
            continue
        g_sub = sub[sub.response_bin == "good"]
        b_sub = sub[sub.response_bin == "bad"]
        all_vals = np.concatenate([sub["pre"].values, sub["post"].values])

        y_max = max(1e-4, float(np.max(all_vals)) * 1.08)
        y_min = 0

        # inner axes
        ax_x = px + Inches(0.35); ax_y = py + Inches(0.16)
        ax_w = sub_w - Inches(0.42); ax_h = sub_h - Inches(0.50)

        def tx(v): return _i(ax_x + v * ax_w)
        def ty(v): return _i(ax_y + ax_h - (v - y_min) / (y_max - y_min) * ax_h)

        # spines
        add_line(s, ax_x, ax_y, ax_x, ax_y + ax_h, LINE, 0.5)
        add_line(s, ax_x, ax_y + ax_h, ax_x + ax_w, ax_y + ax_h, LINE, 0.5)

        # x tick labels (pre / post)
        add_text(s, tx(0) - Inches(0.17), _i(ax_y + ax_h + Inches(0.02)),
                 Inches(0.34), Inches(0.13),
                 "pre", size=6, color=INK, align="center")
        add_text(s, tx(1) - Inches(0.18), _i(ax_y + ax_h + Inches(0.02)),
                 Inches(0.36), Inches(0.13),
                 "post", size=6, color=INK, align="center")

        # y tick max / zero
        add_text(s, _i(ax_x - Inches(0.35)), ty(y_max) - Inches(0.07),
                 Inches(0.32), Inches(0.14),
                 f"{y_max:.3f}", size=5, color=INK, align="right")
        add_text(s, _i(ax_x - Inches(0.35)), ty(0) - Inches(0.08),
                 Inches(0.32), Inches(0.14),
                 "0", size=5, color=INK, align="right")

        # zero reference
        add_line(s, ax_x, ty(0), ax_x + ax_w, ty(0), GREY, 0.3, dashed=True)

        # individual slopes (light)
        LIGHT_GOOD = lighten(GOOD_HEX, 0.55)
        LIGHT_BAD = lighten(BAD_HEX, 0.55)
        for _, row_ in sub.iterrows():
            c_col = LIGHT_GOOD if row_.response_bin == "good" else LIGHT_BAD
            add_line(s, tx(0), ty(row_.pre), tx(1), ty(row_.post),
                     c_col, 0.7)

        # markers at endpoints (hollow)
        for _, row_ in sub.iterrows():
            col = GOOD if row_.response_bin == "good" else BAD
            add_circle(s, tx(0), ty(row_.pre), Emu(18000),
                       fill=WHITE, line_color=col, line_width=0.7)
            add_circle(s, tx(1), ty(row_.post), Emu(18000),
                       fill=WHITE, line_color=col, line_width=0.7)

        # group median slopes (bold)
        g_med_pre = float(np.median(g_sub.pre.values))
        g_med_post = float(np.median(g_sub.post.values))
        b_med_pre = float(np.median(b_sub.pre.values))
        b_med_post = float(np.median(b_sub.post.values))
        add_line(s, tx(0), ty(g_med_pre), tx(1), ty(g_med_post),
                 GOOD, 1.8)
        add_diamond(s, tx(0), ty(g_med_pre), Emu(30000),
                    fill=GOOD, line_color=WHITE, line_width=0.9)
        add_diamond(s, tx(1), ty(g_med_post), Emu(30000),
                    fill=GOOD, line_color=WHITE, line_width=0.9)
        add_line(s, tx(0), ty(b_med_pre), tx(1), ty(b_med_post),
                 BAD, 1.8)
        add_diamond(s, tx(0), ty(b_med_pre), Emu(30000),
                    fill=BAD, line_color=WHITE, line_width=0.9)
        add_diamond(s, tx(1), ty(b_med_post), Emu(30000),
                    fill=BAD, line_color=WHITE, line_width=0.9)

        # V-gene label (below sub-plot)
        star = " ★" if vg == "IGHV6-1" else ""
        add_text(s, px, _i(py + sub_h - Inches(0.28)),
                 sub_w, Inches(0.16),
                 f"{vg}{star}",
                 size=8, bold=True,
                 color=HIGHLIGHT if vg == "IGHV6-1" else INK,
                 align="center")

        # sign counts annotation (in-plot, top-left corner)
        stat = stats_df[stats_df.v_gene == vg]
        if not stat.empty:
            r_ = stat.iloc[0]
            txt = (f"g ↑/↓ {int(r_.good_n_up)}/{int(r_.good_n_down)}  "
                   f"b {int(r_.bad_n_up)}/{int(r_.bad_n_down)}")
            add_text(s, _i(ax_x + Inches(0.02)), _i(ax_y + Inches(0.02)),
                     Inches(1.35), Inches(0.14),
                     txt, size=5, color=RGBColor(0x55, 0x55, 0x55),
                     align="left")
            # Fisher P
            add_text(s, _i(ax_x + Inches(0.02)), _i(ax_y + Inches(0.14)),
                     Inches(1.35), Inches(0.13),
                     f"Fisher P={r_.fisher_P_updown:.3f}",
                     size=5, color=RGBColor(0x55, 0x55, 0x55),
                     align="left")

    # Legend (bottom of slide)
    leg_y = Inches(4.15)
    add_line(s, Inches(0.55), leg_y + Inches(0.06),
             Inches(0.85), leg_y + Inches(0.06), GOOD, 1.8)
    add_diamond(s, Inches(0.70), leg_y + Inches(0.06), Emu(25000),
                fill=GOOD, line_color=WHITE, line_width=0.9)
    add_text(s, Inches(0.90), leg_y - Emu(10000), Inches(1.3), Inches(0.16),
             "good median (n=6)", size=7, color=INK, align="left")
    add_line(s, Inches(2.35), leg_y + Inches(0.06),
             Inches(2.65), leg_y + Inches(0.06), BAD, 1.8)
    add_diamond(s, Inches(2.50), leg_y + Inches(0.06), Emu(25000),
                fill=BAD, line_color=WHITE, line_width=0.9)
    add_text(s, Inches(2.70), leg_y - Emu(10000), Inches(1.3), Inches(0.16),
             "bad median (n=6)", size=7, color=INK, align="left")
    add_circle(s, Inches(4.15), leg_y + Inches(0.06), Emu(18000),
               fill=WHITE, line_color=LINE, line_width=0.7)
    add_text(s, Inches(4.22), leg_y - Emu(10000), Inches(2.0), Inches(0.16),
             "individual subject pre/post",
             size=7, color=INK, align="left")

    # Shared axis-title text (small, below)
    add_text(s, Inches(0.15), Inches(4.32),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "y-axis: IGH-repertoire fraction  ·  x-axis: pre → post RT",
             size=6, color=RGBColor(0x55, 0x55, 0x55), align="center")


build_A()
build_B()
deck_main = f"{OUT}/Fig7_IGHV_coherence_native_editable.pptx"
prs_main.save(deck_main)
print(f"wrote {deck_main}")


# ===========================================================================
# SUPP FIGURE S18
# ===========================================================================
prs_supp = Presentation()


def build_S18A():
    """All 53 V-genes forest: sign counts with Fisher P."""
    s = new_slide(prs_supp)
    draw_panel_letter(s, "A")

    # order by coherence_gap descending
    df = stats_df.sort_values("coherence_gap", ascending=False).reset_index(
        drop=True)
    n = len(df)

    # plot area --- narrow vertical forest
    px = Inches(2.40); py = Inches(0.40)
    pw = Inches(3.60); ph = Inches(3.85)
    row_h = ph / n

    # center-zero axis; bars extend left (bad) and right (good) up to n=6
    def tx(v): return _i(px + pw / 2 + v / 6 * (pw / 2))

    # spines
    add_line(s, px + pw / 2, py, px + pw / 2, py + ph, LINE, 0.6)
    add_line(s, px, py + ph, px + pw, py + ph, LINE, 0.5)
    for v in [-6, -4, -2, 0, 2, 4, 6]:
        xx = tx(v)
        add_line(s, xx, _i(py + ph), xx, _i(py + ph + Inches(0.03)),
                 LINE, 0.4)
        add_text(s, xx - Inches(0.15), _i(py + ph + Inches(0.04)),
                 Inches(0.30), Inches(0.13),
                 f"{v:+d}" if v != 0 else "0",
                 size=6, color=INK, align="center")
    add_text(s, px, _i(py + ph + Inches(0.18)), pw, Inches(0.14),
             "← bad responders (predicted direction)    good responders → ",
             size=7, color=INK, align="center")

    # row labels + bars
    focus = {"IGHV6-1", "IGHV3-7", "IGHV3-74"}
    for i, row in df.iterrows():
        yy = _i(py + (i + 0.5) * row_h)
        bar_h = _i(row_h * 0.55)

        # predicted count = majority direction for each group
        g_pred = max(int(row.good_n_up), int(row.good_n_down))
        b_pred = max(int(row.bad_n_up), int(row.bad_n_down))

        # right-going (good) bar
        add_rect(s, tx(0), yy - bar_h // 2,
                 tx(g_pred) - tx(0), bar_h,
                 fill=GOOD, line_color=WHITE, line_width=0.3)
        # left-going (bad) bar
        add_rect(s, tx(-b_pred), yy - bar_h // 2,
                 tx(0) - tx(-b_pred), bar_h,
                 fill=BAD, line_color=WHITE, line_width=0.3)

        # row label (V-gene name)
        label_color = HIGHLIGHT if row.v_gene in focus else INK
        star = "★ " if row.v_gene in focus else ""
        add_text(s, Inches(0.20), yy - _i(Inches(0.065)),
                 Inches(1.15), Inches(0.12),
                 f"{star}{row.v_gene}",
                 size=5, bold=(row.v_gene in focus),
                 color=label_color, align="right")

        # pattern chip (small rectangle right of V-gene name)
        add_rect(s, Inches(1.40), yy - _i(Inches(0.045)),
                 Inches(0.12), Inches(0.09),
                 fill=PATTERN_COLOR[row.pattern])

        # counts next to their bars
        add_text(s, tx(g_pred) + Inches(0.02), yy - _i(Inches(0.065)),
                 Inches(0.28), Inches(0.12),
                 f"{g_pred}/6",
                 size=5, color=GOOD, bold=True, align="left")
        add_text(s, tx(-b_pred) - Inches(0.28), yy - _i(Inches(0.065)),
                 Inches(0.26), Inches(0.12),
                 f"{b_pred}/6",
                 size=5, color=BAD, bold=True, align="right")

        # Fisher P on the far right
        p_color = HIGHLIGHT if row.fisher_P_updown <= 0.10 else RGBColor(
            0x77, 0x77, 0x77)
        add_text(s, _i(px + pw + Inches(0.05)),
                 yy - _i(Inches(0.065)),
                 Inches(0.50), Inches(0.12),
                 f"P={row.fisher_P_updown:.2f}",
                 size=5, color=p_color, align="left",
                 bold=row.fisher_P_updown <= 0.10)

    # pattern legend (bottom-left)
    leg_y = Inches(4.15)
    add_text(s, Inches(1.40), leg_y - Inches(0.04),
             Inches(0.90), Inches(0.12),
             "pattern chip:",
             size=5, bold=True, color=INK, align="left")
    for i, p in enumerate(["good_coherent_bad_mixed",
                           "bad_coherent_good_mixed",
                           "both_coherent_same",
                           "both_coherent_opposite",
                           "both_mixed"]):
        row_y = leg_y + Inches(0.10) + i * Inches(0.09)
        add_rect(s, Inches(1.40), row_y, Inches(0.10), Inches(0.07),
                 fill=PATTERN_COLOR[p])
        add_text(s, Inches(1.54), row_y - Inches(0.01),
                 Inches(1.1), Inches(0.09),
                 PATTERN_PRETTY[p],
                 size=4, color=INK, align="left")

    # caption
    add_text(s, Inches(0.15), Inches(4.30),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "53 coverage-filtered V-genes, ordered by coherence_gap "
             "(good_majority − bad_majority) descending. "
             "P highlighted if Fisher ≤ 0.10.",
             size=5, color=RGBColor(0x55, 0x55, 0x55), align="center")


def build_S18B():
    """Pattern-class map scatter with all V-genes labelled."""
    s = new_slide(prs_supp)
    draw_panel_letter(s, "B")

    px = Inches(0.80); py = Inches(0.55)
    pw = Inches(4.20); ph = Inches(3.50)

    lo, hi = 0.48, 1.03
    def xm(v): return _i(px + (v - lo) / (hi - lo) * pw)
    def ym(v): return _i(py + ph - (v - lo) / (hi - lo) * ph)

    # spines + ticks
    add_line(s, px, py, px, py + ph, LINE, 0.6)
    add_line(s, px, py + ph, px + pw, py + ph, LINE, 0.6)
    for v in [0.5, 0.6, 0.7, 0.8, 0.9, 1.0]:
        add_line(s, xm(v), _i(py + ph), xm(v), _i(py + ph + Inches(0.04)),
                 LINE, 0.4)
        add_text(s, xm(v) - Inches(0.15), _i(py + ph + Inches(0.05)),
                 Inches(0.30), Inches(0.13),
                 f"{v:.1f}", size=6, color=INK, align="center")
        add_line(s, _i(px - Inches(0.04)), ym(v), px, ym(v), LINE, 0.4)
        add_text(s, _i(px - Inches(0.33)), ym(v) - Inches(0.07),
                 Inches(0.30), Inches(0.13),
                 f"{v:.1f}", size=6, color=INK, align="right")

    add_text(s, px, _i(py + ph + Inches(0.22)), pw, Inches(0.16),
             "bad-group majority fraction",
             size=8, color=INK, align="center")
    yt = add_text(s, _i(px - Inches(0.50)),
                  _i(py + ph / 2 - Inches(0.9)),
                  Inches(0.40), Inches(1.8),
                  "good-group majority fraction",
                  size=8, color=INK, align="center")
    yt.rotation = -90

    # y=x diagonal
    add_line(s, xm(lo + 0.02), ym(lo + 0.02),
             xm(hi - 0.01), ym(hi - 0.01),
             GREY, 0.5, dashed=True)

    # plot all 53 V-genes with labels
    rng = np.random.default_rng(1)
    for _, row in stats_df.iterrows():
        px_pt = xm(row.bad_majority_frac)
        py_pt = ym(row.good_majority_frac)
        # small jitter for label readability
        jx = px_pt + rng.integers(-_i(Inches(0.025)),
                                   _i(Inches(0.025)) + 1)
        jy = py_pt + rng.integers(-_i(Inches(0.025)),
                                   _i(Inches(0.025)) + 1)
        color = PATTERN_COLOR[row.pattern]
        r = Emu(18000)
        focus = row.v_gene in {"IGHV6-1", "IGHV3-7", "IGHV3-74"}
        if focus:
            r = Emu(28000)
            add_circle(s, jx, jy, r, fill=color,
                       line_color=HIGHLIGHT, line_width=1.2)
        else:
            add_circle(s, jx, jy, r, fill=color,
                       line_color=WHITE, line_width=0.3)
        # label
        label_color = HIGHLIGHT if focus else RGBColor(0x44, 0x44, 0x44)
        add_text(s, jx + Inches(0.04), jy - Inches(0.06),
                 Inches(0.65), Inches(0.12),
                 row.v_gene,
                 size=4 if not focus else 5,
                 bold=focus,
                 color=label_color, align="left")

    # pattern legend
    leg_x = px + pw + Inches(0.15); leg_y = py + Inches(0.1)
    add_text(s, leg_x, leg_y - Inches(0.16),
             Inches(1.30), Inches(0.14),
             "Pattern",
             size=7, bold=True, color=INK, align="left")
    for i, p in enumerate(["good_coherent_bad_mixed",
                           "bad_coherent_good_mixed",
                           "both_coherent_same",
                           "both_coherent_opposite",
                           "both_mixed"]):
        row_y = leg_y + i * Inches(0.24)
        add_circle(s, leg_x + Inches(0.08), row_y + Inches(0.06),
                   Emu(18000),
                   fill=PATTERN_COLOR[p], line_color=WHITE, line_width=0.3)
        add_text(s, leg_x + Inches(0.18), row_y - Inches(0.02),
                 Inches(1.1), Inches(0.22),
                 PATTERN_PRETTY[p].replace(", ", ",\n"),
                 size=5, color=INK, align="left", anchor="top")

    # caption
    add_text(s, Inches(0.15), Inches(4.30),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "Each point = one of 53 coverage-filtered V-genes. Points "
             "above the y=x diagonal = more coherent in good responders. "
             "Highlighted V-genes appear in Main Fig 7 B.",
             size=6, color=RGBColor(0x55, 0x55, 0x55), align="center")


build_S18A()
build_S18B()
deck_supp = f"{OUT}/SuppFig_S18_IGHV_coherence_full_native_editable.pptx"
prs_supp.save(deck_supp)
print(f"wrote {deck_supp}")
