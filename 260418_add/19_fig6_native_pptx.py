#!/usr/bin/env python3
"""
19_fig6_native_pptx.py

Build Figure 6 (§3.9 Target engagement of four Thread-1 baseline factors
in paired radiation-phase biopsies) and its supplementary companion
Supp Fig S17 (member-pathway detail + per-subject Δ heatmap) as
native-editable PowerPoint decks.

Rules (same as script 18 / Fig 5):
    * one panel per slide
    * every in-plot text element is a add_textbox TEXT_BOX (editable,
      not grouped)
    * every line / tick / marker / bar / ribbon / cell is a native
      connector or auto-shape
    * Arial throughout, no panel titles inside the plot
    * TNT project palette: GOOD=#0a7d6e (deep teal) / BAD=#c53e1f
      (deep coral)
    * all shapes get empty <a:effectLst/> via kill_shadow to override
      PowerPoint theme's default soft shadow

Main Figure 6 panels
--------------------
  A  2x2 grid of paired pre->post composite-z spaghetti plots for the
     four Thread-1 factors (DSB/HDR repair, Tumor cell-cycle, E2F/MYC,
     EMT); 12 subjects per sub-plot (6 good + 6 bad); annotation box
     per sub-plot with within-group sign counts and binomial P.
  B  Directional-concordance grouped bar chart: for each of the four
     factors, the fraction of subjects moving in the biologically
     predicted direction, good vs bad, with n_pred/n_total labels and
     a chance-level (0.5) reference line. EMT good 6/6 (P=0.016) is
     highlighted because it is the only composite-level nominal
     single-factor significance.

Supp Figure S17 panels
----------------------
  A  Member-pathway sign-count dotted stack: for each of the 4 factors
     and each of its member pathways, show how many of the 6 good
     responders and how many of the 6 bad responders moved in the
     predicted direction.
  B  Per-subject Δ heatmap across all 17 signatures (4 composites + 13
     member pathways) × 12 paired subjects grouped good | bad,
     diverging colour scale with composite rows annotated.

Style references consulted (paired-biopsy pharmacodynamic convention)
---------------------------------------------------------------------
  - Tumeh et al Nature 2014 (PD-1 pre/post CD8 density): hollow-circle
    paired points connected by thin line, group-coloured.
  - Ribas lab (Tumeh / Riaz) Cell 2017 (pre/on-treatment ICI): small
    per-subject spaghetti arranged in sub-panel grids by marker.
  - Cercek et al NEJM 2022 (dostarlimab rectal): per-subject response
    waterfall + paired lesion trajectories.
  - Petitprez et al Nature 2020 / Helmink et al Nature 2020 / Cabrita
    et al Nature 2020 (B-cell TLS trilogy): directional fraction bars
    with n/N annotations.
  - Thorsson et al Immunity 2018 (TCGA immune landscape): grouped bar
    charts with chance-reference lines and small-N transparency.
"""

import os
import numpy as np
import pandas as pd
from scipy import stats
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Shared infrastructure --- same as script 18
# ---------------------------------------------------------------------------
def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


GOOD = RGBColor(0x0A, 0x7D, 0x6E)
BAD = RGBColor(0xC5, 0x3E, 0x1F)
INK = RGBColor(0x22, 0x22, 0x22)
LINE = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0xBB, 0xBB, 0xBB)
LT_GREY = RGBColor(0xDD, 0xDD, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xD4, 0xA3, 0x00)   # gold for EMT 6/6 callout
GOOD_HEX = (0x0A, 0x7D, 0x6E)
BAD_HEX = (0xC5, 0x3E, 0x1F)

FONT = "Arial"
SLIDE_W = Inches(6.5)
SLIDE_H = Inches(4.5)

OUT = "/data/data/TNT/analysis/260418_add/ppt"
DATA = "/data/data/TNT/analysis/260418_add"
os.makedirs(OUT, exist_ok=True)


def new_slide(prs, w=SLIDE_W, h=SLIDE_H):
    prs.slide_width = w
    prs.slide_height = h
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, size=8, bold=False,
             color=INK, align="left", anchor="middle", font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "bottom": MSO_ANCHOR.BOTTOM,
                          "middle": MSO_ANCHOR.MIDDLE}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = str(text)
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.color.rgb = color
    kill_shadow(tb)
    return tb


def _i(v):
    """Coerce to int --- PowerPoint rejects XML width='0.0' type strings,
    so every coordinate / size passed to a shape creator must be an
    integer number of EMU."""
    return int(round(float(v)))


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.5, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        try:
            c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass
    kill_shadow(c)
    return c


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0.5):
    # sizes must be strictly positive for PowerPoint to render
    w = max(_i(w), 1)
    h = max(_i(h), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _i(x), _i(y), w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_circle(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_diamond(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def lighten(rgb_hex, factor=0.72):
    """Blend an RGB triple towards white by 'factor' (0=identity, 1=white)."""
    r0, g0, b0 = rgb_hex
    return RGBColor(int(r0 + (255 - r0) * factor),
                    int(g0 + (255 - g0) * factor),
                    int(b0 + (255 - b0) * factor))


def draw_panel_letter(slide, letter):
    add_text(slide, Inches(0.15), Inches(0.1), Inches(0.4), Inches(0.35),
             letter, size=14, bold=True, color=INK, align="left")


# ---------------------------------------------------------------------------
# Load paired baseline-factor data
# ---------------------------------------------------------------------------
delta_df = pd.read_csv(f"{DATA}/baseline_factor_per_subject_delta.tsv",
                       sep="\t")
sign_df = pd.read_csv(f"{DATA}/baseline_factor_sign_table.tsv", sep="\t")
stats_df = pd.read_csv(f"{DATA}/baseline_factor_pharmacodynamics_stats.tsv",
                       sep="\t")

FACTORS = [
    ("DSB_HDR_repair", "DSB / HDR repair", "down"),
    ("Tumor_cellcycle", "Tumor cell-cycle", "down"),
    ("E2F_MYC_cellcycle", "E2F / MYC", "down"),
    ("EMT", "EMT", "up"),
]

# subset to composite-level rows for Panel A
comp = delta_df[delta_df.member == "composite"].copy()


# ===========================================================================
# MAIN FIGURE 6
# ===========================================================================
prs_main = Presentation()


def build_A():
    """Oriented-Δ distribution plot: single axis, Δ × sign(predicted
    direction), 4 factors × 2 groups side-by-side with IQR box +
    median diamond + jittered hollow points. Shaded predicted-direction
    zone above zero makes the target-engagement message visually
    immediate; good/bad strips sitting at overlapping heights make the
    MW-NS message visually immediate; EMT good gold star marks the
    only composite-level nominal significance; top-right badge reports
    the 40/48 aggregate."""
    s = new_slide(prs_main)
    draw_panel_letter(s, "A")

    # ---- plot area
    px = Inches(0.95); py = Inches(0.55)
    pw = Inches(4.95); ph = Inches(3.00)

    # ---- data preparation: oriented Δ per subject per factor
    #   oriented_Δ = Δ × (+1 if pred=up, -1 if pred=down)
    #   i.e. positive = moved in predicted direction
    records = []
    for fname, pretty, pred in FACTORS:
        sign_mul = -1 if pred == "down" else +1
        sub = comp[comp.factor == fname].copy()
        sub["oriented_delta"] = sub["delta"] * sign_mul
        sub["factor"] = fname
        records.append(sub)
    oriented = pd.concat(records, ignore_index=True)

    y_values = oriented["oriented_delta"].values
    y_abs_max = max(3.5, float(np.ceil(np.max(np.abs(y_values)) * 1.05)))
    y_lo, y_hi = -y_abs_max, y_abs_max

    def ty(v):
        return int(py + ph - (v - y_lo) / (y_hi - y_lo) * ph)

    # ---- x positioning: 4 factor centers, with good / bad strip offset
    n_fac = len(FACTORS)
    group_w = pw / n_fac
    strip_off = Inches(0.28)
    center_x = [px + (i + 0.5) * group_w for i in range(n_fac)]
    good_x = [c - strip_off for c in center_x]
    bad_x = [c + strip_off for c in center_x]

    # =================================================================
    # 1. Shaded "predicted direction" zone (above zero)
    # =================================================================
    SHADE_PRED = RGBColor(0xEE, 0xF6, 0xF3)   # very faint teal tint
    SHADE_OPP = RGBColor(0xFB, 0xF1, 0xEE)    # very faint coral tint
    add_rect(s, px, py, pw, ty(0) - py, fill=SHADE_PRED)
    add_rect(s, px, ty(0), pw, py + ph - ty(0), fill=SHADE_OPP)

    # zone labels (discreet, italic, grey)
    add_text(s, px + Inches(0.08), py + Inches(0.03),
             Inches(3.5), Inches(0.18),
             "↑ predicted direction  (target engaged)",
             size=7, color=RGBColor(0x3A, 0x7A, 0x6B), align="left")
    add_text(s, px + Inches(0.08), py + ph - Inches(0.20),
             Inches(3.5), Inches(0.18),
             "↓ opposite direction",
             size=7, color=RGBColor(0x9B, 0x5A, 0x48), align="left")

    # =================================================================
    # 2. Axes (spines + ticks)
    # =================================================================
    add_line(s, px, py, px, py + ph, LINE, 0.6)                  # y spine
    add_line(s, px, py + ph, px + pw, py + ph, LINE, 0.6)        # x spine
    # y ticks at round values inside range
    tick_step = 2 if y_abs_max > 3 else 1
    yticks = list(range(-int(y_abs_max // tick_step * tick_step),
                        int(y_abs_max // tick_step * tick_step) + 1,
                        tick_step))
    for yv in yticks:
        yy = ty(yv)
        add_line(s, px - Inches(0.05), yy, px, yy, LINE, 0.5)
        add_text(s, px - Inches(0.5), yy - Inches(0.08),
                 Inches(0.45), Inches(0.16),
                 f"{yv:+d}" if yv != 0 else "0",
                 size=7, color=INK, align="right")

    # y-axis title (rotated)
    yt = add_text(s, Inches(0.10), py + ph / 2 - Inches(0.95),
                  Inches(0.40), Inches(1.9),
                  "Oriented Δ (post − pre) · sign(predicted)",
                  size=8, color=INK, align="center")
    yt.rotation = -90

    # =================================================================
    # 3. Zero reference line (bold)
    # =================================================================
    add_line(s, px, ty(0), px + pw, ty(0), INK, 1.1)
    add_text(s, px + pw + Inches(0.02), ty(0) - Inches(0.08),
             Inches(0.55), Inches(0.16),
             "no change",
             size=6, color=INK, align="left")

    # =================================================================
    # 4. Per factor: good + bad strip with IQR box, median, jitter
    # =================================================================
    rng = np.random.default_rng(42)
    LIGHT_GOOD = lighten(GOOD_HEX, 0.65)
    LIGHT_BAD = lighten(BAD_HEX, 0.65)

    for i, (fname, pretty, pred) in enumerate(FACTORS):
        fac_sub = oriented[oriented.factor == fname]
        # factor separator line (vertical faint) between factor groups
        if i > 0:
            sep_x = int(px + i * group_w)
            add_line(s, sep_x, py + Inches(0.12), sep_x,
                     py + ph - Inches(0.12),
                     RGBColor(0xDD, 0xDD, 0xDD), 0.4, dashed=True)

        for grp_name, strip_cx, fill_col, edge_col in [
            ("good", good_x[i], LIGHT_GOOD, GOOD),
            ("bad", bad_x[i], LIGHT_BAD, BAD),
        ]:
            grp = fac_sub[fac_sub.response_bin == grp_name]
            vals = grp["oriented_delta"].values
            q25, q75 = np.percentile(vals, [25, 75])
            med = float(np.median(vals))

            box_w = Inches(0.22)
            half_w = box_w / 2

            # IQR box
            add_rect(s, strip_cx - half_w, ty(q75),
                     box_w, max(ty(q25) - ty(q75), 2),
                     fill=fill_col, line_color=edge_col, line_width=0.8)
            # median bar across the box (thick)
            add_line(s, strip_cx - half_w, ty(med),
                     strip_cx + half_w, ty(med),
                     edge_col, 2.0)
            # filled diamond marker centred on the median (extra emphasis)
            add_diamond(s, strip_cx, ty(med), Emu(40000),
                        fill=edge_col, line_color=WHITE, line_width=1.0)

            # jittered individual points
            jitter_range = int(Inches(0.14))
            for v in vals:
                jx = strip_cx + rng.integers(-jitter_range // 2,
                                             jitter_range // 2 + 1)
                add_circle(s, jx, ty(v), Emu(22000),
                           fill=WHITE, line_color=edge_col, line_width=0.9)

            # below-axis annotations (under each strip)
            sign_row = sign_df[(sign_df.factor == fname)
                               & (sign_df.group == grp_name)].iloc[0]
            n_pred = int(sign_row.n_predicted); n_tot = int(sign_row.n_total)
            p_val = float(sign_row.sign_binomial_one_sided_P)
            add_text(s, strip_cx - Inches(0.28),
                     py + ph + Inches(0.03),
                     Inches(0.56), Inches(0.14),
                     f"{n_pred}/{n_tot}",
                     size=7, bold=True, color=edge_col, align="center")
            add_text(s, strip_cx - Inches(0.35),
                     py + ph + Inches(0.17),
                     Inches(0.70), Inches(0.13),
                     f"P={p_val:.3f}",
                     size=6, color=edge_col, align="center")

            # gold star on EMT good 6/6 (the only composite-level
            # nominal significance)
            if fname == "EMT" and grp_name == "good":
                add_text(s, strip_cx + Inches(0.13),
                         ty(med) - Inches(0.11),
                         Inches(0.22), Inches(0.22),
                         "★", size=14, bold=True,
                         color=HIGHLIGHT, align="center")

        # factor name (below x-axis, centred on factor group)
        arrow = "↓" if pred == "down" else "↑"
        add_text(s, center_x[i] - group_w / 2 + Inches(0.05),
                 py + ph + Inches(0.33),
                 group_w - Inches(0.1), Inches(0.20),
                 f"{pretty} {arrow}",
                 size=8, bold=True, color=INK, align="center")

        # MW P above each factor group (between good and bad)
        try:
            mw_p = float(stats_df[(stats_df.factor == fname)
                                  & (stats_df.level == "composite")]
                         .iloc[0]["mw_p"])
        except Exception:
            mw_p = float("nan")
        add_text(s, center_x[i] - Inches(0.45),
                 py - Inches(0.02),
                 Inches(0.9), Inches(0.16),
                 f"MW P = {mw_p:.2f}",
                 size=6, color=RGBColor(0x66, 0x66, 0x66),
                 align="center")

    # =================================================================
    # 5. Aggregate badge (top-right, next to plot)
    # =================================================================
    # compute 40/48 from the oriented data
    n_pred_total = int((oriented.oriented_delta > 0).sum())
    n_total = int(len(oriented))
    # binomial P(>= n_pred_total out of n_total under p = 0.5)
    agg_p = float(stats.binom.sf(n_pred_total - 1, n_total, 0.5))
    agg_p_txt = (f"P < 10⁻⁵" if agg_p < 1e-5
                 else f"P = {agg_p:.2e}")

    agg_x = px + pw - Inches(1.90)
    agg_y = py - Inches(0.02)
    agg_w = Inches(1.85); agg_h = Inches(0.26)
    # shifted up out of the way; keeping it subtle so it doesn't steal
    # attention from the plot itself
    agg_x = Inches(4.45); agg_y = Inches(0.12)
    agg_w = Inches(1.95); agg_h = Inches(0.33)
    add_rect(s, agg_x, agg_y, agg_w, agg_h,
             fill=WHITE, line_color=HIGHLIGHT, line_width=0.8)
    add_text(s, agg_x + Inches(0.04), agg_y + Inches(0.02),
             agg_w - Inches(0.08), Inches(0.14),
             f"{n_pred_total} / {n_total} moved in predicted direction",
             size=6, bold=True, color=INK, align="center", anchor="top")
    add_text(s, agg_x + Inches(0.04), agg_y + Inches(0.16),
             agg_w - Inches(0.08), Inches(0.14),
             f"({100 * n_pred_total / n_total:.0f} %, binomial {agg_p_txt})",
             size=6, color=HIGHLIGHT, align="center", anchor="top")

    # =================================================================
    # 6. Legend row (bottom)
    # =================================================================
    leg_y = Inches(4.20)
    # good swatch: light teal box + dark teal edge + hollow circle
    add_rect(s, Inches(0.95), leg_y, Inches(0.22), Inches(0.13),
             fill=LIGHT_GOOD, line_color=GOOD, line_width=0.8)
    add_circle(s, Inches(1.06), leg_y + Inches(0.065),
               Emu(18000), fill=WHITE, line_color=GOOD, line_width=0.8)
    add_text(s, Inches(1.22), leg_y - Emu(10000),
             Inches(0.95), Inches(0.16),
             "good (n=6)",
             size=7, color=INK, align="left")
    # bad swatch
    add_rect(s, Inches(2.25), leg_y, Inches(0.22), Inches(0.13),
             fill=LIGHT_BAD, line_color=BAD, line_width=0.8)
    add_circle(s, Inches(2.36), leg_y + Inches(0.065),
               Emu(18000), fill=WHITE, line_color=BAD, line_width=0.8)
    add_text(s, Inches(2.52), leg_y - Emu(10000),
             Inches(0.95), Inches(0.16),
             "bad (n=6)",
             size=7, color=INK, align="left")
    # star legend
    add_text(s, Inches(3.60), leg_y - Emu(10000),
             Inches(0.25), Inches(0.18),
             "★", size=11, bold=True, color=HIGHLIGHT, align="left")
    add_text(s, Inches(3.80), leg_y - Emu(10000),
             Inches(2.6), Inches(0.16),
             "binomial P < 0.05 (EMT good 6/6, P = 0.016)",
             size=7, color=INK, align="left")
    # box/median/diamond legend
    add_text(s, Inches(0.95), leg_y + Inches(0.18),
             Inches(5.3), Inches(0.14),
             "Box = IQR (25th–75th percentile)  │  "
             "filled diamond = group median  │  "
             "hollow circle = individual paired subject (n=12, 6+6)",
             size=6, color=RGBColor(0x55, 0x55, 0x55), align="left")


def build_B():
    """Directional-concordance grouped bar chart."""
    s = new_slide(prs_main)
    draw_panel_letter(s, "B")

    # plot area
    px = Inches(0.95); py = Inches(0.55)
    pw = Inches(4.8); ph = Inches(3.10)

    def tx(i): return int(px + (i + 0.5) / len(FACTORS) * pw)

    y_lo, y_hi = 0.0, 1.12
    def ty(v):
        return int(py + ph - (v - y_lo) / (y_hi - y_lo) * ph)

    # spines
    add_line(s, px, py, px, py + ph, LINE, 0.6)
    add_line(s, px, py + ph, px + pw, py + ph, LINE, 0.6)

    # y-axis ticks (%) -- 0, 25, 50, 75, 100
    for yv, lab in [(0, "0 %"), (0.25, "25 %"), (0.5, "50 %"),
                    (0.75, "75 %"), (1.0, "100 %")]:
        yy = ty(yv)
        add_line(s, px - Inches(0.05), yy, px, yy, LINE, 0.4)
        add_text(s, px - Inches(0.55), yy - Inches(0.09),
                 Inches(0.5), Inches(0.18),
                 lab, size=7, color=INK, align="right")

    # y-axis title (rotated)
    yt = add_text(s, Inches(0.10), py + ph / 2 - Inches(0.8),
                  Inches(0.35), Inches(1.6),
                  "Fraction of subjects moving in predicted direction",
                  size=8, color=INK, align="center")
    yt.rotation = -90

    # chance line at 0.5
    add_line(s, px + Inches(0.02), ty(0.5), px + pw - Inches(0.02),
             ty(0.5), GREY, 0.5, dashed=True)
    add_text(s, px + pw + Inches(0.02), ty(0.5) - Inches(0.08),
             Inches(0.6), Inches(0.16),
             "chance", size=6, color=INK, align="left")

    # bars --- two per factor (good, bad)
    group_w = pw / len(FACTORS)
    bar_w = int(group_w * 0.28)
    for i, (fname, pretty, pred) in enumerate(FACTORS):
        center = tx(i)
        g = sign_df[(sign_df.factor == fname) & (sign_df.group == "good")].iloc[0]
        b = sign_df[(sign_df.factor == fname) & (sign_df.group == "bad")].iloc[0]

        # good bar (left)
        xg = center - int(bar_w * 1.15)
        top_g = ty(float(g.fraction_predicted))
        add_rect(s, xg, top_g, bar_w, int(ty(0) - top_g),
                 fill=GOOD, line_color=WHITE, line_width=0.3)
        add_text(s, xg - Inches(0.15), top_g - Inches(0.24),
                 Inches(0.6), Inches(0.18),
                 f"{int(g.n_predicted)}/{int(g.n_total)}",
                 size=7, bold=True, color=GOOD, align="center")
        # highlight EMT 6/6 P=0.016 with gold asterisk
        if fname == "EMT":
            add_text(s, xg - Inches(0.15), top_g - Inches(0.40),
                     Inches(0.6), Inches(0.16),
                     f"★ P = {float(g.sign_binomial_one_sided_P):.3f}",
                     size=6, bold=True, color=HIGHLIGHT, align="center")

        # bad bar (right)
        xb = center + int(bar_w * 0.15)
        top_b = ty(float(b.fraction_predicted))
        add_rect(s, xb, top_b, bar_w, int(ty(0) - top_b),
                 fill=BAD, line_color=WHITE, line_width=0.3)
        add_text(s, xb - Inches(0.15), top_b - Inches(0.24),
                 Inches(0.6), Inches(0.18),
                 f"{int(b.n_predicted)}/{int(b.n_total)}",
                 size=7, bold=True, color=BAD, align="center")

        # factor name below
        add_text(s, center - group_w / 2 + Inches(0.08),
                 py + ph + Inches(0.12),
                 group_w - Inches(0.15), Inches(0.22),
                 pretty, size=7, color=INK, align="center")
        add_text(s, center - group_w / 2 + Inches(0.08),
                 py + ph + Inches(0.30),
                 group_w - Inches(0.15), Inches(0.16),
                 f"(pred: {pred})", size=6, color=RGBColor(0x66, 0x66, 0x66),
                 align="center")

    # x-axis title
    add_text(s, px, py + ph + Inches(0.55), pw, Inches(0.22),
             "Thread-1 baseline factor",
             size=8, color=INK, align="center")

    # legend top-right
    leg_x = px + Inches(3.6); leg_y = py + Inches(0.08)
    add_rect(s, leg_x, leg_y, Inches(0.20), Inches(0.13), fill=GOOD)
    add_text(s, leg_x + Inches(0.25), leg_y - Emu(10000),
             Inches(0.9), Inches(0.16),
             "good (n=6)", size=7, color=INK, align="left")
    add_rect(s, leg_x, leg_y + Inches(0.20), Inches(0.20),
             Inches(0.13), fill=BAD)
    add_text(s, leg_x + Inches(0.25), leg_y + Inches(0.19),
             Inches(0.9), Inches(0.16),
             "bad (n=6)", size=7, color=INK, align="left")


build_A()
build_B()
deck_main = f"{OUT}/Fig6_target_engagement_native_editable.pptx"
prs_main.save(deck_main)
print(f"wrote {deck_main}")


# ===========================================================================
# SUPP FIGURE S17 --- member-pathway detail and per-subject Δ heatmap
# ===========================================================================
prs_supp = Presentation()


def build_S17A():
    """Member-level sign-count bar: 17 signature rows (4 composites + 13
    members), for each row show the good-6 and bad-6 predicted-direction
    counts as horizontal bars."""
    s = new_slide(prs_supp)
    draw_panel_letter(s, "A")

    # build row list
    rows = []
    for fname, pretty, pred in FACTORS:
        rows.append((fname, "composite", pretty, pred, True))
        mems = (delta_df[(delta_df.factor == fname)
                         & (delta_df.member != "composite")]
                .member.unique().tolist())
        for m in sorted(mems):
            rows.append((fname, m, m, pred, False))

    n = len(rows)
    px = Inches(2.8); py = Inches(0.40)
    pw = Inches(3.45); ph = Inches(3.75)

    def ty(i):
        return int(py + (i + 0.5) / n * ph)
    def tx(v):
        return int(px + (v + 6) / 12 * pw)   # range -6 to +6

    # spines
    add_line(s, px + pw / 2, py, px + pw / 2, py + ph, LINE, 0.8)  # 0 axis
    add_line(s, px, py + ph, px + pw, py + ph, LINE, 0.6)          # x spine
    # x ticks -6..+6 by 2
    for v in [-6, -4, -2, 0, 2, 4, 6]:
        xx = tx(v)
        add_line(s, xx, int(py + ph), xx,
                 int(py + ph + Inches(0.05)), LINE, 0.4)
        add_text(s, xx - Inches(0.15), int(py + ph + Inches(0.06)),
                 Inches(0.3), Inches(0.14),
                 f"{v:+d}", size=6, color=INK, align="center")
    add_text(s, px, py + ph + Inches(0.22), pw, Inches(0.18),
             "← bad (predicted-direction count)     "
             "good (predicted-direction count) →",
             size=7, color=INK, align="center")

    # composite rows get a bold left label, members indented
    for i, (fname, mname, pretty, pred, is_comp) in enumerate(rows):
        yy = ty(i)
        # row label
        label = pretty if is_comp else (
            "  " + (pretty[:28] + "…" if len(pretty) > 28 else pretty))
        if is_comp:
            label = f"▸ {pretty}"
        add_text(s, Inches(0.25), yy - Inches(0.09),
                 px - Inches(0.30), Inches(0.16),
                 label, size=6 if is_comp else 5,
                 bold=is_comp, color=INK, align="right")

        # compute counts for this row
        if is_comp:
            sub = delta_df[(delta_df.factor == fname)
                           & (delta_df.member == "composite")]
        else:
            sub = delta_df[(delta_df.factor == fname)
                           & (delta_df.member == mname)]
        good_deltas = sub[sub.response_bin == "good"]["delta"].values
        bad_deltas = sub[sub.response_bin == "bad"]["delta"].values

        def count_predicted(arr, pred_dir):
            if pred_dir == "down":
                return int((arr < 0).sum())
            return int((arr > 0).sum())

        g_pred = count_predicted(good_deltas, pred)
        b_pred = count_predicted(bad_deltas, pred)

        # good count as right-going bar
        bar_h = int(ph / n * 0.52)
        add_rect(s, tx(0), yy - bar_h // 2,
                 tx(g_pred) - tx(0), bar_h, fill=GOOD,
                 line_color=WHITE, line_width=0.3)
        add_text(s, tx(g_pred) + Inches(0.04), yy - Inches(0.08),
                 Inches(0.3), Inches(0.14),
                 f"{g_pred}/6", size=5, color=GOOD, bold=True, align="left")
        # bad count as left-going bar
        add_rect(s, tx(-b_pred), yy - bar_h // 2,
                 tx(0) - tx(-b_pred), bar_h, fill=BAD,
                 line_color=WHITE, line_width=0.3)
        add_text(s, tx(-b_pred) - Inches(0.28), yy - Inches(0.08),
                 Inches(0.28), Inches(0.14),
                 f"{b_pred}/6", size=5, color=BAD, bold=True, align="right")

    # legend
    leg_y = Inches(4.20)
    add_rect(s, Inches(1.9), leg_y, Inches(0.16), Inches(0.12), fill=GOOD)
    add_text(s, Inches(2.08), leg_y - Emu(10000), Inches(0.9),
             Inches(0.16), "good  →  predicted direction count",
             size=6, color=INK, align="left")
    add_rect(s, Inches(1.9), leg_y + Inches(0.14), Inches(0.16),
             Inches(0.12), fill=BAD)
    add_text(s, Inches(2.08), leg_y + Inches(0.13), Inches(0.9),
             Inches(0.16), "bad  ←  predicted direction count",
             size=6, color=INK, align="left")

    # caption
    add_text(s, Inches(0.15), Inches(4.35), SLIDE_W - Inches(0.3),
             Inches(0.12),
             "▸ = composite (4 Thread-1 factors); unprefixed rows = "
             "member ssGSEA pathways (13 total)",
             size=5, color=RGBColor(0x66, 0x66, 0x66), align="center")


def build_S17B():
    """Per-subject Δ heatmap for all 17 signatures × 12 subjects grouped
    good | bad."""
    s = new_slide(prs_supp)
    draw_panel_letter(s, "B")

    # row list = 17 signatures (same ordering as S17A)
    rows = []
    for fname, pretty, pred in FACTORS:
        rows.append((fname, "composite", pretty, pred, True))
        mems = (delta_df[(delta_df.factor == fname)
                         & (delta_df.member != "composite")]
                .member.unique().tolist())
        for m in sorted(mems):
            rows.append((fname, m, m, pred, False))

    # subject order: 6 good then 6 bad, each sorted by subject_id
    good_subj = sorted(delta_df[delta_df.response_bin == "good"]
                       .subject_id.unique().tolist())
    bad_subj = sorted(delta_df[delta_df.response_bin == "bad"]
                      .subject_id.unique().tolist())
    subj_order = good_subj + bad_subj

    n_rows = len(rows)
    n_cols = len(subj_order)

    px = Inches(2.2); py = Inches(0.45)
    pw = Inches(4.00); ph = Inches(3.55)
    cell_w = pw / n_cols
    cell_h = ph / n_rows

    # diverging ramp: teal for predicted direction, coral for opposite
    def ramp(v, pred_dir):
        v = float(v)
        # orient so positive intensity = predicted direction
        oriented = -v if pred_dir == "down" else v
        lim = 2.0
        t = max(-1, min(1, oriented / lim))
        if t >= 0:
            r = int(255 - t * (255 - GOOD_HEX[0]))
            g = int(255 - t * (255 - GOOD_HEX[1]))
            b = int(255 - t * (255 - GOOD_HEX[2]))
        else:
            t = -t
            r = int(255 - t * (255 - BAD_HEX[0]))
            g = int(255 - t * (255 - BAD_HEX[1]))
            b = int(255 - t * (255 - BAD_HEX[2]))
        return RGBColor(r, g, b)

    for i, (fname, mname, pretty, pred, is_comp) in enumerate(rows):
        for j, sid in enumerate(subj_order):
            sel = delta_df[(delta_df.factor == fname)
                           & (delta_df.member == ("composite" if is_comp else mname))
                           & (delta_df.subject_id == sid)]
            v = float(sel.delta.values[0]) if not sel.empty else 0.0
            x = int(px + j * cell_w)
            y = int(py + i * cell_h)
            add_rect(s, x, y, int(cell_w), int(cell_h),
                     fill=ramp(v, pred))

    # row labels (truncated)
    for i, (fname, mname, pretty, pred, is_comp) in enumerate(rows):
        label = f"▸ {pretty}" if is_comp else pretty
        if len(label) > 30:
            label = label[:29] + "…"
        yy = int(py + (i + 0.5) * cell_h)
        add_text(s, Inches(0.25), yy - Inches(0.07),
                 px - Inches(0.30), Inches(0.14),
                 label, size=5 if not is_comp else 6,
                 bold=is_comp, color=INK, align="right")

    # column labels (subject_id) and group divider
    for j, sid in enumerate(subj_order):
        xx = int(px + (j + 0.5) * cell_w)
        color = GOOD if sid in good_subj else BAD
        add_text(s, xx - Inches(0.18), py + ph + Inches(0.03),
                 Inches(0.36), Inches(0.14),
                 str(sid), size=5, bold=True, color=color, align="center")

    # vertical divider between good and bad columns
    div_x = int(px + len(good_subj) * cell_w)
    add_line(s, div_x, py, div_x, py + ph, INK, 0.8)

    # group labels above the column headers
    add_text(s, px, py - Inches(0.18), cell_w * len(good_subj),
             Inches(0.14), f"good (n={len(good_subj)})",
             size=7, bold=True, color=GOOD, align="center")
    add_text(s, div_x, py - Inches(0.18), cell_w * len(bad_subj),
             Inches(0.14), f"bad (n={len(bad_subj)})",
             size=7, bold=True, color=BAD, align="center")

    # x-axis title
    add_text(s, px, py + ph + Inches(0.20),
             pw, Inches(0.18),
             "Subject ID",
             size=7, color=INK, align="center")

    # colour bar
    cb_x = px + pw + Inches(0.18)
    cb_y = py + Inches(0.10)
    cb_w = Inches(0.18)
    cb_h = Inches(1.0)
    for k in range(40):
        t = 1 - 2 * k / 39
        # simulate predicted-direction-oriented ramp with pred='down'
        oriented = t
        if oriented >= 0:
            r = int(255 - oriented * (255 - GOOD_HEX[0]))
            g = int(255 - oriented * (255 - GOOD_HEX[1]))
            b = int(255 - oriented * (255 - GOOD_HEX[2]))
        else:
            o = -oriented
            r = int(255 - o * (255 - BAD_HEX[0]))
            g = int(255 - o * (255 - BAD_HEX[1]))
            b = int(255 - o * (255 - BAD_HEX[2]))
        add_rect(s, cb_x, int(cb_y + k * cb_h / 40),
                 cb_w, int(cb_h / 40) + Emu(2000),
                 fill=RGBColor(r, g, b))
    add_rect(s, cb_x, cb_y, cb_w, cb_h, line_color=LINE, line_width=0.4)
    add_text(s, cb_x - Inches(0.30), cb_y - Inches(0.12),
             cb_w + Inches(0.6), Inches(0.14),
             "Δ vs predicted", size=6, color=INK, align="center")
    add_text(s, cb_x + cb_w + Inches(0.02), cb_y - Inches(0.04),
             Inches(0.3), Inches(0.12),
             "pred", size=5, color=GOOD, bold=True, align="left")
    add_text(s, cb_x + cb_w + Inches(0.02), cb_y + cb_h - Inches(0.08),
             Inches(0.3), Inches(0.12),
             "opp.", size=5, color=BAD, bold=True, align="left")
    add_text(s, cb_x + cb_w + Inches(0.02), cb_y + cb_h / 2 - Inches(0.06),
             Inches(0.3), Inches(0.12),
             "0", size=5, color=INK, align="left")

    # caption
    add_text(s, Inches(0.15), Inches(4.30),
             SLIDE_W - Inches(0.3), Inches(0.15),
             "Cells coloured by Δ (post − pre) re-oriented by predicted "
             "direction: teal = moved in predicted direction, coral = opposite.",
             size=5, color=RGBColor(0x66, 0x66, 0x66), align="center")


build_S17A()
build_S17B()
deck_supp = f"{OUT}/SuppFig_S17_target_engagement_members_native_editable.pptx"
prs_supp.save(deck_supp)
print(f"wrote {deck_supp}")
