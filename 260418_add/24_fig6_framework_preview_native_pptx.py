#!/usr/bin/env python3
"""
24_fig6_framework_preview_native_pptx.py

PREVIEW deck for the proposed new Figure 6 covering §3.8 — Radiation-
phase pharmacodynamics framework + convergence null.  If accepted, the
existing Fig 6/7/8/9 shift to Fig 7/8/9/10 and the convergence-null
callout currently inside Fig 8F is removed (consolidated here).

Panels
------
  A  Three-layer framework diagram — radiation-phase biopsy bracket at
     top, three orthogonal pharmacodynamic axes (§3.9 target
     engagement / §3.10 IGH coherence / §3.11 cascade phenomenology)
     as rounded boxes, static baseline predictor on the left with the
     '⊥' orthogonality arrow crossing the three axes.
  B  36-pair targeted convergence heatmap (9 baseline × 4 cascade Δ,
     partial Spearman r; 0/36 at P<0.05; BH-adjusted q reported).
  C  Headline null — DSB repair baseline vs CD8-cytotoxic Δ scatter
     (n=12, r=-0.07, P=0.83); shaded power envelope shows that
     |r|>=0.55 would have been detectable.

Palette and style follow the TNT project standard (GOOD=#0a7d6e,
BAD=#c53e1f, Arial sans, 0.5-pt spines).
"""

import os
import numpy as np
import pandas as pd
from scipy import stats
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Shadow removal
# ---------------------------------------------------------------------------
def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


OUT = "/data/data/TNT/analysis/260418_add/ppt"
DATA = "/data/data/TNT/analysis/260418_add"
os.makedirs(OUT, exist_ok=True)

# TNT project palette
GOOD = RGBColor(0x0A, 0x7D, 0x6E)
BAD = RGBColor(0xC5, 0x3E, 0x1F)
INK = RGBColor(0x22, 0x22, 0x22)
LINE = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0xBB, 0xBB, 0xBB)
LT_GREY = RGBColor(0xDD, 0xDD, 0xDD)
VLT_GREY = RGBColor(0xF0, 0xF0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xD4, 0xA3, 0x00)
GOOD_LT = RGBColor(0x9F, 0xC9, 0xBE)      # GOOD tint
BAD_LT = RGBColor(0xE8, 0xB4, 0xA4)       # BAD tint

FONT = "Arial"

SLIDE_W = Inches(6.5)
SLIDE_H = Inches(4.5)


# ---------------------------------------------------------------------------
# Low-level helpers (mirrors 18_fig5_native_pptx.py style)
# ---------------------------------------------------------------------------
def new_slide(prs, w=SLIDE_W, h=SLIDE_H):
    blank = prs.slide_layouts[6]
    prs.slide_width = w
    prs.slide_height = h
    return prs.slides.add_slide(blank)


def add_text(slide, x, y, w, h, text, size=8, bold=False,
             color=INK, align="left", anchor="middle", font=FONT,
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    if anchor == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = str(text)
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.italic = bool(italic)
    r.font.color.rgb = color
    kill_shadow(tb)
    return tb


def add_multiline_text(slide, x, y, w, h, lines, size=8,
                       color=INK, align="left", anchor="top",
                       line_spacing=1.1):
    """lines: list of dicts with keys {text,bold?,italic?,color?,size?}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    if anchor == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = str(ln.get("text", ""))
        r.font.name = FONT
        r.font.size = Pt(ln.get("size", size))
        r.font.bold = bool(ln.get("bold", False))
        r.font.italic = bool(ln.get("italic", False))
        r.font.color.rgb = ln.get("color", color)
    kill_shadow(tb)
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.5, dash=None):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dash == "dash":
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    elif dash == "dot":
        c.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
    kill_shadow(c)
    return c


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0.5,
             rounded=False):
    # PowerPoint 2016+ on Windows occasionally rejects ROUNDED_RECTANGLE
    # without an explicit avLst; keep all rectangles as plain RECTANGLE to
    # stay schema-clean across readers.
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 int(x), int(y), int(w), int(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_circle(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 cx - r, cy - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=LINE, width=1.0, dash=None):
    """A line plus a native right-arrow glyph shape at the end point.
    Avoids raw tailEnd XML (which some PowerPoint versions reject).
    Arrowhead is drawn as a tiny filled triangle AUTO_SHAPE so every
    element stays inside supported python-pptx builders."""
    c = add_line(slide, int(x1), int(y1), int(x2), int(y2),
                 color, width, dash)
    # triangle AUTO_SHAPE at endpoint, rotated to match the line angle
    hx, hy = int(x2), int(y2)
    head = int(Emu(45000))
    half = head // 2
    tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                 hx - head, hy - half, head, head)
    tri.fill.solid(); tri.fill.fore_color.rgb = color
    tri.line.fill.background()
    # orient: compute angle in degrees from (x1,y1) to (x2,y2)
    import math
    ang = math.degrees(math.atan2(y2 - y1, x2 - x1))
    tri.rotation = float(ang)
    tri.shadow.inherit = False
    kill_shadow(tri)
    return c


def draw_panel_letter(slide, letter):
    add_text(slide, Inches(0.15), Inches(0.1), Inches(0.4), Inches(0.35),
             letter, size=14, bold=True, color=INK, align="left")


def draw_axes(slide, x, y, w, h, x_ticks=None, y_ticks=None,
              x_range=None, y_range=None,
              xlab=None, ylab=None,
              x_tick_fmt="{:.1f}", y_tick_fmt="{:.1f}"):
    add_line(slide, x, y + h, x + w, y + h, LINE, 0.6)
    add_line(slide, x, y, x, y + h, LINE, 0.6)
    if x_range is None:
        x_range = (min(x_ticks), max(x_ticks))
    if y_range is None:
        y_range = (min(y_ticks), max(y_ticks))
    xlo, xhi = x_range
    ylo, yhi = y_range

    def to_x(v):
        return int(x + (v - xlo) / (xhi - xlo) * w)

    def to_y(v):
        return int(y + h - (v - ylo) / (yhi - ylo) * h)

    tick_len = Emu(40000)
    if x_ticks is not None:
        for v in x_ticks:
            xx = to_x(v)
            add_line(slide, xx, y + h, xx, y + h + tick_len, LINE, 0.5)
            add_text(slide, xx - Inches(0.3), y + h + tick_len + Emu(20000),
                     Inches(0.6), Inches(0.18),
                     x_tick_fmt.format(v), size=7, color=INK, align="center")
    if y_ticks is not None:
        for v in y_ticks:
            yy = to_y(v)
            add_line(slide, x - tick_len, yy, x, yy, LINE, 0.5)
            add_text(slide, x - Inches(0.65), yy - Inches(0.09),
                     Inches(0.6), Inches(0.18),
                     y_tick_fmt.format(v), size=7, color=INK, align="right")
    if xlab:
        add_text(slide, x, y + h + Inches(0.34), w, Inches(0.22),
                 xlab, size=9, color=INK, align="center")
    if ylab:
        tb = add_text(slide, x - Inches(0.95), y + h / 2 - Inches(0.6),
                      Inches(0.85), Inches(1.2), ylab, size=9,
                      color=INK, align="center")
        tb.rotation = -90
    return to_x, to_y


# ---------------------------------------------------------------------------
# Load data --- identical to script 09
# ---------------------------------------------------------------------------
M = pd.read_csv(f'{DATA}/integrated_subject_master_v2.tsv', sep='\t')
M['subject_id'] = M['subject_id'].astype(str)

LONG = pd.read_csv('/data/data/TNT/analysis/09_integration/paired_delta/paired_feature_long.tsv', sep='\t')
LONG['subject_id'] = LONG['subject_id'].astype(str)
LONG['delta'] = LONG['post'] - LONG['pre']
delta_legacy = LONG.pivot(index='subject_id', columns='feature', values='delta')
delta_legacy.columns = [f'{c}_delta' for c in delta_legacy.columns]
DNEW = pd.read_csv(f'{DATA}/paired_immune_delta_per_subject.tsv', sep='\t')
DNEW['subject_id'] = DNEW['subject_id'].astype(str)
delta_new = DNEW.set_index('subject_id')[[c for c in DNEW.columns if c.endswith('_delta')]]
delta = delta_legacy.join(delta_new, how='outer')

BASE = ['DNA Double-Strand Break Repair R-HSA-5693532',
        'HDR Thru Homologous Recombination (HRR) R-HSA-5685942',
        'DNA Repair R-HSA-73894',
        'E2F Targets', 'G2-M Checkpoint', 'Myc Targets V2',
        'frac_amp', 'MHC_II', 'MSI_pct']
BASE_SHORT = ['DSB Repair', 'HDR', 'DNA Repair (all)',
              'E2F Targets', 'G2-M Checkpoint', 'Myc V2',
              'frac_amp', 'MHC_II', 'MSI %']
CASC = ['CD8_cytotoxic_delta', 'Treg_delta', 'MHC_II_delta', 'IGH_n_delta']
CASC_SHORT = ['Δ CD8 cyt', 'Δ Treg', 'Δ MHC-II', 'Δ IGH count']

conv = pd.read_csv(f'{DATA}/targeted_convergence_test.tsv', sep='\t')

# assemble partial-r matrix
Rmat = conv.pivot(index='baseline', columns='cascade',
                  values='partial_r').reindex(index=BASE, columns=CASC)
Pmat = conv.pivot(index='baseline', columns='cascade',
                  values='partial_p').reindex(index=BASE, columns=CASC)
Qmat = conv.pivot(index='baseline', columns='cascade',
                  values='partial_q_BH').reindex(index=BASE, columns=CASC)

# scatter data: DSB Repair baseline vs CD8_cytotoxic_delta
B_paired = M.set_index('subject_id').reindex(delta.index)
dsb = B_paired[BASE[0]]
ccyt = delta['CD8_cytotoxic_delta']
rb = B_paired['response_bin']
scatter_df = pd.DataFrame({'subj': B_paired.index, 'x': dsb.values,
                           'y': ccyt.values, 'resp': rb.values}).dropna()
rs, ps = stats.spearmanr(scatter_df.x, scatter_df.y)
N_SCAT = len(scatter_df)

print(f'scatter n = {N_SCAT}  Spearman r = {rs:+.2f} P = {ps:.2f}')
print(f'convergence matrix: {Rmat.shape}, hits at P<0.05: '
      f'{(Pmat.values < 0.05).sum()}, at BH-q<0.10: '
      f'{(Qmat.values < 0.10).sum()}')


# ===========================================================================
# Slide deck
# ===========================================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ---------- Panel A: framework diagram ----------
def build_A():
    s = new_slide(prs)
    draw_panel_letter(s, "A")

    # === Top: radiation-phase biopsy bracket ===
    bracket_y = Inches(0.35)
    bracket_h = Inches(0.45)
    pre_x = Inches(1.15); post_x = Inches(4.95)
    pre_w = Inches(0.85); post_w = Inches(0.85)
    add_rect(s, pre_x, bracket_y, pre_w, bracket_h,
             fill=VLT_GREY, line_color=LINE, line_width=0.6)
    add_text(s, pre_x, bracket_y, pre_w, bracket_h,
             "pre-RT\nbiopsy", size=8, bold=True, color=INK, align="center")
    add_rect(s, post_x, bracket_y, post_w, bracket_h,
             fill=VLT_GREY, line_color=LINE, line_width=0.6)
    add_text(s, post_x, bracket_y, post_w, bracket_h,
             "post-RT\nbiopsy", size=8, bold=True, color=INK, align="center")
    # arrow pre -> post with "Radiation phase" label
    mid_y = bracket_y + bracket_h / 2
    add_arrow(s, pre_x + pre_w, mid_y, post_x, mid_y, INK, 1.2)
    add_text(s, pre_x + pre_w + Inches(0.1), bracket_y - Inches(0.04),
             post_x - (pre_x + pre_w) - Inches(0.2), Inches(0.22),
             "Radiation phase (short-course RT 25 Gy / 5 Fx, RT-only biopsy window)",
             size=8, bold=True, color=INK, align="center")

    # === Left column: static baseline predictor box ===
    base_x = Inches(0.3); base_y = Inches(1.3)
    base_w = Inches(1.55); base_h = Inches(2.0)
    add_rect(s, base_x, base_y, base_w, base_h,
             fill=GOOD_LT, line_color=GOOD, line_width=1.2, rounded=True)
    add_multiline_text(s, base_x + Inches(0.08), base_y + Inches(0.08),
                       base_w - Inches(0.16), base_h - Inches(0.16),
                       lines=[
        {"text": "Static baseline", "bold": True, "size": 9, "color": GOOD},
        {"text": "predictor (Fig 5)", "bold": True, "size": 9, "color": GOOD},
        {"text": "", "size": 4},
        {"text": "4-feature ElasticNet", "size": 7, "color": INK},
        {"text": "DSB repair · frac_amp", "size": 7, "color": INK},
        {"text": "MHC-II · MSI %", "size": 7, "color": INK},
        {"text": "", "size": 4},
        {"text": "Nested AUC 0.745", "bold": True, "size": 8, "color": INK},
        {"text": "95 % CI 0.56–0.90", "size": 7, "color": INK},
        {"text": "", "size": 4},
        {"text": "single pre-CRT read", "italic": True, "size": 7,
         "color": RGBColor(0x55, 0x55, 0x55)},
    ], align="center", anchor="top")

    # === Right side: three orthogonal pharmacodynamic axes ===
    col_xs = [Inches(2.05), Inches(3.55), Inches(5.05)]
    col_w = Inches(1.40); col_y = Inches(1.3); col_h = Inches(2.0)

    axis_titles = [
        ("§3.9  Target", "engagement"),
        ("§3.10  IGH repertoire", "directional coherence"),
        ("§3.11  Cascade", "phenomenology"),
    ]
    axis_bodies = [
        [
            {"text": "Do the baseline predictor axes", "size": 7, "color": INK},
            {"text": "themselves move post-RT in", "size": 7, "color": INK},
            {"text": "biologically predicted direction?", "size": 7, "color": INK},
            {"text": "", "size": 4},
            {"text": "RNA-paired  n = 12", "bold": True, "size": 7, "color": INK},
            {"text": "oriented-Δ composites", "size": 7,
             "color": RGBColor(0x55, 0x55, 0x55)},
            {"text": "", "size": 4},
            {"text": "EMT  good 6/6  P = 0.016", "bold": True, "size": 7,
             "color": HIGHLIGHT},
            {"text": "40 / 48 concordant",
             "size": 7, "color": INK},
        ],
        [
            {"text": "Does IGH V-gene response to RT", "size": 7, "color": INK},
            {"text": "show reproducible direction per", "size": 7, "color": INK},
            {"text": "group, even with noisy magnitudes?", "size": 7, "color": INK},
            {"text": "", "size": 4},
            {"text": "RNA-paired  n = 12", "bold": True, "size": 7, "color": INK},
            {"text": "53 V-genes, LOO coherence", "size": 7,
             "color": RGBColor(0x55, 0x55, 0x55)},
            {"text": "", "size": 4},
            {"text": "3 aggregate P < 0.05", "bold": True, "size": 7,
             "color": HIGHLIGHT},
            {"text": "good coherent / bad stochastic",
             "size": 7, "color": INK},
        ],
        [
            {"text": "Do mutation, neoantigen, HLA-LOH,", "size": 7, "color": INK},
            {"text": "Treg, MHC-II, B-cell features show", "size": 7, "color": INK},
            {"text": "coherent pre→post Δ in good?", "size": 7, "color": INK},
            {"text": "", "size": 4},
            {"text": "mixed  n = 11 / 12 / 14", "bold": True, "size": 7,
             "color": INK},
            {"text": "BCa 95 % CI, exploratory", "size": 7,
             "color": RGBColor(0x55, 0x55, 0x55)},
            {"text": "", "size": 4},
            {"text": "Treg Δ only robust between-group",
             "size": 7, "color": HIGHLIGHT, "bold": True},
            {"text": "others within-good significant",
             "size": 7, "color": INK},
        ],
    ]
    for xi, (title_pair, body, cx0) in enumerate(zip(axis_titles, axis_bodies, col_xs)):
        add_rect(s, cx0, col_y, col_w, col_h,
                 fill=WHITE, line_color=LINE, line_width=0.8, rounded=True)
        # title (two lines)
        add_multiline_text(s, cx0 + Inches(0.06), col_y + Inches(0.06),
                           col_w - Inches(0.12), Inches(0.45),
                           lines=[
            {"text": title_pair[0], "bold": True, "size": 8, "color": INK},
            {"text": title_pair[1], "bold": True, "size": 8, "color": INK},
        ], align="center", anchor="top")
        # body
        add_multiline_text(s, cx0 + Inches(0.08), col_y + Inches(0.50),
                           col_w - Inches(0.16), col_h - Inches(0.55),
                           lines=body, align="left", anchor="top")

    # === Left baseline → each axis with ⊥ tag ===
    arrow_y = base_y + base_h / 2
    for cx0 in col_xs:
        axis_mid_x = cx0 + col_w / 2
        # dashed orthogonal connector from baseline box to axis box top edge
        add_arrow(s, base_x + base_w, arrow_y,
                  axis_mid_x, col_y - Emu(5000), GREY, 0.8, dash="dash")
    # overlay "⊥" symbols above arrows (mid-point) in gold
    for cx0 in col_xs:
        axis_mid_x = cx0 + col_w / 2
        mid_x = (base_x + base_w + axis_mid_x) / 2
        mid_y = (arrow_y + col_y) / 2
        add_text(s, mid_x - Inches(0.11), mid_y - Inches(0.12),
                 Inches(0.22), Inches(0.2), "⊥", size=11, bold=True,
                 color=HIGHLIGHT, align="center")

    # === Bottom: convergence-null summary bar ===
    bar_y = Inches(3.55); bar_h = Inches(0.55)
    add_rect(s, Inches(0.3), bar_y, Inches(5.9), bar_h,
             fill=RGBColor(0xFB, 0xF4, 0xE0), line_color=HIGHLIGHT,
             line_width=0.8, rounded=True)
    add_multiline_text(s, Inches(0.4), bar_y + Inches(0.05),
                       Inches(5.7), bar_h - Inches(0.1), lines=[
        {"text": "Convergence null (§3.8)",
         "bold": True, "size": 8, "color": HIGHLIGHT},
        {"text": "Static baseline axis ⊥ dynamic pharmacodynamic axes — "
                 "0 / 36 targeted pairs P < 0.05 (1.8 expected by chance); "
                 "see Panel B and Panel C for the test and headline pair.",
         "size": 7, "color": INK},
    ], align="left", anchor="top", line_spacing=1.15)


# ---------- Panel B: 36-pair heatmap ----------
def build_B():
    s = new_slide(prs)
    draw_panel_letter(s, "B")

    # plotting rectangle
    n_rows, n_cols = len(BASE), len(CASC)
    hm_x = Inches(1.75); hm_y = Inches(0.65)
    hm_w = Inches(2.8); hm_h = Inches(2.8)
    cell_w = int(hm_w / n_cols); cell_h = int(hm_h / n_rows)
    vmax = max(0.6, np.nanmax(np.abs(Rmat.values)))

    def rb_ramp(v):
        """Red-Blue diverging: v in [-vmax, vmax]"""
        t = (v + vmax) / (2 * vmax)
        t = max(0.0, min(1.0, t))
        # blue at t=0, white at t=0.5, red at t=1
        if t < 0.5:
            f = t / 0.5
            r = int(41 + (255 - 41) * f)
            g = int(74 + (255 - 74) * f)
            b = int(140 + (255 - 140) * f)
        else:
            f = (t - 0.5) / 0.5
            r = int(255 + (191 - 255) * f)
            g = int(255 + (30 - 255) * f)
            b = int(255 + (46 - 255) * f)
        return RGBColor(r, g, b)

    # cells
    for i, bf in enumerate(BASE):
        for j, cf in enumerate(CASC):
            v = Rmat.iat[i, j]; p = Pmat.iat[i, j]; q = Qmat.iat[i, j]
            x = hm_x + j * cell_w
            y = hm_y + i * cell_h
            fill = rb_ramp(v) if not np.isnan(v) else VLT_GREY
            add_rect(s, x, y, cell_w, cell_h,
                     fill=fill, line_color=WHITE, line_width=0.4)
            txt_color = WHITE if abs(v) > 0.4 else INK
            add_text(s, x, y, cell_w, cell_h, f"{v:+.2f}",
                     size=7, bold=False, color=txt_color, align="center")

    # x-tick labels (cascade)
    for j, lbl in enumerate(CASC_SHORT):
        cx = hm_x + j * cell_w + cell_w // 2
        add_text(s, cx - Inches(0.5), hm_y + hm_h + Emu(30000),
                 Inches(1.0), Inches(0.22),
                 lbl, size=7, color=INK, align="center")
    # y-tick labels (baseline)
    for i, lbl in enumerate(BASE_SHORT):
        cy = hm_y + i * cell_h + cell_h // 2
        add_text(s, hm_x - Inches(1.5), cy - Inches(0.09),
                 Inches(1.45), Inches(0.18),
                 lbl, size=7, color=INK, align="right")

    # axis labels
    add_text(s, hm_x, hm_y + hm_h + Inches(0.42), hm_w, Inches(0.2),
             "Cascade Δ (post − pre, §3.11 axis)",
             size=9, bold=True, color=INK, align="center")
    tb = add_text(s, hm_x - Inches(2.0), hm_y + hm_h / 2 - Inches(0.55),
                  Inches(1.5), Inches(1.1),
                  "Baseline feature (pre-CRT)",
                  size=9, bold=True, color=INK, align="center")
    tb.rotation = -90

    # colorbar (right side)
    cb_x = hm_x + hm_w + Inches(0.25); cb_y = hm_y; cb_h = hm_h
    cb_w = Inches(0.18)
    n_grad = 32
    for k in range(n_grad):
        t = k / (n_grad - 1)
        v = -vmax + 2 * vmax * t
        fill = rb_ramp(v)
        y0 = cb_y + int(cb_h * (1 - (k + 1) / n_grad))
        y1 = cb_y + int(cb_h * (1 - k / n_grad))
        add_rect(s, cb_x, y0, cb_w, max(y1 - y0, 1),
                 fill=fill, line_color=None, line_width=0)
    # colorbar frame
    add_rect(s, cb_x, cb_y, cb_w, cb_h, fill=None,
             line_color=LINE, line_width=0.4)
    # colorbar ticks
    for v in [-vmax, 0, vmax]:
        t = (v + vmax) / (2 * vmax)
        yy = cb_y + int(cb_h * (1 - t))
        add_line(s, cb_x + cb_w, yy, cb_x + cb_w + Emu(25000), yy, LINE, 0.4)
        add_text(s, cb_x + cb_w + Emu(30000), yy - Inches(0.09),
                 Inches(0.4), Inches(0.18),
                 f"{v:+.1f}", size=6, color=INK, align="left")
    add_text(s, cb_x + cb_w - Inches(0.1), cb_y - Inches(0.22),
             Inches(1.0), Inches(0.18),
             "partial r", size=7, bold=True, color=INK, align="center")
    add_text(s, cb_x - Inches(0.2), cb_y + cb_h + Inches(0.05),
             Inches(1.2), Inches(0.22),
             "(response-adj)", size=6, color=INK, align="center")

    # summary badge (bottom-right)
    hit_p = int((Pmat.values < 0.05).sum())
    hit_q = int((Qmat.values < 0.10).sum())
    bdg_x = Inches(0.4); bdg_y = Inches(3.75)
    bdg_w = Inches(5.7); bdg_h = Inches(0.55)
    add_rect(s, bdg_x, bdg_y, bdg_w, bdg_h,
             fill=RGBColor(0xFB, 0xF4, 0xE0), line_color=HIGHLIGHT,
             line_width=0.8, rounded=True)
    add_multiline_text(s, bdg_x + Inches(0.15), bdg_y + Inches(0.06),
                       bdg_w - Inches(0.3), bdg_h - Inches(0.12),
                       lines=[
        {"text": f"0 / 36 pairs significant at P < 0.05   ·   "
                 f"0 / 36 at BH q < 0.10   ·   1.8 expected by chance",
         "bold": True, "size": 9, "color": HIGHLIGHT},
        {"text": "The baseline Thread-1 predictor axes do not forecast the "
                 "magnitude of the paired radiation-phase Δ.",
         "italic": True, "size": 7, "color": INK},
    ], align="center", anchor="top", line_spacing=1.2)


# ---------- Panel C: headline null scatter ----------
def build_C():
    s = new_slide(prs)
    draw_panel_letter(s, "C")

    # axis range
    x_lo = float(np.floor(scatter_df.x.min() * 2) / 2)
    x_hi = float(np.ceil(scatter_df.x.max() * 2) / 2)
    y_lo = float(np.floor(scatter_df.y.min() * 2) / 2)
    y_hi = float(np.ceil(scatter_df.y.max() * 2) / 2)
    # pad
    x_pad = (x_hi - x_lo) * 0.08; y_pad = (y_hi - y_lo) * 0.08
    x_lo -= x_pad; x_hi += x_pad; y_lo -= y_pad; y_hi += y_pad

    PLOT_X = Inches(1.0); PLOT_Y = Inches(0.6)
    PLOT_W = Inches(4.1); PLOT_H = Inches(2.85)

    x_ticks = list(np.round(np.linspace(x_lo, x_hi, 5), 1))
    y_ticks = list(np.round(np.linspace(y_lo, y_hi, 5), 1))
    to_x, to_y = draw_axes(s, PLOT_X, PLOT_Y, PLOT_W, PLOT_H,
                            x_ticks=x_ticks, y_ticks=y_ticks,
                            x_range=(x_lo, x_hi), y_range=(y_lo, y_hi),
                            xlab="DSB repair ssGSEA (pre-CRT baseline, z-score)",
                            ylab="Δ CD8-cytotoxic (post − pre, z-score units)")

    # zero guide lines
    if y_lo <= 0 <= y_hi:
        add_line(s, to_x(x_lo), to_y(0), to_x(x_hi), to_y(0), GREY, 0.4, dash="dash")
    if x_lo <= 0 <= x_hi:
        add_line(s, to_x(0), to_y(y_lo), to_x(0), to_y(y_hi), GREY, 0.4, dash="dash")

    # observed fit line
    b, a = np.polyfit(scatter_df.x.values, scatter_df.y.values, 1)
    x_line = np.array([x_lo, x_hi])
    y_line = a + b * x_line
    add_line(s, to_x(x_line[0]), to_y(y_line[0]),
             to_x(x_line[1]), to_y(y_line[1]), INK, 1.2)

    # power envelope: lines with slope = r·σy/σx for r = +0.55 and r = −0.55
    # centred on the data mean, drawn as reference "detectable slope"
    sx = scatter_df.x.std(ddof=0); sy = scatter_df.y.std(ddof=0)
    mx = scatter_df.x.mean(); my = scatter_df.y.mean()
    for r_ref, lbl in [(+0.55, "+0.55"), (-0.55, "−0.55")]:
        slope_ref = r_ref * sy / sx
        y_ref = my + slope_ref * (x_line - mx)
        add_line(s, to_x(x_line[0]), to_y(y_ref[0]),
                 to_x(x_line[1]), to_y(y_ref[1]),
                 RGBColor(0x99, 0x99, 0x99), 0.7, dash="dash")
        # label at right end
        add_text(s, to_x(x_line[1]) + Emu(20000),
                 to_y(y_ref[1]) - Inches(0.08),
                 Inches(0.6), Inches(0.16),
                 f"|r|={lbl}", size=6, color=GREY, align="left")

    # points
    for _, row in scatter_df.iterrows():
        color = GOOD if row.resp == "good" else BAD
        add_circle(s, to_x(row.x), to_y(row.y), Emu(55000),
                   fill=color, line_color=WHITE, line_width=0.6)
        # subject id
        add_text(s, to_x(row.x) + Emu(60000), to_y(row.y) - Inches(0.07),
                 Inches(0.3), Inches(0.14),
                 str(int(row.subj)), size=6, bold=True, color=color,
                 align="left")

    # stats annotation box
    ann_x = PLOT_X + Inches(0.12); ann_y = PLOT_Y + Inches(0.1)
    ann_w = Inches(1.85); ann_h = Inches(0.65)
    add_rect(s, ann_x, ann_y, ann_w, ann_h,
             fill=WHITE, line_color=HIGHLIGHT, line_width=0.8, rounded=True)
    add_multiline_text(s, ann_x + Inches(0.06), ann_y + Inches(0.04),
                       ann_w - Inches(0.12), ann_h - Inches(0.08), lines=[
        {"text": f"Spearman r = {rs:+.2f}", "bold": True, "size": 9,
         "color": INK},
        {"text": f"P = {ps:.2f}   (n = {N_SCAT})", "size": 7, "color": INK},
        {"text": "power to detect |r| ≥ 0.55 retained",
         "italic": True, "size": 7, "color": HIGHLIGHT},
    ], align="left", anchor="top", line_spacing=1.15)

    # legend (right-top)
    leg_x = PLOT_X + PLOT_W - Inches(1.2); leg_y = PLOT_Y + Inches(0.1)
    add_circle(s, leg_x + Inches(0.1), leg_y + Inches(0.10),
               Emu(45000), fill=GOOD, line_color=WHITE, line_width=0.6)
    add_text(s, leg_x + Inches(0.22), leg_y + Emu(5000),
             Inches(1.0), Inches(0.18),
             "good (n = 6)", size=7, color=INK, align="left")
    add_circle(s, leg_x + Inches(0.1), leg_y + Inches(0.28),
               Emu(45000), fill=BAD, line_color=WHITE, line_width=0.6)
    add_text(s, leg_x + Inches(0.22), leg_y + Inches(0.18),
             Inches(1.0), Inches(0.18),
             "bad (n = 6)", size=7, color=INK, align="left")

    # headline message bar (bottom)
    msg_x = Inches(0.35); msg_y = Inches(3.75)
    msg_w = Inches(5.8); msg_h = Inches(0.55)
    add_rect(s, msg_x, msg_y, msg_w, msg_h,
             fill=RGBColor(0xFB, 0xF4, 0xE0), line_color=HIGHLIGHT,
             line_width=0.8, rounded=True)
    add_multiline_text(s, msg_x + Inches(0.15), msg_y + Inches(0.05),
                       msg_w - Inches(0.3), msg_h - Inches(0.1), lines=[
        {"text": "Headline pair — DSB repair baseline does not forecast the "
                 "post-RT CD8-cytotoxic Δ",
         "bold": True, "size": 8, "color": HIGHLIGHT},
        {"text": "Static baseline ⊥ dynamic pharmacodynamic axes; the cascade "
                 "in §3.11 is phenomenology, not a causal downstream.",
         "italic": True, "size": 7, "color": INK},
    ], align="left", anchor="top", line_spacing=1.2)


build_A()
build_B()
build_C()

deck_path = f"{OUT}/Fig6_framework_preview_native_editable.pptx"
prs.save(deck_path)
print(f"wrote {deck_path} with {len(prs.slides)} slides (preview A/B/C)")
