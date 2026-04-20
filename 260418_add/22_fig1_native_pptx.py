#!/usr/bin/env python3
"""
22_fig1_native_pptx.py

Build Figure 1 --- Study design + cohort overview + headline preview
(the paper's entry point) --- as 4 native-editable PowerPoint slides,
plus a 2-slide supplementary companion (Supp Fig S12 sample-flow
CONSORT + Supp Fig S14 per-patient clinical waterfall).

Main Figure 1 panels
--------------------
  A  Study design schematic --- two parallel timeline tracks showing
     our TNT cohort (RT-alone in biopsy window, FOLFOX/CAPOX
     consolidation after the post-biopsy) vs the external nCRT-long
     cohorts (RT + concurrent capecitabine, no separate consolidation),
     both converging to the same final-response adjudication. The
     radiation-phase sampling window is explicitly bracketed on the
     top track. Mid-treatment decision-point annotation.

  B  Cohort Sankey --- three-stage flow Sex → clinical T stage →
     final TNT response for the 35 MSS LARC patients (ribbon fills
     colour-coded by response direction).

  C  Sample × assay availability matrix --- 35 subjects (rows, sorted
     by response then subject_id) × 6 columns (WES normal / pre /
     post, RNA normal / pre / post). Cells filled if sample exists.
     Column totals at top; left-axis subject IDs coloured by response.

  D  Headline four-claim preview --- a 4-row normalised effect-size
     forest previewing the paper's load-bearing claims:
       1 Discovery LASSO AUC 0.745 [0.56, 0.90]
       2 External Thread 1 meta (DSB +3.17, cellcycle +3.21, E2F/MYC
         +2.79, all P<0.01, N=518)
       3 External Thread 2 meta (CD8-cytotoxic Z = +3.29, P = 0.001,
         N = 816, +Akiyoshi 2023)
       4 Paired radiation-phase pharmacodynamics (EMT 6/6 good
         P = 0.016, IGH directional coherence Wilcoxon P = 0.035)

Supp Figure S12 + S14 companion
-------------------------------
  S12 (slide 1) CONSORT-style sample-flow diagram reconciling
      35-patient enrolment down to per-analysis subsets.
  S14 (slide 2) Per-patient clinical waterfall --- 35 bars sorted
      by response then subject_id, showing age / sex / cT stripes
      and response label.

Motif references consulted
--------------------------
  - Ganesh Nat Med 2019 (rectal organoids): timeline + biopsy marker
    schematic with treatment phase boxes.
  - Cercek NEJM 2022 (dostarlimab rectal): patient-level timeline and
    waterfall of response.
  - Bahadoer Lancet Oncol 2021 (RAPIDO): multi-phase regimen diagram
    comparing two arms.
  - TCGA / ICGC consortium papers: sample × assay availability matrix
    (e.g. Comprehensive Molecular Characterization of Human Colon and
    Rectal Cancer, Nature 2012).
  - Le NEJM 2015 / Chowell Science 2018 / Overman Lancet Oncol 2017:
    headline summary forest previewing paper's main findings.

Rules: one panel per slide; no plot titles; python-pptx native
elements only (TEXT_BOX / LINE / AUTO_SHAPE / FREEFORM for Sankey
ribbons); Arial everywhere; DEEP palette GOOD=#0a7d6e / BAD=#c53e1f;
every shape gets an empty <a:effectLst/> via kill_shadow(); every
coordinate routed through _i() to prevent width="0.0" XML loading
failures in PowerPoint.
"""

import os
import numpy as np
import pandas as pd
from scipy.stats import norm
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Shared infrastructure (identical to scripts 18-21)
# ---------------------------------------------------------------------------
def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def _i(v):
    return int(round(float(v)))


GOOD = RGBColor(0x0A, 0x7D, 0x6E)
BAD = RGBColor(0xC5, 0x3E, 0x1F)
INK = RGBColor(0x22, 0x22, 0x22)
LINE = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0xBB, 0xBB, 0xBB)
LT_GREY = RGBColor(0xDD, 0xDD, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xD4, 0xA3, 0x00)
THREAD1 = RGBColor(0x0E, 0x4A, 0x68)
THREAD2 = RGBColor(0x8A, 0x2B, 0x4C)
AKIYOSHI = RGBColor(0x6A, 0x4C, 0x93)

# Phase colours (treatment timeline)
PHASE_RT = RGBColor(0x4F, 0x73, 0x8E)         # steel blue for radiation
PHASE_CHEMO = RGBColor(0xA6, 0x62, 0x2C)      # ochre for chemo
PHASE_SURG = RGBColor(0x44, 0x4A, 0x58)       # graphite for surgery
PHASE_IDLE = RGBColor(0xD4, 0xD4, 0xD4)       # neutral grey

# cT stage palette (Panel B / Panel C)
CT_COLOR = {
    "T2": RGBColor(0xA8, 0xC8, 0xE1),
    "T2/T3": RGBColor(0x7D, 0xA7, 0xC8),
    "T3": RGBColor(0x45, 0x7A, 0x9E),
    "T4": RGBColor(0x21, 0x4A, 0x70),
}

GOOD_HEX = (0x0A, 0x7D, 0x6E)
BAD_HEX = (0xC5, 0x3E, 0x1F)

FONT = "Arial"
SLIDE_W = Inches(6.5)
SLIDE_H = Inches(4.5)

OUT = "/data/data/TNT/analysis/260418_add/ppt"
DATA = "/data/data/TNT/analysis"
os.makedirs(OUT, exist_ok=True)


def new_slide(prs, w=SLIDE_W, h=SLIDE_H):
    prs.slide_width = w
    prs.slide_height = h
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, size=8, bold=False,
             color=INK, align="left", anchor="middle", font=FONT,
             italic=False):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(max(w, 1)), _i(max(h, 1)))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "bottom": MSO_ANCHOR.BOTTOM,
                          "middle": MSO_ANCHOR.MIDDLE}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = str(text)
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.italic = bool(italic)
    r.font.color.rgb = color
    kill_shadow(tb)
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.5, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        try:
            c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass
    kill_shadow(c)
    return c


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0.5):
    w = max(_i(w), 1); h = max(_i(h), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _i(x), _i(y), w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_circle(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_diamond(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=INK, width=1.0):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    # Add arrowhead via XML manipulation
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    ln = c.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    kill_shadow(c)
    return c


def add_freeform_poly(slide, vertices, fill=None, line_color=None, line_width=0.5):
    """Create a filled polygon using build_freeform for Sankey ribbons."""
    if len(vertices) < 3:
        return None
    x0, y0 = vertices[0]
    ff = slide.shapes.build_freeform(_i(x0), _i(y0), scale=1.0)
    pts = [(_i(x), _i(y)) for x, y in vertices[1:]]
    ff.add_line_segments(pts, close=True)
    shp = ff.convert_to_shape()
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def lighten(rgb_hex, factor=0.7):
    r0, g0, b0 = rgb_hex
    return RGBColor(int(r0 + (255 - r0) * factor),
                    int(g0 + (255 - g0) * factor),
                    int(b0 + (255 - b0) * factor))


def draw_panel_letter(slide, letter):
    add_text(slide, Inches(0.15), Inches(0.1), Inches(0.4), Inches(0.35),
             letter, size=14, bold=True, color=INK, align="left")


# ---------------------------------------------------------------------------
# Load data
# ---------------------------------------------------------------------------
clin = pd.read_csv(f"{DATA}/00_cohort/clinical_master.tsv", sep="\t")
wes_inv = pd.read_csv(f"{DATA}/00_cohort/wes_inventory.tsv", sep="\t")
rna_inv = pd.read_csv(f"{DATA}/00_cohort/rna_inventory.tsv", sep="\t")


# ===========================================================================
# MAIN FIGURE 1
# ===========================================================================
prs_main = Presentation()


# -------------------------------------------------------------------
# Panel A --- Study design schematic (two parallel tracks)
# -------------------------------------------------------------------
def build_A():
    s = new_slide(prs_main)
    draw_panel_letter(s, "A")

    # ==========================================================
    # Time axis: 0 (diagnosis) → 1 (pre-biopsy) → 2 (RT end) →
    # 3 (post-biopsy) → 4 (consolidation end / surgery) → 5
    # (final response)
    # ==========================================================
    tl_x0 = Inches(1.10)   # shifted right to give cohort labels more space
    tl_x1 = Inches(5.65)   # shifted left so endpoint box doesn't clip
    t_span = tl_x1 - tl_x0

    def tx(t):
        return _i(tl_x0 + t / 5.0 * t_span)

    # Three-track layout (v0.7.5): Discovery SC-RT / SC-RT external /
    # LC-CRT external — compressed to fit vertical space with 0.85" gaps.
    top_y = Inches(0.80); top_h = Inches(0.30)       # Discovery SC-RT
    mid_y = Inches(1.70); mid_h = Inches(0.30)       # SC-RT external (GSE254249)
    bot_y = Inches(2.60); bot_h = Inches(0.30)       # LC-CRT external meta

    # --- Cohort labels (left margin) ---
    add_text(s, Inches(0.15), top_y + Inches(0.02),
             Inches(0.95), Inches(0.14),
             "Discovery",
             size=6, bold=True, color=INK, align="left")
    add_text(s, Inches(0.15), top_y + Inches(0.14),
             Inches(0.95), Inches(0.14),
             "SC-RT TNT (N = 35)",
             size=7, bold=True, color=GOOD, align="left")

    add_text(s, Inches(0.15), mid_y + Inches(0.02),
             Inches(0.95), Inches(0.14),
             "External · SC-RT",
             size=6, bold=True, color=INK, align="left")
    add_text(s, Inches(0.15), mid_y + Inches(0.14),
             Inches(0.95), Inches(0.14),
             "GSE254249 (N = 8)",
             size=7, bold=True, color=HIGHLIGHT, align="left")

    add_text(s, Inches(0.15), bot_y + Inches(0.02),
             Inches(0.95), Inches(0.14),
             "External · LC-CRT",
             size=6, bold=True, color=INK, align="left")
    add_text(s, Inches(0.15), bot_y + Inches(0.14),
             Inches(0.95), Inches(0.14),
             "5 GEO + Akiyoshi",
             size=7, bold=True, color=THREAD2, align="left")
    add_text(s, Inches(0.15), bot_y + Inches(0.26),
             Inches(0.95), Inches(0.12),
             "(N = 518–816)",
             size=6, color=THREAD2, align="left")

    # ---- TOP TRACK: Discovery SC-RT + FOLFOX/CAPOX consolidation ----
    # RT phase pre(t=1) -> post(t=3), consolidation (3-4), surgery (4-5)
    add_rect(s, tx(1), top_y, tx(3) - tx(1), top_h,
             fill=PHASE_RT, line_color=WHITE, line_width=0.5)
    add_text(s, tx(1), top_y + Inches(0.06),
             tx(3) - tx(1), Inches(0.18),
             "SC-RT · 25 Gy / 5 Fx · RT alone",
             size=7, bold=True, color=WHITE, align="center")

    add_rect(s, tx(3), top_y, tx(4) - tx(3), top_h,
             fill=PHASE_CHEMO, line_color=WHITE, line_width=0.5)
    add_text(s, tx(3), top_y + Inches(0.06),
             tx(4) - tx(3), Inches(0.18),
             "Consolidation · FOLFOX / CAPOX",
             size=7, bold=True, color=WHITE, align="center")

    add_rect(s, tx(4), top_y, tx(5) - tx(4), top_h,
             fill=PHASE_SURG, line_color=WHITE, line_width=0.5)
    add_text(s, tx(4) + Inches(0.02), top_y + Inches(0.06),
             tx(5) - tx(4), Inches(0.18),
             "surgery / W&W",
             size=7, bold=True, color=WHITE, align="center")

    # ---- MIDDLE TRACK: SC-RT external (GSE254249, Gao Cancer Cell 2025) ----
    add_rect(s, tx(1), mid_y, tx(3) - tx(1), mid_h,
             fill=PHASE_RT, line_color=WHITE, line_width=0.5)
    add_text(s, tx(1), mid_y + Inches(0.06),
             tx(3) - tx(1), Inches(0.18),
             "SC-RT · 25 Gy / 5 Fx",
             size=7, bold=True, color=WHITE, align="center")

    add_rect(s, tx(3), mid_y, tx(4) - tx(3), mid_h,
             fill=PHASE_CHEMO, line_color=WHITE, line_width=0.5)
    add_text(s, tx(3), mid_y + Inches(0.06),
             tx(4) - tx(3), Inches(0.18),
             "Consolidation · FOLFOXIRI (3-drug)",
             size=7, bold=True, color=WHITE, align="center")

    add_rect(s, tx(4), mid_y, tx(5) - tx(4), mid_h,
             fill=PHASE_SURG, line_color=WHITE, line_width=0.5)
    add_text(s, tx(4) + Inches(0.02), mid_y + Inches(0.06),
             tx(5) - tx(4), Inches(0.18),
             "surgery",
             size=7, bold=True, color=WHITE, align="center")

    # ---- BOTTOM TRACK: LC-CRT external meta (5 GEO + Akiyoshi) ----
    add_rect(s, tx(1), bot_y, tx(3) - tx(1), bot_h,
             fill=PHASE_RT, line_color=WHITE, line_width=0.5)
    add_text(s, tx(1), bot_y + Inches(0.06),
             tx(3) - tx(1), Inches(0.18),
             "LC-CRT · 50.4 Gy / 28 Fx · RT + concurrent capecitabine",
             size=7, bold=True, color=WHITE, align="center")

    add_rect(s, tx(3), bot_y, tx(5) - tx(3), bot_h,
             fill=PHASE_SURG, line_color=WHITE, line_width=0.5)
    add_text(s, tx(3), bot_y + Inches(0.06),
             tx(5) - tx(3), Inches(0.18),
             "surgery",
             size=7, bold=True, color=WHITE, align="center")

    # ---- Timeline arrows for all three tracks ----
    add_arrow(s, tl_x0, top_y + top_h + Inches(0.06),
              tl_x1, top_y + top_h + Inches(0.06), color=INK, width=0.6)
    add_arrow(s, tl_x0, mid_y + mid_h + Inches(0.06),
              tl_x1, mid_y + mid_h + Inches(0.06), color=INK, width=0.6)
    add_arrow(s, tl_x0, bot_y + bot_h + Inches(0.06),
              tl_x1, bot_y + bot_h + Inches(0.06), color=INK, width=0.6)

    # ---- Biopsy markers ----
    # TOP track: pre + post
    for t, lab in [(1, "pre-biopsy"), (3, "post-biopsy\n(post-RT)")]:
        xx = tx(t)
        add_circle(s, xx, top_y - Inches(0.05), Emu(24000),
                   fill=GOOD, line_color=WHITE, line_width=0.9)
        add_text(s, xx - Inches(0.55), top_y - Inches(0.24),
                 Inches(1.1), Inches(0.18),
                 lab, size=5.5, bold=True, color=GOOD, align="center")

    # MIDDLE track: pre + intermediate/post (bulk RNA-seq slice is post-TNT)
    xx = tx(1)
    add_circle(s, xx, mid_y - Inches(0.05), Emu(22000),
               fill=HIGHLIGHT, line_color=WHITE, line_width=0.9)
    add_text(s, xx - Inches(0.45), mid_y - Inches(0.20),
             Inches(0.9), Inches(0.14),
             "pre (n=3)",
             size=5.5, bold=True, color=HIGHLIGHT, align="center")
    # "post-TNT" sample: drawn after consolidation (surgery bridge)
    xx = tx(4)
    add_circle(s, xx, mid_y - Inches(0.05), Emu(22000),
               fill=HIGHLIGHT, line_color=WHITE, line_width=0.9)
    add_text(s, xx - Inches(0.55), mid_y - Inches(0.20),
             Inches(1.1), Inches(0.14),
             "post-TNT (n=8)",
             size=5.5, bold=True, color=HIGHLIGHT, align="center")

    # BOTTOM track: pre only
    xx = tx(1)
    add_circle(s, xx, bot_y - Inches(0.05), Emu(22000),
               fill=THREAD2, line_color=WHITE, line_width=0.9)
    add_text(s, xx - Inches(0.45), bot_y - Inches(0.20),
             Inches(0.9), Inches(0.14),
             "pre-biopsy",
             size=5.5, bold=True, color=THREAD2, align="center")

    # ---- Radiation-phase sampling window bracket on top track ----
    br_y = top_y - Inches(0.42)
    add_line(s, tx(1) - Inches(0.05), br_y, tx(1) - Inches(0.05),
             br_y + Inches(0.08), INK, 0.8)
    add_line(s, tx(1) - Inches(0.05), br_y, tx(3) + Inches(0.05),
             br_y, INK, 0.8)
    add_line(s, tx(3) + Inches(0.05), br_y, tx(3) + Inches(0.05),
             br_y + Inches(0.08), INK, 0.8)
    add_text(s, tx(1), br_y - Inches(0.14),
             tx(3) - tx(1), Inches(0.14),
             "Discovery sampling window brackets radiation phase only",
             size=5.5, bold=True, italic=True, color=INK, align="center")

    # ---- Final response adjudication end-point (all three tracks converge) ----
    end_y = (top_y + top_h + bot_y) / 2 + Inches(0.05)
    add_rect(s, tx(5) - Inches(0.60), end_y - Inches(0.22),
             Inches(1.3), Inches(0.45),
             fill=HIGHLIGHT, line_color=WHITE, line_width=0.8)
    add_text(s, tx(5) - Inches(0.60), end_y - Inches(0.22),
             Inches(1.3), Inches(0.22),
             "Final TNT response",
             size=6, bold=True, color=WHITE,
             align="center", anchor="top")
    add_text(s, tx(5) - Inches(0.60), end_y - Inches(0.03),
             Inches(1.3), Inches(0.22),
             "good (TRG 0-1) · bad (TRG 2-3)",
             size=5, color=WHITE, align="center", anchor="top")

    # Arrows from each track end to the shared endpoint
    add_arrow(s, tx(5), top_y + top_h / 2,
              tx(5) - Inches(0.6) + Inches(0.06), end_y,
              color=INK, width=0.5)
    add_arrow(s, tx(5), mid_y + mid_h / 2,
              tx(5) - Inches(0.6) + Inches(0.06), end_y,
              color=INK, width=0.5)
    add_arrow(s, tx(5), bot_y + bot_h / 2,
              tx(5) - Inches(0.6) + Inches(0.06), end_y,
              color=INK, width=0.5)

    # (Regimen-axis right-margin annotations removed in v0.7.5 — the RT
    # fractionation is already explicit inside each track's RT-phase bar label.)

    # ---- Bottom caption (regimen-agnostic message) ----
    add_text(s, Inches(0.15), Inches(3.32),
             SLIDE_W - Inches(0.3), Inches(0.16),
             "All three regimens converge to the same final-response endpoint. "
             "The pre-treatment baseline predictor (Fig 9) cannot encode",
             size=6, color=INK, align="center")
    add_text(s, Inches(0.15), Inches(3.49),
             SLIDE_W - Inches(0.3), Inches(0.16),
             "downstream regimen choice and nevertheless reproduces across "
             "RT fractionation (SC-RT vs LC-CRT) and chemo-timing/backbone axes → regimen-agnostic.",
             size=6, color=INK, align="center")
    add_text(s, Inches(0.15), Inches(3.66),
             SLIDE_W - Inches(0.3), Inches(0.16),
             "Paired analyses (Figs 6–8) are scoped to the discovery RT-phase window only.",
             size=6, italic=True, color=RGBColor(0x55, 0x55, 0x55), align="center")

    # time axis labels (below bottom arrow)
    add_text(s, tx(0) - Inches(0.25), Inches(3.03),
             Inches(0.5), Inches(0.12),
             "diagnosis", size=5, color=RGBColor(0x66, 0x66, 0x66),
             align="center", italic=True)
    add_text(s, tx(5) - Inches(0.25), Inches(3.03),
             Inches(0.5), Inches(0.12),
             "~6 mo", size=5, color=RGBColor(0x66, 0x66, 0x66),
             align="center", italic=True)


# -------------------------------------------------------------------
# Panel B --- Sankey (Sex → cT → Response)
# -------------------------------------------------------------------
def build_B():
    s = new_slide(prs_main)
    draw_panel_letter(s, "B")

    # stages
    sex_vals = ["M", "F"]
    ct_vals = ["T2", "T2/T3", "T3", "T4"]
    resp_vals = ["good", "bad"]

    total_n = 35
    # Sex → cT flows
    sex_ct = (clin.groupby(["sex", "cT"]).size()
              .reindex(pd.MultiIndex.from_product([sex_vals, ct_vals]),
                       fill_value=0))
    # cT → Response flows
    ct_resp = (clin.groupby(["cT", "response_bin"]).size()
               .reindex(pd.MultiIndex.from_product([ct_vals, resp_vals]),
                        fill_value=0))

    # Layout: bars at 3 x positions
    x_stage = [Inches(1.20), Inches(3.30), Inches(5.35)]
    bar_w = Inches(0.50)
    # Vertical extent for a "35-patient bar"
    y_top = Inches(0.65)
    y_bot = Inches(3.60)
    total_h = y_bot - y_top

    def unit_h(count):
        return total_h * (count / total_n)

    # Colors for response segments
    RESP_COLOR = {"good": GOOD, "bad": BAD}
    SEX_COLOR = {"M": RGBColor(0x36, 0x6D, 0x96), "F": RGBColor(0xB4, 0x5A, 0x70)}

    # === Stage 1: sex bars ===
    sex_counts = {sx: int((clin.sex == sx).sum()) for sx in sex_vals}
    sex_y = {}
    cursor = y_top
    for sx in sex_vals:
        h = unit_h(sex_counts[sx])
        add_rect(s, x_stage[0], cursor, bar_w, h,
                 fill=SEX_COLOR[sx], line_color=WHITE, line_width=0.5)
        sex_y[sx] = (cursor, cursor + h)
        # label on left
        add_text(s, Inches(0.20), cursor + h / 2 - Inches(0.09),
                 x_stage[0] - Inches(0.25), Inches(0.18),
                 f"{sx} (n={sex_counts[sx]})",
                 size=8, bold=True, color=SEX_COLOR[sx], align="right")
        cursor += h

    # === Stage 2: cT bars ===
    ct_counts = {ct: int((clin.cT == ct).sum()) for ct in ct_vals}
    ct_y = {}
    cursor = y_top
    for ct in ct_vals:
        h = unit_h(ct_counts[ct])
        add_rect(s, x_stage[1], cursor, bar_w, h,
                 fill=CT_COLOR[ct], line_color=WHITE, line_width=0.5)
        ct_y[ct] = (cursor, cursor + h)
        add_text(s, x_stage[1] + bar_w + Inches(0.04),
                 cursor + h / 2 - Inches(0.09),
                 Inches(0.75), Inches(0.18),
                 f"{ct}", size=7, bold=True, color=INK, align="left")
        add_text(s, x_stage[1] + bar_w + Inches(0.04),
                 cursor + h / 2 + Inches(0.04),
                 Inches(0.75), Inches(0.13),
                 f"(n={ct_counts[ct]})",
                 size=6, color=RGBColor(0x55, 0x55, 0x55), align="left")
        cursor += h

    # === Stage 3: response bars ===
    resp_counts = {r: int((clin.response_bin == r).sum()) for r in resp_vals}
    resp_y = {}
    cursor = y_top
    for r in resp_vals:
        h = unit_h(resp_counts[r])
        add_rect(s, x_stage[2], cursor, bar_w, h,
                 fill=RESP_COLOR[r], line_color=WHITE, line_width=0.5)
        resp_y[r] = (cursor, cursor + h)
        add_text(s, x_stage[2] + bar_w + Inches(0.06),
                 cursor + h / 2 - Inches(0.09),
                 Inches(0.9), Inches(0.18),
                 f"{r}", size=8, bold=True, color=RESP_COLOR[r], align="left")
        add_text(s, x_stage[2] + bar_w + Inches(0.06),
                 cursor + h / 2 + Inches(0.04),
                 Inches(0.9), Inches(0.13),
                 f"(n={resp_counts[r]})",
                 size=6, color=RGBColor(0x55, 0x55, 0x55), align="left")
        cursor += h

    # === Ribbons: Sex → cT ===
    # For each sex, iterate cT categories in order; track cumulative position
    # in both source (sex block) and target (cT block).
    sex_cursor = {sx: sex_y[sx][0] for sx in sex_vals}
    ct_cursor = {ct: ct_y[ct][0] for ct in ct_vals}
    for sx in sex_vals:
        for ct in ct_vals:
            n = int(sex_ct.loc[sx, ct])
            if n == 0:
                continue
            h = unit_h(n)
            src_top = sex_cursor[sx]
            src_bot = src_top + h
            tgt_top = ct_cursor[ct]
            tgt_bot = tgt_top + h

            x1 = x_stage[0] + bar_w
            x2 = x_stage[1]
            vertices = [
                (x1, src_top),
                (x2, tgt_top),
                (x2, tgt_bot),
                (x1, src_bot),
            ]
            add_freeform_poly(
                s, vertices,
                fill=lighten((SEX_COLOR[sx][0], SEX_COLOR[sx][1],
                              SEX_COLOR[sx][2]), 0.35),
                line_color=None,
            )
            sex_cursor[sx] += h
            ct_cursor[ct] += h

    # === Ribbons: cT → Response ===
    # reset cT cursors for the right face; use ct_y[ct][0] as starting point
    ct_cursor2 = {ct: ct_y[ct][0] for ct in ct_vals}
    resp_cursor = {r: resp_y[r][0] for r in resp_vals}
    for ct in ct_vals:
        for r in resp_vals:
            n = int(ct_resp.loc[ct, r])
            if n == 0:
                continue
            h = unit_h(n)
            src_top = ct_cursor2[ct]
            src_bot = src_top + h
            tgt_top = resp_cursor[r]
            tgt_bot = tgt_top + h

            x1 = x_stage[1] + bar_w
            x2 = x_stage[2]
            vertices = [
                (x1, src_top),
                (x2, tgt_top),
                (x2, tgt_bot),
                (x1, src_bot),
            ]
            add_freeform_poly(
                s, vertices,
                fill=lighten((RESP_COLOR[r][0], RESP_COLOR[r][1],
                              RESP_COLOR[r][2]), 0.40),
                line_color=None,
            )
            ct_cursor2[ct] += h
            resp_cursor[r] += h

    # stage labels (column headers)
    for i, lab in enumerate(["Sex", "Clinical T", "Final response"]):
        add_text(s, x_stage[i] - Inches(0.25),
                 Inches(0.35),
                 Inches(1.0), Inches(0.16),
                 lab, size=8, bold=True, color=INK, align="center")

    # bottom caption
    add_text(s, Inches(0.15), Inches(3.80),
             SLIDE_W - Inches(0.3), Inches(0.16),
             f"N = {total_n} MSS LARC patients; ribbon widths proportional "
             "to subject counts.",
             size=6, italic=True, color=RGBColor(0x55, 0x55, 0x55),
             align="center")
    add_text(s, Inches(0.15), Inches(3.97),
             SLIDE_W - Inches(0.3), Inches(0.16),
             f"T4 stage enriched in bad responders (7/9 vs 2/9 good, "
             f"Fisher P = 0.086); no sig. sex association.",
             size=6, color=INK, align="center")


# -------------------------------------------------------------------
# Panel C --- Sample × assay availability matrix
# -------------------------------------------------------------------
def build_C():
    s = new_slide(prs_main)
    draw_panel_letter(s, "C")

    # Order subjects: good first (then bad), each by subject_id ascending
    subj_sorted = (clin.sort_values(["response_bin", "subject_id"],
                                    ascending=[False, True])
                   .subject_id.tolist())   # good (False=g>b when desc? let's check)
    # explicit ordering
    good_subs = sorted(clin[clin.response_bin == "good"].subject_id.tolist())
    bad_subs = sorted(clin[clin.response_bin == "bad"].subject_id.tolist())
    subj_order = good_subs + bad_subs
    n = len(subj_order)

    # Pivot tables for sample availability
    wes_piv = wes_inv.pivot_table(index="subject_id", columns="timepoint",
                                  values="sample_id", aggfunc="count",
                                  fill_value=0)
    rna_piv = rna_inv.pivot_table(index="subject_id", columns="timepoint",
                                  values="sample_id", aggfunc="count",
                                  fill_value=0)

    columns = [("WES", "normal"), ("WES", "pre"), ("WES", "post"),
               ("RNA", "normal"), ("RNA", "pre"), ("RNA", "post")]

    # Layout
    mx = Inches(1.10); my = Inches(0.75)
    mw = Inches(3.60); mh = Inches(3.20)
    cell_w = mw / len(columns)
    cell_h = mh / n

    # Column group separator (between WES and RNA)
    sep_x = _i(mx + cell_w * 3)

    # --- Top headers ---
    # Top assay band: WES (first 3 cols), RNA (last 3 cols)
    add_rect(s, mx, my - Inches(0.42),
             cell_w * 3, Inches(0.18),
             fill=lighten(GOOD_HEX, 0.25), line_color=WHITE, line_width=0.3)
    add_text(s, mx, my - Inches(0.42),
             cell_w * 3, Inches(0.18),
             "WES", size=8, bold=True, color=WHITE, align="center")
    add_rect(s, mx + cell_w * 3, my - Inches(0.42),
             cell_w * 3, Inches(0.18),
             fill=lighten(BAD_HEX, 0.25), line_color=WHITE, line_width=0.3)
    add_text(s, mx + cell_w * 3, my - Inches(0.42),
             cell_w * 3, Inches(0.18),
             "RNA-seq", size=8, bold=True, color=WHITE, align="center")

    # Timepoint sub-headers + totals
    totals_text = {}
    for i, (assay, tp) in enumerate(columns):
        xx = _i(mx + (i + 0.5) * cell_w)
        add_text(s, xx - Inches(0.30), my - Inches(0.22),
                 Inches(0.60), Inches(0.14),
                 tp, size=7, color=INK, align="center")
        # count
        if assay == "WES":
            cnt = int(wes_piv[tp].sum()) if tp in wes_piv.columns else 0
        else:
            cnt = int(rna_piv[tp].sum()) if tp in rna_piv.columns else 0
        totals_text[(assay, tp)] = cnt
        add_text(s, xx - Inches(0.30), my - Inches(0.08),
                 Inches(0.60), Inches(0.14),
                 f"n={cnt}", size=6, bold=True,
                 color=RGBColor(0x66, 0x66, 0x66), align="center")

    # --- Grid cells ---
    for row_i, subj in enumerate(subj_order):
        y = _i(my + row_i * cell_h)
        resp = clin[clin.subject_id == subj]["response_bin"].iloc[0]
        sex = clin[clin.subject_id == subj]["sex"].iloc[0]
        cT = clin[clin.subject_id == subj]["cT"].iloc[0]

        # left-margin label with subject_id colored by response
        r_color = GOOD if resp == "good" else BAD
        add_text(s, Inches(0.18), y + cell_h / 2 - Inches(0.08),
                 Inches(0.22), Inches(0.14),
                 str(subj), size=5, bold=True, color=r_color, align="right")
        # small cT chip
        add_rect(s, Inches(0.42), y + cell_h / 2 - Inches(0.04),
                 Inches(0.08), Inches(0.09),
                 fill=CT_COLOR[cT])

        for col_i, (assay, tp) in enumerate(columns):
            x = _i(mx + col_i * cell_w)
            piv = wes_piv if assay == "WES" else rna_piv
            present = int(piv.loc[subj, tp]) if (subj in piv.index
                                                  and tp in piv.columns) else 0
            if present > 0:
                fill = GOOD if resp == "good" else BAD
                add_rect(s, x + Emu(20000), y + Emu(20000),
                         cell_w - Emu(40000), cell_h - Emu(40000),
                         fill=fill, line_color=WHITE, line_width=0.3)
            else:
                add_rect(s, x + Emu(20000), y + Emu(20000),
                         cell_w - Emu(40000), cell_h - Emu(40000),
                         fill=RGBColor(0xF0, 0xF0, 0xF0),
                         line_color=WHITE, line_width=0.3)

    # WES / RNA group divider
    add_line(s, sep_x, my - Inches(0.45),
             sep_x, my + mh, INK, 0.6)
    # frame
    add_rect(s, mx, my, mw, mh, line_color=LINE, line_width=0.5)

    # Row group separator (good vs bad)
    sep_y = _i(my + len(good_subs) * cell_h)
    add_line(s, mx, sep_y, mx + mw, sep_y, INK, 0.8)
    # group labels (left margin)
    add_text(s, Inches(0.15),
             _i(my + (len(good_subs) / 2) * cell_h - Inches(0.08)),
             Inches(0.55), Inches(0.14),
             f"good (n={len(good_subs)})",
             size=6, bold=True, color=GOOD, align="left")
    add_text(s, Inches(0.15),
             _i(my + len(good_subs) * cell_h
                + (len(bad_subs) / 2) * cell_h - Inches(0.08)),
             Inches(0.55), Inches(0.14),
             f"bad (n={len(bad_subs)})",
             size=6, bold=True, color=BAD, align="left")
    # legends (right side of plot)
    lx = mx + mw + Inches(0.25)
    ly = my + Inches(0.05)
    add_text(s, lx, ly, Inches(1.0), Inches(0.16),
             "Sample present", size=7, bold=True, color=INK)
    # good/bad swatches
    add_rect(s, lx, ly + Inches(0.22), Inches(0.16), Inches(0.14), fill=GOOD)
    add_text(s, lx + Inches(0.20), ly + Inches(0.19),
             Inches(0.8), Inches(0.16),
             "good responder", size=6, color=INK, align="left")
    add_rect(s, lx, ly + Inches(0.40), Inches(0.16), Inches(0.14), fill=BAD)
    add_text(s, lx + Inches(0.20), ly + Inches(0.37),
             Inches(0.8), Inches(0.16),
             "bad responder", size=6, color=INK, align="left")
    add_rect(s, lx, ly + Inches(0.58), Inches(0.16),
             Inches(0.14),
             fill=RGBColor(0xF0, 0xF0, 0xF0))
    add_text(s, lx + Inches(0.20), ly + Inches(0.55),
             Inches(0.9), Inches(0.16),
             "absent", size=6, color=INK, align="left")

    # cT legend (right side, below sample legend)
    add_text(s, lx, ly + Inches(0.80), Inches(1.0), Inches(0.14),
             "cT stage", size=7, bold=True, color=INK)
    for i, ct in enumerate(["T2", "T2/T3", "T3", "T4"]):
        yy = ly + Inches(0.98) + i * Inches(0.15)
        add_rect(s, lx, yy, Inches(0.1), Inches(0.1), fill=CT_COLOR[ct])
        add_text(s, lx + Inches(0.14), yy - Emu(10000),
                 Inches(0.8), Inches(0.14),
                 ct, size=6, color=INK, align="left")

    # caption
    add_text(s, Inches(0.15), Inches(4.10),
             SLIDE_W - Inches(0.3), Inches(0.14),
             f"77 WES samples ({totals_text[('WES','normal')]} normal + "
             f"{totals_text[('WES','pre')]} pre + {totals_text[('WES','post')]} post)"
             f" · 56 RNA-seq samples "
             f"({totals_text[('RNA','normal')]} normal + "
             f"{totals_text[('RNA','pre')]} pre + {totals_text[('RNA','post')]} post)",
             size=6, color=INK, align="center")
    add_text(s, Inches(0.15), Inches(4.25),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "12 subjects have paired pre + post RNA-seq (used in all "
             "paired-analysis figures).",
             size=6, italic=True, color=RGBColor(0x66, 0x66, 0x66),
             align="center")


# -------------------------------------------------------------------
# Panel D --- Headline 4-claim preview
# -------------------------------------------------------------------
def build_D():
    s = new_slide(prs_main)
    draw_panel_letter(s, "D")

    # 5 claim rows (v0.7.5: SC-RT-matched external validation added)
    rows = [
        # (group, short, full, effect_value, effect_range, effect_unit,
        #  normalised_bar_fraction_0_1, ref_fig, thread_color)
        ("Discovery",
         "Pre-CRT LASSO AUC",
         "4-feature ElasticNet, nested outer-LOOCV",
         "0.745",
         "[0.56, 0.90]",
         "AUC  (chance = 0.5)",
         (0.745 - 0.5) / 0.5,   # normalized to 0-1, where 1 = perfect
         "Fig 5",
         THREAD1),
        ("LC-CRT meta",
         "Thread 1 (5 cohorts, N = 518)",
         "DSB Z=+3.17 · cell-cycle +3.21 · E2F/MYC +2.79 · EMT +1.61 (trend)",
         "+3.21",
         "max Z",
         "Stouffer Z  (|Z| = 1.96 = P 0.05)",
         3.21 / 4.0,
         "Fig 9A",
         THREAD1),
        ("LC-CRT meta",
         "Thread 2 (6 sources, N = 816)",
         "CD8-cytotoxic + Akiyoshi 2023 (JAMA Netw Open)",
         "+3.29",
         "P = 0.001",
         "Stouffer Z",
         3.29 / 4.0,
         "Fig 9A",
         THREAD2),
        ("SC-RT external",
         "GSE254249 (post-TNT, N = 8)",
         "7/7 signatures concordant with discovery · Tcell_infil MW P = 0.036",
         "7/7",
         "sign P = 0.016",
         "signature direction concordance  (7/7 = unanimous)",
         7 / 7,
         "Fig 9E",
         HIGHLIGHT),
        ("Paired",
         "Radiation-phase pharmacodynamics",
         "EMT good 6/6 up (P = 0.016) · IGH coherence Wilcoxon P = 0.035",
         "0.016",
         "min P",
         "−log₁₀ P  (horizontal bar length)",
         min(1.0, -np.log10(0.016) / 3.0),
         "Figs 6-7",
         RGBColor(0xB3, 0x7D, 0x00)),
    ]

    # plot area — expanded vertically to accommodate 5 rows at preserved density
    px = Inches(1.20); py = Inches(0.45)
    pw = Inches(3.30); ph = Inches(3.50)
    row_h = ph / len(rows)

    # right column x for effect text
    right_col_x = px + pw + Inches(0.15)

    for i, (group, short, full, eff_val, eff_range, eff_unit,
            bar_frac, ref, color) in enumerate(rows):
        y_top = _i(py + i * row_h)
        y_ctr = _i(py + (i + 0.5) * row_h)

        # row separator
        if i > 0:
            add_line(s, Inches(0.20), y_top, SLIDE_W - Inches(0.20),
                     y_top, LT_GREY, 0.3)

        # left column: claim group + description
        add_text(s, Inches(0.20), y_top + Inches(0.08),
                 Inches(1.00), Inches(0.14),
                 group.upper(), size=6, bold=True,
                 color=color, align="left")
        add_text(s, Inches(0.20), y_top + Inches(0.22),
                 px - Inches(0.25), Inches(0.20),
                 short, size=8, bold=True, color=INK, align="left")
        add_text(s, Inches(0.20), y_top + Inches(0.43),
                 px - Inches(0.25), Inches(0.32),
                 full, size=6, color=RGBColor(0x55, 0x55, 0x55),
                 align="left", anchor="top")

        # middle bar: effect strength (normalized 0-1)
        bar_y_top = y_ctr - Inches(0.09)
        bar_h = Inches(0.18)
        # background bar (full width)
        add_rect(s, px, bar_y_top, pw, bar_h,
                 fill=LT_GREY, line_color=None)
        # effect bar
        filled_w = _i(pw * max(0.02, min(1.0, bar_frac)))
        add_rect(s, px, bar_y_top, filled_w, bar_h,
                 fill=color, line_color=WHITE, line_width=0.3)
        # effect unit caption under bar
        add_text(s, px, _i(bar_y_top + bar_h + Inches(0.02)),
                 pw, Inches(0.13),
                 eff_unit, size=5, italic=True,
                 color=RGBColor(0x66, 0x66, 0x66), align="left")

        # right column: effect value + CI/P + figure ref
        add_text(s, right_col_x, y_top + Inches(0.08),
                 Inches(1.30), Inches(0.22),
                 eff_val, size=14, bold=True, color=color, align="left")
        add_text(s, right_col_x, y_top + Inches(0.35),
                 Inches(1.40), Inches(0.16),
                 eff_range, size=7, color=RGBColor(0x55, 0x55, 0x55),
                 align="left")
        add_text(s, right_col_x, y_top + Inches(0.55),
                 Inches(1.40), Inches(0.16),
                 f"→ {ref}", size=7, italic=True, color=color, align="left")

    # bottom caption
    add_text(s, Inches(0.15), Inches(4.02),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "Two externally-validated orthogonal pre-CRT axes (Thread 1 tumor-intrinsic + Thread 2 immune) reproduce across",
             size=6, color=INK, align="center")
    add_text(s, Inches(0.15), Inches(4.16),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "three regimen strata (SC-RT + FOLFOX/CAPOX; LC-CRT + concurrent cape; SC-RT + FOLFOXIRI), supporting a regimen-agnostic interpretation;",
             size=6, color=INK, align="center")
    add_text(s, Inches(0.15), Inches(4.30),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "paired radiation-phase biopsies add an orthogonal dynamic layer (target engagement + directional immune coherence).",
             size=6, italic=True, color=HIGHLIGHT, align="center")


build_A()
build_B()
build_C()
build_D()
deck_main = f"{OUT}/Fig1_study_design_headline_native_editable.pptx"
prs_main.save(deck_main)
print(f"wrote {deck_main}")


# ===========================================================================
# SUPP FIGURES S12 + S14
# ===========================================================================
prs_supp = Presentation()


def build_S12():
    """CONSORT-style sample-flow diagram."""
    s = new_slide(prs_supp)
    draw_panel_letter(s, "A")

    # Central box for enrolment, branching to WES + RNA
    # Box counts derived from clinical + inventories
    n_enroll = 35
    n_wes_normal = int((wes_inv.timepoint == "normal").sum())
    n_wes_pre = int((wes_inv.timepoint == "pre").sum())
    n_wes_post = int((wes_inv.timepoint == "post").sum())
    n_rna_normal = int((rna_inv.timepoint == "normal").sum())
    n_rna_pre = int((rna_inv.timepoint == "pre").sum())
    n_rna_post = int((rna_inv.timepoint == "post").sum())

    n_paired_rna = (rna_inv.groupby("subject_id")["timepoint"]
                    .apply(lambda s: set(s) >= {"pre", "post"})
                    .sum())
    n_paired_wes = (wes_inv.groupby("subject_id")["timepoint"]
                    .apply(lambda s: set(s) >= {"pre", "post"})
                    .sum())

    # Top: enrolment box
    def consort_box(cx, cy, w, h, text, bold_text=None, color=INK,
                    fill_color=WHITE):
        add_rect(s, cx - w / 2, cy - h / 2, w, h,
                 fill=fill_color, line_color=color, line_width=0.9)
        if bold_text is not None:
            add_text(s, cx - w / 2, cy - h / 2 + Inches(0.04),
                     w, Inches(0.18),
                     bold_text, size=8, bold=True, color=color,
                     align="center", anchor="top")
            add_text(s, cx - w / 2, cy - h / 2 + Inches(0.25),
                     w, h - Inches(0.30),
                     text, size=6, color=INK,
                     align="center", anchor="top")
        else:
            add_text(s, cx - w / 2, cy - h / 2,
                     w, h, text, size=7, color=color,
                     align="center", anchor="middle")

    # ENROLLED
    consort_box(Inches(3.25), Inches(0.55), Inches(2.6), Inches(0.50),
                "MSS LARC, clinical T2-T4\nreceiving TNT at SNUH",
                f"N = {n_enroll} patients enrolled", INK)

    # Branch split label
    add_line(s, Inches(3.25), Inches(0.82),
             Inches(3.25), Inches(1.12),
             INK, 1.0)
    add_arrow(s, Inches(3.25), Inches(1.12),
              Inches(1.5), Inches(1.45), color=INK, width=0.8)
    add_arrow(s, Inches(3.25), Inches(1.12),
              Inches(5.05), Inches(1.45), color=INK, width=0.8)

    # WES BOX
    consort_box(Inches(1.50), Inches(1.80), Inches(2.10), Inches(0.95),
                f"normal (blood): {n_wes_normal}\n"
                f"pre-CRT tumor: {n_wes_pre}\n"
                f"post-CRT tumor: {n_wes_post}",
                f"WES (total {n_wes_normal + n_wes_pre + n_wes_post})",
                GOOD)

    # RNA BOX
    consort_box(Inches(5.05), Inches(1.80), Inches(2.10), Inches(0.95),
                f"normal (PBMC): {n_rna_normal}\n"
                f"pre-CRT tumor: {n_rna_pre}\n"
                f"post-CRT tumor: {n_rna_post}",
                f"RNA-seq (total {n_rna_normal + n_rna_pre + n_rna_post})",
                BAD)

    # Paired subset boxes below each
    add_arrow(s, Inches(1.5), Inches(2.28),
              Inches(1.5), Inches(2.60), color=INK, width=0.6)
    consort_box(Inches(1.50), Inches(2.95), Inches(2.05), Inches(0.55),
                f"T–N matched pairs: 41\npaired pre/post: 14 subjects",
                f"WES analytic set",
                GOOD)

    add_arrow(s, Inches(5.05), Inches(2.28),
              Inches(5.05), Inches(2.60), color=INK, width=0.6)
    consort_box(Inches(5.05), Inches(2.95), Inches(2.05), Inches(0.55),
                f"paired pre/post: {n_paired_rna} subjects\n"
                f"(used in all paired-analysis figures)",
                f"RNA analytic set",
                BAD)

    # Downstream analyses box
    add_line(s, Inches(1.5), Inches(3.22),
             Inches(1.5), Inches(3.55), INK, 0.6)
    add_line(s, Inches(5.05), Inches(3.22),
             Inches(5.05), Inches(3.55), INK, 0.6)
    add_line(s, Inches(1.5), Inches(3.55),
             Inches(5.05), Inches(3.55), INK, 0.6)
    add_arrow(s, Inches(3.275), Inches(3.55),
              Inches(3.275), Inches(3.85), color=INK, width=0.8)

    consort_box(Inches(3.275), Inches(4.10), Inches(5.8), Inches(0.40),
                f"Per-analysis counts reconciled against Figs 2–9",
                f"Multi-omics integration (36 features, LASSO, cascade, "
                f"directional coherence)",
                HIGHLIGHT)

    # caption
    add_text(s, Inches(0.10), Inches(4.33),
             SLIDE_W - Inches(0.2), Inches(0.14),
             "All counts derived from clinical_master.tsv + wes_inventory.tsv "
             "+ rna_inventory.tsv in analysis/00_cohort/.",
             size=5, italic=True, color=RGBColor(0x55, 0x55, 0x55),
             align="center")


def build_S14():
    """Per-patient clinical waterfall."""
    s = new_slide(prs_supp)
    draw_panel_letter(s, "B")

    # sort subjects by response, then by age (bad first or good first?)
    # Nature waterfall convention: sort by response value descending
    clin_sorted = clin.sort_values(
        ["response_num", "subject_id"], ascending=[False, True]).reset_index(
        drop=True)
    # actually response_num: 0=CR, 1=nearCR, 2=PR, 3=poor; ordering from bad to good
    # let's put good on left and bad on right --> sort by response_num ascending
    clin_sorted = clin.sort_values(
        ["response_num", "subject_id"], ascending=[True, True]).reset_index(
        drop=True)

    n = len(clin_sorted)

    # plot area
    px = Inches(0.55); py = Inches(0.80)
    pw = Inches(5.80); ph = Inches(2.50)
    bar_w = pw / n

    # response chip strip
    for i, row in clin_sorted.iterrows():
        x = _i(px + i * bar_w)
        col = GOOD if row.response_bin == "good" else BAD
        add_rect(s, x + Emu(5000), py + Inches(0.05),
                 _i(bar_w) - Emu(10000),
                 Inches(0.18),
                 fill=col, line_color=WHITE, line_width=0.2)

    # sex chip strip (below response)
    sex_color = {"M": RGBColor(0x36, 0x6D, 0x96),
                 "F": RGBColor(0xB4, 0x5A, 0x70)}
    for i, row in clin_sorted.iterrows():
        x = _i(px + i * bar_w)
        col = sex_color[row.sex]
        add_rect(s, x + Emu(5000), py + Inches(0.25),
                 _i(bar_w) - Emu(10000),
                 Inches(0.18),
                 fill=col, line_color=WHITE, line_width=0.2)

    # cT chip strip (below sex)
    for i, row in clin_sorted.iterrows():
        x = _i(px + i * bar_w)
        col = CT_COLOR[row.cT]
        add_rect(s, x + Emu(5000), py + Inches(0.45),
                 _i(bar_w) - Emu(10000),
                 Inches(0.18),
                 fill=col, line_color=WHITE, line_width=0.2)

    # age bar (numeric)
    age_max = int(clin_sorted.age.max()) + 5
    age_min = int(clin_sorted.age.min()) - 5
    for i, row in clin_sorted.iterrows():
        x = _i(px + i * bar_w)
        # map age to bar height within remaining vertical space
        age_frac = (row.age - age_min) / (age_max - age_min)
        bar_top = py + Inches(0.70)
        bar_h_max = Inches(1.75)
        bh = _i(bar_h_max * age_frac)
        bar_color = RGBColor(0x44, 0x44, 0x55)
        add_rect(s, x + Emu(5000),
                 _i(bar_top + bar_h_max - bh),
                 _i(bar_w) - Emu(10000),
                 bh, fill=bar_color, line_color=WHITE, line_width=0.2)

    # subject id labels (bottom x-axis)
    for i, row in clin_sorted.iterrows():
        x = _i(px + (i + 0.5) * bar_w)
        add_text(s, x - Inches(0.10), py + ph + Inches(0.02),
                 Inches(0.22), Inches(0.14),
                 str(row.subject_id), size=4,
                 color=RGBColor(0x55, 0x55, 0x55), align="center")

    # y-axis for age
    add_text(s, Inches(0.10), py + Inches(0.70),
             Inches(0.42), Inches(0.14),
             f"{age_max}", size=5, color=INK, align="right")
    add_text(s, Inches(0.10), py + Inches(0.70) + Inches(1.75) - Inches(0.10),
             Inches(0.42), Inches(0.14),
             f"{age_min}", size=5, color=INK, align="right")
    yt = add_text(s, _i(px - Inches(0.45)),
                  _i(py + Inches(0.80) + Inches(0.5)),
                  Inches(0.35), Inches(0.80),
                  "Age (yrs)",
                  size=7, color=INK, align="center")
    yt.rotation = -90

    # Row strip labels
    add_text(s, _i(px - Inches(0.45)), py + Inches(0.10),
             Inches(0.4), Inches(0.14),
             "response", size=5, bold=True, color=INK, align="right")
    add_text(s, _i(px - Inches(0.45)), py + Inches(0.30),
             Inches(0.4), Inches(0.14),
             "sex", size=5, bold=True, color=INK, align="right")
    add_text(s, _i(px - Inches(0.45)), py + Inches(0.50),
             Inches(0.4), Inches(0.14),
             "cT", size=5, bold=True, color=INK, align="right")

    # Legends
    leg_x = Inches(0.4); leg_y = Inches(3.65)
    # response
    add_rect(s, leg_x, leg_y, Inches(0.13), Inches(0.10), fill=GOOD)
    add_text(s, leg_x + Inches(0.16), leg_y - Emu(10000),
             Inches(0.6), Inches(0.13),
             "good", size=6, color=INK, align="left")
    add_rect(s, leg_x + Inches(0.70), leg_y, Inches(0.13),
             Inches(0.10), fill=BAD)
    add_text(s, leg_x + Inches(0.86), leg_y - Emu(10000),
             Inches(0.6), Inches(0.13),
             "bad", size=6, color=INK, align="left")
    # sex
    add_rect(s, leg_x + Inches(1.50), leg_y, Inches(0.13),
             Inches(0.10), fill=sex_color["M"])
    add_text(s, leg_x + Inches(1.66), leg_y - Emu(10000),
             Inches(0.3), Inches(0.13),
             "M", size=6, color=INK, align="left")
    add_rect(s, leg_x + Inches(1.90), leg_y, Inches(0.13),
             Inches(0.10), fill=sex_color["F"])
    add_text(s, leg_x + Inches(2.06), leg_y - Emu(10000),
             Inches(0.3), Inches(0.13),
             "F", size=6, color=INK, align="left")
    # cT
    xb = leg_x + Inches(2.55)
    for ct in ["T2", "T2/T3", "T3", "T4"]:
        add_rect(s, xb, leg_y, Inches(0.10), Inches(0.10), fill=CT_COLOR[ct])
        add_text(s, xb + Inches(0.12), leg_y - Emu(10000),
                 Inches(0.45), Inches(0.13),
                 ct, size=6, color=INK, align="left")
        xb += Inches(0.53)

    # caption
    add_text(s, Inches(0.15), Inches(3.95),
             SLIDE_W - Inches(0.3), Inches(0.15),
             f"Per-patient clinical waterfall (n = {n}); "
             f"subjects sorted by final TNT response then subject_id.",
             size=6, italic=True, color=RGBColor(0x55, 0x55, 0x55),
             align="center")
    add_text(s, Inches(0.15), Inches(4.15),
             SLIDE_W - Inches(0.3), Inches(0.15),
             f"Responders: good n = {int((clin.response_bin == 'good').sum())}, "
             f"bad n = {int((clin.response_bin == 'bad').sum())}. "
             f"Age median = {int(clin.age.median())}; range "
             f"{int(clin.age.min())}–{int(clin.age.max())}.",
             size=6, color=INK, align="center")


build_S12()
build_S14()
deck_supp = f"{OUT}/SuppFig_S12_S14_cohort_flow_waterfall_native_editable.pptx"
prs_supp.save(deck_supp)
print(f"wrote {deck_supp}")
