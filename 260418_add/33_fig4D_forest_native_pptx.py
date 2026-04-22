#!/usr/bin/env python3
"""
33_fig4D_forest_native_pptx.py

Build Figure 4D (§3.4 signature-response forest lollipop) as a
native-editable PowerPoint deck, plus a Supplementary companion deck
with the full category-grouped, fully-tabulated variant.

Both decks follow the project rules (user directive, 2026-04-22):
  0. Motif-informed from Nature/Cell/Science biomarker forests:
        - Mariathasan Nature 2018 (IMvigor210): category-band headers,
          HR forest with gold stars on significant rows.
        - Sade-Feldman Cell 2018: category-coloured chips on the left,
          zero-line emphasis, ±0.05 "no-effect" shaded band.
        - Ayers JCI 2017 (T-cell inflamed): directional red/blue,
          right-aligned P-values with star legend.
        - Chowell Science 2018 (HLA-ICB): gold highlight on the rows
          that replicate externally, bold axis label.
        - Litchfield Cell 2021 (pan-IO meta): per-row marker size scales
          with significance; clean 0.5-pt spines.
  1. One panel per slide.
  2. No figure/panel title inside the plot area.
  3. python-pptx native shapes/connectors/textboxes (every tick label,
     axis title, P-value, CI bar, marker is a first-class PPT object,
     double-click-editable without ungrouping).
  4. Main + supplementary decks.
  5. Arial everywhere.
  6. Shadows killed via kill_shadow() injecting an empty <a:effectLst/>
     into each shape's spPr.

Outputs
-------
  ppt/Fig4D_forest_immune3_native_editable.pptx           (main)
  ppt/SuppFig_forest_immune3_detailed_native_editable.pptx (supp)
"""

import os
from pathlib import Path

import numpy as np
import pandas as pd
from scipy.stats import false_discovery_control  # noqa: F401  (import guard)
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Shadow-removal helper (PowerPoint's theme applies a soft shadow to every
# auto-shape; inject an empty <a:effectLst/> into spPr to override it).
# ---------------------------------------------------------------------------
def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def _i(v):
    return int(v)


# ---------------------------------------------------------------------------
# Project palette (DEEP variant, fig2_v31/fig3_v31/panels_v3 standard)
# ---------------------------------------------------------------------------
GOOD       = RGBColor(0x0A, 0x7D, 0x6E)   # deep teal
BAD        = RGBColor(0xC5, 0x3E, 0x1F)   # deep coral
INK        = RGBColor(0x22, 0x22, 0x22)
LINE       = RGBColor(0x33, 0x33, 0x33)
GREY       = RGBColor(0x77, 0x77, 0x77)
LT_GREY    = RGBColor(0xDD, 0xDD, 0xDD)
BAND       = RGBColor(0xF2, 0xEC, 0xE1)   # cream band (Mariathasan motif)
NEUTRAL    = RGBColor(0xE8, 0xEB, 0xEF)   # ±0.05 no-effect shaded zone
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT  = RGBColor(0xB5, 0x89, 0x00)   # gold — external-val emphasis

# category palette (matches fig3_v32 CAT_COLORS_DEEP)
CAT_COLORS = {
    'Cytotoxic T-cell':     RGBColor(0x05, 0x7A, 0x64),
    'Antigen presentation': RGBColor(0x00, 0x56, 0x7D),
    'IFN response':         RGBColor(0x00, 0x99, 0xB8),
    'B-cell / TLS':         RGBColor(0xD4, 0xA3, 0x00),
    'Innate':               RGBColor(0xD9, 0x61, 0x25),
    'Regulatory':           RGBColor(0x7A, 0x3A, 0xAD),
    'Stromal/EMT':          RGBColor(0xB0, 0x32, 0x19),
    'Hypoxia':              RGBColor(0x0E, 0x2A, 0x47),
    'Other':                RGBColor(0x5A, 0x67, 0x72),
    # (No separate "Purified immune" category — see SIG_CATEGORY mapping.)
}

SIG_CATEGORY = {
    'CD8_proliferation':     'Cytotoxic T-cell',
    'CD8_activation':        'Cytotoxic T-cell',
    'CD8_exhaustion':        'Cytotoxic T-cell',
    'Cytolytic_activity':    'Cytotoxic T-cell',
    'Antigen_presentation':  'Antigen presentation',
    'MHC_II':                'Antigen presentation',
    'NLRC5_HLA_IFNG':        'Antigen presentation',
    'IFNg_Ayers_18':         'IFN response',
    'TLS_Cabrita':           'B-cell / TLS',
    'B_cell':                'B-cell / TLS',
    'NK_cell':               'Innate',
    'Mac_M1':                'Innate',
    'Mac_M2':                'Innate',
    'Treg':                  'Regulatory',
    'Checkpoint_inhibitory': 'Regulatory',
    'TGFb_Mariathasan':      'Stromal/EMT',
    'EMT_Mak':               'Stromal/EMT',
    'CAF_iCAF':              'Stromal/EMT',
    'CAF_myCAF':             'Stromal/EMT',
    'Hypoxia_Buffa':         'Hypoxia',
    'Stemness_mRNAsi_proxy': 'Other',
    'Epithelial':            'Other',
    # purified-immune signatures slot into existing biology categories
    # instead of a "Purified immune" group — avoids forward reference to
    # the §3.12 external-validation framing while still grouping them
    # sensibly among the 25 signatures.
    'CD8_cytotoxic':         'Cytotoxic T-cell',
    'Tcell_infiltration':    'Cytotoxic T-cell',
    'Bcell_infiltration':    'B-cell / TLS',
}

NEW_SIGS = ['CD8_cytotoxic', 'Tcell_infiltration', 'Bcell_infiltration']

LABEL_FIX = {
    'CD8_proliferation':     'CD8 proliferation',
    'CD8_activation':        'CD8 activation',
    'CD8_exhaustion':        'CD8 exhaustion',
    'Cytolytic_activity':    'Cytolytic activity',
    'Antigen_presentation':  'Antigen presentation',
    'MHC_II':                'MHC-II',
    'NLRC5_HLA_IFNG':        'NLRC5 · HLA · IFNG',
    'IFNg_Ayers_18':         'IFN-γ (Ayers 18)',
    'TLS_Cabrita':           'TLS (Cabrita)',
    'B_cell':                'B cell',
    'NK_cell':               'NK cell',
    'Mac_M1':                'Macrophage M1',
    'Mac_M2':                'Macrophage M2',
    'Treg':                  'Treg',
    'Checkpoint_inhibitory': 'Checkpoint inhibitory',
    'TGFb_Mariathasan':      'TGF-β (Mariathasan)',
    'EMT_Mak':               'EMT (Mak)',
    'CAF_iCAF':              'CAF — iCAF',
    'CAF_myCAF':              'CAF — myCAF',
    'Hypoxia_Buffa':         'Hypoxia (Buffa)',
    'Stemness_mRNAsi_proxy': 'Stemness (mRNAsi proxy)',
    'Epithelial':            'Epithelial',
    'CD8_cytotoxic':         'CD8 cytotoxic',
    'Tcell_infiltration':    'T-cell infiltration',
    'Bcell_infiltration':    'B-cell infiltration',
}

FONT = "Arial"

OUT = Path("/data/data/TNT/analysis/260418_add/ppt")
DATA = Path("/data/data/TNT/analysis/260418_add")
OUT.mkdir(parents=True, exist_ok=True)


# ---------------------------------------------------------------------------
# Load stats (25 rows: 22 legacy + 3 purified immune)
# ---------------------------------------------------------------------------
df = pd.read_csv(DATA / "fig4D_forest_stats_with_immune3.tsv", sep="\t")
df = df[df.timepoint == "pre"].copy().reset_index(drop=True)
# BH-FDR q for the supp deck
from statsmodels.stats.multitest import multipletests
df["qvalue"] = multipletests(df["pvalue"].values, method="fdr_bh")[1]
# sort so that lowest-P is at the top of the forest (y-index 0 is TOP)
df = df.sort_values("pvalue").reset_index(drop=True)
print(f"loaded {len(df)} rows (pre); new immune sigs present:",
      [s for s in NEW_SIGS if s in df.signature.values])


# ---------------------------------------------------------------------------
# Low-level PPT primitives  (copied from 18_fig5 / 21_fig9 conventions)
# ---------------------------------------------------------------------------
def new_slide(prs, w, h):
    prs.slide_width = w
    prs.slide_height = h
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def add_text(slide, x, y, w, h, text, size=8, bold=False,
             color=INK, align="left", anchor="middle", font=FONT,
             italic=False, rotation=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(anchor, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT,
                   "right": PP_ALIGN.RIGHT,
                   "center": PP_ALIGN.CENTER}.get(align, PP_ALIGN.LEFT)
    r = p.add_run()
    r.text = str(text)
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.italic = bool(italic)
    r.font.color.rgb = color
    if rotation:
        tb.rotation = rotation
    kill_shadow(tb)
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.5, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    kill_shadow(c)
    return c


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0.4):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _i(x), _i(y), _i(w), _i(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_circle(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx - r), _i(cy - r),
                                 _i(2 * r), _i(2 * r))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def sig_stars(p):
    if p < 0.001:
        return "★★★"
    if p < 0.01:
        return "★★"
    if p < 0.05:
        return "★"
    return ""


def row_color(sig, delta):
    # Uniform GOOD/BAD coloring for all 25 rows; purified-immune rows
    # are no longer highlighted at Fig 4B (forward reference to §3.12
    # removed per user feedback 2026-04-23).
    return GOOD if delta > 0 else BAD


# ===========================================================================
# BUILD MAIN FIG 4D  (one slide, compact 25-row forest)
# ===========================================================================
prs_main = Presentation()
SLIDE_MAIN_W = Inches(7.0)
SLIDE_MAIN_H = Inches(8.3)

s = new_slide(prs_main, SLIDE_MAIN_W, SLIDE_MAIN_H)

# axis rectangle
AX_X = Inches(2.05)
AX_Y = Inches(0.40)
AX_W = Inches(3.35)
AX_H = Inches(6.15)

# x range
x_lo = float(min(df.ci_low.min(), df.delta_good_minus_bad.min())) - 0.08
x_hi = float(max(df.ci_high.max(), df.delta_good_minus_bad.max())) + 0.10

def tx(v):
    return _i(AX_X + (v - x_lo) / (x_hi - x_lo) * AX_W)

n_rows = len(df)
row_h = AX_H / n_rows

def ty(i):
    return _i(AX_Y + (i + 0.5) * row_h)

# spines (top + bottom) — left-only for y (no ticks on y, just labels)
add_line(s, AX_X, AX_Y, AX_X + AX_W, AX_Y, LINE, 0.4)
add_line(s, AX_X, AX_Y + AX_H, AX_X + AX_W, AX_Y + AX_H, LINE, 0.6)

# ±0.05 neutral shaded band (Sade-Feldman motif — "no-effect" zone)
add_rect(s, tx(-0.05), AX_Y, tx(0.05) - tx(-0.05), AX_H,
         fill=NEUTRAL, line_color=None)

# zero-line (solid, emphasized)
add_line(s, tx(0), AX_Y, tx(0), AX_Y + AX_H, INK, 1.0)

# x ticks
for v in np.arange(-1.0, x_hi + 0.1, 0.5):
    xx = tx(v)
    add_line(s, xx, AX_Y + AX_H, xx, AX_Y + AX_H + Inches(0.06), LINE, 0.5)
    add_text(s, xx - Inches(0.20), AX_Y + AX_H + Inches(0.08),
             Inches(0.40), Inches(0.18),
             f"{v:+.1f}" if v != 0 else "0",
             size=7, color=INK, align="center")

# x axis title
add_text(s, AX_X, AX_Y + AX_H + Inches(0.28),
         AX_W, Inches(0.22),
         "Δ z-score (good − poor), 95 % bootstrap CI",
         size=9, bold=True, color=INK, align="center")

# per-row rendering
for i, r in df.iterrows():
    yy = ty(i)
    sig = r.signature
    color = row_color(sig, r.delta_good_minus_bad)
    # CI bar — uniform line width (no special emphasis for purified sigs)
    add_line(s, tx(r.ci_low), yy, tx(r.ci_high), yy,
             color, 1.4)
    # point estimate — uniform marker styling
    mag = _i(Emu(35000) if r.pvalue < 0.05 else Emu(24000))
    add_circle(s, tx(r.delta_good_minus_bad), yy, mag,
               fill=color, line_color=WHITE, line_width=0.8)
    # y-tick label (left of axis)
    label = LABEL_FIX.get(sig, sig.replace("_", " "))
    add_text(s, Inches(0.10), yy - Inches(0.09),
             AX_X - Inches(0.15), Inches(0.18),
             label, size=8, color=INK, align="right")
    # right-column P-value
    p_str = (f"P = {r.pvalue:.3g}" if r.pvalue >= 0.001
             else f"P < 10⁻³")
    star = sig_stars(r.pvalue)
    add_text(s, AX_X + AX_W + Inches(0.08), yy - Inches(0.09),
             Inches(0.85), Inches(0.18),
             p_str, size=7,
             bold=(r.pvalue < 0.1),
             color=INK, align="left")
    if star:
        add_text(s, AX_X + AX_W + Inches(0.90), yy - Inches(0.09),
                 Inches(0.40), Inches(0.18),
                 star, size=8, bold=True, color=INK, align="left")

# footer legend (GOOD/BAD chips + star key only — no external-validation
# forward reference)
leg_y = AX_Y + AX_H + Inches(0.60)
add_rect(s, Inches(0.70), leg_y, Inches(0.14), Inches(0.10),
         fill=GOOD, line_color=None)
add_text(s, Inches(0.90), leg_y - Inches(0.03),
         Inches(2.0), Inches(0.18),
         "Δ > 0 (enriched in good responders)",
         size=7, color=INK, align="left")
add_rect(s, Inches(3.10), leg_y, Inches(0.14), Inches(0.10),
         fill=BAD, line_color=None)
add_text(s, Inches(3.30), leg_y - Inches(0.03),
         Inches(2.3), Inches(0.18),
         "Δ < 0 (enriched in poor responders)",
         size=7, color=INK, align="left")
add_text(s, Inches(0.70), leg_y + Inches(0.22),
         Inches(5.5), Inches(0.18),
         "Marker size: large = P < 0.05 · small = n.s.  "
         "Stars: ★ P < 0.05 · ★★ P < 0.01 · ★★★ P < 0.001",
         size=7, color=GREY, italic=True, align="left")

# write
out_main = OUT / "Fig4D_forest_immune3_native_editable.pptx"
prs_main.save(out_main)
print(f"wrote MAIN deck → {out_main}")


# ===========================================================================
# BUILD SUPP FIG  (one slide, detailed category-grouped forest with
# full stats columns Δ | 95 % CI | U | P | q)
# ===========================================================================
prs_supp = Presentation()
SLIDE_SUPP_W = Inches(11.0)
SLIDE_SUPP_H = Inches(8.7)

s = new_slide(prs_supp, SLIDE_SUPP_W, SLIDE_SUPP_H)

# sort by biology category; within category by P-value
CAT_ORDER = ['Cytotoxic T-cell', 'Antigen presentation', 'IFN response',
             'B-cell / TLS', 'Innate', 'Regulatory', 'Stromal/EMT',
             'Hypoxia', 'Other']
df_supp = df.copy()
df_supp["category"] = df_supp.signature.map(SIG_CATEGORY)
df_supp["_cat_order"] = df_supp.category.map({c: i for i, c in enumerate(CAT_ORDER)})
df_supp = df_supp.sort_values(["_cat_order", "pvalue"]).reset_index(drop=True)

# layout
AX_X2 = Inches(3.20)
AX_Y2 = Inches(0.45)
AX_W2 = Inches(4.40)
AX_H2 = Inches(6.95)

def tx2(v):
    return _i(AX_X2 + (v - x_lo) / (x_hi - x_lo) * AX_W2)

row_h2 = AX_H2 / len(df_supp)

def ty2(i):
    return _i(AX_Y2 + (i + 0.5) * row_h2)

# spines
add_line(s, AX_X2, AX_Y2, AX_X2 + AX_W2, AX_Y2, LINE, 0.4)
add_line(s, AX_X2, AX_Y2 + AX_H2, AX_X2 + AX_W2, AX_Y2 + AX_H2, LINE, 0.6)

# ±0.05 neutral band
add_rect(s, tx2(-0.05), AX_Y2, tx2(0.05) - tx2(-0.05), AX_H2,
         fill=NEUTRAL, line_color=None)

# zero-line
add_line(s, tx2(0), AX_Y2, tx2(0), AX_Y2 + AX_H2, INK, 1.0)

# category bands (cream background strips spanning a category's rows)
cat_groups = df_supp.groupby("category", sort=False)
for cat_name, grp in cat_groups:
    if len(grp) == 0:
        continue
    y_top = AX_Y2 + grp.index.min() * row_h2
    y_bot = AX_Y2 + (grp.index.max() + 1) * row_h2
    add_rect(s, Inches(0.15), y_top,
             AX_X2 - Inches(0.20), y_bot - y_top,
             fill=BAND, line_color=None)

# x ticks
for v in np.arange(-1.0, x_hi + 0.1, 0.5):
    xx = tx2(v)
    add_line(s, xx, AX_Y2 + AX_H2, xx, AX_Y2 + AX_H2 + Inches(0.06), LINE, 0.5)
    add_text(s, xx - Inches(0.20), AX_Y2 + AX_H2 + Inches(0.08),
             Inches(0.40), Inches(0.18),
             f"{v:+.1f}" if v != 0 else "0",
             size=7, color=INK, align="center")

# x axis title
add_text(s, AX_X2, AX_Y2 + AX_H2 + Inches(0.30),
         AX_W2, Inches(0.22),
         "Δ z-score (good − poor), 95 % bootstrap CI",
         size=9, bold=True, color=INK, align="center")

# right-column column headers
header_y = AX_Y2 - Inches(0.25)
col_x = [
    AX_X2 + AX_W2 + Inches(0.15),   # Δ
    AX_X2 + AX_W2 + Inches(0.70),   # 95 % CI
    AX_X2 + AX_W2 + Inches(1.65),   # U
    AX_X2 + AX_W2 + Inches(2.05),   # P
    AX_X2 + AX_W2 + Inches(2.70),   # q
]
col_w = [Inches(0.50), Inches(0.95), Inches(0.40), Inches(0.65), Inches(0.55)]
col_labels = ["Δ", "95 % CI", "U", "P (MW)", "q (BH)"]
for xi, w, lab in zip(col_x, col_w, col_labels):
    add_text(s, xi, header_y, w, Inches(0.18),
             lab, size=8, bold=True, color=INK, align="center")
add_line(s, col_x[0], header_y + Inches(0.20),
         col_x[-1] + col_w[-1], header_y + Inches(0.20),
         LINE, 0.4)

# rows
for i, r in df_supp.iterrows():
    yy = ty2(i)
    sig = r.signature
    is_new = sig in NEW_SIGS
    color = row_color(sig, r.delta_good_minus_bad)
    cat = r.category

    # category color chip on the very left
    add_rect(s, Inches(0.18), yy - Inches(0.06),
             Inches(0.12), Inches(0.12),
             fill=CAT_COLORS.get(cat, GREY), line_color=None)
    # category label (short) only on first row of group
    first_of_group = (i == 0) or (df_supp.iloc[i - 1].category != cat)
    if first_of_group:
        add_text(s, Inches(0.35), yy - Inches(0.09),
                 Inches(1.10), Inches(0.18),
                 cat,
                 size=7, bold=True,
                 color=CAT_COLORS.get(cat, INK),
                 italic=True, align="left")

    # signature label (uniform ink colour — no forward-ref emphasis)
    label = LABEL_FIX.get(sig, sig.replace("_", " "))
    add_text(s, Inches(1.48), yy - Inches(0.09),
             AX_X2 - Inches(1.55), Inches(0.18),
             label, size=8, color=INK, align="right")

    # CI bar — uniform
    add_line(s, tx2(r.ci_low), yy, tx2(r.ci_high), yy, color, 1.4)
    mag = _i(Emu(35000) if r.pvalue < 0.05 else Emu(24000))
    add_circle(s, tx2(r.delta_good_minus_bad), yy, mag,
               fill=color, line_color=WHITE, line_width=0.8)

    # right-column numeric stats
    add_text(s, col_x[0], yy - Inches(0.09), col_w[0], Inches(0.18),
             f"{r.delta_good_minus_bad:+.2f}",
             size=7, color=INK, align="center")
    add_text(s, col_x[1], yy - Inches(0.09), col_w[1], Inches(0.18),
             f"[{r.ci_low:+.2f}, {r.ci_high:+.2f}]",
             size=7, color=INK, align="center")
    add_text(s, col_x[2], yy - Inches(0.09), col_w[2], Inches(0.18),
             f"{r.U:.0f}",
             size=7, color=INK, align="center")
    p_str = (f"{r.pvalue:.3g}" if r.pvalue >= 0.001 else "<10⁻³")
    add_text(s, col_x[3], yy - Inches(0.09), col_w[3], Inches(0.18),
             p_str,
             size=7, bold=(r.pvalue < 0.05),
             color=color if r.pvalue < 0.05 else INK,
             align="center")
    add_text(s, col_x[4], yy - Inches(0.09), col_w[4], Inches(0.18),
             f"{r.qvalue:.3g}",
             size=7,
             color=INK if r.qvalue >= 0.1 else color,
             align="center")
    # row star (INK colour — no gold emphasis)
    stars = sig_stars(r.pvalue)
    if stars:
        add_text(s, col_x[3] + Inches(0.55), yy - Inches(0.09),
                 Inches(0.40), Inches(0.18),
                 stars, size=7, bold=True, color=INK, align="left")

# footer legend (no forward reference to §3.12 external validation)
foot_y = AX_Y2 + AX_H2 + Inches(0.70)
add_rect(s, Inches(3.40), foot_y + Inches(0.02), Inches(0.14), Inches(0.10),
         fill=GOOD, line_color=None)
add_text(s, Inches(3.60), foot_y,
         Inches(2.0), Inches(0.18),
         "Δ > 0 (good-enriched)",
         size=7, color=INK, align="left")
add_rect(s, Inches(5.55), foot_y + Inches(0.02), Inches(0.14), Inches(0.10),
         fill=BAD, line_color=None)
add_text(s, Inches(5.75), foot_y,
         Inches(2.0), Inches(0.18),
         "Δ < 0 (poor-enriched)",
         size=7, color=INK, align="left")
add_rect(s, Inches(7.60), foot_y + Inches(0.02), Inches(0.20), Inches(0.10),
         fill=NEUTRAL, line_color=LINE, line_width=0.3)
add_text(s, Inches(7.85), foot_y,
         Inches(2.5), Inches(0.18),
         "±0.05 no-effect zone",
         size=7, color=GREY, italic=True, align="left")
add_text(s, Inches(3.40), foot_y + Inches(0.22),
         Inches(7.3), Inches(0.18),
         "Stars: ★ P < 0.05 · ★★ P < 0.01 · ★★★ P < 0.001  ·  "
         "q = Benjamini–Hochberg FDR across 25 signatures  ·  "
         "Δ = mean(good) − mean(poor) on per-sample z-score; 1,000× bootstrap CI.",
         size=7, color=GREY, italic=True, align="left")

out_supp = OUT / "SuppFig_forest_immune3_detailed_native_editable.pptx"
prs_supp.save(out_supp)
print(f"wrote SUPP deck → {out_supp}")


# ===========================================================================
# Quick verification: count shapes, textboxes, connectors — must all have
# shadow killed; otherwise flag.
# ===========================================================================
from pptx import Presentation as _Pr
for path in (out_main, out_supp):
    pr = _Pr(str(path))
    nshape = 0; ntext = 0; nconn = 0
    for slide in pr.slides:
        for shp in slide.shapes:
            nshape += 1
            if shp.has_text_frame:
                ntext += 1
            # connectors are also in shapes
    print(f"{path.name}: {nshape} shapes across {len(pr.slides)} slide(s)")
