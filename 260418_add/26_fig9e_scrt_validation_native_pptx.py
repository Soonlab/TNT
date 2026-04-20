#!/usr/bin/env python3
"""
26_fig9e_scrt_validation_native_pptx.py

Build Figure 9 Panel E (new in v0.7.5) as a native-editable PowerPoint
slide: SC-RT-matched external validation on GSE254249 (Gao Cancer Cell
2025).

Panel E content
---------------
  1. Per-signature horizontal bar — post-TNT Δ(good − bad) in GSE254249
     (n = 8, 5 CR vs 3 non-CR), 7 Thread 1 + Thread 2 signatures, with
     expected-direction shading (teal positive zone = good-up direction).
  2. Right-margin MW P annotation per row, with ★ gold-highlight on
     Tcell_infiltration P = 0.036 (the one nominal hit at n = 8).
  3. Discovery (n = 33) and LC-CRT meta (N = 518-816) effect-size
     reference markers as small hollow symbols behind each bar, so all
     three evidence streams are visible on a common axis.
  4. Bottom callout box summarising 7/7 concordance with discovery
     direction (binomial sign P = 0.016), the SC-RT + FOLFOXIRI regimen,
     and the three-regimen-stratum regimen-agnostic framing (SC-RT /
     LC-CRT / SC-RT-FOLFOXIRI all concordant).

Native PPT only: python-pptx AUTO_SHAPE / TEXT_BOX / LINE connectors,
Arial, kill_shadow, integer EMU, TNT palette.
"""
import os
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---- Shared infra (mirrors 21_fig9_native_pptx.py) ----
def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def _i(v):
    return int(round(float(v)))


GOOD = RGBColor(0x0A, 0x7D, 0x6E)
BAD = RGBColor(0xC5, 0x3E, 0x1F)
INK = RGBColor(0x22, 0x22, 0x22)
LINE = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0xBB, 0xBB, 0xBB)
LT_GREY = RGBColor(0xDD, 0xDD, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xD4, 0xA3, 0x00)
THREAD1 = RGBColor(0x0E, 0x4A, 0x68)
THREAD2 = RGBColor(0x8A, 0x2B, 0x4C)
GOOD_TINT = RGBColor(0xD6, 0xEC, 0xE7)
BAD_TINT = RGBColor(0xF4, 0xD9, 0xD3)

FONT = "Arial"
SLIDE_W = Inches(6.5)
SLIDE_H = Inches(4.5)

OUT = "/data/data/TNT/analysis/260418_add/ppt"
DATA = "/data/data/TNT/analysis/260418_add"
os.makedirs(OUT, exist_ok=True)


def new_slide(prs, w=SLIDE_W, h=SLIDE_H):
    prs.slide_width = w
    prs.slide_height = h
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, size=8, bold=False,
             color=INK, align="left", anchor="middle", font=FONT,
             italic=False):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(max(w, 1)), _i(max(h, 1)))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "bottom": MSO_ANCHOR.BOTTOM,
                          "middle": MSO_ANCHOR.MIDDLE}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = str(text)
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.italic = bool(italic)
    r.font.color.rgb = color
    kill_shadow(tb)
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.5, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        try:
            c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass
    kill_shadow(c)
    return c


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0.5):
    w = max(_i(w), 1); h = max(_i(h), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _i(x), _i(y), w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_circle(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_diamond(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def draw_panel_letter(slide, letter):
    add_text(slide, Inches(0.15), Inches(0.1), Inches(0.4), Inches(0.35),
             letter, size=14, bold=True, color=INK, align="left")


# ---- Load data ----
scrt = pd.read_csv(f"{DATA}/gse254249_post_response_stats.tsv", sep="\t")
discovery = pd.read_csv(f"{DATA}/discovery_signature_stats.tsv", sep="\t")
lc_meta = pd.read_csv(f"{DATA}/FINAL_meta_with_akiyoshi.tsv", sep="\t")

# Signature order: Thread 1 first (4 sigs), then Thread 2 (3 sigs)
SIG_ORDER = [
    ("DSB_HDR_repair", "DSB / HDR repair", 1, "up"),
    ("Tumor_cellcycle", "Tumor cell-cycle", 1, "up"),
    ("E2F_MYC_cellcycle", "E2F / MYC", 1, "up"),
    ("EMT", "EMT", 1, "down"),
    ("CD8_cytotoxic", "CD8-cytotoxic", 2, "up"),
    ("Tcell_infiltration", "T-cell infiltration", 2, "up"),
    ("Bcell_infiltration", "B-cell infiltration", 2, "up"),
]


def get_scrt(sig):
    r = scrt[scrt.signature == sig].iloc[0]
    return float(r.delta), float(r.mw_p), int(r.concordant)


def get_discovery(sig):
    r = discovery[discovery.signature == sig].iloc[0]
    return float(r.discovery_delta), float(r.discovery_Z)


def get_lc_Z(sig):
    r = lc_meta[lc_meta.signature == sig]
    if r.empty:
        return None
    return float(r.iloc[0].Z)


# ---- Build Panel E ----
prs = Presentation()
s = new_slide(prs)
draw_panel_letter(s, "E")

# Panel title (italic, below letter)
add_text(s, Inches(0.55), Inches(0.14),
         Inches(5.8), Inches(0.25),
         "SC-RT-matched external validation — GSE254249 (Gao et al. Cancer Cell 2025)",
         size=9, bold=True, color=INK, italic=False, align="left")
add_text(s, Inches(0.55), Inches(0.32),
         Inches(5.8), Inches(0.20),
         "n = 8 post-TNT bulk RNA-seq  ·  SC-RT 5×5 Gy + FOLFOXIRI  ·  5 CR vs 3 non-CR",
         size=6.5, italic=True, color=RGBColor(0x66, 0x66, 0x66), align="left")

# Plot axes
ax_x = Inches(2.10)
ax_y = Inches(0.70)
ax_w = Inches(3.35)
ax_h = Inches(2.55)
x_lo, x_hi = -1.6, 1.8


def tx(v):
    return _i(ax_x + (v - x_lo) / (x_hi - x_lo) * ax_w)


n_rows = len(SIG_ORDER)
row_h = ax_h / n_rows

def ty(i):
    return _i(ax_y + (i + 0.5) * row_h)


# Expected-direction shading (teal on good-up side)
# For "up"-expected signatures: positive x is good-direction -> teal
# For "down"-expected signature (EMT): negative x is good-direction -> teal
# Single banner: faint teal on +x, faint coral on -x
add_rect(s, tx(0), ax_y, ax_w - (tx(0) - ax_x), ax_h,
         fill=GOOD_TINT, line_color=None)
add_rect(s, ax_x, ax_y, tx(0) - ax_x, ax_h,
         fill=BAD_TINT, line_color=None)

# x-axis
add_line(s, ax_x, ax_y + ax_h, ax_x + ax_w, ax_y + ax_h, LINE, 0.6)

# Zero line + ticks
add_line(s, tx(0), ax_y, tx(0), ax_y + ax_h, INK, 1.0)
for v in [-1.5, -1.0, -0.5, 0, 0.5, 1.0, 1.5]:
    xx = tx(v)
    add_line(s, xx, _i(ax_y + ax_h), xx,
             _i(ax_y + ax_h + Inches(0.05)), LINE, 0.4)
    lbl = f"{v:+.1f}" if v != 0 else "0"
    add_text(s, xx - Inches(0.20), _i(ax_y + ax_h + Inches(0.06)),
             Inches(0.40), Inches(0.13),
             lbl, size=6, color=INK, align="center")
add_text(s, ax_x, _i(ax_y + ax_h + Inches(0.20)),
         ax_w, Inches(0.15),
         "Post-TNT Δ = z-score(good) − z-score(bad)   (predicted-direction-oriented)",
         size=7, color=INK, align="center")

# Top banner for predicted-direction shading
add_text(s, tx(1.2), ax_y - Inches(0.10),
         Inches(1.5), Inches(0.13),
         "predicted-direction zone →",
         size=5.5, italic=True, color=GOOD, align="left")
add_text(s, ax_x + Inches(0.02), ax_y - Inches(0.10),
         Inches(1.3), Inches(0.13),
         "← opposite direction",
         size=5.5, italic=True, color=BAD, align="left")

# Rows
for i, (sig, pretty, thread, pred) in enumerate(SIG_ORDER):
    y = ty(i)
    delta, p, conc = get_scrt(sig)
    # Re-orient delta so positive = predicted direction on x-axis
    dx = delta if pred == "up" else -delta
    dx_clamped = max(x_lo, min(x_hi, dx))

    # Row background separator (alternating light grey on odd rows)
    if i % 2 == 1:
        add_rect(s, ax_x, y - row_h / 2, ax_w, row_h,
                 fill=RGBColor(0xF8, 0xF8, 0xF8), line_color=None)

    # Signature label (left)
    row_color = THREAD1 if thread == 1 else THREAD2
    add_text(s, Inches(0.16), y - Inches(0.08),
             ax_x - Inches(0.20), Inches(0.16),
             pretty, size=7, bold=False, color=row_color, align="right")

    # SC-RT Δ bar
    bar_h = Emu(60000)
    if dx >= 0:
        # Positive bar: grow from 0 to tx(dx_clamped)
        bar_fill = GOOD if pred == "up" else BAD
        add_rect(s, tx(0), y - bar_h / 2,
                 tx(dx_clamped) - tx(0), bar_h,
                 fill=bar_fill, line_color=INK, line_width=0.4)
    else:
        bar_fill = BAD
        add_rect(s, tx(dx_clamped), y - bar_h / 2,
                 tx(0) - tx(dx_clamped), bar_h,
                 fill=bar_fill, line_color=INK, line_width=0.4)

    # Discovery and LC-CRT Z-score reference markers
    # Overlay as small hollow markers on SAME re-oriented x-axis
    #   discovery_Z:   from discovery tsv (direct z, already "good-up" oriented
    #                  because delta sign in TSV is good-minus-bad with expected
    #                  convention; for EMT discovery_Z uses signed delta
    #                  which is already negative = good<bad).
    disc_delta, disc_Z = get_discovery(sig)
    disc_dx = disc_delta if pred == "up" else -disc_delta
    disc_dx_clamped = max(x_lo, min(x_hi, disc_dx))
    add_circle(s, tx(disc_dx_clamped), y, Emu(18000),
               fill=WHITE, line_color=INK, line_width=0.9)

    # LC-CRT pooled Z -> displayed at a scaled position on the same x-axis
    # LC-CRT Z is in standard-error units, not z-score-delta units.
    # For a visual reference only, we map Z onto the x-axis by dividing
    # by a heuristic scaling factor 3 so Z = +3 maps to ~+1 on the delta
    # x-axis (this is a legend-level reference, exact numbers are in the
    # right-margin annotation).
    lc_Z = get_lc_Z(sig)
    if lc_Z is not None:
        lc_x_plot = max(x_lo, min(x_hi, lc_Z / 3.0))
        add_diamond(s, tx(lc_x_plot), y, Emu(22000),
                    fill=WHITE, line_color=THREAD1 if thread == 1 else THREAD2,
                    line_width=1.0)

    # Right-margin MW P annotation
    right_x = ax_x + ax_w + Inches(0.05)
    p_str = f"P = {p:.3f}" if p < 1 else "P = 1.000"
    star = "★" if p < 0.05 else ("·" if p < 0.10 else "")
    is_hit = p < 0.05
    add_text(s, right_x, y - Inches(0.13),
             Inches(0.72), Inches(0.13),
             f"Δ = {delta:+.2f}", size=6.5,
             bold=is_hit, color=HIGHLIGHT if is_hit else INK,
             align="left")
    add_text(s, right_x, y - Inches(0.01),
             Inches(0.72), Inches(0.13),
             p_str, size=6,
             color=HIGHLIGHT if is_hit else RGBColor(0x66, 0x66, 0x66),
             align="left")
    if star:
        add_text(s, right_x + Inches(0.55), y - Inches(0.13),
                 Inches(0.20), Inches(0.13),
                 star, size=9, bold=True, color=HIGHLIGHT, align="left")

# Bottom callout box: three-stream convergence summary
callout_y = Inches(3.42)
callout_h = Inches(0.85)
callout = add_rect(s, Inches(0.35), callout_y,
                   Inches(5.80), callout_h,
                   fill=RGBColor(0xFF, 0xF8, 0xE1), line_color=HIGHLIGHT, line_width=1.0)

add_text(s, Inches(0.45), callout_y + Inches(0.03),
         Inches(5.6), Inches(0.20),
         "Three-regimen regimen-agnostic convergence",
         size=8, bold=True, color=INK, align="left")
add_text(s, Inches(0.45), callout_y + Inches(0.22),
         Inches(5.6), Inches(0.18),
         "• Discovery SC-RT + FOLFOX/CAPOX (N=33 pre-RT): LASSO AUC 0.745 [0.56, 0.90]",
         size=6.5, color=INK, align="left")
add_text(s, Inches(0.45), callout_y + Inches(0.37),
         Inches(5.6), Inches(0.18),
         "• LC-CRT + concurrent capecitabine (N=518–816): T1 DSB Z=+3.17, cellcycle Z=+3.21, E2F/MYC Z=+2.79; T2 CD8 Z=+3.29",
         size=6.5, color=INK, align="left")
add_text(s, Inches(0.45), callout_y + Inches(0.52),
         Inches(5.6), Inches(0.18),
         "• SC-RT + FOLFOXIRI (N=8 post-TNT): 7/7 signatures concordant  —  binomial sign P = 0.016  —  Tcell_infil MW P = 0.036 ★",
         size=6.5, bold=True, color=HIGHLIGHT, align="left")
add_text(s, Inches(0.45), callout_y + Inches(0.67),
         Inches(5.6), Inches(0.18),
         "Baseline predictor biology reproduces across both RT fractionation and chemo-timing/backbone axes.",
         size=6.5, italic=True, color=INK, align="left")

# Legend
leg_y = Inches(4.30)
add_text(s, Inches(0.15), leg_y - Inches(0.03),
         Inches(1.0), Inches(0.13),
         "Markers:", size=6, bold=True, color=INK, align="left")
# Filled bar
add_rect(s, Inches(0.78), leg_y - Inches(0.01), Inches(0.16), Inches(0.07),
         fill=GOOD, line_color=INK, line_width=0.3)
add_text(s, Inches(0.98), leg_y - Inches(0.04), Inches(1.10), Inches(0.14),
         "SC-RT Δ (bar)", size=5.5, color=INK, align="left")
add_circle(s, Inches(2.00), leg_y + Inches(0.03), Emu(15000),
           fill=WHITE, line_color=INK, line_width=0.8)
add_text(s, Inches(2.08), leg_y - Inches(0.04), Inches(1.20), Inches(0.14),
         "Discovery Δ (circle)", size=5.5, color=INK, align="left")
add_diamond(s, Inches(3.30), leg_y + Inches(0.03), Emu(18000),
            fill=WHITE, line_color=THREAD1, line_width=1.0)
add_text(s, Inches(3.40), leg_y - Inches(0.04), Inches(1.80), Inches(0.14),
         "LC-CRT pooled Z / 3 (diamond)", size=5.5, color=INK, align="left")

# Save
out_path = f"{OUT}/Fig9E_scrt_validation_native_editable.pptx"
prs.save(out_path)
print(f"Wrote {out_path}")

# Also save a PDF export for inclusion in main composite
import subprocess
try:
    subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf",
                    "--outdir", OUT, out_path],
                   check=True, capture_output=True, timeout=60)
    print(f"Wrote {OUT}/Fig9E_scrt_validation_native_editable.pdf")
except Exception as e:
    print(f"(PDF conversion skipped: {e})")
