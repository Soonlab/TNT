#!/usr/bin/env python3
"""
37_fig4D_pathway_volcano_native_pptx.py

Build Figure 3C / Figure 4D PATHWAY-LEVEL volcano as native-editable
PowerPoint decks — main + supplementary.

User directive (2026-04-22): the volcano should be at pathway level,
divided into cell cycle / DNA repair / immune / EMT-Stromal categories.

Data: combined Hallmark (50) + curated Reactome (18) = 68 pathways,
scored by fgsea on the good-vs-poor pre-CRT DEG ranking.

Motif-informed from:
  - Subramanian et al PNAS 2005 + Liberzon Cell Syst 2015 (Hallmark
    paper): dot-per-pathway plots, NES on x, significance threshold
    horizontal.
  - Mariathasan 2018 Nature (IMvigor210, Fig 3b pathway volcano):
    category-coloured lollipop volcano, gold emphasis on top hits.
  - Sade-Feldman 2018 Cell (Fig 2g pathway volcano): direction
    arrows, horizontal threshold, per-category marker clustering.
  - Cristescu 2018 Science (T-cell inflamed pathway volcano): top
    pathway labels bold + coloured.
  - Rouillard Database 2016 / Enrichr: NES × -log10(FDR) with gold
    emphasis for |NES| ≥ 2 at publication-grade.

Rules followed:
  1. One panel per slide.
  2. No title inside plot.
  3. python-pptx native: axes/ticks/text/markers are all first-class
     shapes or textboxes (double-click editable).
  4. Main + supp decks.
  5. Arial.
  6. Shadows killed.

Outputs
-------
  ppt/Fig4D_pathway_volcano_native_editable.pptx              (main)
  ppt/SuppFig_pathway_volcano_detailed_native_editable.pptx   (supp)
"""

from pathlib import Path

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn
from lxml import etree


def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def _i(v):
    return int(v)


# Project palette + per-category colours (match matplotlib panel)
GOOD      = RGBColor(0x0A, 0x7D, 0x6E)
BAD       = RGBColor(0xC5, 0x3E, 0x1F)
INK       = RGBColor(0x0E, 0x2A, 0x47)
LINE      = RGBColor(0x33, 0x33, 0x33)
GREY      = RGBColor(0x77, 0x77, 0x77)
LT_GREY   = RGBColor(0xC0, 0xC6, 0xCF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xD4, 0xA3, 0x00)

CAT_COLOR = {
    'Cell cycle':  RGBColor(0x05, 0x7A, 0x64),
    'DNA repair':  RGBColor(0x00, 0x56, 0x7D),
    'Immune':      RGBColor(0xC1, 0x14, 0x56),
    'EMT/Stromal': RGBColor(0xB0, 0x32, 0x19),
    'Other':       RGBColor(0xA6, 0xAD, 0xB5),
}

FONT = "Arial"

OUT = Path("/data/data/TNT/analysis/260418_add/ppt")
DATA_ROOT = Path("/data/data/TNT/analysis")
OUT.mkdir(parents=True, exist_ok=True)


# ---------------------------------------------------------------------------
# Load + categorise (same logic as 36_fig4D_pathway_volcano_rebuild.py)
# ---------------------------------------------------------------------------
hm = pd.read_csv(DATA_ROOT / "05_rna_deg_gsea/GSEA_Hallmark_pre.tsv", sep="\t")
re = pd.read_csv(DATA_ROOT / "05_rna_deg_gsea/GSEA_Reactome_pre.tsv", sep="\t")

HM_CELL_CYCLE = {
    'HALLMARK_E2F_TARGETS','HALLMARK_G2M_CHECKPOINT','HALLMARK_MYC_TARGETS_V1',
    'HALLMARK_MYC_TARGETS_V2','HALLMARK_MITOTIC_SPINDLE'}
HM_DNA_REPAIR = {'HALLMARK_DNA_REPAIR'}
HM_IMMUNE = {
    'HALLMARK_INTERFERON_ALPHA_RESPONSE','HALLMARK_INTERFERON_GAMMA_RESPONSE',
    'HALLMARK_INFLAMMATORY_RESPONSE','HALLMARK_COMPLEMENT',
    'HALLMARK_IL2_STAT5_SIGNALING','HALLMARK_IL6_JAK_STAT3_SIGNALING',
    'HALLMARK_ALLOGRAFT_REJECTION','HALLMARK_TNFA_SIGNALING_VIA_NFKB',
    'HALLMARK_COAGULATION'}
HM_EMT_STROMAL = {
    'HALLMARK_EPITHELIAL_MESENCHYMAL_TRANSITION','HALLMARK_MYOGENESIS',
    'HALLMARK_APICAL_JUNCTION','HALLMARK_APICAL_SURFACE',
    'HALLMARK_ANGIOGENESIS','HALLMARK_HEDGEHOG_SIGNALING'}
RE_CELL_CYCLE = {
    'REACTOME_MITOTIC_PROMETAPHASE','REACTOME_MITOTIC_METAPHASE_AND_ANAPHASE',
    'REACTOME_MITOTIC_SPINDLE_CHECKPOINT',
    'REACTOME_MITOTIC_G1_PHASE_AND_G1_S_TRANSITION',
    'REACTOME_MITOTIC_G2_G2_M_PHASES'}
RE_DNA_REPAIR = {'REACTOME_HOMOLOGY_DIRECTED_REPAIR','REACTOME_MISMATCH_REPAIR'}
RE_IMMUNE = {
    'REACTOME_CD22_MEDIATED_BCR_REGULATION',
    'REACTOME_SIGNALING_BY_THE_B_CELL_RECEPTOR_BCR',
    'REACTOME_ANTIGEN_ACTIVATES_B_CELL_RECEPTOR_BCR_LEADING_TO_GENERATION_OF_SECOND_MESSENGERS',
    'REACTOME_TCR_SIGNALING','REACTOME_INTERFERON_ALPHA_BETA_SIGNALING'}
RE_EMT_STROMAL = {
    'REACTOME_ASSEMBLY_OF_COLLAGEN_FIBRILS_AND_OTHER_MULTIMERIC_STRUCTURES',
    'REACTOME_ELASTIC_FIBRE_FORMATION','REACTOME_ECM_PROTEOGLYCANS',
    'REACTOME_COLLAGEN_DEGRADATION','REACTOME_COLLAGEN_FORMATION',
    'REACTOME_INTEGRIN_CELL_SURFACE_INTERACTIONS'}

def categorise(name):
    if name in HM_CELL_CYCLE | RE_CELL_CYCLE: return 'Cell cycle'
    if name in HM_DNA_REPAIR | RE_DNA_REPAIR: return 'DNA repair'
    if name in HM_IMMUNE | RE_IMMUNE:         return 'Immune'
    if name in HM_EMT_STROMAL | RE_EMT_STROMAL: return 'EMT/Stromal'
    return 'Other'

re_keep = RE_CELL_CYCLE | RE_DNA_REPAIR | RE_IMMUNE | RE_EMT_STROMAL
re_sel = re[re.pathway.isin(re_keep)].copy()
hm['source'] = 'Hallmark'; re_sel['source'] = 'Reactome'
df = pd.concat([hm, re_sel], ignore_index=True)
df['cat'] = df.pathway.apply(categorise)
df['-log10padj'] = -np.log10(df['padj'].replace(0, 1e-300))

SHORT_MAP = {
    'Antigen Activates B Cell Receptor Bcr Leading To Generation Of Second Messengers': 'BCR → 2nd messengers',
    'Assembly Of Collagen Fibrils And Other Multimeric Structures': 'Collagen assembly',
    'Signaling By The B Cell Receptor Bcr': 'BCR signaling',
    'Cd22 Mediated Bcr Regulation': 'CD22 BCR reg',
    'Homology Directed Repair': 'HDR',
    'Mismatch Repair': 'MMR',
    'Mitotic Prometaphase': 'Mitotic prometaphase',
    'Mitotic Metaphase And Anaphase': 'Metaphase-anaphase',
    'Mitotic Spindle Checkpoint': 'Spindle checkpoint',
    'Mitotic G1 Phase And G1 S Transition': 'G1-S transition',
    'Mitotic G2 G2 M Phases': 'G2/M phase',
    'Interferon Alpha Beta Signaling': 'IFN α/β signaling',
    'Ecm Proteoglycans': 'ECM proteoglycans',
    'Elastic Fibre Formation': 'Elastic fibre',
    'Collagen Degradation': 'Collagen degradation',
    'Collagen Formation': 'Collagen formation',
    'Integrin Cell Surface Interactions': 'Integrin–cell surface',
    'Tcr Signaling': 'TCR signaling',
}
def short_name(p):
    s = (p.replace('HALLMARK_', '').replace('REACTOME_', '')
         .replace('_', ' ').title())
    return SHORT_MAP.get(s, s)

df['short'] = df.pathway.map(short_name)

NES_EMPH = 2.0
PADJ_STRICT = 0.05
PADJ_LENIENT = 0.25
print(f"N={len(df)} pathways, padj<{PADJ_STRICT}: {(df.padj<PADJ_STRICT).sum()}, "
      f"padj<{PADJ_LENIENT}: {(df.padj<PADJ_LENIENT).sum()}")


# ---------------------------------------------------------------------------
# Low-level PPT primitives (same as 35_fig4D_volcano)
# ---------------------------------------------------------------------------
def new_slide(prs, w, h):
    prs.slide_width = w; prs.slide_height = h
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, size=8, bold=False,
             color=INK, align="left", anchor="middle", italic=False,
             rotation=0):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0;  tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {"top":MSO_ANCHOR.TOP,"middle":MSO_ANCHOR.MIDDLE,
                          "bottom":MSO_ANCHOR.BOTTOM}.get(anchor, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = {"left":PP_ALIGN.LEFT,"right":PP_ALIGN.RIGHT,
                   "center":PP_ALIGN.CENTER}.get(align, PP_ALIGN.LEFT)
    r = p.add_run(); r.text = str(text)
    r.font.name = FONT; r.font.size = Pt(size)
    r.font.bold = bool(bold); r.font.italic = bool(italic)
    r.font.color.rgb = color
    if rotation: tb.rotation = rotation
    kill_shadow(tb)
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.5, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color; c.line.width = Pt(width)
    if dashed: c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    kill_shadow(c)
    return c


def add_circle(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx - r), _i(cy - r),
                                 _i(2 * r), _i(2 * r))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None: shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_hollow_circle(slide, cx, cy, r, line_color, line_width=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx - r), _i(cy - r),
                                 _i(2 * r), _i(2 * r))
    shp.fill.background()
    shp.line.color.rgb = line_color
    shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_triangle_arrow(slide, x1, y1, x2, y2, color, lw=1.8):
    add_line(slide, x1, y1, x2, y2, color=color, width=lw)
    dx, dy = x2 - x1, y2 - y1
    ang = np.degrees(np.arctan2(dy, dx))
    head = Emu(60000)
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                 _i(x2 - head / 2), _i(y2 - head / 2),
                                 _i(head), _i(head))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.rotation = ang
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


# ---------------------------------------------------------------------------
# Slide builder
# ---------------------------------------------------------------------------
def build_slide(prs, slide_w, slide_h, *,
                ax_x, ax_y, ax_w, ax_h,
                label_top_n, variant="main", legend_cols=2):
    s = new_slide(prs, slide_w, slide_h)

    x_lo, x_hi = -3.2, 3.2
    y_lo, y_hi = 0.0, max(26, float(df['-log10padj'].max()) * 1.04)

    def tx(v): return _i(ax_x + (v - x_lo) / (x_hi - x_lo) * ax_w)
    def ty(v): return _i(ax_y + ax_h - (v - y_lo) / (y_hi - y_lo) * ax_h)

    # spines
    add_line(s, ax_x, ax_y + ax_h, ax_x + ax_w, ax_y + ax_h, LINE, 0.6)
    add_line(s, ax_x, ax_y,        ax_x,        ax_y + ax_h, LINE, 0.6)

    # x ticks
    for v in range(-3, 4):
        xx = tx(v)
        add_line(s, xx, ax_y + ax_h, xx, ax_y + ax_h + Inches(0.05),
                 LINE, 0.4)
        add_text(s, xx - Inches(0.20), ax_y + ax_h + Inches(0.06),
                 Inches(0.40), Inches(0.18),
                 f"{v:+d}" if v != 0 else "0",
                 size=7, color=INK, align="center")
    # y ticks — every 5 units
    for v in range(0, int(y_hi) + 1, 5):
        yy = ty(v)
        add_line(s, ax_x - Inches(0.05), yy, ax_x, yy, LINE, 0.4)
        add_text(s, ax_x - Inches(0.45), yy - Inches(0.09),
                 Inches(0.38), Inches(0.18),
                 f"{v}", size=7, color=INK, align="right")

    # axis labels
    add_text(s, ax_x, ax_y + ax_h + Inches(0.28),
             ax_w, Inches(0.22),
             "NES (normalised enrichment score)",
             size=9, bold=True, color=INK, align="center")
    yl_W = Inches(1.40); yl_H = Inches(0.22)
    yl_cx = ax_x - Inches(0.55); yl_cy = ax_y + ax_h / 2
    lb = add_text(s, yl_cx - yl_W / 2, yl_cy - yl_H / 2, yl_W, yl_H,
                  "−log₁₀(BH padj)",
                  size=9, bold=True, color=INK, align="center")
    lb.rotation = -90

    # zero NES vertical + padj thresholds
    add_line(s, tx(0), ax_y, tx(0), ax_y + ax_h, INK, 0.6)
    add_line(s, ax_x, ty(-np.log10(PADJ_STRICT)),
             ax_x + ax_w, ty(-np.log10(PADJ_STRICT)),
             GREY, 0.6, dashed=True)
    add_text(s, ax_x + Inches(0.08), ty(-np.log10(PADJ_STRICT)) - Inches(0.20),
             Inches(1.0), Inches(0.18),
             f"BH padj = {PADJ_STRICT}",
             size=7, color=GREY, italic=True, align="left")
    add_line(s, ax_x, ty(-np.log10(PADJ_LENIENT)),
             ax_x + ax_w, ty(-np.log10(PADJ_LENIENT)),
             LT_GREY, 0.4, dashed=True)
    add_text(s, ax_x + Inches(0.08), ty(-np.log10(PADJ_LENIENT)) - Inches(0.20),
             Inches(1.0), Inches(0.18),
             f"BH padj = {PADJ_LENIENT}",
             size=7, color=LT_GREY, italic=True, align="left")

    # dots: draw non-sig first (faded), then sig (vivid)
    ns_r  = Emu(36000) if variant == "main" else Emu(30000)
    sig_r = Emu(56000) if variant == "main" else Emu(46000)
    emph_r = Emu(90000) if variant == "main" else Emu(72000)

    # non-sig
    for _, r in df[df['padj'] >= PADJ_LENIENT].iterrows():
        add_circle(s, tx(r.NES), ty(r['-log10padj']), ns_r,
                   fill=CAT_COLOR[r['cat']], line_color=None)
    # sig
    for _, r in df[df['padj'] < PADJ_LENIENT].iterrows():
        add_circle(s, tx(r.NES), ty(r['-log10padj']), sig_r,
                   fill=CAT_COLOR[r['cat']],
                   line_color=INK, line_width=0.4)
    # |NES| >= 2 emphasis ring
    for _, r in df[df['NES'].abs() >= NES_EMPH].iterrows():
        add_hollow_circle(s, tx(r.NES), ty(r['-log10padj']),
                          emph_r, HIGHLIGHT, 1.1)

    # pathway labels — top by -log10(padj) * |NES| among padj<lenient
    sig_df = df[df['padj'] < PADJ_LENIENT].copy()
    sig_df['score'] = sig_df['-log10padj'] * sig_df['NES'].abs()
    to_label = sig_df.sort_values('score', ascending=False).head(label_top_n)
    rng = np.random.default_rng(3)
    for _, r in to_label.iterrows():
        is_right = r.NES > 0
        color_ = CAT_COLOR[r['cat']]
        box_w = Inches(1.7)
        x_text = (tx(r.NES) + Inches(0.08)) if is_right else (tx(r.NES) - box_w - Inches(0.08))
        y_off = Inches(-0.08) + Emu(int(rng.uniform(-45000, 45000)))
        add_text(s, x_text, ty(r['-log10padj']) + y_off,
                 box_w, Inches(0.16),
                 r['short'],
                 size=7.5 if variant == "main" else 8,
                 bold=True, color=color_,
                 align="left" if is_right else "right")

    # direction arrows top
    y_arrow = ty(y_hi * 0.96)
    add_triangle_arrow(s, tx(0.6), y_arrow, tx(x_hi * 0.92), y_arrow,
                       GOOD, lw=2.2)
    add_text(s, tx(0.6), y_arrow - Inches(0.30),
             tx(x_hi * 0.92) - tx(0.6), Inches(0.22),
             "Enriched in Good responders",
             size=9, bold=True, color=GOOD, align="center")
    add_triangle_arrow(s, tx(-0.6), y_arrow, tx(x_lo * 0.92), y_arrow,
                       BAD, lw=2.2)
    add_text(s, tx(x_lo * 0.92), y_arrow - Inches(0.30),
             tx(-0.6) - tx(x_lo * 0.92), Inches(0.22),
             "Enriched in Poor responders",
             size=9, bold=True, color=BAD, align="center")

    # legend below x-axis title
    leg_y = ax_y + ax_h + Inches(0.62)
    entries = [
        (CAT_COLOR['Cell cycle'],  None,      0.26, "Cell cycle"),
        (CAT_COLOR['DNA repair'],  None,      0.26, "DNA repair"),
        (CAT_COLOR['Immune'],      None,      0.26, "Immune"),
        (CAT_COLOR['EMT/Stromal'], None,      0.26, "EMT / Stromal"),
        (CAT_COLOR['Other'],       None,      0.22, "Other Hallmark"),
        (None,                     HIGHLIGHT, 0.30, f"|NES| ≥ {NES_EMPH:.0f}"),
    ]
    row_gap = Inches(0.28)
    entry_w = ax_w / legend_cols
    for idx, (fill_, edge_, r_in, label) in enumerate(entries):
        r_i = idx // legend_cols
        c_i = idx % legend_cols
        cx = ax_x + entry_w * c_i + Inches(0.15)
        cy = leg_y + r_i * row_gap + Inches(0.09)
        if fill_ is not None:
            add_circle(s, cx, cy, Emu(int(r_in * 100000)),
                       fill=fill_, line_color=INK, line_width=0.3)
        else:
            add_hollow_circle(s, cx, cy, Emu(int(r_in * 100000)),
                              edge_, line_width=1.0)
        add_text(s, cx + Inches(0.15), leg_y + r_i * row_gap,
                 entry_w - Inches(0.30), Inches(0.18),
                 label, size=7.5, color=INK, align="left")


# ---------------------------------------------------------------------------
# MAIN deck
# ---------------------------------------------------------------------------
prs_main = Presentation()
build_slide(
    prs_main,
    slide_w=Inches(8.5), slide_h=Inches(7.0),
    ax_x=Inches(1.30), ax_y=Inches(0.45),
    ax_w=Inches(6.80), ax_h=Inches(4.80),
    label_top_n=15,
    variant="main",
    legend_cols=3,
)
out_main = OUT / "Fig4D_pathway_volcano_native_editable.pptx"
prs_main.save(out_main)
print(f"wrote MAIN → {out_main}")


# ---------------------------------------------------------------------------
# SUPP deck (wider, more labels)
# ---------------------------------------------------------------------------
prs_supp = Presentation()
build_slide(
    prs_supp,
    slide_w=Inches(11.5), slide_h=Inches(7.8),
    ax_x=Inches(1.45), ax_y=Inches(0.55),
    ax_w=Inches(9.40), ax_h=Inches(5.90),
    label_top_n=35,
    variant="supp",
    legend_cols=6,
)
out_supp = OUT / "SuppFig_pathway_volcano_detailed_native_editable.pptx"
prs_supp.save(out_supp)
print(f"wrote SUPP → {out_supp}")


# count shapes
from pptx import Presentation as _Pr
for path in (out_main, out_supp):
    pr = _Pr(str(path))
    nshape = sum(len(list(s.shapes)) for s in pr.slides)
    print(f"{path.name}: {nshape} shapes across {len(pr.slides)} slide(s)")
