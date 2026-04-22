#!/usr/bin/env python3
"""
35_fig4D_volcano_native_pptx.py

Build Figure 3C / Figure 4D (DESeq2 good-vs-poor DEG volcano) as a
native-editable PowerPoint deck — main + supplementary variants.

User directives (2026-04-22):
  0. Motif-informed from Nature / Cell / Science volcano-style panels:
       - Mariathasan 2018 Nature (IMvigor210): teal/coral direction
         coloring, top-gene text labels with adjustText-style nudging,
         dashed P-threshold line + italic "P=xx" tag at the left margin.
       - Sade-Feldman 2018 Cell: direction-arrow headers ("Up in X
         responders"), outlined marker edges, density hexbin underlay
         for non-significant genes.
       - Love et al DESeq2 / EnhancedVolcano (Blighe Bioconductor,
         widely cited in Nat Commun / Cell Rep): quadruple classifier
         style — (ns, |FC|-only, P-only, sig+|FC|), with gold ring
         emphasis for the high-effect hits.
       - Litchfield 2021 Cell (pan-IO meta): gold outer ring on the
         most biologically striking hits; 0.5-pt minimalist spines.
       - Lu 2021 Nature: per-quadrant running counts as subtle
         in-panel annotations (supp version only).
  1. One panel per slide.
  2. No figure/panel title inside the plot.
  3. python-pptx native shapes/connectors/textboxes (double-click
     editable, no grouping).
  4. Main + supplementary decks in 260418_add/ppt/.
  5. Arial throughout.
  6. Shadows killed via kill_shadow() empty <a:effectLst/> injection.

Practical note on the 22,838 gene point cloud: rendering every gene as
a native OVAL shape would bloat the deck past PowerPoint's responsive
editing threshold. Standard publication-PPT convention (IMvigor, CIT,
NatGen) is to render the non-significant density as a downsampled dot
cloud (here ~500 points for main, ~1200 for supp) while keeping every
biologically-relevant significant gene (n = 228) as a first-class
native shape; gene name labels for the top 30 (main) / 60 (supp) are
individual add_textbox elements.

Outputs
-------
  ppt/Fig4D_volcano_native_editable.pptx              (main)
  ppt/SuppFig_volcano_detailed_native_editable.pptx   (supp)
"""

import os
from pathlib import Path

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Shadow kill helper — injects an empty <a:effectLst/> to override theme
# ---------------------------------------------------------------------------
def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def _i(v):
    return int(v)


# ---------------------------------------------------------------------------
# Project palette — DEEP variant (fig2_v31 / fig3_v32 / panels_v3)
# ---------------------------------------------------------------------------
GOOD       = RGBColor(0x0A, 0x7D, 0x6E)
BAD        = RGBColor(0xC5, 0x3E, 0x1F)
INK        = RGBColor(0x0E, 0x2A, 0x47)
LINE       = RGBColor(0x33, 0x33, 0x33)
GREY       = RGBColor(0x77, 0x77, 0x77)
LT_GREY    = RGBColor(0xC0, 0xC6, 0xCF)
NS_CLOUD   = RGBColor(0xDD, 0xE2, 0xE8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT  = RGBColor(0xD4, 0xA3, 0x00)

FONT = "Arial"

OUT = Path("/data/data/TNT/analysis/260418_add/ppt")
DATA_ROOT = Path("/data/data/TNT/analysis")
OUT.mkdir(parents=True, exist_ok=True)


# ---------------------------------------------------------------------------
# Load DEG data + define axis bounds
# ---------------------------------------------------------------------------
deg = pd.read_csv(DATA_ROOT / "05_rna_deg_gsea/DEG_good_vs_bad_pre.tsv", sep="\t")
deg['-log10p'] = -np.log10(deg['pvalue'].replace(0, 1e-300))
deg = deg.dropna(subset=['log2FoldChange', '-log10p']).copy()
deg['log2FC_clip'] = deg['log2FoldChange'].clip(-5, 5)
deg['nlp_clip'] = deg['-log10p'].clip(0, 8)

P_THRESH = 0.01
FC_THRESH = 1.0
NLP_THRESH = -np.log10(P_THRESH)  # 2.0

sig_mask = deg['pvalue'] < P_THRESH
sig_df = deg[sig_mask].copy()
ns_df = deg[~sig_mask].copy()

up_good = sig_df[sig_df['log2FoldChange'] > 0]
up_poor = sig_df[sig_df['log2FoldChange'] < 0]
fc_big = sig_df[sig_df['log2FoldChange'].abs() >= FC_THRESH]

print(f"DEG total={len(deg):,}  sig={len(sig_df)} "
      f"(good-up {len(up_good)} / poor-up {len(up_poor)})  "
      f"|log2FC|>={FC_THRESH}: {len(fc_big)}")


# ---------------------------------------------------------------------------
# Low-level PPT primitives (copied from 33_fig4D_forest / 18_fig5 / 21_fig9)
# ---------------------------------------------------------------------------
def new_slide(prs, w, h):
    prs.slide_width = w
    prs.slide_height = h
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def add_text(slide, x, y, w, h, text, size=8, bold=False,
             color=INK, align="left", anchor="middle", font=FONT,
             italic=False, rotation=0):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top  = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(anchor, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT,
                   "center": PP_ALIGN.CENTER}.get(align, PP_ALIGN.LEFT)
    r = p.add_run()
    r.text = str(text)
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.italic = bool(italic)
    r.font.color.rgb = color
    if rotation:
        tb.rotation = rotation
    kill_shadow(tb)
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.5, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dashed:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    kill_shadow(c)
    return c


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0.4):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 _i(x), _i(y), _i(w), _i(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_circle(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx - r), _i(cy - r),
                                 _i(2 * r), _i(2 * r))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_triangle_arrow(slide, x1, y1, x2, y2, color, lw=1.5):
    """Thin arrow using a line + a small rotated RIGHT_TRIANGLE head.
    Avoids raw XML <a:tailEnd/> that corrupts PowerPoint compat."""
    add_line(slide, x1, y1, x2, y2, color=color, width=lw)
    dx, dy = x2 - x1, y2 - y1
    L = max(1e-6, (dx * dx + dy * dy) ** 0.5)
    ang = np.degrees(np.arctan2(dy, dx))
    head = Emu(60000)
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                 _i(x2 - head / 2), _i(y2 - head / 2),
                                 _i(head), _i(head))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.rotation = ang
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


# ---------------------------------------------------------------------------
# Core slide builder (reused for main and supp with different sizing)
# ---------------------------------------------------------------------------
def build_volcano_slide(prs, slide_w, slide_h, *,
                        ax_x, ax_y, ax_w, ax_h,
                        ns_sample_n, label_top_n,
                        show_fdr=False, show_quadrant_counts=False,
                        variant="main", legend_cols=4):
    s = new_slide(prs, slide_w, slide_h)

    x_lo, x_hi = -5.2, 5.2
    y_lo, y_hi = 0.0, 8.5

    def tx(v):
        return _i(ax_x + (v - x_lo) / (x_hi - x_lo) * ax_w)

    def ty(v):
        return _i(ax_y + ax_h - (v - y_lo) / (y_hi - y_lo) * ax_h)

    # ===== spines =====
    add_line(s, ax_x, ax_y + ax_h, ax_x + ax_w, ax_y + ax_h, LINE, 0.6)
    add_line(s, ax_x, ax_y,        ax_x,        ax_y + ax_h, LINE, 0.6)

    # ===== x ticks =====
    for v in range(-5, 6):
        xx = tx(v)
        add_line(s, xx, ax_y + ax_h, xx, ax_y + ax_h + Inches(0.05),
                 LINE, 0.4)
        add_text(s, xx - Inches(0.20), ax_y + ax_h + Inches(0.06),
                 Inches(0.40), Inches(0.18),
                 f"{v:+d}" if v != 0 else "0",
                 size=7, color=INK, align="center")
    # ===== y ticks =====
    for v in range(0, 9):
        yy = ty(v)
        add_line(s, ax_x - Inches(0.05), yy, ax_x, yy, LINE, 0.4)
        add_text(s, ax_x - Inches(0.45), yy - Inches(0.09),
                 Inches(0.38), Inches(0.18),
                 f"{v}", size=7, color=INK, align="right")

    # ===== axis labels =====
    add_text(s, ax_x, ax_y + ax_h + Inches(0.28),
             ax_w, Inches(0.22),
             "log₂ fold change (good vs poor)",
             size=9, bold=True, color=INK, align="center")
    # y-axis label: rotated textbox. Pre-rotation (W, H); after -90°
    # the shape occupies (H, W) in screen space around its own centre.
    yl_W = Inches(1.40)       # text-direction span (enough for label)
    yl_H = Inches(0.22)       # text height
    yl_cx = ax_x - Inches(0.55)
    yl_cy = ax_y + ax_h / 2
    lb = add_text(s, yl_cx - yl_W / 2, yl_cy - yl_H / 2,
                  yl_W, yl_H,
                  "−log₁₀(p-value)",
                  size=9, bold=True, color=INK, align="center")
    lb.rotation = -90

    # ===== reference lines =====
    add_line(s, tx(0), ax_y, tx(0), ax_y + ax_h, INK, 0.6)
    add_line(s, ax_x, ty(NLP_THRESH), ax_x + ax_w, ty(NLP_THRESH),
             GREY, 0.6, dashed=True)
    add_text(s, ax_x + Inches(0.05), ty(NLP_THRESH) - Inches(0.20),
             Inches(0.9), Inches(0.18),
             f"P = {P_THRESH}",
             size=7, color=GREY, italic=True, align="left")

    # optional FDR line (supp only)
    if show_fdr:
        padj_min = deg['padj'].dropna().min()
        if np.isfinite(padj_min):
            # draw q=0.1 horizontal only if any gene reaches it
            nlp_fdr = None
            qcut = deg[deg['padj'].fillna(1) < 0.1]
            if len(qcut):
                nlp_fdr = qcut['-log10p'].min()
            if nlp_fdr and nlp_fdr < y_hi:
                add_line(s, ax_x, ty(nlp_fdr),
                         ax_x + ax_w, ty(nlp_fdr),
                         HIGHLIGHT, 0.6, dashed=True)
                add_text(s, ax_x + Inches(0.05),
                         ty(nlp_fdr) - Inches(0.20),
                         Inches(0.9), Inches(0.18),
                         "BH q = 0.1",
                         size=7, color=HIGHLIGHT, italic=True,
                         align="left")
        else:
            # explicitly annotate that no gene reaches q<0.1 (paper's stance)
            add_text(s, ax_x + ax_w - Inches(2.6),
                     ax_y + Inches(0.10),
                     Inches(2.5), Inches(0.18),
                     "No gene reaches BH q < 0.1",
                     size=7, color=GREY, italic=True, align="right")

    # ===== non-significant cloud =====
    rng = np.random.default_rng(0)
    ns_sample = ns_df.sample(min(ns_sample_n, len(ns_df)),
                             random_state=42)
    r_ns = Emu(16000) if variant == "main" else Emu(14000)
    for _, g in ns_sample.iterrows():
        add_circle(s, tx(g.log2FC_clip), ty(g.nlp_clip), r_ns,
                   fill=NS_CLOUD, line_color=None)

    # ===== significant genes =====
    r_sig = Emu(28000) if variant == "main" else Emu(22000)
    for _, g in up_good.iterrows():
        add_circle(s, tx(g.log2FC_clip), ty(g.nlp_clip), r_sig,
                   fill=GOOD, line_color=INK, line_width=0.4)
    for _, g in up_poor.iterrows():
        add_circle(s, tx(g.log2FC_clip), ty(g.nlp_clip), r_sig,
                   fill=BAD, line_color=INK, line_width=0.4)

    # ===== |log2FC|>=1 gold emphasis ring =====
    r_ring = Emu(48000) if variant == "main" else Emu(38000)
    for _, g in fc_big.iterrows():
        shp = slide_add_hollow_circle(s, tx(g.log2FC_clip),
                                      ty(g.nlp_clip), r_ring,
                                      HIGHLIGHT, 1.2)

    # ===== gene labels (top by -log10p × |log2FC|) =====
    to_label = (sig_df.assign(score=sig_df['-log10p'] * sig_df['log2FoldChange'].abs())
                      .sort_values("score", ascending=False)
                      .head(label_top_n))
    # deterministic jitter so labels don't collide: radial offset
    rng2 = np.random.default_rng(7)
    for _, g in to_label.iterrows():
        color_ = GOOD if g.log2FoldChange > 0 else BAD
        lx = tx(g.log2FC_clip)
        ly = ty(g.nlp_clip)
        # label to the right for up-good, left for up-poor
        if g.log2FoldChange > 0:
            off_x = Inches(0.10)
            align = "left"
            box_w = Inches(0.85)
            text_x = lx + off_x
        else:
            off_x = Inches(0.10)
            align = "right"
            box_w = Inches(0.85)
            text_x = lx - box_w - off_x
        off_y = Inches(-0.08) + Emu(int(rng2.uniform(-50000, 50000)))
        add_text(s, text_x, ly + off_y,
                 box_w, Inches(0.16),
                 g.gene, size=7 if variant == "main" else 7.5,
                 bold=True, color=color_, align=align)

    # ===== direction arrows at top =====
    y_arrow = ty(7.85)
    add_triangle_arrow(s, tx(0.7), y_arrow, tx(4.6), y_arrow,
                       GOOD, lw=2.2)
    add_text(s, tx(0.7), y_arrow - Inches(0.30),
             tx(4.6) - tx(0.7), Inches(0.20),
             "Up in Good responders",
             size=9, bold=True, color=GOOD, align="center")
    add_triangle_arrow(s, tx(-0.7), y_arrow, tx(-4.6), y_arrow,
                       BAD, lw=2.2)
    add_text(s, tx(-4.6), y_arrow - Inches(0.30),
             tx(-0.7) - tx(-4.6), Inches(0.20),
             "Up in Poor responders",
             size=9, bold=True, color=BAD, align="center")

    # ===== per-quadrant significant-count annotations (supp only) =====
    if show_quadrant_counts:
        add_text(s, tx(3.5), ty(NLP_THRESH + 0.4),
                 Inches(1.5), Inches(0.16),
                 f"n = {len(up_good)}",
                 size=8, bold=True, color=GOOD, align="left")
        add_text(s, tx(-4.5), ty(NLP_THRESH + 0.4),
                 Inches(1.5), Inches(0.16),
                 f"n = {len(up_poor)}",
                 size=8, bold=True, color=BAD, align="left")

    # ===== legend below x-axis title; supp=1 row x 4 cols, main=2 rows x 2 cols =====
    leg_y = ax_y + ax_h + Inches(0.62)
    entries = [
        (GOOD,     None,      0.25,
         f"Up in Good, P < {P_THRESH}  (n = {len(up_good)})"),
        (BAD,      None,      0.25,
         f"Up in Poor, P < {P_THRESH}  (n = {len(up_poor)})"),
        (None,     HIGHLIGHT, 0.30,
         f"|log₂ FC| ≥ {FC_THRESH:.0f}  (n = {len(fc_big)})"),
        (NS_CLOUD, None,      0.20,
         f"Non-significant (P ≥ {P_THRESH})"),
    ]
    n_rows = (len(entries) + legend_cols - 1) // legend_cols
    entry_w = ax_w / legend_cols
    row_gap = Inches(0.28)
    for idx, (fill_, edge_, r_in, label) in enumerate(entries):
        r_i = idx // legend_cols
        c_i = idx % legend_cols
        cx = ax_x + entry_w * c_i + Inches(0.15)
        cy = leg_y + r_i * row_gap + Inches(0.09)
        if fill_ is not None:
            add_circle(s, cx, cy, Emu(int(r_in * 100000)),
                       fill=fill_, line_color=INK, line_width=0.3)
        else:
            add_circle(s, cx, cy, Emu(int(r_in * 100000)),
                       fill=WHITE, line_color=edge_, line_width=1.0)
        add_text(s, cx + Inches(0.15), leg_y + r_i * row_gap,
                 entry_w - Inches(0.30), Inches(0.18),
                 label, size=7.5, color=INK, align="left")


def slide_add_hollow_circle(slide, cx, cy, r, line_color, line_width=1.0):
    """Hollow ring without fill (emphasis ring for |log2FC|>=1)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx - r), _i(cy - r),
                                 _i(2 * r), _i(2 * r))
    shp.fill.background()
    shp.line.color.rgb = line_color
    shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


# ===========================================================================
# MAIN deck
# ===========================================================================
prs_main = Presentation()
build_volcano_slide(
    prs_main,
    slide_w=Inches(8.2), slide_h=Inches(6.5),
    ax_x=Inches(1.30), ax_y=Inches(0.45),
    ax_w=Inches(6.50), ax_h=Inches(4.30),
    ns_sample_n=500,
    label_top_n=30,
    show_fdr=False,
    show_quadrant_counts=False,
    variant="main",
    legend_cols=2,
)
out_main = OUT / "Fig4D_volcano_native_editable.pptx"
prs_main.save(out_main)
print(f"wrote MAIN deck → {out_main}")


# ===========================================================================
# SUPP deck
# ===========================================================================
prs_supp = Presentation()
build_volcano_slide(
    prs_supp,
    slide_w=Inches(11.5), slide_h=Inches(7.8),
    ax_x=Inches(1.45), ax_y=Inches(0.55),
    ax_w=Inches(9.40), ax_h=Inches(6.10),
    ns_sample_n=1200,
    label_top_n=60,
    show_fdr=True,
    show_quadrant_counts=True,
    variant="supp",
    legend_cols=4,
)
out_supp = OUT / "SuppFig_volcano_detailed_native_editable.pptx"
prs_supp.save(out_supp)
print(f"wrote SUPP deck → {out_supp}")


# ===========================================================================
# Verify shape / text counts
# ===========================================================================
from pptx import Presentation as _Pr
for path in (out_main, out_supp):
    pr = _Pr(str(path))
    nshape = sum(len(list(s.shapes)) for s in pr.slides)
    print(f"{path.name}: {nshape} shapes across {len(pr.slides)} slide(s)")
