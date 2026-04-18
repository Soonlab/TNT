#!/usr/bin/env python3
"""
21_fig9_native_pptx.py

Build Figure 9 --- External validation / regimen-agnostic two-thread
meta-analysis (§3.12, the paper's most load-bearing figure) --- and
its supplementary companion Supp Fig S19 as native-editable PowerPoint
decks.

Main Figure 9 panels
--------------------
  A  Horizontal meta-analysis forest plot (headline): 7 signatures
     grouped by thread, per-cohort individual Z values as small
     hollow circles + pooled Stouffer Z as filled diamond; Thread 2
     CD8 shown at two levels (5-cohort and 6-source with Akiyoshi
     2023 paper-level); null line at Z=0; P<0.05 threshold dashed
     reference at |Z|=1.96; right-column text = pooled Z, P and N.

  B  Per-cohort × per-signature effect-size heatmap: rows = 7
     signatures grouped by thread; columns = 5 primary long-course
     nCRT cohorts (+ a highlighted Akiyoshi column for CD8 row);
     cell = signed delta (good − bad) with diverging teal/coral
     ramp; cell text = delta value; right-margin column = pooled Z
     (gold if significant); top margin = N per cohort.

  C  Concordance stacked bar per signature: for each of the 7
     signatures, a horizontal bar with segments encoding how many
     of the 5 primary cohorts show the signature in the predicted
     direction (with P<0.05 in dark teal, as a trend in light teal)
     vs in the opposite direction (dark/light coral); chance line
     at 2.5/5; annotation per signature: pooled Z and summary
     verdict.

  D  Discovery-vs-external concordance scatter: x = discovery
     cohort Stouffer-style Z (from n=33 Mann-Whitney P) for each
     signature; y = external pooled Stouffer Z. Quadrant-based
     interpretation with threshold lines at Z=±1.96 (dashed):
       · upper-right: reproducible biomarker (both sig)
       · upper-left: external-only (discovery under-powered)
       · lower-right: discovery-only (external heterogeneous)
     Each signature labelled in place; shapes coloured by thread.
     Shows the under-powered-immune-discovery finding transparently.

Supp Figure S19 panels
----------------------
  A  9-cohort unrestricted sensitivity meta forest (for transparency
     around the primary-meta 5-cohort restriction). Same forest
     layout as Main A but including GSE119409, GSE94104, GSE46862,
     GSE150082; shows how the direction and significance weaken or
     flip for each excluded cohort, validating their exclusion on
     pre-specified ≥3/4 concordance grounds.

  B  Akiyoshi 2023 alternative-statistic sensitivity analysis for
     the Thread-2 CD8 augmentation: multiple published statistics
     from the JAMA Network Open 2023 paper converted to Z and
     combined with the 5-cohort CD8 meta; bars showing each
     combination and its pooled Z (6-source range 2.90 to 3.60,
     all P<0.004).

Rules (same as scripts 18-20):
  * one panel per slide
  * no plot titles
  * python-pptx native elements only (TEXT_BOX / LINE connector /
    AUTO_SHAPE); Arial throughout
  * TNT palette GOOD=#0a7d6e / BAD=#c53e1f; kill_shadow on every
    shape; all coordinates through _i()

Motif references consulted
--------------------------
  - Rooney et al Cell 2015 (pan-cancer cytolytic activity meta):
    horizontal forest with per-cohort markers and pooled diamond.
  - Bahadoer et al Lancet Oncol 2021 (RAPIDO): Cochrane-style
    subgroup forest with reference line.
  - Chowell et al Science 2018 (HLA-ICB meta): cohort forest +
    pooled pooled estimate with paper-level evidence row.
  - Le et al NEJM 2015 / Overman et al Lancet Oncol 2017 / Andre
    et al NEJM 2020 (MSI-PD1 trials meta): forest with CI bars.
  - Thorsson et al Immunity 2018: signature × cohort heatmap with
    marginal pooled-Z.
  - Newman et al Nat Methods 2015: signature × cohort diverging
    heatmap with pooled row.
"""

import os
import numpy as np
import pandas as pd
from scipy.stats import norm
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Shared infrastructure
# ---------------------------------------------------------------------------
def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def _i(v):
    return int(round(float(v)))


GOOD = RGBColor(0x0A, 0x7D, 0x6E)
BAD = RGBColor(0xC5, 0x3E, 0x1F)
INK = RGBColor(0x22, 0x22, 0x22)
LINE = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0xBB, 0xBB, 0xBB)
LT_GREY = RGBColor(0xDD, 0xDD, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xD4, 0xA3, 0x00)

# Thread colours
THREAD1 = RGBColor(0x0E, 0x4A, 0x68)        # deep navy-teal for Thread 1
THREAD2 = RGBColor(0x8A, 0x2B, 0x4C)        # deep rose for Thread 2
AKIYOSHI = RGBColor(0x6A, 0x4C, 0x93)       # purple for paper-level

# Cohort colours
COHORT_COLOR = {
    "GSE35452": RGBColor(0x2E, 0x86, 0xAB),
    "GSE45404": RGBColor(0xA6, 0x50, 0x2C),
    "GSE56699": RGBColor(0x5A, 0x82, 0x61),
    "GSE133057": RGBColor(0xE3, 0x9A, 0x1E),
    "GSE87211": RGBColor(0x2F, 0x30, 0x80),
    "GSE216616": AKIYOSHI,                      # Akiyoshi
    # excluded cohorts
    "GSE119409": RGBColor(0x88, 0x88, 0x88),
    "GSE94104": RGBColor(0x88, 0x88, 0x88),
    "GSE46862": RGBColor(0x88, 0x88, 0x88),
    "GSE150082": RGBColor(0x88, 0x88, 0x88),
}

GOOD_HEX = (0x0A, 0x7D, 0x6E)
BAD_HEX = (0xC5, 0x3E, 0x1F)

FONT = "Arial"
SLIDE_W = Inches(6.5)
SLIDE_H = Inches(4.5)

OUT = "/data/data/TNT/analysis/260418_add/ppt"
DATA = "/data/data/TNT/analysis/260418_add"
os.makedirs(OUT, exist_ok=True)


def new_slide(prs, w=SLIDE_W, h=SLIDE_H):
    prs.slide_width = w
    prs.slide_height = h
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, size=8, bold=False,
             color=INK, align="left", anchor="middle", font=FONT,
             italic=False):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(max(w, 1)), _i(max(h, 1)))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "bottom": MSO_ANCHOR.BOTTOM,
                          "middle": MSO_ANCHOR.MIDDLE}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = str(text)
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.italic = bool(italic)
    r.font.color.rgb = color
    kill_shadow(tb)
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.5, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        try:
            c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass
    kill_shadow(c)
    return c


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0.5):
    w = max(_i(w), 1); h = max(_i(h), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _i(x), _i(y), w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_circle(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_diamond(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def lighten(rgb_hex, factor=0.7):
    r0, g0, b0 = rgb_hex
    return RGBColor(int(r0 + (255 - r0) * factor),
                    int(g0 + (255 - g0) * factor),
                    int(b0 + (255 - b0) * factor))


def draw_panel_letter(slide, letter):
    add_text(slide, Inches(0.15), Inches(0.1), Inches(0.4), Inches(0.35),
             letter, size=14, bold=True, color=INK, align="left")


# ---------------------------------------------------------------------------
# Load data
# ---------------------------------------------------------------------------
final_meta = pd.read_csv(f"{DATA}/FINAL_meta_with_akiyoshi.tsv", sep="\t")
restricted5 = pd.read_csv(f"{DATA}/restricted5_meta_combined.tsv", sep="\t")
per_cohort = pd.read_csv(f"{DATA}/restricted5_per_cohort_detail.tsv",
                         sep="\t")
thread1_summary = pd.read_csv(f"{DATA}/thread1_per_cohort_summary.tsv",
                              sep="\t")
discovery = pd.read_csv(f"{DATA}/discovery_signature_stats.tsv", sep="\t")

# Primary cohort order + N per cohort
PRIMARY_COHORTS = ["GSE35452", "GSE45404", "GSE56699", "GSE133057", "GSE87211"]
COHORT_N = {}
for gse in PRIMARY_COHORTS:
    row = thread1_summary[thread1_summary.gse == gse]
    if not row.empty:
        COHORT_N[gse] = int(row.iloc[0]["n_samples"])
COHORT_N["GSE216616"] = 298

# Signature ordering for all panels
THREAD1_SIGS = [
    ("DSB_HDR_repair", "DSB / HDR repair", "up"),
    ("Tumor_cellcycle", "Tumor cell-cycle", "up"),
    ("E2F_MYC_cellcycle", "E2F / MYC", "up"),
    ("EMT", "EMT", "down"),
]
THREAD2_SIGS = [
    ("CD8_cytotoxic", "CD8-cytotoxic", "up"),
    ("Tcell_infiltration", "T-cell infiltration", "up"),
    ("Bcell_infiltration", "B-cell infiltration", "up"),
]


def get_pooled(sig, source="5cohort"):
    if source == "5cohort":
        r = restricted5[restricted5.signature == sig]
        if r.empty:
            return None, None, None
        r = r.iloc[0]
        return float(r.Z), float(r.p_meta), int(r.total_n)
    elif source == "6source":
        r = final_meta[final_meta.signature == sig]
        if r.empty:
            return None, None, None
        r = r.iloc[0]
        return float(r.Z), float(r.p_meta), int(r.n_total)
    return None, None, None


def per_cohort_delta(sig, cohort):
    r = per_cohort[(per_cohort.signature == sig)
                   & (per_cohort.gse == cohort)]
    if r.empty:
        return None, None
    r = r.iloc[0]
    return float(r.delta), float(r.pvalue)


def per_cohort_Z(sig, cohort):
    delta, p = per_cohort_delta(sig, cohort)
    if delta is None:
        return None
    sgn = 1 if delta > 0 else (-1 if delta < 0 else 0)
    # EMT predicted direction is "good < bad" (i.e. delta negative in good).
    # Re-orient so positive Z means predicted direction.
    if sig == "EMT":
        sgn = -sgn
    # map two-sided p to z
    if p <= 0 or not np.isfinite(p):
        return None
    z = norm.ppf(1 - p / 2) * sgn
    return float(z)


# ===========================================================================
# MAIN FIGURE 9
# ===========================================================================
prs_main = Presentation()


# -------------------------------------------------------------------
# Panel A  (Slide 1) --- horizontal meta-analysis forest plot
# -------------------------------------------------------------------
def build_A():
    s = new_slide(prs_main)
    draw_panel_letter(s, "A")

    # rows layout (top to bottom):
    #  0 Thread 1 section header
    #  1-4 Thread 1 signatures
    #  5 blank separator
    #  6 Thread 2 section header
    #  7 CD8 5-cohort
    #  8 CD8 6-source (Akiyoshi) -- highlighted
    #  9 Tcell_infiltration
    # 10 Bcell_infiltration

    rows = [
        ("header", "Thread 1 — tumor-intrinsic  (N = 518, 5 long-course nCRT cohorts)", THREAD1),
        ("sig", "DSB_HDR_repair", "DSB / HDR repair", "5cohort", "up"),
        ("sig", "Tumor_cellcycle", "Tumor cell-cycle", "5cohort", "up"),
        ("sig", "E2F_MYC_cellcycle", "E2F / MYC", "5cohort", "up"),
        ("sig", "EMT", "EMT", "5cohort", "down"),
        ("gap", ""),
        ("header", "Thread 2 — immune  (N = 518, 5 cohorts  ·  +Akiyoshi N = 816, 6 sources)", THREAD2),
        ("sig", "CD8_cytotoxic", "CD8-cytotoxic (5 cohorts)", "5cohort", "up"),
        ("sig", "CD8_cytotoxic", "CD8-cytotoxic + Akiyoshi 2023", "6source", "up",
         "★ 6-source augmented"),
        ("sig", "Tcell_infiltration", "T-cell infiltration", "5cohort", "up"),
        ("sig", "Bcell_infiltration", "B-cell infiltration", "5cohort", "up"),
    ]

    # plot area
    ax_x = Inches(2.30)
    ax_y = Inches(0.60)
    ax_w = Inches(3.35)
    ax_h = Inches(3.30)

    # x range: -1 to 4 (Stouffer Z)
    x_lo, x_hi = -1.2, 4.0

    def tx(v):
        return _i(ax_x + (v - x_lo) / (x_hi - x_lo) * ax_w)

    # y positions per row
    n_rows = len(rows)
    row_h = ax_h / n_rows

    def ty(i):
        return _i(ax_y + (i + 0.5) * row_h)

    # spines + x ticks
    add_line(s, ax_x, ax_y, ax_x + ax_w, ax_y, LINE, 0.5)          # top
    add_line(s, ax_x, ax_y + ax_h, ax_x + ax_w,
             ax_y + ax_h, LINE, 0.6)                                # x spine
    for v in [-1, 0, 1, 2, 3, 4]:
        xx = tx(v)
        add_line(s, xx, _i(ax_y + ax_h), xx,
                 _i(ax_y + ax_h + Inches(0.05)), LINE, 0.4)
        add_text(s, xx - Inches(0.15), _i(ax_y + ax_h + Inches(0.06)),
                 Inches(0.30), Inches(0.14),
                 f"{v:+d}" if v != 0 else "0",
                 size=6, color=INK, align="center")
    add_text(s, ax_x, _i(ax_y + ax_h + Inches(0.20)),
             ax_w, Inches(0.15),
             "Stouffer pooled Z (predicted-direction-oriented)",
             size=8, color=INK, align="center")

    # reference lines: Z=0 (solid), Z=±1.96 (dashed at P=0.05)
    add_line(s, tx(0), ax_y, tx(0), ax_y + ax_h, INK, 1.0)
    add_line(s, tx(1.96), ax_y, tx(1.96), ax_y + ax_h,
             GREY, 0.5, dashed=True)
    add_line(s, tx(-1.96), ax_y, tx(-1.96), ax_y + ax_h,
             GREY, 0.5, dashed=True)
    add_text(s, tx(1.96) - Inches(0.20), _i(ax_y + Inches(0.02)),
             Inches(0.4), Inches(0.12),
             "P=0.05", size=5, color=RGBColor(0x99, 0x99, 0x99),
             italic=True, align="center")

    # render rows
    for i, row in enumerate(rows):
        y = ty(i)
        if row[0] == "header":
            # italic header strip
            add_rect(s, Inches(0.15), y - Inches(0.10),
                     ax_x - Inches(0.15) + ax_w, Inches(0.22),
                     fill=RGBColor(0xF2, 0xEC, 0xE1), line_color=None)
            add_text(s, Inches(0.20), y - Inches(0.08),
                     ax_x - Inches(0.25) + ax_w, Inches(0.18),
                     row[1], size=7, bold=True, color=row[2],
                     italic=True, align="left")
            continue
        if row[0] == "gap":
            continue

        # row layout
        _, sig, pretty, source, pred = row[:5]
        highlight = len(row) > 5  # CD8 6-source augmented row

        # signature label (left column)
        label_color = (HIGHLIGHT if highlight
                       else (THREAD1 if sig in [s0 for s0, _, _ in THREAD1_SIGS]
                             else THREAD2))
        star = "  ★" if highlight else ""
        add_text(s, Inches(0.18), y - Inches(0.09),
                 ax_x - Inches(0.20), Inches(0.18),
                 pretty + star,
                 size=7, bold=highlight, color=label_color,
                 align="right")

        # pooled Z and per-cohort markers
        Z, p, N = get_pooled(sig, source=source)
        if Z is None:
            continue

        # per-cohort small circles (only for 5-cohort Z)
        cohorts_here = PRIMARY_COHORTS[:]
        if source == "6source":
            cohorts_here = PRIMARY_COHORTS + ["GSE216616"]

        for gse in cohorts_here:
            if gse == "GSE216616":
                # use Akiyoshi's paper-level Z = +2.81
                pz = 2.81
            else:
                pz = per_cohort_Z(sig, gse)
            if pz is None:
                continue
            pz_clamped = max(x_lo, min(x_hi, pz))
            xx = tx(pz_clamped)
            add_circle(s, xx, y, Emu(18000),
                       fill=WHITE, line_color=COHORT_COLOR[gse],
                       line_width=0.9)

        # pooled diamond
        pooled_x = tx(max(x_lo, min(x_hi, Z)))
        d_color = (HIGHLIGHT if highlight
                   else (THREAD1 if sig in [s0 for s0, _, _ in THREAD1_SIGS]
                         else THREAD2))
        dia_size = Emu(55000) if highlight else Emu(45000)
        add_diamond(s, pooled_x, y, dia_size,
                    fill=d_color, line_color=WHITE, line_width=1.2)

        # right-column text: Z, P, N
        sig_str = f"Z = {Z:+.2f}"
        p_str = (f"P < 10⁻³" if p < 0.001
                 else f"P = {p:.3f}")
        n_str = f"N = {N:,}"

        right_x = ax_x + ax_w + Inches(0.05)
        # mini-star for P<0.05
        is_sig = abs(Z) >= 1.96
        star_count = ("★★★" if p < 0.005 and is_sig
                      else ("★★" if p < 0.05 and is_sig
                            else ("★" if p < 0.1 else "")))
        add_text(s, right_x, y - Inches(0.13),
                 Inches(0.80), Inches(0.13),
                 sig_str, size=7,
                 bold=is_sig,
                 color=d_color if is_sig else RGBColor(0x66, 0x66, 0x66),
                 align="left")
        add_text(s, right_x, y - Inches(0.01),
                 Inches(0.80), Inches(0.13),
                 p_str, size=6,
                 color=d_color if is_sig else RGBColor(0x66, 0x66, 0x66),
                 align="left")
        # star markers next to P
        add_text(s, right_x + Inches(0.45), y - Inches(0.13),
                 Inches(0.30), Inches(0.13),
                 star_count, size=8, bold=True, color=HIGHLIGHT,
                 align="left")

    # cohort legend (bottom)
    leg_y = Inches(4.08)
    add_text(s, Inches(0.15), leg_y - Inches(0.03),
             Inches(0.9), Inches(0.14),
             "Per-cohort Z:", size=6, bold=True, color=INK, align="left")
    cohort_order = PRIMARY_COHORTS + ["GSE216616"]
    cohort_labels = {
        "GSE35452": "GSE35452 (46)",
        "GSE45404": "GSE45404 (80)",
        "GSE56699": "GSE56699 (44)",
        "GSE133057": "GSE133057 (33)",
        "GSE87211": "GSE87211 (353)",
        "GSE216616": "GSE216616 Akiyoshi '23 (298)",
    }
    xb = Inches(0.90)
    for gse in cohort_order:
        add_circle(s, xb, leg_y + Inches(0.05), Emu(15000),
                   fill=WHITE, line_color=COHORT_COLOR[gse], line_width=0.8)
        add_text(s, xb + Inches(0.05), leg_y - Inches(0.03),
                 Inches(1.20), Inches(0.16),
                 cohort_labels[gse], size=5, color=INK, align="left")
        xb += Inches(0.95)

    # pooled diamond legend
    add_diamond(s, Inches(0.90), leg_y + Inches(0.22), Emu(35000),
                fill=THREAD1, line_color=WHITE, line_width=1.0)
    add_text(s, Inches(1.00), leg_y + Inches(0.14),
             Inches(1.3), Inches(0.16),
             "pooled Stouffer Z (T1)",
             size=6, color=INK, align="left")
    add_diamond(s, Inches(2.25), leg_y + Inches(0.22), Emu(35000),
                fill=THREAD2, line_color=WHITE, line_width=1.0)
    add_text(s, Inches(2.35), leg_y + Inches(0.14),
             Inches(1.3), Inches(0.16),
             "pooled Z (T2)",
             size=6, color=INK, align="left")
    add_diamond(s, Inches(3.30), leg_y + Inches(0.22), Emu(45000),
                fill=HIGHLIGHT, line_color=WHITE, line_width=1.0)
    add_text(s, Inches(3.45), leg_y + Inches(0.14),
             Inches(2.3), Inches(0.16),
             "Akiyoshi-augmented",
             size=6, color=INK, align="left")


# -------------------------------------------------------------------
# Panel B  (Slide 2) --- per-cohort x per-signature heatmap
# -------------------------------------------------------------------
def build_B():
    s = new_slide(prs_main)
    draw_panel_letter(s, "B")

    signatures = THREAD1_SIGS + THREAD2_SIGS
    n_rows = len(signatures)
    cohort_cols = PRIMARY_COHORTS[:]

    hm_x = Inches(2.60); hm_y = Inches(0.80)
    hm_w = Inches(2.25); hm_h = Inches(3.00)
    cell_w = hm_w / len(cohort_cols)
    cell_h = hm_h / n_rows

    # diverging ramp: teal (+ in predicted direction) to coral (- opposite)
    def ramp(v, predicted):
        # normalize v into [-1, 1] with saturation at 0.6
        sat = 0.6
        vv = max(-1, min(1, float(v) / sat))
        if predicted == "down":
            vv = -vv  # re-orient so positive = predicted direction
        if vv >= 0:
            r = int(255 - vv * (255 - GOOD_HEX[0]))
            g = int(255 - vv * (255 - GOOD_HEX[1]))
            b = int(255 - vv * (255 - GOOD_HEX[2]))
        else:
            t = -vv
            r = int(255 - t * (255 - BAD_HEX[0]))
            g = int(255 - t * (255 - BAD_HEX[1]))
            b = int(255 - t * (255 - BAD_HEX[2]))
        return RGBColor(r, g, b)

    # render cells
    for i, (sig, pretty, pred) in enumerate(signatures):
        for j, gse in enumerate(cohort_cols):
            delta, p = per_cohort_delta(sig, gse)
            if delta is None:
                delta = 0; p = 1.0
            x = _i(hm_x + j * cell_w)
            y = _i(hm_y + i * cell_h)
            add_rect(s, x, y, _i(cell_w), _i(cell_h),
                     fill=ramp(delta, pred),
                     line_color=WHITE, line_width=0.5)
            # delta value text
            add_text(s, x, y + _i(cell_h) / 2 - Inches(0.08),
                     _i(cell_w), Inches(0.14),
                     f"{delta:+.2f}", size=6,
                     bold=(p < 0.05),
                     color=INK if abs(delta) < 0.4 else WHITE,
                     align="center")
            # significance asterisk (bottom-right of cell)
            if p < 0.05:
                add_text(s, x + _i(cell_w) - Inches(0.16),
                         y + _i(cell_h) - Inches(0.14),
                         Inches(0.15), Inches(0.12),
                         "*", size=8, bold=True, color=INK,
                         align="right")

    # row labels
    for i, (sig, pretty, pred) in enumerate(signatures):
        y = _i(hm_y + (i + 0.5) * cell_h)
        arrow = "↓" if pred == "down" else "↑"
        thread_c = THREAD1 if sig in [s0 for s0, _, _ in THREAD1_SIGS] else THREAD2
        add_text(s, Inches(0.20), y - Inches(0.08),
                 Inches(2.35), Inches(0.16),
                 f"{pretty} {arrow}",
                 size=7, color=thread_c, align="right")

    # column labels (cohort name) + N sub-label
    for j, gse in enumerate(cohort_cols):
        x = _i(hm_x + (j + 0.5) * cell_w)
        add_text(s, x - Inches(0.45), _i(hm_y - Inches(0.35)),
                 Inches(0.9), Inches(0.16),
                 gse, size=6, bold=True, color=COHORT_COLOR[gse],
                 align="center")
        n = COHORT_N.get(gse, 0)
        add_text(s, x - Inches(0.4), _i(hm_y - Inches(0.20)),
                 Inches(0.8), Inches(0.14),
                 f"N = {n}", size=5,
                 color=RGBColor(0x66, 0x66, 0x66),
                 align="center")

    # right margin: pooled Z
    rZ_x = hm_x + hm_w + Inches(0.15)
    add_text(s, rZ_x, _i(hm_y - Inches(0.35)),
             Inches(1.4), Inches(0.16),
             "Pooled Z (5-cohort)", size=6, bold=True, color=INK,
             align="left")
    add_text(s, rZ_x, _i(hm_y - Inches(0.20)),
             Inches(1.4), Inches(0.14),
             "(5-cohort N = 518)", size=5,
             color=RGBColor(0x66, 0x66, 0x66), align="left")

    for i, (sig, pretty, pred) in enumerate(signatures):
        y = _i(hm_y + (i + 0.5) * cell_h)
        Z, p, _ = get_pooled(sig, source="5cohort")
        is_sig = abs(Z) >= 1.96
        label = f"Z = {Z:+.2f}"
        sub_label = (f"P < 10⁻³" if p < 0.001 else f"P = {p:.3f}")
        add_text(s, rZ_x, y - Inches(0.12),
                 Inches(1.0), Inches(0.13),
                 label,
                 size=7, bold=is_sig,
                 color=HIGHLIGHT if is_sig else RGBColor(0x66, 0x66, 0x66),
                 align="left")
        add_text(s, rZ_x, y + Inches(0.01),
                 Inches(1.0), Inches(0.13),
                 sub_label, size=5,
                 color=HIGHLIGHT if is_sig else RGBColor(0x66, 0x66, 0x66),
                 align="left")

    # colour bar
    cb_x = hm_x; cb_y = hm_y + hm_h + Inches(0.22)
    cb_w = hm_w; cb_h = Inches(0.14)
    for k in range(40):
        t = -1 + 2 * k / 39
        add_rect(s, _i(cb_x + k * cb_w / 40), cb_y,
                 _i(cb_w / 40) + Emu(2000), cb_h,
                 fill=ramp(t * 0.6, "up"))
    add_rect(s, cb_x, cb_y, cb_w, cb_h, line_color=LINE, line_width=0.3)
    for v, lab in [(-0.6, "−0.6"), (0, "0"), (0.6, "+0.6")]:
        xx = _i(cb_x + (v + 0.6) / 1.2 * cb_w)
        add_text(s, xx - Inches(0.18), cb_y + cb_h + Emu(10000),
                 Inches(0.36), Inches(0.14),
                 lab, size=5, color=INK, align="center")
    add_text(s, cb_x, cb_y - Inches(0.14),
             cb_w, Inches(0.12),
             "Δ (good − bad) re-oriented by predicted direction",
             size=6, color=INK, align="center")

    # caption
    add_text(s, Inches(0.15), Inches(4.20),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "* = per-cohort Mann-Whitney P < 0.05. "
             "Pooled Z values gold-highlighted when |Z| ≥ 1.96 (P < 0.05).",
             size=6, color=RGBColor(0x55, 0x55, 0x55), align="center")


# -------------------------------------------------------------------
# Panel C  (Slide 3) --- concordance stacked bar
# -------------------------------------------------------------------
def build_C():
    s = new_slide(prs_main)
    draw_panel_letter(s, "C")

    signatures = THREAD1_SIGS + THREAD2_SIGS
    n_rows = len(signatures)

    # plot area
    bx = Inches(2.05); by = Inches(0.75)
    bw = Inches(3.30); bh = Inches(3.10)
    row_h = bh / n_rows

    # x-range: 0 to 5 cohorts
    def tx(v):
        return _i(bx + v / 5 * bw)

    # spines + ticks
    add_line(s, bx, by + bh, bx + bw, by + bh, LINE, 0.6)
    for v in range(6):
        add_line(s, tx(v), _i(by + bh), tx(v),
                 _i(by + bh + Inches(0.04)), LINE, 0.4)
        add_text(s, tx(v) - Inches(0.12), _i(by + bh + Inches(0.05)),
                 Inches(0.24), Inches(0.14),
                 f"{v}", size=7, color=INK, align="center")
    add_text(s, bx, _i(by + bh + Inches(0.20)),
             bw, Inches(0.15),
             "Number of primary cohorts (of 5)",
             size=8, color=INK, align="center")

    # chance line at 2.5
    add_line(s, tx(2.5), by, tx(2.5), by + bh, INK, 0.8, dashed=True)
    add_text(s, tx(2.5) - Inches(0.35), _i(by + Inches(0.02)),
             Inches(0.7), Inches(0.13),
             "chance (2.5)", size=5,
             color=RGBColor(0x88, 0x88, 0x88),
             align="center", italic=True)

    # bar segments:
    #   pred_sig  (dark teal, P<0.05 in predicted direction)
    #   pred_trend (light teal, same direction but NS)
    #   opp_trend  (light coral, opposite direction but NS)
    #   opp_sig   (dark coral, opposite direction and P<0.05)

    LIGHT_GOOD = lighten(GOOD_HEX, 0.55)
    LIGHT_BAD = lighten(BAD_HEX, 0.55)

    for i, (sig, pretty, pred) in enumerate(signatures):
        y = _i(by + (i + 0.5) * row_h)
        bar_h = _i(row_h * 0.60)

        counts = {"pred_sig": 0, "pred_trend": 0, "opp_trend": 0, "opp_sig": 0}
        for gse in PRIMARY_COHORTS:
            delta, p = per_cohort_delta(sig, gse)
            if delta is None:
                continue
            # re-orient: positive = predicted direction
            oriented = -delta if pred == "down" else delta
            if oriented > 0:
                if p < 0.05:
                    counts["pred_sig"] += 1
                else:
                    counts["pred_trend"] += 1
            else:
                if p < 0.05:
                    counts["opp_sig"] += 1
                else:
                    counts["opp_trend"] += 1

        # stack: pred_sig + pred_trend + opp_trend + opp_sig = 5
        segs = [("pred_sig", counts["pred_sig"], GOOD),
                ("pred_trend", counts["pred_trend"], LIGHT_GOOD),
                ("opp_trend", counts["opp_trend"], LIGHT_BAD),
                ("opp_sig", counts["opp_sig"], BAD)]

        x_cursor = tx(0)
        for label, cnt, color in segs:
            if cnt == 0:
                continue
            seg_w = tx(cnt) - tx(0)
            add_rect(s, x_cursor, y - bar_h // 2, seg_w, bar_h,
                     fill=color, line_color=WHITE, line_width=0.4)
            # label inside segment if big enough
            if cnt >= 2:
                add_text(s, x_cursor, y - Inches(0.08),
                         seg_w, Inches(0.14),
                         str(cnt), size=7, bold=True, color=WHITE,
                         align="center")
            x_cursor += seg_w

        # signature label on left
        arrow = "↓" if pred == "down" else "↑"
        thread_c = (THREAD1 if sig in [s0 for s0, _, _ in THREAD1_SIGS]
                    else THREAD2)
        add_text(s, Inches(0.20), y - Inches(0.09),
                 bx - Inches(0.22), Inches(0.18),
                 f"{pretty} {arrow}",
                 size=7, color=thread_c, align="right")

        # right-side pooled Z
        Z, p, _ = get_pooled(sig, source="5cohort")
        is_sig = abs(Z) >= 1.96
        add_text(s, _i(bx + bw + Inches(0.05)),
                 y - Inches(0.12), Inches(0.8), Inches(0.13),
                 f"Z = {Z:+.2f}", size=7, bold=is_sig,
                 color=HIGHLIGHT if is_sig else RGBColor(0x66, 0x66, 0x66),
                 align="left")
        p_str = "P < 10⁻³" if p < 0.001 else f"P = {p:.3f}"
        add_text(s, _i(bx + bw + Inches(0.05)),
                 y + Inches(0.01), Inches(0.8), Inches(0.13),
                 p_str, size=6,
                 color=HIGHLIGHT if is_sig else RGBColor(0x66, 0x66, 0x66),
                 align="left")

    # legend
    leg_y = Inches(4.10)
    items = [(GOOD, "predicted dir., P<0.05"),
             (LIGHT_GOOD, "predicted dir., trend"),
             (LIGHT_BAD, "opposite dir., trend"),
             (BAD, "opposite dir., P<0.05")]
    xb = Inches(0.15)
    for color, lab in items:
        add_rect(s, xb, leg_y, Inches(0.16), Inches(0.12), fill=color)
        add_text(s, xb + Inches(0.19), leg_y - Emu(10000),
                 Inches(1.45), Inches(0.16),
                 lab, size=5, color=INK, align="left")
        xb += Inches(1.60)


# -------------------------------------------------------------------
# Panel D  (Slide 4) --- discovery vs external scatter
# -------------------------------------------------------------------
def build_D():
    s = new_slide(prs_main)
    draw_panel_letter(s, "D")

    signatures = THREAD1_SIGS + THREAD2_SIGS

    # plot area
    px = Inches(1.20); py = Inches(0.55)
    pw = Inches(3.50); ph = Inches(3.30)

    x_lo, x_hi = -1.5, 4.0
    y_lo, y_hi = -1.5, 4.0

    def tx(v): return _i(px + (v - x_lo) / (x_hi - x_lo) * pw)
    def ty(v):
        return _i(py + ph - (v - y_lo) / (y_hi - y_lo) * ph)

    # ------ Quadrant shading (very faint) ------
    # upper-right: reproducible (both sig)
    add_rect(s, tx(1.96), ty(y_hi), tx(x_hi) - tx(1.96),
             ty(1.96) - ty(y_hi),
             fill=RGBColor(0xEE, 0xF6, 0xF3))
    # upper-left: external-only (discovery under-powered)
    add_rect(s, tx(x_lo), ty(y_hi), tx(1.96) - tx(x_lo),
             ty(1.96) - ty(y_hi),
             fill=RGBColor(0xF5, 0xEE, 0xDC))
    # lower-right: discovery-only (external heterogeneous)
    add_rect(s, tx(1.96), ty(1.96), tx(x_hi) - tx(1.96),
             ty(y_lo) - ty(1.96),
             fill=RGBColor(0xE9, 0xF0, 0xF5))

    # spines
    add_line(s, px, py, px, py + ph, LINE, 0.6)
    add_line(s, px, py + ph, px + pw, py + ph, LINE, 0.6)

    # axis ticks
    for v in [-1, 0, 1, 2, 3, 4]:
        xx = tx(v); yy = ty(v)
        add_line(s, xx, _i(py + ph), xx,
                 _i(py + ph + Inches(0.04)), LINE, 0.4)
        add_text(s, xx - Inches(0.13), _i(py + ph + Inches(0.05)),
                 Inches(0.26), Inches(0.14),
                 f"{v:+d}" if v != 0 else "0",
                 size=7, color=INK, align="center")
        add_line(s, _i(px - Inches(0.04)), yy, px, yy, LINE, 0.4)
        add_text(s, _i(px - Inches(0.32)), yy - Inches(0.08),
                 Inches(0.28), Inches(0.14),
                 f"{v:+d}" if v != 0 else "0",
                 size=7, color=INK, align="right")

    # reference lines: Z=0 (solid), Z=±1.96 (dashed), diagonal y=x
    add_line(s, tx(0), py, tx(0), py + ph, INK, 0.6)
    add_line(s, px, ty(0), px + pw, ty(0), INK, 0.6)
    add_line(s, tx(1.96), py, tx(1.96), py + ph, GREY, 0.4, dashed=True)
    add_line(s, px, ty(1.96), px + pw, ty(1.96), GREY, 0.4, dashed=True)
    add_line(s, tx(x_lo), ty(x_lo), tx(x_hi), ty(x_hi),
             RGBColor(0x99, 0x99, 0x99), 0.4, dashed=True)

    # quadrant labels (corners, small italic)
    add_text(s, tx(2.2), ty(3.8), Inches(1.6), Inches(0.18),
             "UR  reproducible",
             size=6, bold=True, color=GOOD, italic=True, align="left")
    add_text(s, tx(-1.3), ty(3.8), Inches(1.6), Inches(0.18),
             "UL  external-only",
             size=6, bold=True, color=HIGHLIGHT, italic=True, align="left")
    add_text(s, tx(2.2), ty(-1.3), Inches(1.6), Inches(0.18),
             "LR  discovery-only",
             size=6, bold=True, color=RGBColor(0x00, 0x66, 0xCC),
             italic=True, align="left")

    # axis titles
    add_text(s, px, _i(py + ph + Inches(0.22)),
             pw, Inches(0.15),
             "Discovery Stouffer-style Z  (n = 33 pre, MW-derived)",
             size=8, color=INK, align="center")
    yt = add_text(s, _i(px - Inches(0.85)),
                  _i(py + ph / 2 - Inches(1.0)),
                  Inches(0.40), Inches(2.0),
                  "External pooled Stouffer Z (5-cohort primary)",
                  size=8, color=INK, align="center")
    yt.rotation = -90

    # plot 7 signatures
    for sig, pretty, pred in signatures:
        disc_row = discovery[discovery.signature == sig]
        if disc_row.empty:
            continue
        disc_z = float(disc_row.iloc[0]["discovery_Z"])
        ext_z, _, _ = get_pooled(sig, source="5cohort")
        if ext_z is None:
            continue
        thread_c = (THREAD1 if sig in [s0 for s0, _, _ in THREAD1_SIGS]
                    else THREAD2)
        # clamp
        dx = max(x_lo, min(x_hi, disc_z))
        dy = max(y_lo, min(y_hi, ext_z))
        xx = tx(dx); yy = ty(dy)

        is_thread1 = sig in [s0 for s0, _, _ in THREAD1_SIGS]
        if is_thread1:
            add_circle(s, xx, yy, Emu(42000),
                       fill=thread_c, line_color=WHITE, line_width=1.2)
        else:
            add_diamond(s, xx, yy, Emu(45000),
                        fill=thread_c, line_color=WHITE, line_width=1.2)

        # label
        # offset the label to avoid overlap with marker
        label_offset_x = Inches(0.12)
        label_offset_y = -Inches(0.08)
        add_text(s, xx + label_offset_x, yy + label_offset_y,
                 Inches(1.4), Inches(0.16),
                 pretty, size=6,
                 bold=True, color=thread_c, align="left")

        # Z values on second line
        add_text(s, xx + label_offset_x, yy + label_offset_y + Inches(0.13),
                 Inches(1.4), Inches(0.14),
                 f"(disc {disc_z:+.1f}, ext {ext_z:+.1f})",
                 size=5, color=RGBColor(0x55, 0x55, 0x55), align="left")

    # also plot the Akiyoshi-augmented CD8 as a star marker
    # discovery for CD8 is the same (disc_z = -0.20)
    cd8_disc = float(discovery[discovery.signature == "CD8_cytotoxic"]
                     .iloc[0]["discovery_Z"])
    cd8_ext6, _, _ = get_pooled("CD8_cytotoxic", source="6source")
    xx = tx(max(x_lo, min(x_hi, cd8_disc)))
    yy = ty(max(y_lo, min(y_hi, cd8_ext6)))
    add_diamond(s, xx, yy, Emu(60000),
                fill=HIGHLIGHT, line_color=WHITE, line_width=1.5)
    add_text(s, xx + Inches(0.14), yy - Inches(0.08),
             Inches(1.6), Inches(0.16),
             "CD8-cyt ★ +Akiyoshi",
             size=6, bold=True, color=HIGHLIGHT, align="left")
    add_text(s, xx + Inches(0.14), yy + Inches(0.05),
             Inches(1.6), Inches(0.13),
             f"(disc {cd8_disc:+.1f}, 6-src ext {cd8_ext6:+.1f})",
             size=5, color=HIGHLIGHT, align="left")

    # legend (right side)
    leg_x = Inches(5.25); leg_y = py + Inches(0.10)
    add_text(s, leg_x, leg_y, Inches(1.20), Inches(0.14),
             "Marker", size=7, bold=True, color=INK, align="left")
    add_circle(s, leg_x + Inches(0.08), leg_y + Inches(0.22), Emu(35000),
               fill=THREAD1, line_color=WHITE, line_width=1.0)
    add_text(s, leg_x + Inches(0.22), leg_y + Inches(0.15),
             Inches(1.0), Inches(0.16),
             "Thread 1", size=6, color=INK, align="left")
    add_diamond(s, leg_x + Inches(0.08), leg_y + Inches(0.42), Emu(35000),
                fill=THREAD2, line_color=WHITE, line_width=1.0)
    add_text(s, leg_x + Inches(0.22), leg_y + Inches(0.35),
             Inches(1.0), Inches(0.16),
             "Thread 2", size=6, color=INK, align="left")
    add_diamond(s, leg_x + Inches(0.08), leg_y + Inches(0.62), Emu(45000),
                fill=HIGHLIGHT, line_color=WHITE, line_width=1.2)
    add_text(s, leg_x + Inches(0.22), leg_y + Inches(0.55),
             Inches(1.2), Inches(0.16),
             "Akiyoshi-aug.", size=6, color=INK, align="left")

    # bottom caption
    add_text(s, Inches(0.15), Inches(4.15),
             SLIDE_W - Inches(0.3), Inches(0.16),
             "Dashed vertical/horizontal lines = Z = ±1.96 (P = 0.05 threshold). "
             "Dashed diagonal = y = x.",
             size=6, color=RGBColor(0x55, 0x55, 0x55), align="center")
    add_text(s, Inches(0.15), Inches(4.32),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "CD8-cytotoxic sits in the UL (external-only) quadrant in discovery (n=33 "
             "under-powered for immune), but Akiyoshi-augmented moves toward UR.",
             size=6, italic=True, color=HIGHLIGHT, align="center")


build_A()
build_B()
build_C()
build_D()
deck_main = f"{OUT}/Fig9_external_validation_native_editable.pptx"
prs_main.save(deck_main)
print(f"wrote {deck_main}")


# ===========================================================================
# SUPP FIGURE S19
# ===========================================================================
prs_supp = Presentation()


def build_S19A():
    """9-cohort unrestricted sensitivity meta forest."""
    s = new_slide(prs_supp)
    draw_panel_letter(s, "A")

    # recomputed 9-cohort meta using Stouffer's Z √N-weighted
    # we can get individual cohort Z from per_cohort stats for each of 9 GSEs
    # and compute pooled Z
    from scipy.stats import norm as norm_dist

    # load the 9-cohort wide per-cohort thread1 file
    t1_wide = pd.read_csv(f"{DATA}/thread1_per_cohort_wide.tsv", sep="\t",
                          header=[0, 1])
    t1_wide.columns = ["_".join(map(str, c)).strip("_") for c in t1_wide.columns]
    t1_wide = t1_wide.rename(columns={t1_wide.columns[0]: "gse"})

    all_cohorts = PRIMARY_COHORTS + ["GSE46862", "GSE94104",
                                     "GSE119409", "GSE150082"]
    n_map = {}
    for gse in all_cohorts:
        row = thread1_summary[thread1_summary.gse == gse]
        if not row.empty:
            n_map[gse] = int(row.iloc[0]["n_samples"])

    # Thread 1 signatures in t1_wide: DSB_HDR_repair, E2F_MYC_cellcycle, EMT,
    #   Tumor_cellcycle
    # Signature ordering: show both thread 1 and per_cohort
    signatures = THREAD1_SIGS + THREAD2_SIGS

    def compute_pooled_z(sig, cohorts):
        zs = []; weights = []
        for gse in cohorts:
            delta, p = per_cohort_delta(sig, gse)
            if delta is None or not np.isfinite(p) or p <= 0:
                continue
            sgn = 1 if delta > 0 else -1
            if sig == "EMT":
                sgn = -sgn
            n = n_map.get(gse, 0)
            if n == 0:
                continue
            z = norm_dist.ppf(1 - p / 2) * sgn
            zs.append(z); weights.append(np.sqrt(n))
        if not zs:
            return None, None, None
        zs = np.array(zs); w = np.array(weights)
        pooled = (w * zs).sum() / np.sqrt((w ** 2).sum())
        p_pool = 2 * (1 - norm_dist.cdf(abs(pooled)))
        total_n = sum(n_map.get(c, 0) for c in cohorts)
        return pooled, p_pool, total_n

    # render
    ax_x = Inches(2.30); ax_y = Inches(0.55)
    ax_w = Inches(3.40); ax_h = Inches(3.10)

    x_lo, x_hi = -3.0, 4.0

    def tx(v):
        return _i(ax_x + (v - x_lo) / (x_hi - x_lo) * ax_w)

    n_sig_rows = len(signatures)
    row_h = ax_h / (n_sig_rows + 1)

    # spines + ticks
    add_line(s, ax_x, ax_y + ax_h, ax_x + ax_w,
             ax_y + ax_h, LINE, 0.6)
    for v in [-3, -2, -1, 0, 1, 2, 3, 4]:
        xx = tx(v)
        add_line(s, xx, _i(ax_y + ax_h), xx,
                 _i(ax_y + ax_h + Inches(0.04)), LINE, 0.4)
        add_text(s, xx - Inches(0.13), _i(ax_y + ax_h + Inches(0.05)),
                 Inches(0.26), Inches(0.14),
                 f"{v:+d}" if v != 0 else "0",
                 size=6, color=INK, align="center")
    add_text(s, ax_x, _i(ax_y + ax_h + Inches(0.20)),
             ax_w, Inches(0.15),
             "Stouffer pooled Z (oriented)",
             size=8, color=INK, align="center")

    add_line(s, tx(0), ax_y, tx(0), ax_y + ax_h, INK, 1.0)
    add_line(s, tx(1.96), ax_y, tx(1.96), ax_y + ax_h,
             GREY, 0.5, dashed=True)
    add_line(s, tx(-1.96), ax_y, tx(-1.96), ax_y + ax_h,
             GREY, 0.5, dashed=True)

    # paired rows per signature: 5-cohort + 9-cohort
    # too tall --- simpler: two diamonds per row (5-cohort solid, 9-cohort hollow)
    for i, (sig, pretty, pred) in enumerate(signatures):
        y = _i(ax_y + (i + 0.5) * row_h + Inches(0.15))

        # 5-cohort
        Z5, p5, N5 = get_pooled(sig, source="5cohort")
        # 9-cohort
        Z9, p9, N9 = compute_pooled_z(sig, all_cohorts)

        thread_c = (THREAD1 if sig in [s0 for s0, _, _ in THREAD1_SIGS]
                    else THREAD2)
        arrow = "↓" if pred == "down" else "↑"

        # label
        add_text(s, Inches(0.20), y - Inches(0.09),
                 ax_x - Inches(0.22), Inches(0.18),
                 f"{pretty} {arrow}",
                 size=7, color=thread_c, align="right")

        # 5-cohort diamond (solid)
        xx5 = tx(max(x_lo, min(x_hi, Z5)))
        add_diamond(s, xx5, y - Inches(0.09), Emu(40000),
                    fill=thread_c, line_color=WHITE, line_width=1.0)

        # 9-cohort diamond (hollow)
        xx9 = tx(max(x_lo, min(x_hi, Z9)))
        add_diamond(s, xx9, y + Inches(0.09), Emu(40000),
                    fill=WHITE, line_color=thread_c, line_width=1.2)

        # right text: Z (5-cohort), Z (9-cohort)
        rx = ax_x + ax_w + Inches(0.05)
        add_text(s, rx, y - Inches(0.17), Inches(0.85), Inches(0.13),
                 f"5-c Z = {Z5:+.2f}", size=6, bold=abs(Z5) >= 1.96,
                 color=thread_c if abs(Z5) >= 1.96 else RGBColor(0x66, 0x66, 0x66),
                 align="left")
        add_text(s, rx, y, Inches(0.85), Inches(0.13),
                 f"9-c Z = {Z9:+.2f}", size=6, bold=abs(Z9) >= 1.96,
                 color=RGBColor(0x66, 0x66, 0x66),
                 align="left")

    # legend (bottom)
    leg_y = Inches(3.95)
    add_diamond(s, Inches(0.55), leg_y + Inches(0.05), Emu(30000),
                fill=THREAD1, line_color=WHITE, line_width=1.0)
    add_text(s, Inches(0.70), leg_y - Emu(5000),
             Inches(1.60), Inches(0.15),
             "5-cohort primary (N = 518)",
             size=6, color=INK, align="left")
    add_diamond(s, Inches(0.55), leg_y + Inches(0.23), Emu(30000),
                fill=WHITE, line_color=THREAD1, line_width=1.2)
    add_text(s, Inches(0.70), leg_y + Inches(0.17),
             Inches(1.6), Inches(0.15),
             "9-cohort unrestricted (N = 721)",
             size=6, color=INK, align="left")

    # caption
    add_text(s, Inches(0.15), Inches(4.35),
             SLIDE_W - Inches(0.3), Inches(0.15),
             "Restriction to 5 concordant cohorts (≥3/4 Thread-1 features in discovery "
             "direction) sharpens the pooled Z for every signature.",
             size=6, italic=True, color=RGBColor(0x55, 0x55, 0x55),
             align="center")


def build_S19B():
    """Akiyoshi 2023 alternative-statistic sensitivity."""
    s = new_slide(prs_supp)
    draw_panel_letter(s, "B")

    # alternative stats from Akiyoshi 2023 paper, each combined with
    # the 5-cohort Thread-2 CD8 meta
    alt_stats = [
        ("cytolytic activity (GZMA×PRF1)",
         "eFig 4B, TRG1/2 vs TRG3/4 P=0.005", 2.81, "★ primary"),
        ("effector memory CD8 ssGSEA",
         "eFig 8, P<0.001", 3.35, ""),
        ("MCP-counter cytotoxic lymphocyte",
         "eFig 4A, P=0.005", 2.81, ""),
        ("activated CD8 ssGSEA",
         "eFig 8, P=0.03", 2.17, ""),
    ]

    # 5-cohort only Z (primary CD8 meta without Akiyoshi)
    base_z, base_p, _ = get_pooled("CD8_cytotoxic", source="5cohort")
    w_akiyoshi = np.sqrt(298)

    bx = Inches(0.70); by = Inches(0.50)
    bw = Inches(4.70); bh = Inches(2.90)

    x_lo, x_hi = 1.5, 4.5
    row_h = bh / (len(alt_stats) + 2)

    def tx(v):
        return _i(bx + (v - x_lo) / (x_hi - x_lo) * bw)

    # spines
    add_line(s, bx, by + bh, bx + bw, by + bh, LINE, 0.5)
    for v in [1.5, 2.0, 2.5, 3.0, 3.5, 4.0, 4.5]:
        xx = tx(v)
        add_line(s, xx, _i(by + bh), xx,
                 _i(by + bh + Inches(0.04)), LINE, 0.4)
        add_text(s, xx - Inches(0.14), _i(by + bh + Inches(0.05)),
                 Inches(0.28), Inches(0.14),
                 f"{v:+.1f}", size=6, color=INK, align="center")
    add_text(s, bx, _i(by + bh + Inches(0.20)),
             bw, Inches(0.15),
             "6-source pooled Stouffer Z (CD8-cytotoxic)",
             size=8, color=INK, align="center")

    add_line(s, tx(1.96), by, tx(1.96), by + bh, GREY, 0.5, dashed=True)
    add_text(s, tx(1.96) - Inches(0.30), _i(by + Inches(0.02)),
             Inches(0.6), Inches(0.12),
             "P=0.05", size=5, color=RGBColor(0x99, 0x99, 0x99),
             italic=True, align="center")

    # Reference: base 5-cohort only
    y_base = _i(by + (0.5) * row_h)
    add_text(s, Inches(0.20), y_base - Inches(0.08),
             Inches(2.0), Inches(0.14),
             "5-cohort only (no Akiyoshi)",
             size=7, bold=True, color=RGBColor(0x55, 0x55, 0x55),
             align="right", italic=True)
    xx5 = tx(max(x_lo, min(x_hi, base_z)))
    add_diamond(s, xx5, y_base, Emu(40000),
                fill=RGBColor(0x99, 0x99, 0x99), line_color=WHITE,
                line_width=1.0)
    add_text(s, _i(bx + bw + Inches(0.05)), y_base - Inches(0.08),
             Inches(1.0), Inches(0.14),
             f"Z = {base_z:.2f}",
             size=6, color=RGBColor(0x66, 0x66, 0x66), align="left")

    # each Akiyoshi alternative stat
    for i, (stat_name, src_desc, z_ak, tag) in enumerate(alt_stats):
        y = _i(by + (i + 1.5) * row_h)
        # combine with 5-cohort primary using Stouffer (5 cohorts treated as
        # aggregate with their √N weights; plus Akiyoshi paper-level Z with
        # weight √298).
        # Since 5-cohort pool already gives Z5 with total weight W5, to add
        # Akiyoshi we need W5 as sqrt(sum N_i) implicitly. We'll compute by
        # reconstruction:
        w5 = np.sqrt(518)
        combined = (base_z * w5 + z_ak * w_akiyoshi) / np.sqrt(w5 ** 2 + w_akiyoshi ** 2)
        p_combined = 2 * (1 - norm.cdf(abs(combined)))

        # label
        label_color = HIGHLIGHT if "primary" in tag else INK
        add_text(s, Inches(0.20), y - Inches(0.15),
                 Inches(2.0), Inches(0.13),
                 stat_name,
                 size=6, bold=True, color=label_color, align="right")
        add_text(s, Inches(0.20), y - Inches(0.01),
                 Inches(2.0), Inches(0.13),
                 src_desc,
                 size=5, italic=True, color=RGBColor(0x77, 0x77, 0x77),
                 align="right")

        # diamond
        xx = tx(max(x_lo, min(x_hi, combined)))
        dia_color = HIGHLIGHT if "primary" in tag else AKIYOSHI
        size = Emu(50000) if "primary" in tag else Emu(40000)
        add_diamond(s, xx, y, size,
                    fill=dia_color, line_color=WHITE, line_width=1.1)

        # right text: Z and P
        rx = bx + bw + Inches(0.05)
        add_text(s, rx, y - Inches(0.12),
                 Inches(1.0), Inches(0.13),
                 f"Z = {combined:+.2f}",
                 size=7, bold=True,
                 color=dia_color, align="left")
        p_str = "P < 10⁻³" if p_combined < 0.001 else f"P = {p_combined:.3f}"
        add_text(s, rx, y + Inches(0.01),
                 Inches(1.0), Inches(0.13),
                 p_str,
                 size=6, color=dia_color, align="left")

    # bottom caption
    add_text(s, Inches(0.15), Inches(3.90),
             SLIDE_W - Inches(0.3), Inches(0.15),
             "All four Akiyoshi alternative statistics give 6-source pooled Z > +2.15; "
             "primary (cytolytic activity) gives Z = +3.29 (P = 0.001).",
             size=7, bold=True, color=INK, align="center")
    add_text(s, Inches(0.15), Inches(4.10),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "Result is robust to choice of Akiyoshi statistic.",
             size=6, italic=True, color=HIGHLIGHT, align="center")


build_S19A()
build_S19B()
deck_supp = f"{OUT}/SuppFig_S19_external_validation_sensitivity_native_editable.pptx"
prs_supp.save(deck_supp)
print(f"wrote {deck_supp}")
