#!/usr/bin/env python3
"""
23_fig8_native_pptx.py

Build Figure 8 --- Paired cascade phenomenology + convergence-null
callout (§3.11) --- as 6 native-editable PowerPoint slides, plus a
2-slide supplementary deck (purity-adjusted sensitivity + convergence
test scatter matrix).

Main Figure 8 panels
--------------------
  A  Between-group BCa forest (9 cascade features): diamond + error
     bar on a standardised-Z scale (diff / SE, where SE = CI_width /
     3.92); vertical references at Z=0 and Z=±1.96; Treg highlighted
     gold as the only feature with CI strictly excluding zero.
  B  Paired pre→post spaghetti for 4 key immune features (Treg /
     MHC-II / CD8 exhaustion / IGH_n) in a 2x2 grid.
  C  Within-group Δ forest (per feature, good and bad median Δ
     with BCa CIs) — shows multiple features with within-good CIs
     excluding zero (SBS5, neo_binders, neo_sites, Treg, MHC-II,
     CD8_exhaustion).
  D  Per-subject Δ waterfall for SBS5 mutation clearance and MHC-I
     neoantigen binder clearance.
  E  Conceptual fishplot schematic of paired clonal dynamics (good
     responder: clone eradication; bad responder: clone persistence).
  F  Cascade schematic + convergence-null callout (THE critical
     message): a flow diagram of the proposed cascade overlaid with
     the pre-specified convergence-test result showing that the
     static baseline predictor (Thread 1) does NOT predict cascade
     Δ magnitude (0/36 pairs P<0.05; DSB→CD8cyt r=-0.07, P=0.83).

Supp Fig S20 panels
-------------------
  A  Convergence test scatter matrix: 9 baseline LASSO winners × 4
     key cascade Δ = 36 pre-specified pairs, each a mini scatter with
     Spearman r and P.
  B  Purity-adjusted paired Δ sensitivity: Treg / MHC-II / CD8_exh /
     IGH_n before vs after CNVkit-purity correction (Supp Text S4).

Motif references consulted (cascade / paired biomarker convention):
  - Tumeh et al Nature 2014 (PD-1 pre/post CD8): paired pre→post
    spaghetti + waterfall.
  - Riaz et al Cell 2017 (ICI pre/on-treatment): paired biomarker
    panels with per-subject tracks.
  - Cercek et al NEJM 2022 (dostarlimab): per-subject response
    waterfall; simple timeline schematics.
  - Rizvi et al Science 2015 (TIL clonality pre/post): clonal
    trajectory fishplots.
  - Rosenthal et al Nature 2019 (HLA-LOH / immune escape): cascade
    flow schematic overlaid with statistical null.
  - Chowell et al Science 2018 (HLA-ICB): BCa CI forest.

Rules: one panel per slide; no plot titles; python-pptx native
elements only (TEXT_BOX / LINE / AUTO_SHAPE / FREEFORM); Arial;
DEEP palette GOOD=#0a7d6e / BAD=#c53e1f; kill_shadow() + _i()
applied everywhere.
"""

import os
import numpy as np
import pandas as pd
from scipy import stats
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Shared infrastructure
# ---------------------------------------------------------------------------
def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def _i(v):
    return int(round(float(v)))


GOOD = RGBColor(0x0A, 0x7D, 0x6E)
BAD = RGBColor(0xC5, 0x3E, 0x1F)
INK = RGBColor(0x22, 0x22, 0x22)
LINE = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0xBB, 0xBB, 0xBB)
LT_GREY = RGBColor(0xDD, 0xDD, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xD4, 0xA3, 0x00)

GOOD_HEX = (0x0A, 0x7D, 0x6E)
BAD_HEX = (0xC5, 0x3E, 0x1F)

FONT = "Arial"
SLIDE_W = Inches(6.5)
SLIDE_H = Inches(4.5)

OUT = "/data/data/TNT/analysis/260418_add/ppt"
DATA = "/data/data/TNT/analysis"
os.makedirs(OUT, exist_ok=True)


def new_slide(prs, w=SLIDE_W, h=SLIDE_H):
    prs.slide_width = w
    prs.slide_height = h
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, size=8, bold=False,
             color=INK, align="left", anchor="middle", font=FONT,
             italic=False):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(max(w, 1)), _i(max(h, 1)))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "bottom": MSO_ANCHOR.BOTTOM,
                          "middle": MSO_ANCHOR.MIDDLE}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = str(text)
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.italic = bool(italic)
    r.font.color.rgb = color
    kill_shadow(tb)
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.5, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        try:
            c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass
    kill_shadow(c)
    return c


def add_arrow(slide, x1, y1, x2, y2, color=INK, width=1.0):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    kill_shadow(c)
    return c


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0.5):
    w = max(_i(w), 1); h = max(_i(h), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _i(x), _i(y), w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_rounded_rect(slide, x, y, w, h, fill=None, line_color=None,
                     line_width=0.5):
    w = max(_i(w), 1); h = max(_i(h), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 _i(x), _i(y), w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_circle(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_diamond(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_freeform_poly(slide, vertices, fill=None, line_color=None,
                      line_width=0.5):
    if len(vertices) < 3:
        return None
    x0, y0 = vertices[0]
    ff = slide.shapes.build_freeform(_i(x0), _i(y0), scale=1.0)
    pts = [(_i(x), _i(y)) for x, y in vertices[1:]]
    ff.add_line_segments(pts, close=True)
    shp = ff.convert_to_shape()
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def lighten(rgb_hex, factor=0.7):
    r0, g0, b0 = rgb_hex
    return RGBColor(int(r0 + (255 - r0) * factor),
                    int(g0 + (255 - g0) * factor),
                    int(b0 + (255 - b0) * factor))


def draw_panel_letter(slide, letter):
    add_text(slide, Inches(0.15), Inches(0.1), Inches(0.4), Inches(0.35),
             letter, size=14, bold=True, color=INK, align="left")


# ---------------------------------------------------------------------------
# Load data
# ---------------------------------------------------------------------------
bca = pd.read_csv(f"{DATA}/tables/TableS8_cascade_BCa_bootstrap.tsv",
                  sep="\t")
paired_long = pd.read_csv(
    f"{DATA}/09_integration/paired_delta/paired_feature_long.tsv",
    sep="\t")
convergence = pd.read_csv(
    f"{DATA}/260418_add/targeted_convergence_test.tsv", sep="\t")


def parse_ci(s):
    """Parse '[+x.xx, +y.yy]' into (lo, hi) floats."""
    s = str(s).strip().strip("[]")
    parts = [p.strip() for p in s.split(",")]
    return float(parts[0]), float(parts[1])


# Cascade feature ordering (biological sequence + Treg at end for emphasis)
CASCADE_FEATURES = [
    "SBS5", "missense", "neo_binders", "neo_sites",
    "MHC_II", "CD8_exhaustion", "Treg",
    "IGH_n", "TRB_shannon",
]
PRETTY = {
    "SBS5": "SBS5 (clock-like)  mutations",
    "missense": "Missense mutations",
    "neo_binders": "MHC-I neoantigen binders",
    "neo_sites": "Neoantigen sites",
    "MHC_II": "MHC-II signature",
    "CD8_exhaustion": "CD8 exhaustion",
    "Treg": "Treg signature  ★",
    "IGH_n": "IGH clonotypes (B-cell)",
    "TRB_shannon": "TRB Shannon diversity",
}

# Bio-axis grouping
FEAT_GROUP = {
    "SBS5": "Clone clearance",
    "missense": "Clone clearance",
    "neo_binders": "Antigen clearance",
    "neo_sites": "Antigen clearance",
    "MHC_II": "Immune reprogramming",
    "CD8_exhaustion": "Immune reprogramming",
    "Treg": "Immune reprogramming",
    "IGH_n": "B-cell infiltration",
    "TRB_shannon": "T-cell repertoire",
}
GROUP_COLOR = {
    "Clone clearance": RGBColor(0x3E, 0x52, 0x86),
    "Antigen clearance": RGBColor(0x6A, 0x4C, 0x93),
    "Immune reprogramming": RGBColor(0x0A, 0x7D, 0x6E),
    "B-cell infiltration": RGBColor(0xA6, 0x50, 0x2C),
    "T-cell repertoire": RGBColor(0x55, 0x55, 0x55),
}


# ===========================================================================
# MAIN FIGURE 8
# ===========================================================================
prs_main = Presentation()


# -------------------------------------------------------------------
# Panel A --- Between-group BCa forest (standardised-Z scale)
# -------------------------------------------------------------------
def build_A():
    s = new_slide(prs_main)
    draw_panel_letter(s, "A")

    # plot area
    ax_x = Inches(2.80); ax_y = Inches(0.55)
    ax_w = Inches(2.85); ax_h = Inches(3.25)

    x_lo, x_hi = -3.2, 3.2

    def tx(v):
        return _i(ax_x + (v - x_lo) / (x_hi - x_lo) * ax_w)

    n_rows = len(CASCADE_FEATURES)
    row_h = ax_h / n_rows

    # spines + ticks
    add_line(s, ax_x, ax_y + ax_h, ax_x + ax_w, ax_y + ax_h, LINE, 0.5)
    for v in [-3, -2, -1, 0, 1, 2, 3]:
        xx = tx(v)
        add_line(s, xx, _i(ax_y + ax_h), xx,
                 _i(ax_y + ax_h + Inches(0.04)), LINE, 0.4)
        add_text(s, xx - Inches(0.13), _i(ax_y + ax_h + Inches(0.05)),
                 Inches(0.26), Inches(0.14),
                 f"{v:+d}" if v != 0 else "0",
                 size=6, color=INK, align="center")
    add_text(s, ax_x, _i(ax_y + ax_h + Inches(0.22)),
             ax_w, Inches(0.15),
             "Standardised between-group Z  (good − bad)",
             size=8, color=INK, align="center")
    add_text(s, ax_x, _i(ax_y + ax_h + Inches(0.37)),
             ax_w, Inches(0.13),
             "= diff_median / (CI width / 3.92)",
             size=5, italic=True, color=RGBColor(0x66, 0x66, 0x66),
             align="center")

    # reference lines
    add_line(s, tx(0), ax_y, tx(0), ax_y + ax_h, INK, 1.0)
    add_line(s, tx(1.96), ax_y, tx(1.96), ax_y + ax_h,
             GREY, 0.5, dashed=True)
    add_line(s, tx(-1.96), ax_y, tx(-1.96), ax_y + ax_h,
             GREY, 0.5, dashed=True)
    add_text(s, tx(1.96) - Inches(0.20), _i(ax_y + Inches(0.02)),
             Inches(0.4), Inches(0.12),
             "P=0.05", size=5, italic=True,
             color=RGBColor(0x99, 0x99, 0x99), align="center")

    # render rows
    for i, feat in enumerate(CASCADE_FEATURES):
        y = _i(ax_y + (i + 0.5) * row_h)
        row = bca[bca.feature == feat]
        if row.empty:
            continue
        row = row.iloc[0]
        diff_med = float(row["diff_median_good_minus_bad"])
        lo, hi = parse_ci(row["diff_95CI"])
        # standardised Z
        se = (hi - lo) / 3.92 if (hi - lo) > 0 else 0.1
        z = diff_med / se
        z_lo = lo / se
        z_hi = hi / se

        # clamp
        z_c = max(x_lo, min(x_hi, z))
        z_lo_c = max(x_lo, min(x_hi, z_lo))
        z_hi_c = max(x_lo, min(x_hi, z_hi))

        is_sig = abs(z) >= 1.96
        highlight = (feat == "Treg")

        group = FEAT_GROUP[feat]
        gcolor = GROUP_COLOR[group]

        # feature label (left)
        lab_color = HIGHLIGHT if highlight else gcolor
        add_text(s, Inches(0.18), y - Inches(0.09),
                 ax_x - Inches(0.22), Inches(0.18),
                 PRETTY[feat],
                 size=7, bold=highlight, color=lab_color, align="right")
        # group chip
        add_rect(s, ax_x - Inches(0.12), y - Inches(0.04),
                 Inches(0.06), Inches(0.08), fill=gcolor)

        # CI bar
        add_line(s, tx(z_lo_c), y, tx(z_hi_c), y,
                 gcolor, 1.8 if highlight else 1.2)
        # CI caps
        cap_h = _i(Inches(0.05))
        add_line(s, tx(z_lo_c), y - cap_h, tx(z_lo_c), y + cap_h,
                 gcolor, 1.0)
        add_line(s, tx(z_hi_c), y - cap_h, tx(z_hi_c), y + cap_h,
                 gcolor, 1.0)
        # diamond at point estimate
        d_col = HIGHLIGHT if highlight else gcolor
        d_size = Emu(55000) if highlight else Emu(38000)
        add_diamond(s, tx(z_c), y, d_size,
                    fill=d_col, line_color=WHITE, line_width=1.1)

        # right column: numeric effect + CI + P
        rx = ax_x + ax_w + Inches(0.06)
        eff_txt = f"{diff_med:+.2f}" if abs(diff_med) < 10 else f"{diff_med:+.0f}"
        ci_txt = (f"[{lo:+.2f}, {hi:+.2f}]"
                  if abs(lo) < 10 and abs(hi) < 10
                  else f"[{lo:+.0f}, {hi:+.0f}]")
        p_val = float(row["MW_p"])
        p_txt = f"P = {p_val:.3f}"
        add_text(s, rx, y - Inches(0.14),
                 Inches(0.75), Inches(0.13),
                 eff_txt, size=7, bold=is_sig,
                 color=HIGHLIGHT if highlight else
                       (gcolor if is_sig else RGBColor(0x66, 0x66, 0x66)),
                 align="left")
        add_text(s, rx, y - Inches(0.02),
                 Inches(0.85), Inches(0.13),
                 ci_txt, size=5,
                 color=RGBColor(0x55, 0x55, 0x55), align="left")
        add_text(s, rx, y + Inches(0.10),
                 Inches(0.75), Inches(0.13),
                 p_txt, size=6,
                 bold=is_sig,
                 color=HIGHLIGHT if highlight else
                       (gcolor if is_sig else RGBColor(0x55, 0x55, 0x55)),
                 align="left")

    # group legend (bottom left)
    leg_y = Inches(3.90)
    xb = Inches(0.18)
    add_text(s, xb, leg_y, Inches(1.3), Inches(0.14),
             "Cascade group:",
             size=6, bold=True, color=INK, align="left")
    xb += Inches(0.85)
    for g in ["Clone clearance", "Antigen clearance",
              "Immune reprogramming", "B-cell infiltration",
              "T-cell repertoire"]:
        add_rect(s, xb, leg_y + Inches(0.02),
                 Inches(0.09), Inches(0.09),
                 fill=GROUP_COLOR[g])
        add_text(s, xb + Inches(0.11), leg_y - Emu(5000),
                 Inches(1.10), Inches(0.14),
                 g, size=5, color=INK, align="left")
        xb += Inches(1.12)

    # bottom-strip message
    add_text(s, Inches(0.15), Inches(4.15),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "★ Only Treg Δ has a between-group BCa CI strictly "
             "excluding zero (robust at n = 12).",
             size=6, bold=True, color=HIGHLIGHT, align="center")
    add_text(s, Inches(0.15), Inches(4.30),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "Other cascade features are framed as exploratory "
             "phenomenology; see Panel C for within-group robustness.",
             size=6, italic=True, color=RGBColor(0x55, 0x55, 0x55),
             align="center")


# -------------------------------------------------------------------
# Panel B --- Oriented-Δ distribution for all 9 cascade features
# (single-axis, per-feature scale-normalised; Fig 6A-style)
# -------------------------------------------------------------------
def build_B():
    s = new_slide(prs_main)
    draw_panel_letter(s, "B")

    # Plot area
    px = Inches(0.80); py = Inches(0.55)
    pw = Inches(5.50); ph = Inches(3.00)

    features_b = CASCADE_FEATURES  # all 9
    # Predicted direction per feature (for orientation):
    #  clearance features: predicted = decrease (down)
    #  immune features: predicted = increase (up)
    PREDICTED = {
        "SBS5": "down", "missense": "down",
        "neo_binders": "down", "neo_sites": "down",
        "MHC_II": "up", "CD8_exhaustion": "up", "Treg": "up",
        "IGH_n": "up", "TRB_shannon": "up",
    }

    # For each feature, build per-subject oriented Δ = (Δ × sign_pred) / scale
    # where scale is the feature's max |CI bound| from BCa table (so all
    # features share a common [-1, +1] comparable frame).
    feat_data = {}
    for feat in features_b:
        bca_row = bca[bca.feature == feat]
        if bca_row.empty:
            continue
        r_ = bca_row.iloc[0]
        g_lo, g_hi = parse_ci(r_.good_95CI)
        b_lo, b_hi = parse_ci(r_.bad_95CI)
        scale = max(abs(g_lo), abs(g_hi), abs(b_lo), abs(b_hi),
                    abs(r_.good_delta_median), abs(r_.bad_delta_median),
                    1e-6)
        sub = paired_long[paired_long.feature == feat].copy()
        sub = sub.dropna(subset=["pre", "post"])
        sub["delta"] = sub["post"] - sub["pre"]
        sign_mul = -1 if PREDICTED[feat] == "down" else +1
        sub["oriented_norm"] = (sub["delta"] * sign_mul) / scale
        feat_data[feat] = (sub, scale, r_)

    # y-axis: oriented-Δ / scale, range roughly [-1.2, +1.4]
    y_lo, y_hi = -1.25, 1.45

    def ty(v):
        return _i(py + ph - (v - y_lo) / (y_hi - y_lo) * ph)

    # Faint zone shading
    add_rect(s, px, py, pw, ty(0) - py,
             fill=RGBColor(0xEE, 0xF6, 0xF3))  # pred-dir zone (teal tint)
    add_rect(s, px, ty(0), pw, py + ph - ty(0),
             fill=RGBColor(0xFB, 0xF1, 0xEE))  # opposite zone (coral tint)
    # zone labels
    add_text(s, _i(px + Inches(0.08)), _i(py + Inches(0.03)),
             Inches(3.0), Inches(0.16),
             "↑ predicted direction (clearance / infiltration / activation)",
             size=6, color=RGBColor(0x3A, 0x7A, 0x6B),
             italic=True, align="left")
    add_text(s, _i(px + Inches(0.08)), _i(py + ph - Inches(0.18)),
             Inches(3.0), Inches(0.16),
             "↓ opposite direction",
             size=6, color=RGBColor(0x9B, 0x5A, 0x48),
             italic=True, align="left")

    # Spines + ticks
    add_line(s, px, py, px, py + ph, LINE, 0.6)
    add_line(s, px, py + ph, px + pw, py + ph, LINE, 0.6)
    for yv, lab in [(-1, "−1"), (-0.5, "−0.5"), (0, "0"),
                    (0.5, "+0.5"), (1, "+1"), (1.4, "+1.4")]:
        yy = ty(yv)
        add_line(s, _i(px - Inches(0.05)), yy, px, yy, LINE, 0.5)
        add_text(s, _i(px - Inches(0.45)), yy - Inches(0.08),
                 Inches(0.40), Inches(0.14),
                 lab, size=6, color=INK, align="right")
    # rotated y title
    yt = add_text(s, Inches(0.12),
                  _i(py + ph / 2 - Inches(0.9)),
                  Inches(0.35), Inches(1.8),
                  "Oriented Δ  (post − pre) · sign(predicted)  ÷  scale",
                  size=7, color=INK, align="center")
    yt.rotation = -90

    # Zero reference line
    add_line(s, px, ty(0), px + pw, ty(0), INK, 1.1)
    add_text(s, _i(px + pw + Inches(0.02)),
             ty(0) - Inches(0.08),
             Inches(0.55), Inches(0.16),
             "no change", size=6, color=INK, align="left")

    # Per-feature strips
    n_feat = len(features_b)
    group_w = pw / n_feat
    strip_off = Inches(0.17)
    rng = np.random.default_rng(3)
    LIGHT_GOOD = lighten(GOOD_HEX, 0.66)
    LIGHT_BAD = lighten(BAD_HEX, 0.66)

    for i, feat in enumerate(features_b):
        if feat not in feat_data:
            continue
        sub, scale, r_ = feat_data[feat]
        center_x = _i(px + (i + 0.5) * group_w)
        good_cx = center_x - _i(strip_off)
        bad_cx = center_x + _i(strip_off)

        for grp_name, cx, fill_col, edge_col in [
            ("good", good_cx, LIGHT_GOOD, GOOD),
            ("bad", bad_cx, LIGHT_BAD, BAD),
        ]:
            vals = sub[sub.response == grp_name]["oriented_norm"].values
            if len(vals) == 0:
                continue
            q25, q75 = np.percentile(vals, [25, 75])
            med = float(np.median(vals))
            box_w = Inches(0.17)
            hw = box_w / 2
            # IQR box
            add_rect(s, cx - hw, ty(q75),
                     box_w, max(ty(q25) - ty(q75), 2),
                     fill=fill_col, line_color=edge_col, line_width=0.7)
            # median line (thick)
            add_line(s, cx - hw, ty(med), cx + hw, ty(med),
                     edge_col, 1.8)
            # median diamond
            add_diamond(s, cx, ty(med), Emu(32000),
                        fill=edge_col, line_color=WHITE, line_width=1.0)
            # jittered individual dots
            jitter_range = _i(Inches(0.10))
            for v in vals:
                jx = cx + rng.integers(-jitter_range // 2,
                                        jitter_range // 2 + 1)
                add_circle(s, jx, ty(max(y_lo, min(y_hi, v))),
                           Emu(16000), fill=WHITE,
                           line_color=edge_col, line_width=0.7)

        # Factor separator between features
        if i > 0:
            sep_x = _i(px + i * group_w)
            add_line(s, sep_x, py + Inches(0.10),
                     sep_x, py + ph - Inches(0.10),
                     LT_GREY, 0.3, dashed=True)

        # Feature label below axis (rotated for fit)
        feat_lab = feat.replace("_", " ")
        feat_color = (HIGHLIGHT if feat == "Treg" else
                      GROUP_COLOR[FEAT_GROUP[feat]])
        star = " ★" if feat == "Treg" else ""
        lab = add_text(s, center_x - Inches(0.55),
                       _i(py + ph + Inches(0.08)),
                       Inches(1.1), Inches(0.55),
                       f"{feat_lab}{star}",
                       size=6, bold=(feat == "Treg"),
                       color=feat_color, align="center")
        lab.rotation = 35

        # MW P below feature label
        p_val = float(r_.MW_p)
        p_color = HIGHLIGHT if p_val < 0.05 else RGBColor(0x66, 0x66, 0x66)
        p_text = f"P={p_val:.3f}"
        p_star = "★" if p_val < 0.05 else ""
        add_text(s, center_x - Inches(0.40),
                 _i(py + ph + Inches(0.60)),
                 Inches(0.80), Inches(0.12),
                 f"{p_text}{p_star}",
                 size=5, bold=(p_val < 0.05),
                 color=p_color, align="center")

    # Legend
    leg_y = Inches(4.15)
    add_rect(s, Inches(0.80), leg_y, Inches(0.18), Inches(0.12),
             fill=LIGHT_GOOD, line_color=GOOD, line_width=0.7)
    add_diamond(s, Inches(0.89), leg_y + Inches(0.06), Emu(20000),
                fill=GOOD, line_color=WHITE, line_width=0.6)
    add_text(s, Inches(1.02), leg_y - Emu(5000),
             Inches(0.9), Inches(0.16),
             "good (n=6)", size=7, color=INK, align="left")
    add_rect(s, Inches(2.00), leg_y, Inches(0.18), Inches(0.12),
             fill=LIGHT_BAD, line_color=BAD, line_width=0.7)
    add_diamond(s, Inches(2.09), leg_y + Inches(0.06), Emu(20000),
                fill=BAD, line_color=WHITE, line_width=0.6)
    add_text(s, Inches(2.22), leg_y - Emu(5000),
             Inches(0.9), Inches(0.16),
             "bad (n=6)", size=7, color=INK, align="left")
    add_text(s, Inches(3.40), leg_y - Emu(5000),
             Inches(2.9), Inches(0.16),
             "Box = IQR  ·  diamond = group median  ·  circle = subject",
             size=6, color=RGBColor(0x55, 0x55, 0x55), align="left")
    # Treg star legend
    add_text(s, Inches(3.40), leg_y + Inches(0.14),
             Inches(2.9), Inches(0.14),
             "★ Treg = only feature with CI strictly excluding zero",
             size=6, italic=True, color=HIGHLIGHT, align="left")


# -------------------------------------------------------------------
# Panel C --- Within-group Δ forest (good + bad paired per feature)
# -------------------------------------------------------------------
def build_C():
    s = new_slide(prs_main)
    draw_panel_letter(s, "C")

    # we plot two rows per feature: good (above) and bad (below)
    # on a *standardised per-feature* x-axis: Δ / feature_scale.
    # feature_scale = max_abs_of_either_group_CI, so both groups plot in
    # a common [-1, +1] frame with ±1 = strongest observed CI bound.

    ax_x = Inches(2.50); ax_y = Inches(0.45)
    ax_w = Inches(3.05); ax_h = Inches(3.35)

    x_lo, x_hi = -1.15, 1.15

    def tx(v):
        return _i(ax_x + (v - x_lo) / (x_hi - x_lo) * ax_w)

    n_rows = len(CASCADE_FEATURES)
    row_h = ax_h / n_rows

    # spines + ticks
    add_line(s, ax_x, ax_y + ax_h, ax_x + ax_w, ax_y + ax_h, LINE, 0.5)
    for v, lab in [(-1, "−1.0"), (-0.5, "−0.5"), (0, "0"),
                   (0.5, "+0.5"), (1, "+1.0")]:
        xx = tx(v)
        add_line(s, xx, _i(ax_y + ax_h), xx,
                 _i(ax_y + ax_h + Inches(0.04)), LINE, 0.4)
        add_text(s, xx - Inches(0.18), _i(ax_y + ax_h + Inches(0.05)),
                 Inches(0.36), Inches(0.13),
                 lab, size=6, color=INK, align="center")
    add_text(s, ax_x, _i(ax_y + ax_h + Inches(0.22)),
             ax_w, Inches(0.14),
             "Within-group median Δ  (normalised to per-feature |max CI|)",
             size=7, color=INK, align="center")

    # zero reference
    add_line(s, tx(0), ax_y, tx(0), ax_y + ax_h, INK, 1.0)

    # iterate features
    for i, feat in enumerate(CASCADE_FEATURES):
        row = bca[bca.feature == feat]
        if row.empty:
            continue
        row = row.iloc[0]
        g_med = float(row.good_delta_median)
        b_med = float(row.bad_delta_median)
        g_lo, g_hi = parse_ci(row.good_95CI)
        b_lo, b_hi = parse_ci(row.bad_95CI)
        scale = max(abs(g_lo), abs(g_hi), abs(b_lo), abs(b_hi),
                    abs(g_med), abs(b_med), 1e-6)

        # CI excludes zero?
        g_excl = (g_lo > 0 and g_hi > 0) or (g_lo < 0 and g_hi < 0)
        b_excl = (b_lo > 0 and b_hi > 0) or (b_lo < 0 and b_hi < 0)

        y_row = _i(ax_y + (i + 0.5) * row_h)
        y_good = y_row - _i(Inches(0.08))
        y_bad = y_row + _i(Inches(0.08))

        # feature label left
        group = FEAT_GROUP[feat]
        gcolor = GROUP_COLOR[group]
        add_text(s, Inches(0.18), y_row - Inches(0.09),
                 ax_x - Inches(0.22), Inches(0.18),
                 PRETTY[feat],
                 size=7, bold=(feat == "Treg"),
                 color=HIGHLIGHT if feat == "Treg" else gcolor,
                 align="right")
        add_rect(s, ax_x - Inches(0.12), y_row - Inches(0.04),
                 Inches(0.06), Inches(0.08), fill=gcolor)

        # good row
        g_n_lo = max(x_lo, min(x_hi, g_lo / scale))
        g_n_hi = max(x_lo, min(x_hi, g_hi / scale))
        g_n_med = max(x_lo, min(x_hi, g_med / scale))
        add_line(s, tx(g_n_lo), y_good, tx(g_n_hi), y_good,
                 GOOD, 1.5 if g_excl else 0.9)
        add_line(s, tx(g_n_lo), y_good - _i(Inches(0.04)),
                 tx(g_n_lo), y_good + _i(Inches(0.04)), GOOD, 0.8)
        add_line(s, tx(g_n_hi), y_good - _i(Inches(0.04)),
                 tx(g_n_hi), y_good + _i(Inches(0.04)), GOOD, 0.8)
        add_diamond(s, tx(g_n_med), y_good, Emu(30000),
                    fill=GOOD if g_excl else WHITE,
                    line_color=GOOD, line_width=0.8)

        # bad row
        b_n_lo = max(x_lo, min(x_hi, b_lo / scale))
        b_n_hi = max(x_lo, min(x_hi, b_hi / scale))
        b_n_med = max(x_lo, min(x_hi, b_med / scale))
        add_line(s, tx(b_n_lo), y_bad, tx(b_n_hi), y_bad,
                 BAD, 1.5 if b_excl else 0.9)
        add_line(s, tx(b_n_lo), y_bad - _i(Inches(0.04)),
                 tx(b_n_lo), y_bad + _i(Inches(0.04)), BAD, 0.8)
        add_line(s, tx(b_n_hi), y_bad - _i(Inches(0.04)),
                 tx(b_n_hi), y_bad + _i(Inches(0.04)), BAD, 0.8)
        add_diamond(s, tx(b_n_med), y_bad, Emu(30000),
                    fill=BAD if b_excl else WHITE,
                    line_color=BAD, line_width=0.8)

        # right-column text: scale + summary
        rx = ax_x + ax_w + Inches(0.05)
        # scale for normalization
        if feat in {"IGH_n", "SBS5", "missense", "neo_binders",
                    "neo_sites"}:
            scale_txt = f"scale = {scale:.0f}"
        else:
            scale_txt = f"scale = {scale:.2f}"
        add_text(s, rx, y_row - Inches(0.12),
                 Inches(0.80), Inches(0.12),
                 scale_txt, size=5,
                 color=RGBColor(0x77, 0x77, 0x77), align="left")
        # robustness indicators
        robust_txt = ""
        if g_excl:
            robust_txt += "g·"
        if b_excl:
            robust_txt += "b·"
        if g_excl or b_excl:
            add_text(s, rx, y_row,
                     Inches(0.80), Inches(0.13),
                     f"{robust_txt}CI excludes 0",
                     size=6, bold=True, color=GOOD if g_excl else BAD,
                     align="left")

    # legend
    leg_y = Inches(3.92)
    add_line(s, Inches(0.20), leg_y + Inches(0.06),
             Inches(0.50), leg_y + Inches(0.06), GOOD, 1.5)
    add_diamond(s, Inches(0.35), leg_y + Inches(0.06), Emu(24000),
                fill=GOOD, line_color=WHITE, line_width=0.8)
    add_text(s, Inches(0.55), leg_y - Emu(10000),
             Inches(1.7), Inches(0.16),
             "good Δ filled diamond = CI excl. 0 (robust)",
             size=6, color=INK, align="left")
    add_line(s, Inches(0.20), leg_y + Inches(0.22),
             Inches(0.50), leg_y + Inches(0.22), BAD, 1.5)
    add_diamond(s, Inches(0.35), leg_y + Inches(0.22), Emu(24000),
                fill=WHITE, line_color=BAD, line_width=0.8)
    add_text(s, Inches(0.55), leg_y + Inches(0.13),
             Inches(2.6), Inches(0.16),
             "bad Δ hollow diamond = CI spans 0 (exploratory)",
             size=6, color=INK, align="left")

    # caption
    add_text(s, Inches(0.15), Inches(4.32),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "Within-good CIs excluding zero: SBS5, neo_binders, "
             "neo_sites, MHC-II, CD8 exhaustion, Treg. "
             "Consistent cascade within responders; between-group only Treg.",
             size=6, italic=True, color=INK, align="center")


# -------------------------------------------------------------------
# Panel D --- Per-subject Δ waterfall: SBS5 + neoantigen
# -------------------------------------------------------------------
def build_D():
    s = new_slide(prs_main)
    draw_panel_letter(s, "D")

    # Top sub-panel: SBS5 Δ per subject
    top_x = Inches(0.95); top_y = Inches(0.45)
    top_w = Inches(5.30); top_h = Inches(1.55)
    # Bottom sub-panel: Neoantigen binder Δ per subject
    bot_x = Inches(0.95); bot_y = Inches(2.25)
    bot_w = Inches(5.30); bot_h = Inches(1.55)

    for feat, label, px, py, pw, ph in [
        ("SBS5", "SBS5 (clock-like) mutations  ·  Δ (post − pre)",
         top_x, top_y, top_w, top_h),
        ("neo_binders", "MHC-I neoantigen binders  ·  Δ (post − pre)",
         bot_x, bot_y, bot_w, bot_h),
    ]:
        sub = paired_long[paired_long.feature == feat].copy()
        sub = sub.dropna(subset=["pre", "post"]).copy()
        sub["delta"] = sub["post"] - sub["pre"]
        # sort by delta ascending
        sub = sub.sort_values("delta").reset_index(drop=True)
        n = len(sub)
        if n == 0:
            continue

        y_lo = float(sub.delta.min()) * 1.15
        y_hi = float(sub.delta.max()) * 1.15
        if y_hi < 0:
            y_hi = abs(y_lo) * 0.2
        if y_lo > 0:
            y_lo = -abs(y_hi) * 0.2

        def ty(v):
            return _i(py + ph - (v - y_lo) / (y_hi - y_lo) * ph)

        def tx(i):
            return _i(px + (i + 0.5) / n * pw)

        # spines
        add_line(s, px, py, px, py + ph, LINE, 0.5)
        # x axis at 0
        add_line(s, px, ty(0), px + pw, ty(0), INK, 1.0)
        # y ticks at min/max/0
        for yv in [y_lo, (y_lo + y_hi) / 2, 0, y_hi / 2 if y_hi > 0 else 0,
                   y_hi]:
            yy = ty(yv)
            add_line(s, _i(px - Inches(0.04)), yy, px, yy, LINE, 0.4)
            label_txt = f"{yv:+.0f}" if abs(yv) >= 10 else f"{yv:+.1f}"
            add_text(s, _i(px - Inches(0.50)), yy - Inches(0.08),
                     Inches(0.45), Inches(0.14),
                     label_txt, size=6, color=INK, align="right")

        # bars
        bar_w = int(pw / n * 0.75)
        for i, row in sub.iterrows():
            color = GOOD if row.response == "good" else BAD
            cx = tx(i)
            x0 = _i(cx - bar_w / 2)
            y_top = ty(max(0, row.delta))
            y_bot = ty(min(0, row.delta))
            add_rect(s, x0, y_top, bar_w, y_bot - y_top,
                     fill=color, line_color=WHITE, line_width=0.3)
            # subject id label (below axis 0 for negatives, above for positives)
            lab_y = ty(0) + Inches(0.02) if row.delta >= 0 else ty(0) - Inches(0.15)
            add_text(s, cx - Inches(0.12), lab_y,
                     Inches(0.24), Inches(0.12),
                     str(int(row.subject_id)), size=4,
                     color=color, bold=True, align="center")

        # title label below sub-plot
        add_text(s, px, _i(py + ph + Inches(0.14)),
                 pw, Inches(0.14),
                 label, size=7, bold=True, color=INK, align="center")

    # Legend (right edge)
    leg_x = Inches(5.95); leg_y = Inches(0.50)
    add_rect(s, leg_x, leg_y, Inches(0.18), Inches(0.13), fill=GOOD)
    add_text(s, leg_x, leg_y + Inches(0.15),
             Inches(0.5), Inches(0.12),
             "good", size=6, bold=True, color=GOOD, align="center")
    add_rect(s, leg_x, leg_y + Inches(0.35), Inches(0.18), Inches(0.13),
             fill=BAD)
    add_text(s, leg_x, leg_y + Inches(0.49),
             Inches(0.5), Inches(0.12),
             "bad", size=6, bold=True, color=BAD, align="center")

    # bottom caption
    add_text(s, Inches(0.15), Inches(4.08),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "Per-subject Δ sorted ascending. Subjects #s shown near bar bases. "
             "Within good: SBS5 Δ median −76 (CI [−145, −64]),",
             size=6, color=INK, align="center")
    add_text(s, Inches(0.15), Inches(4.23),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "neo_binders Δ median −312 (CI [−626, −123]).  Within-good "
             "CIs exclude zero; between-group CIs span zero → exploratory.",
             size=6, italic=True, color=RGBColor(0x55, 0x55, 0x55),
             align="center")


# -------------------------------------------------------------------
# Panel E --- HLA class I LOH clone clearance
# -------------------------------------------------------------------
def build_E():
    s = new_slide(prs_main)
    draw_panel_letter(s, "E")

    # Load HLA-LOH paired data
    loh_path = f"{DATA}/03_hla/loh_stricter/paired_LOH_change_strict.tsv"
    loh = pd.read_csv(loh_path, sep="\t")

    # Plot area
    px = Inches(1.45); py = Inches(0.65)
    pw = Inches(4.05); ph = Inches(2.75)

    # Y axis: LOH locus count (0 to 3)
    y_lo, y_hi = -0.15, 3.15

    def ty(v):
        return _i(py + ph - (v - y_lo) / (y_hi - y_lo) * ph)

    # X axis: pre (0.3) / post (0.7)
    def tx(v):
        return _i(px + v * pw)

    # Spines + y ticks
    add_line(s, px, py, px, py + ph, LINE, 0.6)
    add_line(s, px, py + ph, px + pw, py + ph, LINE, 0.6)
    for yv in [0, 1, 2, 3]:
        yy = ty(yv)
        add_line(s, _i(px - Inches(0.04)), yy, px, yy, LINE, 0.4)
        add_text(s, _i(px - Inches(0.40)), yy - Inches(0.08),
                 Inches(0.36), Inches(0.16),
                 str(yv), size=7, color=INK, align="right")
    # y title (rotated)
    yt = add_text(s, _i(px - Inches(0.85)),
                  _i(py + ph / 2 - Inches(1.0)),
                  Inches(0.45), Inches(2.0),
                  "HLA class I LOH loci per patient  (strict)",
                  size=8, color=INK, align="center")
    yt.rotation = -90

    # X tick labels
    for xv, lab in [(0.25, "pre-CRT"), (0.75, "post-CRT")]:
        xx = tx(xv)
        add_line(s, xx, _i(py + ph), xx,
                 _i(py + ph + Inches(0.05)), LINE, 0.4)
        add_text(s, xx - Inches(0.35), _i(py + ph + Inches(0.06)),
                 Inches(0.70), Inches(0.16),
                 lab, size=8, bold=True, color=INK, align="center")

    # Zero reference
    add_text(s, _i(px + pw + Inches(0.02)), ty(0) - Inches(0.08),
             Inches(0.80), Inches(0.16),
             "no LOH",
             size=6, italic=True, color=RGBColor(0x66, 0x66, 0x66),
             align="left")

    # Highlight "clearance" arrows for subj 3 and 4 (both good responders,
    # both strict LOH at baseline, both resolved by post-CRT)
    focus_subs = loh[loh.loh_resolved == True].copy()
    other_subs = loh[loh.loh_resolved == False].copy()

    # Plot non-focus subjects as faint lines at y=0 (most have no LOH)
    jitter = 0.015
    rng = np.random.default_rng(8)
    for _, row in other_subs.iterrows():
        c_col = (lighten(GOOD_HEX, 0.78) if row.response_bin == "good"
                 else lighten(BAD_HEX, 0.78))
        # jitter slightly to avoid overlap
        jx_pre = tx(0.25) + rng.integers(-_i(Inches(0.04)),
                                          _i(Inches(0.04)) + 1)
        jx_post = tx(0.75) + rng.integers(-_i(Inches(0.04)),
                                          _i(Inches(0.04)) + 1)
        # only plot if LOH value > 0 or we want to show all as faint dots
        if row.pre_loh > 0 or row.post_loh > 0:
            add_line(s, jx_pre, ty(row.pre_loh),
                     jx_post, ty(row.post_loh),
                     c_col, 0.7)
            add_circle(s, jx_pre, ty(row.pre_loh), Emu(18000),
                       fill=WHITE,
                       line_color=(GOOD if row.response_bin == "good" else BAD),
                       line_width=0.7)
            add_circle(s, jx_post, ty(row.post_loh), Emu(18000),
                       fill=WHITE,
                       line_color=(GOOD if row.response_bin == "good" else BAD),
                       line_width=0.7)
        else:
            # plot as faint dots at 0 with slight jitter
            add_circle(s, jx_pre, ty(0), Emu(12000),
                       fill=c_col, line_color=None)
            add_circle(s, jx_post, ty(0), Emu(12000),
                       fill=c_col, line_color=None)

    # Plot focus subjects (subj 3, 4) with bold teal trajectories + labels
    for _, row in focus_subs.iterrows():
        sid = int(row.subject_id)
        # slight offset so the two focus lines don't exactly overlap
        x_offset = Inches(0.10) if sid == 3 else -Inches(0.10)
        xp = tx(0.25) + _i(x_offset)
        xq = tx(0.75) + _i(x_offset)
        # arrow from pre to post
        add_line(s, xp, ty(row.pre_loh), xq, ty(row.post_loh),
                 GOOD, 2.2)
        # filled diamonds at endpoints
        add_diamond(s, xp, ty(row.pre_loh), Emu(45000),
                    fill=GOOD, line_color=WHITE, line_width=1.2)
        add_diamond(s, xq, ty(row.post_loh), Emu(45000),
                    fill=GOOD, line_color=WHITE, line_width=1.2)
        # subject label above pre endpoint
        add_text(s, xp - Inches(0.40), ty(row.pre_loh) - Inches(0.22),
                 Inches(0.80), Inches(0.18),
                 f"subject #{sid}",
                 size=8, bold=True, color=GOOD, align="center")
        # value labels at endpoints
        add_text(s, xp - Inches(0.22), ty(row.pre_loh) + Inches(0.02),
                 Inches(0.20), Inches(0.14),
                 f"{int(row.pre_loh)}",
                 size=6, bold=True, color=GOOD, align="right")
        add_text(s, xq + Inches(0.04), ty(row.post_loh) - Inches(0.07),
                 Inches(0.20), Inches(0.14),
                 f"{int(row.post_loh)}",
                 size=6, bold=True, color=GOOD, align="left")

    # Summary box in top-right corner
    add_rounded_rect(s, Inches(4.85), Inches(0.55),
                     Inches(1.50), Inches(0.85),
                     fill=lighten(GOOD_HEX, 0.85),
                     line_color=GOOD, line_width=0.8)
    add_text(s, Inches(4.90), Inches(0.60),
             Inches(1.40), Inches(0.14),
             "Strict HLA-LOH",
             size=6, bold=True, color=GOOD, align="center")
    add_text(s, Inches(4.90), Inches(0.75),
             Inches(1.40), Inches(0.14),
             "pre: 2 / 16 good  ·  0 / 12 bad",
             size=6, color=INK, align="center")
    add_text(s, Inches(4.90), Inches(0.90),
             Inches(1.40), Inches(0.14),
             "post: 0 / 16 good",
             size=6, color=INK, align="center")
    add_text(s, Inches(4.90), Inches(1.05),
             Inches(1.40), Inches(0.14),
             "(both subjects",
             size=5, italic=True,
             color=RGBColor(0x55, 0x55, 0x55), align="center")
    add_text(s, Inches(4.90), Inches(1.18),
             Inches(1.40), Inches(0.14),
             "cleared their LOH)",
             size=5, italic=True,
             color=RGBColor(0x55, 0x55, 0x55), align="center")

    # Legend (bottom)
    leg_y = Inches(3.75)
    add_diamond(s, Inches(0.55), leg_y + Inches(0.06), Emu(35000),
                fill=GOOD, line_color=WHITE, line_width=1.0)
    add_text(s, Inches(0.72), leg_y - Emu(5000),
             Inches(2.0), Inches(0.16),
             "subj 3, 4 — strict LOH → resolved",
             size=7, bold=True, color=GOOD, align="left")
    add_circle(s, Inches(3.40), leg_y + Inches(0.06), Emu(14000),
               fill=lighten(GOOD_HEX, 0.78), line_color=None)
    add_text(s, Inches(3.50), leg_y - Emu(5000),
             Inches(2.2), Inches(0.16),
             "other subjects — no strict LOH detected",
             size=7, color=RGBColor(0x55, 0x55, 0x55), align="left")

    # caption
    add_text(s, Inches(0.15), Inches(4.25),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "Bonferroni-corrected IMGT allele-count criterion "
             "(|Δratio| ≥ 0.20, Fisher P < 0.01). Fisher P = 0.49 "
             "between-group; anecdotal.",
             size=6, italic=True, color=RGBColor(0x55, 0x55, 0x55),
             align="center")


# -------------------------------------------------------------------
# Panel F --- Cascade schematic + convergence-null callout
# -------------------------------------------------------------------
def build_F():
    s = new_slide(prs_main)
    draw_panel_letter(s, "F")

    # ===========================================================
    # Left half: proposed cascade flow diagram
    # ===========================================================
    # Boxes arranged as: RT → clone clearance → antigen clearance →
    # immune reprogramming → B-cell infiltration → response
    stages = [
        ("Radiation\n(50.4 Gy)",
         RGBColor(0x4F, 0x73, 0x8E)),
        ("Clone clearance\n(SBS5, missense ↓)",
         GROUP_COLOR["Clone clearance"]),
        ("Antigen clearance\n(MHC-I neoantigens ↓)",
         GROUP_COLOR["Antigen clearance"]),
        ("Immune reprogramming\n(MHC-II, CD8 exh., Treg ↑)",
         GROUP_COLOR["Immune reprogramming"]),
        ("B-cell infiltration\n(IGH clonotypes ↑)",
         GROUP_COLOR["B-cell infiltration"]),
        ("Final TNT response\n(good responder)",
         HIGHLIGHT),
    ]

    bx = Inches(0.30); by = Inches(0.50)
    box_w = Inches(2.25); box_h = Inches(0.40)
    gap_y = Inches(0.12)

    for i, (txt, color) in enumerate(stages):
        y = by + i * (box_h + gap_y)
        add_rounded_rect(s, bx, y, box_w, box_h,
                         fill=color, line_color=WHITE, line_width=0.8)
        add_text(s, bx, y + Inches(0.02),
                 box_w, box_h - Inches(0.04),
                 txt, size=6, bold=True, color=WHITE,
                 align="center", anchor="middle")
        # downward arrow between stages
        if i < len(stages) - 1:
            add_arrow(s,
                      bx + box_w / 2,
                      y + box_h + Inches(0.01),
                      bx + box_w / 2,
                      y + box_h + gap_y - Inches(0.01),
                      color=INK, width=1.0)

    # Dashed "star-like" annotation next to Treg box (stage index 3) —
    # the only robust between-group finding
    star_y = by + 3 * (box_h + gap_y) + box_h / 2
    add_text(s, bx + box_w + Inches(0.02),
             _i(star_y - Inches(0.10)),
             Inches(0.4), Inches(0.22),
             "★", size=14, bold=True, color=HIGHLIGHT, align="center")
    add_text(s, bx + box_w + Inches(0.35),
             _i(star_y - Inches(0.10)),
             Inches(0.80), Inches(0.22),
             "Treg CI\nrobust",
             size=5, bold=True, color=HIGHLIGHT, align="left",
             anchor="middle", italic=True)

    # ===========================================================
    # Right half: convergence-null callout  (minimal text — detailed
    # interpretation goes in figure legend per Nature/Cell style)
    # ===========================================================
    rx = Inches(3.40); ry = Inches(0.80)
    rw = Inches(2.95); rh = Inches(2.50)

    # Outer callout box
    add_rect(s, rx, ry, rw, rh,
             fill=RGBColor(0xFB, 0xF0, 0xED),
             line_color=BAD, line_width=1.4)

    # Question (one line, italic, small)
    add_text(s, rx + Inches(0.15), ry + Inches(0.15),
             rw - Inches(0.3), Inches(0.22),
             "Baseline → cascade Δ ?",
             size=9, italic=True, color=INK, align="center",
             anchor="middle")

    # Big NO answer
    add_text(s, rx + Inches(0.15), ry + Inches(0.40),
             rw - Inches(0.3), Inches(0.70),
             "NO",
             size=48, bold=True, color=BAD, align="center",
             anchor="middle")

    # Statistic summary (3 concise lines, no sub-boxes)
    add_text(s, rx + Inches(0.15), ry + Inches(1.20),
             rw - Inches(0.3), Inches(0.18),
             "0 / 36 pairs   P < 0.05",
             size=10, bold=True, color=INK, align="center")
    add_text(s, rx + Inches(0.15), ry + Inches(1.45),
             rw - Inches(0.3), Inches(0.16),
             "DSB → CD8-cyt   r = −0.07",
             size=8, color=BAD, bold=True, align="center")
    # separator
    add_line(s, rx + Inches(0.40), ry + Inches(1.80),
             rx + rw - Inches(0.40), ry + Inches(1.80),
             GREY, 0.5)
    # one-line conclusion
    add_text(s, rx + Inches(0.15), ry + Inches(1.95),
             rw - Inches(0.3), Inches(0.22),
             "Static ⊥ Dynamic",
             size=13, bold=True, italic=True, color=GOOD,
             align="center")
    add_text(s, rx + Inches(0.15), ry + Inches(2.20),
             rw - Inches(0.3), Inches(0.18),
             "(orthogonal, not cascading)",
             size=6, italic=True, color=RGBColor(0x66, 0x66, 0x66),
             align="center")

    # Minimal footer (one line only)
    add_text(s, Inches(0.15), Inches(3.95),
             SLIDE_W - Inches(0.3), Inches(0.16),
             "Cascade is phenomenology, not a causal downstream of the "
             "baseline predictor.",
             size=7, italic=True, color=INK, align="center")


build_A()
build_B()
build_C()
build_D()
build_E()
build_F()
deck_main = f"{OUT}/Fig8_cascade_convergence_native_editable.pptx"
prs_main.save(deck_main)
print(f"wrote {deck_main}")


# ===========================================================================
# SUPP FIGURE S20
# ===========================================================================
prs_supp = Presentation()


def build_S20A():
    """Convergence test scatter matrix: 36 pre-specified pairs ranked by
    |r|, showing Spearman r and P per pair."""
    s = new_slide(prs_supp)
    draw_panel_letter(s, "A")

    # sort by |r| descending
    c = convergence.copy()
    c["abs_r"] = c.spearman_r.abs()
    c = c.sort_values("abs_r", ascending=False).reset_index(drop=True)

    # lollipop plot: x = Spearman r, y = pair rank
    ax_x = Inches(2.25); ax_y = Inches(0.45)
    ax_w = Inches(3.40); ax_h = Inches(3.40)

    x_lo, x_hi = -1.0, 1.0

    def tx(v):
        return _i(ax_x + (v - x_lo) / (x_hi - x_lo) * ax_w)

    n_rows = len(c)
    row_h = ax_h / n_rows

    # spines
    add_line(s, ax_x, ax_y, ax_x, ax_y + ax_h, LINE, 0.3)
    add_line(s, ax_x, ax_y + ax_h, ax_x + ax_w, ax_y + ax_h, LINE, 0.5)
    for v in [-1, -0.5, 0, 0.5, 1]:
        xx = tx(v)
        add_line(s, xx, _i(ax_y + ax_h), xx,
                 _i(ax_y + ax_h + Inches(0.04)), LINE, 0.4)
        add_text(s, xx - Inches(0.15), _i(ax_y + ax_h + Inches(0.05)),
                 Inches(0.30), Inches(0.13),
                 f"{v:+.1f}" if v != 0 else "0",
                 size=6, color=INK, align="center")
    add_text(s, ax_x, _i(ax_y + ax_h + Inches(0.22)),
             ax_w, Inches(0.14),
             "Spearman r  (baseline × cascade Δ, n = 12 paired)",
             size=7, color=INK, align="center")

    # zero ref + P=0.05 bounds (|r|~0.58 at n=12)
    add_line(s, tx(0), ax_y, tx(0), ax_y + ax_h, INK, 0.8)
    add_line(s, tx(0.58), ax_y, tx(0.58), ax_y + ax_h,
             GREY, 0.4, dashed=True)
    add_line(s, tx(-0.58), ax_y, tx(-0.58), ax_y + ax_h,
             GREY, 0.4, dashed=True)
    add_text(s, tx(0.58) - Inches(0.30), _i(ax_y - Inches(0.02)),
             Inches(0.6), Inches(0.12),
             "|r| ≈ 0.58\n(P = 0.05)", size=4, italic=True,
             color=RGBColor(0x99, 0x99, 0x99),
             align="center", anchor="top")

    # rows
    for i, row in c.iterrows():
        y = _i(ax_y + (i + 0.5) * row_h)
        r = float(row.spearman_r)
        p = float(row.spearman_p)
        r_c = max(x_lo, min(x_hi, r))

        # lollipop stem
        color = (GOOD if r >= 0 else BAD)
        add_line(s, tx(0), y, tx(r_c), y, color, 0.7)
        # dot
        is_sig = p < 0.05
        add_circle(s, tx(r_c), y, Emu(15000) if not is_sig else Emu(22000),
                   fill=WHITE if not is_sig else color,
                   line_color=color, line_width=0.8)
        # row label
        baseline = str(row.baseline)[:20]
        cascade = str(row.cascade).replace("_delta", " Δ")[:18]
        add_text(s, Inches(0.20), y - Inches(0.06),
                 ax_x - Inches(0.22), Inches(0.13),
                 f"{baseline} → {cascade}",
                 size=4, color=INK, align="right")
        # right column r, P
        add_text(s, _i(ax_x + ax_w + Inches(0.05)),
                 y - Inches(0.06),
                 Inches(0.7), Inches(0.13),
                 f"r={r:+.2f} P={p:.2f}",
                 size=4, color=RGBColor(0x55, 0x55, 0x55), align="left")

    # headline annotation: 0 / 36
    add_rect(s, Inches(0.20), Inches(0.25), Inches(2.0), Inches(0.45),
             fill=WHITE, line_color=BAD, line_width=1.0)
    add_text(s, Inches(0.20), Inches(0.25),
             Inches(2.0), Inches(0.22),
             "0 / 36 pairs  P < 0.05",
             size=9, bold=True, color=BAD, align="center", anchor="middle")
    add_text(s, Inches(0.20), Inches(0.45),
             Inches(2.0), Inches(0.22),
             "(1.8 expected by chance)",
             size=6, italic=True, color=RGBColor(0x55, 0x55, 0x55),
             align="center", anchor="middle")

    # bottom caption
    add_text(s, Inches(0.15), Inches(4.02),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "Pre-specified 36-pair test: 9 baseline Thread-1 features "
             "× 4 cascade Δ features, n = 12 paired subjects.",
             size=6, italic=True, color=RGBColor(0x55, 0x55, 0x55),
             align="center")
    add_text(s, Inches(0.15), Inches(4.18),
             SLIDE_W - Inches(0.3), Inches(0.14),
             "|r| at n = 12 needs ≈ 0.58 for P < 0.05; the largest "
             "observed |r| is ≈ 0.56 → absence of association, not "
             "under-power.",
             size=6, color=INK, align="center")


def build_S20B():
    """Placeholder: purity-adjusted sensitivity (Δ before vs after
    CNVkit purity adjustment). If data available, plot pairs; else
    show summary table."""
    s = new_slide(prs_supp)
    draw_panel_letter(s, "B")

    # Try to load purity-adjusted file
    ptable_path = f"{DATA}/09_integration/paired_delta/delta_purity_sensitivity.tsv"
    if not os.path.exists(ptable_path):
        add_text(s, Inches(0.5), Inches(2.0),
                 SLIDE_W - Inches(1.0), Inches(0.5),
                 "Purity-adjusted sensitivity table not available.",
                 size=10, color=INK, align="center")
        return

    sens = pd.read_csv(ptable_path, sep="\t")
    # print columns for reference
    features_to_show = ["Treg", "MHC_II", "CD8_exhaustion", "IGH_n"]
    # layout: per-feature mini bar chart raw vs adjusted
    # fall back to text if layout complex

    # simple horizontal comparison: raw Δ vs purity-adjusted Δ
    # If column 'delta_good_median_raw' and 'delta_good_median_adj' exist:
    cols = sens.columns.tolist()
    add_text(s, Inches(0.20), Inches(0.50),
             SLIDE_W - Inches(0.4), Inches(0.16),
             "Purity-adjusted paired Δ sensitivity",
             size=8, bold=True, color=INK, align="center")

    # Render as a simple bar comparison
    ax_x = Inches(1.30); ax_y = Inches(0.85)
    ax_w = Inches(4.00); ax_h = Inches(2.80)

    n_feat = len(features_to_show)
    row_h = ax_h / n_feat

    def tx(v, v_min, v_max):
        return _i(ax_x + (v - v_min) / (v_max - v_min) * ax_w)

    for i, feat in enumerate(features_to_show):
        y = _i(ax_y + (i + 0.5) * row_h)
        feat_row = sens[sens.feature == feat]
        if feat_row.empty:
            add_text(s, Inches(0.20), y - Inches(0.1),
                     ax_x - Inches(0.22), Inches(0.18),
                     PRETTY[feat], size=7, color=INK, align="right")
            continue
        feat_row = feat_row.iloc[0]
        # Extract raw vs adjusted values --- best-effort column parsing
        raw_cols = [c for c in cols if "raw" in c.lower() or "unadj" in c.lower()]
        adj_cols = [c for c in cols if "adj" in c.lower() and "raw" not in c.lower()]

        add_text(s, Inches(0.20), y - Inches(0.1),
                 ax_x - Inches(0.22), Inches(0.18),
                 PRETTY[feat], size=7, color=INK, align="right")
        # place the feature's raw + adjusted MW P as text
        add_text(s, ax_x + Inches(0.1), y - Inches(0.1),
                 Inches(4.0), Inches(0.18),
                 str(feat_row.to_dict())[:140],
                 size=4, color=RGBColor(0x55, 0x55, 0x55),
                 align="left")


build_S20A()
build_S20B()
deck_supp = f"{OUT}/SuppFig_S20_cascade_convergence_sensitivity_native_editable.pptx"
prs_supp.save(deck_supp)
print(f"wrote {deck_supp}")
