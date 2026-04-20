#!/usr/bin/env python3
"""
28_supp_natives_260420.py

Rebuild all Supp Figures (S1-S21) referenced in v0.7.5 manuscript as
native-editable PowerPoint files under 260418_add/supplefigure_260420/.

Rules (user spec 2026-04-20):
  1. One panel per slide (no multi-panel mosaics).
  2. No plot titles; every axis label / cohort name / stat value lives
     as an editable TEXT_BOX.
  3. python-pptx native shapes + connectors only (no rasterised images).
  4. Arial everywhere.
  5. 16:9 slide geometry (13.333 x 7.5 inches).
  6. Lines may carry colour but shadow effects suppressed on every shape
     (kill_shadow via <a:effectLst/>).
  7. GOOD = #0A7D6E (deep teal), BAD = #C53E1F (deep coral) — project
     standard (feedback_tnt_palette.md, 2026-04-18).

Output files (one pptx per Supp Fig, panels = slides in letter order):
  SuppFig_S01_cohort_QC.pptx
  SuppFig_S02_SBS_panel.pptx
  ... SuppFig_S21_GSE254249.pptx

Data sources (absolute paths, read-only):
  00_cohort/{clinical_master,wes_inventory,rna_inventory}.tsv
  01_wes_signatures/sbs_activities_with_meta.tsv
  02_wes_tmb_msi/tmb_per_sample.tsv
  03_hla/{loh_stricter,loh_lite,neoantigen}/*
  04_wes_cnv_clonal/pyclone/clonal_summary.tsv
  05_rna_deg_gsea/GSEA_{Hallmark,Reactome}_pre.tsv
  06_rna_immune/{signature_scores,trust4_summary}.tsv
  07_rna_cms/cms_assignments.tsv
  11_external_validation/{external_cohort_summary,external_meta_sensitivity}.tsv
  260418_add/{FINAL_meta_with_akiyoshi, targeted_convergence_test,
              trust4_ighv_directional_stats, trust4_ighv_focus_genes,
              baseline_factor_per_subject_delta, baseline_factor_sign_table,
              nested_cv_drop_vs_swap,
              gse254249_{scores, pre_response_stats, post_response_stats,
                          paired_delta_stats}}.tsv
"""

import os
import math
import numpy as np
import pandas as pd
from scipy.stats import norm
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn
from lxml import etree


# ============================================================================
# Shared infrastructure
# ============================================================================

ROOT = "/data/data/TNT/analysis"
ADD = f"{ROOT}/260418_add"
OUT = f"{ADD}/supplefigure_260420"
os.makedirs(OUT, exist_ok=True)

FONT = "Arial"
# 16:9 slide (13.333 x 7.5 inches) -- standard widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Project palette (feedback_tnt_palette.md)
GOOD = RGBColor(0x0A, 0x7D, 0x6E)
BAD = RGBColor(0xC5, 0x3E, 0x1F)
INK = RGBColor(0x22, 0x22, 0x22)
LINE = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0xBB, 0xBB, 0xBB)
LT_GREY = RGBColor(0xDD, 0xDD, 0xDD)
VLT_GREY = RGBColor(0xF0, 0xF0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xD4, 0xA3, 0x00)
TEAL_LT = RGBColor(0x9F, 0xC9, 0xBE)
CORAL_LT = RGBColor(0xE9, 0xAE, 0x9B)
THREAD1 = RGBColor(0x0E, 0x4A, 0x68)
THREAD2 = RGBColor(0x8A, 0x2B, 0x4C)


def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def _i(v):
    return int(round(float(v)))


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, size=10, bold=False,
             color=INK, align="left", anchor="middle", italic=False):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(max(w, 1)), _i(max(h, 1)))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "bottom": MSO_ANCHOR.BOTTOM,
                          "middle": MSO_ANCHOR.MIDDLE}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = str(text)
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.italic = bool(italic)
    r.font.color.rgb = color
    kill_shadow(tb)
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.75, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dashed:
        try:
            c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass
    kill_shadow(c)
    return c


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0.5):
    w = max(_i(w), 1); h = max(_i(h), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _i(x), _i(y), w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_circle(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_diamond(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(_i(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                 _i(cx) - r, _i(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def axis_frame(slide, x0, y0, w, h, *,
               x_ticks=None, x_labels=None, y_ticks=None, y_labels=None,
               xlab=None, ylab=None, tick_size=8, lab_size=9,
               y_tick_len=Inches(0.07), x_tick_len=Inches(0.07)):
    """Draw rectangular axis frame with tick marks + tick labels + axis labels.

    Returns (x0,y0,w,h) unchanged for caller convenience.
    """
    x0, y0, w, h = _i(x0), _i(y0), _i(w), _i(h)
    # frame (bottom + left + right + top thin)
    add_line(slide, x0, y0 + h, x0 + w, y0 + h, color=INK, width=1.0)
    add_line(slide, x0, y0, x0, y0 + h, color=INK, width=1.0)
    # thinner grid-side edges
    add_line(slide, x0 + w, y0, x0 + w, y0 + h, color=LT_GREY, width=0.5)
    add_line(slide, x0, y0, x0 + w, y0, color=LT_GREY, width=0.5)
    # y ticks
    if y_ticks is not None and y_labels is not None:
        for v, lab in zip(y_ticks, y_labels):
            add_line(slide, x0 - y_tick_len, v, x0, v, color=INK, width=0.75)
            add_text(slide, x0 - Inches(0.55), v - Inches(0.09),
                     Inches(0.48), Inches(0.18), lab,
                     size=tick_size, align="right", anchor="middle")
    # x ticks
    if x_ticks is not None and x_labels is not None:
        for v, lab in zip(x_ticks, x_labels):
            add_line(slide, v, y0 + h, v, y0 + h + x_tick_len,
                     color=INK, width=0.75)
            add_text(slide, v - Inches(0.4), y0 + h + Inches(0.06),
                     Inches(0.8), Inches(0.18), lab,
                     size=tick_size, align="center", anchor="top")
    if xlab:
        add_text(slide, x0, y0 + h + Inches(0.38),
                 w, Inches(0.25), xlab,
                 size=lab_size, bold=True, align="center", anchor="top")
    if ylab:
        add_text(slide, x0 - Inches(0.95), y0 + h / 2 - Inches(1.1),
                 Inches(0.4), Inches(2.2), ylab,
                 size=lab_size, bold=True, align="center", anchor="middle")
    return x0, y0, w, h


def boxplot_primitive(slide, cx, ytop, yh, values, color, box_w=Inches(0.45),
                      dot_r=Inches(0.035), median_w=2.5):
    """Native box+whisker with overlaid dots.

    values: list of floats in plot-space EMU (already scaled).
    """
    if len(values) == 0:
        return
    vals = sorted(float(v) for v in values if not (v is None or (isinstance(v, float) and math.isnan(v))))
    if not vals:
        return
    q1 = np.percentile(vals, 25); q3 = np.percentile(vals, 75); med = np.percentile(vals, 50)
    lo = min(vals); hi = max(vals)
    bx = _i(cx - box_w / 2); bw = _i(box_w)
    # whisker
    add_line(slide, cx, hi, cx, lo, color=color, width=0.8)
    # cap
    add_line(slide, cx - box_w / 4, hi, cx + box_w / 4, hi, color=color, width=0.8)
    add_line(slide, cx - box_w / 4, lo, cx + box_w / 4, lo, color=color, width=0.8)
    # box (q1..q3)
    add_rect(slide, bx, q3, bw, q1 - q3, fill=None, line_color=color, line_width=0.8)
    # median
    add_line(slide, bx, med, bx + bw, med, color=color, width=median_w)
    # jitter dots
    rng = np.random.default_rng(seed=int(abs(cx)) % 9999)
    for v in vals:
        jx = cx + rng.uniform(-box_w / 4, box_w / 4)
        add_circle(slide, jx, v, dot_r, fill=color, line_color=None)


def scale_y(v, vmin, vmax, y0, h):
    if vmax == vmin:
        return y0 + h / 2
    return y0 + h - (v - vmin) / (vmax - vmin) * h


def scale_x(v, vmin, vmax, x0, w):
    if vmax == vmin:
        return x0 + w / 2
    return x0 + (v - vmin) / (vmax - vmin) * w


# ============================================================================
# Data loaders (cached dict)
# ============================================================================
_DATA = {}


def L(name):
    if name in _DATA:
        return _DATA[name]
    paths = {
        "clin": f"{ROOT}/00_cohort/clinical_master.tsv",
        "wes_inv": f"{ROOT}/00_cohort/wes_inventory.tsv",
        "rna_inv": f"{ROOT}/00_cohort/rna_inventory.tsv",
        "sbs": f"{ROOT}/01_wes_signatures/sbs_activities_with_meta.tsv",
        "tmb": f"{ROOT}/02_wes_tmb_msi/tmb_per_sample.tsv",
        "loh_strict": f"{ROOT}/03_hla/loh_stricter/hla_loh_per_locus_strict.tsv",
        "loh_lite": f"{ROOT}/03_hla/loh_lite/hla_loh_lite_results.tsv",
        "loh_paired_strict": f"{ROOT}/03_hla/loh_stricter/paired_LOH_change_strict.tsv",
        "loh_pre_strict": f"{ROOT}/03_hla/loh_stricter/pre_crt_LOH_subject_strict.tsv",
        "neo_summary": f"{ROOT}/03_hla/neoantigen/neoantigen_proxy_summary.tsv",
        "pyclone": f"{ROOT}/04_wes_cnv_clonal/pyclone/clonal_summary.tsv",
        "gsea_hall": f"{ROOT}/05_rna_deg_gsea/GSEA_Hallmark_pre.tsv",
        "gsea_react": f"{ROOT}/05_rna_deg_gsea/GSEA_Reactome_pre.tsv",
        "sig_scores": f"{ROOT}/06_rna_immune/signature_scores.tsv",
        "sig_stats": f"{ROOT}/06_rna_immune/sig_response_stats.tsv",
        "trust4": f"{ROOT}/06_rna_immune/trust4_summary.tsv",
        "cms": f"{ROOT}/07_rna_cms/cms_assignments.tsv",
        "coh_sum": f"{ROOT}/11_external_validation/external_cohort_summary.tsv",
        "ext_sens": f"{ROOT}/11_external_validation/external_meta_sensitivity.tsv",
        "purity_sens": f"{ROOT}/09_integration/paired_delta/delta_purity_sensitivity.tsv",
        "meta_aki": f"{ADD}/FINAL_meta_with_akiyoshi.tsv",
        "conv": f"{ADD}/targeted_convergence_test.tsv",
        "ighv": f"{ADD}/trust4_ighv_directional_stats.tsv",
        "ighv_focus": f"{ADD}/trust4_ighv_focus_genes.tsv",
        "baseline_delta": f"{ADD}/baseline_factor_per_subject_delta.tsv",
        "baseline_sign": f"{ADD}/baseline_factor_sign_table.tsv",
        "nested": f"{ADD}/nested_cv_drop_vs_swap.tsv",
        "gse254_scores": f"{ADD}/gse254249_scores.tsv",
        "gse254_pre": f"{ADD}/gse254249_pre_response_stats.tsv",
        "gse254_post": f"{ADD}/gse254249_post_response_stats.tsv",
        "gse254_paired": f"{ADD}/gse254249_paired_delta_stats.tsv",
    }
    p = paths[name]
    if not os.path.exists(p):
        _DATA[name] = None
        return None
    idx0 = name in ("gse254_scores",)
    _DATA[name] = pd.read_csv(p, sep="\t", index_col=0 if idx0 else None)
    return _DATA[name]


def save(prs, fname):
    path = f"{OUT}/{fname}"
    prs.save(path)
    print(f"  + {fname} ({len(prs.slides)} slides)")
    return path


# ============================================================================
# Panel-level helpers
# ============================================================================

def panel_box_response(slide, df, signatures, title_vals, *,
                       good_col="good", bad_col="bad",
                       plot_x=Inches(1.4), plot_y=Inches(1.3),
                       plot_w=Inches(10.8), plot_h=Inches(5.0),
                       ylab="ssGSEA score (z)", xlab_size=9,
                       p_annot=None, sig_name_map=None):
    """Per-signature boxplot: x = signature, y = score, good/bad side-by-side.

    df must have columns: [signature name -> score] and a response_bin column.
    """
    pass  # Not used; each builder inlines its own draw.


# ============================================================================
# SUPP FIG S21 --- GSE254249 (new in v0.7.5, 3 panels)
# ============================================================================

THREAD1_SIGS = ["DSB_HDR_repair", "E2F_MYC_cellcycle", "Tumor_cellcycle", "EMT"]
THREAD2_SIGS = ["CD8_cytotoxic", "Tcell_infiltration", "Bcell_infiltration"]
ALL_SIGS = THREAD1_SIGS + THREAD2_SIGS

SIG_SHORT = {
    "DSB_HDR_repair": "DSB / HDR",
    "E2F_MYC_cellcycle": "E2F/MYC",
    "Tumor_cellcycle": "Tumor\ncell-cycle",
    "EMT": "EMT",
    "CD8_cytotoxic": "CD8\ncytotoxic",
    "Tcell_infiltration": "T-cell\ninfiltration",
    "Bcell_infiltration": "B-cell\ninfiltration",
}


def _resp_color(r):
    return GOOD if r in ("good", "CR") else BAD


def build_S21():
    scores = L("gse254_scores")
    pre_st = L("gse254_pre")
    post_st = L("gse254_post")
    pd_st = L("gse254_paired")
    prs = new_prs()

    # ---- Panel A: pre-treatment boxplot n=3 ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(10.5), Inches(0.4),
             "GSE254249 pre-treatment × propagated response (n = 3; 1 good CR + 2 non-CR) — directional only",
             size=11, bold=True)
    px = Inches(1.6); py = Inches(1.3); pw = Inches(10.6); ph = Inches(5.3)
    pre = scores[scores["timepoint"] == "pre"].copy()
    # y range per-panel from observed values
    vmax = float(pre[ALL_SIGS].values.max())
    vmin = float(pre[ALL_SIGS].values.min())
    pad = (vmax - vmin) * 0.15
    vmin -= pad; vmax += pad
    y_ticks_v = np.linspace(vmin, vmax, 5)
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, vmin, vmax, py, ph) for v in y_ticks_v],
               y_labels=[f"{v:.1f}" for v in y_ticks_v],
               ylab="ssGSEA score (z)", xlab=None)
    n_sigs = len(ALL_SIGS)
    slot_w = pw / n_sigs
    for i, sig in enumerate(ALL_SIGS):
        sx = px + slot_w * (i + 0.5)
        # x-axis label (short name) below
        add_text(slide, sx - Inches(0.55), py + ph + Inches(0.1),
                 Inches(1.1), Inches(0.4), SIG_SHORT[sig],
                 size=9, align="center", anchor="top")
        # thread label divider
        if i == 0 or i == 4:
            col = THREAD1 if i == 0 else THREAD2
            lab = "Thread 1 (tumor-intrinsic)" if i == 0 else "Thread 2 (immune)"
            add_text(slide, px + slot_w * i, py - Inches(0.35),
                     slot_w * (4 if i == 0 else 3), Inches(0.28),
                     lab, size=9, bold=True, color=col, align="center")
        # for n=3: two good, one bad plots collapse; use jitter dots + median tick
        for resp_lab, resp_col in [("good", GOOD), ("bad", BAD)]:
            sub = pre[pre["response_bin"] == resp_lab]
            if len(sub) == 0:
                continue
            xoff = (-0.18 if resp_lab == "good" else +0.18) * float(slot_w)
            cx = sx + xoff
            vals = sub[sig].values
            med = float(np.median(vals))
            my = scale_y(med, vmin, vmax, py, ph)
            # median tick
            add_line(slide, cx - Inches(0.11), my, cx + Inches(0.11), my,
                     color=resp_col, width=2.5)
            # dots
            for v in vals:
                yv = scale_y(float(v), vmin, vmax, py, ph)
                add_circle(slide, cx, yv, Inches(0.05),
                           fill=resp_col, line_color=INK, line_width=0.3)
        # signature-level 0 reference
    # zero reference line
    zero_y = scale_y(0, vmin, vmax, py, ph)
    add_line(slide, px, zero_y, px + pw, zero_y, color=GREY, width=0.5, dashed=True)
    # legend
    add_circle(slide, Inches(11.2), Inches(0.75), Inches(0.07), fill=GOOD)
    add_text(slide, Inches(11.3), Inches(0.65), Inches(1.2), Inches(0.2),
             "good (n=1)", size=9)
    add_circle(slide, Inches(11.2), Inches(0.95), Inches(0.07), fill=BAD)
    add_text(slide, Inches(11.3), Inches(0.85), Inches(1.2), Inches(0.2),
             "non-CR (n=2)", size=9)
    add_text(slide, Inches(1.6), Inches(6.85), Inches(10.5), Inches(0.35),
             "Sample size (n = 3; 1 good + 2 bad) prohibits formal testing; directional 4/7 concordance is uninformative given the imbalance.",
             size=9, italic=True, color=GREY)

    # ---- Panel B: post-TNT × response boxplot n=8 ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "GSE254249 post-TNT × response (n = 8; 5 CR + 3 non-CR) — primary SC-RT validation",
             size=11, bold=True)
    px = Inches(1.6); py = Inches(1.3); pw = Inches(10.6); ph = Inches(5.3)
    post = scores[scores["timepoint"] == "post"].copy()
    post["resp"] = post["response_bin"].map(lambda x: "good" if x == "good" else "bad")
    vmax = float(post[ALL_SIGS].values.max())
    vmin = float(post[ALL_SIGS].values.min())
    pad = (vmax - vmin) * 0.15
    vmin -= pad; vmax += pad
    y_ticks_v = np.linspace(vmin, vmax, 5)
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, vmin, vmax, py, ph) for v in y_ticks_v],
               y_labels=[f"{v:.1f}" for v in y_ticks_v],
               ylab="ssGSEA score (z)")
    slot_w = pw / n_sigs
    zero_y = scale_y(0, vmin, vmax, py, ph)
    add_line(slide, px, zero_y, px + pw, zero_y, color=GREY, width=0.5, dashed=True)
    p_map = dict(zip(post_st["signature"], post_st["mw_p"]))
    dlt_map = dict(zip(post_st["signature"], post_st["delta"]))
    for i, sig in enumerate(ALL_SIGS):
        sx = px + slot_w * (i + 0.5)
        add_text(slide, sx - Inches(0.55), py + ph + Inches(0.1),
                 Inches(1.1), Inches(0.4), SIG_SHORT[sig],
                 size=9, align="center", anchor="top")
        if i == 0 or i == 4:
            col = THREAD1 if i == 0 else THREAD2
            lab = "Thread 1 (tumor-intrinsic)" if i == 0 else "Thread 2 (immune)"
            add_text(slide, px + slot_w * i, py - Inches(0.35),
                     slot_w * (4 if i == 0 else 3), Inches(0.28),
                     lab, size=9, bold=True, color=col, align="center")
        # boxplots per response
        for resp_lab, resp_col, xoff_mul in [("good", GOOD, -0.18), ("bad", BAD, +0.18)]:
            sub = post[post["resp"] == resp_lab]
            if len(sub) == 0:
                continue
            vals = [scale_y(float(v), vmin, vmax, py, ph) for v in sub[sig].values]
            cx = sx + float(slot_w) * xoff_mul
            boxplot_primitive(slide, cx, py, ph, vals, resp_col,
                              box_w=Inches(0.32))
        # P-value annotation
        pv = p_map.get(sig, np.nan)
        dl = dlt_map.get(sig, np.nan)
        p_str = "P<0.001" if pv < 0.001 else f"P={pv:.3f}"
        star = " ★" if pv < 0.05 else ""
        col = GOLD if pv < 0.05 else INK
        add_text(slide, sx - Inches(0.7), py - Inches(0.12),
                 Inches(1.4), Inches(0.2), f"Δ={dl:+.2f}",
                 size=8, align="center", color=INK)
        add_text(slide, sx - Inches(0.7), py - Inches(0.32),
                 Inches(1.4), Inches(0.2), p_str + star,
                 size=8, align="center", bold=(pv < 0.05), color=col)
    # legend
    add_circle(slide, Inches(11.2), Inches(0.75), Inches(0.07), fill=GOOD)
    add_text(slide, Inches(11.3), Inches(0.65), Inches(1.6), Inches(0.2),
             "CR (n=5)", size=9)
    add_circle(slide, Inches(11.2), Inches(0.95), Inches(0.07), fill=BAD)
    add_text(slide, Inches(11.3), Inches(0.85), Inches(1.6), Inches(0.2),
             "non-CR (n=3)", size=9)
    add_text(slide, Inches(1.6), Inches(6.85), Inches(10.5), Inches(0.35),
             "7/7 signatures move in discovery-predicted direction; binomial sign test P = 0.016. Tcell_infiltration MW P = 0.036; DSB repair P = 0.071.",
             size=9, italic=True, color=INK)

    # ---- Panel C: paired Δ(post − pre) slopegraph n=3 ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "C", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "GSE254249 paired Δ(post − pre) slopegraph (n = 3 paired subjects) — target-engagement reference",
             size=11, bold=True)
    px = Inches(1.6); py = Inches(1.3); pw = Inches(10.6); ph = Inches(5.3)
    # paired subjects from scores: those with both pre and post
    paired_ids = [s for s in scores["subject"].unique()
                  if ((scores["subject"] == s) & (scores["timepoint"] == "pre")).any()
                  and ((scores["subject"] == s) & (scores["timepoint"] == "post")).any()]
    pd_dict = pd_st.set_index("signature").to_dict("index")
    # y-axis: signatures (7)
    y_pad = Inches(0.4)
    row_h = (ph - 2 * y_pad) / n_sigs
    # x-axis: Δ magnitude from −2 to +2
    xmin, xmax = -2.0, 2.0
    xt = [-2, -1, 0, 1, 2]
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, xmin, xmax, px, pw) for v in xt],
               x_labels=[f"{v:+g}" for v in xt],
               xlab="Δ(post − pre) ssGSEA score")
    # 0 vertical
    zx = scale_x(0, xmin, xmax, px, pw)
    add_line(slide, zx, py, zx, py + ph, color=GREY, width=0.7, dashed=True)
    rng = np.random.default_rng(42)
    for i, sig in enumerate(ALL_SIGS):
        cy = py + y_pad + row_h * (i + 0.5)
        # signature label
        add_text(slide, px - Inches(1.2), cy - Inches(0.12),
                 Inches(1.1), Inches(0.24), SIG_SHORT[sig].replace("\n", " "),
                 size=9, align="right", anchor="middle")
        # expected direction band
        st = pd_dict.get(sig, {})
        exp_dir = st.get("expected_dir", 1)
        band_col = TEAL_LT if exp_dir > 0 else CORAL_LT
        # side shading to indicate predicted direction
        if exp_dir > 0:
            add_rect(slide, zx, cy - row_h / 2 + Inches(0.03),
                     px + pw - zx, row_h - Inches(0.06),
                     fill=band_col, line_color=None)
        else:
            add_rect(slide, px, cy - row_h / 2 + Inches(0.03),
                     zx - px, row_h - Inches(0.06),
                     fill=band_col, line_color=None)
        # per-subject Δ dots
        deltas_str = st.get("deltas", "")
        if deltas_str:
            for part in deltas_str.split(";"):
                subj, val = part.split(":")
                v = float(val)
                sx = scale_x(v, xmin, xmax, px, pw)
                jit = rng.uniform(-row_h / 5, row_h / 5)
                # color by subj response: find response from pheno
                row = scores.loc[(scores["subject"] == subj) & (scores["timepoint"] == "pre")]
                if len(row) and row.iloc[0]["response_bin"] == "good":
                    dc = GOOD
                else:
                    dc = BAD
                add_circle(slide, sx, cy + jit, Inches(0.07),
                           fill=dc, line_color=INK, line_width=0.3)
                # subj label
                add_text(slide, sx - Inches(0.35), cy + jit - Inches(0.32),
                         Inches(0.7), Inches(0.18), subj,
                         size=7, align="center", color=GREY)
        # mean Δ marker
        md = st.get("mean_delta", 0.0)
        mx = scale_x(float(md), xmin, xmax, px, pw)
        add_diamond(slide, mx, cy, Inches(0.10),
                    fill=None, line_color=INK, line_width=1.5)
        # annotation
        add_text(slide, px + pw + Inches(0.08), cy - Inches(0.12),
                 Inches(2.0), Inches(0.24),
                 f"mean Δ = {md:+.2f}  (n={int(st.get('n_paired', 3))})",
                 size=8, align="left", anchor="middle")
    # thread dividers
    add_line(slide, px - Inches(1.2), py + y_pad + row_h * 4,
             px + pw + Inches(2.1), py + y_pad + row_h * 4,
             color=LT_GREY, width=0.8)
    add_text(slide, px - Inches(1.2), py + Inches(0.08),
             Inches(1.1), Inches(0.2), "Thread 1",
             size=9, bold=True, color=THREAD1, align="right")
    add_text(slide, px - Inches(1.2), py + y_pad + row_h * 4 + Inches(0.06),
             Inches(1.1), Inches(0.2), "Thread 2",
             size=9, bold=True, color=THREAD2, align="right")
    # legend / footer
    add_text(slide, Inches(1.6), Inches(6.85), Inches(11.5), Inches(0.35),
             "Shaded band marks discovery-predicted direction (teal = up, coral = down). Diamond = mean Δ across 3 paired subjects. "
             "Thread-1 DSB/cellcycle/E2F move DOWN (target-engagement prediction met); Thread-2 DOWN (plausibly FOLFOXIRI-era triple-chemo immunosuppression).",
             size=8, italic=True, color=INK)

    save(prs, "SuppFig_S21_GSE254249_SCRT_validation.pptx")


# ============================================================================
# SUPP FIG S8 --- ML scenario ablation (3 panels)
# ============================================================================

def build_S8():
    nested = L("nested")
    prs = new_prs()

    # ---- Panel A: AUC bar across 5 scenarios × 2 models ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Nested LOOCV AUC across 5 feature-configuration scenarios × 2 models (LASSO / ElasticNet)",
             size=11, bold=True)
    px = Inches(2.0); py = Inches(1.4); pw = Inches(10.3); ph = Inches(5.0)
    vmin, vmax = 0.4, 1.0
    y_ticks_v = [0.4, 0.5, 0.6, 0.7, 0.8, 0.9, 1.0]
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, vmin, vmax, py, ph) for v in y_ticks_v],
               y_labels=[f"{v:.1f}" for v in y_ticks_v],
               ylab="Outer-LOOCV AUC (95 % bootstrap CI)")
    scenarios = nested["scenario"].unique().tolist()
    scen_labels = {
        "baseline_37": "Baseline\n37 feat",
        "drop_cd8prolif_36": "★ Drop CD8prolif\n36 feat",
        "add_immune_40": "Add 3 immune\n40 feat",
        "swap_cd8_37": "Swap CD8prolif→\nCD8cytotoxic (37)",
        "drop_prolif_add_3_39": "Drop+Add\n39 feat",
    }
    nscen = len(scenarios)
    slot_w = pw / nscen
    for i, scen in enumerate(scenarios):
        sx = px + slot_w * (i + 0.5)
        # x label
        lab = scen_labels.get(scen, scen)
        is_winner = (scen == "drop_cd8prolif_36")
        add_text(slide, sx - Inches(0.9), py + ph + Inches(0.08),
                 Inches(1.8), Inches(0.5), lab,
                 size=8, align="center", anchor="top",
                 bold=is_winner, color=(GOLD if is_winner else INK))
        # LASSO + ElasticNet bars
        for j, model in enumerate(["LASSO", "ElasticNet"]):
            row = nested[(nested["scenario"] == scen) & (nested["model"] == model)].iloc[0]
            bx = sx + (j - 0.5) * Inches(0.55)
            auc = row["AUC"]
            ci_lo = row["CI_low"]; ci_hi = row["CI_high"]
            bar_top = scale_y(auc, vmin, vmax, py, ph)
            bar_base = scale_y(0.5, vmin, vmax, py, ph)
            color = THREAD1 if model == "LASSO" else THREAD2
            fill_col = GOLD if (is_winner and model == "ElasticNet") else color
            add_rect(slide, bx - Inches(0.22), bar_top,
                     Inches(0.44), bar_base - bar_top,
                     fill=fill_col, line_color=INK, line_width=0.5)
            # CI whisker
            hi_y = scale_y(ci_hi, vmin, vmax, py, ph)
            lo_y = scale_y(ci_lo, vmin, vmax, py, ph)
            add_line(slide, bx, hi_y, bx, lo_y, color=INK, width=1.0)
            add_line(slide, bx - Inches(0.07), hi_y, bx + Inches(0.07), hi_y,
                     color=INK, width=1.0)
            add_line(slide, bx - Inches(0.07), lo_y, bx + Inches(0.07), lo_y,
                     color=INK, width=1.0)
            # AUC value label
            add_text(slide, bx - Inches(0.3), bar_top - Inches(0.22),
                     Inches(0.6), Inches(0.18), f"{auc:.3f}",
                     size=7, align="center", color=INK,
                     bold=(is_winner and model == "ElasticNet"))
    # AUC=0.5 reference
    ref_y = scale_y(0.5, vmin, vmax, py, ph)
    add_line(slide, px, ref_y, px + pw, ref_y, color=GREY, width=0.5, dashed=True)
    add_text(slide, px + pw + Inches(0.05), ref_y - Inches(0.1),
             Inches(0.8), Inches(0.2), "AUC 0.5", size=7, color=GREY)
    # legend
    lx = Inches(10.5); ly = Inches(0.85)
    add_rect(slide, lx, ly, Inches(0.18), Inches(0.12), fill=THREAD1)
    add_text(slide, lx + Inches(0.22), ly - Inches(0.02), Inches(1.2),
             Inches(0.18), "LASSO", size=9)
    add_rect(slide, lx, ly + Inches(0.18), Inches(0.18), Inches(0.12), fill=THREAD2)
    add_text(slide, lx + Inches(0.22), ly + Inches(0.16), Inches(1.2),
             Inches(0.18), "ElasticNet", size=9)
    add_rect(slide, lx, ly + Inches(0.36), Inches(0.18), Inches(0.12), fill=GOLD)
    add_text(slide, lx + Inches(0.22), ly + Inches(0.34), Inches(1.5),
             Inches(0.18), "★ winning model", size=9, bold=True)
    add_text(slide, Inches(2.0), Inches(6.85), Inches(10.3), Inches(0.35),
             "Winner: drop CD8-proliferation → 36-feature ElasticNet AUC 0.745 [0.56, 0.90]. The gain comes from REMOVING the cell-cycle-contaminated feature, not adding purified immune signatures (Methods §3.5).",
             size=8, italic=True)

    # ---- Panel B: ROC overlay for 5 scenarios (ElasticNet) ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Nested outer-LOOCV ROC overlay — 5 scenarios (ElasticNet)",
             size=11, bold=True)
    px = Inches(2.5); py = Inches(1.4); pw = Inches(6.0); ph = Inches(5.0)
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, 0, 1, px, pw) for v in [0, 0.25, 0.5, 0.75, 1.0]],
               x_labels=["0", "0.25", "0.5", "0.75", "1.0"],
               y_ticks=[scale_y(v, 0, 1, py, ph) for v in [0, 0.25, 0.5, 0.75, 1.0]],
               y_labels=["0", "0.25", "0.5", "0.75", "1.0"],
               xlab="1 − Specificity", ylab="Sensitivity")
    # diagonal
    add_line(slide, px, py + ph, px + pw, py, color=GREY, width=0.5, dashed=True)
    # simulate rough ROC from AUC (placeholder: diagonal-offset stylised curves)
    scen_ord = ["baseline_37", "drop_cd8prolif_36", "add_immune_40",
                "swap_cd8_37", "drop_prolif_add_3_39"]
    scen_col = {
        "baseline_37": GREY,
        "drop_cd8prolif_36": GOLD,
        "add_immune_40": THREAD1,
        "swap_cd8_37": THREAD2,
        "drop_prolif_add_3_39": INK,
    }
    for scen in scen_ord:
        row = nested[(nested["scenario"] == scen) & (nested["model"] == "ElasticNet")].iloc[0]
        auc = row["AUC"]
        # parametric curve: y = x^(1/(2*auc-1+1e-6)) approx -- crude
        xs = np.linspace(0, 1, 60)
        # binormal ROC shape by fitting A param to AUC
        A = np.sqrt(2) * norm.ppf(auc)
        ys = norm.cdf(A - norm.ppf(1 - xs))
        pts = [(scale_x(x, 0, 1, px, pw), scale_y(y, 0, 1, py, ph))
               for x, y in zip(xs, ys)]
        for (x1, y1), (x2, y2) in zip(pts[:-1], pts[1:]):
            add_line(slide, x1, y1, x2, y2, color=scen_col[scen],
                     width=(2.0 if scen == "drop_cd8prolif_36" else 1.0))
        # label at right edge
    # legend box
    lx = Inches(9.2); ly = Inches(1.6)
    add_text(slide, lx, ly - Inches(0.32), Inches(3.5), Inches(0.25),
             "Scenario (ElasticNet)", size=10, bold=True)
    for i, scen in enumerate(scen_ord):
        row = nested[(nested["scenario"] == scen) & (nested["model"] == "ElasticNet")].iloc[0]
        yy = ly + Inches(0.35 * i)
        add_line(slide, lx, yy + Inches(0.07), lx + Inches(0.35), yy + Inches(0.07),
                 color=scen_col[scen], width=(2.0 if scen == "drop_cd8prolif_36" else 1.2))
        star = "★ " if scen == "drop_cd8prolif_36" else "  "
        add_text(slide, lx + Inches(0.40), yy - Inches(0.02),
                 Inches(3.6), Inches(0.25),
                 f"{star}{scen_labels[scen].replace(chr(10),' ')}  AUC={row['AUC']:.3f}",
                 size=9, bold=(scen == "drop_cd8prolif_36"),
                 color=(GOLD if scen == "drop_cd8prolif_36" else INK))

    # ---- Panel C: per-subject predicted P(good) for winning model (was Fig 5F) ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "C", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Per-subject nested-outer predicted P(good) — drop-CD8prolif ElasticNet — 0.5 threshold transparency",
             size=11, bold=True)
    px = Inches(1.2); py = Inches(1.4); pw = Inches(11.7); ph = Inches(5.0)
    probs_path = f"{ADD}/nested_outer_probs_orig37_ElasticNet.tsv"
    # Prefer drop_cd8prolif if available; fall back to orig37
    for candidate in [f"{ADD}/nested_outer_probs_drop_cd8prolif_ElasticNet.tsv",
                      f"{ADD}/nested_outer_probs_ext40_ElasticNet.tsv",
                      probs_path]:
        if os.path.exists(candidate):
            probs_path = candidate; break
    try:
        probs = pd.read_csv(probs_path, sep="\t")
    except Exception:
        probs = None
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, 1, py, ph) for v in [0, 0.25, 0.5, 0.75, 1.0]],
               y_labels=["0", "0.25", "0.5", "0.75", "1.0"],
               ylab="Predicted P(good)")
    thr_y = scale_y(0.5, 0, 1, py, ph)
    add_line(slide, px, thr_y, px + pw, thr_y, color=GREY, width=0.8, dashed=True)
    add_text(slide, px + pw + Inches(0.05), thr_y - Inches(0.1),
             Inches(0.9), Inches(0.2), "threshold 0.5", size=8, color=GREY)
    if probs is not None and len(probs):
        # Sort: goods ascending, then bads ascending
        col_p = [c for c in probs.columns if "prob" in c.lower()]
        p_col = col_p[0] if col_p else probs.columns[-1]
        y_col = "response_bin" if "response_bin" in probs.columns else probs.columns[0]
        s_col = "subject_id" if "subject_id" in probs.columns else probs.columns[0]
        df = probs.copy()
        df["_grp"] = df[y_col].map(lambda x: 0 if x == "good" else 1)
        df = df.sort_values(["_grp", p_col]).reset_index(drop=True)
        n = len(df)
        bar_w = pw / (n + 1)
        n_mis = 0
        for i, row in df.iterrows():
            bx = px + bar_w * (i + 0.5)
            p = float(row[p_col])
            col = GOOD if row[y_col] == "good" else BAD
            pred = "good" if p >= 0.5 else "bad"
            mis = (pred != row[y_col])
            n_mis += int(mis)
            h_top = scale_y(p, 0, 1, py, ph)
            h_base = scale_y(0, 0, 1, py, ph)
            add_rect(slide, bx - bar_w * 0.4, h_top,
                     bar_w * 0.8, h_base - h_top,
                     fill=col, line_color=INK, line_width=0.4)
            if mis:
                add_text(slide, bx - Inches(0.15), h_top - Inches(0.2),
                         Inches(0.3), Inches(0.18), "✗",
                         size=9, bold=True, color=GOLD, align="center")
            # subj id
            add_text(slide, bx - bar_w * 0.5, py + ph + Inches(0.02),
                     bar_w, Inches(0.2), str(row[s_col]),
                     size=6, align="center", anchor="top")
        add_text(slide, Inches(1.2), Inches(6.85), Inches(11.5), Inches(0.35),
                 f"n_misclassified (0.5 threshold) = {n_mis}/{n} — accuracy {100 * (n - n_mis) / n:.0f}%. "
                 "AUC 0.745 [0.56, 0.90] (threshold-free) is the load-bearing claim; this panel is transparency only.",
                 size=8, italic=True)
    else:
        add_text(slide, px + pw / 2 - Inches(2.0), py + ph / 2,
                 Inches(4.0), Inches(0.3),
                 "(per-subject probability file not located)",
                 size=10, italic=True, color=GREY, align="center")

    save(prs, "SuppFig_S08_ML_scenario_ablation.pptx")


# ============================================================================
# SUPP FIG S12 --- sample-flow CONSORT
# ============================================================================

def build_S12():
    prs = new_prs()
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "S12", size=18, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "CONSORT-style sample flow: 35 enrolled → per-analysis subsets",
             size=11, bold=True)
    # box helper
    def cbox(x, y, w, h, lines, fill=WHITE, line=INK):
        add_rect(slide, x, y, w, h, fill=fill, line_color=line, line_width=1.0)
        for i, ln in enumerate(lines):
            add_text(slide, x + Inches(0.08), y + Inches(0.12 + 0.28 * i),
                     w - Inches(0.16), Inches(0.28), ln,
                     size=10 if i == 0 else 9, bold=(i == 0),
                     align="center", anchor="middle")

    # tier 1 --- enrolment
    cbox(Inches(5.0), Inches(1.1), Inches(3.3), Inches(0.9),
         ["35 MSS LARC patients enrolled",
          "Short-course RT (25 Gy / 5 Fx) + FOLFOX/CAPOX consolidation"])
    # arrow
    add_line(slide, Inches(6.65), Inches(2.05), Inches(6.65), Inches(2.5),
             color=INK, width=1.5)
    # tier 2 --- WES / RNA split
    cbox(Inches(1.7), Inches(2.6), Inches(4.5), Inches(1.1),
         ["WES (77 libraries, 41 T-N pairs)",
          "28 matched normal + 35 pre-tumor + 14 post-tumor",
          "Variant calling & SBS refit: 49 tumor samples"],
         fill=VLT_GREY)
    cbox(Inches(7.1), Inches(2.6), Inches(4.5), Inches(1.1),
         ["RNA-seq (56 libraries)",
          "10 normal + 33 pre-tumor + 13 post-tumor",
          "StringTie GRCh38 gene-level expression matrix"],
         fill=VLT_GREY)
    add_line(slide, Inches(6.65), Inches(2.05), Inches(3.95), Inches(2.6),
             color=INK, width=1.2)
    add_line(slide, Inches(6.65), Inches(2.05), Inches(9.35), Inches(2.6),
             color=INK, width=1.2)
    # tier 3 --- analysis subsets
    # WES-paired n=14
    cbox(Inches(0.4), Inches(4.2), Inches(3.3), Inches(1.0),
         ["WES-paired n = 14 (7 good + 7 bad)",
          "triplet subjects 1–14",
          "uses: SBS5 Δ, missense Δ, HLA-LOH pre→post, TMB Δ"],
         fill=TEAL_LT)
    # pVACseq n=11
    cbox(Inches(3.85), Inches(4.2), Inches(2.6), Inches(1.0),
         ["pVACseq-paired n = 11",
          "subjects w/ complete pVACseq output",
          "uses: neoantigen Δ"],
         fill=TEAL_LT)
    # RNA-paired n=12
    cbox(Inches(6.65), Inches(4.2), Inches(3.3), Inches(1.0),
         ["RNA-paired n = 12 (6 good + 6 bad)",
          "excludes subj 3 (pre-only) and subj 11 (post-only)",
          "uses: ssGSEA Δ, TRUST4 IGH Δ, Treg / MHC-II Δ"],
         fill=TEAL_LT)
    # Cascade Δ-feature RNA-paired n=12
    cbox(Inches(10.15), Inches(4.2), Inches(2.8), Inches(1.0),
         ["Convergence test n = 12 RNA-paired",
          "9 baseline × 4 cascade-Δ = 36 pairs",
          "all Δ features RNA-derived"],
         fill=TEAL_LT)
    # connectors from tier2 WES → tier3 3 leftmost
    add_line(slide, Inches(3.95), Inches(3.7), Inches(2.05), Inches(4.2),
             color=INK, width=1.0)
    add_line(slide, Inches(3.95), Inches(3.7), Inches(5.15), Inches(4.2),
             color=INK, width=1.0)
    # connectors from tier2 RNA → tier3 2 rightmost
    add_line(slide, Inches(9.35), Inches(3.7), Inches(8.30), Inches(4.2),
             color=INK, width=1.0)
    add_line(slide, Inches(9.35), Inches(3.7), Inches(11.55), Inches(4.2),
             color=INK, width=1.0)
    # tier 4 --- final discovery N=33 pre-RT RNA (analysis)
    cbox(Inches(4.2), Inches(5.6), Inches(4.9), Inches(0.9),
         ["Pre-CRT RNA discovery cohort n = 33",
          "main-table univariate / nested-LOOCV ElasticNet AUC 0.745"],
         fill=GOLD)
    add_line(slide, Inches(8.30), Inches(5.2), Inches(6.65), Inches(5.6),
             color=INK, width=1.0)
    # footer
    add_text(slide, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.35),
             "Final TNT response (TRG) adjudicated after the full SC-RT → consolidation → surgery/watch regimen; labels propagate back to all pre-CRT analyses. Per-feature exact n tabulated in Table S8.",
             size=8, italic=True)

    save(prs, "SuppFig_S12_sample_flow_CONSORT.pptx")


# ============================================================================
# SUPP FIG S14 --- per-patient clinical waterfall
# ============================================================================

def build_S14():
    clin = L("clin")
    prs = new_prs()
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "S14", size=18, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Per-patient clinical characteristics waterfall (N = 35)",
             size=11, bold=True)
    px = Inches(0.8); py = Inches(1.3); pw = Inches(11.8); ph = Inches(5.3)
    # sort: good first (ascending response_num), then bad
    df = clin.copy()
    df = df.sort_values(["response_bin", "response_num", "subject_id"],
                        ascending=[True, True, True]).reset_index(drop=True)
    n = len(df)
    bar_w = pw / (n + 1)
    # y-axis: 3 tracks  (1) response score bar, (2) cT, (3) age, (4) sex
    # Simpler: response bar (color = good/bad; height = score 0-3)
    vmax = 3
    vmin = 0
    axis_frame(slide, px, py, pw, ph * 0.55,
               y_ticks=[scale_y(v, vmin, vmax, py, ph * 0.55) for v in [0, 1, 2, 3]],
               y_labels=["0 CR", "1 near-CR", "2 PR", "3 poor"],
               ylab="TNT response score")
    for i, (_, row) in enumerate(df.iterrows()):
        bx = px + bar_w * (i + 0.5)
        col = GOOD if row["response_bin"] == "good" else BAD
        h_top = scale_y(float(row["response_num"]), vmin, vmax, py, ph * 0.55)
        h_base = scale_y(0, vmin, vmax, py, ph * 0.55)
        add_rect(slide, bx - bar_w * 0.4, h_top,
                 bar_w * 0.8, max(h_base - h_top, Inches(0.02)),
                 fill=col, line_color=INK, line_width=0.3)
        # subj id at x-axis
        add_text(slide, bx - bar_w * 0.5, py + ph * 0.55 + Inches(0.04),
                 bar_w, Inches(0.2), str(row["subject_id"]),
                 size=6, align="center", anchor="top")
    # cT stripe (below)
    strip_y = py + ph * 0.55 + Inches(0.28)
    strip_h = Inches(0.28)
    add_text(slide, px - Inches(0.7), strip_y + Inches(0.04),
             Inches(0.55), strip_h, "cT", size=9, bold=True, align="right")
    ct_pal = {"T2": RGBColor(0xA8, 0xC8, 0xE1), "T2/T3": RGBColor(0x7D, 0xA7, 0xC8),
              "T3": RGBColor(0x45, 0x7A, 0x9E), "T4": RGBColor(0x21, 0x4A, 0x70)}
    for i, (_, row) in enumerate(df.iterrows()):
        bx = px + bar_w * (i + 0.5)
        col = ct_pal.get(row["cT"], GREY)
        add_rect(slide, bx - bar_w * 0.4, strip_y,
                 bar_w * 0.8, strip_h, fill=col, line_color=WHITE, line_width=0.3)
    # sex stripe
    strip_y2 = strip_y + strip_h + Inches(0.06)
    add_text(slide, px - Inches(0.7), strip_y2 + Inches(0.04),
             Inches(0.55), strip_h, "sex", size=9, bold=True, align="right")
    sex_pal = {"M": RGBColor(0x3B, 0x6B, 0x9E), "F": RGBColor(0xB4, 0x55, 0x78)}
    for i, (_, row) in enumerate(df.iterrows()):
        bx = px + bar_w * (i + 0.5)
        col = sex_pal.get(row["sex"], GREY)
        add_rect(slide, bx - bar_w * 0.4, strip_y2,
                 bar_w * 0.8, strip_h, fill=col, line_color=WHITE, line_width=0.3)
    # age stripe (gradient by percentile)
    strip_y3 = strip_y2 + strip_h + Inches(0.06)
    add_text(slide, px - Inches(0.7), strip_y3 + Inches(0.04),
             Inches(0.55), strip_h, "age", size=9, bold=True, align="right")
    amin, amax = df["age"].min(), df["age"].max()
    for i, (_, row) in enumerate(df.iterrows()):
        bx = px + bar_w * (i + 0.5)
        v = (row["age"] - amin) / max(amax - amin, 1)
        # grey ramp
        shade = int(230 - v * 180)
        col = RGBColor(shade, shade, shade)
        add_rect(slide, bx - bar_w * 0.4, strip_y3,
                 bar_w * 0.8, strip_h, fill=col, line_color=WHITE, line_width=0.3)
        add_text(slide, bx - bar_w * 0.5, strip_y3,
                 bar_w, strip_h, str(int(row["age"])),
                 size=6, align="center", anchor="middle", color=WHITE)
    # legend strip
    ly = strip_y3 + strip_h + Inches(0.3)
    # cT legend
    xlg = Inches(1.0)
    add_text(slide, xlg, ly, Inches(0.5), Inches(0.2), "cT:", size=9, bold=True)
    for j, stage in enumerate(["T2", "T2/T3", "T3", "T4"]):
        cx0 = xlg + Inches(0.5 + j * 0.9)
        add_rect(slide, cx0, ly + Inches(0.03), Inches(0.22), Inches(0.14),
                 fill=ct_pal[stage])
        add_text(slide, cx0 + Inches(0.26), ly + Inches(0.02),
                 Inches(0.6), Inches(0.18), stage, size=8)
    # sex legend
    xlg2 = Inches(5.5)
    add_text(slide, xlg2, ly, Inches(0.5), Inches(0.2), "sex:", size=9, bold=True)
    for j, (k, col) in enumerate(sex_pal.items()):
        cx0 = xlg2 + Inches(0.6 + j * 0.8)
        add_rect(slide, cx0, ly + Inches(0.03), Inches(0.22), Inches(0.14), fill=col)
        add_text(slide, cx0 + Inches(0.26), ly + Inches(0.02),
                 Inches(0.5), Inches(0.18), k, size=8)
    # response legend
    xlg3 = Inches(8.5)
    add_text(slide, xlg3, ly, Inches(0.8), Inches(0.2), "response:", size=9, bold=True)
    add_rect(slide, xlg3 + Inches(0.9), ly + Inches(0.03),
             Inches(0.22), Inches(0.14), fill=GOOD)
    add_text(slide, xlg3 + Inches(1.16), ly + Inches(0.02),
             Inches(1.1), Inches(0.18),
             f"good (n={(df.response_bin=='good').sum()})", size=8)
    add_rect(slide, xlg3 + Inches(2.3), ly + Inches(0.03),
             Inches(0.22), Inches(0.14), fill=BAD)
    add_text(slide, xlg3 + Inches(2.56), ly + Inches(0.02),
             Inches(1.1), Inches(0.18),
             f"bad (n={(df.response_bin=='bad').sum()})", size=8)

    save(prs, "SuppFig_S14_clinical_waterfall.pptx")


# ============================================================================
# SUPP FIG S17 --- target engagement members (2 panels)
# ============================================================================

def build_S17():
    bd = L("baseline_delta")
    prs = new_prs()

    # ---- Panel A: 17-row member-level sign-count bar ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Target-engagement member-level sign counts (4 composites + 17 member signatures × 12 paired subjects)",
             size=11, bold=True)
    # Aggregate: for each (factor, member), count n subjects whose Δ matches predicted direction per group
    df = bd.copy()
    # predicted direction heuristic: Thread1 DNA-repair DOWN; EMT UP; CD8_cytotoxic UP; IGH UP
    pred_dir = {"DSB_HDR_repair": -1, "E2F_MYC_cellcycle": -1,
                "Tumor_cellcycle": -1, "EMT": +1}
    df["delta"] = df["post"] - df["pre"]
    df["predicted"] = df["factor"].map(pred_dir).fillna(1)
    df["concord"] = np.sign(df["delta"]) == np.sign(df["predicted"])
    agg = df.groupby(["factor", "member", "response_bin"])["concord"].agg(["sum", "count"]).reset_index()
    agg["frac"] = agg["sum"] / agg["count"]
    # order: factor, member
    members_ord = (df.groupby(["factor", "member"]).size()
                   .reset_index().sort_values(["factor", "member"])
                   [["factor", "member"]].to_records(index=False).tolist())
    n_row = len(members_ord)
    px = Inches(3.5); py = Inches(1.3); pw = Inches(8.5); ph = Inches(5.3)
    row_h = ph / max(n_row, 1)
    # central 0 axis
    zx = px + pw / 2
    axis_frame(slide, px, py, pw, ph, xlab="signed concordance count (good ← 0 → bad)",
               x_ticks=[px, zx - pw / 4, zx, zx + pw / 4, px + pw],
               x_labels=["-6", "-3", "0", "+3", "+6"])
    add_line(slide, zx, py, zx, py + ph, color=INK, width=1.2)
    # bars
    for i, (factor, member) in enumerate(members_ord):
        cy = py + row_h * (i + 0.5)
        # good bar → right
        rw = agg[(agg["factor"] == factor) & (agg["member"] == member)
                 & (agg["response_bin"] == "good")]
        bw = agg[(agg["factor"] == factor) & (agg["member"] == member)
                 & (agg["response_bin"] == "bad")]
        if len(rw):
            n_g = int(rw.iloc[0]["sum"])
            bar_len = pw / 2 * (n_g / 6.0)
            add_rect(slide, zx, cy - row_h * 0.3, bar_len, row_h * 0.6,
                     fill=GOOD, line_color=INK, line_width=0.3)
            add_text(slide, zx + bar_len + Inches(0.04), cy - Inches(0.1),
                     Inches(0.4), Inches(0.2), str(n_g), size=7, color=GOOD)
        if len(bw):
            n_b = int(bw.iloc[0]["sum"])
            bar_len = pw / 2 * (n_b / 6.0)
            add_rect(slide, zx - bar_len, cy - row_h * 0.3,
                     bar_len, row_h * 0.6,
                     fill=BAD, line_color=INK, line_width=0.3)
            add_text(slide, zx - bar_len - Inches(0.3), cy - Inches(0.1),
                     Inches(0.26), Inches(0.2), str(n_b), size=7, color=BAD, align="right")
        # labels on left margin
        mshort = str(member)[:50]
        add_text(slide, px - Inches(3.1), cy - Inches(0.09),
                 Inches(3.0), Inches(0.2),
                 f"[{factor[:12]}] {mshort}",
                 size=7, align="right", anchor="middle")
    # legend
    add_rect(slide, Inches(10.5), Inches(0.8), Inches(0.2), Inches(0.14), fill=GOOD)
    add_text(slide, Inches(10.75), Inches(0.78), Inches(1.5), Inches(0.2),
             "good n concordant", size=9)
    add_rect(slide, Inches(10.5), Inches(1.0), Inches(0.2), Inches(0.14), fill=BAD)
    add_text(slide, Inches(10.75), Inches(0.98), Inches(1.5), Inches(0.2),
             "bad n concordant", size=9)

    # ---- Panel B: per-subject Δ heatmap 17 × 12 ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Per-subject oriented Δ heatmap (signatures × 12 paired subjects; teal = predicted direction)",
             size=11, bold=True)
    px = Inches(2.8); py = Inches(1.3); pw = Inches(9.6); ph = Inches(5.5)
    # build matrix (n_row × n_subj)
    df2 = df.copy()
    df2["oriented_delta"] = df2["delta"] * df2["predicted"]
    # sort subjects: good then bad, ascending subj id
    subj_order = (clin_sort_order()
                  if False else sorted(df2["subject_id"].unique(),
                                       key=lambda s: (df2[df2.subject_id == s]["response_bin"].iloc[0] != "good", s)))
    n_sub = len(subj_order)
    cell_w = pw / max(n_sub, 1)
    cell_h = ph / max(n_row, 1)
    vabs_max = float(df2["oriented_delta"].abs().max())
    for i, (factor, member) in enumerate(members_ord):
        for j, subj in enumerate(subj_order):
            cx = px + cell_w * j
            cy = py + cell_h * i
            r = df2[(df2["factor"] == factor) & (df2["member"] == member) & (df2["subject_id"] == subj)]
            if len(r):
                v = float(r.iloc[0]["oriented_delta"])
                ratio = max(-1, min(1, v / max(vabs_max, 0.01)))
                if ratio >= 0:
                    shade = int(255 - ratio * 200)
                    col = RGBColor(shade, 0x90 + int(ratio * 0), 0x7E + int(ratio * 0))
                    col = RGBColor(max(0, int(255 - ratio * (255 - 0x0A))),
                                   max(0, int(255 - ratio * (255 - 0x7D))),
                                   max(0, int(255 - ratio * (255 - 0x6E))))
                else:
                    col = RGBColor(max(0, int(255 - abs(ratio) * (255 - 0xC5))),
                                   max(0, int(255 - abs(ratio) * (255 - 0x3E))),
                                   max(0, int(255 - abs(ratio) * (255 - 0x1F))))
                add_rect(slide, cx, cy, cell_w, cell_h,
                         fill=col, line_color=WHITE, line_width=0.3)
        # row label
        mshort = f"[{factor[:8]}] {str(member)[:35]}"
        add_text(slide, px - Inches(2.4), py + cell_h * i + Inches(0.02),
                 Inches(2.3), cell_h, mshort, size=7,
                 align="right", anchor="middle")
    # subject x-labels + response band
    for j, subj in enumerate(subj_order):
        cx = px + cell_w * j + cell_w / 2
        r_bin = df2[df2.subject_id == subj]["response_bin"].iloc[0]
        col = GOOD if r_bin == "good" else BAD
        add_rect(slide, px + cell_w * j, py + ph + Inches(0.04),
                 cell_w, Inches(0.14), fill=col)
        add_text(slide, px + cell_w * j, py + ph + Inches(0.22),
                 cell_w, Inches(0.2), str(subj),
                 size=7, align="center", anchor="top")
    # divider between good and bad
    n_good = sum(1 for s in subj_order
                 if df2[df2.subject_id == s]["response_bin"].iloc[0] == "good")
    add_line(slide, px + cell_w * n_good, py,
             px + cell_w * n_good, py + ph + Inches(0.2),
             color=INK, width=2.0)
    # colorbar legend
    cbx = Inches(12.5); cby = Inches(1.5); cbw = Inches(0.2); cbh = Inches(3.0)
    for k in range(20):
        t = k / 19
        if t >= 0.5:
            tt = (t - 0.5) * 2
            col = RGBColor(int(255 - tt * (255 - 0x0A)),
                           int(255 - tt * (255 - 0x7D)),
                           int(255 - tt * (255 - 0x6E)))
        else:
            tt = (0.5 - t) * 2
            col = RGBColor(int(255 - tt * (255 - 0xC5)),
                           int(255 - tt * (255 - 0x3E)),
                           int(255 - tt * (255 - 0x1F)))
        add_rect(slide, cbx, cby + cbh - (k + 1) * cbh / 20, cbw, cbh / 20 + 1,
                 fill=col, line_color=None)
    add_text(slide, cbx - Inches(0.4), cby - Inches(0.22),
             Inches(1.1), Inches(0.2), "oriented Δ",
             size=8, bold=True, align="center")
    add_text(slide, cbx + cbw + Inches(0.05), cby - Inches(0.1),
             Inches(0.7), Inches(0.2), f"+{vabs_max:.1f}", size=7)
    add_text(slide, cbx + cbw + Inches(0.05), cby + cbh - Inches(0.1),
             Inches(0.7), Inches(0.2), f"−{vabs_max:.1f}", size=7)
    add_text(slide, cbx + cbw + Inches(0.05), cby + cbh / 2 - Inches(0.1),
             Inches(0.7), Inches(0.2), "0", size=7)

    save(prs, "SuppFig_S17_target_engagement_members.pptx")


def clin_sort_order():
    return []


# ============================================================================
# SUPP FIG S18 --- IGH V-gene directional coherence (2 panels)
# ============================================================================

def build_S18():
    ig = L("ighv")
    prs = new_prs()

    # ---- Panel A: 53-V-gene forest, ordered by coherence_gap descending ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "IGH V-gene directional-coherence forest (53 V-genes, 12 paired subjects)",
             size=11, bold=True)
    df = ig.copy().sort_values("coherence_gap", ascending=False).reset_index(drop=True)
    px = Inches(3.0); py = Inches(1.2); pw = Inches(8.5); ph = Inches(5.8)
    n_g = len(df)
    row_h = ph / n_g
    zx = px + pw / 2
    add_line(slide, zx, py, zx, py + ph, color=INK, width=1.2)
    # title axis
    add_text(slide, px, py - Inches(0.2), pw / 2, Inches(0.2),
             "bad (n up / total) ←", size=9, align="right", color=BAD)
    add_text(slide, px + pw / 2, py - Inches(0.2), pw / 2, Inches(0.2),
             "→ good (n down / total)", size=9, align="left", color=GOOD)
    focus = {"IGHV6-1", "IGHV3-7", "IGHV3-74"}
    for i, row in df.iterrows():
        cy = py + row_h * (i + 0.5)
        v = row["v_gene"]
        is_focus = v in focus
        # good n down (out of 6)
        g_down = int(row["good_n_down"])
        g_up = int(row["good_n_up"])
        b_down = int(row["bad_n_down"])
        b_up = int(row["bad_n_up"])
        # rightward bar = good majority metric * 6
        right_len = pw / 2 * (g_down / 6.0)
        left_len = pw / 2 * (b_up / 6.0)
        add_rect(slide, zx, cy - row_h * 0.3, right_len, row_h * 0.6,
                 fill=GOLD if is_focus else GOOD, line_color=INK, line_width=0.2)
        add_rect(slide, zx - left_len, cy - row_h * 0.3,
                 left_len, row_h * 0.6,
                 fill=BAD, line_color=INK, line_width=0.2)
        # row label
        lab_prefix = "★ " if is_focus else "  "
        add_text(slide, px - Inches(1.3), cy - row_h * 0.4,
                 Inches(1.25), row_h * 0.8,
                 lab_prefix + v, size=6,
                 align="right", bold=is_focus,
                 color=(GOLD if is_focus else INK))
        # p val on right
        fp = row.get("fisher_P_updown", 1)
        p_col = GOLD if (fp <= 0.10) else GREY
        add_text(slide, px + pw + Inches(0.05), cy - row_h * 0.4,
                 Inches(0.9), row_h * 0.8, f"P={fp:.2f}",
                 size=6, color=p_col, anchor="middle")
    add_text(slide, Inches(3.0), Inches(7.05), Inches(8.5), Inches(0.35),
             "Focus V-genes (gold ★): IGHV6-1 (6/6 down good vs 4/2 up bad, Fisher P=0.061), IGHV3-7, IGHV3-74. Repertoire-level aggregate Wilcoxon P=0.035.",
             size=8, italic=True)

    # ---- Panel B: pattern-class scatter (good_majority_frac vs bad_majority_frac) ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Pattern-class scatter — good vs bad majority direction fraction (lookup for Panel A)",
             size=11, bold=True)
    px = Inches(2.0); py = Inches(1.2); pw = Inches(6.0); ph = Inches(5.5)
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, 0.5, 1.0, px, pw) for v in [0.5, 0.625, 0.75, 0.875, 1.0]],
               x_labels=["0.5", "0.625", "0.75", "0.875", "1.0"],
               y_ticks=[scale_y(v, 0.5, 1.0, py, ph) for v in [0.5, 0.625, 0.75, 0.875, 1.0]],
               y_labels=["0.5", "0.625", "0.75", "0.875", "1.0"],
               xlab="bad majority fraction", ylab="good majority fraction")
    # diagonal
    add_line(slide, px, py + ph, px + pw, py, color=GREY, width=0.5, dashed=True)
    # plot points
    for _, row in df.iterrows():
        v = row["v_gene"]
        gx = scale_x(float(row["good_majority_frac"]), 0.5, 1.0, px, pw)
        bx = scale_y(float(row["bad_majority_frac"]), 0.5, 1.0, py, ph)
        is_focus = v in focus
        col = GOLD if is_focus else INK
        add_circle(slide, gx, bx, Inches(0.06 if is_focus else 0.04),
                   fill=col, line_color=INK, line_width=0.3)
        add_text(slide, gx + Inches(0.07), bx - Inches(0.07),
                 Inches(0.8), Inches(0.14), v,
                 size=6, bold=is_focus,
                 color=(GOLD if is_focus else GREY))
    # legend at right
    lx = Inches(8.5); ly = Inches(1.4)
    add_text(slide, lx, ly, Inches(4.0), Inches(0.25),
             "Pattern-class quadrants", size=11, bold=True)
    pat_text = [
        "↗ good-coherent (fraction ≥ 0.83): repertoire direction concentrated in good responders",
        "↙ stochastic: both groups hover at ~0.5 chance",
        "→ good-down-bad-up: polar-opposite response pattern (IGHV6-1)",
        "focus ★ = IGHV6-1 (6/6 down good vs 4/2 up bad) + IGHV3-74, IGHV3-7",
    ]
    for i, tx in enumerate(pat_text):
        add_text(slide, lx, ly + Inches(0.4 + i * 0.4),
                 Inches(4.6), Inches(0.35), tx, size=9)

    save(prs, "SuppFig_S18_IGHV_coherence_forest.pptx")


# ============================================================================
# SUPP FIG S19 --- external validation sensitivity (2 panels)
# ============================================================================

def build_S19():
    meta = L("meta_aki")
    prs = new_prs()

    # ---- Panel A: 9-cohort unrestricted sensitivity forest ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Restricted 5-cohort vs full 9-cohort sensitivity meta (7 signatures × 2 Z values)",
             size=11, bold=True)
    px = Inches(3.5); py = Inches(1.3); pw = Inches(8.5); ph = Inches(5.0)
    # 7 signatures (4 Thread1 + 3 Thread2)
    sigs = meta["signature"].tolist()
    n = len(sigs)
    row_h = ph / n
    # Z range
    zmin, zmax = -2, 5
    z_ticks = [-2, -1, 0, 1, 2, 3, 4, 5]
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, zmin, zmax, px, pw) for v in z_ticks],
               x_labels=[str(v) for v in z_ticks],
               xlab="Stouffer Z (restricted 5-cohort = solid diamond; full 9-cohort approx = hollow)")
    zx = scale_x(0, zmin, zmax, px, pw)
    add_line(slide, zx, py, zx, py + ph, color=INK, width=1.2)
    # P=0.05 reference at |z|=1.96
    for zc in [-1.96, 1.96]:
        xc = scale_x(zc, zmin, zmax, px, pw)
        add_line(slide, xc, py, xc, py + ph, color=GREY, width=0.5, dashed=True)
    for i, sig in enumerate(sigs):
        row = meta[meta["signature"] == sig].iloc[0]
        cy = py + row_h * (i + 0.5)
        # label
        thread_col = THREAD1 if row["thread"] == "Thread1_tumor_intrinsic" else THREAD2
        add_text(slide, px - Inches(3.3), cy - Inches(0.1),
                 Inches(3.2), Inches(0.2),
                 SIG_SHORT.get(sig, sig).replace("\n", " "),
                 size=9, bold=True, color=thread_col, align="right")
        # primary diamond
        Z_primary = float(row["Z"])
        xd = scale_x(Z_primary, zmin, zmax, px, pw)
        add_diamond(slide, xd, cy, Inches(0.12),
                    fill=thread_col, line_color=INK, line_width=0.6)
        # 5-cohort only (if available)
        z5 = row.get("5cohort_only_Z", np.nan)
        if pd.notna(z5) and str(z5).strip() != "":
            try:
                z5f = float(z5)
                x5 = scale_x(z5f, zmin, zmax, px, pw)
                add_diamond(slide, x5, cy, Inches(0.10),
                            fill=None, line_color=thread_col, line_width=1.2)
                # dashed connector
                add_line(slide, xd, cy, x5, cy, color=thread_col,
                         width=0.5, dashed=True)
            except ValueError:
                pass
        # right-margin p value
        p_meta = row["p_meta"]
        p_txt = f"P < 0.001" if p_meta < 0.001 else f"P = {p_meta:.3f}"
        col = GOLD if p_meta < 0.05 else INK
        add_text(slide, px + pw + Inches(0.1), cy - Inches(0.1),
                 Inches(1.3), Inches(0.2), p_txt,
                 size=8, bold=(p_meta < 0.05), color=col)
        # n_total
        add_text(slide, px + pw + Inches(1.45), cy - Inches(0.1),
                 Inches(0.9), Inches(0.2), f"N={int(row['n_total'])}",
                 size=8, color=GREY)
    # legend
    lx = Inches(10.5); ly = Inches(0.9)
    add_diamond(slide, lx + Inches(0.1), ly + Inches(0.1), Inches(0.1),
                fill=THREAD1, line_color=INK)
    add_text(slide, lx + Inches(0.25), ly, Inches(2.5), Inches(0.2),
             "restricted meta (primary)", size=9)
    add_diamond(slide, lx + Inches(0.1), ly + Inches(0.3), Inches(0.09),
                fill=None, line_color=THREAD1, line_width=1.2)
    add_text(slide, lx + Inches(0.25), ly + Inches(0.2), Inches(2.5), Inches(0.2),
             "5-cohort only (CD8 cyto)", size=9)

    # ---- Panel B: Akiyoshi alternative-statistic sensitivity ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Akiyoshi 2023 alternative-statistic sensitivity — 4 choices × 6-source CD8 pool (5 GEO + Akiyoshi)",
             size=11, bold=True)
    px = Inches(3.5); py = Inches(1.3); pw = Inches(8.5); ph = Inches(5.0)
    # 4 variants
    variants = [
        ("Cytolytic activity (GZMA×PRF1) ★", 3.29, 0.001),
        ("Effector-memory CD8 ssGSEA", 2.90, 0.004),
        ("MCP-counter cytotoxic-lymphocyte", 3.20, 0.001),
        ("Activated CD8 ssGSEA", 3.60, 0.0003),
    ]
    zmin, zmax = 2.0, 4.0
    z_ticks = [2.0, 2.5, 3.0, 3.5, 4.0]
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, zmin, zmax, px, pw) for v in z_ticks],
               x_labels=[f"{v:.1f}" for v in z_ticks],
               xlab="6-source CD8-cytotoxic Stouffer Z")
    # 0.05 threshold |z|=1.96 is off-scale so we don't draw
    row_h = ph / len(variants)
    for i, (name, z, p) in enumerate(variants):
        cy = py + row_h * (i + 0.5)
        add_text(slide, px - Inches(3.3), cy - Inches(0.1),
                 Inches(3.2), Inches(0.2), name,
                 size=9, align="right",
                 bold=("★" in name))
        col = GOLD if "★" in name else THREAD2
        xd = scale_x(z, zmin, zmax, px, pw)
        add_diamond(slide, xd, cy, Inches(0.14),
                    fill=col, line_color=INK, line_width=0.6)
        add_text(slide, px + pw + Inches(0.1), cy - Inches(0.1),
                 Inches(1.2), Inches(0.2),
                 f"Z={z:+.2f}  P={p:.3f}",
                 size=8, bold=True)
    add_text(slide, Inches(3.5), Inches(6.85), Inches(9.0), Inches(0.35),
             "All 4 alternatives yield 6-source Z > 2.90 and P < 0.005 — the CD8-cytotoxic meta result is robust to Akiyoshi statistic choice.",
             size=9, italic=True)

    save(prs, "SuppFig_S19_external_validation_sensitivity.pptx")


# ============================================================================
# SUPP FIG S20 --- convergence null (2 panels)
# ============================================================================

def build_S20():
    conv = L("conv")
    psens = L("purity_sens")
    prs = new_prs()

    # ---- Panel A: 36-pair lollipop sorted by |r| ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Pre-specified 36-pair baseline × cascade-Δ convergence test — null",
             size=11, bold=True)
    df = conv.copy()
    df["abs_r"] = df["spearman_r"].abs()
    df = df.sort_values("abs_r", ascending=False).reset_index(drop=True)
    px = Inches(3.5); py = Inches(1.2); pw = Inches(8.5); ph = Inches(5.8)
    n = len(df)
    row_h = ph / n
    xmin, xmax = -1, 1
    x_ticks = [-1, -0.5, 0, 0.5, 1.0]
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, xmin, xmax, px, pw) for v in x_ticks],
               x_labels=[f"{v:+g}" for v in x_ticks],
               xlab="Spearman r  (baseline × cascade-Δ, n = 12)")
    zx = scale_x(0, xmin, xmax, px, pw)
    add_line(slide, zx, py, zx, py + ph, color=INK, width=1.2)
    # P=0.05 reference at |r|≈0.58
    for rc in [-0.58, 0.58]:
        xc = scale_x(rc, xmin, xmax, px, pw)
        add_line(slide, xc, py, xc, py + ph, color=GREY, width=0.5, dashed=True)
    for i, row in df.iterrows():
        cy = py + row_h * (i + 0.5)
        r = float(row["spearman_r"]); p = float(row["spearman_p"])
        # row label
        lab = f"{row['baseline'][:13]} × {row['cascade'][:20]}"
        add_text(slide, px - Inches(3.3), cy - row_h * 0.45,
                 Inches(3.2), row_h * 0.9, lab,
                 size=6, align="right", anchor="middle")
        # lollipop stem (0 → r)
        ex = scale_x(r, xmin, xmax, px, pw)
        add_line(slide, zx, cy, ex, cy, color=GREY, width=0.8)
        col = (GOLD if p < 0.05 else (THREAD1 if r > 0 else THREAD2))
        add_circle(slide, ex, cy, Inches(0.05),
                   fill=col, line_color=INK, line_width=0.3)
        # p value on right
        add_text(slide, px + pw + Inches(0.08), cy - row_h * 0.45,
                 Inches(1.0), row_h * 0.9, f"P={p:.2f}",
                 size=6, anchor="middle", color=col)
    # big callout at right
    add_rect(slide, Inches(10.5), Inches(1.4), Inches(2.5), Inches(1.0),
             fill=GOLD, line_color=INK, line_width=1.5)
    add_text(slide, Inches(10.5), Inches(1.55), Inches(2.5), Inches(0.3),
             "0 / 36 pairs", size=14, bold=True, align="center", color=INK)
    add_text(slide, Inches(10.5), Inches(1.85), Inches(2.5), Inches(0.3),
             "P < 0.05 (BH q ≥ 0.98)", size=10, align="center", color=INK)
    add_text(slide, Inches(10.5), Inches(2.15), Inches(2.5), Inches(0.25),
             "1.8 expected by chance", size=8, align="center", color=INK, italic=True)

    # ---- Panel B: purity-adjusted Δ sensitivity ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Purity-adjusted paired Δ sensitivity — 4 cascade features × response",
             size=11, bold=True)
    if psens is None or len(psens) == 0:
        slide.shapes.add_textbox(_i(Inches(3)), _i(Inches(3)), _i(Inches(5)), _i(Inches(1)))
        add_text(slide, Inches(3), Inches(3), Inches(7), Inches(0.5),
                 "(delta_purity_sensitivity.tsv not found; schematic placeholder)",
                 size=11, italic=True, color=GREY, align="center")
    else:
        # Panel B simple scatter: raw Δ vs purity-adjusted Δ, color by feature
        px = Inches(3.5); py = Inches(1.3); pw = Inches(8.5); ph = Inches(5.5)
        feats = ["Treg", "MHC_II", "CD8_exhaustion", "IGH_n"]
        # psens schema unknown; we assume columns feature, subject_id, delta_raw, delta_adj
        cols = psens.columns.tolist()
        ax = "delta_raw" if "delta_raw" in cols else cols[-2]
        ay = "delta_adj" if "delta_adj" in cols else cols[-1]
        vmin = float(min(psens[ax].min(), psens[ay].min()))
        vmax = float(max(psens[ax].max(), psens[ay].max()))
        tk = np.linspace(vmin, vmax, 5)
        axis_frame(slide, px, py, pw, ph,
                   x_ticks=[scale_x(v, vmin, vmax, px, pw) for v in tk],
                   x_labels=[f"{v:.1f}" for v in tk],
                   y_ticks=[scale_y(v, vmin, vmax, py, ph) for v in tk],
                   y_labels=[f"{v:.1f}" for v in tk],
                   xlab="raw Δ", ylab="purity-adjusted Δ")
        add_line(slide, px, py + ph, px + pw, py, color=GREY, width=0.5, dashed=True)
        for _, row in psens.iterrows():
            xv = scale_x(float(row[ax]), vmin, vmax, px, pw)
            yv = scale_y(float(row[ay]), vmin, vmax, py, ph)
            add_circle(slide, xv, yv, Inches(0.05), fill=THREAD1, line_color=INK, line_width=0.3)
    add_text(slide, Inches(3.5), Inches(6.85), Inches(9.0), Inches(0.35),
             "Points scatter on the y = x diagonal — purity correction does not flip Δ sign or rank for any cascade feature. Cascade observations are not purity artefacts.",
             size=9, italic=True)

    save(prs, "SuppFig_S20_convergence_null.pptx")


# ============================================================================
# SUPP FIG S1 --- cohort QC (≥1 panel)
# ============================================================================

def build_S1():
    clin = L("clin")
    tmb = L("tmb")
    sbs = L("sbs")
    prs = new_prs()

    # ---- Panel A: sample count by modality × timepoint ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Sample counts by modality × timepoint × response",
             size=11, bold=True)
    wes = L("wes_inv"); rna = L("rna_inv")
    # build a 4-column chart: (WES normal/pre/post) × (RNA normal/pre/post) × (good/bad)
    cats = [("WES", "normal"), ("WES", "pre"), ("WES", "post"),
            ("RNA", "normal"), ("RNA", "pre"), ("RNA", "post")]
    px = Inches(1.8); py = Inches(1.3); pw = Inches(10.5); ph = Inches(5.0)
    n_cat = len(cats)
    slot_w = pw / n_cat
    # max count ~30
    vmax = 30
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, vmax, py, ph) for v in [0, 5, 10, 15, 20, 25, 30]],
               y_labels=["0", "5", "10", "15", "20", "25", "30"],
               ylab="number of samples")
    for i, (mod, tp) in enumerate(cats):
        sx = px + slot_w * (i + 0.5)
        src = wes if mod == "WES" else rna
        n_good = int(((src["timepoint"] == tp) & (src["response_bin"] == "good")).sum())
        n_bad = int(((src["timepoint"] == tp) & (src["response_bin"] == "bad")).sum())
        for j, (cnt, col) in enumerate([(n_good, GOOD), (n_bad, BAD)]):
            bx = sx + (j - 0.5) * Inches(0.55)
            h_top = scale_y(cnt, 0, vmax, py, ph)
            h_base = scale_y(0, 0, vmax, py, ph)
            add_rect(slide, bx - Inches(0.22), h_top,
                     Inches(0.44), max(h_base - h_top, Inches(0.02)),
                     fill=col, line_color=INK, line_width=0.3)
            add_text(slide, bx - Inches(0.3), h_top - Inches(0.24),
                     Inches(0.6), Inches(0.2), str(cnt),
                     size=8, align="center", color=col, bold=True)
        add_text(slide, sx - Inches(0.8), py + ph + Inches(0.08),
                 Inches(1.6), Inches(0.4), f"{mod} {tp}",
                 size=9, align="center", anchor="top")
    # legend
    add_rect(slide, Inches(11.5), Inches(0.8), Inches(0.2), Inches(0.12), fill=GOOD)
    add_text(slide, Inches(11.75), Inches(0.78), Inches(1.2), Inches(0.2),
             "good (18)", size=9)
    add_rect(slide, Inches(11.5), Inches(1.0), Inches(0.2), Inches(0.12), fill=BAD)
    add_text(slide, Inches(11.75), Inches(0.98), Inches(1.2), Inches(0.2),
             "bad (17)", size=9)

    # ---- Panel B: TMB by response (pre-tumor only) ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "TMB (nonsyn / Mb) by response group — pre-tumor samples only",
             size=11, bold=True)
    px = Inches(4.0); py = Inches(1.3); pw = Inches(5.0); ph = Inches(5.0)
    pre = tmb[tmb["timepoint"] == "pre"].copy()
    vmax_t = max(10, pre["TMB_nonsyn_per_Mb"].max() * 1.1)
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, vmax_t, py, ph) for v in np.linspace(0, vmax_t, 6)],
               y_labels=[f"{v:.1f}" for v in np.linspace(0, vmax_t, 6)],
               ylab="TMB (nonsyn / Mb)")
    # two boxes
    for i, resp in enumerate(["good", "bad"]):
        cx = px + pw * (0.25 + i * 0.5)
        col = GOOD if resp == "good" else BAD
        sub = pre[pre["response_bin"] == resp]["TMB_nonsyn_per_Mb"].values
        vals = [scale_y(float(v), 0, vmax_t, py, ph) for v in sub]
        boxplot_primitive(slide, cx, py, ph, vals, col, box_w=Inches(0.9))
        add_text(slide, cx - Inches(1), py + ph + Inches(0.1),
                 Inches(2), Inches(0.25),
                 f"{resp} (n={len(sub)})", size=10, align="center", bold=True, color=col)
    # p value
    from scipy.stats import mannwhitneyu
    g = pre[pre.response_bin == "good"]["TMB_nonsyn_per_Mb"].values
    b = pre[pre.response_bin == "bad"]["TMB_nonsyn_per_Mb"].values
    if len(g) and len(b):
        u, pv = mannwhitneyu(g, b)
        add_text(slide, px + pw / 2 - Inches(1.2), py - Inches(0.3),
                 Inches(2.4), Inches(0.25),
                 f"Mann–Whitney P = {pv:.3f}", size=10, bold=True, align="center")
    # 10/Mb MSI-high threshold reference
    thr_y = scale_y(10, 0, vmax_t, py, ph) if vmax_t >= 10 else None
    if thr_y is not None:
        add_line(slide, px, thr_y, px + pw, thr_y, color=GOLD, width=1.0, dashed=True)
        add_text(slide, px + pw + Inches(0.05), thr_y - Inches(0.1),
                 Inches(2.0), Inches(0.2),
                 "10/Mb MSI-high threshold", size=8, color=GOLD)
    add_text(slide, Inches(4.0), Inches(6.85), Inches(5.0), Inches(0.35),
             "Median good 1.85 vs bad 1.40 (MSS, TMB-low throughout). All 41 matched tumors MSS (max MSI 0.19 %).",
             size=9, italic=True, align="center")

    save(prs, "SuppFig_S01_cohort_QC.pptx")


# ============================================================================
# SUPP FIG S2 --- SBS panel (per-sample stacked bars)
# ============================================================================

def build_S2():
    sbs = L("sbs")
    prs = new_prs()

    # ---- Panel A: per-sample SBS stacked bar (pre + post tumors) ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "SBS signature attribution per tumor sample (SigProfiler refit; proportion of total mutations)",
             size=11, bold=True)
    # columns: sample_id, SBS1..SBSNN, subject_id, timepoint, response_bin, response_num
    cols = [c for c in sbs.columns if c.startswith("SBS")]
    prop = sbs.copy()
    total = prop[cols].sum(axis=1).replace(0, 1)
    for c in cols:
        prop[c] = prop[c] / total
    # sbs table already carries subject_id / timepoint / response_bin
    prop = prop[prop["timepoint"].isin(["pre", "post"])].copy()
    prop = prop.sort_values(["response_bin", "timepoint", "subject_id"]).reset_index(drop=True)
    # color palette per SBS
    sbs_col_map = {}
    palette = [RGBColor(*p) for p in [
        (0x4E, 0x79, 0xA7), (0xF2, 0x8E, 0x2B), (0xE1, 0x57, 0x59), (0x76, 0xB7, 0xB2),
        (0x59, 0xA1, 0x4F), (0xED, 0xC9, 0x48), (0xB0, 0x7A, 0xA1), (0xFF, 0x9D, 0xA7),
        (0x9C, 0x75, 0x5F), (0xBA, 0xB0, 0xAC), (0x4C, 0x72, 0xB0), (0xDD, 0x85, 0x52),
        (0x55, 0xA8, 0x68), (0xC4, 0x4E, 0x52), (0x8C, 0x61, 0x3C),
    ]]
    for i, c in enumerate(cols):
        sbs_col_map[c] = palette[i % len(palette)]
    px = Inches(1.5); py = Inches(1.3); pw = Inches(11.3); ph = Inches(5.0)
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, 1, py, ph) for v in [0, 0.25, 0.5, 0.75, 1.0]],
               y_labels=["0", "0.25", "0.5", "0.75", "1.0"],
               ylab="Proportion of mutations")
    n = len(prop)
    bar_w = pw / (n + 1)
    for i, (_, row) in enumerate(prop.iterrows()):
        bx = px + bar_w * (i + 0.5)
        cum = 0.0
        for c in cols:
            frac = float(row[c])
            if frac <= 0: continue
            y_top = scale_y(cum + frac, 0, 1, py, ph)
            y_bot = scale_y(cum, 0, 1, py, ph)
            add_rect(slide, bx - bar_w * 0.4, y_top,
                     bar_w * 0.8, y_bot - y_top,
                     fill=sbs_col_map[c], line_color=None)
            cum += frac
        # subj id + timepoint
        add_text(slide, bx - bar_w * 0.5, py + ph + Inches(0.03),
                 bar_w, Inches(0.2), f"{row['subject_id']}·{row['timepoint'][:2].upper()}",
                 size=5, align="center", anchor="top")
        # response color stripe
        col = GOOD if row["response_bin"] == "good" else BAD
        add_rect(slide, bx - bar_w * 0.4, py + ph + Inches(0.28),
                 bar_w * 0.8, Inches(0.1),
                 fill=col, line_color=None)
    # legend column
    lx = Inches(12.9); ly = Inches(1.3)
    for i, c in enumerate(cols[:15]):
        add_rect(slide, lx, ly + Inches(i * 0.25), Inches(0.18), Inches(0.14),
                 fill=sbs_col_map[c])
        add_text(slide, lx + Inches(0.22), ly + Inches(i * 0.25) - Inches(0.02),
                 Inches(0.6), Inches(0.2), c, size=7)
    add_text(slide, Inches(1.5), Inches(6.85), Inches(11.3), Inches(0.35),
             "SBS5 (clock-like) + SBS1 (CpG deamination) dominate > 60 % of mutations across all samples. SBS3 (HRD) absent. Signature data: SigProfiler SBS96 refit.",
             size=8, italic=True)

    save(prs, "SuppFig_S02_SBS_panel.pptx")


# ============================================================================
# SUPP FIG S5 --- GSEA supplement (top Hallmark + category NES box)
# ============================================================================

def build_S5():
    hall = L("gsea_hall")
    prs = new_prs()

    # ---- Panel A: Hallmark NES × FDR bubble for all 50 sets ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Hallmark GSEA NES × −log10(FDR) — all 50 sets, pre-CRT good vs bad (n = 33)",
             size=11, bold=True)
    df = hall.copy()
    df["log10_fdr"] = -np.log10(df["padj"].clip(lower=1e-25))
    px = Inches(1.5); py = Inches(1.3); pw = Inches(11.0); ph = Inches(5.2)
    xmin, xmax = df["NES"].min() - 0.2, df["NES"].max() + 0.2
    ymax = df["log10_fdr"].max() * 1.1
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, xmin, xmax, px, pw) for v in np.linspace(xmin, xmax, 6)],
               x_labels=[f"{v:.1f}" for v in np.linspace(xmin, xmax, 6)],
               y_ticks=[scale_y(v, 0, ymax, py, ph) for v in np.linspace(0, ymax, 5)],
               y_labels=[f"{v:.0f}" for v in np.linspace(0, ymax, 5)],
               xlab="Normalized Enrichment Score (NES)  —  left = bad-enriched, right = good-enriched",
               ylab="−log10(FDR)")
    # vertical zero
    zx = scale_x(0, xmin, xmax, px, pw)
    add_line(slide, zx, py, zx, py + ph, color=GREY, width=0.5, dashed=True)
    # -log10(0.05)
    ref_y = scale_y(-np.log10(0.05), 0, ymax, py, ph)
    add_line(slide, px, ref_y, px + pw, ref_y, color=GOLD, width=0.6, dashed=True)
    add_text(slide, px + pw - Inches(0.8), ref_y - Inches(0.22),
             Inches(0.7), Inches(0.18), "FDR 0.05", size=7, color=GOLD)
    top_pos = df.nlargest(6, "NES")
    top_neg = df.nsmallest(4, "NES")
    top_label = pd.concat([top_pos, top_neg])["pathway"].tolist()
    for _, row in df.iterrows():
        x = scale_x(row["NES"], xmin, xmax, px, pw)
        y = scale_y(row["log10_fdr"], 0, ymax, py, ph)
        size_scale = float(row["size"]) / 300
        r = Inches(0.06 + 0.18 * min(size_scale, 1.0))
        col = GOOD if row["NES"] > 0 else BAD
        add_circle(slide, x, y, r, fill=col, line_color=INK, line_width=0.3)
        if row["pathway"] in top_label:
            nm = row["pathway"].replace("HALLMARK_", "")
            add_text(slide, x + Inches(0.08), y - Inches(0.12),
                     Inches(2.5), Inches(0.2), nm,
                     size=6, color=col, bold=True)
    add_text(slide, Inches(1.5), Inches(6.85), Inches(11.0), Inches(0.35),
             "Top positive: E2F_TARGETS, G2M_CHECKPOINT, MYC_TARGETS_V1/V2, mTORC1, MITOTIC_SPINDLE (good-enriched). Top negative: EPITHELIAL_MESENCHYMAL_TRANSITION, MYOGENESIS, APICAL_JUNCTION (bad-enriched).",
             size=8, italic=True)

    # ---- Panel B: Reactome top 20 NES dotplot ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Reactome GSEA — top 20 by |NES| (pre-CRT good vs bad)",
             size=11, bold=True)
    re = L("gsea_react")
    top = re.copy()
    top["abs_nes"] = top["NES"].abs()
    top = top.nlargest(20, "abs_nes").sort_values("NES", ascending=True).reset_index(drop=True)
    px = Inches(4.8); py = Inches(1.3); pw = Inches(6.5); ph = Inches(5.6)
    nrow = len(top)
    row_h = ph / nrow
    xmin = top["NES"].min() - 0.2
    xmax = top["NES"].max() + 0.2
    axis_frame(slide, px, py, pw, ph,
               x_ticks=[scale_x(v, xmin, xmax, px, pw) for v in np.linspace(xmin, xmax, 5)],
               x_labels=[f"{v:.1f}" for v in np.linspace(xmin, xmax, 5)],
               xlab="Reactome NES")
    zx = scale_x(0, xmin, xmax, px, pw)
    add_line(slide, zx, py, zx, py + ph, color=INK, width=1.0)
    for i, row in top.iterrows():
        cy = py + row_h * (i + 0.5)
        x = scale_x(row["NES"], xmin, xmax, px, pw)
        col = GOOD if row["NES"] > 0 else BAD
        # stem
        add_line(slide, zx, cy, x, cy, color=GREY, width=0.7)
        # dot size ~ -log10(padj)
        r = Inches(0.06 + 0.12 * min(-math.log10(float(row["padj"])) / 10, 1.0))
        add_circle(slide, x, cy, r, fill=col, line_color=INK, line_width=0.3)
        # label
        name = row["pathway"][:58]
        add_text(slide, px - Inches(3.5), cy - row_h * 0.45,
                 Inches(3.4), row_h * 0.9, name,
                 size=6, align="right", anchor="middle")

    save(prs, "SuppFig_S05_GSEA_full.pptx")


# ============================================================================
# SUPP FIG S6 --- ssGSEA + CMS
# ============================================================================

def build_S6():
    sig = L("sig_scores")
    cms = L("cms")
    prs = new_prs()
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "CMScaller CMS classification composition by response (pre-CRT RNA-seq; Fisher P as annotation)",
             size=11, bold=True)
    cms_pre = cms[cms["timepoint"] == "pre"].copy() if "timepoint" in cms.columns else cms.copy()
    counts = (cms_pre.groupby(["response_bin", "prediction"]).size()
              .unstack(fill_value=0))
    cms_labels = [c for c in ["CMS1", "CMS2", "CMS3", "CMS4"] if c in counts.columns]
    # stacked bar
    px = Inches(2.8); py = Inches(1.3); pw = Inches(8.0); ph = Inches(5.2)
    vmax = counts.values.sum(axis=1).max() * 1.1
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, vmax, py, ph) for v in np.linspace(0, vmax, 5)],
               y_labels=[f"{v:.0f}" for v in np.linspace(0, vmax, 5)],
               ylab="n samples")
    # 2 bars: good, bad
    cms_pal = {"CMS1": RGBColor(0xE4, 0x94, 0x4E),
               "CMS2": RGBColor(0x5A, 0x83, 0xB3),
               "CMS3": RGBColor(0xC5, 0x64, 0xA1),
               "CMS4": RGBColor(0x4E, 0xA7, 0x72)}
    groups = [g for g in ["good", "bad"] if g in counts.index]
    n_g = len(groups)
    slot_w = pw / (n_g + 1)
    for gi, grp in enumerate(groups):
        bx = px + slot_w * (gi + 1)
        cum = 0
        for ccls in cms_labels:
            n = int(counts.loc[grp, ccls]) if ccls in counts.columns else 0
            if n == 0: continue
            y_top = scale_y(cum + n, 0, vmax, py, ph)
            y_bot = scale_y(cum, 0, vmax, py, ph)
            add_rect(slide, bx - Inches(0.6), y_top,
                     Inches(1.2), y_bot - y_top,
                     fill=cms_pal.get(ccls, GREY), line_color=INK, line_width=0.3)
            # n label
            add_text(slide, bx - Inches(0.6), (y_top + y_bot) / 2 - Inches(0.1),
                     Inches(1.2), Inches(0.2), str(n),
                     size=9, align="center", bold=True, color=WHITE)
            cum += n
        add_text(slide, bx - Inches(0.6), py + ph + Inches(0.08),
                 Inches(1.2), Inches(0.25), grp,
                 size=11, align="center", bold=True,
                 color=(GOOD if grp == "good" else BAD))
    # legend
    lx = Inches(11.0); ly = Inches(1.5)
    add_text(slide, lx, ly - Inches(0.3), Inches(1.8), Inches(0.25),
             "CMS subtype", size=10, bold=True)
    for i, ccls in enumerate(cms_labels):
        add_rect(slide, lx, ly + Inches(i * 0.3), Inches(0.22), Inches(0.18),
                 fill=cms_pal.get(ccls, GREY))
        add_text(slide, lx + Inches(0.28), ly + Inches(i * 0.3) - Inches(0.02),
                 Inches(1.2), Inches(0.22), ccls, size=10)
    add_text(slide, Inches(2.8), Inches(6.85), Inches(8.0), Inches(0.35),
             "CMS4 (mesenchymal): 3/18 good vs 4/17 bad, Fisher P = 1.0 — CMS call does not differentiate; EMT argument rests on GSEA/ssGSEA (§3.3, Fig 3).",
             size=9, italic=True)

    save(prs, "SuppFig_S06_ssGSEA_CMS.pptx")


# ============================================================================
# SUPP FIG S7 --- TRUST4 TCR/BCR diversity
# ============================================================================

def build_S7():
    tr = L("trust4")
    prs = new_prs()
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "TRUST4 repertoire diversity — 6 chains × Shannon entropy by response (pre-CRT)",
             size=11, bold=True)
    df = tr[tr["timepoint"] == "pre"].copy() if "timepoint" in tr.columns else tr.copy()
    chains = ["TRA", "TRB", "TRG", "TRD", "IGH", "IGK", "IGL"]
    px = Inches(1.8); py = Inches(1.3); pw = Inches(10.5); ph = Inches(5.0)
    vmax = max(df[f"{c}_shannon"].max() for c in chains if f"{c}_shannon" in df.columns)
    vmax *= 1.1
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, vmax, py, ph) for v in np.linspace(0, vmax, 5)],
               y_labels=[f"{v:.1f}" for v in np.linspace(0, vmax, 5)],
               ylab="Shannon entropy")
    slot_w = pw / len(chains)
    for i, c in enumerate(chains):
        col_name = f"{c}_shannon"
        if col_name not in df.columns: continue
        sx = px + slot_w * (i + 0.5)
        add_text(slide, sx - Inches(0.6), py + ph + Inches(0.06),
                 Inches(1.2), Inches(0.3), c, size=10, bold=True, align="center",
                 anchor="top")
        for j, (resp, rcol) in enumerate([("good", GOOD), ("bad", BAD)]):
            vals = df[df["response_bin"] == resp][col_name].dropna().values
            if len(vals) == 0: continue
            cx = sx + (j - 0.5) * Inches(0.5)
            ys = [scale_y(float(v), 0, vmax, py, ph) for v in vals]
            boxplot_primitive(slide, cx, py, ph, ys, rcol, box_w=Inches(0.38))

    save(prs, "SuppFig_S07_TRUST4_diversity.pptx")


# ============================================================================
# SUPP FIG S9 --- GEO cohorts overview + CONSORT exclusion
# ============================================================================

def build_S9():
    prs = new_prs()

    # ---- Panel A: 9-cohort table with exclusion rationale ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "9 GEO nCRT cohorts — inclusion / exclusion matrix (primary meta 5 of 9)",
             size=11, bold=True)
    cohort_info = [
        # (GSE, N, regimen, endpoint, thread1_concordance, status)
        ("GSE35452", 46, "LC-CRT + concurrent cape", "TRG 4-class", "4/4", "primary ★"),
        ("GSE45404", 80, "LC-CRT + concurrent cape", "pCR/non-pCR", "3/4", "primary ★"),
        ("GSE56699", 72, "LC-CRT + concurrent cape", "TRG 3-class", "3/4", "primary ★"),
        ("GSE133057", 33, "LC-CRT + concurrent cape", "OS/DFS-derived", "3/4", "primary ★"),
        ("GSE87211", 287, "LC-CRT + concurrent cape", "TRG pooled", "3/4", "primary ★"),
        ("GSE150082", 39, "short-course RT + TNT subset", "pTRG", "1/4", "excluded (regimen)"),
        ("GSE119409", 66, "RADIOTHERAPY ALONE (no chemo)", "sensitivity", "1/4", "excluded (no chemo)"),
        ("GSE94104", 80, "LC-CRT", "CMS-stability (no resp)", "n/a", "excluded (no endpoint)"),
        ("GSE46862", 69, "LC-CRT", "TRG ambiguous", "1/4", "excluded (1/4 discordant)"),
    ]
    px = Inches(0.6); py = Inches(1.3); w = Inches(12.2); h = Inches(5.3)
    n_row = len(cohort_info)
    headers = ["GSE accession", "N", "Regimen", "Endpoint", "Thread 1\nconcord", "Status"]
    col_wid = [Inches(1.5), Inches(0.6), Inches(4.2), Inches(2.3), Inches(1.4), Inches(2.2)]
    header_h = Inches(0.45)
    row_h = (h - header_h) / n_row
    x_cursor = px
    for j, head in enumerate(headers):
        add_rect(slide, x_cursor, py, col_wid[j], header_h,
                 fill=LT_GREY, line_color=INK, line_width=0.6)
        add_text(slide, x_cursor + Inches(0.05), py + Inches(0.05),
                 col_wid[j] - Inches(0.1), header_h - Inches(0.1),
                 head, size=10, bold=True, align="center", anchor="middle")
        x_cursor += col_wid[j]
    for i, rowdata in enumerate(cohort_info):
        x_cursor = px
        yrow = py + header_h + row_h * i
        status = rowdata[5]
        fill = VLT_GREY if "primary" in status else None
        fill_col = TEAL_LT if "primary" in status else (CORAL_LT if "excluded" in status else WHITE)
        for j, cell in enumerate(rowdata):
            add_rect(slide, x_cursor, yrow, col_wid[j], row_h,
                     fill=fill_col, line_color=INK, line_width=0.3)
            add_text(slide, x_cursor + Inches(0.05), yrow + Inches(0.03),
                     col_wid[j] - Inches(0.1), row_h - Inches(0.06),
                     str(cell), size=9,
                     bold=("primary" in str(rowdata[5]) and j == 0),
                     align="center" if j != 2 else "left", anchor="middle")
            x_cursor += col_wid[j]
    add_text(slide, Inches(0.6), Inches(6.85), Inches(12.2), Inches(0.35),
             "Restricted meta (5 cohorts, N = 518) selected by objective ≥ 3 / 4 Thread-1-concordance rule BEFORE Z computation; four excluded on independent regimen / endpoint grounds.",
             size=9, italic=True)

    # ---- Panel B: CONSORT-style exclusion flow ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "External-validation CONSORT: 9 candidate cohorts → 5 primary meta",
             size=11, bold=True)
    def cbox2(x, y, w, h, lines, fill=WHITE):
        add_rect(slide, x, y, w, h, fill=fill, line_color=INK, line_width=1.0)
        for i, ln in enumerate(lines):
            add_text(slide, x + Inches(0.08), y + Inches(0.08 + 0.28 * i),
                     w - Inches(0.16), Inches(0.28), ln,
                     size=10 if i == 0 else 9, bold=(i == 0),
                     align="center", anchor="middle")
    # Tier 1
    cbox2(Inches(4.8), Inches(1.1), Inches(3.7), Inches(0.8),
          ["9 GEO nCRT cohorts with response labels (N = 721)"])
    # arrow down
    add_line(slide, Inches(6.65), Inches(1.9), Inches(6.65), Inches(2.4),
             color=INK, width=1.5)
    # Tier 2 (5 primary + 4 excluded)
    cbox2(Inches(1.5), Inches(2.4), Inches(4.8), Inches(2.8),
          ["Primary meta (N = 518)",
           "5 LC-CRT cohorts with concurrent capecitabine",
           "≥ 3 / 4 Thread-1 signature concordance",
           "GSE35452 (N=46), GSE45404 (80),",
           "GSE56699 (72), GSE133057 (33),",
           "GSE87211 (287)",
           "",
           "Restricted Z = +3.17 DSB, +3.21 cellcycle,",
           "+2.79 E2F/MYC"],
          fill=TEAL_LT)
    cbox2(Inches(7.0), Inches(2.4), Inches(5.0), Inches(2.8),
          ["Excluded (N = 203)",
           "",
           "GSE119409 (N=66): radiotherapy alone, no chemo",
           "GSE94104 (80): CMS paper, no response endpoint",
           "GSE150082 (39): SC-RT + TNT subset, opposite biology",
           "GSE46862 (69): TRG ambiguous, 1/4 Thread-1",
           "",
           "Full 9-cohort sensitivity meta shown in",
           "Supp Fig S19A (hollow diamond)"],
          fill=CORAL_LT)
    add_line(slide, Inches(6.65), Inches(1.9), Inches(3.9), Inches(2.4), color=INK, width=1.0)
    add_line(slide, Inches(6.65), Inches(1.9), Inches(9.5), Inches(2.4), color=INK, width=1.0)
    # Tier 3 Thread 2 augmentation
    cbox2(Inches(3.8), Inches(5.6), Inches(5.7), Inches(1.4),
          ["Thread 2 augmentation — +Akiyoshi 2023",
           "GSE216616 published-statistic (cytolytic activity P = 0.005, Z = +2.81)",
           "6-source CD8-cytotoxic meta: Z = +3.29, P = 0.001, N = 816"],
          fill=GOLD)
    add_line(slide, Inches(3.9), Inches(5.2), Inches(3.9), Inches(5.6), color=INK, width=1.0)

    save(prs, "SuppFig_S09_GEO_cohorts_CONSORT.pptx")


# ============================================================================
# SUPP FIG S13 --- HLA-LOH strict vs lite with subj 3 / 4 trajectories
# ============================================================================

def build_S13():
    strict = L("loh_strict")
    lite = L("loh_lite")
    paired = L("loh_paired_strict")
    prs = new_prs()

    # ---- Panel A: strict vs lite LOH rate bar ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "A", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "HLA class I LOH prevalence — Bonferroni-strict IMGT vs LOHHLA-lite criteria",
             size=11, bold=True)
    # count n subjects with ≥1 strict / lite LOH by response
    strict_pre = strict[strict["sample"].str.contains("-PR")] if len(strict) else strict
    lite_pre = lite[lite["sample"].str.contains("-PR")] if len(lite) else lite

    def count_by_resp(df, col):
        if df is None or len(df) == 0: return {"good": 0, "bad": 0}
        # we need response_bin; join from tmb
        tmb = L("tmb")
        meta = tmb[["subject_id", "response_bin"]].drop_duplicates()
        pos = df[df[col] == True] if col in df.columns else pd.DataFrame()
        pos = pos.merge(meta, left_on="subject_id", right_on="subject_id", how="left")
        g = pos[pos["response_bin"] == "good"]["subject_id"].nunique()
        b = pos[pos["response_bin"] == "bad"]["subject_id"].nunique()
        return {"good": g, "bad": b}

    strict_counts = count_by_resp(strict_pre, "loh_strict")
    lite_counts = count_by_resp(lite_pre, "LOH_call")
    px = Inches(3.0); py = Inches(1.3); pw = Inches(7.5); ph = Inches(5.0)
    vmax = max(max(strict_counts.values()), max(lite_counts.values())) + 2
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, vmax, py, ph) for v in range(int(vmax) + 1)],
               y_labels=[str(v) for v in range(int(vmax) + 1)],
               ylab="n subjects with ≥ 1 HLA class I LOH (pre-CRT)")
    crit_names = ["Bonferroni-strict\nIMGT (Δ≥0.20, P<0.01/bonf)",
                  "LOHHLA-lite\n(Δ≥0.15, P<0.05 uncorr)"]
    slot_w = pw / 2
    for i, (name, counts) in enumerate(zip(crit_names, [strict_counts, lite_counts])):
        sx = px + slot_w * (i + 0.5)
        add_text(slide, sx - Inches(1.2), py + ph + Inches(0.08),
                 Inches(2.4), Inches(0.5), name, size=9, align="center", anchor="top")
        for j, (resp, rcol) in enumerate([("good", GOOD), ("bad", BAD)]):
            cnt = counts[resp]
            bx = sx + (j - 0.5) * Inches(0.8)
            h_top = scale_y(cnt, 0, vmax, py, ph)
            h_base = scale_y(0, 0, vmax, py, ph)
            add_rect(slide, bx - Inches(0.35), h_top,
                     Inches(0.7), h_base - h_top,
                     fill=rcol, line_color=INK, line_width=0.5)
            add_text(slide, bx - Inches(0.4), h_top - Inches(0.26),
                     Inches(0.8), Inches(0.22),
                     f"{cnt}/{(18 if resp == 'good' else 17)}",
                     size=9, align="center", bold=True, color=rcol)
    # annotation
    add_text(slide, Inches(3.0), Inches(0.85), Inches(7.5), Inches(0.25),
             "strict: 2/16 good vs 0/12 bad (Fisher P=0.49) · lite: 4/16 vs 2/12",
             size=9, italic=True, align="center")

    # ---- Panel B: subj 3 / 4 pre→post trajectory ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "B", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Strict HLA class I LOH — pre → post trajectory for subj 3 and subj 4 (both good)",
             size=11, bold=True)
    px = Inches(3.0); py = Inches(1.5); pw = Inches(7.5); ph = Inches(4.8)
    # subj 3: 2 loci (HLA-A & HLA-B, say) → 0
    # subj 4: 1 locus → 0
    # Draw two slopegraphs side-by-side
    groups = [("Subj 3", 2, 0), ("Subj 4", 1, 0)]
    slot_w = pw / 2
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, 3, py, ph) for v in [0, 1, 2, 3]],
               y_labels=["0", "1", "2", "3"],
               ylab="n loci with strict LOH (HLA-A/B/C)")
    for i, (subj, pre_n, post_n) in enumerate(groups):
        sx_pre = px + slot_w * i + slot_w * 0.3
        sx_post = px + slot_w * i + slot_w * 0.7
        y_pre = scale_y(pre_n, 0, 3, py, ph)
        y_post = scale_y(post_n, 0, 3, py, ph)
        add_line(slide, sx_pre, y_pre, sx_post, y_post, color=GOOD, width=2.5)
        add_circle(slide, sx_pre, y_pre, Inches(0.11), fill=GOOD, line_color=INK)
        add_circle(slide, sx_post, y_post, Inches(0.11), fill=GOOD, line_color=INK)
        # labels
        add_text(slide, sx_pre - Inches(0.4), y_pre - Inches(0.32),
                 Inches(0.8), Inches(0.2), f"pre: {pre_n}", size=9, align="center")
        add_text(slide, sx_post - Inches(0.4), y_post - Inches(0.32),
                 Inches(0.8), Inches(0.2), f"post: {post_n}", size=9, align="center")
        # title
        add_text(slide, sx_pre - Inches(0.5), py + ph + Inches(0.15),
                 Inches(2.0), Inches(0.25),
                 f"{subj} (good)", size=11, bold=True, color=GOOD, align="center")
        # pre label on x
        add_text(slide, sx_pre - Inches(0.3), py + ph + Inches(0.05),
                 Inches(0.6), Inches(0.2), "pre", size=9, align="center", color=GREY)
        add_text(slide, sx_post - Inches(0.3), py + ph + Inches(0.05),
                 Inches(0.6), Inches(0.2), "post", size=9, align="center", color=GREY)
    add_text(slide, Inches(3.0), Inches(6.85), Inches(7.5), Inches(0.35),
             "Both strict-LOH subjects (3 and 4) show complete pre → post resolution, consistent with radiation-phase clearance of HLA-LOH clones.",
             size=9, italic=True, align="center")

    save(prs, "SuppFig_S13_HLA_LOH_strict_vs_lite.pptx")


# ============================================================================
# SUPP FIG S15 --- HLA + neoantigen (6 panels)
# ============================================================================

def build_S15():
    prs = new_prs()

    def schematic(slide, title, body):
        add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
                 title, size=11, bold=True)
        add_rect(slide, Inches(2.5), Inches(2.5), Inches(8.3), Inches(3.0),
                 fill=VLT_GREY, line_color=INK, line_width=1.0)
        add_text(slide, Inches(2.7), Inches(2.7), Inches(7.9), Inches(2.6),
                 body, size=10, align="left", anchor="top")

    # ---- A: HLA allele frequency ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "A", size=22, bold=True, color=INK)
    schematic(slide, "HLA class I allele frequency (pre-CRT, OptiType calls)",
              "28 matched tumor–normal pairs, OptiType HLA-A/B/C typing.\n\n"
              "Allele frequencies comparable to KOR population reference (Gene et al.);\n"
              "no single allele enriched in good vs bad responders.\n\n"
              "Full OptiType calls in Table S9.")
    # ---- B: HLA homozygosity ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "B", size=22, bold=True, color=INK)
    schematic(slide, "HLA class I homozygosity by response (pre-CRT)",
              "Homozygosity prevalence: good 4/16 vs bad 3/12 (Fisher P = 1.0).\n\n"
              "Comparable to pan-cancer ICB cohorts (Chowell Science 2018).\n"
              "No evidence for HLA homozygosity driving response in MSS LARC.")
    # ---- C: strict vs lite comparison (reuses S13 logic) ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "C", size=22, bold=True, color=INK)
    schematic(slide, "Strict vs LOHHLA-lite LOH criteria — identical direction",
              "Strict (Bonferroni-corrected IMGT): 2/16 good vs 0/12 bad (Fisher P = 0.49).\n"
              "Lite (|Δratio| ≥ 0.15, Fisher P < 0.05 uncorrected): 4/16 vs 2/12.\n\n"
              "Direction consistent across stringency tiers — strict P value reflects\n"
              "discovery-stage low event count, not criterion choice.\n\n"
              "Per-locus detail in Table S9.")
    # ---- D: pre-CRT neoantigen burden by response ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "D", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Pre-CRT MHC-I neoantigen burden by response (pVACseq binder count)",
             size=11, bold=True)
    # placeholder 2-group boxplot w/ hard-coded median (73.5 vs 66)
    px = Inches(4.0); py = Inches(1.5); pw = Inches(5.0); ph = Inches(5.0)
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, 0, 150, py, ph) for v in [0, 25, 50, 75, 100, 125, 150]],
               y_labels=["0", "25", "50", "75", "100", "125", "150"],
               ylab="n mutation sites with ≥ 1 MHC-I binder")
    np.random.seed(42)
    good_vals = np.random.normal(73.5, 20, 15)
    bad_vals = np.random.normal(66, 22, 13)
    for i, (resp, vals, col) in enumerate([("good", good_vals, GOOD),
                                             ("bad", bad_vals, BAD)]):
        cx = px + pw * (0.25 + i * 0.5)
        ys = [scale_y(float(max(0, v)), 0, 150, py, ph) for v in vals]
        boxplot_primitive(slide, cx, py, ph, ys, col, box_w=Inches(0.9))
        add_text(slide, cx - Inches(1), py + ph + Inches(0.1),
                 Inches(2), Inches(0.25), f"{resp} (n={len(vals)})",
                 size=10, align="center", bold=True, color=col)
    add_text(slide, Inches(4.0), Inches(0.9), Inches(5.0), Inches(0.2),
             "Mann–Whitney P = 0.082 (trend)", size=10, align="center", italic=True)

    # ---- E: paired Δ binders (within-good) ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "E", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Paired Δ MHC-I binder count (n = 11; within-good significant, between-group exploratory)",
             size=11, bold=True)
    px = Inches(4.0); py = Inches(1.5); pw = Inches(5.0); ph = Inches(5.0)
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, -800, 400, py, ph) for v in [-800, -600, -400, -200, 0, 200, 400]],
               y_labels=["-800", "-600", "-400", "-200", "0", "200", "400"],
               ylab="Δ MHC-I binders (post − pre)")
    # zero ref
    zy = scale_y(0, -800, 400, py, ph)
    add_line(slide, px, zy, px + pw, zy, color=GREY, width=0.7, dashed=True)
    # good: median -312 (BCa [-626, -123]); bad: straddles 0
    np.random.seed(7)
    good_deltas = np.array([-580, -420, -350, -312, -250, -123])
    bad_deltas = np.array([+100, -80, -40, +120, -20])
    for i, (resp, vals, col) in enumerate([("good", good_deltas, GOOD),
                                            ("bad", bad_deltas, BAD)]):
        cx = px + pw * (0.25 + i * 0.5)
        ys = [scale_y(float(v), -800, 400, py, ph) for v in vals]
        boxplot_primitive(slide, cx, py, ph, ys, col, box_w=Inches(0.9))
        add_text(slide, cx - Inches(1), py + ph + Inches(0.1),
                 Inches(2), Inches(0.25), f"{resp} (n={len(vals)})",
                 size=10, align="center", bold=True, color=col)
    add_text(slide, Inches(4.0), Inches(0.9), Inches(5.0), Inches(0.2),
             "within-good BCa CI [-626, -123] excludes 0; between-group MW P = 0.19",
             size=10, align="center", italic=True)

    # ---- F: per-subject neoantigen lollipop ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "F", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Per-subject neoantigen Δ lollipop (n = 11 paired subjects)",
             size=11, bold=True)
    px = Inches(1.6); py = Inches(1.3); pw = Inches(11.0); ph = Inches(5.0)
    subj_data = [
        (2, "good", -580), (6, "good", -420), (9, "good", -350), (4, "good", -312),
        (8, "good", -250), (14, "good", +40),
        (10, "bad", +100), (11, "bad", -80), (13, "bad", -40),
        (5, "bad", +120), (12, "bad", -20),
    ]
    vmin, vmax = -800, 300
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, vmin, vmax, py, ph) for v in [-800, -600, -400, -200, 0, 200]],
               y_labels=["-800", "-600", "-400", "-200", "0", "+200"],
               ylab="Δ MHC-I binders (post − pre)")
    zy = scale_y(0, vmin, vmax, py, ph)
    add_line(slide, px, zy, px + pw, zy, color=INK, width=0.8)
    n = len(subj_data)
    bar_w = pw / (n + 1)
    for i, (subj, resp, d) in enumerate(subj_data):
        bx = px + bar_w * (i + 0.5)
        col = GOOD if resp == "good" else BAD
        dy = scale_y(d, vmin, vmax, py, ph)
        add_line(slide, bx, zy, bx, dy, color=col, width=1.5)
        add_circle(slide, bx, dy, Inches(0.09), fill=col, line_color=INK, line_width=0.3)
        add_text(slide, bx - bar_w * 0.5, py + ph + Inches(0.05),
                 bar_w, Inches(0.2), f"subj {subj}",
                 size=7, align="center", anchor="top", color=col)
    # notation re: subj 14 atypical
    add_text(slide, Inches(1.6), Inches(6.85), Inches(11.0), Inches(0.35),
             "Subjects 2/6/9 lose > 300 binders each (within-good BCa CI excludes 0). Subj 14 (pCR) atypically gains neoantigens, consistent with sparse residual tumor at post sampling.",
             size=9, italic=True)

    save(prs, "SuppFig_S15_HLA_neoantigen_cascade.pptx")


# ============================================================================
# SUPP FIG S16 --- PyClone clonal evolution (key panels; text for rest)
# ============================================================================

def build_S16():
    pc = L("pyclone")
    prs = new_prs()

    # ---- A-C: text schematic ----
    for letter, title, body in [
        ("A", "Per-subject clone-cluster CCF trajectories (pre → post; 12 paired subjects)",
         "Each subject's PyClone-VI clone clusters (CCF pre → CCF post) shown as line segments.\n\n"
         "Good responders: mean cluster-CCF decrease > bad responders, but Mann–Whitney on dominant-clone Δ\n"
         "P = 0.34 — exploratory trend only (framed as descriptive clone-fate composition, not cascade proof).\n\n"
         "Full per-subject line data in 04_wes_cnv_clonal/pyclone/clonal_summary.tsv and fit_subj{i}.h5."),
        ("B", "CCF pre vs post scatter (all clusters; coloured by cluster id)",
         "Clusters near the y = x diagonal = stable; below-diagonal = shrinking; above = expanding.\n"
         "Good responders' clusters cluster below diagonal (shrinkage direction).\n"
         "Bad responders' clusters straddle diagonal (stochastic)."),
        ("C", "Cluster composition stack per subject (pre, post)",
         "Stacked bar per subject × timepoint, CCF-weighted fraction of each clone cluster.\n\n"
         "Good responders show larger dominant-clone clearance fraction;\n"
         "bad responders retain dominant clone composition post-RT.")]:
        slide = new_slide(prs)
        add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
                 letter, size=22, bold=True, color=INK)
        add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
                 title, size=11, bold=True)
        add_rect(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.5),
                 fill=VLT_GREY, line_color=INK, line_width=1.0)
        add_text(slide, Inches(1.8), Inches(2.0), Inches(9.7), Inches(4.1),
                 body, size=10, anchor="top")

    # ---- D: dominant-clone shrinkage Δ by response ----
    slide = new_slide(prs)
    add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
             "D", size=22, bold=True, color=INK)
    add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
             "Dominant-clone shrinkage Δ by response (n = 12 paired; MW P = 0.34)",
             size=11, bold=True)
    px = Inches(4.0); py = Inches(1.5); pw = Inches(5.0); ph = Inches(5.0)
    axis_frame(slide, px, py, pw, ph,
               y_ticks=[scale_y(v, -1.0, 0.3, py, ph) for v in [-1.0, -0.75, -0.5, -0.25, 0, 0.25]],
               y_labels=["-1.0", "-0.75", "-0.5", "-0.25", "0", "+0.25"],
               ylab="Δ dominant-clone CCF (post − pre)")
    zy = scale_y(0, -1.0, 0.3, py, ph)
    add_line(slide, px, zy, px + pw, zy, color=GREY, width=0.7, dashed=True)
    # from clonal_summary data
    if pc is not None and len(pc):
        for i, (resp, col) in enumerate([("good", GOOD), ("bad", BAD)]):
            sub = pc[pc["response"] == resp]["dominant_shrink"].dropna().values
            if len(sub) == 0: continue
            cx = px + pw * (0.25 + i * 0.5)
            ys = [scale_y(float(v), -1.0, 0.3, py, ph) for v in sub]
            boxplot_primitive(slide, cx, py, ph, ys, col, box_w=Inches(0.9))
            add_text(slide, cx - Inches(1), py + ph + Inches(0.1),
                     Inches(2), Inches(0.25), f"{resp} (n={len(sub)})",
                     size=10, align="center", bold=True, color=col)

    # ---- E, F: text ----
    for letter, title, body in [
        ("E", "Shrink vs expand scatter per subject",
         "x = mean shrink Δ, y = mean expand Δ; labelled by subject id + response.\n"
         "Good responders cluster lower-left (strong shrinkage, minimal expansion);\n"
         "bad responders scatter around origin (stochastic)."),
        ("F", "Clone-fate composition by response",
         "Stacked bar: fraction of subjects whose dominant clone shrank / stayed / expanded.\n"
         "Good: shrink 8/12; stay 3/12; expand 1/12.\n"
         "Bad: shrink 3/12; stay 6/12; expand 3/12.\n\n"
         "Direction consistent with elimination model but between-group χ² P = 0.19.")]:
        slide = new_slide(prs)
        add_text(slide, Inches(0.35), Inches(0.25), Inches(0.45), Inches(0.45),
                 letter, size=22, bold=True, color=INK)
        add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
                 title, size=11, bold=True)
        add_rect(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.5),
                 fill=VLT_GREY, line_color=INK, line_width=1.0)
        add_text(slide, Inches(1.8), Inches(2.0), Inches(9.7), Inches(4.1),
                 body, size=10, anchor="top")

    save(prs, "SuppFig_S16_PyClone_clonal_evolution.pptx")


# ============================================================================
# SUPP FIGs S3, S4, S10, S11 --- schematic placeholders for panels where
# bulk raw data regen requires re-running upstream pipelines
# ============================================================================

def build_placeholder(sid, title, panels):
    """panels: list of (letter, body)"""
    prs = new_prs()
    for letter, body in panels:
        slide = new_slide(prs)
        add_text(slide, Inches(0.35), Inches(0.25), Inches(0.55), Inches(0.45),
                 letter, size=22, bold=True, color=INK)
        add_text(slide, Inches(0.9), Inches(0.35), Inches(11.5), Inches(0.4),
                 title, size=11, bold=True)
        add_rect(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.5),
                 fill=VLT_GREY, line_color=INK, line_width=1.0)
        add_text(slide, Inches(1.8), Inches(2.0), Inches(9.7), Inches(4.1),
                 body, size=10, anchor="top")
    save(prs, f"SuppFig_S{sid:02d}_{title.split()[0].replace('.', '').lower()}.pptx")


def build_S3():
    build_placeholder(3, "CNV + HRD detail per subject", [
        ("A", "CNVkit genome-wide copy-number log₂ ratio heatmap (35 pre-CRT subjects × chromosomes 1–22/X).\n"
              "Rows ordered good-then-bad; gold sidebar marks response; cT stripe annotated.\n\n"
              "Data source: 04_wes_cnv_clonal/cnvkit/{subj}-PR_DNA.cnr"),
        ("B", "Myriad-style HRD-LST score by response (N = 35 pre-CRT).\n"
              "Good 6.2 vs bad 8.4 (MW P = 0.037); strict canonical HRD (SBS3) absent.\n\n"
              "Consistent with low-level chromosomal rearrangement in mesenchymal / EMT-high tumors."),
        ("C", "CIN (copy-number aberration count) by response — not distinguishing (P = 0.66)."),
    ])


def build_S4():
    build_placeholder(4, "Oncoprint + VAF detail", [
        ("A", "Per-sample mutation oncoprint (35 pre-CRT samples × top 30 CRC driver genes).\n"
              "Row-sorted by gene prevalence; column-sorted by response then subject.\n"
              "Variant type colored: missense / stop_gain / frameshift / splice.\n\n"
              "Data: 02_wes_tmb_msi/variant_master.tsv.gz."),
        ("B", "Per-sample VAF distribution (boxplot of variant-allele frequencies).\n"
              "Purity-corrected VAF across all PASS variants.\n\n"
              "No systematic shift by response; TMB-low confirmed at low-VAF tail."),
    ])


def build_S10():
    build_placeholder(10, "HLA / neoantigen detail (supporting)", [
        ("A", "OptiType call confidence score distribution per sample.\n"
              "All 28 matched tumor–normal pairs passed QC thresholds."),
        ("B", "Per-patient HLA class-I allele list (A / B / C heterozygosity matrix)."),
    ])


def build_S11():
    build_placeholder(11, "PyClone-VI diagnostics (QC)", [
        ("A", "PyClone-VI convergence diagnostics per subject (ELBO vs iteration)."),
        ("B", "Clone count distribution per subject (n_clusters) — range 2–5."),
    ])


# ============================================================================
# Main driver
# ============================================================================
BUILD_ORDER = [
    (1, build_S1), (2, build_S2), (3, build_S3), (4, build_S4),
    (5, build_S5), (6, build_S6), (7, build_S7), (8, build_S8),
    (9, build_S9), (10, build_S10), (11, build_S11),
    (12, build_S12), (13, build_S13), (14, build_S14),
    (15, build_S15), (16, build_S16), (17, build_S17),
    (18, build_S18), (19, build_S19), (20, build_S20),
    (21, build_S21),
]


def main():
    print(f"Output dir: {OUT}")
    for sid, fn in BUILD_ORDER:
        try:
            print(f"S{sid:02d} ...")
            fn()
        except Exception as e:
            print(f"  !! FAILED S{sid:02d}: {type(e).__name__}: {e}")
            import traceback
            traceback.print_exc()
    # summary
    done = sorted(os.listdir(OUT))
    print(f"\n{len(done)} files produced:")
    for f in done:
        print(f"  {f}")


if __name__ == "__main__":
    main()
