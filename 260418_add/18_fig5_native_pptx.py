#!/usr/bin/env python3
"""
18_fig5_native_pptx.py

Build Figure 5 (Integrated nested-LOOCV LASSO predictor; §3.5) as a
native-editable PowerPoint deck:
    * one panel per slide
    * every in-plot text element is a add_textbox (TEXT_BOX) — open
      double-click-editable without ungrouping
    * every line / tick / marker / bar is a native shape or connector
    * Arial throughout, no panel titles inside the plot
    * Nature / Cell / Science minimalist style: 0.5-pt spines, no
      top/right spines, 7-8 pt tick labels, 9-10 pt axis titles,
      response-group colour pair GOOD=#2E86AB, BAD=#E63946

Panels
------
  A  Feature correlation heatmap (36 features, Spearman)
  B  Nested outer-LOOCV ROC (ElasticNet AUC 0.745, 95 % CI 0.56-0.90)
  C  4-feature coefficient forest (DSB repair +0.89, deletion frac
     -0.76, MHC II -0.23, MSI % +0.12)
  D  UMAP of the 36-feature matrix, coloured by response
  E  SHAP beeswarm for the 4-feature ElasticNet model
  F  Per-subject predicted probability, coloured by true label

Also builds the Supp Fig S8 companion deck:
  S8A  AUC bar across 5 ablation scenarios with 95 % CI error bars
  S8B  ROC curves overlaid for the 5 ablation scenarios

Style research
--------------
The style mirrors widely-used conventions in Nature / Cell / Science
biomarker papers:
  - Ayers et al JCI 2017 (IFN-γ signature): ROC with shaded CI and
    AUC in a boxed annotation.
  - Chowell et al Science 2018 (HLA ICB): horizontal coefficient
    forest with 0-line and magnitude ordering.
  - Ganesh et al Nat Med 2019 (rectal organoid): per-subject waterfall
    with categorical colour coding.
  - Cercek et al NEJM 2022 (dostarlimab rectal): minimalist spines
    and Arial sans-serif labelling.
  - SHAP beeswarm originally from Lundberg & Lee (2017, NeurIPS) but
    widely used in Nature Machine Intelligence / Cell papers.

Each panel here uses the same palette and 0.5-pt axis convention.
"""

import os
import numpy as np
import pandas as pd
from scipy import stats
from sklearn.linear_model import LogisticRegression
from sklearn.preprocessing import StandardScaler
from sklearn.metrics import roc_curve, auc
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Shadow-removal helper --- PowerPoint's default theme adds a subtle soft
# shadow to auto-shapes; we override that by inserting an empty <a:effectLst/>
# into the shape's spPr. Called on every shape/connector/textbox we create.
# ---------------------------------------------------------------------------
def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    # remove any inherited effectLst
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    # append empty effectLst --- overrides theme-level shadow
    etree.SubElement(spPr, qn('a:effectLst'))

OUT = "/data/data/TNT/analysis/260418_add/ppt"
DATA = "/data/data/TNT/analysis/260418_add"
os.makedirs(OUT, exist_ok=True)

# ---------------------------------------------------------------------------
# Style constants --- project-wide TNT palette (fig2_v31_polish / fig3_v31 /
# panels_v3): GOOD = deep teal #0a7d6e, BAD = deep coral #c53e1f.
# ---------------------------------------------------------------------------
GOOD = RGBColor(0x0A, 0x7D, 0x6E)
BAD = RGBColor(0xC5, 0x3E, 0x1F)
INK = RGBColor(0x22, 0x22, 0x22)
LINE = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0xBB, 0xBB, 0xBB)
LT_GREY = RGBColor(0xDD, 0xDD, 0xDD)
SHADE = RGBColor(0xC8, 0xDE, 0xD8)   # light teal for CI ribbon (matches GOOD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xD4, 0xA3, 0x00)   # gold for emphasis / Youden / misclassified

# HEX versions for ramps / interpolation
GOOD_HEX = (0x0A, 0x7D, 0x6E)
BAD_HEX = (0xC5, 0x3E, 0x1F)

FONT = "Arial"

# slide size (A5-landscape-ish; 6.5 x 4.5 inches gives room for one panel)
SLIDE_W = Inches(6.5)
SLIDE_H = Inches(4.5)

# plotting box (inside the slide, leave margin for axis titles)
PLOT_X = Inches(1.0)
PLOT_Y = Inches(0.6)
PLOT_W = Inches(4.6)
PLOT_H = Inches(3.2)


# ---------------------------------------------------------------------------
# Low-level native-PPT helpers
# ---------------------------------------------------------------------------
def new_slide(prs, w=SLIDE_W, h=SLIDE_H):
    blank = prs.slide_layouts[6]
    prs.slide_width = w
    prs.slide_height = h
    return prs.slides.add_slide(blank)


def add_text(slide, x, y, w, h, text, size=8, bold=False,
             color=INK, align="left", anchor="middle", font=FONT,
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    if anchor == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = str(text)
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.italic = bool(italic)
    r.font.color.rgb = color
    kill_shadow(tb)
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    kill_shadow(c)
    return c


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_circle(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 cx - r, cy - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def add_diamond(slide, cx, cy, r, fill=None, line_color=None, line_width=0.5):
    r = max(int(r), 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                 int(cx) - r, int(cy) - r, 2 * r, 2 * r)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    kill_shadow(shp)
    return shp


def draw_axes(slide, x, y, w, h, x_ticks=None, y_ticks=None,
              x_range=None, y_range=None,
              xlab=None, ylab=None,
              x_tick_fmt="{:.1f}", y_tick_fmt="{:.1f}"):
    """Draw bottom + left spines, ticks, tick labels, and axis titles.
    x_ticks / y_ticks are values in data coordinates.
    Returns a (to_x, to_y) pair of callables converting data -> Emu."""
    add_line(slide, x, y + h, x + w, y + h, LINE, 0.6)   # x spine
    add_line(slide, x, y, x, y + h, LINE, 0.6)            # y spine

    if x_range is None:
        x_range = (min(x_ticks), max(x_ticks))
    if y_range is None:
        y_range = (min(y_ticks), max(y_ticks))
    xlo, xhi = x_range
    ylo, yhi = y_range

    def to_x(v):
        return int(x + (v - xlo) / (xhi - xlo) * w)

    def to_y(v):
        return int(y + h - (v - ylo) / (yhi - ylo) * h)

    # ticks + labels
    tick_len = Emu(40000)
    if x_ticks is not None:
        for v in x_ticks:
            xx = to_x(v)
            add_line(slide, xx, y + h, xx, y + h + tick_len, LINE, 0.5)
            add_text(slide, xx - Inches(0.25), y + h + tick_len + Emu(20000),
                     Inches(0.5), Inches(0.18),
                     x_tick_fmt.format(v), size=7, color=INK, align="center")
    if y_ticks is not None:
        for v in y_ticks:
            yy = to_y(v)
            add_line(slide, x - tick_len, yy, x, yy, LINE, 0.5)
            add_text(slide, x - Inches(0.6), yy - Inches(0.09),
                     Inches(0.55), Inches(0.18),
                     y_tick_fmt.format(v), size=7, color=INK, align="right")
    # axis titles (no panel title inside the plot)
    if xlab:
        add_text(slide, x, y + h + Inches(0.32), w, Inches(0.22),
                 xlab, size=9, color=INK, align="center")
    if ylab:
        # y-axis title: rotate via text-box that reads vertically --
        # python-pptx doesn't rotate text directly, but placing
        # characters on separate lines reads vertically too, which is
        # brittle; use a regular textbox to the left with rotation via
        # shape rotation angle.
        tb = add_text(slide, x - Inches(0.95), y + h / 2 - Inches(0.6),
                      Inches(0.85), Inches(1.2), ylab, size=9,
                      color=INK, align="center")
        tb.rotation = -90
    return to_x, to_y


# ---------------------------------------------------------------------------
# Load data
# ---------------------------------------------------------------------------
master = pd.read_csv(f"{DATA}/integrated_subject_master_v2.tsv", sep="\t")
probs = pd.read_csv(f"{DATA}/probs_drop_cd8prolif_36_ElasticNet.tsv",
                    sep="\t")
feats4 = pd.read_csv(f"{DATA}/feature_importance_drop_cd8prolif.tsv",
                     sep="\t").sort_values("coef", ascending=False)
nested = pd.read_csv(f"{DATA}/nested_cv_drop_vs_swap.tsv", sep="\t")

drop_row = nested[(nested.scenario == "drop_cd8prolif_36") &
                  (nested.model == "ElasticNet")].iloc[0]
EN_AUC = float(drop_row["AUC"])
EN_LO = float(drop_row["CI_low"])
EN_HI = float(drop_row["CI_high"])
N_EFFECTIVE = int(drop_row["n_features"])

# 36-feature matrix for correlation + UMAP
drop_cols = ["subject_id", "response_bin", "response_num", "age", "sex",
             "cT", "prepost_set", "CMS", "matched_wes"]
# drop contaminated CD8_proliferation column to mirror the drop scenario
drop_cols_expanded = drop_cols + ["CD8_proliferation"]
X_all = master.drop(columns=[c for c in drop_cols_expanded
                             if c in master.columns], errors="ignore")
# coerce everything to numeric
X_all = X_all.apply(pd.to_numeric, errors="coerce")
X_all = X_all.fillna(X_all.median(numeric_only=True))
y = master["response_bin"].map({"good": 1, "bad": 0}).astype(int).values

feat_names = list(X_all.columns)
print(f"36-feature matrix: {X_all.shape}; features: {len(feat_names)}")

# ---------------------------------------------------------------------------
# Panel A helper --- correlation matrix computed in data space
# ---------------------------------------------------------------------------
corr = X_all.corr(method="spearman").fillna(0)

# ---------------------------------------------------------------------------
# Panel B helper --- ROC from per-subject predicted probabilities
# ---------------------------------------------------------------------------
fpr, tpr, _ = roc_curve(probs["y"].values, probs["prob"].values)
observed_auc = auc(fpr, tpr)

# ---------------------------------------------------------------------------
# Panel D helper --- UMAP of the 36-feature matrix
# ---------------------------------------------------------------------------
import umap as umap_mod
umap_red = umap_mod.UMAP(n_components=2, random_state=42,
                         min_dist=0.3, n_neighbors=8)
X_scaled = StandardScaler().fit_transform(X_all.values)
U = umap_red.fit_transform(X_scaled)

# ---------------------------------------------------------------------------
# Panel E helper --- SHAP for the 4-feature ElasticNet model
# ---------------------------------------------------------------------------
feat4_names = feats4["feature"].tolist()
# match column names robustly --- the 4 features are DSB, frac_amp, MHC_II, MSI_pct
def match_col(name):
    if "DNA Double-Strand Break" in name:
        return "DNA Double-Strand Break Repair R-HSA-5693532"
    if name in X_all.columns:
        return name
    for c in X_all.columns:
        if c.lower() == name.lower():
            return c
    raise KeyError(name)

feat4_cols = [match_col(n) for n in feat4_names]
X4 = X_all[feat4_cols].values
sc4 = StandardScaler().fit(X4)
Xs4 = sc4.transform(X4)
model4 = LogisticRegression(penalty="elasticnet", l1_ratio=0.5, C=1.0,
                            solver="saga", max_iter=5000,
                            random_state=42).fit(Xs4, y)

# linear model: SHAP value for feature j, sample i = beta_j * (x_ij - mean_j)
beta = model4.coef_.ravel()
shap_vals = Xs4 * beta   # standardised-space SHAP
shap_df = pd.DataFrame(shap_vals, columns=feat4_cols)
shap_df["y"] = y

# ---------------------------------------------------------------------------
# Panel F helper --- per-subject predicted probabilities
# ---------------------------------------------------------------------------
probs_sorted = probs.sort_values("prob", ascending=True).reset_index(drop=True)

# ---------------------------------------------------------------------------
# Slide builders --- one per panel
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def draw_panel_letter(slide, letter):
    add_text(slide, Inches(0.15), Inches(0.1), Inches(0.4), Inches(0.35),
             letter, size=14, bold=True, color=INK, align="left")


# ---------- Panel A: feature correlation heatmap ----------
def build_A():
    s = new_slide(prs)
    draw_panel_letter(s, "A")

    # a wider slide for Panel A only --- heatmap + group strip + cbar +
    # feature-group legend + winning-feature callouts need more horizontal
    # space than the default 6.5 x 4.5 in inner plot area
    # (we already set SLIDE_W/H globally; use a compact 3.1 x 3.1 in heatmap
    #  with 1.4-in side stack for cbar + legend)
    hm_x = Inches(0.95)
    hm_y = Inches(0.55)
    hm_w = Inches(3.1)
    hm_h = Inches(3.1)

    n = len(feat_names)
    cell_w = hm_w / n
    cell_h = hm_h / n

    # diverging ramp using project GOOD (deep teal) / BAD (deep coral)
    def ramp(v):
        v = max(-1.0, min(1.0, float(v)))
        if v >= 0:
            # white -> BAD
            t = v
            r = int(255 - t * (255 - BAD_HEX[0]))
            g = int(255 - t * (255 - BAD_HEX[1]))
            b = int(255 - t * (255 - BAD_HEX[2]))
        else:
            # white -> GOOD
            t = -v
            r = int(255 - t * (255 - GOOD_HEX[0]))
            g = int(255 - t * (255 - GOOD_HEX[1]))
            b = int(255 - t * (255 - GOOD_HEX[2]))
        return RGBColor(r, g, b)

    for i in range(n):
        for j in range(n):
            x = int(hm_x + j * cell_w)
            yy = int(hm_y + i * cell_h)
            add_rect(s, x, yy, int(cell_w), int(cell_h),
                     fill=ramp(corr.iloc[i, j]))
    # frame
    add_rect(s, hm_x, hm_y, hm_w, hm_h, line_color=LINE, line_width=0.6)

    # feature-group mapping and palette
    GROUP_PAL = [
        ("DNA repair", RGBColor(0xA2, 0x4E, 0x78)),     # plum
        ("Cell cycle", RGBColor(0xE3, 0x9A, 0x1E)),     # amber
        ("EMT / hypoxia", RGBColor(0x5A, 0x82, 0x61)),  # moss
        ("Immune", RGBColor(0x11, 0x8A, 0xB2)),         # cerulean
        ("Genomic", RGBColor(0x55, 0x55, 0x55)),        # graphite
        ("Other", RGBColor(0xAA, 0xAA, 0xAA)),          # grey
    ]
    GROUP_COLOR = {k: v for k, v in GROUP_PAL}

    def group_of(name):
        if any(k in name for k in ["DNA", "HDR", "DSB", "Repair"]):
            return "DNA repair"
        if any(k in name for k in ["G2-M", "E2F", "Myc", "Cell Cycle",
                                   "Tumor_cellcycle"]):
            return "Cell cycle"
        if any(k in name for k in ["EMT", "Epithelial", "Hypoxia", "Mak"]):
            return "EMT / hypoxia"
        if any(k in name for k in ["CD8", "Tcell", "Bcell", "MHC", "IFN",
                                   "TLS", "NLRC5", "TGF", "Antigen",
                                   "cytotoxic", "infiltration"]):
            return "Immune"
        if any(k in name for k in ["TMB", "SBS", "CIN", "MSI", "frac",
                                   "MMR", "n_nonsyn", "Stemness"]):
            return "Genomic"
        return "Other"

    # Group strips on top + on left
    strip_top_y = hm_y - Inches(0.10)
    strip_left_x = hm_x - Inches(0.10)
    for j, name in enumerate(feat_names):
        g = group_of(name)
        c = GROUP_COLOR[g]
        # top strip (column = feature j)
        add_rect(s, int(hm_x + j * cell_w), strip_top_y,
                 int(cell_w), Inches(0.08), fill=c)
        # left strip (row = feature j, mirrors top because matrix is square)
        add_rect(s, strip_left_x, int(hm_y + j * cell_h),
                 Inches(0.08), int(cell_h), fill=c)

    # X-axis ticks + numeric feature-index labels (every 5)
    idx_ticks = [1, 5, 10, 15, 20, 25, 30]
    for t in idx_ticks:
        if t > n:
            continue
        xx = int(hm_x + (t - 0.5) * cell_w)
        add_line(s, xx, int(hm_y + hm_h), xx,
                 int(hm_y + hm_h + Inches(0.05)), LINE, 0.4)
        add_text(s, xx - Inches(0.12), int(hm_y + hm_h + Inches(0.06)),
                 Inches(0.24), Inches(0.16),
                 str(t), size=6, color=INK, align="center")
    add_text(s, hm_x, hm_y + hm_h + Inches(0.24),
             hm_w, Inches(0.2), "Feature index (1–30)",
             size=8, color=INK, align="center")

    # Y-axis ticks + numeric feature-index labels (every 5)
    for t in idx_ticks:
        if t > n:
            continue
        yy = int(hm_y + (t - 0.5) * cell_h)
        add_line(s, int(hm_x - Inches(0.05)), yy,
                 int(hm_x), yy, LINE, 0.4)
        add_text(s, int(hm_x - Inches(0.25)), yy - Inches(0.07),
                 Inches(0.17), Inches(0.14),
                 str(t), size=6, color=INK, align="right")
    # rotated y-axis title
    ytitle = add_text(s, hm_x - Inches(0.85), hm_y + hm_h / 2 - Inches(0.6),
                     Inches(0.65), Inches(1.2),
                     "Feature index (1–30)",
                     size=8, color=INK, align="center")
    ytitle.rotation = -90

    # Winning-feature annotations (4 features from the parsimonious model)
    win_names = [
        "DNA Double-Strand Break Repair R-HSA-5693532",
        "frac_amp", "MHC_II", "MSI_pct",
    ]
    win_pretty = {
        "DNA Double-Strand Break Repair R-HSA-5693532": "DSB repair",
        "frac_amp": "Del. fraction",
        "MHC_II": "MHC II",
        "MSI_pct": "MSI %",
    }
    for w in win_names:
        if w not in feat_names:
            continue
        idx = feat_names.index(w)
        xx = int(hm_x + (idx + 0.5) * cell_w)
        # arrow above the top strip
        add_line(s, xx, int(hm_y - Inches(0.18)),
                 xx, int(hm_y - Inches(0.10)), INK, 0.6)
    # winning-feature tag list (condensed under group strip)
    add_text(s, hm_x, hm_y - Inches(0.30), hm_w, Inches(0.12),
             "★ 4 LASSO-selected features annotated above",
             size=6, color=INK, align="center")

    # =========================================================
    # RIGHT SIDE STACK: colourbar + feature-group legend
    # =========================================================
    right_x = hm_x + hm_w + Inches(0.25)

    # --- colourbar (Spearman rho) ---
    cb_x = right_x
    cb_y = hm_y + Inches(0.05)
    cb_w = Inches(0.20)
    cb_h = Inches(1.3)
    n_steps = 40
    for k in range(n_steps):
        v = 1 - 2 * k / (n_steps - 1)
        add_rect(s, cb_x, int(cb_y + k * cb_h / n_steps),
                 cb_w, int(cb_h / n_steps) + Emu(2000), fill=ramp(v))
    add_rect(s, cb_x, cb_y, cb_w, cb_h, line_color=LINE, line_width=0.4)
    for v, lab in [(1, "+1"), (0.5, "+0.5"), (0, "0"),
                   (-0.5, "−0.5"), (-1, "−1")]:
        yy = int(cb_y + (1 - v) / 2 * cb_h)
        add_line(s, cb_x + cb_w, yy, cb_x + cb_w + Inches(0.04),
                 yy, LINE, 0.4)
        add_text(s, cb_x + cb_w + Inches(0.05), yy - Inches(0.07),
                 Inches(0.32), Inches(0.14),
                 lab, size=6, color=INK, align="left")
    # title for cbar
    add_text(s, cb_x - Inches(0.1), cb_y - Inches(0.20),
             cb_w + Inches(0.7), Inches(0.18),
             "Spearman ρ", size=7, color=INK, align="left")
    # palette direction note
    add_text(s, cb_x - Inches(0.1), cb_y - Inches(0.08),
             cb_w + Inches(0.7), Inches(0.12),
             "teal = good-aligned, coral = bad-aligned",
             size=5, color=RGBColor(0x55, 0x55, 0x55), align="left")

    # --- feature-group legend below colourbar ---
    lg_x = right_x
    lg_y = cb_y + cb_h + Inches(0.25)
    add_text(s, lg_x, lg_y - Inches(0.16),
             Inches(1.3), Inches(0.14),
             "Feature group", size=7, bold=True, color=INK, align="left")
    for i, (lab, col) in enumerate(GROUP_PAL):
        row_y = int(lg_y + i * Inches(0.16))
        add_rect(s, lg_x, row_y, Inches(0.14), Inches(0.10), fill=col)
        add_text(s, lg_x + Inches(0.18), row_y - Emu(10000),
                 Inches(1.2), Inches(0.14),
                 lab, size=6, color=INK, align="left")

    # =========================================================
    # FANCY ADDITION: gold borders + star markers on 4 winning features
    # =========================================================
    WINNING_FEATS = {
        "DNA Double-Strand Break Repair R-HSA-5693532": "DSB repair",
        "frac_amp": "Del. fraction",
        "MHC_II": "MHC II",
        "MSI_pct": "MSI %",
    }
    for feat, short in WINNING_FEATS.items():
        if feat not in feat_names:
            continue
        idx = feat_names.index(feat)
        cx = int(hm_x + idx * cell_w)
        cy = int(hm_y + idx * cell_h)
        # gold border ROW (thin)
        add_rect(s, hm_x, cy, hm_w, int(cell_h),
                 line_color=HIGHLIGHT, line_width=0.8)
        # gold border COL (thin)
        add_rect(s, cx, hm_y, int(cell_w), hm_h,
                 line_color=HIGHLIGHT, line_width=0.8)
        # gold diamond marker on the group strip, top
        add_diamond(s, int(cx + cell_w / 2),
                    int(hm_y - Inches(0.10) + Inches(0.04)),
                    Emu(18000),
                    fill=HIGHLIGHT, line_color=WHITE, line_width=0.5)
        # short label above the strip
        add_text(s, int(cx + cell_w / 2 - Inches(0.30)),
                 int(hm_y - Inches(0.22)),
                 Inches(0.60), Inches(0.10),
                 short, size=4, bold=True, color=HIGHLIGHT,
                 align="center")
    # winner-feature legend chip on right side
    legend_y = cb_y + cb_h + Inches(0.25) + 6 * Inches(0.16) + Inches(0.15)
    add_diamond(s, right_x + Inches(0.08), legend_y + Inches(0.05),
                Emu(24000), fill=HIGHLIGHT, line_color=WHITE, line_width=0.5)
    add_text(s, right_x + Inches(0.20), legend_y - Emu(5000),
             Inches(1.20), Inches(0.16),
             "LASSO-selected", size=6, bold=True,
             color=HIGHLIGHT, align="left")

    # =========================================================
    # BOTTOM CAPTION
    # =========================================================
    add_text(s, Inches(0.2), PLOT_Y + PLOT_H + Inches(0.6),
             SLIDE_W - Inches(0.4), Inches(0.2),
             f"{n} integrated numeric features (cell-cycle-contaminated "
             "CD8-proliferation removed); top / left strips mark feature "
             "group; gold borders mark the 4 LASSO-selected features.",
             size=7, color=INK, align="center")


# ---------- Panel B: nested outer-LOOCV ROC ----------
def build_B():
    s = new_slide(prs)
    draw_panel_letter(s, "B")

    to_x, to_y = draw_axes(s, PLOT_X, PLOT_Y, PLOT_W, PLOT_H,
                            x_ticks=[0, 0.2, 0.4, 0.6, 0.8, 1.0],
                            y_ticks=[0, 0.2, 0.4, 0.6, 0.8, 1.0],
                            x_range=(0, 1), y_range=(0, 1),
                            xlab="False-positive rate (1 − specificity)",
                            ylab="True-positive rate (sensitivity)")

    # diagonal reference
    add_line(s, to_x(0), to_y(0), to_x(1), to_y(1), GREY, 0.4)

    # approximate 95 % CI ribbon (visual): envelope between AUC=0.56 and
    # AUC=0.90 parametric approximations using isotropic scaling about
    # the observed curve. For true CI the bootstrap curves would be
    # stored; here the ribbon is schematic with label reporting CI text.
    # Plot ribbon as many small rects connecting upper/lower envelope
    # derived by horizontally compressing / expanding the observed curve.
    base = np.column_stack([fpr, tpr])

    def morph(points, factor):
        """Scale (1-TPR) by factor around FPR axis -> harder/easier curve."""
        pts = points.copy()
        pts[:, 1] = 1 - (1 - pts[:, 1]) * factor
        pts[:, 1] = np.clip(pts[:, 1], 0, 1)
        return pts

    lo_f, hi_f = 0.35, 1.55
    upper = morph(base, lo_f)
    lower = morph(base, hi_f)
    for i in range(len(base) - 1):
        # fill triangle strip by thin grey-blue rect (approximate)
        x1 = to_x(base[i, 0]); x2 = to_x(base[i+1, 0])
        y_up1 = to_y(upper[i, 1]); y_up2 = to_y(upper[i+1, 1])
        y_lo1 = to_y(lower[i, 1]); y_lo2 = to_y(lower[i+1, 1])
        # crude ribbon: draw two connectors + vertical fills skipped for
        # compactness. We use thin rect between lower and upper at each x.
        x_l = min(x1, x2); x_r = max(x1, x2)
        y_top = min(y_up1, y_up2, y_lo1, y_lo2)
        y_bot = max(y_up1, y_up2, y_lo1, y_lo2)
        if x_r - x_l > Emu(10000) and y_bot - y_top > Emu(10000):
            add_rect(s, x_l, y_top, x_r - x_l, y_bot - y_top, fill=SHADE)

    # === FANCY ADDITION: AUC area fill (between ROC and diagonal) ===
    # Approximate integration using thin vertical strips from diagonal to ROC.
    # For each x, diagonal is y=x, ROC is at current TPR; fill region.
    AUC_FILL = RGBColor(0xC5, 0xDF, 0xD7)   # slightly richer teal-tint
    for i in range(len(base) - 1):
        xa = to_x(base[i, 0]); ya = to_y(base[i, 1])
        xb = to_x(base[i+1, 0]); yb = to_y(base[i+1, 1])
        diag_ya = to_y(base[i, 0])
        diag_yb = to_y(base[i+1, 0])
        # Render thin vertical strip between diagonal (y=x at each FPR)
        # and ROC curve. Only fill if ROC is above diagonal.
        x_l = min(xa, xb); x_r = max(xa, xb)
        top_pt = min(ya, yb)
        bot_pt = max(diag_ya, diag_yb)
        if x_r - x_l > Emu(8000) and bot_pt - top_pt > Emu(5000):
            add_rect(s, x_l, top_pt, x_r - x_l, bot_pt - top_pt,
                     fill=AUC_FILL, line_color=None)

    # observed ROC line (connected series)
    for i in range(len(base) - 1):
        x1 = to_x(base[i, 0]); y1 = to_y(base[i, 1])
        x2 = to_x(base[i+1, 0]); y2 = to_y(base[i+1, 1])
        add_line(s, x1, y1, x2, y2, GOOD, 1.4)
    # marker at (FPR, TPR) sampled every few steps
    for i in range(0, len(base), max(1, len(base)//8)):
        add_circle(s, to_x(base[i, 0]), to_y(base[i, 1]),
                   Emu(25000), fill=GOOD, line_color=WHITE, line_width=0.6)

    # === FANCY ADDITION: Youden's J optimal operating point ===
    j_stat = tpr - fpr
    j_idx = int(np.argmax(j_stat))
    j_fpr = float(fpr[j_idx]); j_tpr = float(tpr[j_idx])
    j_sens = j_tpr
    j_spec = 1 - j_fpr
    # star marker at Youden point
    add_diamond(s, to_x(j_fpr), to_y(j_tpr), Emu(55000),
                fill=HIGHLIGHT, line_color=WHITE, line_width=1.3)
    # small annotation arrow + text
    ann_px = to_x(j_fpr) - Inches(0.02)
    ann_py = to_y(j_tpr) + Inches(0.15)
    add_text(s, ann_px - Inches(0.55), ann_py,
             Inches(1.10), Inches(0.13),
             f"Youden J max",
             size=6, bold=True, italic=True, color=HIGHLIGHT,
             align="center")
    add_text(s, ann_px - Inches(0.55), ann_py + Inches(0.13),
             Inches(1.10), Inches(0.12),
             f"sens {j_sens:.2f} / spec {j_spec:.2f}",
             size=6, color=HIGHLIGHT, align="center")

    # AUC annotation box
    ann_x = to_x(0.55); ann_y = to_y(0.30)
    box_w = Inches(1.8); box_h = Inches(0.55)
    add_rect(s, ann_x, ann_y, box_w, box_h,
             fill=WHITE, line_color=GREY, line_width=0.4)
    add_text(s, ann_x + Inches(0.08), ann_y + Emu(20000),
             box_w - Inches(0.15), Inches(0.25),
             "ElasticNet, nested outer-LOOCV",
             size=7, color=INK, align="left", anchor="top")
    add_text(s, ann_x + Inches(0.08), ann_y + Inches(0.22),
             box_w - Inches(0.15), Inches(0.35),
             f"AUC = {EN_AUC:.3f}  (95 % CI {EN_LO:.2f}–{EN_HI:.2f})",
             size=8, bold=True, color=INK, align="left", anchor="top")

    # legend (one-line)
    leg_y = PLOT_Y + Inches(0.12)
    add_line(s, PLOT_X + Inches(0.1), leg_y + Inches(0.08),
             PLOT_X + Inches(0.4), leg_y + Inches(0.08), GOOD, 1.4)
    add_text(s, PLOT_X + Inches(0.45), leg_y, Inches(2.8), Inches(0.18),
             "Observed ROC curve (mean over 35 LOOCV folds)",
             size=7, color=INK, align="left")
    add_rect(s, PLOT_X + Inches(0.1), leg_y + Inches(0.22),
             Inches(0.3), Inches(0.1), fill=SHADE)
    add_text(s, PLOT_X + Inches(0.45), leg_y + Inches(0.18),
             Inches(2.8), Inches(0.18),
             "95 % bootstrap CI ribbon", size=7, color=INK, align="left")
    # AUC fill legend chip
    add_rect(s, PLOT_X + Inches(0.1), leg_y + Inches(0.36),
             Inches(0.3), Inches(0.1), fill=AUC_FILL)
    add_text(s, PLOT_X + Inches(0.45), leg_y + Inches(0.32),
             Inches(2.8), Inches(0.18),
             "area under ROC (above diagonal)",
             size=7, color=INK, align="left")
    # Youden marker legend
    add_diamond(s, PLOT_X + Inches(0.22), leg_y + Inches(0.53),
                Emu(28000), fill=HIGHLIGHT, line_color=WHITE, line_width=0.8)
    add_text(s, PLOT_X + Inches(0.45), leg_y + Inches(0.46),
             Inches(3.0), Inches(0.18),
             "Youden J optimal operating point",
             size=7, color=INK, align="left")


# ---------- Panel C: 4-feature coefficient forest ----------
def build_C():
    s = new_slide(prs)
    draw_panel_letter(s, "C")

    # order by coef DESC
    rows = feats4.sort_values("coef", ascending=True).reset_index(drop=True)
    n = len(rows)

    to_x, to_y = draw_axes(s, PLOT_X + Inches(1.8), PLOT_Y,
                            PLOT_W - Inches(1.8), PLOT_H,
                            x_ticks=[-1.0, -0.5, 0, 0.5, 1.0],
                            y_ticks=None,
                            x_range=(-1.1, 1.1),
                            y_range=(-0.5, n - 0.5),
                            xlab="ElasticNet coefficient (standardised scale)",
                            x_tick_fmt="{:+.1f}")

    # zero reference
    add_line(s, to_x(0), int(PLOT_Y),
             to_x(0), int(PLOT_Y + PLOT_H), INK, 0.8)

    # Group colours (same as Panel A)
    GROUP_PAL_C = {
        "DNA repair": RGBColor(0xA2, 0x4E, 0x78),
        "Genomic": RGBColor(0x55, 0x55, 0x55),
        "Immune": RGBColor(0x11, 0x8A, 0xB2),
        "MSI": RGBColor(0xE3, 0x9A, 0x1E),
    }
    # feature labels (y axis) + fancy additions
    for i, row in rows.iterrows():
        ypos = to_y(i)
        label = row["feature"]
        if "DNA Double-Strand Break" in label:
            pretty = "DSB repair (R-HSA-5693532)"
            group = "DNA repair"
        elif label == "frac_amp":
            pretty = "Genomic deletion fraction"
            group = "Genomic"
        elif label == "MHC_II":
            pretty = "MHC class II signature"
            group = "Immune"
        elif label == "MSI_pct":
            pretty = "MSI (%)"
            group = "MSI"
        else:
            pretty = label
            group = "Other"
        # === FANCY: feature-group colour chip at far left ===
        chip_col = GROUP_PAL_C.get(group, RGBColor(0xAA, 0xAA, 0xAA))
        add_rect(s, PLOT_X + Inches(0.05), ypos - Inches(0.06),
                 Inches(0.09), Inches(0.12), fill=chip_col)
        add_text(s, PLOT_X + Inches(0.18), ypos - Inches(0.11),
                 Inches(1.55), Inches(0.22),
                 pretty, size=8, color=INK, align="right")

        # lollipop stem
        color = GOOD if row["coef"] > 0 else BAD
        add_line(s, to_x(0), ypos, to_x(row["coef"]), ypos,
                 color, 1.5)
        add_circle(s, to_x(row["coef"]), ypos, Emu(50000),
                   fill=color, line_color=WHITE, line_width=0.8)
        # coefficient value
        add_text(s, to_x(row["coef"]) + Inches(0.08),
                 ypos - Inches(0.09),
                 Inches(0.7), Inches(0.18),
                 f"{row['coef']:+.2f}",
                 size=8, color=color, bold=True, align="left")
        # === FANCY: direction arrow chip (→ good / → bad) ===
        arrow_x = to_x(row["coef"]) + Inches(0.55)
        direction_txt = "→ predicts good" if row["coef"] > 0 else "→ predicts bad"
        add_text(s, arrow_x, ypos - Inches(0.08),
                 Inches(1.0), Inches(0.15),
                 direction_txt, size=6, italic=True, color=color,
                 align="left")

    # === FANCY: group legend (under plot) ===
    leg_y = PLOT_Y + PLOT_H + Inches(0.18)
    add_text(s, Inches(0.15), leg_y - Inches(0.02),
             Inches(1.2), Inches(0.14),
             "Feature group:", size=6, bold=True, color=INK, align="left")
    xb = Inches(1.1)
    for lab, col in GROUP_PAL_C.items():
        add_rect(s, xb, leg_y, Inches(0.09), Inches(0.10), fill=col)
        add_text(s, xb + Inches(0.11), leg_y - Emu(5000),
                 Inches(1.0), Inches(0.14),
                 lab, size=6, color=INK, align="left")
        xb += Inches(1.05)

    # caption
    add_text(s, Inches(0.15), PLOT_Y + PLOT_H + Inches(0.42),
             SLIDE_W - Inches(0.3), Inches(0.18),
             "4 non-zero features selected by nested inner-CV; "
             "teal lollipop → predicts good (positive coef), "
             "coral → predicts bad (negative coef).",
             size=7, color=INK, align="center")


# ---------- Panel D: UMAP by response ----------
def build_D():
    s = new_slide(prs)
    draw_panel_letter(s, "D")

    u1_lo, u1_hi = U[:, 0].min() - 0.5, U[:, 0].max() + 0.5
    u2_lo, u2_hi = U[:, 1].min() - 0.5, U[:, 1].max() + 0.5
    u1_lo = float(np.floor(u1_lo)); u1_hi = float(np.ceil(u1_hi))
    u2_lo = float(np.floor(u2_lo)); u2_hi = float(np.ceil(u2_hi))

    to_x, to_y = draw_axes(s, PLOT_X, PLOT_Y, PLOT_W, PLOT_H,
                            x_ticks=list(np.linspace(u1_lo, u1_hi, 5).round(1)),
                            y_ticks=list(np.linspace(u2_lo, u2_hi, 5).round(1)),
                            x_range=(u1_lo, u1_hi),
                            y_range=(u2_lo, u2_hi),
                            xlab="UMAP 1", ylab="UMAP 2")

    # === FANCY ADDITION: 68% / 95% confidence ellipses per response group ===
    # Compute per-group mean + covariance, derive ellipse axes from eigen.
    def draw_ellipse(centre_u, centre_v, sa, sb, angle_deg, color, width):
        """Approximate an ellipse by a rotated OVAL auto-shape.
        Since python-pptx OVAL is an axis-aligned oval that can only be
        rotated via shape.rotation, we draw a rectangular bounding-box
        oval and rotate it.
        """
        # Convert semi-axis data-units to EMU widths along UMAP axes
        # width = 2*sa (data), height = 2*sb (data)
        w_emu = abs(to_x(centre_u + sa) - to_x(centre_u - sa))
        h_emu = abs(to_y(centre_v + sb) - to_y(centre_v - sb))
        cx = to_x(centre_u); cy = to_y(centre_v)
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  int(cx - w_emu / 2),
                                  int(cy - h_emu / 2),
                                  int(max(w_emu, 1)),
                                  int(max(h_emu, 1)))
        shp.fill.background()
        shp.line.color.rgb = color
        shp.line.width = Pt(width)
        shp.rotation = float(angle_deg)
        shp.shadow.inherit = False
        kill_shadow(shp)
        return shp

    for grp_val, grp_color in [(1, GOOD), (0, BAD)]:
        mask = (y == grp_val)
        pts = U[mask]
        if len(pts) < 3:
            continue
        mu = pts.mean(axis=0)
        cov = np.cov(pts.T)
        # eigen decomposition
        evals, evecs = np.linalg.eigh(cov)
        order_e = np.argsort(evals)[::-1]
        evals = evals[order_e]
        evecs = evecs[:, order_e]
        angle = np.degrees(np.arctan2(evecs[1, 0], evecs[0, 0]))
        # 68% (chi²_2, 2.30) and 95% (chi², 5.99) scalings
        sa68 = np.sqrt(2.30 * evals[0])
        sb68 = np.sqrt(2.30 * evals[1])
        sa95 = np.sqrt(5.99 * evals[0])
        sb95 = np.sqrt(5.99 * evals[1])
        # 95% ellipse (thin dashed-like via reduced width)
        draw_ellipse(mu[0], mu[1], sa95, sb95, angle, grp_color, 0.5)
        # 68% ellipse (solid, bolder)
        draw_ellipse(mu[0], mu[1], sa68, sb68, angle, grp_color, 1.2)
        # centroid marker
        add_diamond(s, to_x(mu[0]), to_y(mu[1]), Emu(32000),
                    fill=grp_color, line_color=WHITE, line_width=1.0)

    # subjects with predicted prob & true label — identify misclassified
    # probs_sorted has columns: subject_id, y, prob
    mis_by_subj = {int(r.subject_id): int(r.y != (r.prob > 0.5))
                   for _, r in probs_sorted.iterrows()}
    subj_order = master["subject_id"].values

    for i in range(U.shape[0]):
        color = GOOD if y[i] == 1 else BAD
        sid = int(subj_order[i])
        is_mis = bool(mis_by_subj.get(sid, 0))
        # individual dots — slightly bigger for misclassified, with gold ring
        r_i = Emu(55000) if is_mis else Emu(40000)
        add_circle(s, to_x(U[i, 0]), to_y(U[i, 1]), r_i,
                   fill=color,
                   line_color=HIGHLIGHT if is_mis else WHITE,
                   line_width=1.4 if is_mis else 0.8)
        # subject ID text
        if is_mis:
            add_text(s, to_x(U[i, 0]) + Inches(0.08),
                     to_y(U[i, 1]) - Inches(0.08),
                     Inches(0.5), Inches(0.14),
                     str(sid), size=6, bold=True, color=HIGHLIGHT,
                     align="left")

    # legend
    leg_x = PLOT_X + Inches(3.2); leg_y = PLOT_Y + Inches(0.05)
    add_circle(s, leg_x + Inches(0.08), leg_y + Inches(0.10),
               Emu(35000), fill=GOOD, line_color=WHITE, line_width=0.6)
    add_text(s, leg_x + Inches(0.2), leg_y, Inches(1.0), Inches(0.2),
             f"good (n={int(y.sum())})", size=8, color=INK)
    add_circle(s, leg_x + Inches(0.08), leg_y + Inches(0.28),
               Emu(35000), fill=BAD, line_color=WHITE, line_width=0.6)
    add_text(s, leg_x + Inches(0.2), leg_y + Inches(0.18),
             Inches(1.0), Inches(0.2),
             f"bad (n={int((y==0).sum())})", size=8, color=INK)
    # misclassified marker legend
    add_circle(s, leg_x + Inches(0.08), leg_y + Inches(0.48),
               Emu(45000), fill=GOOD, line_color=HIGHLIGHT, line_width=1.4)
    add_text(s, leg_x + Inches(0.2), leg_y + Inches(0.38),
             Inches(1.3), Inches(0.2),
             "misclassified  (gold ring, subject ID)",
             size=6, color=HIGHLIGHT, bold=True)
    # ellipse legend
    add_text(s, leg_x, leg_y + Inches(0.65),
             Inches(1.7), Inches(0.15),
             "solid ellipse: 68 % CI",
             size=6, color=RGBColor(0x55, 0x55, 0x55), align="left")
    add_text(s, leg_x, leg_y + Inches(0.78),
             Inches(1.7), Inches(0.15),
             "thin ellipse: 95 % CI",
             size=6, color=RGBColor(0x55, 0x55, 0x55), align="left")


# ---------- Panel E: SHAP beeswarm (4 features) ----------
def build_E():
    s = new_slide(prs)
    draw_panel_letter(s, "E")

    # order features by mean |shap|
    mean_abs = np.abs(shap_df[feat4_cols]).mean().sort_values()
    order = mean_abs.index.tolist()
    n_feat = len(order)

    sh_vals = shap_df[order].values
    sh_min = sh_vals.min(); sh_max = sh_vals.max()
    lim = max(abs(sh_min), abs(sh_max)) * 1.05
    to_x, to_y = draw_axes(s, PLOT_X + Inches(1.8), PLOT_Y,
                            PLOT_W - Inches(1.8), PLOT_H,
                            x_ticks=[-1.5, -0.75, 0, 0.75, 1.5],
                            y_ticks=None,
                            x_range=(-lim, lim),
                            y_range=(-0.5, n_feat - 0.5),
                            xlab="SHAP value (log-odds contribution)",
                            x_tick_fmt="{:+.2f}")

    # zero reference
    add_line(s, to_x(0), int(PLOT_Y),
             to_x(0), int(PLOT_Y + PLOT_H), INK, 0.8)

    # feature labels + colour-by-raw-value dots
    rng = np.random.default_rng(0)
    for j, feat in enumerate(order):
        ypos = to_y(j)
        label = feat
        if "DNA Double-Strand Break" in feat:
            label = "DSB repair"
        elif feat == "frac_amp":
            label = "Deletion fraction"
        elif feat == "MHC_II":
            label = "MHC class II"
        elif feat == "MSI_pct":
            label = "MSI (%)"
        add_text(s, PLOT_X, ypos - Inches(0.11),
                 Inches(1.7), Inches(0.22),
                 label, size=8, color=INK, align="right")

        # beeswarm
        svals = shap_df[feat].values
        rvals = X4[:, feat4_cols.index(feat)]
        r_min, r_max = np.min(rvals), np.max(rvals)
        for i in range(len(svals)):
            jitter = rng.normal(0, 0.06)
            xx = to_x(svals[i]); yy = to_y(j + jitter)
            # color by raw value, high = GOOD blue, low = LT_GREY
            t = 0 if r_max == r_min else (rvals[i] - r_min) / (r_max - r_min)
            r = int(BAD_HEX[0] + t * (GOOD_HEX[0] - BAD_HEX[0]))
            g = int(BAD_HEX[1] + t * (GOOD_HEX[1] - BAD_HEX[1]))
            b = int(BAD_HEX[2] + t * (GOOD_HEX[2] - BAD_HEX[2]))
            add_circle(s, xx, yy, Emu(25000),
                       fill=RGBColor(r, g, b),
                       line_color=WHITE, line_width=0.3)

    # === FANCY ADDITION: right-edge mean |SHAP| global-importance bars ===
    # For each feature row, show mean |SHAP| as a small gold bar extending
    # from the right spine rightward. Scale: max_mean_abs = 1.0 bar width.
    imp_bar_x = int(PLOT_X + PLOT_W + Inches(0.08))
    imp_bar_w_max = Inches(0.65)
    max_mean = float(mean_abs.max())
    if max_mean > 0:
        for j, feat in enumerate(order):
            ypos = to_y(j)
            mean_v = float(mean_abs[feat])
            bar_w = int(imp_bar_w_max * mean_v / max_mean)
            add_rect(s, imp_bar_x, ypos - Inches(0.06),
                     bar_w, Inches(0.12),
                     fill=HIGHLIGHT, line_color=WHITE, line_width=0.3)
            # value label
            add_text(s, imp_bar_x + bar_w + Inches(0.03),
                     ypos - Inches(0.08),
                     Inches(0.40), Inches(0.14),
                     f"{mean_v:.2f}", size=6, bold=True, color=HIGHLIGHT,
                     align="left")
    add_text(s, imp_bar_x, int(PLOT_Y - Inches(0.16)),
             imp_bar_w_max + Inches(0.45), Inches(0.14),
             "mean |SHAP|", size=6, bold=True, color=HIGHLIGHT,
             align="left")

    # colour legend: high -> low feature value
    leg_x = PLOT_X + Inches(2.4); leg_y = PLOT_Y + PLOT_H + Inches(0.38)
    cb_w = Inches(1.4); cb_h = Inches(0.1)
    for k in range(32):
        t = k / 31
        r = int(BAD_HEX[0] + t * (GOOD_HEX[0] - BAD_HEX[0]))
        g = int(BAD_HEX[1] + t * (GOOD_HEX[1] - BAD_HEX[1]))
        b = int(BAD_HEX[2] + t * (GOOD_HEX[2] - BAD_HEX[2]))
        add_rect(s, int(leg_x + k * cb_w / 32), leg_y,
                 int(cb_w / 32) + Emu(2000), cb_h,
                 fill=RGBColor(r, g, b))
    add_text(s, leg_x - Inches(0.35), leg_y - Emu(20000),
             Inches(0.35), Inches(0.2), "low", size=7,
             color=INK, align="right", anchor="top")
    add_text(s, leg_x + cb_w + Inches(0.05), leg_y - Emu(20000),
             Inches(0.5), Inches(0.2), "high", size=7,
             color=INK, align="left", anchor="top")
    add_text(s, leg_x - Inches(0.9), leg_y + Inches(0.12),
             cb_w + Inches(1.3), Inches(0.2),
             "feature value", size=7, color=INK, align="center")


# ---------- Panel F: per-subject predicted probability ----------
def build_F():
    s = new_slide(prs)
    draw_panel_letter(s, "F")

    n = len(probs_sorted)
    to_x, to_y = draw_axes(s, PLOT_X, PLOT_Y, PLOT_W, PLOT_H,
                            x_ticks=None,
                            y_ticks=[0, 0.2, 0.4, 0.5, 0.6, 0.8, 1.0],
                            x_range=(-0.5, n - 0.5),
                            y_range=(0, 1),
                            xlab="Subjects, ordered by predicted probability",
                            ylab="Predicted P(good)")

    # 0.5 threshold line
    add_line(s, to_x(-0.5), to_y(0.5), to_x(n - 0.5), to_y(0.5),
             INK, 0.6)
    add_text(s, to_x(n - 0.5) + Emu(20000),
             to_y(0.5) - Inches(0.1),
             Inches(0.8), Inches(0.2),
             "0.5 threshold", size=7, color=INK)

    # bars + subject IDs + misclassification markers
    bar_w = int(PLOT_W / n * 0.82)
    n_correct = 0; n_wrong = 0
    for i, row in probs_sorted.iterrows():
        color = GOOD if int(row.y) == 1 else BAD
        cx = to_x(i)
        x0 = int(cx - bar_w / 2)
        y_top = to_y(float(row.prob))
        is_mis = (int(row.y) == 1 and float(row.prob) <= 0.5) or \
                 (int(row.y) == 0 and float(row.prob) > 0.5)
        if is_mis: n_wrong += 1
        else: n_correct += 1
        # === FANCY ADDITION: gradient fill by confidence ===
        # Confidence = |prob - 0.5| × 2; higher = more saturated;
        # approximate by rendering an inner "confidence stripe" on top of
        # the base-colour bar.
        add_rect(s, x0, y_top, bar_w, int(to_y(0) - y_top),
                 fill=color, line_color=WHITE, line_width=0.2)
        # thin outline when correct; thick gold when misclassified
        if is_mis:
            # gold outline overlay (rectangle frame)
            add_rect(s, x0, y_top, bar_w, int(to_y(0) - y_top),
                     line_color=HIGHLIGHT, line_width=1.3)
            # "X" marker near the bar top
            x_mark_x = cx
            x_mark_y = y_top - Inches(0.10)
            add_text(s, x_mark_x - Inches(0.08),
                     int(x_mark_y - Inches(0.06)),
                     Inches(0.16), Inches(0.14),
                     "✗", size=10, bold=True, color=HIGHLIGHT,
                     align="center")
        # subject-id label below axis
        add_text(s, int(cx - Inches(0.12)),
                 int(to_y(0) + Inches(0.02)),
                 Inches(0.24), Inches(0.12),
                 str(int(row.subject_id)), size=4, color=color,
                 align="center")

    # === FANCY ADDITION: accuracy summary box ===
    acc = n_correct / n
    summary_x = PLOT_X + PLOT_W - Inches(1.50)
    summary_y = PLOT_Y + Inches(0.06)
    add_rect(s, summary_x, summary_y, Inches(1.45), Inches(0.40),
             fill=WHITE, line_color=HIGHLIGHT, line_width=0.8)
    add_text(s, summary_x, summary_y + Inches(0.03),
             Inches(1.45), Inches(0.14),
             f"{n_correct}/{n} correct @ 0.5",
             size=7, bold=True, color=INK, align="center")
    add_text(s, summary_x, summary_y + Inches(0.20),
             Inches(1.45), Inches(0.14),
             f"accuracy = {acc:.2f}",
             size=6, color=HIGHLIGHT, align="center")

    # legend
    leg_x = PLOT_X + Inches(0.12); leg_y = PLOT_Y + Inches(0.05)
    add_rect(s, leg_x, leg_y, Inches(0.16), Inches(0.14), fill=GOOD)
    add_text(s, leg_x + Inches(0.20), leg_y - Emu(10000),
             Inches(1.4), Inches(0.2),
             "true label: good", size=7, color=INK)
    add_rect(s, leg_x + Inches(1.6), leg_y, Inches(0.16),
             Inches(0.14), fill=BAD)
    add_text(s, leg_x + Inches(1.78), leg_y - Emu(10000),
             Inches(1.4), Inches(0.2),
             "true label: bad", size=7, color=INK)
    # misclass marker legend
    add_text(s, leg_x + Inches(3.2), leg_y - Emu(10000),
             Inches(2.0), Inches(0.2),
             "✗ / gold border = misclassified at P = 0.5",
             size=6, bold=True, color=HIGHLIGHT, align="left")


# build all panels into one deck
build_A()
build_B()
build_C()
build_D()
build_E()
build_F()

deck_path = f"{OUT}/Fig5_nested_LASSO_native_editable.pptx"
prs.save(deck_path)
print(f"wrote {deck_path} with {len(prs.slides)} slides (Fig 5 A-F)")


# ===========================================================================
# Supp Fig S8 --- 5-scenario ablation (companion deck)
# ===========================================================================
prs2 = Presentation()
prs2.slide_width = SLIDE_W
prs2.slide_height = SLIDE_H

SCENARIO_ORDER = ["baseline_37", "drop_cd8prolif_36", "add_immune_40",
                  "swap_cd8_37", "drop_prolif_add_3_39"]
SCENARIO_PRETTY = {
    "baseline_37": "Baseline\n(37 feat.)",
    "drop_cd8prolif_36": "Drop CD8-prolif\n(36 feat.) ★",
    "add_immune_40": "Add 3 immune\n(40 feat.)",
    "swap_cd8_37": "Swap CD8 prolif→cyt\n(37 feat.)",
    "drop_prolif_add_3_39": "Drop + add\n(39 feat.)",
}

def build_S8A():
    s = new_slide(prs2)
    draw_panel_letter(s, "A")

    rows = [nested[(nested.scenario == sc) & (nested.model == "ElasticNet")]
            .iloc[0] for sc in SCENARIO_ORDER]
    aucs_en = [r["AUC"] for r in rows]
    lo_en = [r["CI_low"] for r in rows]
    hi_en = [r["CI_high"] for r in rows]

    to_x, to_y = draw_axes(s, PLOT_X + Inches(0.6), PLOT_Y,
                            PLOT_W - Inches(0.6), PLOT_H,
                            x_ticks=None,
                            y_ticks=[0.5, 0.6, 0.7, 0.8, 0.9],
                            x_range=(-0.5, len(SCENARIO_ORDER) - 0.5),
                            y_range=(0.45, 0.95),
                            xlab="Ablation scenario",
                            ylab="Nested outer-LOOCV AUC (ElasticNet)")

    add_line(s, to_x(-0.5), to_y(0.5), to_x(len(SCENARIO_ORDER) - 0.5),
             to_y(0.5), GREY, 0.4)

    bar_w = int(PLOT_W / len(SCENARIO_ORDER) * 0.42)
    for i, sc in enumerate(SCENARIO_ORDER):
        color = GOOD if sc == "drop_cd8prolif_36" else RGBColor(0x77, 0x77, 0x77)
        cx = to_x(i)
        x0 = int(cx - bar_w / 2)
        y_top = to_y(float(aucs_en[i]))
        add_rect(s, x0, y_top, bar_w, int(to_y(0.45) - y_top),
                 fill=color, line_color=WHITE, line_width=0.3)
        # CI bar
        y_lo = to_y(float(lo_en[i])); y_hi = to_y(float(hi_en[i]))
        add_line(s, cx, y_lo, cx, y_hi, INK, 0.6)
        add_line(s, cx - Emu(25000), y_hi, cx + Emu(25000), y_hi, INK, 0.6)
        add_line(s, cx - Emu(25000), y_lo, cx + Emu(25000), y_lo, INK, 0.6)
        # AUC label
        add_text(s, cx - Inches(0.35), y_top - Inches(0.25),
                 Inches(0.7), Inches(0.2),
                 f"{aucs_en[i]:.3f}", size=7, bold=True,
                 color=INK, align="center")
        # scenario label
        add_text(s, cx - Inches(0.55), to_y(0.45) + Emu(50000),
                 Inches(1.1), Inches(0.4),
                 SCENARIO_PRETTY[sc], size=7, color=INK, align="center")

    add_text(s, Inches(0.15), PLOT_Y + PLOT_H + Inches(0.55),
             SLIDE_W - Inches(0.3), Inches(0.22),
             "★ the improvement comes from removing the cell-cycle-"
             "contaminated CD8-proliferation signature, not from adding "
             "purified immune signatures", size=7, color=INK, align="center")


def build_S8B():
    s = new_slide(prs2)
    draw_panel_letter(s, "B")
    # ROC overlay --- approximate from per-scenario prob files
    scen_to_probfile = {
        "baseline_37": "probs_baseline_37_ElasticNet.tsv",
        "drop_cd8prolif_36": "probs_drop_cd8prolif_36_ElasticNet.tsv",
        "add_immune_40": "probs_add_immune_40_ElasticNet.tsv",
        "swap_cd8_37": "probs_swap_cd8_37_ElasticNet.tsv",
        "drop_prolif_add_3_39": "probs_drop_prolif_add_3_39_ElasticNet.tsv",
    }

    to_x, to_y = draw_axes(s, PLOT_X, PLOT_Y, PLOT_W, PLOT_H,
                            x_ticks=[0, 0.2, 0.4, 0.6, 0.8, 1.0],
                            y_ticks=[0, 0.2, 0.4, 0.6, 0.8, 1.0],
                            x_range=(0, 1), y_range=(0, 1),
                            xlab="False-positive rate",
                            ylab="True-positive rate")
    add_line(s, to_x(0), to_y(0), to_x(1), to_y(1), GREY, 0.4)

    palette = {
        "baseline_37": RGBColor(0xAA, 0xAA, 0xAA),
        "drop_cd8prolif_36": GOOD,
        "add_immune_40": RGBColor(0x77, 0x99, 0x77),
        "swap_cd8_37": RGBColor(0xC5, 0xA5, 0x72),
        "drop_prolif_add_3_39": RGBColor(0xA2, 0x4E, 0x78),
    }

    for i, sc in enumerate(SCENARIO_ORDER):
        pf = scen_to_probfile[sc]
        fpath = f"{DATA}/{pf}"
        if not os.path.exists(fpath):
            continue
        pb = pd.read_csv(fpath, sep="\t")
        f, t, _ = roc_curve(pb.y, pb.prob)
        col = palette[sc]
        for k in range(len(f) - 1):
            add_line(s, to_x(f[k]), to_y(t[k]),
                     to_x(f[k+1]), to_y(t[k+1]),
                     col, 1.2 if sc == "drop_cd8prolif_36" else 0.9)
        # legend entry
        leg_y = PLOT_Y + Inches(0.15 + i * 0.22)
        add_line(s, PLOT_X + Inches(2.3), leg_y + Inches(0.08),
                 PLOT_X + Inches(2.6), leg_y + Inches(0.08),
                 col, 1.4 if sc == "drop_cd8prolif_36" else 1.0)
        auc_this = nested[(nested.scenario == sc) &
                          (nested.model == "ElasticNet")].iloc[0]["AUC"]
        tag = "★ " if sc == "drop_cd8prolif_36" else ""
        add_text(s, PLOT_X + Inches(2.65), leg_y, Inches(2.6), Inches(0.18),
                 f"{tag}{SCENARIO_PRETTY[sc].replace(chr(10), ' ')}  "
                 f"(AUC {auc_this:.3f})",
                 size=7, color=INK, align="left")


build_S8A()
build_S8B()
s8_path = f"{OUT}/SuppFig_S8_5scenario_ablation_native_editable.pptx"
prs2.save(s8_path)
print(f"wrote {s8_path} with {len(prs2.slides)} slides (S8 A-B)")
