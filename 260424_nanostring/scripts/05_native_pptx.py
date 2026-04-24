#!/usr/bin/env python3
"""
Native-editable PowerPoint decks for Supp Figs S22–S28 (NanoString orthogonal
validation, §3.14).

One slide = one panel. Every in-plot text element is a TEXT_BOX (Arial, 8 pt
axis ticks, 9-10 pt titles, 6.5-pt annotation). Every line is a CONNECTOR,
every bar / marker is an AUTO_SHAPE. Shadow removed from every shape via
`kill_shadow(...)`.

Palette (project standard 2026-04-18, per MEMORY.md):
  GOOD = #0A7D6E  (deep teal)
  BAD  = #C53E1F  (deep coral)
  GOLD = #D4A300  (highlight / ceiling hit star)

Decks produced:
  FigS22_NanoString_composite_native_editable.pptx       (6 panels)
  FigS23_NanoString_prepostdelta_heatmap_native_editable.pptx  (1 panel)
  FigS24_NanoString_canonical_signatures_native_editable.pptx  (6 panels)
  FigS25_IBI_vs_IAE_fingerprint_native_editable.pptx     (2 panels)
  FigS26_subject_radar_native_editable.pptx              (3 panels)  -- linear alternative
  FigS27_NanoString_prespec_primary_paired_native_editable.pptx   (4 panels)
  FigS28_NanoString_platform_concordance_native_editable.pptx     (2 panels)
"""

from __future__ import annotations
import os
from pathlib import Path

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path("/mnt/sda1/data/TNT/analysis/260424_nanostring")
TABLES = ROOT / "tables"
OUT = ROOT / "manuscript" / "figures"
OUT.mkdir(parents=True, exist_ok=True)

# ---- palette ----
GOOD = RGBColor(0x0A, 0x7D, 0x6E)
BAD = RGBColor(0xC5, 0x3E, 0x1F)
GOLD = RGBColor(0xD4, 0xA3, 0x00)
INK = RGBColor(0x22, 0x22, 0x22)
LINE = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0x88, 0x88, 0x88)
LTGRY = RGBColor(0xDD, 0xDD, 0xDD)
PAPER = RGBColor(0xF6, 0xF6, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x1F, 0x5B, 0x82)
FONT = "Arial"

# slide size: 10 x 6 inches (landscape, fits one panel)
SLIDE_W_IN = 10.0
SLIDE_H_IN = 6.0


# =====================  helpers  =====================
def kill_shadow(shape):
    elem = shape._element
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        return
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def add_textbox(slide, x_in, y_in, w_in, h_in, text, *,
                size=8, bold=False, italic=False,
                color=INK, align="left", anchor="top", family=FONT):
    tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in),
                                   Inches(w_in), Inches(h_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.text = ""
    para = tf.paragraphs[0]
    para.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
    r = para.add_run()
    r.text = str(text)
    r.font.name = family
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                            "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    kill_shadow(tb)
    return tb


def add_line(slide, x1_in, y1_in, x2_in, y2_in, *,
              color=LINE, width_pt=0.6, dash=False):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     Inches(x1_in), Inches(y1_in),
                                     Inches(x2_in), Inches(y2_in))
    ln.line.color.rgb = color
    ln.line.width = Pt(width_pt)
    if dash:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        ln.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    kill_shadow(ln)
    return ln


def add_rect(slide, x_in, y_in, w_in, h_in, *, fill=None, line_color=LINE,
             line_width_pt=0.4, line=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x_in), Inches(y_in),
                                  Inches(w_in), Inches(h_in))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width_pt)
    else:
        shp.line.fill.background()
    kill_shadow(shp)
    return shp


def add_oval(slide, cx_in, cy_in, d_in=0.12, *, fill=INK, line_color=LINE,
              line_width_pt=0.3):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(cx_in - d_in / 2),
                                  Inches(cy_in - d_in / 2),
                                  Inches(d_in), Inches(d_in))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color
    shp.line.width = Pt(line_width_pt)
    kill_shadow(shp)
    return shp


def add_diamond(slide, cx_in, cy_in, d_in=0.14, *, fill=GOLD, line_color=LINE,
                 line_width_pt=0.3):
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                  Inches(cx_in - d_in / 2),
                                  Inches(cy_in - d_in / 2),
                                  Inches(d_in), Inches(d_in))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color
    shp.line.width = Pt(line_width_pt)
    kill_shadow(shp)
    return shp


def interp_color(v, vmin, vmax):
    """Interpolate between BAD (negative) and GOOD (positive) through WHITE."""
    if vmax == vmin:
        return WHITE
    t = (v - vmin) / (vmax - vmin)
    t = max(0.0, min(1.0, t))
    if t < 0.5:
        # BAD to WHITE
        f = t * 2
        r = int(0xC5 + (0xFF - 0xC5) * f)
        g = int(0x3E + (0xFF - 0x3E) * f)
        b = int(0x1F + (0xFF - 0x1F) * f)
    else:
        # WHITE to GOOD
        f = (t - 0.5) * 2
        r = int(0xFF - (0xFF - 0x0A) * f)
        g = int(0xFF - (0xFF - 0x7D) * f)
        b = int(0xFF - (0xFF - 0x6E) * f)
    return RGBColor(r, g, b)


def panel_frame(slide, px, py, pw, ph, *, title=None, subtitle=None,
                 xlab=None, ylab=None):
    """Draw an axis rectangle and title/subtitle/axis labels."""
    # plot box spines
    add_line(slide, px, py + ph, px + pw, py + ph, color=LINE, width_pt=0.6)   # bottom
    add_line(slide, px, py, px, py + ph, color=LINE, width_pt=0.6)             # left
    if title:
        add_textbox(slide, px - 0.3, py - 0.45, pw + 0.6, 0.3,
                     title, size=11, bold=True, align="left")
    if subtitle:
        add_textbox(slide, px - 0.3, py - 0.18, pw + 0.6, 0.2,
                     subtitle, size=8, color=GREY, align="left", italic=True)
    if xlab:
        add_textbox(slide, px, py + ph + 0.30, pw, 0.25,
                     xlab, size=9, align="center")
    if ylab:
        # vertical label simulated by rotated textbox
        tb = add_textbox(slide, px - 0.75, py + ph / 2 - 0.4, 0.6, 0.8,
                          ylab, size=9, align="center")
        tb.rotation = -90


def add_legend_chip(slide, x, y, label, color, size=8):
    add_rect(slide, x, y, 0.18, 0.12, fill=color, line=True,
             line_color=LINE, line_width_pt=0.3)
    add_textbox(slide, x + 0.22, y - 0.02, 1.8, 0.18, label,
                 size=size, align="left")


# =====================  Fig S22 — 6-panel composite  =====================
def build_figS22():
    prs = new_prs()
    # Panel layouts: grid 2x3, each slide = 1 panel + 1 title slide
    # To reduce slide count, we'll compress the 6 panels to 1 overview slide + 6 panel slides
    # Header slide
    hdr = blank_slide(prs)
    add_textbox(hdr, 0.5, 1.5, 9.0, 1.0,
                 "Supp Fig S22. NanoString PanCancer Immune orthogonal validation",
                 size=16, bold=True, align="center")
    add_textbox(hdr, 0.5, 2.6, 9.0, 0.8,
                 "6-subject extreme-phenotype (3 pCR vs 3 poor) × paired pre/post × 730 probes.\n"
                 "Six panels: A pre/post/Δ heatmap · B pre-treatment direction waterfall "
                 "· C regulatory-grade signatures · D IAE vs IBI composite · E IAE vs IBI genes "
                 "· F Pre-spec cascade arrow 4 (ΔCD8-exh × ΔTLS-8, r=+0.82, P=0.046).",
                 size=10, color=GREY, align="center")

    # --- Panel A: Heatmap (composite × timepoint) ---
    pre_mw = pd.read_csv(TABLES / "v2_pre_MW.tsv", sep="\t").set_index("composite")
    post_mw = pd.read_csv(TABLES / "v2_post_MW.tsv", sep="\t").set_index("composite")
    dlt_mw = pd.read_csv(TABLES / "v2_delta_MW.tsv", sep="\t").set_index("composite")
    order = pre_mw.sort_values("MW_P_1s_good_gt_bad").index.tolist()
    pre_delta = (pre_mw.loc[order, "good_mean"] - pre_mw.loc[order, "bad_mean"]).values
    post_delta = (post_mw.loc[order, "good_mean"] - post_mw.loc[order, "bad_mean"]).values
    dlt_delta = (dlt_mw.loc[order, "good_mean"] - dlt_mw.loc[order, "bad_mean"]).values
    M = np.vstack([pre_delta, post_delta, dlt_delta]).T
    P = np.vstack([pre_mw.loc[order, "MW_P_1s_good_gt_bad"].values,
                    post_mw.loc[order, "MW_P_1s_good_gt_bad"].values,
                    dlt_mw.loc[order, "MW_P_1s_good_gt_bad"].values]).T
    vmax = float(np.abs(M).max())
    vmin = -vmax

    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.15, 9.0, 0.3,
                 "A  Pre/post/Δ composite heatmap (good − bad)", size=13, bold=True)
    add_textbox(s, 0.5, 0.50, 9.0, 0.25,
                 "Star ★ = one-sided MW P ≤ 0.05 (ceiling at n=3 vs 3). Dot · = P ≤ 0.10.",
                 size=8, color=GREY, italic=True)
    n_rows = len(order)
    col_labels = ["pre", "post", "Δ"]
    # Heatmap origin
    hx0 = 3.4; hy0 = 0.95
    cell_w = 0.8; cell_h = min(0.21, 4.6 / n_rows)
    # col headers
    for j, cl in enumerate(col_labels):
        add_textbox(s, hx0 + j * cell_w, hy0 - 0.3, cell_w, 0.25,
                     cl, size=10, bold=True, align="center")
    for i, sig in enumerate(order):
        add_textbox(s, hx0 - 2.7, hy0 + i * cell_h - 0.02, 2.6, cell_h + 0.02,
                     sig, size=7, align="right", anchor="middle")
        for j in range(3):
            v = M[i, j]
            fill = interp_color(v, vmin, vmax)
            add_rect(s, hx0 + j * cell_w, hy0 + i * cell_h,
                      cell_w, cell_h, fill=fill, line=True, line_color=GREY,
                      line_width_pt=0.2)
            # text
            mark = "★" if P[i, j] <= 0.05 else ("·" if P[i, j] <= 0.10 else "")
            val = f"{v:+.2f} {mark}".strip()
            txt_col = WHITE if abs(v) > vmax * 0.6 else INK
            add_textbox(s, hx0 + j * cell_w, hy0 + i * cell_h, cell_w, cell_h,
                         val, size=6, color=txt_col, align="center", anchor="middle")
    # colorbar (vertical strip)
    cbx = hx0 + 3 * cell_w + 0.4; cby = hy0
    cbh = n_rows * cell_h; cbw = 0.25
    n_steps = 40
    for k in range(n_steps):
        v = vmin + (vmax - vmin) * k / (n_steps - 1)
        add_rect(s, cbx, cby + cbh - (k + 1) * cbh / n_steps,
                  cbw, cbh / n_steps, fill=interp_color(v, vmin, vmax),
                  line=False)
    add_textbox(s, cbx, cby - 0.3, cbw + 0.6, 0.25,
                 f"+{vmax:.2f}", size=7, color=GREY)
    add_textbox(s, cbx, cby + cbh + 0.04, cbw + 0.6, 0.25,
                 f"{-vmax:.2f}", size=7, color=GREY)
    add_textbox(s, cbx + cbw + 0.1, cby + cbh / 2 - 0.12, 1.0, 0.25,
                 "good − bad (z)", size=7, color=GREY)

    # --- Panel B: Pre waterfall ---
    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.15, 9.0, 0.3,
                 "B  Pre-treatment direction waterfall — 23/23 good > bad",
                 size=13, bold=True)
    add_textbox(s, 0.5, 0.50, 9.0, 0.25,
                 "Six composites reach ceiling P = 0.05 (Ayers TIS, IFN-γ 6/10, CD8 cytotoxic, M1 macrophage, GC-TF).",
                 size=8, color=GREY, italic=True)
    # bar axis 1x: use full width
    bx0 = 3.5; by0 = 0.95; bh = 0.19
    bar_len_px = 4.0
    bar_max = max(abs(pre_delta.min()), abs(pre_delta.max()))
    zero_x = bx0
    add_line(s, zero_x, by0 - 0.1, zero_x, by0 + n_rows * bh + 0.1,
              color=LINE, width_pt=0.6)
    for i, sig in enumerate(order):
        v = pre_delta[i]
        p = pre_mw.loc[sig, "MW_P_1s_good_gt_bad"]
        add_textbox(s, bx0 - 2.9, by0 + i * bh - 0.02, 2.7, bh + 0.04,
                     sig, size=7, align="right", anchor="middle")
        bar_w = abs(v) / bar_max * bar_len_px
        color = GOOD if v > 0 else BAD
        bx_start = zero_x if v > 0 else zero_x - bar_w
        add_rect(s, bx_start, by0 + i * bh + 0.02, bar_w, bh - 0.05,
                  fill=color, line=True, line_color=LINE, line_width_pt=0.25)
        star = " ★" if p <= 0.05 else ""
        tx = bx_start + bar_w + 0.05 if v > 0 else bx_start - 0.4
        add_textbox(s, tx, by0 + i * bh - 0.02, 0.7, bh + 0.04,
                     f"{v:+.2f}{star}", size=6.5,
                     color=(GOLD if star else GREY),
                     align=("left" if v > 0 else "right"), anchor="middle")
    add_textbox(s, zero_x - 2.0, by0 + n_rows * bh + 0.1, 4.0, 0.25,
                 "good − bad (composite z)", size=9, align="center")
    add_textbox(s, zero_x - 2.0, by0 + n_rows * bh + 0.4, 4.0, 0.25,
                 "★ = one-sided MW P ≤ 0.05 (ceiling)", size=7.5,
                 color=GOLD, align="center", italic=True)

    # --- Panel C: canonical signatures ---
    build_six_panel_canonical(prs, slide_title="C  Regulatory-grade signatures")

    # --- Panel D: IAE vs IBI composite ---
    iae = pd.read_csv(TABLES / "v2_IAE_vs_IBI_descriptive.tsv", sep="\t")
    build_bar_horizontal_panel(
        prs, iae.head(14), title="D  IAE (n=2) vs IBI (n=3) — composite fingerprint",
        subtitle="Positive (teal) = IAE > IBI; negative (coral) = IBI > IAE.",
        label_col="feature", value_col="IAE_minus_IBI")

    # --- Panel E: IAE vs IBI genes ---
    genes = pd.read_csv(TABLES / "v2_IAE_vs_IBI_gene_descriptive.tsv", sep="\t").head(18)
    build_bar_horizontal_panel(
        prs, genes, title="E  Top 18 gene discriminators",
        subtitle="IAE productive cascade (JAK3, IKBKB, TAPBP, ZAP70, BATF) vs IBI suppressive (KIR_Inh, ARG2, RAG1).",
        label_col="gene", value_col="IAE_minus_IBI", font_mono=True)

    # --- Panel F: Pre-spec T2 ---
    cd = pd.read_csv(TABLES / "composite_subject_delta.tsv", sep="\t", index_col=0)
    cohort = [(2,"good"),(4,"good"),(14,"good"),(10,"bad"),(11,"bad"),(13,"bad")]
    xs = cd.loc[[f"subj_{s}" for s,_ in cohort], "CD8_exh"].values
    ys = cd.loc[[f"subj_{s}" for s,_ in cohort], "TLS_8"].values
    build_scatter_panel(
        prs, xs, ys, cohort,
        title="F  Pre-spec cascade arrow 4 (T2): ΔCD8-exh × ΔTLS-8",
        subtitle="Pearson r = +0.82, P = 0.046 · Spearman ρ = +0.77, P = 0.072 (n = 6 paired)",
        xlab="Δ CD8 exhaustion (z)", ylab="Δ TLS-8 (z)")

    fpath = OUT / "FigS22_NanoString_composite_native_editable.pptx"
    prs.save(str(fpath))
    print(f"wrote {fpath}")


def build_six_panel_canonical(prs, slide_title):
    """Panel C: 2×3 mini-scatter per signature (pre vs post, good vs bad)."""
    pre_df = pd.read_csv(TABLES / "v2_composite_pre.tsv", sep="\t", index_col="subject")
    post_df = pd.read_csv(TABLES / "v2_composite_post.tsv", sep="\t", index_col="subject")
    pre_mw = pd.read_csv(TABLES / "v2_pre_MW.tsv", sep="\t").set_index("composite")
    post_mw = pd.read_csv(TABLES / "v2_post_MW.tsv", sep="\t").set_index("composite")
    targets = [("Ayers_TIS", "Ayers TIS"), ("IFNg_6", "IFN-γ 6-gene"),
               ("IFNg_10_Ayers", "IFN-γ 10-gene"), ("CD8_cytotoxic", "CD8 cytotoxic"),
               ("IMPRES_pos", "IMPRES"), ("M1_macro", "M1 macrophage")]
    cohort = [(2,"good"),(4,"good"),(14,"good"),(10,"bad"),(11,"bad"),(13,"bad")]

    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.15, 9.0, 0.3, slide_title, size=13, bold=True)
    add_textbox(s, 0.5, 0.50, 9.0, 0.25,
                 "Six FDA/regulatory-grade immune signatures: pre-RT (left) vs post-RT (right), good (teal) vs bad (coral).",
                 size=8, color=GREY, italic=True)

    rng = np.random.default_rng(42)
    # layout 2 rows x 3 cols
    px0, py0 = 0.6, 1.0
    pw, ph = 2.9, 2.25
    gap_x, gap_y = 0.3, 0.4
    for i, (key, title) in enumerate(targets):
        r = i // 3; c = i % 3
        x = px0 + c * (pw + gap_x)
        y = py0 + r * (ph + gap_y)
        # plot box
        add_rect(s, x, y, pw, ph, fill=None, line=True,
                  line_color=GREY, line_width_pt=0.4)
        add_textbox(s, x, y - 0.22, pw, 0.2, title,
                     size=9, bold=True, align="center")
        # scale
        all_vals = list(pre_df[key].values) + list(post_df[key].values)
        ymin, ymax = min(all_vals), max(all_vals)
        yr = ymax - ymin if ymax > ymin else 1.0
        y_pad = yr * 0.15
        ymin -= y_pad; ymax += y_pad
        # draw zero line
        y_zero = y + ph * (1 - (0 - ymin) / (ymax - ymin))
        add_line(s, x, y_zero, x + pw, y_zero, color=GREY, width_pt=0.3, dash=True)
        # 4 x positions: pre-good (x+0.3pw), pre-bad (x+0.5pw), post-good (x+0.7pw), post-bad (x+0.9pw)
        x_positions = {("good","pre"): x + pw * 0.18,
                        ("bad","pre"): x + pw * 0.38,
                        ("good","post"): x + pw * 0.62,
                        ("bad","post"): x + pw * 0.82}
        # subject points
        for subj, bn in cohort:
            for tag, df_ in [("pre", pre_df), ("post", post_df)]:
                xx = x_positions[(bn, tag)] + rng.uniform(-0.05, 0.05)
                yy = y + ph * (1 - (df_.loc[subj, key] - ymin) / (ymax - ymin))
                add_oval(s, xx, yy, d_in=0.10, fill=(GOOD if bn=="good" else BAD))
        # group mean bars
        for tag, df_ in [("pre", pre_df), ("post", post_df)]:
            for bn in ["good", "bad"]:
                mx = x_positions[(bn, tag)]
                mean_val = df_.loc[[s for s,b in cohort if b==bn], key].mean()
                my = y + ph * (1 - (mean_val - ymin) / (ymax - ymin))
                add_line(s, mx - 0.11, my, mx + 0.11, my,
                          color=(GOOD if bn=="good" else BAD), width_pt=1.8)
        # divider between pre and post
        add_line(s, x + pw * 0.5, y, x + pw * 0.5, y + ph,
                  color=GREY, width_pt=0.4, dash=True)
        # P-value annotations
        p_pre = pre_mw.loc[key, "MW_P_1s_good_gt_bad"]
        p_post = post_mw.loc[key, "MW_P_1s_good_gt_bad"]
        add_textbox(s, x + pw * 0.05, y + 0.05, pw * 0.45, 0.18,
                     f"pre P = {p_pre:.2f}",
                     size=7, bold=(p_pre <= 0.05),
                     color=(GOLD if p_pre <= 0.05 else GREY), align="center")
        add_textbox(s, x + pw * 0.55, y + 0.05, pw * 0.40, 0.18,
                     f"post P = {p_post:.2f}",
                     size=7, bold=(p_post <= 0.05),
                     color=(GOLD if p_post <= 0.05 else GREY), align="center")
        # x labels
        add_textbox(s, x + pw * 0.14, y + ph + 0.02, pw * 0.25, 0.2,
                     "g / b  pre", size=7, color=GREY, align="center")
        add_textbox(s, x + pw * 0.6, y + ph + 0.02, pw * 0.25, 0.2,
                     "g / b  post", size=7, color=GREY, align="center")
        # y axis minmax
        add_textbox(s, x - 0.35, y - 0.05, 0.3, 0.18, f"{ymax:+.1f}",
                     size=6, color=GREY, align="right")
        add_textbox(s, x - 0.35, y + ph - 0.12, 0.3, 0.18, f"{ymin:+.1f}",
                     size=6, color=GREY, align="right")


def build_bar_horizontal_panel(prs, df, *, title, subtitle, label_col,
                                 value_col, font_mono=False):
    """Panel D/E: horizontal diverging bar, labeled."""
    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.15, 9.0, 0.3, title, size=13, bold=True)
    add_textbox(s, 0.5, 0.50, 9.0, 0.25, subtitle, size=8, color=GREY, italic=True)
    n = len(df)
    vals = df[value_col].values
    labels = df[label_col].values
    vmax = max(abs(vals.min()), abs(vals.max()))
    if vmax == 0: vmax = 1
    bx0 = 4.5; by0 = 1.0
    bh = min(0.25, 4.5 / n)
    bar_len_px = 3.5
    zero_x = bx0
    add_line(s, zero_x, by0 - 0.08, zero_x, by0 + n * bh + 0.08,
              color=LINE, width_pt=0.6)
    fam = "Courier New" if font_mono else FONT
    for i in range(n):
        v = vals[i]
        lab = str(labels[i])
        add_textbox(s, bx0 - 3.9, by0 + i * bh, 3.7, bh,
                     lab, size=7, align="right", anchor="middle", family=fam)
        bar_w = abs(v) / vmax * bar_len_px
        color = GOOD if v > 0 else BAD
        bx_start = zero_x if v > 0 else zero_x - bar_w
        add_rect(s, bx_start, by0 + i * bh + 0.04, bar_w, bh - 0.08,
                  fill=color, line=True, line_color=LINE, line_width_pt=0.25)
        tx = bx_start + bar_w + 0.05 if v > 0 else bx_start - 0.55
        add_textbox(s, tx, by0 + i * bh, 0.55, bh,
                     f"{v:+.2f}", size=6.5, color=GREY,
                     align=("left" if v > 0 else "right"), anchor="middle")
    add_textbox(s, bx0 - 2.0, by0 + n * bh + 0.2, 4.0, 0.25,
                 f"IAE − IBI (mean Δ z)", size=9, align="center")
    # legend
    add_legend_chip(s, 1.0, 5.5, "IAE > IBI (good inflamed)", GOOD)
    add_legend_chip(s, 3.6, 5.5, "IBI > IAE (bad inflamed)", BAD)


def build_scatter_panel(prs, xs, ys, cohort, *, title, subtitle, xlab, ylab):
    """Panel F: scatter with regression line, subject labels."""
    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.15, 9.0, 0.3, title, size=13, bold=True)
    add_textbox(s, 0.5, 0.50, 9.0, 0.25, subtitle, size=8, color=GREY, italic=True)
    # plot box
    px, py, pw, ph = 2.5, 1.1, 5.0, 4.2
    add_rect(s, px, py, pw, ph, fill=None, line=True,
              line_color=LINE, line_width_pt=0.5)
    xmin, xmax = xs.min() - 0.3, xs.max() + 0.3
    ymin, ymax = ys.min() - 0.3, ys.max() + 0.3
    # zero axes
    zx = px + pw * (0 - xmin) / (xmax - xmin)
    zy = py + ph * (1 - (0 - ymin) / (ymax - ymin))
    add_line(s, zx, py, zx, py + ph, color=GREY, width_pt=0.3, dash=True)
    add_line(s, px, zy, px + pw, zy, color=GREY, width_pt=0.3, dash=True)
    # regression line
    m, c = np.polyfit(xs, ys, 1)
    x_line = np.array([xmin, xmax])
    y_line = m * x_line + c
    x_line_clip = np.clip(x_line, xmin, xmax)
    y_line_clip = np.clip(y_line, ymin, ymax)
    rx1 = px + pw * (x_line_clip[0] - xmin) / (xmax - xmin)
    ry1 = py + ph * (1 - (y_line_clip[0] - ymin) / (ymax - ymin))
    rx2 = px + pw * (x_line_clip[1] - xmin) / (xmax - xmin)
    ry2 = py + ph * (1 - (y_line_clip[1] - ymin) / (ymax - ymin))
    add_line(s, rx1, ry1, rx2, ry2, color=GREY, width_pt=0.6, dash=True)
    # points
    for (subj, bn), xv, yv in zip(cohort, xs, ys):
        px_ = px + pw * (xv - xmin) / (xmax - xmin)
        py_ = py + ph * (1 - (yv - ymin) / (ymax - ymin))
        add_oval(s, px_, py_, d_in=0.18, fill=(GOOD if bn=="good" else BAD))
        add_textbox(s, px_ + 0.1, py_ - 0.1, 0.5, 0.18,
                     f"s{subj}", size=7, color=INK)
    # axis titles
    add_textbox(s, px, py + ph + 0.35, pw, 0.25, xlab, size=9, align="center")
    tb = add_textbox(s, px - 0.95, py + ph / 2 - 0.3, 0.8, 0.6, ylab,
                      size=9, align="center")
    tb.rotation = -90
    # min/max on axes
    add_textbox(s, px - 0.3, py + ph - 0.12, 0.3, 0.2, f"{ymin:+.1f}",
                 size=7, color=GREY, align="right")
    add_textbox(s, px - 0.3, py - 0.05, 0.3, 0.2, f"{ymax:+.1f}",
                 size=7, color=GREY, align="right")
    add_textbox(s, px - 0.2, py + ph + 0.05, 0.5, 0.2, f"{xmin:+.1f}",
                 size=7, color=GREY, align="center")
    add_textbox(s, px + pw - 0.3, py + ph + 0.05, 0.5, 0.2, f"{xmax:+.1f}",
                 size=7, color=GREY, align="center")
    # legend
    add_legend_chip(s, 7.8, 1.3, "good (pCR)", GOOD)
    add_legend_chip(s, 7.8, 1.5, "bad (poor)", BAD)


# =====================  Fig S23 — pre/post/Δ heatmap (full-page)  =====================
def build_figS23():
    prs = new_prs()
    pre_mw = pd.read_csv(TABLES / "v2_pre_MW.tsv", sep="\t").set_index("composite")
    post_mw = pd.read_csv(TABLES / "v2_post_MW.tsv", sep="\t").set_index("composite")
    dlt_mw = pd.read_csv(TABLES / "v2_delta_MW.tsv", sep="\t").set_index("composite")
    order = pre_mw.sort_values("MW_P_1s_good_gt_bad").index.tolist()
    pre_delta = (pre_mw.loc[order, "good_mean"] - pre_mw.loc[order, "bad_mean"]).values
    post_delta = (post_mw.loc[order, "good_mean"] - post_mw.loc[order, "bad_mean"]).values
    dlt_delta = (dlt_mw.loc[order, "good_mean"] - dlt_mw.loc[order, "bad_mean"]).values
    M = np.vstack([pre_delta, post_delta, dlt_delta]).T
    P = np.vstack([pre_mw.loc[order, "MW_P_1s_good_gt_bad"].values,
                    post_mw.loc[order, "MW_P_1s_good_gt_bad"].values,
                    dlt_mw.loc[order, "MW_P_1s_good_gt_bad"].values]).T
    vmax = float(np.abs(M).max())
    vmin = -vmax

    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.15, 9.0, 0.3,
                 "Supp Fig S23. NanoString pre / post / Δ composite heatmap (full-page)",
                 size=13, bold=True)
    add_textbox(s, 0.5, 0.45, 9.0, 0.25,
                 "All 23 composites (rows) × 3 timepoints (cols). Star ★ = MW 1s P ≤ 0.05 (ceiling). "
                 "Dot · = P ≤ 0.10. Row order: pre-treatment MW P ascending.",
                 size=8, color=GREY, italic=True)
    n_rows = len(order)
    hx0 = 4.0; hy0 = 0.95
    cell_w = 1.0
    cell_h = min(0.21, 4.8 / n_rows)
    for j, cl in enumerate(["pre", "post", "Δ (post − pre)"]):
        add_textbox(s, hx0 + j * cell_w, hy0 - 0.32, cell_w, 0.25,
                     cl, size=10, bold=True, align="center")
    for i, sig in enumerate(order):
        add_textbox(s, hx0 - 3.3, hy0 + i * cell_h - 0.02, 3.2, cell_h + 0.02,
                     sig, size=7.5, align="right", anchor="middle")
        for j in range(3):
            v = M[i, j]
            fill = interp_color(v, vmin, vmax)
            add_rect(s, hx0 + j * cell_w, hy0 + i * cell_h,
                      cell_w, cell_h, fill=fill, line=True, line_color=GREY,
                      line_width_pt=0.15)
            mark = " ★" if P[i, j] <= 0.05 else (" ·" if P[i, j] <= 0.10 else "")
            val = f"{v:+.2f}{mark}"
            txt_col = WHITE if abs(v) > vmax * 0.6 else INK
            add_textbox(s, hx0 + j * cell_w, hy0 + i * cell_h, cell_w, cell_h,
                         val, size=6.5, color=txt_col, align="center", anchor="middle")
    # colorbar
    cbx = hx0 + 3 * cell_w + 0.4; cby = hy0
    cbh = n_rows * cell_h; cbw = 0.3
    n_steps = 60
    for k in range(n_steps):
        v = vmin + (vmax - vmin) * k / (n_steps - 1)
        add_rect(s, cbx, cby + cbh - (k + 1) * cbh / n_steps,
                  cbw, cbh / n_steps, fill=interp_color(v, vmin, vmax), line=False)
    add_rect(s, cbx, cby, cbw, cbh, fill=None, line=True,
              line_color=LINE, line_width_pt=0.4)
    add_textbox(s, cbx + cbw + 0.05, cby - 0.05, 0.8, 0.2, f"+{vmax:.2f}",
                 size=7, color=GREY)
    add_textbox(s, cbx + cbw + 0.05, cby + cbh - 0.12, 0.8, 0.2, f"{-vmax:.2f}",
                 size=7, color=GREY)
    add_textbox(s, cbx + cbw + 0.05, cby + cbh / 2 - 0.08, 1.0, 0.2,
                 "good − bad (z)", size=7, color=GREY)
    fpath = OUT / "FigS23_NanoString_prepostdelta_heatmap_native_editable.pptx"
    prs.save(str(fpath))
    print(f"wrote {fpath}")


# =====================  Fig S24 — canonical signatures (6 panels)  =====================
def build_figS24():
    prs = new_prs()
    hdr = blank_slide(prs)
    add_textbox(hdr, 0.5, 1.5, 9.0, 1.0,
                 "Supp Fig S24. NanoString regulatory-grade signatures — pre vs post",
                 size=16, bold=True, align="center")
    add_textbox(hdr, 0.5, 2.6, 9.0, 0.6,
                 "Ayers TIS, IFN-γ 6-gene, IFN-γ 10-gene, CD8 cytotoxic, IMPRES, M1 macrophage.\n"
                 "Six panels (one slide each), pre-RT and post-RT side-by-side per signature. Group means as solid bars.",
                 size=10, color=GREY, align="center")

    pre_df = pd.read_csv(TABLES / "v2_composite_pre.tsv", sep="\t", index_col="subject")
    post_df = pd.read_csv(TABLES / "v2_composite_post.tsv", sep="\t", index_col="subject")
    pre_mw = pd.read_csv(TABLES / "v2_pre_MW.tsv", sep="\t").set_index("composite")
    post_mw = pd.read_csv(TABLES / "v2_post_MW.tsv", sep="\t").set_index("composite")
    cohort = [(2,"good"),(4,"good"),(14,"good"),(10,"bad"),(11,"bad"),(13,"bad")]
    targets = [("Ayers_TIS", "Ayers TIS (Keytruda companion; 16 of 18 genes)", "Ayers et al JCI 2017 [63]"),
               ("IFNg_6", "IFN-γ 6-gene", "Ayers et al JCI 2017 [64]"),
               ("IFNg_10_Ayers", "IFN-γ 10-gene (expanded)", "Ayers et al JCI 2017 [64]"),
               ("CD8_cytotoxic", "CD8 cytotoxic (GZMA/B/H/K, PRF1, IFNG, CD8A/B, GNLY)",
                "in-house pure effector panel, matches external LC-CRT meta Z=+2.74"),
               ("IMPRES_pos", "IMPRES (14 of 15 positive features)", "Auslander et al Nat Med 2018 [65]"),
               ("M1_macro", "M1 macrophage (CXCL9/10/11, IL12B, CD80/86, IFNG, TNF, IL1B)",
                "in-house panel; classically activated polarization")]
    rng = np.random.default_rng(7)
    for key, title, prov in targets:
        s = blank_slide(prs)
        add_textbox(s, 0.5, 0.15, 9.0, 0.3, title, size=12, bold=True)
        add_textbox(s, 0.5, 0.42, 9.0, 0.22, prov, size=8, color=GREY, italic=True)
        # plot box
        px, py, pw, ph = 1.8, 1.0, 6.5, 4.3
        add_rect(s, px, py, pw, ph, fill=None, line=True,
                  line_color=LINE, line_width_pt=0.5)
        all_vals = list(pre_df[key].values) + list(post_df[key].values)
        ymin, ymax = min(all_vals), max(all_vals)
        yr = ymax - ymin; y_pad = yr * 0.2
        ymin -= y_pad; ymax += y_pad
        # zero line
        y_zero = py + ph * (1 - (0 - ymin) / (ymax - ymin))
        add_line(s, px, y_zero, px + pw, y_zero, color=GREY, width_pt=0.3, dash=True)
        # two columns: pre (left half) and post (right half)
        # within each column: good (x_off 0.25), bad (x_off 0.75)
        x_positions = {("good","pre"): px + pw * 0.15, ("bad","pre"): px + pw * 0.35,
                        ("good","post"): px + pw * 0.65, ("bad","post"): px + pw * 0.85}
        for subj, bn in cohort:
            for tag, df_ in [("pre", pre_df), ("post", post_df)]:
                xx = x_positions[(bn, tag)] + rng.uniform(-0.08, 0.08)
                yy = py + ph * (1 - (df_.loc[subj, key] - ymin) / (ymax - ymin))
                add_oval(s, xx, yy, d_in=0.16, fill=(GOOD if bn=="good" else BAD))
                add_textbox(s, xx + 0.08, yy - 0.08, 0.4, 0.15, f"s{subj}",
                             size=6, color=GREY)
        # group-mean bars
        for tag, df_ in [("pre", pre_df), ("post", post_df)]:
            for bn in ["good", "bad"]:
                mx = x_positions[(bn, tag)]
                mean_val = df_.loc[[s for s,b in cohort if b==bn], key].mean()
                my = py + ph * (1 - (mean_val - ymin) / (ymax - ymin))
                add_line(s, mx - 0.18, my, mx + 0.18, my,
                          color=(GOOD if bn=="good" else BAD), width_pt=2.5)
        # pre/post divider
        add_line(s, px + pw * 0.5, py, px + pw * 0.5, py + ph,
                  color=LINE, width_pt=0.5, dash=True)
        # P annotations
        p_pre = pre_mw.loc[key, "MW_P_1s_good_gt_bad"]
        p_post = post_mw.loc[key, "MW_P_1s_good_gt_bad"]
        star_pre = "★ " if p_pre <= 0.05 else ""
        star_post = "★ " if p_post <= 0.05 else ""
        add_textbox(s, px + pw * 0.05, py + 0.05, pw * 0.40, 0.25,
                     f"{star_pre}pre MW P = {p_pre:.3f}", size=10,
                     bold=(p_pre <= 0.05), color=(GOLD if p_pre <= 0.05 else GREY),
                     align="center")
        add_textbox(s, px + pw * 0.55, py + 0.05, pw * 0.40, 0.25,
                     f"{star_post}post MW P = {p_post:.3f}", size=10,
                     bold=(p_post <= 0.05), color=(GOLD if p_post <= 0.05 else GREY),
                     align="center")
        # x tick labels
        add_textbox(s, px + pw * 0.05, py + ph + 0.06, pw * 0.4, 0.25,
                     "good     ·     bad    (pre)", size=8, color=GREY, align="center")
        add_textbox(s, px + pw * 0.55, py + ph + 0.06, pw * 0.4, 0.25,
                     "good     ·     bad    (post)", size=8, color=GREY, align="center")
        # y axis labels
        add_textbox(s, px - 0.55, py - 0.08, 0.45, 0.25, f"{ymax:+.2f}",
                     size=7, color=GREY, align="right")
        add_textbox(s, px - 0.55, py + ph - 0.13, 0.45, 0.25, f"{ymin:+.2f}",
                     size=7, color=GREY, align="right")
        add_textbox(s, px - 0.55, py + ph / 2 - 0.12, 0.45, 0.25, "0",
                     size=7, color=GREY, align="right")
        tb = add_textbox(s, px - 1.35, py + ph / 2 - 0.3, 0.8, 0.6,
                          "composite z", size=10, align="center")
        tb.rotation = -90
        # legend chips
        add_legend_chip(s, 8.6, 1.1, "good (pCR)", GOOD)
        add_legend_chip(s, 8.6, 1.3, "bad (poor)", BAD)
    fpath = OUT / "FigS24_NanoString_canonical_signatures_native_editable.pptx"
    prs.save(str(fpath))
    print(f"wrote {fpath}")


# =====================  Fig S25 — IAE vs IBI (composite + genes)  =====================
def build_figS25():
    prs = new_prs()
    hdr = blank_slide(prs)
    add_textbox(hdr, 0.5, 1.3, 9.0, 1.5,
                 "Supp Fig S25. Inflamed-but-Ineffective (IBI, n=3 bad)\n"
                 "vs Inflamed-Active-Effective (IAE, n=2 good)",
                 size=15, bold=True, align="center")
    add_textbox(hdr, 0.5, 2.9, 9.0, 1.0,
                 "IAE = subjects 2, 4 (inflamed good pCR). IBI = subjects 10, 11, 13 (inflamed bad poor).\n"
                 "Inflamed = ΔAyers TIS > 0 (regulatory-grade post-RT). IAE productive cascade "
                 "(TCR → JAK-STAT → NF-κB → BATF → TAPBP)\nvs IBI naive-B / NK-inhibitory / ARG2-suppression "
                 "without productive signalling.",
                 size=9, color=GREY, align="center", italic=True)
    iae = pd.read_csv(TABLES / "v2_IAE_vs_IBI_descriptive.tsv", sep="\t")
    build_bar_horizontal_panel(prs, iae.head(16),
                                title="A  Composite-level IAE vs IBI fingerprint",
                                subtitle="Top 16 composites / ratios by |IAE − IBI|. "
                                         "Positive (teal) = IAE > IBI; negative (coral) = IBI > IAE.",
                                label_col="feature", value_col="IAE_minus_IBI")
    genes = pd.read_csv(TABLES / "v2_IAE_vs_IBI_gene_descriptive.tsv", sep="\t").head(24)
    build_bar_horizontal_panel(prs, genes,
                                title="B  Gene-level top 24 discriminators",
                                subtitle="IAE productive: PLA2G6, JAK3, IKBKB, TAPBP, ZAP70, IFIH1, "
                                         "STAT6, PYCARD, BATF, LAG3, TLR8. IBI unproductive: "
                                         "KIR_Inh_Subgroup_2, ARG2, RAG1, MS4A1, XCL2.",
                                label_col="gene", value_col="IAE_minus_IBI", font_mono=True)
    fpath = OUT / "FigS25_IBI_vs_IAE_fingerprint_native_editable.pptx"
    prs.save(str(fpath))
    print(f"wrote {fpath}")


# =====================  Fig S26 — subject radar (linear alt.)  =====================
def build_figS26():
    prs = new_prs()
    hdr = blank_slide(prs)
    add_textbox(hdr, 0.5, 1.3, 9.0, 1.5,
                 "Supp Fig S26. Subject-level deep-dive — s4, s2, s11",
                 size=15, bold=True, align="center")
    add_textbox(hdr, 0.5, 2.9, 9.0, 1.2,
                 "For each subject, 17 composite axes shown as horizontal bars (z-score units); "
                 "pre-RT as hollow bar, post-RT as filled bar.\n"
                 "s2 = textbook good (inflamed-expansion route). "
                 "s4 = atypical good (antigen-presentation-centric route, B-cell ↓).\n"
                 "s11 = IBI archetype (bad, cold baseline + moderate post Δ but "
                 "CD8 effector stays flat).",
                 size=9, color=GREY, align="center", italic=True)

    pre_df = pd.read_csv(TABLES / "v2_composite_pre.tsv", sep="\t", index_col="subject")
    post_df = pd.read_csv(TABLES / "v2_composite_post.tsv", sep="\t", index_col="subject")
    axes_cats = ["TLS_8","Plasma_proxy","GC_TF","Naive_B","Memory_B","Ayers_TIS",
                  "IFNg_6","CD8_cytotoxic","Teff_cytotoxic","Treg","CD8_exh",
                  "HLA_II","HLA_I_machinery_narrow","M1_macro","M2_macro",
                  "NK_activating","DC_mature"]
    subjects = [(4, "s4  good (atypical pCR)"),
                 (2, "s2  good (textbook pCR, inflamed-expansion)"),
                 (11, "s11  bad (poor, IBI archetype)")]
    for subj, title in subjects:
        s = blank_slide(prs)
        add_textbox(s, 0.5, 0.15, 9.0, 0.3, title, size=12, bold=True)
        add_textbox(s, 0.5, 0.45, 9.0, 0.25,
                     "Pre-RT: hollow teal/coral bar.  Post-RT: filled.  "
                     "z = NanoString composite z-score across 12-sample matrix.",
                     size=8, color=GREY, italic=True)
        bn = "good" if subj in (2, 4, 14) else "bad"
        base_col = GOOD if bn == "good" else BAD
        px, py, pw, ph = 4.5, 0.85, 4.5, 4.8
        add_rect(s, px, py, pw, ph, fill=None, line=True,
                  line_color=LINE, line_width_pt=0.4)
        # axis range
        all_vals = list(pre_df.loc[subj, axes_cats].values) + \
                     list(post_df.loc[subj, axes_cats].values)
        vmin, vmax = min(all_vals), max(all_vals)
        vpad = max(abs(vmin), abs(vmax)) * 0.15
        vmin -= vpad; vmax += vpad
        # zero line
        zero_x = px + pw * (0 - vmin) / (vmax - vmin)
        add_line(s, zero_x, py, zero_x, py + ph, color=GREY, width_pt=0.5, dash=True)
        bh = ph / len(axes_cats)
        for i, cat in enumerate(axes_cats):
            pre_v = pre_df.loc[subj, cat]
            post_v = post_df.loc[subj, cat]
            y_row = py + i * bh
            add_textbox(s, px - 2.3, y_row + bh * 0.1, 2.2, bh - 0.03,
                         cat, size=7, align="right", anchor="middle")
            # pre bar (hollow)
            px_w = abs(pre_v) / (vmax - vmin) * pw
            bx = zero_x - px_w if pre_v < 0 else zero_x
            add_rect(s, bx, y_row + bh * 0.1, px_w, bh * 0.35,
                      fill=None, line=True, line_color=base_col, line_width_pt=0.8)
            # post bar (filled)
            po_w = abs(post_v) / (vmax - vmin) * pw
            bx2 = zero_x - po_w if post_v < 0 else zero_x
            add_rect(s, bx2, y_row + bh * 0.5, po_w, bh * 0.35,
                      fill=base_col, line=True, line_color=LINE, line_width_pt=0.3)
            # Δ label on right
            delta = post_v - pre_v
            tx = px + pw + 0.05
            add_textbox(s, tx, y_row + bh * 0.1, 0.6, bh - 0.03,
                         f"Δ{delta:+.2f}",
                         size=6.5, color=(GOLD if abs(delta) >= 1 else GREY),
                         bold=(abs(delta) >= 1),
                         anchor="middle")
        add_textbox(s, px - 1.0, py + ph + 0.2, pw + 2.0, 0.2,
                     "Pre (hollow)           Δ label           Post (filled)",
                     size=8, color=GREY, align="center", italic=True)
        add_textbox(s, px - 0.4, py + ph + 0.02, 0.5, 0.2, f"{vmin:+.1f}",
                     size=7, color=GREY, align="right")
        add_textbox(s, px + pw - 0.3, py + ph + 0.02, 0.5, 0.2, f"{vmax:+.1f}",
                     size=7, color=GREY, align="left")
        add_textbox(s, zero_x - 0.1, py + ph + 0.02, 0.2, 0.2, "0",
                     size=7, color=GREY, align="center")
    fpath = OUT / "FigS26_subject_radar_native_editable.pptx"
    prs.save(str(fpath))
    print(f"wrote {fpath}")


# =====================  Fig S27 — pre-spec primary paired (4 panels)  =====================
def build_figS27():
    prs = new_prs()
    hdr = blank_slide(prs)
    add_textbox(hdr, 0.5, 1.3, 9.0, 1.5,
                 "Supp Fig S27. Pre-registered Arrow 5 rescue primary paired Δ",
                 size=15, bold=True, align="center")
    add_textbox(hdr, 0.5, 2.9, 9.0, 1.5,
                 "P1 CXCL13, P2 TLS-8, P3 Plasma-proxy, P4 GC-TF. "
                 "One-sided good > bad MW with 3-vs-3 ceiling P = 0.050.\n"
                 "All four primaries show bad ≥ good direction (ceiling not reached).\n"
                 "Pre-spec decision rule: NULL + direction flip — Fig 8F arrow 5 remains dashed / qualified.",
                 size=9, color=GREY, align="center", italic=True)

    zmat = pd.read_csv(TABLES / "logz_matrix.tsv", sep="\t", index_col=0)
    comp = pd.read_csv(TABLES / "composite_scores.tsv", sep="\t")
    pairs = [(2,"TNT RNA 5","TNT RNA 6","good"),(4,"TNT RNA 11","TNT RNA 12","good"),
             (14,"TNT RNA 41","TNT RNA 42","good"),(10,"TNT RNA 29","TNT RNA 30","bad"),
             (11,"TNT RNA 32","TNT RNA 33","bad"),(13,"TNT RNA 38","TNT RNA 39","bad")]
    pri = pd.read_csv(TABLES / "P1_P4_primary.tsv", sep="\t").set_index("target")
    targets = [("P1_CXCL13", "CXCL13"),
               ("P2_TLS_8", "TLS_8"),
               ("P3_Plasma_proxy", "Plasma_proxy"),
               ("P4_GC_TF", "GC_TF")]
    for tgt_name, key in targets:
        s = blank_slide(prs)
        pri_row = pri.loc[tgt_name]
        if key == "CXCL13":
            pre_vals = {subj: zmat.loc["CXCL13", p] for (subj,p,q,_) in pairs}
            post_vals = {subj: zmat.loc["CXCL13", q] for (subj,p,q,_) in pairs}
            title = f"{tgt_name}   (single-gene CXCL13)"
        else:
            sub = comp[comp["composite"] == key]
            pre_vals = {subj: float(sub[sub["sample"] == p]["score"].iloc[0]) for (subj,p,q,_) in pairs}
            post_vals = {subj: float(sub[sub["sample"] == q]["score"].iloc[0]) for (subj,p,q,_) in pairs}
            title = f"{tgt_name}  ({key} composite)"

        add_textbox(s, 0.5, 0.15, 9.0, 0.3, title, size=13, bold=True)
        add_textbox(s, 0.5, 0.45, 9.0, 0.25,
                     f"Good mean Δ = {pri_row['good_mean']:+.2f}  ·  "
                     f"Bad mean Δ = {pri_row['bad_mean']:+.2f}  ·  "
                     f"MW one-sided P (good > bad) = {pri_row['MW_P_1s_good_gt_bad']:.3f}",
                     size=10, color=GREY, italic=True)
        # plot box
        px, py, pw, ph = 2.2, 1.0, 5.8, 4.4
        add_rect(s, px, py, pw, ph, fill=None, line=True,
                  line_color=LINE, line_width_pt=0.5)
        allv = list(pre_vals.values()) + list(post_vals.values())
        ymin, ymax = min(allv), max(allv)
        yr = ymax - ymin; pad = yr * 0.2
        ymin -= pad; ymax += pad
        y_zero = py + ph * (1 - (0 - ymin) / (ymax - ymin))
        add_line(s, px, y_zero, px + pw, y_zero, color=GREY, width_pt=0.3, dash=True)
        # x-axis: pre at 0.25, post at 0.75 (relative)
        x_pre = px + pw * 0.25
        x_post = px + pw * 0.75
        add_textbox(s, x_pre - 0.3, py + ph + 0.08, 0.6, 0.25,
                     "pre-RT", size=10, align="center")
        add_textbox(s, x_post - 0.3, py + ph + 0.08, 0.6, 0.25,
                     "post-RT", size=10, align="center")
        for subj, _p, _q, bn in pairs:
            y1 = py + ph * (1 - (pre_vals[subj] - ymin) / (ymax - ymin))
            y2 = py + ph * (1 - (post_vals[subj] - ymin) / (ymax - ymin))
            col = GOOD if bn == "good" else BAD
            add_line(s, x_pre, y1, x_post, y2, color=col, width_pt=1.2)
            add_oval(s, x_pre, y1, d_in=0.18, fill=col)
            add_oval(s, x_post, y2, d_in=0.18, fill=col)
            add_textbox(s, x_post + 0.12, y2 - 0.10, 0.5, 0.2, f"s{subj}",
                         size=7, color=INK)
        # group-mean slope
        for bn, col in [("good", GOOD), ("bad", BAD)]:
            pre_m = np.mean([pre_vals[s] for (s,_p,_q,b) in pairs if b == bn])
            post_m = np.mean([post_vals[s] for (s,_p,_q,b) in pairs if b == bn])
            y1 = py + ph * (1 - (pre_m - ymin) / (ymax - ymin))
            y2 = py + ph * (1 - (post_m - ymin) / (ymax - ymin))
            add_line(s, x_pre - 0.2, y1, x_pre + 0.2, y1, color=col, width_pt=3)
            add_line(s, x_post - 0.2, y2, x_post + 0.2, y2, color=col, width_pt=3)
        # y axis ticks
        add_textbox(s, px - 0.5, py - 0.08, 0.45, 0.2, f"{ymax:+.2f}",
                     size=7, color=GREY, align="right")
        add_textbox(s, px - 0.5, py + ph - 0.12, 0.45, 0.2, f"{ymin:+.2f}",
                     size=7, color=GREY, align="right")
        add_textbox(s, px - 0.5, y_zero - 0.1, 0.45, 0.2, "0",
                     size=7, color=GREY, align="right")
        tb = add_textbox(s, px - 1.25, py + ph / 2 - 0.3, 0.8, 0.6,
                          "z-score", size=10, align="center")
        tb.rotation = -90
        # legend
        add_legend_chip(s, 8.3, 1.1, "good (pCR)", GOOD)
        add_legend_chip(s, 8.3, 1.3, "bad (poor)", BAD)
    fpath = OUT / "FigS27_NanoString_prespec_primary_paired_native_editable.pptx"
    prs.save(str(fpath))
    print(f"wrote {fpath}")


# =====================  Fig S28 — platform concordance  =====================
def build_figS28():
    prs = new_prs()
    hdr = blank_slide(prs)
    add_textbox(hdr, 0.5, 1.3, 9.0, 1.5,
                 "Supp Fig S28. NanoString vs RNA-seq platform concordance",
                 size=15, bold=True, align="center")
    add_textbox(hdr, 0.5, 2.9, 9.0, 1.5,
                 "Per-gene Pearson correlation of NanoString Δ × RNA-seq Δ across 5 subjects (s2/4/14/10/13; "
                 "subj 11 RNA-seq pre missing).\n"
                 "694 expressed shared genes. Median r = +0.754. 90.1 % of genes r > 0.\n"
                 "Rules out platform artefact as explanation for the direction flip seen in the pre-spec Δ analysis.",
                 size=9, color=GREY, align="center", italic=True)

    import numpy as np
    s4 = pd.read_csv(TABLES / "S4_platform_concordance.tsv", sep="\t")

    # Panel A: histogram
    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.15, 9.0, 0.3,
                 "A  Histogram of NanoString Δ × RNA-seq Δ Pearson r (n = 694 genes)",
                 size=13, bold=True)
    add_textbox(s, 0.5, 0.45, 9.0, 0.25,
                 f"Median r = +{s4.pearson_r.median():.3f}. Gold dashed line = median.",
                 size=8, color=GREY, italic=True)
    px, py, pw, ph = 1.5, 1.0, 7.5, 4.3
    add_rect(s, px, py, pw, ph, fill=None, line=True,
              line_color=LINE, line_width_pt=0.5)
    bins = np.linspace(-1, 1, 41)
    counts, edges = np.histogram(s4["pearson_r"].values, bins=bins)
    bar_w = pw / len(counts)
    max_c = counts.max()
    for i, c in enumerate(counts):
        if c <= 0:
            continue
        h = (c / max_c) * ph
        add_rect(s, px + i * bar_w, py + ph - h, bar_w * 0.92, h,
                  fill=GOOD, line=True, line_color=LINE, line_width_pt=0.2)
    med = s4.pearson_r.median()
    med_x = px + pw * (med - (-1)) / 2
    add_line(s, med_x, py, med_x, py + ph, color=GOLD, width_pt=1.5, dash=True)
    add_textbox(s, med_x - 0.7, py - 0.05, 1.4, 0.2,
                 f"median = +{med:.2f}",
                 size=9, bold=True, color=GOLD, align="center")
    # zero line
    zero_x = px + pw * 0.5
    add_line(s, zero_x, py, zero_x, py + ph, color=GREY, width_pt=0.3, dash=True)
    # x ticks
    for xv, lbl in [(-1, "−1"), (-0.5, "−0.5"), (0, "0"), (0.5, "0.5"), (1, "+1")]:
        tx = px + pw * (xv - (-1)) / 2
        add_line(s, tx, py + ph, tx, py + ph + 0.08, color=LINE, width_pt=0.4)
        add_textbox(s, tx - 0.2, py + ph + 0.1, 0.4, 0.18, lbl,
                     size=8, color=GREY, align="center")
    add_textbox(s, px, py + ph + 0.35, pw, 0.25,
                 "Pearson r (NanoString Δ × RNA-seq Δ, n = 5 subjects)",
                 size=9, align="center")
    tb = add_textbox(s, px - 1.0, py + ph / 2 - 0.2, 0.8, 0.5,
                      "Gene count", size=10, align="center")
    tb.rotation = -90

    # Panel B: top-15 concordant
    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.15, 9.0, 0.3,
                 "B  Top 15 concordant genes by Pearson r",
                 size=13, bold=True)
    add_textbox(s, 0.5, 0.45, 9.0, 0.25,
                 "Gene-level NanoString Δ and RNA-seq Δ agree for top cytokines, "
                 "Ig markers, cell-cycle markers (BIRC5), and chemokines (CXCL5, CXCL3, CCL20).",
                 size=8, color=GREY, italic=True)
    top = s4.sort_values("pearson_r", ascending=False).head(15)
    bx0 = 4.2; by0 = 0.9
    bh = 0.28
    bar_max_px = 4.0
    for i, (_, row) in enumerate(top.iterrows()):
        r = row["pearson_r"]
        add_textbox(s, bx0 - 2.2, by0 + i * bh - 0.02, 2.0, bh + 0.04,
                     row["gene"], size=8, align="right", anchor="middle",
                     family="Courier New")
        bar_w = (r - 0.9) / 0.1 * bar_max_px  # zoom to [0.9, 1.0]
        add_rect(s, bx0, by0 + i * bh + 0.05, bar_w, bh - 0.1,
                  fill=GOOD, line=True, line_color=LINE, line_width_pt=0.25)
        add_textbox(s, bx0 + bar_w + 0.1, by0 + i * bh - 0.02, 0.5, bh + 0.04,
                     f"{r:.3f}", size=7, color=GREY, align="left", anchor="middle")
    add_textbox(s, bx0 + bar_max_px / 2 - 1, by0 + len(top) * bh + 0.1, 2.5, 0.25,
                 "Pearson r (zoomed to 0.9–1.0)", size=9, align="center")
    fpath = OUT / "FigS28_NanoString_platform_concordance_native_editable.pptx"
    prs.save(str(fpath))
    print(f"wrote {fpath}")


# =====================  main  =====================
if __name__ == "__main__":
    build_figS22()
    build_figS23()
    build_figS24()
    build_figS25()
    build_figS26()
    build_figS27()
    build_figS28()
    print("\nAll 7 native-editable PPTs written to:", OUT)
    for p in sorted(OUT.glob("FigS*_native_editable.pptx")):
        print(f"  {p.name}  ({p.stat().st_size // 1024} KB)")
