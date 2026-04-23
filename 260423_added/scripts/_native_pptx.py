"""Native PowerPoint helper — every element is an editable shape/textbox.
All text uses Arial; no shadow effects anywhere.

Usage:
    from _native_pptx import NativePanel, GOOD, BAD, GREY, NEUTRAL
    p = NativePanel(slide_in=(10, 7.5))
    p.axes(origin_in=(1.5, 1.0), size_in=(7, 5), xlim=(0,4), ylim=(0,100),
           xticks=[(0.5,'a'),(1.5,'b')], yticks=[(0,'0'),(50,'50'),(100,'100')],
           xlabel='group', ylabel='count')
    p.boxplot(x=1, values=[...], color=GOOD)
    p.scatter([...], [...], color=GOOD)
    p.hline(y=0, dash=True)
    p.text(x_in=5, y_in=0.5, text='P = 0.43', fontsize=9)
    p.save_slide(prs, layout_idx=6)  # blank

Coordinate system: we map *data* coordinates to *inches* inside the axis box,
then convert to EMU. Caller sets an `axes()` region and uses x/y in data units.
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from dataclasses import dataclass, field
from typing import Iterable, Sequence, Tuple, Optional, List
import numpy as np

# ---- Palette (TNT project standard: GOOD=#0a7d6e deep teal, BAD=#c53e1f deep coral)
GOOD      = RGBColor(0x0A, 0x7D, 0x6E)
BAD       = RGBColor(0xC5, 0x3E, 0x1F)
GOOD_FILL = RGBColor(0xCF, 0xE4, 0xE1)   # light teal (fill alpha equivalent)
BAD_FILL  = RGBColor(0xF4, 0xD8, 0xD2)   # light coral
GREY      = RGBColor(0x6C, 0x75, 0x7D)
GREY_LT   = RGBColor(0xCC, 0xCC, 0xCC)
NEUTRAL   = RGBColor(0x33, 0x33, 0x33)
BLACK     = RGBColor(0x11, 0x11, 0x11)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

FONT = 'Arial'

NSMAP = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}


def kill_shadow(shape):
    """Force-remove any shadow/effect: empty effectLst + zero effectRef idx."""
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
    sp_xml = shape._element
    # 1. Empty <a:effectLst/> inside every spPr (kills inline shadows)
    for spPr in sp_xml.xpath('.//a:spPr'):
        for e in spPr.findall(qn('a:effectLst')):
            spPr.remove(e)
        etree.SubElement(spPr, qn('a:effectLst'))
    # 2. Zero out any effectRef idx in <p:style> (kills theme shadow)
    for eref in sp_xml.xpath('.//a:effectRef'):
        eref.set('idx', '0')


def set_line(shape, color: RGBColor, width_pt: float = 0.75, dash: bool = False):
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width_pt)
    if dash:
        lnPr = line._get_or_add_ln()
        # remove existing prstDash
        for d in lnPr.findall(qn('a:prstDash')):
            lnPr.remove(d)
        dash_el = etree.SubElement(lnPr, qn('a:prstDash'))
        dash_el.set('val', 'dash')


def set_fill(shape, color: Optional[RGBColor]):
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color


@dataclass
class Axes:
    origin_in: Tuple[float, float]    # top-left of plot area (inches from slide top-left)
    size_in: Tuple[float, float]      # (width_in, height_in)
    xlim: Tuple[float, float]
    ylim: Tuple[float, float]

    def x2in(self, x: float) -> float:
        return self.origin_in[0] + (x - self.xlim[0]) / (self.xlim[1] - self.xlim[0]) * self.size_in[0]

    def y2in(self, y: float) -> float:
        # inches from slide top: higher y => smaller inches
        return self.origin_in[1] + self.size_in[1] - (y - self.ylim[0]) / (self.ylim[1] - self.ylim[0]) * self.size_in[1]


class NativePanel:
    """One panel == one slide. Coordinates in inches relative to slide top-left."""

    def __init__(self, prs: Presentation, slide_in: Tuple[float, float] = (10.0, 7.5)):
        self.prs = prs
        prs.slide_width  = Inches(slide_in[0])
        prs.slide_height = Inches(slide_in[1])
        blank_layout = prs.slide_layouts[6]  # Blank
        self.slide = prs.slides.add_slide(blank_layout)
        self.ax: Optional[Axes] = None
        self.slide_in = slide_in

    # ---- low-level primitives ----
    def text(self, x_in: float, y_in: float, text: str, fontsize: float = 9,
             bold: bool = False, italic: bool = False, color: RGBColor = BLACK,
             w_in: float = 2.5, h_in: float = 0.3, align: str = 'left',
             anchor: str = 'top'):
        tb = self.slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        tf.word_wrap = True
        anchor_map = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}
        tf.vertical_anchor = anchor_map.get(anchor, MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]
        align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(fontsize)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        kill_shadow(tb)
        return tb

    def text_rotated(self, x_in: float, y_in: float, text: str, fontsize: float = 9,
                     bold: bool = False, color: RGBColor = BLACK,
                     w_in: float = 2.5, h_in: float = 0.3):
        """Rotated -90° textbox (for y-axis label)."""
        tb = self.slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
        tb.rotation = -90
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = text
        run.font.name = FONT; run.font.size = Pt(fontsize); run.font.bold = bold
        run.font.color.rgb = color
        kill_shadow(tb)
        return tb

    def line_in(self, x0: float, y0: float, x1: float, y1: float,
                color: RGBColor = NEUTRAL, width_pt: float = 0.75, dash: bool = False):
        conn = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                               Inches(x0), Inches(y0), Inches(x1), Inches(y1))
        set_line(conn, color, width_pt, dash=dash)
        kill_shadow(conn)
        return conn

    def rect_in(self, x_in: float, y_in: float, w_in: float, h_in: float,
                fill: Optional[RGBColor] = None, line: Optional[RGBColor] = NEUTRAL,
                line_width_pt: float = 0.75):
        r = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
        set_fill(r, fill)
        if line is None:
            r.line.fill.background()
        else:
            set_line(r, line, line_width_pt)
        r.text_frame.text = ''  # no label
        kill_shadow(r)
        return r

    def ellipse_in(self, cx_in: float, cy_in: float, d_in: float,
                   fill: RGBColor = GOOD, line: Optional[RGBColor] = WHITE,
                   line_width_pt: float = 0.6):
        r = self.slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(cx_in - d_in/2), Inches(cy_in - d_in/2),
                                        Inches(d_in), Inches(d_in))
        set_fill(r, fill)
        if line is None:
            r.line.fill.background()
        else:
            set_line(r, line, line_width_pt)
        r.text_frame.text = ''
        kill_shadow(r)
        return r

    def diamond_in(self, cx_in: float, cy_in: float, d_in: float,
                   fill: RGBColor = NEUTRAL, line: Optional[RGBColor] = WHITE,
                   line_width_pt: float = 0.6):
        r = self.slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                        Inches(cx_in - d_in/2), Inches(cy_in - d_in/2),
                                        Inches(d_in), Inches(d_in))
        set_fill(r, fill)
        if line is None:
            r.line.fill.background()
        else:
            set_line(r, line, line_width_pt)
        r.text_frame.text = ''
        kill_shadow(r)
        return r

    # ---- axes ----
    def axes(self, origin_in: Tuple[float, float], size_in: Tuple[float, float],
             xlim: Tuple[float, float], ylim: Tuple[float, float],
             xticks: Optional[List[Tuple[float, str]]] = None,
             yticks: Optional[List[Tuple[float, str]]] = None,
             xlabel: str = '', ylabel: str = '',
             xlabel_fontsize: float = 9, ylabel_fontsize: float = 9,
             tick_fontsize: float = 8, axis_lw_pt: float = 0.75,
             draw_x_axis: bool = True, draw_y_axis: bool = True,
             y_grid_values: Optional[List[float]] = None):
        self.ax = Axes(origin_in, size_in, xlim, ylim)
        x0, y0 = origin_in
        w, h  = size_in
        # axis spines
        if draw_x_axis:
            self.line_in(x0, y0 + h, x0 + w, y0 + h, BLACK, width_pt=axis_lw_pt)
        if draw_y_axis:
            self.line_in(x0, y0, x0, y0 + h, BLACK, width_pt=axis_lw_pt)
        # y grid lines (light)
        if y_grid_values:
            for yv in y_grid_values:
                y_in = self.ax.y2in(yv)
                self.line_in(x0, y_in, x0 + w, y_in, GREY_LT, width_pt=0.4)
        # x ticks + labels (below axis)
        if xticks:
            for xv, lbl in xticks:
                x_in = self.ax.x2in(xv)
                self.line_in(x_in, y0 + h, x_in, y0 + h + 0.05, BLACK, width_pt=axis_lw_pt)
                if lbl:
                    self.text(x_in - 0.6, y0 + h + 0.07, lbl, fontsize=tick_fontsize,
                              w_in=1.2, h_in=0.25, align='center')
        # y ticks + labels (left of axis) — narrower box so it doesn't overlap rotated label
        if yticks:
            for yv, lbl in yticks:
                y_in = self.ax.y2in(yv)
                self.line_in(x0 - 0.05, y_in, x0, y_in, BLACK, width_pt=axis_lw_pt)
                if lbl:
                    self.text(x0 - 0.58, y_in - 0.09, lbl, fontsize=tick_fontsize,
                              w_in=0.5, h_in=0.22, align='right')
        # x-axis label
        if xlabel:
            self.text(x0, y0 + h + 0.50, xlabel, fontsize=xlabel_fontsize,
                      bold=True, w_in=w, h_in=0.3, align='center')
        # y-axis label (rotated) — moved further left to clear tick labels
        if ylabel:
            self.text_rotated(x0 - 1.55, y0 + h/2 - 0.9, ylabel,
                              fontsize=ylabel_fontsize, bold=True,
                              w_in=1.8, h_in=0.3)

    # ---- data-coord primitives ----
    def dot(self, x: float, y: float, d_in: float = 0.10,
            color: RGBColor = GOOD, edge: Optional[RGBColor] = WHITE):
        return self.ellipse_in(self.ax.x2in(x), self.ax.y2in(y), d_in, fill=color, line=edge)

    def diamond(self, x: float, y: float, d_in: float = 0.14,
                color: RGBColor = NEUTRAL, edge: Optional[RGBColor] = WHITE):
        return self.diamond_in(self.ax.x2in(x), self.ax.y2in(y), d_in, fill=color, line=edge)

    def line(self, x0: float, y0: float, x1: float, y1: float,
             color: RGBColor = NEUTRAL, width_pt: float = 0.8, dash: bool = False):
        return self.line_in(self.ax.x2in(x0), self.ax.y2in(y0),
                            self.ax.x2in(x1), self.ax.y2in(y1),
                            color=color, width_pt=width_pt, dash=dash)

    def hline(self, y: float, color: RGBColor = GREY, width_pt: float = 0.5, dash: bool = True):
        x0_in = self.ax.origin_in[0]
        x1_in = self.ax.origin_in[0] + self.ax.size_in[0]
        y_in = self.ax.y2in(y)
        return self.line_in(x0_in, y_in, x1_in, y_in, color=color, width_pt=width_pt, dash=dash)

    def vline(self, x: float, color: RGBColor = GREY, width_pt: float = 0.5, dash: bool = True):
        y0_in = self.ax.origin_in[1]
        y1_in = self.ax.origin_in[1] + self.ax.size_in[1]
        x_in = self.ax.x2in(x)
        return self.line_in(x_in, y0_in, x_in, y1_in, color=color, width_pt=width_pt, dash=dash)

    def rect_data(self, x0: float, y0: float, x1: float, y1: float,
                  fill: Optional[RGBColor] = None, line: Optional[RGBColor] = NEUTRAL,
                  line_width_pt: float = 0.75):
        x_in = self.ax.x2in(min(x0, x1))
        w_in = abs(self.ax.x2in(x1) - self.ax.x2in(x0))
        y_in = self.ax.y2in(max(y0, y1))
        h_in = abs(self.ax.y2in(y0) - self.ax.y2in(y1))
        return self.rect_in(x_in, y_in, w_in, h_in, fill=fill, line=line, line_width_pt=line_width_pt)

    # ---- composite: boxplot ----
    def boxplot(self, x_center: float, values: Sequence[float], width: float = 0.35,
                fill: RGBColor = GOOD_FILL, edge: RGBColor = GOOD, line_width_pt: float = 0.9,
                median_color: RGBColor = BLACK, show_whiskers: bool = True):
        """Box (Q1-Q3), median line, whiskers (1.5*IQR or min/max within fence)."""
        v = np.asarray(values, dtype=float)
        v = v[~np.isnan(v)]
        if len(v) < 1: return
        q1, med, q3 = np.percentile(v, [25, 50, 75])
        iqr = q3 - q1
        lo_fence = max(v.min(), q1 - 1.5 * iqr)
        hi_fence = min(v.max(), q3 + 1.5 * iqr)
        # box
        self.rect_data(x_center - width/2, q1, x_center + width/2, q3, fill=fill, line=edge, line_width_pt=line_width_pt)
        # median
        self.line(x_center - width/2, med, x_center + width/2, med, color=median_color, width_pt=1.3)
        # whiskers
        if show_whiskers:
            self.line(x_center, q3, x_center, hi_fence, color=edge, width_pt=line_width_pt)
            self.line(x_center, q1, x_center, lo_fence, color=edge, width_pt=line_width_pt)
            self.line(x_center - width/4, hi_fence, x_center + width/4, hi_fence, color=edge, width_pt=line_width_pt)
            self.line(x_center - width/4, lo_fence, x_center + width/4, lo_fence, color=edge, width_pt=line_width_pt)

    def jitter_scatter(self, x_center: float, values: Sequence[float], width: float = 0.18,
                       color: RGBColor = GOOD, d_in: float = 0.09, seed: int = 7):
        rng = np.random.default_rng(seed)
        vals = np.asarray(values, dtype=float)
        vals = vals[~np.isnan(vals)]
        jitter = rng.uniform(-width, width, size=len(vals))
        for xi, v in zip(jitter, vals):
            self.dot(x_center + xi, v, d_in=d_in, color=color, edge=WHITE)

    # ---- save ----
    def save(self, path: str):
        self.prs.save(path)


def half_violin(panel, x_center: float, values, side: str = 'left',
                max_width: float = 0.25, color=None, fill=None,
                edge_width_pt: float = 0.75, n_kde: int = 60):
    """Freeform-based half-violin (fully editable polygon).

    Args:
        panel        : NativePanel with active axes
        x_center     : data-space x coord of the violin axis
        values       : raw data values
        side         : 'left' or 'right' (half offset direction)
        max_width    : max half-width in data x units
        color        : edge color (RGBColor)
        fill         : fill color (RGBColor) or None
    """
    from scipy import stats as _st
    v = np.asarray(values, dtype=float)
    v = v[~np.isnan(v)]
    if len(v) < 2:
        return None
    ymin, ymax = float(v.min()), float(v.max())
    if ymax - ymin < 1e-9: ymax = ymin + 1e-6
    try:
        kde = _st.gaussian_kde(v, bw_method='scott')
    except Exception:
        return None
    y_samples = np.linspace(ymin, ymax, n_kde)
    d = kde(y_samples)
    d = d / d.max() * max_width
    sign = -1 if side == 'left' else +1
    # polygon points in (data_x, data_y) order:
    pts_data = [(x_center, ymin)]
    for ys, ds in zip(y_samples, d):
        pts_data.append((x_center + sign * ds, ys))
    pts_data.append((x_center, ymax))
    # convert to EMU
    ax = panel.ax
    pts_emu = []
    for xd, yd in pts_data:
        x_in = ax.x2in(xd); y_in = ax.y2in(yd)
        pts_emu.append((Emu(int(x_in * 914400)), Emu(int(y_in * 914400))))
    # first point as start, rest as offsets from start... actually
    # pptx freeform expects ABSOLUTE coordinates (not deltas)
    start_x, start_y = pts_emu[0]
    builder = panel.slide.shapes.build_freeform(start_x, start_y, scale=1)
    # add line segments: list of (x, y) absolute
    builder.add_line_segments(pts_emu[1:], close=True)
    shape = builder.convert_to_shape()
    if fill is not None:
        set_fill(shape, fill)
    else:
        shape.fill.background()
    if color is not None:
        set_line(shape, color, edge_width_pt)
    else:
        shape.line.fill.background()
    kill_shadow(shape)
    return shape


def stacked_hbar(panel, x: float, y: float, h_data: float,
                 seg_vals: list, seg_colors: list, y_width: float = 0.8):
    """Horizontal stacked bar in data coords (used for HLA allele freq Fig 8A).

    Args:
        x, y     : bar position (x = 0 baseline usually, y = row index)
        h_data   : (not used - use y_width)
        seg_vals : list of segment lengths (in data x units)
        seg_colors : list of RGBColor matching seg_vals
    """
    cx = 0.0
    for v, col in zip(seg_vals, seg_colors):
        if v <= 0: continue
        panel.rect_data(cx, y - y_width/2, cx + v, y + y_width/2,
                        fill=col, line=col, line_width_pt=0.4)
        cx += v
    return cx


def make_prs_blank():
    prs = Presentation()
    prs.slide_width  = Inches(10.0)
    prs.slide_height = Inches(7.5)
    return prs


def panel_letter(slide, text: str = 'A', x_in: float = 0.35, y_in: float = 0.30):
    """Top-left panel letter (big bold Arial)."""
    tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(0.5), Inches(0.5))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = BLACK
    kill_shadow(tb)
    return tb


def footer_note(slide, text: str, y_in: float = 7.05, fontsize: float = 7.5):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(y_in), Inches(9.0), Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = text
    run.font.name = FONT; run.font.size = Pt(fontsize); run.font.italic = True
    run.font.color.rgb = GREY
    kill_shadow(tb)
    return tb
