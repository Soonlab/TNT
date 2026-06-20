"""
Fig 7C v2 — Convergence test heatmap (9 baseline × 4 cascade Δ)

Differences from v1:
- Cascade Δ list switched to the user-specified v2 set
  (Δ SBS5, Δ MHC-I neoantigen binders, Δ Treg, Δ IGH clonotypes).
- Cell value = **plain Spearman r** (the manuscript-quoted convention; matches
  the headline DSB × CD8-cytotoxic Δ = −0.07 / P = 0.83 exactly).
- No gold "headline" border on any cell — Δ CD8-cytotoxic is no longer one of
  the four cascade columns, so the manuscript headline pair is conveyed via a
  separate annotation textbox below the heatmap.
- Per-cell `n` is annotated subtly (n varies 10–14 across pairs because the v2
  cascade list mixes WES-paired SBS5/neo_binders with RNA-paired Treg/IGH_n).

Project figure rules: see make_fig7c.py docstring (Rule 0–6 carried over).
"""
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path(__file__).resolve().parent.parent
TSV_IN = ROOT / "tables" / "convergence_36pair_used_v2.tsv"
PPTX_OUT = ROOT / "figures" / "Fig7C_convergence_heatmap_v2.pptx"

# --- feature lists (must match compute_convergence_v2.py order) ---
BASELINES = [
    ("DNA Double-Strand Break Repair R-HSA-5693532", "DSB repair (Reactome)"),
    ("DNA Repair R-HSA-73894",                       "DNA repair (Reactome)"),
    ("HDR Thru Homologous Recombination (HRR) R-HSA-5685942",
                                                     "HRR (Reactome)"),
    ("E2F Targets",                                  "E2F targets (Hallmark)"),
    ("G2-M Checkpoint",                              "G2-M checkpoint (Hallmark)"),
    ("Myc Targets V2",                               "Myc targets V2 (Hallmark)"),
    ("MHC_II",                                       "MHC-II (baseline)"),
    ("MSI_pct",                                      "MSI (%)"),
    ("frac_amp",                                     "Genomic amp fraction"),
]
CASCADES = [
    ("SBS5_delta",         "Δ SBS5"),
    ("neo_binders_delta",  "Δ MHC-I neoantigen\nbinders"),
    ("Treg_delta",         "Δ Treg"),
    ("IGH_n_delta",        "Δ IGH clonotypes"),
]

# Palette anchors (ColorBrewer RdBu_r)
BLUE_DEEP = (33, 102, 172)
RED_DEEP  = (178,  24,  43)
WHITE     = (255, 255, 255)
GREY_BORD = RGBColor(200, 200, 200)
BLACK     = RGBColor(0, 0, 0)
GOLD      = RGBColor(218, 165, 32)

VMIN, VMAX = -0.6, 0.6  # v2 range expanded to accommodate observed |r| up to 0.56


# --- helpers ---
def diverging(r: float, vmin: float = VMIN, vmax: float = VMAX) -> RGBColor:
    r = max(min(r, vmax), vmin)
    if r >= 0:
        t = r / vmax
        return RGBColor(
            int(WHITE[0] - t * (WHITE[0] - RED_DEEP[0])),
            int(WHITE[1] - t * (WHITE[1] - RED_DEEP[1])),
            int(WHITE[2] - t * (WHITE[2] - RED_DEEP[2])),
        )
    else:
        t = -r / (-vmin)
        return RGBColor(
            int(WHITE[0] - t * (WHITE[0] - BLUE_DEEP[0])),
            int(WHITE[1] - t * (WHITE[1] - BLUE_DEEP[1])),
            int(WHITE[2] - t * (WHITE[2] - BLUE_DEEP[2])),
        )


def kill_shadow(shape) -> None:
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    etree.SubElement(spPr, qn("a:effectLst"))


def set_text(tb, text, size=8, bold=False, color=BLACK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    kill_shadow(tb)


def fmt_r(r): return f"{r:+.2f}".replace("-", "−")


# --- load data ---
df = pd.read_csv(TSV_IN, sep="\t")
rmap = {(row.baseline_feature, row.cascade_feature): row for row in df.itertuples()}

# --- build slide ---
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

GRID_LEFT = Inches(3.30)
GRID_TOP  = Inches(1.75)
CELL_W    = Inches(0.92)
CELL_H    = Inches(0.50)
N_ROWS    = len(BASELINES)
N_COLS    = len(CASCADES)
GRID_W    = CELL_W * N_COLS
GRID_H    = CELL_H * N_ROWS

# column axis title
col_title = slide.shapes.add_textbox(
    GRID_LEFT, GRID_TOP - Inches(1.15), GRID_W, Inches(0.28),
)
set_text(col_title,
         "Cascade Δ features (radiation-phase dynamics)",
         size=10, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# column labels
for j, (_, lbl) in enumerate(CASCADES):
    tb = slide.shapes.add_textbox(
        GRID_LEFT + CELL_W * j, GRID_TOP - Inches(0.85),
        CELL_W, Inches(0.78),
    )
    set_text(tb, lbl, size=9, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

# row axis title
row_title = slide.shapes.add_textbox(
    GRID_LEFT - Inches(2.20), GRID_TOP - Inches(0.34),
    Inches(2.10), Inches(0.26),
)
set_text(row_title, "Baseline tumor-intrinsic features",
         size=10, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)

# row labels
for i, (_, lbl) in enumerate(BASELINES):
    tb = slide.shapes.add_textbox(
        GRID_LEFT - Inches(2.20), GRID_TOP + CELL_H * i,
        Inches(2.10), CELL_H,
    )
    set_text(tb, lbl, size=9, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# --- cells ---
for i, (b_col, _) in enumerate(BASELINES):
    for j, (c_col, _) in enumerate(CASCADES):
        row = rmap.get((b_col, c_col))
        if row is None:
            continue
        r_val = row.spearman_r       # plain Spearman (matches manuscript convention)
        p_val = row.spearman_p
        q_val = row.BH_q_plain
        n_val = row.n

        x = GRID_LEFT + CELL_W * j
        y = GRID_TOP  + CELL_H * i

        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, CELL_W, CELL_H)
        rect.fill.solid()
        rect.fill.fore_color.rgb = diverging(r_val)
        rect.line.color.rgb = GREY_BORD
        rect.line.width = Pt(0.5)
        kill_shadow(rect)

        suffix = ""
        if q_val < 0.05:
            suffix = "**"
        elif p_val < 0.05:
            suffix = "*"

        text_color = RGBColor(255, 255, 255) if abs(r_val) >= 0.45 else BLACK

        # r value (centred)
        ann = slide.shapes.add_textbox(x, y - Inches(0.04), CELL_W, CELL_H)
        set_text(ann, fmt_r(r_val) + suffix,
                 size=9, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # n subscript (small, bottom-right corner of cell)
        ntb = slide.shapes.add_textbox(
            x + CELL_W - Inches(0.30), y + CELL_H - Inches(0.18),
            Inches(0.28), Inches(0.16),
        )
        set_text(ntb, f"n={n_val}", size=6, color=text_color,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)

# --- colorbar ---
CB_LEFT = GRID_LEFT + GRID_W + Inches(0.55)
CB_TOP  = GRID_TOP
CB_W    = Inches(0.26)
CB_H    = GRID_H

cblab = slide.shapes.add_textbox(
    CB_LEFT - Inches(0.25), CB_TOP - Inches(0.34),
    Inches(1.7), Inches(0.26),
)
set_text(cblab, "Plain Spearman r",
         size=9, bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

N_STOPS = 100
stop_h = CB_H / N_STOPS
for k in range(N_STOPS):
    frac = k / (N_STOPS - 1)
    r_at = VMAX - frac * (VMAX - VMIN)
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        CB_LEFT, CB_TOP + stop_h * k,
        CB_W, stop_h + Emu(900),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = diverging(r_at)
    rect.line.fill.background()
    kill_shadow(rect)

border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, CB_LEFT, CB_TOP, CB_W, CB_H)
border.fill.background()
border.line.color.rgb = RGBColor(64, 64, 64)
border.line.width = Pt(0.75)
kill_shadow(border)

for r_val in [0.60, 0.30, 0.0, -0.30, -0.60]:
    frac = (VMAX - r_val) / (VMAX - VMIN)
    ty = CB_TOP + CB_H * frac
    tick = slide.shapes.add_connector(
        1,
        CB_LEFT + CB_W, ty,
        CB_LEFT + CB_W + Inches(0.08), ty,
    )
    tick.line.color.rgb = RGBColor(64, 64, 64)
    tick.line.width = Pt(0.5)
    kill_shadow(tick)
    txt = "0" if r_val == 0 else fmt_r(r_val)
    tb = slide.shapes.add_textbox(
        CB_LEFT + CB_W + Inches(0.12), ty - Inches(0.10),
        Inches(0.55), Inches(0.20),
    )
    set_text(tb, txt, size=8, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# --- bottom annotations ---
LEG_LEFT = GRID_LEFT - Inches(2.20)
LEG_TOP  = GRID_TOP + GRID_H + Inches(0.30)
LEG_W    = Inches(8.30)

leg1 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP, LEG_W, Inches(0.22))
set_text(leg1,
         "n = 10–14 paired subjects (varies by cascade Δ type) · "
         "plain Spearman ρ (manuscript-quoted convention) · "
         "BH-corrected across 36 pairs",
         size=8, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

leg2 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP + Inches(0.22), LEG_W, Inches(0.22))
set_text(leg2,
         "Cell value = plain Spearman r;  small n=… in lower-right corner of each cell.   "
         "* nominal P < 0.05    ** BH q < 0.05.",
         size=8, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

leg3 = slide.shapes.add_textbox(LEG_LEFT, LEG_TOP + Inches(0.44), LEG_W, Inches(0.22))
set_text(leg3,
         "Convergence result: 0/36 plain P < 0.05    ·    0/36 BH q < 0.05    "
         "(static and dynamic layers observationally independent).",
         size=8, bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# Headline-pair annotation (gold-bordered call-out box, distinct from heatmap)
HP_LEFT = LEG_LEFT
HP_TOP  = LEG_TOP + Inches(0.78)
HP_W    = LEG_W
HP_H    = Inches(0.42)

hp_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, HP_LEFT, HP_TOP, HP_W, HP_H)
hp_bg.fill.solid()
hp_bg.fill.fore_color.rgb = RGBColor(252, 247, 232)  # very pale gold
hp_bg.line.color.rgb = GOLD
hp_bg.line.width = Pt(1.25)
kill_shadow(hp_bg)

hp_tb = slide.shapes.add_textbox(
    HP_LEFT + Inches(0.10), HP_TOP + Inches(0.03),
    HP_W - Inches(0.20), HP_H - Inches(0.06),
)
set_text(hp_tb,
         "Manuscript headline pair (§3.10, plain Spearman, n = 12):  "
         "DSB repair × Δ CD8-cytotoxic,  r = −0.07,  P = 0.83.   "
         "Δ CD8-cytotoxic is not one of the four v2 cascade columns; see Supp Fig S20A "
         "for the full 36-pair forest including CD8-cytotoxic Δ.",
         size=8, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# --- save ---
PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(PPTX_OUT)
print(f"[ok] wrote {PPTX_OUT}")
